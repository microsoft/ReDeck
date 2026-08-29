"""Tests for the refactored codegen modules: code_transforms, code_executor."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.backends.python_pptx import code_transforms, code_executor


# --- code_transforms tests ---

class TestExtractCode:

    def test_extracts_from_code_block(self):
        response = '''Here is the code:
```python
def build_slide(prs, slide, image_dir):
    pass
```
Done.'''
        code = code_transforms.extract_code(response)
        assert code is not None
        assert "def build_slide" in code

    def test_returns_none_for_no_code(self):
        response = "I cannot generate the code."
        assert code_transforms.extract_code(response) is None

    def test_extracts_without_fences(self):
        response = """def build_slide(prs, slide, image_dir):
    pass"""
        code = code_transforms.extract_code(response)
        assert code is not None
        assert "def build_slide" in code


class TestSanitizeCode:

    def test_strips_markdown_bold(self):
        code = 'p.text = "**Bold text** normal"'
        result = code_transforms.sanitize_code(code)
        assert "**" not in result
        assert "Bold text" in result

    def test_blocks_slide_creation(self):
        code = '''def build_slide(prs, slide, image_dir):
    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
    new_slide.shapes.add_textbox(0, 0, 100, 100)
'''
        result = code_transforms.sanitize_code(code)
        assert "slides.add_slide" not in result
        assert "BLOCKED" in result

    def test_strips_redundant_pptx_imports(self):
        code = '''from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
def build_slide(prs, slide, image_dir):
    pass'''
        result = code_transforms.sanitize_code(code)
        assert "import handled by runtime" in result

    def test_preserves_runs0_pattern(self):
        """After cleanup, runs[0] patterns are NOT rewritten — LLM must avoid them."""
        code = 'p.runs[0].font.size = Pt(18)'
        result = code_transforms.sanitize_code(code)
        # No longer rewritten
        assert "runs[0]" in result


class TestInjectImageHeightCap:

    def test_injects_helper(self):
        code = '''def build_slide(prs, slide, image_dir):
    slide.shapes.add_picture("img.png", Inches(1), Inches(1), width=Inches(4), height=Inches(3))
'''
        result = code_transforms.inject_image_height_cap(code)
        assert "_add_pic" in result
        assert "slide.shapes.add_picture" not in result

    def test_handles_code_without_add_picture(self):
        code = '''def build_slide(prs, slide, image_dir):
    tf = slide.shapes.add_textbox(0, 0, 100, 100)
'''
        result = code_transforms.inject_image_height_cap(code)
        assert "_add_pic" in result  # helper is injected regardless
        assert "add_textbox" in result  # other calls not replaced


class TestStripImageCode:

    def test_strips_add_picture(self):
        code = '''x = 1
slide.shapes.add_picture("img.png", 0, 0)
y = 2'''
        result = code_transforms.strip_image_code(code)
        assert "add_picture" not in result
        assert "x = 1" in result
        assert "y = 2" in result

    def test_strips_img_path_assignments(self):
        code = '''img_path = os.path.join(image_dir, "fig.png")
slide.shapes.add_picture(img_path, 0, 0)
text = "hello"'''
        result = code_transforms.strip_image_code(code)
        assert "img_path" not in result
        assert 'text = "hello"' in result


class TestStripComments:

    def test_strips_comments(self):
        code = '''# This is a comment
x = 1  # inline comment
y = 2
# Another comment'''
        result = code_transforms.strip_comments(code)
        assert "# This is a comment" not in result
        assert "x = 1" in result


class TestEstimateTextContent:

    def test_estimates_text(self):
        code = '''p.text = "Hello World"
p2.text = "Another line of text"'''
        count = code_transforms.estimate_text_content(code)
        assert count > 0

    def test_zero_for_no_text(self):
        code = "x = 1\ny = 2"
        count = code_transforms.estimate_text_content(code)
        assert count == 0


class TestExtractCoordinates:

    def test_extracts_inches(self):
        code = 'slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.0), Inches(1.5))'
        coords = code_transforms.extract_coordinates(code)
        assert len(coords) > 0


class TestBuildShapeNameMap:

    def test_builds_map(self):
        code = '''title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.0), Inches(1.0))
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.0), Inches(5.0))'''
        result = code_transforms.build_shape_name_map(code)
        assert "title_box" in result or len(result) > 0


# --- code_executor tests ---

class TestExecuteCode:

    def test_successful_execution(self):
        code = '''def build_slide(prs, slide, image_dir):
    from pptx.util import Inches, Pt
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(5))
    tf.text_frame.paragraphs[0].text = "Hello from codegen"
    tf.text_frame.paragraphs[0].font.size = Pt(18)
'''
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        success, error = code_executor.execute_code(code, prs, slide, "/tmp")
        assert success is True
        assert error == ""
        # Verify the textbox was added
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Hello from codegen" in t for t in texts)

    def test_missing_build_slide(self):
        """Code without build_slide should fail (either missing function or exec error)."""
        code = '''x = 1
y = 2
'''
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        success, error = code_executor.execute_code(code, prs, slide, "/tmp")
        assert success is False
        assert len(error) > 0  # some error reported

    def test_runtime_error_caught(self):
        code = '''def build_slide(prs, slide, image_dir):
    raise ValueError("Something went wrong")
'''
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        success, error = code_executor.execute_code(code, prs, slide, "/tmp")
        assert success is False
        assert "ValueError" in error

    def test_disallowed_import_blocked(self):
        code = '''def build_slide(prs, slide, image_dir):
    import subprocess
    subprocess.run(["ls"])
'''
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        success, error = code_executor.execute_code(code, prs, slide, "/tmp")
        assert success is False
        assert "ImportError" in error or "not allowed" in error
