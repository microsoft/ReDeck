"""Tests for PptxSpatialAnalyzer."""
from __future__ import annotations

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.backends.python_pptx.refiner.font_metrics import FontMetrics
from app.backends.python_pptx.refiner.spatial_analyzer import (
    PptxSpatialAnalyzer,
    SpatialReport,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BLANK_LAYOUT = 6


def _make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _save(prs) -> str:
    f = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(f.name)
    return f.name


def _analyzer():
    return PptxSpatialAnalyzer(FontMetrics())


# ── tests ─────────────────────────────────────────────────────────────────────

def test_analyze_returns_report():
    prs = _make_prs()
    layout = prs.slide_layouts[BLANK_LAYOUT]
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Hello World"
    path = _save(prs)
    try:
        report = _analyzer().analyze(path)
        assert isinstance(report, SpatialReport)
        assert len(report.extractions) == 1
        assert report.slide_width_emu == int(SLIDE_W)
        assert report.slide_height_emu == int(SLIDE_H)
    finally:
        os.unlink(path)


def test_detects_text_overflow():
    prs = _make_prs()
    layout = prs.slide_layouts[BLANK_LAYOUT]
    slide = prs.slides.add_slide(layout)
    # Very narrow box, large text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(0.5), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "This is a very long text that definitely overflows the box width"
    run.font.size = Pt(24)
    path = _save(prs)
    try:
        report = _analyzer().analyze(path)
        assert len(report.overflows) > 0, "Expected at least one overflow"
    finally:
        os.unlink(path)


def test_detects_overlap():
    prs = _make_prs()
    layout = prs.slide_layouts[BLANK_LAYOUT]
    slide = prs.slides.add_slide(layout)
    # Two boxes that overlap significantly
    b1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
    b1.text_frame.text = "Box A"
    b2 = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(3), Inches(2))
    b2.text_frame.text = "Box B"
    path = _save(prs)
    try:
        report = _analyzer().analyze(path)
        assert len(report.overlaps) > 0, "Expected at least one overlap pair"
    finally:
        os.unlink(path)


REAL_PPTX = os.environ.get(
    "REDECK_REAL_PPTX",
    str(Path(__file__).resolve().parent.parent / "final_merged.pptx"),
)


@pytest.mark.skipif(not os.path.exists(REAL_PPTX), reason="real pptx not found")
def test_real_pptx():
    report = _analyzer().analyze(REAL_PPTX)
    assert len(report.extractions) >= 10, "Expected at least 10 slides"
    total_issues = len(report.overflows) + len(report.overlaps)
    assert total_issues > 0, "Expected some issues in real presentation"
