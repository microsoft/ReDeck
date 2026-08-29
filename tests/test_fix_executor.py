"""Tests for PptxFixExecutor — no LLM required."""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.backends.python_pptx.refiner.fix_executor import PptxFixExecutor
from app.backends.python_pptx.refiner.fix_planner import FixOp, FixPlan


# ── Fixture helpers ───────────────────────────────────────────────────────────

def _make_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX with one textbox named 'test_box'."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
    txBox.name = "test_box"
    tf = txBox.text_frame
    tf.text = "Sample text"
    run = tf.paragraphs[0].runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(24)

    out = tmp_path / "test_slide.pptx"
    prs.save(out)
    return out


def _load_shape(pptx_path: Path, shape_name: str = "test_box"):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    for s in slide.shapes:
        if s.name == shape_name:
            return s
    raise KeyError(f"Shape '{shape_name}' not found")


# ── Tests ─────────────────────────────────────────────────────────────────────

def test_shrink_font(tmp_path):
    pptx = _make_pptx(tmp_path)
    plan = FixPlan(
        ops=[FixOp(slide_index=0, shape_name="test_box", op_type="shrink_font",
                   params={"target_size_pt": 14.0}, reason="test")],
        rationale="shrink test",
    )
    out = tmp_path / "out.pptx"
    PptxFixExecutor().apply(pptx, plan, output_path=out)

    shape = _load_shape(out)
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(14), f"Expected 14pt, got {run.font.size}"


def test_move(tmp_path):
    pptx = _make_pptx(tmp_path)
    target_left = int(Inches(5))
    plan = FixPlan(
        ops=[FixOp(slide_index=0, shape_name="test_box", op_type="move",
                   params={"left_emu": target_left}, reason="test")],
    )
    out = tmp_path / "out.pptx"
    PptxFixExecutor().apply(pptx, plan, output_path=out)

    shape = _load_shape(out)
    assert shape.left == target_left, f"Expected left={target_left}, got {shape.left}"


def test_resize(tmp_path):
    pptx = _make_pptx(tmp_path)
    new_w = int(Inches(6))
    new_h = int(Inches(1))
    plan = FixPlan(
        ops=[FixOp(slide_index=0, shape_name="test_box", op_type="resize",
                   params={"width_emu": new_w, "height_emu": new_h}, reason="test")],
    )
    out = tmp_path / "out.pptx"
    PptxFixExecutor().apply(pptx, plan, output_path=out)

    shape = _load_shape(out)
    assert shape.width == new_w
    assert shape.height == new_h


def test_set_word_wrap(tmp_path):
    pptx = _make_pptx(tmp_path)
    plan = FixPlan(
        ops=[FixOp(slide_index=0, shape_name="test_box", op_type="set_word_wrap",
                   params={"wrap": True}, reason="test")],
    )
    out = tmp_path / "out.pptx"
    PptxFixExecutor().apply(pptx, plan, output_path=out)

    shape = _load_shape(out)
    assert shape.text_frame.word_wrap is True


def test_unknown_shape_skipped(tmp_path):
    pptx = _make_pptx(tmp_path)
    plan = FixPlan(
        ops=[FixOp(slide_index=0, shape_name="nonexistent", op_type="shrink_font",
                   params={"target_size_pt": 10.0}, reason="test")],
    )
    out = tmp_path / "out.pptx"
    # Should not raise
    PptxFixExecutor().apply(pptx, plan, output_path=out)
    # Original shape unchanged
    shape = _load_shape(out, "test_box")
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(24)
