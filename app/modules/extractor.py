"""StructuralExtractor - extracts structure from compiled PPTX."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from ..schemas.extraction import ExtractedObject, SlideExtraction

logger = logging.getLogger(__name__)


class StructuralExtractor:
    """Extracts structural information from a compiled PPTX file."""

    def extract(self, pptx_path: str | Path) -> list[SlideExtraction]:
        """Extract structure from all slides in a PPTX."""
        prs = Presentation(str(pptx_path))
        extractions = []

        for idx, slide in enumerate(prs.slides):
            extraction = self._extract_slide(slide, idx)
            extractions.append(extraction)

        logger.info("Extracted %d slides from %s", len(extractions), pptx_path)
        return extractions

    def _extract_slide(self, slide, slide_index: int) -> SlideExtraction:
        """Extract structure from a single slide."""
        objects = []
        title = ""
        total_text = 0

        for shape in slide.shapes:
            obj = self._extract_shape(shape, slide_index)
            objects.append(obj)
            total_text += len(obj.text_content)

            # Try to identify title
            if shape.name and "title" in shape.name.lower():
                title = obj.text_content
            elif not title and shape.has_text_frame:
                # First text frame as fallback title
                if not title:
                    title = obj.text_content[:100]

        # Check notes
        has_notes = False
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
                has_notes = bool(notes_text.strip())

        # Fallback for image-based slides (e.g., HTML codegen):
        # If slide has no text objects but has notes with [TITLE]/[CONTENT]
        # markers, use the notes text as content so evaluators can read it.
        if total_text == 0 and has_notes and "[TITLE]" in notes_text:
            for line in notes_text.split("\n"):
                line = line.strip()
                if line.startswith("[TITLE] "):
                    title = line[len("[TITLE] "):]
                elif line.startswith("[CONTENT] "):
                    content_text = line[len("[CONTENT] "):]
                    # Create a synthetic text object so evaluators see content
                    objects.append(ExtractedObject(
                        object_id=f"notes_content_{slide_index}",
                        shape_name="notes_content",
                        object_type="text_box",
                        bbox_emu=[0, 0, 0, 0],
                        text_content=content_text,
                        font_sizes_pt=[],
                        has_image=False,
                        z_order=-1,
                    ))
                    total_text += len(content_text)

        return SlideExtraction(
            slide_id=slide_index + 1,
            slide_index=slide_index,
            title=title,
            objects=objects,
            total_text_length=total_text,
            total_objects=len(objects),
            has_notes=has_notes,
            notes_text=notes_text,
        )

    def _extract_shape(self, shape, slide_index: int) -> ExtractedObject:
        """Extract information from a single shape."""
        # Get bounding box in EMU
        bbox_emu = [
            shape.left or 0,
            shape.top or 0,
            shape.width or 0,
            shape.height or 0,
        ]

        # Get text content
        text_content = ""
        font_sizes = []

        if shape.has_text_frame:
            parts = []
            for para in shape.text_frame.paragraphs:
                parts.append(para.text)
                # Check paragraph-level font size
                if para.font and para.font.size:
                    font_sizes.append(para.font.size.pt)
                # Check run-level font sizes
                for run in para.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
            text_content = "\n".join(parts)

            # If no font sizes found from paragraphs/runs, check the text
            # frame's default paragraph format
            if not font_sizes and text_content.strip():
                try:
                    tf = shape.text_frame
                    if (hasattr(tf, '_txBody') and tf._txBody is not None):
                        # Try default paragraph font
                        for para in tf.paragraphs:
                            if para._pPr is not None:
                                defRPr = para._pPr.find(
                                    '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr'
                                )
                                if defRPr is not None and defRPr.get('sz'):
                                    sz = int(defRPr.get('sz')) / 100  # hundredths of point
                                    font_sizes.append(sz)
                except Exception:
                    pass

        # Determine object type
        obj_type = self._classify_shape(shape)

        # Check for image
        has_image = False
        image_path = ""
        if shape.shape_type is not None:
            try:
                if hasattr(shape, "image"):
                    has_image = True
            except Exception:
                pass

        return ExtractedObject(
            object_id=shape.name or f"shape_{shape.shape_id}",
            shape_name=shape.name or "",
            object_type=obj_type,
            bbox_emu=bbox_emu,
            text_content=text_content,
            font_sizes_pt=font_sizes,
            has_image=has_image,
            image_path=image_path,
            z_order=0,
        )

    def _classify_shape(self, shape) -> str:
        """Classify a shape into our type taxonomy."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        elif shape.has_text_frame:
            return "text_box"
        else:
            return "shape"
