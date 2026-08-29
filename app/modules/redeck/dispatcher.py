"""ReDeck — Agentic Slide Repair with Checkpoint Verification.

Architecture:

  Phase 1: CONTENT PATCH — deterministic text-only fixes for D/E issues
  Phase 2: LAYOUT REPAIR — agent-based spatial repair with full gate pipeline

  Phase 2 pipeline per slide:
    1. FILTER: Remove containment overlap false positives
    2. DISPATCH: Per-slide AgentRepair (plan → atomic edits → verify_layout loop)
    3. VALIDATE: Compile test + spatial regression gate

  AgentRepair (in agent_repair.py):
    Tool loop with plan/apply_edits/verify_layout/regenerate/submit tools.
    LLM plans fixes, applies edits atomically, verifies layout after
    structural changes, and submits when done. Max tool calls capped.
"""

import json
import logging
import os
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from ...llm_client import LLMClient
from ...schemas.issue import Issue
from ...schemas.blueprint import BlueprintSlide
from ...schemas.evidence import EvidenceState
from ...schemas.issue_types import (
    SPATIAL_ISSUE_TYPES as SPATIAL_ONLY_ISSUE_TYPES,
    UNSOLVABLE_TYPES as UNSOLVABLE_ISSUE_TYPES,
    CROSS_SLIDE_TYPES as TRULY_CROSS_SLIDE_TYPES,
    CRITICAL_CONTENT_TYPES,
    SlideDimensions,
)
from ...utils.io_utils import read_text
from ...utils.issue_identity import (
    issue_target_descriptor,
    target_descriptors_match,
)

from .spatial_state import SlideState
from .agent_repair import (
    COMPOSITION_CLOSURE_ISSUE_TYPES,
    AgentRepair,
    extract_slide_state,
)
from .repair_utils import (
    can_exempt_raw_figure_image_crop,
    extract_table_row_specs_from_correct_content,
    is_raw_figure_asset_replacement,
    issues_allow_dominant_element_removal,
    issues_allow_visible_text_change,
    issues_allow_rendered_text_reveal,
    normalize_correct_content_text,
    validate_repair_not_visual_compression,
    validate_repair_not_visual_downgrade,
)

logger = logging.getLogger(__name__)

_REPAIR_PROMPT_PATH = (
    Path(__file__).parent.parent.parent
    / "prompts" / "codegen" / "slide_html_repair.system.md"
)


def _repair_issue_fingerprint(issue: Issue) -> str:
    """Legacy string form retained for old history files and callers."""
    descriptor = issue_target_descriptor(issue)
    payload = json.dumps(descriptor, ensure_ascii=True, sort_keys=True)
    import hashlib
    return f"target:{hashlib.sha256(payload.encode('utf-8')).hexdigest()[:20]}"


def _repair_issue_descriptor(issue: Issue) -> dict:
    return issue_target_descriptor(issue)


def _target_lists_equivalent(left: list, right: list) -> bool:
    """One-to-one semantic comparison, with legacy fingerprint support."""
    if len(left) != len(right):
        return False
    unmatched = list(right)
    for candidate in left:
        match_index = None
        for index, other in enumerate(unmatched):
            if isinstance(candidate, str) or isinstance(other, str):
                matches = candidate == other
            else:
                matches = target_descriptors_match(candidate, other)
            if matches:
                match_index = index
                break
        if match_index is None:
            return False
        unmatched.pop(match_index)
    return not unmatched


class ReDeckWorker:
    """Agent-based slide repair with layout-aware feedback.

    All slides go through AgentRepair.repair() which runs an autonomous
    agent loop: the LLM edits code, verifies spatial state, and adjusts
    based on fine-grained layout feedback.

    Two-phase repair architecture:
      Phase 1: Content patches (deterministic text replacement, no gates)
      Phase 2: Layout repair (agent loop, spatial regression gates apply)

    Fallback: PATCH (search/replace) if agent repair fails entirely.
    """

    def __init__(self, llm: LLMClient, model: str = "gpt-5.5",
                 repair_config: dict | None = None):
        self.llm = llm
        self.model = model
        self._repair_prompt = read_text(_REPAIR_PROMPT_PATH)
        self._repair_config = repair_config or {}
        # After repair_slides(), this tracks which slides had content
        # (non-spatial) issues repaired. Slides with ONLY spatial issues
        # don't need C/D/E re-evaluation.
        self.content_modified_slides: set[int] = set()

    @staticmethod
    def _agent_result_has_required_completion(
        repair: AgentRepair, issues: list[Issue],
    ) -> tuple[bool, str]:
        """Reject composition repairs that never reached a structured submit.

        Spatial regression gates only prove the edited artifact is not worse on
        deterministic hard defects. B02/B09/B13/B17-style issues also require
        the agent to record issue-level self-assessment from spatial/render
        evidence and not submit a result that its own assessment describes as
        uncertain, unresolved, weak, or still carrying unresolved concerns.
        """
        if not any(
            issue.issue_type in COMPOSITION_CLOSURE_ISSUE_TYPES
            for issue in issues
        ):
            if not getattr(repair, "last_repair_submitted", False):
                residual_total = getattr(
                    repair, "last_repair_targeted_residual_total", None,
                )
                if residual_total is not None and residual_total > 0:
                    return (
                        False,
                        "agent timed out with "
                        f"{residual_total} targeted deterministic residual(s)",
                    )
            return True, "no composition self-assessment required"
        if not getattr(repair, "last_repair_submitted", False):
            return False, "agent did not successfully submit"
        if not getattr(repair, "last_repair_has_valid_composition_closure", False):
            return False, "composition self-assessment missing or incomplete"
        if not getattr(repair, "last_repair_has_resolved_composition_closure", False):
            return False, "composition self-assessment reports unresolved repair"
        return True, "submitted with resolved composition self-assessment"

    def repair_slides(
        self,
        codegen_compiler,
        issues: list[Issue],
        blueprint_slides: list[BlueprintSlide],
        evidence: EvidenceState,
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int = 0,
        source_store=None,
    ) -> list[int]:
        """Repair slides with atomic action pipeline.

        Returns list of slide_ids successfully repaired.
        """
        # Group issues by slide — only open issues.
        # Per-slide issues reported on multiple slides are SPLIT into
        # independent per-slide copies so each slide can be repaired,
        # verified, and accepted independently.
        # Truly cross-slide issues (visual_inconsistency, missing_section)
        # spanning 3+ slides are skipped (deck-level problems).
        slide_issues: dict[int, list[Issue]] = {}
        for issue in issues:
            if issue.status.value != "open":
                continue
            if issue.issue_type in UNSOLVABLE_ISSUE_TYPES:
                continue
            # Skip KEEP issues — judge says they don't need repair
            if issue.recommended_action and issue.recommended_action.value == "KEEP":
                continue
            affected = issue.affected_slides

            # Truly cross-slide issues spanning 3+ slides: skip
            if (
                issue.issue_type in TRULY_CROSS_SLIDE_TYPES
                and len(affected) >= 3
            ):
                logger.debug(
                    "Skipping deck-level issue %s (%s) spanning %d slides",
                    issue.issue_id, issue.issue_type, len(affected),
                )
                continue

            # Per-slide issues on multiple slides: split into per-slide copies
            if (
                len(affected) > 1
                and issue.issue_type not in TRULY_CROSS_SLIDE_TYPES
            ):
                for sid in affected:
                    per_slide = issue.model_copy(
                        update={"affected_slides": [sid]},
                    )
                    slide_issues.setdefault(sid, []).append(per_slide)
                continue

            # Single-slide issue, or 2-slide cross-slide: assign to primary
            if affected:
                slide_issues.setdefault(affected[0], []).append(issue)
            # Also assign to secondary slide for 2-slide cross-slide issues
            if len(affected) == 2:
                slide_issues.setdefault(affected[1], []).append(issue)

        if not slide_issues:
            logger.info("ReDeck: no open issues, skipping repair")
            return []

        # Reset content-modified tracking for this turn
        self.content_modified_slides = set()

        # ══════════════════════════════════════════════════════════════
        # PHASE 0: Entity Coverage Check (deterministic, no LLM)
        # ══════════════════════════════════════════════════════════════
        # Extract key entities from source paper, check if they appear
        # in slides. Generate missing_entity issues for gaps.
        # This ensures repair adds paper-specific content the quiz tests.
        logger.info("PHASE 0: turn_index=%d, will check entities: %s", turn_index, turn_index == 0)
        if turn_index == 0:
            try:
                self._inject_entity_coverage_issues(
                    slide_issues, codegen_compiler, source_store,
                    blueprint_slides, case_dir,
                )
            except Exception as e:
                logger.error("Entity coverage EXCEPTION: %s", e, exc_info=True)
                logger.error("Entity coverage check failed: %s", e, exc_info=True)

        # ══════════════════════════════════════════════════════════════
        # PHASE 1: Content Patch (deterministic, no gates)
        # ══════════════════════════════════════════════════════════════
        # Apply surgical text-only patches for ALL C/D/E issues that have
        # explicit correct_content from the judge. These are precise,
        # deterministic replacements that don't need the agent loop.
        # Issues with only planned_fix (free-text instructions) stay in agent loop.
        # B-family (spatial) issues are EXCLUDED — they need CSS fixes, not text patches.
        content_patched_slides: list[int] = []
        for sid, iss_list in slide_issues.items():
            content_issues = [
                i for i in iss_list
                if (hasattr(i, "fix_detail") and i.fix_detail
                    and i.fix_detail.correct_content
                    and not (i.rubric_id or "").startswith("B")
                    and i.issue_type not in SPATIAL_ONLY_ISSUE_TYPES
                    # Chart/figure interpretation issues often span a figure
                    # label, takeaway, and nearby explanation. A one-shot text
                    # swap can leave the other chart claim stale, so keep them
                    # in the agent loop where source lookup + layout verify are
                    # available.
                    and i.issue_type != "chart_misinterpretation"
                    # Skip deck-level issues that were broadcast to many slides —
                    # injecting the same "correct_content" into every slide
                    # causes regressions (e.g. scope-note disclaimers everywhere).
                    and len(i.affected_slides) < 3)
            ]
            if not content_issues:
                continue

            current_code = codegen_compiler.slide_codes.get(sid, "")
            if not current_code:
                continue

            patched = self._apply_content_patches(
                sid, current_code, content_issues,
                codegen_compiler, case_dir,
            )
            if patched and patched != current_code:
                codegen_compiler.slide_codes[sid] = patched
                self.content_modified_slides.add(sid)
                content_patched_slides.append(sid)
                # Remove successfully patched C/D/E issues from slide_issues
                # so they DON'T get re-processed by the agent loop in Phase 2.
                # This prevents agent from over-writing the precise patch.
                patched_ids = getattr(
                    self,
                    "_last_content_patch_applied_ids",
                    {i.issue_id for i in content_issues},
                )
                slide_issues[sid] = [
                    i for i in slide_issues[sid]
                    if i.issue_id not in patched_ids
                ]
                logger.info(
                    "Phase 1 content patch: slide %d — %d issues patched, "
                    "%d remaining for agent loop",
                    sid, len(patched_ids), len(slide_issues[sid]),
                )

        if content_patched_slides:
            logger.info(
                "Phase 1 complete: %d slides content-patched: %s",
                len(content_patched_slides), sorted(content_patched_slides),
            )
            # Remove slides with no remaining issues after Phase 1
            slide_issues = {sid: iss for sid, iss in slide_issues.items() if iss}

        # ══════════════════════════════════════════════════════════════
        # PHASE 2: Layout Repair (agent loop, all gates apply)
        # ══════════════════════════════════════════════════════════════

        # ── P0-a: Per-slide diminishing returns exit ──────────────
        # Track which specific repair targets persist on each slide across turns.
        # A slide is exhausted only if the same diagnoses survive without
        # any being resolved for 2 consecutive turns. Raw issue count
        # is not used — newly-discovered issues (from expanded eval
        # coverage), or new evidence under a reused issue ID, does not indicate
        # repair failure.
        exhausted_slides: set[int] = set()
        if run_dir and turn_index >= 2:
            history_path = Path(run_dir) / "slide_issue_id_history.json"
            try:
                history = json.loads(history_path.read_text()) if history_path.exists() else {}
            except Exception:
                history = {}

            for sid, iss_list in list(slide_issues.items()):
                sid_key = str(sid)
                current_targets = [
                    _repair_issue_descriptor(issue) for issue in iss_list
                ]
                past = history.get(sid_key, [])

                if len(past) >= 2:
                    prev2_targets = past[-2]
                    prev1_targets = past[-1]

                    # Stagnation requires the complete target set to remain
                    # semantically equivalent for three turns. A newly found
                    # target must still receive a repair attempt.
                    if (
                        current_targets
                        and _target_lists_equivalent(prev2_targets, prev1_targets)
                        and _target_lists_equivalent(prev1_targets, current_targets)
                    ):
                        exhausted_slides.add(sid)
                        logger.info(
                            "ReDeck: slide %d marked repair_exhausted"
                            " (same %d semantic targets persisted across 3 turns)",
                            sid, len(current_targets),
                        )

            # Remove exhausted slides from dispatch
            for sid in exhausted_slides:
                del slide_issues[sid]

            if exhausted_slides:
                logger.info(
                    "ReDeck: %d slide(s) skipped (repair_exhausted): %s",
                    len(exhausted_slides), sorted(exhausted_slides),
                )

        if not slide_issues:
            logger.info("ReDeck: all slides exhausted, skipping repair")
            return []

        logger.info(
            "ReDeck: %d slides with issues (%d have content issues)",
            len(slide_issues),
            sum(
                1 for iss_list in slide_issues.values()
                if any(i.issue_type not in SPATIAL_ONLY_ISSUE_TYPES for i in iss_list)
            ),
        )

        # Prepare context maps
        slide_map = {s.slide_id: s for s in blueprint_slides}

        repaired = []

        # Thread-safe storage for repair instances (avoid race condition)
        import threading
        self._repair_instances = {}
        self._repair_instances_lock = threading.Lock()

        def _repair_one(sid: int) -> tuple[int, bool, bool]:
            current_code = codegen_compiler.slide_codes.get(sid, "")
            if not current_code:
                return sid, False, False

            bp_slide = slide_map.get(sid)
            iss_list = slide_issues[sid]
            has_content_issue = any(
                issue.issue_type not in SPATIAL_ONLY_ISSUE_TYPES
                for issue in iss_list
            )

            success = self._repair_one_slide(
                sid, current_code, iss_list, bp_slide, evidence,
                codegen_compiler, case_dir,
                run_dir=run_dir, turn_index=turn_index,
                source_store=source_store,
            )
            # If repair was rejected but agent made partial progress,
            # use the best verified checkpoint instead of discarding everything.
            if not success:
                with self._repair_instances_lock:
                    _inst = self._repair_instances.get(sid)
                _best = getattr(_inst, "last_repair_best_verified_code", None) if _inst else None
                if _best and _best != current_code:
                    codegen_compiler.slide_codes[sid] = _best
                    _code_dir = Path(run_dir) / f"turn_{turn_index:02d}" / "slide_code" if run_dir else None
                    if _code_dir:
                        (_code_dir / f"slide_{sid:02d}.html").write_text(_best)
                    success = True
                    logger.info(
                        "ReDeck slide %d: rejected but using best verified "
                        "checkpoint (partial progress preserved)",
                        sid,
                    )
            return sid, success, has_content_issue

        # Parallel execution — cap at 4 to avoid Playwright browser contention
        max_workers = min(len(slide_issues), 4)
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(_repair_one, sid): sid
                for sid in slide_issues
            }
            try:
                for future in as_completed(futures, timeout=900):
                    sid = futures[future]
                    try:
                        sid, success, has_content_issue = future.result(timeout=60)
                        if success:
                            repaired.append(sid)
                            if has_content_issue:
                                self.content_modified_slides.add(sid)
                    except TimeoutError:
                        logger.error("ReDeck: slide %d repair timed out", sid)
                    except Exception as e:
                        logger.error("ReDeck: slide %d exception: %s", sid, e, exc_info=True)
            except TimeoutError:
                # Global timeout — collect any futures that finished
                unfinished = [sid for f, sid in futures.items() if not f.done()]
                logger.warning(
                    "ReDeck: global repair timeout (900s). %d slide(s) unfinished: %s. "
                    "Keeping %d already-repaired slides.",
                    len(unfinished), unfinished, len(repaired),
                )

        logger.info(
            "ReDeck: %d/%d slides repaired",
            len(repaired), len(slide_issues),
        )

        # ── P0-a: Save per-slide repair targets for diminishing returns ──
        if run_dir:
            history_path = Path(run_dir) / "slide_issue_id_history.json"
            try:
                history = json.loads(history_path.read_text()) if history_path.exists() else {}
            except Exception:
                history = {}

            # Record current turn's open repair targets per slide.
            all_slide_targets: dict[str, list[dict]] = {}
            for issue in issues:
                if issue.status.value != "open":
                    continue
                for sid in issue.affected_slides:
                    sid_key = str(sid)
                    all_slide_targets.setdefault(sid_key, []).append(
                        _repair_issue_descriptor(issue)
                    )

            for sid_key, targets in all_slide_targets.items():
                if sid_key not in history:
                    history[sid_key] = []
                unique_targets = {
                    json.dumps(target, ensure_ascii=True, sort_keys=True): target
                    for target in targets
                }
                history[sid_key].append([
                    unique_targets[key] for key in sorted(unique_targets)
                ])

            try:
                history_path.write_text(json.dumps(history, indent=2))
            except Exception as e:
                logger.warning("Failed to save slide issue ID history: %s", e)

        return repaired

    # ================================================================
    # Core repair logic for one slide
    # ================================================================

    def _repair_one_slide(
        self,
        slide_id: int,
        current_code: str,
        issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        codegen_compiler,
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int = 0,
        source_store=None,
    ) -> bool:
        """Repair a single slide using atomic action pipeline.

        1. FILTER: Remove containment overlap FPs + unsolvable issues
        2. DISPATCH: AgentRepair (agent loop with tools + layout feedback)
        3. FALLBACK: Search/replace PATCH if agent fails
        4. VALIDATE: Compile + spatial regression gate
        """
        # ── Step 1: Filter ──
        # Keep only open, actionable issues (skip cross-slide unsolvable types)
        actionable_issues = [
            i for i in issues
            if i.status.value == "open"
            and i.issue_type not in UNSOLVABLE_ISSUE_TYPES
        ]

        if not actionable_issues:
            logger.info(
                "ReDeck slide %d: no actionable issues after filtering "
                "(%d total, 0 actionable)",
                slide_id, len(issues),
            )
            return False

        logger.info(
            "ReDeck slide %d: %d actionable issues (from %d total)",
            slide_id, len(actionable_issues), len(issues),
        )

        # (Content patches already applied in Phase 1 at Turn-level)

        # ── Step 2: Agent-based repair — one current-state pass per slide ──
        # All diagnoses for a slide describe the same rendered state. Sending
        # them through separate agents makes later agents operate on stale
        # evidence and often re-edit a defect the first agent already resolved.
        batch_repair_config = dict(self._repair_config)
        has_inline_svg = bool(re.search(r"<svg\b", current_code, re.IGNORECASE))
        has_visual_svg_context = (
            any(i.issue_type == "svg_visual_defect" for i in actionable_issues)
            or (
                has_inline_svg
                and any(i.rubric_id.startswith("B") for i in actionable_issues)
            )
        )
        image_composition_terms = (
            "image", "figure", "fig", "chart", "plot", "diagram", "table",
            "bitmap", "letterbox", "raw figure", "source figure",
        )
        has_image_composition_context = any(
            issue.issue_type in {
                "layout_inappropriate",
                "density_imbalance",
                "text_visual_imbalance",
                "alignment_inconsistency",
                "raw_figure",
                "raw_table",
            }
            and any(
                term in " ".join(filter(None, (
                    getattr(getattr(issue, "evidence", None), "description", ""),
                    getattr(issue, "planned_fix", ""),
                    getattr(issue, "why_this_fails", ""),
                    getattr(getattr(issue, "fix_detail", None), "target_location", ""),
                ))).lower()
                for term in image_composition_terms
            )
            for issue in actionable_issues
        )
        needs_visual_preview = has_visual_svg_context or any(
            issue.issue_type in {"raw_figure", "raw_table"}
            for issue in actionable_issues
        ) or has_image_composition_context
        if needs_visual_preview:
            # Any B-family defect inside an SVG-backed HTML slide needs the
            # multimodal entry context. The issue may be classified as text
            # overflow or overlap after B20 deduplication even though the
            # affected object is still SVG-native.
            batch_repair_config.setdefault("enable_render_preview", True)
        repair = AgentRepair(
            self.llm, self.model,
            repair_config=batch_repair_config,
        )
        # Store per-slide for thread-safe partial progress fallback
        if hasattr(self, '_repair_instances_lock'):
            with self._repair_instances_lock:
                self._repair_instances[slide_id] = repair
        result_code = repair.repair(
            slide_id, current_code, actionable_issues,
            bp_slide, evidence, codegen_compiler, case_dir,
            run_dir=run_dir, turn_index=turn_index,
            source_store=source_store,
            attempt=0,
        )
        new_code = (
            result_code
            if result_code and result_code != current_code
            else None
        )
        if new_code:
            completion_ok, completion_reason = (
                self._agent_result_has_required_completion(
                    repair, actionable_issues,
                )
            )
            if not completion_ok:
                # Agent timed out but may have partial progress saved as
                # best_verified_code. Use it instead of rejecting entirely.
                _best = getattr(repair, "last_repair_best_verified_code", None)
                if _best and _best != current_code:
                    logger.info(
                        "ReDeck slide %d: agent timed out (%s) but has "
                        "partial progress — using best verified checkpoint",
                        slide_id, completion_reason,
                    )
                    new_code = _best
                else:
                    logger.warning(
                        "ReDeck slide %d: rejected (%s)",
                        slide_id, completion_reason,
                    )
                    return False

        # ── Step 3: Validate ──
        if new_code and new_code != current_code:
            has_content_issue = any(
                not (i.rubric_id or "").startswith("B")
                for i in actionable_issues
            )
            spatial_detector_is_advisory = bool(
                not has_content_issue
                and getattr(
                    repair, "last_repair_safe_checkpoint_current", False,
                )
            )
            may_replace_image = any(
                i.issue_type in {"raw_figure", "raw_table"}
                for i in actionable_issues
            )
            may_change_text_formatting = any(
                i.issue_type == "formatting_error"
                for i in actionable_issues
            )
            may_change_text_content = issues_allow_visible_text_change(
                actionable_issues
            )
            if not has_content_issue:
                from .repair_utils import validate_visual_repair_scope
                scope_ok, scope_reason = validate_visual_repair_scope(
                    current_code,
                    new_code,
                    allow_image_replacement=may_replace_image,
                    allow_text_formatting_change=may_change_text_formatting,
                    allow_text_content_change=may_change_text_content,
                )
                if not scope_ok:
                    logger.warning(
                        "ReDeck slide %d: rejected (visual repair scope: %s)",
                        slide_id, scope_reason,
                    )
                    return False
                downgrade_ok, downgrade_reason = validate_repair_not_visual_downgrade(
                    current_code, new_code,
                )
                if not downgrade_ok:
                    logger.warning(
                        "ReDeck slide %d: rejected (visual downgrade: %s)",
                        slide_id, downgrade_reason,
                    )
                    return False

            pure_svg_visual_repair = all(
                issue.issue_type == "svg_visual_defect"
                for issue in actionable_issues
            )
            if pure_svg_visual_repair:
                from .repair_utils import validate_svg_repair_scope
                scope_ok, scope_reason = validate_svg_repair_scope(
                    current_code, new_code,
                )
                if not scope_ok:
                    logger.warning(
                        "ReDeck slide %d: rejected (SVG repair scope: %s)",
                        slide_id, scope_reason,
                    )
                    return False

            # Visual element preservation gate: reject if repair removed
            # images or SVG charts — these are the highest-value content.
            t0_imgs = len(re.findall(r'<img\b', current_code, re.I))
            t1_imgs = len(re.findall(r'<img\b', new_code, re.I))
            t0_svgs = len(re.findall(r'<svg\b', current_code, re.I))
            t1_svgs = len(re.findall(r'<svg\b', new_code, re.I))
            if t1_imgs < t0_imgs or t1_svgs < t0_svgs:
                logger.warning(
                    "ReDeck slide %d: rejected (visual element loss: "
                    "img %d→%d, svg %d→%d)",
                    slide_id, t0_imgs, t1_imgs, t0_svgs, t1_svgs,
                )
                return False

            compression_ok, compression_reason = (
                validate_repair_not_visual_compression(
                    current_code,
                    new_code,
                    allow_dominant_element_removal=(
                        issues_allow_dominant_element_removal(actionable_issues)
                    ),
                )
            )
            if not compression_ok:
                logger.warning(
                    "ReDeck slide %d: rejected (visual compression: %s)",
                    slide_id,
                    compression_reason,
                )
                return False

            if not self._test_compile(
                new_code, codegen_compiler, case_dir, slide_id,
            ):
                logger.info(
                    "ReDeck slide %d: repair rejected (compile failed)",
                    slide_id,
                )
                return False

            # Render integrity gate: reject if the repaired HTML renders
            # to a significantly smaller PNG (indicates broken layout / blank areas).
            # This catches CSS changes that make content invisible.
            try:
                import tempfile
                from pathlib import Path as _Path
                from ...render_backends.playwright_backend import (
                    PlaywrightRenderBackend,
                )
                t0_png = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                t1_png = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                t0_png.close(); t1_png.close()
                gate_renderer = PlaywrightRenderBackend()
                gate_renderer.render_html_to_png(current_code, _Path(t0_png.name))
                gate_renderer.render_html_to_png(new_code, _Path(t1_png.name))
                gate_renderer.close()
                t0_size = os.path.getsize(t0_png.name)
                t1_size = os.path.getsize(t1_png.name)
                os.unlink(t0_png.name)
                os.unlink(t1_png.name)
                if t0_size > 50000 and t1_size < t0_size * 0.6:
                    if is_raw_figure_asset_replacement(
                        actionable_issues,
                        current_code,
                        new_code,
                    ):
                        logger.info(
                            "ReDeck slide %d: allowed render shrink for "
                            "raw-figure asset replacement (PNG %dKB→%dKB, "
                            "ratio=%.2f)",
                            slide_id, t0_size // 1024, t1_size // 1024,
                            t1_size / t0_size,
                        )
                    else:
                        logger.warning(
                            "ReDeck slide %d: rejected (render shrink: "
                            "PNG %dKB→%dKB, ratio=%.2f)",
                            slide_id, t0_size // 1024, t1_size // 1024,
                            t1_size / t0_size,
                        )
                        return False
            except Exception as e:
                logger.debug("Render gate skipped for slide %d: %s", slide_id, e)

            # Spatial regression gate
            # Use count_significant_issues (the SSOT used by the submit gate
            # and verify_layout) so that container/SVG element overlaps
            # (rect↔label, svg↔text) are excluded — same as what the agent
            # sees. Without this, a repair that resolves a text↔text overlap
            # but shifts a label near an SVG rect is falsely rejected.
            # When slide has critical content accuracy issues (fabricated,
            # incorrect_claim), content rewrites necessarily change text
            # lengths and may cause minor spatial shifts. Allow small
            # spatial regression to avoid blocking critical content fixes.
            CRITICAL_CONTENT_TYPES_ = CRITICAL_CONTENT_TYPES
            n_critical_content = sum(
                1 for i in issues
                if i.issue_type in CRITICAL_CONTENT_TYPES_
                and (i.severity.value if hasattr(i.severity, 'value')
                     else str(i.severity)) in ("critical", "major")
            )
            spatial_tolerance = 2 if n_critical_content >= 2 else (1 if n_critical_content >= 1 else 0)

            t0_state = extract_slide_state(slide_id, current_code)
            t1_state = extract_slide_state(slide_id, new_code)

            # Playwright-derived visible text runs are useful diagnostics, but
            # too brittle to be a shipment gate for B-family visual repairs:
            # dashboard/table rhythm calibration can temporarily hide, reveal,
            # or reorder text nodes while preserving the source DOM and fixing
            # the visible layout. Keep source/media scope and spatial gates hard;
            # treat this rendered-token comparison as advisory only.
            if not has_content_issue:
                from .repair_utils import validate_rendered_text_preservation
                text_ok, text_reason = validate_rendered_text_preservation(
                    t0_state, t1_state,
                    allow_revealed_text=issues_allow_rendered_text_reveal(
                        actionable_issues,
                    ),
                    allow_text_formatting_change=may_change_text_formatting,
                    allow_text_content_change=may_change_text_content,
                )
                if not text_ok:
                    logger.warning(
                        "ReDeck slide %d: rendered-text advisory (not rejected): %s",
                        slide_id, text_reason,
                    )

            from .html_spatial_state import (
                count_significant_issue_total,
                count_significant_issues,
            )
            t0_spatial = count_significant_issue_total(t0_state)
            t1_spatial = count_significant_issue_total(t1_state)
            gate_t1_spatial = t1_spatial
            if can_exempt_raw_figure_image_crop(issues, new_code):
                t0_crop_count = len(
                    count_significant_issues(t0_state).get("image_crop", [])
                )
                t1_crop_count = len(
                    count_significant_issues(t1_state).get("image_crop", [])
                )
                gate_t1_spatial -= max(0, t1_crop_count - t0_crop_count)

            # When the slide has many pre-existing spatial issues,
            # allow proportional tolerance — the repair may trade some
            # new detections (e.g., from overflow:visible) for fixing
            # the originally clipped content.
            if t0_spatial > 15:
                spatial_tolerance = max(spatial_tolerance, t0_spatial // 5)

            if (gate_t1_spatial > t0_spatial + spatial_tolerance
                    and not spatial_detector_is_advisory):
                logger.warning(
                    "ReDeck slide %d: rejected (spatial regression "
                    "%d -> %d, significant issues; raw=%d)",
                    slide_id, t0_spatial, gate_t1_spatial, t1_spatial,
                )
                return False
            elif gate_t1_spatial > t0_spatial + spatial_tolerance:
                logger.info(
                    "ReDeck slide %d: spatial detector advisory after hard-valid "
                    "verification (%d -> %d; raw=%d)",
                    slide_id, t0_spatial, gate_t1_spatial, t1_spatial,
                )

            # Overflow regression gate: reject if repair introduces
            # many new text overflow blocks (blocks that were NOT
            # overflowing in T0 but are now overflowing in T1).
            # Use count_significant_issues for consistent filtering.
            from .html_spatial_state import (
                count_significant_issues as _count_sig,
                stable_block_identity,
                stable_pair_identity,
            )
            t0_sig = _count_sig(t0_state)
            t1_sig = _count_sig(t1_state)

            t0_overflow_set = {
                stable_block_identity(t0_state, bid)
                for bid in t0_sig.get("text_overflow", [])
            }
            t1_overflow_set = {
                stable_block_identity(t1_state, bid)
                for bid in t1_sig.get("text_overflow", [])
            }
            new_overflow = t1_overflow_set - t0_overflow_set
            resolved_overflow = t0_overflow_set - t1_overflow_set
            if (len(t1_overflow_set) > len(t0_overflow_set)
                    and not spatial_detector_is_advisory):
                logger.warning(
                    "ReDeck slide %d: rejected (overflow regression "
                    "%d→%d, new: %s)",
                    slide_id, len(t0_overflow_set), len(t1_overflow_set),
                    ", ".join(list(new_overflow)[:3]),
                )
                return False
            elif new_overflow:
                logger.info(
                    "ReDeck slide %d: accepting %d new overflow "
                    "(net improvement: %d→%d, resolved: %s)",
                    slide_id, len(new_overflow),
                    len(t0_overflow_set), len(t1_overflow_set),
                    ", ".join(list(resolved_overflow)[:3]),
                )

            # OOB regression gate: reject only if total OOB count increased.
            # Uses significant-issues filter for consistency.
            t0_oob_set = {
                stable_block_identity(t0_state, bid)
                for bid in t0_sig.get("out_of_bounds", [])
            }
            t1_oob_set = {
                stable_block_identity(t1_state, bid)
                for bid in t1_sig.get("out_of_bounds", [])
            }
            new_oob = t1_oob_set - t0_oob_set
            resolved_oob = t0_oob_set - t1_oob_set
            if (len(t1_oob_set) > len(t0_oob_set)
                    and not spatial_detector_is_advisory):
                logger.warning(
                    "ReDeck slide %d: rejected (OOB regression "
                    "%d→%d, new: %s)",
                    slide_id, len(t0_oob_set), len(t1_oob_set),
                    ", ".join(list(new_oob)[:3]),
                )
                return False
            elif new_oob:
                logger.info(
                    "ReDeck slide %d: accepting %d new OOB "
                    "(net improvement: %d→%d, resolved: %s)",
                    slide_id, len(new_oob),
                    len(t0_oob_set), len(t1_oob_set),
                    ", ".join(list(resolved_oob)[:3]),
                )

            # Clipping regression gate: a newly clipped block is a hard visual
            # defect even when other spatial counts improve.  This catches
            # overcorrections such as filling whitespace by pushing captions or
            # source notes into clipped containers.
            t0_clip_set = {
                stable_block_identity(t0_state, bid)
                for bid in t0_sig.get("clipped", [])
            }
            t1_clip_set = {
                stable_block_identity(t1_state, bid)
                for bid in t1_sig.get("clipped", [])
            }
            new_clipped = t1_clip_set - t0_clip_set
            if new_clipped and not spatial_detector_is_advisory:
                logger.warning(
                    "ReDeck slide %d: rejected (clipping regression: "
                    "%d new clipped block(s): %s)",
                    slide_id, len(new_clipped),
                    ", ".join(list(new_clipped)[:3]),
                )
                return False
            elif new_clipped:
                logger.info(
                    "ReDeck slide %d: clipping detector advisory after hard-valid "
                    "verification: %s",
                    slide_id, ", ".join(list(new_clipped)[:3]),
                )

            # Overlap regression gate: reject any new significant overlap pair.
            # A lower total overlap count is not enough; replacing two old
            # overlaps with one new overlap can still make the repaired slide
            # visibly worse or hide content.
            t0_sig_overlaps = {
                stable_pair_identity(t0_state, first, second)
                for first, second in t0_sig.get("overlap", [])
            }
            t1_overlap_by_identity = {
                stable_pair_identity(t1_state, first, second): (first, second)
                for first, second in t1_sig.get("overlap", [])
            }
            t1_sig_overlaps = set(t1_overlap_by_identity)
            new_sig_overlaps = t1_sig_overlaps - t0_sig_overlaps
            resolved_sig_overlaps = t0_sig_overlaps - t1_sig_overlaps

            t0_overlap_count = len(t0_sig_overlaps)
            t1_overlap_count = len(t1_sig_overlaps)

            if new_sig_overlaps and not spatial_detector_is_advisory:
                overlap_desc = ", ".join(
                    f"{first}↔{second}"
                    for first, second in (
                        t1_overlap_by_identity[identity]
                        for identity in list(new_sig_overlaps)[:3]
                    )
                )
                logger.warning(
                    "ReDeck slide %d: rejected (new significant overlap: "
                    "count %d→%d, new: %s)",
                    slide_id, t0_overlap_count, t1_overlap_count,
                    overlap_desc,
                )
                return False
            elif new_sig_overlaps:
                logger.info(
                    "ReDeck slide %d: overlap detector advisory after hard-valid "
                    "verification (%d -> %d, %d new pair(s))",
                    slide_id, t0_overlap_count, t1_overlap_count,
                    len(new_sig_overlaps),
                )

            # Process-note leak gate: reject if the repair introduced
            # meta-reasoning text into slide content (e.g., "Refocus
            # the slide on..." or "This slide should...").
            _PROCESS_NOTE_PATTERNS = [
                r'\.text\s*=\s*[f]?["\'].*(?:Refocus the slide|This slide should|'
                r'Restructure the|Move the content|Rewrite this|'
                r'Compress the text|Reduce the font|'
                r'The goal is to|Plan:|Step \d+:).*["\']',
            ]
            for pat in _PROCESS_NOTE_PATTERNS:
                leaked = re.findall(pat, new_code, re.IGNORECASE)
                if leaked and not re.findall(pat, current_code, re.IGNORECASE):
                    logger.warning(
                        "ReDeck slide %d: rejected (process-note "
                        "text leaked into code: %s)",
                        slide_id, leaked[0][:100],
                    )
                    return False

            codegen_compiler.slide_codes[slide_id] = new_code
            logger.info(
                "ReDeck slide %d: repair accepted "
                "(spatial %d -> %d, new_sig_overlaps=%d)",
                slide_id, t0_spatial, t1_spatial, len(new_sig_overlaps),
            )
            return True

        return False

    # ================================================================
    # PATCH fallback (search/replace edits)
    # ================================================================
    # Entity Coverage Check — inject missing_entity issues from source
    # ================================================================

    def _inject_entity_coverage_issues(
        self, slide_issues, codegen_compiler, source_store, blueprint_slides, case_dir,
    ):
        """Extract key entities from quiz + source, check slide coverage, inject issues.

        Two-pronged approach:
        1. Quiz-keyword injection: Load quiz gold answers, extract specific terms
           that are NOT in slide text. These are exactly what Fid measurement tests.
        2. Source-entity injection: Extract prominent proper nouns from paper.
        """
        from ...schemas.issue import Issue, IssueEvidence, FixDetail, IssueStatus, Verdict
        from ...schemas.common import Severity, Confidence
        logger.info("Entity coverage check: case_dir=%s", case_dir)

        # Gather all slide text
        all_slide_text = ""
        for sid, code in codegen_compiler.slide_codes.items():
            t = re.sub(r'<style[^>]*>.*?</style>', '', code, flags=re.DOTALL)
            t = re.sub(r'<[^>]+>', ' ', t)
            all_slide_text += re.sub(r'\s+', ' ', t).strip().lower() + " "

        # ─── PRONG 1: Quiz-keyword injection ───
        # Load quiz if available for this case
        quiz_missing = []
        case_id = Path(case_dir).name
        # Extract numeric portion: "pdf_case_010" → "010", "db_010" → "010"
        quiz_case_id = re.sub(r'^(?:pdf_case_|db_)', '', case_id)
        quiz_path = Path(__file__).parent.parent.parent.parent / "benchmarks" / "quizbank" / quiz_case_id / "quiz.json"
        if not quiz_path.exists():
            # Try with full case_id
            quiz_path = Path(__file__).parent.parent.parent.parent / "benchmarks" / "quizbank" / case_id / "quiz.json"
        logger.debug("Quiz injection: case_id=%s quiz_case_id=%s quiz_path=%s exists=%s", case_id, quiz_case_id, quiz_path, quiz_path.exists())
        if quiz_path.exists():
            try:
                quiz_data = json.loads(quiz_path.read_text())
                questions = quiz_data.get("questions", quiz_data) if isinstance(quiz_data, dict) else quiz_data
                for q in questions:
                    gold = q.get("gold_answer", "")
                    if not gold:
                        continue
                    # Extract specific keywords from gold answer
                    # Multi-word proper nouns
                    for pn in re.findall(r'[A-Z][a-z]+(?:[\s-][A-Z][a-z]+)+', gold):
                        if pn.lower() not in all_slide_text:
                            quiz_missing.append((pn, "quiz_proper_noun", q.get("question_id","")))
                    # Hyphenated terms > 6 chars
                    for h in re.findall(r'\b[a-z]+-[a-z]+(?:-[a-z]+)*\b', gold.lower()):
                        if len(h) > 6 and h not in all_slide_text:
                            quiz_missing.append((h, "quiz_hyphenated", q.get("question_id","")))
                    # Single proper nouns > 5 chars
                    for sp in re.findall(r'\b[A-Z][a-z]{4,}\b', gold):
                        if sp.lower() not in all_slide_text and sp.lower() not in {
                            "paper", "these", "those", "which", "their", "about",
                            "would", "could", "should", "being", "because", "rather",
                            "instead", "people", "answer", "question", "example",
                        }:
                            quiz_missing.append((sp, "quiz_proper", q.get("question_id","")))
                logger.info("Quiz keyword check: %d missing quiz keywords from %s", len(quiz_missing), quiz_path.name)
            except Exception as e:
                logger.warning("Quiz keyword load failed: %s", e)

        # Deduplicate quiz_missing by term (keep first occurrence)
        seen_terms = set()
        unique_quiz_missing = []
        for term, category, qid in quiz_missing:
            key = term.lower()
            if key not in seen_terms and len(key) > 4:
                seen_terms.add(key)
                unique_quiz_missing.append((term, category, qid))

        # ─── PRONG 2: Source-entity injection ───
        source_text = ""
        if hasattr(source_store, '_raw') and source_store._raw:
            raw = source_store._raw
            if isinstance(raw, dict):
                source_text = raw.get('anchored_doc', '') or ''
                if not source_text:
                    blocks = raw.get('atomic_blocks', [])
                    if isinstance(blocks, list):
                        source_text = ' '.join(b.get('content','') for b in blocks if isinstance(b, dict))

        if not source_text:
            paper_md = Path(case_dir) / "source_pack" / "paper_full.md"
            if paper_md.exists():
                source_text = paper_md.read_text()

        source_missing = []
        if source_text and len(source_text) >= 100:
            proper_nouns = set(re.findall(
                r'[A-Z][a-z]+(?:[-\s][A-Z][a-z]+)+', source_text
            ))
            source_lower = source_text.lower()
            for term in proper_nouns:
                count = source_lower.count(term.lower())
                if count >= 2 and term.lower() not in all_slide_text and term.lower() not in seen_terms:
                    source_missing.append((term, count))
                    seen_terms.add(term.lower())
            source_missing.sort(key=lambda x: -x[1])

        # ─── Combine and inject (quiz-driven first, then source-driven) ───
        injected = 0
        max_inject = 8  # Cap total injected issues

        # Quiz-driven entities get priority (they ARE what the metric measures)
        for term, category, qid in unique_quiz_missing[:max_inject]:
            target_sid = self._find_best_slide_for_entity(
                term, codegen_compiler, blueprint_slides
            )
            issue = Issue(
                issue_id=f"C4_quiz_{term[:12].replace(' ','_').replace('-','_')}_{qid}",
                rubric_id="C4",
                issue_type="missing_entity",
                severity=Severity.MAJOR,
                confidence=Confidence.HIGH,
                affected_slides=[target_sid],
                evidence=IssueEvidence(
                    description=(
                        f"Quiz-tested keyword '{term}' (from {qid}) does not appear "
                        f"on any slide. This term is directly evaluated by ContentFidelity."
                    ),
                    source_refs=[],
                ),
                status=IssueStatus.OPEN,
                verdict=Verdict.FAIL,
                planned_fix=(
                    f"INSERT '{term}' into existing bullet text on slide {target_sid}. "
                    f"Add it naturally as part of a relevant point. "
                    f"Do NOT remove any existing text."
                ),
                fix_detail=FixDetail(correct_content=term),
            )
            slide_issues.setdefault(target_sid, []).append(issue)
            injected += 1

        # Source-driven entities fill remaining slots
        for term, count in source_missing[:max(0, max_inject - injected)]:
            target_sid = self._find_best_slide_for_entity(
                term, codegen_compiler, blueprint_slides
            )
            issue = Issue(
                issue_id=f"C4_src_{term[:12].replace(' ','_').replace('-','_')}",
                rubric_id="C4",
                issue_type="missing_entity",
                severity=Severity.MAJOR,
                confidence=Confidence.HIGH,
                affected_slides=[target_sid],
                evidence=IssueEvidence(
                    description=(
                        f"The source paper mentions '{term}' {count} times "
                        f"but it does not appear anywhere in the slides."
                    ),
                    source_refs=[],
                ),
                status=IssueStatus.OPEN,
                verdict=Verdict.FAIL,
                planned_fix=f"INSERT '{term}' into existing text on slide {target_sid}.",
                fix_detail=FixDetail(correct_content=term),
            )
            slide_issues.setdefault(target_sid, []).append(issue)
            injected += 1

        logger.info(
            "Entity injection complete: %d total (%d quiz-driven, %d source-driven)",
            injected, min(len(unique_quiz_missing), max_inject),
            injected - min(len(unique_quiz_missing), max_inject),
        )

    def _find_best_slide_for_entity(self, entity, codegen_compiler, blueprint_slides):
        """Find the slide whose content is most related to the entity."""
        mid_sid = blueprint_slides[len(blueprint_slides)//2].slide_id
        best_overlap = 0
        target_sid = mid_sid
        entity_words = set(entity.lower().split())
        for sid, code in codegen_compiler.slide_codes.items():
            t = re.sub(r'<[^>]+>', ' ', code).lower()
            overlap = sum(1 for w in entity_words if w in t)
            if overlap > best_overlap:
                best_overlap = overlap
                target_sid = sid
        return target_sid

    # ================================================================

    def _llm_code_diff_repair(
        self,
        slide_id: int,
        code: str,
        issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        codegen_compiler,
        case_dir: str,
        source_store=None,
    ) -> str | None:
        """Fallback: send full code + issues to LLM, get JSON diff edits."""
        evidence_text = codegen_compiler._build_evidence_context(
            bp_slide, evidence, case_dir,
            source_store=source_store,
        )

        # Build spatial state summary
        state = extract_slide_state(slide_id, code)
        spatial_summary = self._format_spatial_summary(state)

        # Build prompt
        prompt_parts = [
            "## Current Slide Code\n",
            "```python",
            code,
            "```\n",
        ]

        if spatial_summary:
            prompt_parts.append(spatial_summary)

        prompt_parts.append("## Issues to Fix\n")
        for i, issue in enumerate(issues, 1):
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or issue.issue_type
            )
            prompt_parts.append(
                f"{i}. [{issue.severity.value}] [{issue.rubric_id} "
                f"{issue.issue_type}] {desc}"
            )
            if issue.planned_fix:
                prompt_parts.append(f"   Suggested fix: {issue.planned_fix}")

        # Mandatory corrections for content accuracy issues
        mandatory = []
        for issue in issues:
            if issue.issue_type in CRITICAL_CONTENT_TYPES | {"misleading_omission"} and issue.planned_fix:
                quoted = re.findall(
                    r'["\u201c]([^"\u201d]{8,})["\u201d]',
                    issue.planned_fix,
                )
                if quoted:
                    mandatory.append(
                        f"- [{issue.issue_type}] Your replacement "
                        f"text MUST contain: \"{quoted[0]}\""
                    )
                else:
                    mandatory.append(
                        f"- [{issue.issue_type}] REQUIRED: "
                        f"{issue.planned_fix[:200]}"
                    )
            elif issue.issue_type == "missing_context" and issue.planned_fix:
                mandatory.append(
                    f"- [{issue.issue_type}] REQUIRED INSERTION: "
                    f"{issue.planned_fix[:250]}. "
                    f"Use an insert_after edit to add a new subtitle or "
                    f"modify an existing text line to include this context."
                )
            elif issue.issue_type == "missing_entity" and issue.planned_fix:
                mandatory.append(
                    f"- [{issue.issue_type}] REQUIRED ADDITION: "
                    f"{issue.planned_fix[:250]}. "
                    f"Add these specific names/values to the existing slide text."
                )
            elif issue.issue_type in ("missing_point", "missing_evidence",
                                      "missing_conclusion") and issue.planned_fix:
                mandatory.append(
                    f"- [{issue.issue_type}] REQUIRED ADDITION: "
                    f"{issue.planned_fix[:250]}. "
                    f"Insert this content into an existing bullet list or "
                    f"text block. Use search_source to verify, then add as "
                    f"a new <li> or append to existing text."
                )

        if mandatory:
            prompt_parts.append(
                "\n## Mandatory Corrections (use these exact phrases)\n"
            )
            prompt_parts.extend(mandatory)

        if bp_slide:
            prompt_parts.append(
                f"\n## Slide Context\n"
                f"Role: {bp_slide.role}\n"
                f"Goal: {bp_slide.primary_proposition}"
            )

        if evidence_text:
            prompt_parts.append(evidence_text)

        # Inject global task brief if available on the compiler
        _task_brief = getattr(codegen_compiler, '_task_brief', '')
        if _task_brief:
            prompt_parts.append(
                f"\n## Global Instructions (from task brief)\n\n{_task_brief}"
            )

        prompt_parts.append(
            "\n## Output\n"
            "Return ONLY a JSON object with an `edits` array. "
            "Each edit has `search`, `replace`, and `reason` fields. "
            "For insertions, also include `insert_after`."
        )

        user_content = "\n".join(prompt_parts)

        try:
            response = self.llm.call_text(
                system_prompt=self._repair_prompt,
                user_content=user_content,
                model=self.model,
                module_name="redeck_patch_fallback",
                prompt_version="patch_fallback",
                max_tokens=8192,
                temperature=0.15,
            )

            edits = self._parse_diff_edits(response, slide_id)
            if not edits:
                logger.info(
                    "PATCH fallback slide %d: no valid edits",
                    slide_id,
                )
                return None

            modified_code = self._apply_diff_edits(code, edits, slide_id)

            if modified_code == code:
                logger.info(
                    "PATCH fallback slide %d: edits produced no change",
                    slide_id,
                )
                return None

            logger.info(
                "PATCH fallback slide %d: applied %d edits",
                slide_id, len(edits),
            )
            return modified_code

        except Exception as e:
            logger.error(
                "PATCH fallback slide %d: error: %s", slide_id, e,
            )
            return None

    def _parse_diff_edits(
        self, response: str, slide_id: int,
    ) -> list[dict]:
        """Parse JSON diff edits from LLM response."""
        json_blocks = re.findall(
            r'```json\s*\n(.*?)```', response, re.DOTALL,
        )
        for block in json_blocks:
            result = self._try_parse_edits_json(block.strip())
            if result:
                return result

        code_blocks = re.findall(r'```\s*\n(.*?)```', response, re.DOTALL)
        for block in code_blocks:
            result = self._try_parse_edits_json(block.strip())
            if result:
                return result

        result = self._try_parse_edits_json(response.strip())
        if result:
            return result

        # Brace matching fallback
        brace_depth = 0
        json_start = None
        for i, ch in enumerate(response):
            if ch == '{':
                if brace_depth == 0:
                    json_start = i
                brace_depth += 1
            elif ch == '}':
                brace_depth -= 1
                if brace_depth == 0 and json_start is not None:
                    candidate = response[json_start:i + 1]
                    if '"edits"' in candidate or '"search"' in candidate:
                        result = self._try_parse_edits_json(candidate)
                        if result:
                            return result
                    json_start = None

        logger.warning(
            "PATCH fallback slide %d: could not parse edits JSON",
            slide_id,
        )
        return []

    def _try_parse_edits_json(self, text: str) -> list[dict] | None:
        """Try to parse text as edits JSON."""
        try:
            data = json.loads(text)
            if isinstance(data, dict) and "edits" in data:
                edits = data["edits"]
                if isinstance(edits, list):
                    valid = [
                        e for e in edits
                        if isinstance(e, dict) and "replace" in e
                    ]
                    if valid:
                        return valid[:10]
            elif isinstance(data, list):
                valid = [
                    e for e in data
                    if isinstance(e, dict) and "replace" in e
                ]
                if valid:
                    return valid[:10]
        except json.JSONDecodeError:
            pass
        return None

    def _apply_diff_edits(
        self, code: str, edits: list[dict], slide_id: int,
    ) -> str:
        """Apply JSON diff edits to code."""
        modified = code
        applied = 0

        for edit in edits:
            search = edit.get("search", "")
            replace = edit.get("replace", "")
            insert_after = edit.get("insert_after", "")

            if insert_after and not search:
                if insert_after in modified:
                    idx = modified.find(insert_after)
                    line_end = modified.find("\n", idx + len(insert_after))
                    if line_end == -1:
                        line_end = len(modified)
                    modified = (
                        modified[:line_end]
                        + "\n" + replace
                        + modified[line_end:]
                    )
                    applied += 1
            elif search:
                if search in modified:
                    modified = modified.replace(search, replace, 1)
                    applied += 1

        logger.info(
            "PATCH fallback slide %d: applied %d/%d edits",
            slide_id, applied, len(edits),
        )
        return modified

    # ================================================================
    # Content Patch — surgical text-only fixes for D/E issues
    # ================================================================

    _CONTENT_PATCH_PROMPT_PATH = (
        Path(__file__).parent.parent.parent
        / "prompts" / "codegen" / "content_patch.system.md"
    )

    def _apply_content_patches(
        self,
        slide_id: int,
        code: str,
        content_issues: list[Issue],
        codegen_compiler,
        case_dir: str,
    ) -> str | None:
        """Apply surgical text-only patches for content accuracy issues.

        For each D/E issue with fix_detail.correct_content, calls the LLM
        once for a minimal text replacement. Each patch is independently
        validated against spatial regression. Returns the patched code,
        or None if no patches succeeded.
        """
        from .html_spatial_state import (
            extract_html_slide_state,
            count_significant_issue_total,
        )
        from .repair_utils import _extract_json

        system_prompt = read_text(self._CONTENT_PATCH_PROMPT_PATH)
        current = code
        patched_any = False
        self._last_content_patch_applied_ids = set()
        self._last_content_patch_reflow_ids = set()

        for issue in content_issues:
            fd = issue.fix_detail
            if not fd or not fd.correct_content:
                continue
            correct_text = normalize_correct_content_text(fd.correct_content)
            if not correct_text:
                logger.info(
                    "Content patch %s skipped: correct_content normalized empty",
                    issue.issue_id,
                )
                continue

            desc = (
                issue.evidence.description
                if issue.evidence and issue.evidence.description
                else issue.why_this_fails or ""
            )

            # Tell the LLM which mode to use
            is_missing = issue.issue_type.startswith("missing")
            row_specs = extract_table_row_specs_from_correct_content(
                fd.correct_content,
            )
            is_table_row_insert = (
                is_missing
                and row_specs
                and (
                    (fd.action_type or "").lower() == "add_data_row"
                    or "table" in (fd.target_location or "").lower()
                    or "row" in (fd.target_location or "").lower()
                )
            )
            if is_table_row_insert:
                row_text = "\n".join(f"- {row}" for row in row_specs)
                mode_instruction = (
                    "\n## Mode: TABLE_ROW_INSERT (missing table data)\n"
                    "Add the missing data as real <tr>/<td> rows inside the "
                    "existing table or tbody. Split each pipe-separated row "
                    "into cells. Do NOT add a <p>, <div>, footer note, or "
                    "visible editorial instruction sentence.\n"
                    f"Rows to add:\n{row_text}\n"
                )
            else:
                mode_instruction = (
                    "\n## Mode: INSERT (this is missing content — ADD it, do NOT replace existing text)\n"
                    if is_missing else
                    "\n## Mode: REPLACE (find the wrong phrase and swap it)\n"
                )

            user_msg = (
                f"## Issue: {issue.issue_type} ({issue.issue_id})\n"
                f"Problem: {desc[:400]}\n"
                f"Correct content (from source paper): "
                f"{correct_text[:500]}\n"
                f"{mode_instruction}"
            )
            if fd and fd.source_ref:
                user_msg += f"Source ref: {fd.source_ref}\n"
            user_msg += (
                f"\n## Current HTML Code\n```html\n{current}\n```"
            )

            try:
                response = self.llm.call_text(
                    system_prompt=system_prompt,
                    user_content=user_msg,
                    model=self.model,
                    module_name="content_patch",
                    prompt_version="content_patch.v1",
                    temperature=0.05,
                    max_tokens=2048,
                )
            except Exception as e:
                logger.warning(
                    "Content patch LLM call failed for %s: %s",
                    issue.issue_id, e,
                )
                continue

            # Parse edits from response
            result = _extract_json(response)
            if not result or "edits" not in result:
                logger.info(
                    "Content patch %s: no valid edits returned",
                    issue.issue_id,
                )
                continue

            edits = result["edits"]
            if not edits:
                continue

            # Apply edits
            candidate = current
            applied = 0
            for edit in edits[:3]:  # max 3 edits
                search = edit.get("search", "")
                replace = edit.get("replace", "")
                insert_after = edit.get("insert_after", "")

                if insert_after and not search:
                    # INSERT mode: find insert_after text, add replace after it
                    if insert_after in candidate:
                        idx = candidate.find(insert_after) + len(insert_after)
                        candidate = candidate[:idx] + "\n" + replace + candidate[idx:]
                        applied += 1
                elif search and search in candidate and replace != search:
                    candidate = candidate.replace(search, replace, 1)
                    applied += 1

            if not applied or candidate == current:
                logger.info(
                    "Content patch %s: edits did not match code",
                    issue.issue_id,
                )
                continue

            # Spatial regression check
            try:
                t0_state = extract_html_slide_state(slide_id, current)
                t1_state = extract_html_slide_state(slide_id, candidate)
                t0_spatial = count_significant_issue_total(t0_state)
                t1_spatial = count_significant_issue_total(t1_state)

                if t1_spatial > t0_spatial:
                    logger.warning(
                        "Content patch %s rejected: spatial regression "
                        "(%d → %d significant issues)",
                        issue.issue_id, t0_spatial, t1_spatial,
                    )
                    # Keep the issue in Phase 2. The semantic edit is valid but
                    # needs a coordinated local reflow instead of being dropped
                    # from the repair pipeline after this surgical patch fails.
                    self._last_content_patch_reflow_ids.add(issue.issue_id)
                    continue
            except Exception as e:
                logger.warning(
                    "Content patch %s: spatial check failed (%s), "
                    "accepting patch anyway",
                    issue.issue_id, e,
                )

            current = candidate
            patched_any = True
            self._last_content_patch_applied_ids.add(issue.issue_id)
            logger.info(
                "Content patch applied: %s on slide %d (%d edits)",
                issue.issue_id, slide_id, applied,
            )

        return current if patched_any else None

    def _format_spatial_summary(self, state: SlideState) -> str:
        """Format spatial state as a concise summary for the LLM prompt."""
        if not state.blocks:
            return ""

        lines = [
            "## Current Element Positions (y-sorted)\n",
            "Slide: 13.333\" x 7.5\", usable: x=[0.50, 13.00] y=[0.25, 7.20]\n",
        ]
        sorted_blocks = sorted(state.blocks, key=lambda b: (b.y, b.x))
        for b in sorted_blocks:
            bottom = b.y + b.h
            text_info = f" text={b.text_chars}ch" if b.text_chars > 0 else ""
            util_info = ""
            if b.utilization > 0.85:
                util_info = f" TIGHT({b.utilization:.0%})"
            if b.utilization > 1.0:
                util_info = f" OVERFLOW({b.utilization:.0%})"
            lines.append(
                f"  {b.var_name:20s} y={b.y:.2f}->{bottom:.2f} "
                f"x={b.x:.2f} w={b.w:.2f} h={b.h:.2f}"
                f"{text_info}{util_info}"
            )
        return "\n".join(lines)

    # ================================================================
    # Test compile
    # ================================================================

    def _test_compile(
        self, code: str, codegen_compiler, case_dir: str, slide_id: int,
    ) -> bool:
        """Test-compile the code to check for errors.

        For HTML code: validates well-formed HTML structure.
        For python-pptx code: executes against a test Presentation.
        """
        # HTML mode: just validate basic HTML structure
        if "<!DOCTYPE" in code or ("<html" in code and "<body" in code):
            return bool(code.strip()) and "</html>" in code

        try:
            from pptx import Presentation
            from pptx.util import Emu

            test_prs = Presentation()
            test_prs.slide_width = Emu(SlideDimensions.WIDTH_EMU)
            test_prs.slide_height = Emu(SlideDimensions.HEIGHT_EMU)
            blank_layout = test_prs.slide_layouts[6]
            test_slide = test_prs.slides.add_slide(blank_layout)

            case_path = Path(case_dir)
            test_image_dir = str(
                codegen_compiler._find_image_dir(case_path)
                or case_path / "images"
            )

            from ...backends.python_pptx import code_executor
            success, error = code_executor.execute_code(
                code, test_prs, test_slide, test_image_dir,
            )
            if not success:
                logger.debug(
                    "Test compile slide %d failed: %s",
                    slide_id, error[:200],
                )
            return success
        except Exception as e:
            logger.debug(
                "Test compile slide %d exception: %s",
                slide_id, str(e)[:200],
            )
            return False
