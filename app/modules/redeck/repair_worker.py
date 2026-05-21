"""Agentic slide repair with checkpoint verification.

Architecture:

  1. FILTER: Remove containment overlap false positives
  2. DISPATCH: Per-slide AgentRepair (plan → atomic edits → verify_layout loop)
  3. VALIDATE: Compile test + spatial regression gate

  AgentRepair (in agent_repair.py):
    Tool loop with plan/apply_edits/verify_layout/regenerate/submit tools.
    LLM plans fixes, applies edits atomically, verifies layout after
    structural changes, and submits when done. Max tool calls capped.
"""

import json
import logging
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from ...llm_client import LLMClient
from ...schemas.issue import Issue
from ...schemas.blueprint import BlueprintSlide
from ...schemas.evidence import EvidenceState
from ...schemas.issue_types import (
    SPATIAL_ISSUE_TYPES as SPATIAL_ONLY_ISSUE_TYPES,
    UNSOLVABLE_TYPES as UNSOLVABLE_ISSUE_TYPES,
    CROSS_SLIDE_TYPES as TRULY_CROSS_SLIDE_TYPES,
    CRITICAL_CONTENT_TYPES,
    SlideDimensions,
)
from ...utils.io_utils import read_text

from .spatial_state import SlideState
from .agent_repair import AgentRepair, extract_slide_state

logger = logging.getLogger(__name__)

_REPAIR_PROMPT_PATH = (
    Path(__file__).parent.parent.parent
    / "prompts" / "codegen" / "slide_html_repair.system.md"
)


class ReDeckRepairWorker:
    """Agent-based repair with layout-aware feedback.

    All slides go through AgentRepair.repair() which runs an autonomous
    agent loop: the LLM edits code, verifies spatial state, and adjusts
    based on fine-grained layout feedback.

    Falls back to targeted search/replace when agent repair fails entirely.
    """

    def __init__(self, llm: LLMClient, model: str = "gpt-4o",
                 repair_config: dict | None = None):
        self.llm = llm
        self.model = model
        self._repair_prompt = read_text(_REPAIR_PROMPT_PATH)
        self._repair_config = repair_config or {}
        # After repair_slides(), this tracks which slides had content
        # (non-spatial) issues repaired. Slides with ONLY spatial issues
        # don't need C/D/E re-evaluation.
        self.content_modified_slides: set[int] = set()

    def repair_slides(
        self,
        codegen_compiler,
        issues: list[Issue],
        blueprint_slides: list[BlueprintSlide],
        evidence: EvidenceState,
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int = 0,
        source_store=None,
    ) -> list[int]:
        """Repair slides with atomic action pipeline.

        Returns list of slide_ids successfully repaired.
        """
        # Group issues by slide — only open issues.
        # Per-slide issues reported on multiple slides are SPLIT into
        # independent per-slide copies so each slide can be repaired,
        # verified, and accepted independently.
        # Truly cross-slide issues (visual_inconsistency, missing_section)
        # spanning 3+ slides are skipped (deck-level problems).
        slide_issues: dict[int, list[Issue]] = {}
        for issue in issues:
            if issue.status.value != "open":
                continue
            if issue.issue_type in UNSOLVABLE_ISSUE_TYPES:
                continue
            # Skip KEEP issues — judge says they don't need repair
            if issue.recommended_action and issue.recommended_action.value == "KEEP":
                continue
            affected = issue.affected_slides

            # Truly cross-slide issues spanning 3+ slides: skip
            if (
                issue.issue_type in TRULY_CROSS_SLIDE_TYPES
                and len(affected) >= 3
            ):
                logger.debug(
                    "Skipping deck-level issue %s (%s) spanning %d slides",
                    issue.issue_id, issue.issue_type, len(affected),
                )
                continue

            # Per-slide issues on multiple slides: split into per-slide copies
            if (
                len(affected) > 1
                and issue.issue_type not in TRULY_CROSS_SLIDE_TYPES
            ):
                for sid in affected:
                    per_slide = issue.model_copy(
                        update={"affected_slides": [sid]},
                    )
                    slide_issues.setdefault(sid, []).append(per_slide)
                continue

            # Single-slide issue, or 2-slide cross-slide: assign to primary
            if affected:
                slide_issues.setdefault(affected[0], []).append(issue)
            # Also assign to secondary slide for 2-slide cross-slide issues
            if len(affected) == 2:
                slide_issues.setdefault(affected[1], []).append(issue)

        if not slide_issues:
            logger.info("ReDeck: no open issues, skipping repair")
            return []

        # ── P0-a: Per-slide diminishing returns exit ──────────────
        # If a slide had N issues last turn and still has ≥N issues this turn
        # for 2 consecutive turns, mark as repair_exhausted and skip.
        # This prevents the "getting worse with each fix" death spiral on slides like 2/4/6/9/14.
        exhausted_slides: set[int] = set()
        if run_dir and turn_index >= 2:
            history_path = Path(run_dir) / "slide_issue_history.json"
            try:
                history = json.loads(history_path.read_text()) if history_path.exists() else {}
            except Exception:
                history = {}

            for sid, iss_list in list(slide_issues.items()):
                sid_key = str(sid)
                past = history.get(sid_key, [])  # list of issue counts per turn
                if len(past) >= 2:
                    prev2 = past[-2]
                    prev1 = past[-1]
                    current = len(iss_list)
                    # Stagnation: 2 consecutive turns without improvement
                    if current >= prev1 and prev1 >= prev2 and prev2 > 0:
                        exhausted_slides.add(sid)
                        logger.info(
                            "ReDeck: slide %d marked repair_exhausted "
                            "(issue counts: T-%d=%d, T-%d=%d, T-%d=%d — no improvement)",
                            sid, turn_index - 2, prev2, turn_index - 1, prev1,
                            turn_index, current,
                        )

            # Remove exhausted slides from dispatch
            for sid in exhausted_slides:
                del slide_issues[sid]

            if exhausted_slides:
                logger.info(
                    "ReDeck: %d slide(s) skipped (repair_exhausted): %s",
                    len(exhausted_slides), sorted(exhausted_slides),
                )

        if not slide_issues:
            logger.info("ReDeck: all slides exhausted, skipping repair")
            return []

        # Track which slides have content (non-spatial) issues.
        # Slides with ONLY spatial issues don't need C/D/E re-evaluation.
        self.content_modified_slides = set()
        for sid, iss_list in slide_issues.items():
            has_content_issue = any(
                i.issue_type not in SPATIAL_ONLY_ISSUE_TYPES
                for i in iss_list
            )
            if has_content_issue:
                self.content_modified_slides.add(sid)

        logger.info(
            "ReDeck: %d slides with issues (%d have content issues)",
            len(slide_issues), len(self.content_modified_slides),
        )

        # Prepare context maps
        slide_map = {s.slide_id: s for s in blueprint_slides}

        repaired = []

        def _repair_one(sid: int) -> tuple[int, bool]:
            current_code = codegen_compiler.slide_codes.get(sid, "")
            if not current_code:
                return sid, False

            bp_slide = slide_map.get(sid)
            iss_list = slide_issues[sid]

            success = self._repair_one_slide(
                sid, current_code, iss_list, bp_slide, evidence,
                codegen_compiler, case_dir,
                run_dir=run_dir, turn_index=turn_index,
                source_store=source_store,
            )
            return sid, success

        # Parallel execution
        max_workers = min(len(slide_issues), 4)  # parallel slide repairs
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(_repair_one, sid): sid
                for sid in slide_issues
            }
            for future in as_completed(futures):
                sid = futures[future]
                try:
                    sid, success = future.result()
                    if success:
                        repaired.append(sid)
                        logger.info("ReDeck: slide %d repaired", sid)
                except Exception as e:
                    logger.error("ReDeck: slide %d exception: %s", sid, e, exc_info=True)

        logger.info(
            "ReDeck: %d/%d slides repaired",
            len(repaired), len(slide_issues),
        )

        # ── P0-a: Save per-slide issue counts for diminishing returns tracking ──
        if run_dir:
            history_path = Path(run_dir) / "slide_issue_history.json"
            try:
                history = json.loads(history_path.read_text()) if history_path.exists() else {}
            except Exception:
                history = {}

            # Record current turn's issue counts (including exhausted slides)
            # Re-count from original issues list (before exhausted filtering)
            all_slide_counts: dict[int, int] = {}
            for issue in issues:
                if issue.status.value != "open":
                    continue
                for sid in issue.affected_slides:
                    all_slide_counts[sid] = all_slide_counts.get(sid, 0) + 1

            for sid, count in all_slide_counts.items():
                sid_key = str(sid)
                if sid_key not in history:
                    history[sid_key] = []
                history[sid_key].append(count)

            try:
                history_path.write_text(json.dumps(history, indent=2))
            except Exception as e:
                logger.warning("Failed to save slide issue history: %s", e)

        return repaired

    # ================================================================
    # Core repair logic for one slide
    # ================================================================

    def _repair_one_slide(
        self,
        slide_id: int,
        current_code: str,
        issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        codegen_compiler,
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int = 0,
        source_store=None,
    ) -> bool:
        """Repair a single slide using atomic action pipeline.

        1. FILTER: Remove containment overlap FPs + unsolvable issues
        2. DISPATCH: AgentRepair (agent loop with tools + layout feedback)
        3. FALLBACK: targeted patch if agent fails
        4. VALIDATE: Compile + spatial regression gate
        """
        # ── Step 1: Filter ──
        # Keep only open, actionable issues (skip cross-slide unsolvable types)
        actionable_issues = [
            i for i in issues
            if i.status.value == "open"
            and i.issue_type not in UNSOLVABLE_ISSUE_TYPES
        ]

        if not actionable_issues:
            logger.info(
                "ReDeck slide %d: no actionable issues after filtering "
                "(%d total, 0 actionable)",
                slide_id, len(issues),
            )
            return False

        logger.info(
            "ReDeck slide %d: %d actionable issues (from %d total)",
            slide_id, len(actionable_issues), len(issues),
        )

        # ── Step 2: Agent-based repair (tool-calling loop) ──
        repair = AgentRepair(
            self.llm, self.model,
            repair_config=self._repair_config,
        )
        new_code = repair.repair(
            slide_id, current_code, actionable_issues,
            bp_slide, evidence, codegen_compiler, case_dir,
            run_dir=run_dir, turn_index=turn_index,
            source_store=source_store,
        )

        # ── Step 3: Fallback to targeted patch if agent failed ──
        # Only fall back if the agent actually attempted a repair
        # (i.e., had issues to work on). If repair() returned None
        # because all issues were filtered out (minor structural,
        # skip types), do NOT fall back — nothing to fix.
        if not new_code and repair._current_issues:
            logger.info(
                "ReDeck slide %d: agent repair failed, "
                "trying targeted patch fallback",
                slide_id,
            )
            new_code = self._llm_code_diff_repair(
                slide_id, current_code, actionable_issues,
                bp_slide, evidence, codegen_compiler, case_dir,
                source_store=source_store,
            )

        # ── Step 4: Validate ──
        if new_code and new_code != current_code:
            if not self._test_compile(
                new_code, codegen_compiler, case_dir, slide_id,
            ):
                logger.info(
                    "ReDeck slide %d: repair rejected (compile failed)",
                    slide_id,
                )
                return False

            # Spatial regression gate
            # When slide has critical content accuracy issues (fabricated,
            # incorrect_claim), content rewrites necessarily change text
            # lengths and may cause minor spatial shifts. Allow small
            # spatial regression to avoid blocking critical content fixes.
            CRITICAL_CONTENT_TYPES_ = CRITICAL_CONTENT_TYPES
            n_critical_content = sum(
                1 for i in issues
                if i.issue_type in CRITICAL_CONTENT_TYPES_
                and (i.severity.value if hasattr(i.severity, 'value')
                     else str(i.severity)) in ("critical", "major")
            )
            spatial_tolerance = 2 if n_critical_content >= 2 else (1 if n_critical_content >= 1 else 0)

            t0_state = extract_slide_state(slide_id, current_code)
            t1_state = extract_slide_state(slide_id, new_code)
            t0_spatial = (
                len(t0_state.overlap_pairs) + len(t0_state.oob_blocks)
            )
            t1_spatial = (
                len(t1_state.overlap_pairs) + len(t1_state.oob_blocks)
            )

            if t1_spatial > t0_spatial + spatial_tolerance:
                logger.warning(
                    "ReDeck slide %d: rejected (spatial regression "
                    "%d -> %d)",
                    slide_id, t0_spatial, t1_spatial,
                )
                return False

            # Overflow regression gate: reject if repair introduces
            # many new text overflow blocks (blocks that were NOT
            # overflowing in T0 but are now overflowing in T1).
            t0_overflow_set = set(t0_state.overflow_blocks)
            t1_overflow_set = set(t1_state.overflow_blocks)
            new_overflow = t1_overflow_set - t0_overflow_set
            resolved_overflow = t0_overflow_set - t1_overflow_set
            if len(t1_overflow_set) > len(t0_overflow_set):
                logger.warning(
                    "ReDeck slide %d: rejected (overflow regression "
                    "%d→%d, new: %s)",
                    slide_id, len(t0_overflow_set), len(t1_overflow_set),
                    ", ".join(list(new_overflow)[:3]),
                )
                return False
            elif new_overflow:
                logger.info(
                    "ReDeck slide %d: accepting %d new overflow "
                    "(net improvement: %d→%d, resolved: %s)",
                    slide_id, len(new_overflow),
                    len(t0_overflow_set), len(t1_overflow_set),
                    ", ".join(list(resolved_overflow)[:3]),
                )

            # OOB regression gate: reject only if total OOB count increased.
            # Layout restructuring may shift elements so some elements go
            # OOB while others return in-bounds.
            t0_oob_set = set(t0_state.oob_blocks)
            t1_oob_set = set(t1_state.oob_blocks)
            new_oob = t1_oob_set - t0_oob_set
            resolved_oob = t0_oob_set - t1_oob_set
            if len(t1_oob_set) > len(t0_oob_set):
                logger.warning(
                    "ReDeck slide %d: rejected (OOB regression "
                    "%d→%d, new: %s)",
                    slide_id, len(t0_oob_set), len(t1_oob_set),
                    ", ".join(list(new_oob)[:3]),
                )
                return False
            elif new_oob:
                logger.info(
                    "ReDeck slide %d: accepting %d new OOB "
                    "(net improvement: %d→%d, resolved: %s)",
                    slide_id, len(new_oob),
                    len(t0_oob_set), len(t1_oob_set),
                    ", ".join(list(resolved_oob)[:3]),
                )

            # Net-improvement overlap gate: allow new overlap pairs as
            # long as the TOTAL overlap area decreased (layout restructuring
            # may create transient new pairs while resolving worse ones).
            t0_overlap_set = {
                (min(a, b), max(a, b))
                for a, b, _ in t0_state.overlap_pairs
            }
            t0_total_overlap_area = sum(
                area for _, _, area in t0_state.overlap_pairs
            )
            t1_total_overlap_area = sum(
                area for _, _, area in t1_state.overlap_pairs
            )
            new_overlaps = [
                (a, b, area)
                for a, b, area in t1_state.overlap_pairs
                if (min(a, b), max(a, b)) not in t0_overlap_set
            ]
            # Accept if: (a) total overlap area did not increase, OR
            # (b) total overlap count decreased (net improvement).
            # Reject only if both overlap count AND total area got worse.
            t0_overlap_count = len(t0_state.overlap_pairs)
            t1_overlap_count = len(t1_state.overlap_pairs)
            net_count_improved = t1_overlap_count < t0_overlap_count
            net_area_improved = t1_total_overlap_area <= t0_total_overlap_area + 0.1

            if new_overlaps and not net_count_improved and not net_area_improved:
                # Both count and area got worse — genuine regression
                total_new_area = sum(area for _, _, area in new_overlaps)
                # Still tolerate if critical content issues exist and
                # the regression is small
                if n_critical_content >= 1 and total_new_area < 1.0:
                    logger.info(
                        "ReDeck slide %d: tolerating %d new overlaps "
                        "(total %.1f sq in) due to %d critical content issues",
                        slide_id, len(new_overlaps), total_new_area,
                        n_critical_content,
                    )
                else:
                    overlap_desc = ", ".join(
                        f"{a}↔{b}({area:.1f}sq)"
                        for a, b, area in new_overlaps[:3]
                    )
                    logger.warning(
                        "ReDeck slide %d: rejected (overlap regression: "
                        "count %d→%d, area %.1f→%.1f, new: %s)",
                        slide_id, t0_overlap_count, t1_overlap_count,
                        t0_total_overlap_area, t1_total_overlap_area,
                        overlap_desc,
                    )
                    return False
            elif new_overlaps:
                logger.info(
                    "ReDeck slide %d: accepting %d new overlap pairs "
                    "(net improvement: count %d→%d, area %.1f→%.1f)",
                    slide_id, len(new_overlaps),
                    t0_overlap_count, t1_overlap_count,
                    t0_total_overlap_area, t1_total_overlap_area,
                )

            import re

            # Process-note leak gate: reject if the repair introduced
            # meta-reasoning text into slide content (e.g., "Refocus
            # the slide on..." or "This slide should...").
            _PROCESS_NOTE_PATTERNS = [
                r'\.text\s*=\s*[f]?["\'].*(?:Refocus the slide|This slide should|'
                r'Restructure the|Move the content|Rewrite this|'
                r'Compress the text|Reduce the font|'
                r'The goal is to|Plan:|Step \d+:).*["\']',
            ]
            for pat in _PROCESS_NOTE_PATTERNS:
                leaked = re.findall(pat, new_code, re.IGNORECASE)
                if leaked and not re.findall(pat, current_code, re.IGNORECASE):
                    logger.warning(
                        "ReDeck slide %d: rejected (process-note "
                        "text leaked into code: %s)",
                        slide_id, leaked[0][:100],
                    )
                    return False

            # ── Step 5: Post-repair spatial enforcement ──
            # Apply deterministic CSS fixes for any remaining overflow/OOB.
            # This catches rebound cases where content edits re-introduce
            # spatial issues.
            enforced_code = self._enforce_spatial_clean(
                slide_id, new_code, codegen_compiler, case_dir,
            )
            if enforced_code and enforced_code != new_code:
                logger.info(
                    "ReDeck slide %d: post-repair spatial enforcement "
                    "applied CSS fixes",
                    slide_id,
                )
                new_code = enforced_code

            codegen_compiler.slide_codes[slide_id] = new_code
            logger.info(
                "ReDeck slide %d: repair accepted "
                "(spatial %d -> %d, new_overlaps=%d)",
                slide_id, t0_spatial, t1_spatial, len(new_overlaps),
            )
            return True

        return False

    # ================================================================
    # Post-repair spatial enforcement (deterministic CSS fixes)
    # ================================================================

    def _enforce_spatial_clean(
        self, slide_id: int, code: str, codegen_compiler, case_dir: str,
    ) -> str | None:
        """Apply deterministic CSS fixes for remaining overflow/OOB/clipping.

        Returns modified code, or None if no fixes needed or fix failed.
        Iterates up to 3 rounds of detection → fix → re-check.
        """
        from .html_spatial_state import extract_html_slide_state
        import re

        MAX_ITER = 3
        current = code
        for iteration in range(MAX_ITER):
            try:
                st = extract_html_slide_state(slide_id, current)
            except Exception as e:
                logger.warning(
                    "Slide %d: spatial enforcement extraction failed: %s",
                    slide_id, e,
                )
                return current if current != code else None

            # Collect blocks with significant overflow or clipping
            problem_bids = set()
            for bid in st.overflow_blocks:
                blk = next((b for b in st.blocks if b.block_id == bid), None)
                if blk and max(blk.overflow_bottom_px, blk.overflow_right_px) > 5:
                    problem_bids.add(bid)
            for bid in st.clipped_blocks:
                blk = next((b for b in st.blocks if b.block_id == bid), None)
                if blk and blk.clipped_bottom_px > 5:
                    problem_bids.add(bid)
            # OOB blocks
            for bid in st.oob_blocks:
                problem_bids.add(bid)

            if not problem_bids:
                return current if current != code else None

            logger.info(
                "Slide %d: spatial enforcement iter %d — %d problem blocks: %s",
                slide_id, iteration + 1, len(problem_bids),
                ", ".join(sorted(problem_bids)[:5]),
            )

            # Apply CSS fixes: add overflow:hidden + word-break to
            # overflowing containers, reduce font-size for overflow blocks
            modified = current
            for bid in problem_bids:
                blk = next((b for b in st.blocks if b.block_id == bid), None)
                if not blk:
                    continue

                # Strategy: inject overflow safety CSS into the element
                # Find the element by its data-block-id or class
                # Add overflow:hidden as a last resort to prevent visual overflow
                eid = bid  # block_id typically matches element id

                # Try to find and patch inline style for this element
                # Pattern: id="<bid>" ... style="..."
                pat_id = re.compile(
                    rf'(id\s*=\s*["\']){re.escape(eid)}(["\'][^>]*style\s*=\s*["\'])([^"\']*?)(["\'])',
                    re.IGNORECASE,
                )
                m = pat_id.search(modified)
                if m:
                    existing_style = m.group(3)
                    fixes = []
                    if bid in st.overflow_blocks or bid in st.clipped_blocks:
                        if 'overflow' not in existing_style:
                            fixes.append('overflow:hidden')
                        if 'word-break' not in existing_style:
                            fixes.append('word-break:break-word')
                        # Reduce font-size if present
                        fs_match = re.search(r'font-size\s*:\s*(\d+(?:\.\d+)?)\s*px', existing_style)
                        if fs_match:
                            old_fs = float(fs_match.group(1))
                            new_fs = max(12, old_fs * 0.85)
                            if new_fs < old_fs:
                                existing_style = re.sub(
                                    r'font-size\s*:\s*\d+(?:\.\d+)?\s*px',
                                    f'font-size:{new_fs:.1f}px',
                                    existing_style,
                                )
                    if bid in st.oob_blocks:
                        # Clamp position
                        for prop in ('left', 'top'):
                            prop_match = re.search(rf'{prop}\s*:\s*(-?\d+(?:\.\d+)?)\s*px', existing_style)
                            if prop_match:
                                val = float(prop_match.group(1))
                                if val < 0:
                                    existing_style = re.sub(
                                        rf'{prop}\s*:\s*-?\d+(?:\.\d+)?\s*px',
                                        f'{prop}:0px',
                                        existing_style,
                                    )

                    new_style = existing_style
                    if fixes:
                        new_style = existing_style.rstrip(';') + ';' + ';'.join(fixes)
                    if new_style != m.group(3):
                        modified = modified[:m.start(3)] + new_style + modified[m.end(3):]
                else:
                    # Element not found by id+style pattern — try adding
                    # a global <style> block for overflow safety
                    if bid in st.overflow_blocks or bid in st.clipped_blocks:
                        safety_css = (
                            f'\n<style>#{re.escape(eid)} {{ '
                            f'overflow:hidden; word-break:break-word; }}</style>\n'
                        )
                        # Insert before </body> or at end
                        if '</body>' in modified:
                            modified = modified.replace(
                                '</body>', safety_css + '</body>', 1,
                            )
                        else:
                            modified += safety_css

            if modified == current:
                # No CSS changes possible — stop iteration
                return current if current != code else None
            current = modified

            # Verify the fix compiles
            if not self._test_compile(current, codegen_compiler, case_dir, slide_id):
                logger.warning(
                    "Slide %d: spatial enforcement CSS fix broke compile, reverting",
                    slide_id,
                )
                return code  # revert to pre-enforcement

        return current if current != code else None

    # ================================================================
    # Targeted patch fallback (search/replace edits)
    # ================================================================

    def _llm_code_diff_repair(
        self,
        slide_id: int,
        code: str,
        issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        codegen_compiler,
        case_dir: str,
        source_store=None,
    ) -> str | None:
        """Fallback: send full code + issues to LLM, get JSON diff edits."""
        evidence_text = codegen_compiler._build_evidence_context(
            bp_slide, evidence, case_dir,
            source_store=source_store,
        )

        # Build spatial state summary
        state = extract_slide_state(slide_id, code)
        spatial_summary = self._format_spatial_summary(state)

        # Build prompt
        prompt_parts = [
            "## Current Slide Code\n",
            "```python",
            code,
            "```\n",
        ]

        if spatial_summary:
            prompt_parts.append(spatial_summary)

        prompt_parts.append("## Issues to Fix\n")
        for i, issue in enumerate(issues, 1):
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or issue.issue_type
            )
            prompt_parts.append(
                f"{i}. [{issue.severity.value}] [{issue.rubric_id} "
                f"{issue.issue_type}] {desc}"
            )
            if issue.planned_fix:
                prompt_parts.append(f"   Suggested fix: {issue.planned_fix}")

        # Mandatory corrections for content accuracy issues
        mandatory = []
        for issue in issues:
            if issue.issue_type in CRITICAL_CONTENT_TYPES | {"misleading_omission"} and issue.planned_fix:
                quoted = re.findall(
                    r'["\u201c]([^"\u201d]{8,})["\u201d]',
                    issue.planned_fix,
                )
                if quoted:
                    mandatory.append(
                        f"- [{issue.issue_type}] Your replacement "
                        f"text MUST contain: \"{quoted[0]}\""
                    )
                else:
                    mandatory.append(
                        f"- [{issue.issue_type}] REQUIRED: "
                        f"{issue.planned_fix[:200]}"
                    )
            elif issue.issue_type == "missing_context" and issue.planned_fix:
                mandatory.append(
                    f"- [{issue.issue_type}] REQUIRED INSERTION: "
                    f"{issue.planned_fix[:250]}. "
                    f"Use an insert_after edit to add a new subtitle or "
                    f"modify an existing text line to include this context."
                )
            elif issue.issue_type == "missing_entity" and issue.planned_fix:
                mandatory.append(
                    f"- [{issue.issue_type}] REQUIRED ADDITION: "
                    f"{issue.planned_fix[:250]}. "
                    f"Add these specific names/values to the existing slide text."
                )

        if mandatory:
            prompt_parts.append(
                "\n## Mandatory Corrections (use these exact phrases)\n"
            )
            prompt_parts.extend(mandatory)

        if bp_slide:
            prompt_parts.append(
                f"\n## Slide Context\n"
                f"Role: {bp_slide.role}\n"
                f"Goal: {bp_slide.primary_proposition}"
            )

        if evidence_text:
            prompt_parts.append(evidence_text)

        # Inject global task brief if available on the compiler
        _task_brief = getattr(codegen_compiler, '_task_brief', '')
        if _task_brief:
            prompt_parts.append(
                f"\n## Global Instructions (from task brief)\n\n{_task_brief}"
            )

        prompt_parts.append(
            "\n## Output\n"
            "Return ONLY a JSON object with an `edits` array. "
            "Each edit has `search`, `replace`, and `reason` fields. "
            "For insertions, also include `insert_after`."
        )

        user_content = "\n".join(prompt_parts)

        try:
            response = self.llm.call_text(
                system_prompt=self._repair_prompt,
                user_content=user_content,
                model=self.model,
                module_name="redeck_code_diff",
                prompt_version="repair_diff.v1",
                max_tokens=8192,
                temperature=0.15,
            )

            edits = self._parse_diff_edits(response, slide_id)
            if not edits:
                logger.info(
                    "targeted patch slide %d: no valid edits",
                    slide_id,
                )
                return None

            modified_code = self._apply_diff_edits(code, edits, slide_id)

            if modified_code == code:
                logger.info(
                    "targeted patch slide %d: edits produced no change",
                    slide_id,
                )
                return None

            logger.info(
                "targeted patch slide %d: applied %d edits",
                slide_id, len(edits),
            )
            return modified_code

        except Exception as e:
            logger.error(
                "targeted patch slide %d: error: %s", slide_id, e,
            )
            return None

    def _parse_diff_edits(
        self, response: str, slide_id: int,
    ) -> list[dict]:
        """Parse JSON diff edits from LLM response."""
        json_blocks = re.findall(
            r'```json\s*\n(.*?)```', response, re.DOTALL,
        )
        for block in json_blocks:
            result = self._try_parse_edits_json(block.strip())
            if result:
                return result

        code_blocks = re.findall(r'```\s*\n(.*?)```', response, re.DOTALL)
        for block in code_blocks:
            result = self._try_parse_edits_json(block.strip())
            if result:
                return result

        result = self._try_parse_edits_json(response.strip())
        if result:
            return result

        # Brace matching fallback
        brace_depth = 0
        json_start = None
        for i, ch in enumerate(response):
            if ch == '{':
                if brace_depth == 0:
                    json_start = i
                brace_depth += 1
            elif ch == '}':
                brace_depth -= 1
                if brace_depth == 0 and json_start is not None:
                    candidate = response[json_start:i + 1]
                    if '"edits"' in candidate or '"search"' in candidate:
                        result = self._try_parse_edits_json(candidate)
                        if result:
                            return result
                    json_start = None

        logger.warning(
            "targeted patch slide %d: could not parse edits JSON",
            slide_id,
        )
        return []

    def _try_parse_edits_json(self, text: str) -> list[dict] | None:
        """Try to parse text as edits JSON."""
        try:
            data = json.loads(text)
            if isinstance(data, dict) and "edits" in data:
                edits = data["edits"]
                if isinstance(edits, list):
                    valid = [
                        e for e in edits
                        if isinstance(e, dict) and "replace" in e
                    ]
                    if valid:
                        return valid[:10]
            elif isinstance(data, list):
                valid = [
                    e for e in data
                    if isinstance(e, dict) and "replace" in e
                ]
                if valid:
                    return valid[:10]
        except json.JSONDecodeError:
            pass
        return None

    def _apply_diff_edits(
        self, code: str, edits: list[dict], slide_id: int,
    ) -> str:
        """Apply JSON diff edits to code."""
        modified = code
        applied = 0

        for edit in edits:
            search = edit.get("search", "")
            replace = edit.get("replace", "")
            insert_after = edit.get("insert_after", "")

            if insert_after and not search:
                if insert_after in modified:
                    idx = modified.find(insert_after)
                    line_end = modified.find("\n", idx + len(insert_after))
                    if line_end == -1:
                        line_end = len(modified)
                    modified = (
                        modified[:line_end]
                        + "\n" + replace
                        + modified[line_end:]
                    )
                    applied += 1
            elif search:
                if search in modified:
                    modified = modified.replace(search, replace, 1)
                    applied += 1

        logger.info(
            "targeted patch slide %d: applied %d/%d edits",
            slide_id, applied, len(edits),
        )
        return modified

    def _format_spatial_summary(self, state: SlideState) -> str:
        """Format spatial state as a concise summary for the LLM prompt."""
        if not state.blocks:
            return ""

        lines = [
            "## Current Element Positions (y-sorted)\n",
            "Slide: 13.333\" x 7.5\", usable: x=[0.50, 13.00] y=[0.25, 7.20]\n",
        ]
        sorted_blocks = sorted(state.blocks, key=lambda b: (b.y, b.x))
        for b in sorted_blocks:
            bottom = b.y + b.h
            text_info = f" text={b.text_chars}ch" if b.text_chars > 0 else ""
            util_info = ""
            if b.utilization > 0.85:
                util_info = f" TIGHT({b.utilization:.0%})"
            if b.utilization > 1.0:
                util_info = f" OVERFLOW({b.utilization:.0%})"
            lines.append(
                f"  {b.var_name:20s} y={b.y:.2f}->{bottom:.2f} "
                f"x={b.x:.2f} w={b.w:.2f} h={b.h:.2f}"
                f"{text_info}{util_info}"
            )
        return "\n".join(lines)

    # ================================================================
    # Test compile
    # ================================================================

    def _test_compile(
        self, code: str, codegen_compiler, case_dir: str, slide_id: int,
    ) -> bool:
        """Test-compile the code to check for errors.

        For HTML code: validates well-formed HTML structure.
        For generated Python slide code: executes against a test Presentation.
        """
        # HTML mode: just validate basic HTML structure
        if "<!DOCTYPE" in code or ("<html" in code and "<body" in code):
            return bool(code.strip()) and "</html>" in code

        try:
            from pptx import Presentation
            from pptx.util import Emu

            test_prs = Presentation()
            test_prs.slide_width = Emu(SlideDimensions.WIDTH_EMU)
            test_prs.slide_height = Emu(SlideDimensions.HEIGHT_EMU)
            blank_layout = test_prs.slide_layouts[6]
            test_slide = test_prs.slides.add_slide(blank_layout)

            case_path = Path(case_dir)
            test_image_dir = str(
                codegen_compiler._find_image_dir(case_path)
                or case_path / "images"
            )

            from ...backends.html_codegen import code_executor
            success, error = code_executor.execute_code(
                code, test_prs, test_slide, test_image_dir,
            )
            if not success:
                logger.debug(
                    "Test compile slide %d failed: %s",
                    slide_id, error[:200],
                )
            return success
        except Exception as e:
            logger.debug(
                "Test compile slide %d exception: %s",
                slide_id, str(e)[:200],
            )
            return False
