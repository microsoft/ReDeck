"""Agent-Based Slide Repair — Tool-Calling Loop with Layout Feedback.

Architecture:
  The LLM acts as an autonomous agent with access to tools:
  - apply_edits:    search/replace edits on code
  - verify_layout:  compile + extract spatial state for feedback
  - rollback:       undo last N edits
  - delete_shape:   remove a shape by variable name
  - get_current_code: return line-numbered code
  - search_source:  search original paper for claims/numbers
  - lookup_table:   look up tables from source materials
  - submit:         finalize repair

  The agent decides what to edit, when to verify, and when to rollback,
  based on fine-grained spatial feedback. The model gets layout feedback
  DURING editing, not just at the end, allowing it to see overlaps/overflow
  caused by its edits and adjust accordingly.
"""

import json
import difflib
import logging
import os
import re
import xml.etree.ElementTree as ET
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

from ...llm_client import LLMClient
from ...schemas.blueprint import BlueprintSlide
from ...schemas.evidence import EvidenceState
from ...schemas.issue import Issue
from ...schemas.issue_types import (
    ISSUE_TYPE_DEFS, SlideDimensions, SpatialThresholds,
    UNSOLVABLE_TYPES, CRITICAL_CONTENT_TYPES, LAYOUT_REPAIR_TYPES,
)
from ...utils.io_utils import read_text

from .spatial_state import (
    ContentBlock,
    format_spatial_state,
    format_checkpoint_result,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    USABLE_LEFT,
    USABLE_RIGHT,
    USABLE_TOP,
    USABLE_BOTTOM,
)


def extract_slide_state(slide_id: int, code: str):
    """Extract slide state via Playwright HTML rendering."""
    from .html_spatial_state import extract_html_slide_state
    return extract_html_slide_state(slide_id, code)

# Reuse helpers from repair_utils (extracted from legacy multi_action)
from .repair_utils import (
    _extract_json,
    _extract_all_json,
    _has_extra_json,
    _apply_edits,
    _parse_viz_data,
    CONTENT_ACCURACY_ISSUE_TYPES,
    can_exempt_raw_figure_image_crop,
    compute_overflow_px,
    dom_parent_path,
    extract_table_row_specs_from_correct_content,
    html_image_css_crop_hints,
    issues_allow_dominant_element_removal,
    issues_allow_rendered_text_reveal,
    issues_allow_support_copy_compression,
    issues_allow_visible_text_change,
    normalize_correct_content_text,
    validate_repair_not_visual_compression,
    validate_repair_not_visual_downgrade,
    validate_visual_repair_scope,
)

logger = logging.getLogger(__name__)

_REPAIR_PROMPT_PATH = (
    Path(__file__).parent.parent.parent
    / "prompts" / "codegen" / "slide_html_repair.system.md"
)

# Structural issue types that require layout changes (Inches/Pt/RGBColor)
# and may warrant more lenient content retention.
# NOTE: must be a subset of ALL_VALID_TYPES from the registry.
STRUCTURAL_ISSUE_TYPES = {
    "layout_inappropriate", "density_imbalance", "text_visual_imbalance",
    "text_overflow", "overlap", "missing_data_visualization",
    "low_contrast", "alignment_inconsistency",
}

# Visual issues that benefit from coarse whole-slide distribution evidence.
# Other visual issues still receive the neutral per-element LAYOUT ANCHOR, but
# not the SPACE MAP: occupancy cells are useful for composition and balance,
# not for contrast, SVG topology, formatting, or other local visual defects.
COMPOSITIONAL_SPATIAL_ISSUE_TYPES = frozenset({
    "layout_inappropriate",
    "text_visual_imbalance",
    "density_imbalance",
    "alignment_inconsistency",
    # Kept for compatibility with older issue payloads that used these names.
    "whitespace_imbalance",
    "whitespace_asymmetry",
})

# Issues where a clean hard-defect check is not enough. The agent must record
# an issue-level self-assessment against the original visual/compositional
# complaint; deterministic tools provide evidence, not the final visual verdict.
COMPOSITION_CLOSURE_ISSUE_TYPES = COMPOSITIONAL_SPATIAL_ISSUE_TYPES | frozenset({
    "raw_figure",
    "raw_table",
})

# Composition issues are judged visually, but a named composition region still
# owns objective visibility failures inside that region. These measurements do
# not decide whether the composition is attractive; they prevent clipped,
# out-of-canvas, or covered protected content from being mistaken for a pass.
COMPOSITION_TARGET_HARD_SPATIAL_CATEGORIES = frozenset({
    "overlap",
    "occlusion",
    "text_overflow",
    "clipped",
    "canvas_truncation",
    "out_of_bounds",
})

# These issue types can also use conservative same-parent peer hypotheses.
# Keeping this separate from occupancy routing makes the information contract
# explicit even though the current sets intentionally overlap.
RELATIONAL_SPATIAL_ISSUE_TYPES = COMPOSITIONAL_SPATIAL_ISSUE_TYPES

EXACT_CONTENT_ISSUE_TYPES = frozenset({
    "numeric_error",
    "entity_error",
    "missing_entity",
})


def _is_long_content_target(text: str) -> bool:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return False
    sentence_breaks = len(re.findall(r"[.!?](?:\s+|$)", text))
    separators = text.count(";") + text.count(":")
    return len(text) > 170 or (len(text) > 120 and sentence_breaks >= 2) or separators >= 3


def _source_target_label(issue_type: str, content: str, *, limit: int = 320) -> str:
    """Format judge source text as a slide-safe semantic target.

    correct_content is often a source-derived meaning target, not a string that
    should be pasted verbatim into a fixed slide region. Numeric/entity fixes
    still need exact values; longer D/E/C claims should be compressed or merged
    into the closest same-topic visible sentence.
    """
    content = re.sub(r"\s+", " ", content or "").strip()
    if not content:
        return ""
    clipped = content[:limit]
    if issue_type in EXACT_CONTENT_ISSUE_TYPES and not _is_long_content_target(content):
        return f'Exact source value/text: "{clipped}"'
    return (
        f'Source-backed semantic target: "{clipped}" '
        "Cover these facts in compact presentation wording; paraphrase or split "
        "across existing same-topic body text if needed, but do not paste a long "
        "source-style sentence wholesale."
    )


def _count_html_words(html: str) -> int:
    """Count visible words in HTML, stripping style/script/tags."""
    t = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL)
    t = re.sub(r'<script[^>]*>.*?</script>', '', t, flags=re.DOTALL)
    t = re.sub(r'<[^>]+>', ' ', t)
    t = re.sub(r'&[a-zA-Z]+;', ' ', t)
    return len(t.split())


def _fixed_format_text_budget_warnings(
    html: str,
    *,
    baseline_html: str | None = None,
) -> list[str]:
    """Return new title/footer budget violations in HTML slide text.

    This catches a recurring content-repair failure mode where long source
    sentences are pasted into a fixed-format title or footer. DOM overflow can
    miss this when the parent clips text, so the repair loop needs a simple
    text-budget guard before submit.
    """
    try:
        from bs4 import BeautifulSoup
    except Exception:
        return []

    def visible_text(tag) -> str:
        return re.sub(r"\s+", " ", " ".join(tag.stripped_strings)).strip()

    def class_tokens(tag) -> set[str]:
        raw = tag.get("class", []) or []
        if isinstance(raw, str):
            raw = raw.split()
        return {str(item).strip().lower() for item in raw if str(item).strip()}

    def role_for(tag) -> str:
        classes = class_tokens(tag)
        class_text = " ".join(classes)
        name = (tag.name or "").lower()
        if "source" in classes or "credit" in classes:
            return "source note"
        if classes & {"bottom-bar", "bottom_bar", "takeaway", "footer", "footnote"}:
            return "bottom/footer bar"
        if (
            name in {"h1", "h2"}
            or classes & {"title", "slide-title", "page-title", "headline"}
            or class_text in {"slide title", "page title"}
        ):
            return "title/header"
        return ""

    budgets = {
        "title/header": (115, 16),
        "bottom/footer bar": (220, 34),
        "source note": (90, 14),
    }

    def collect(src: str) -> list[tuple[str, str]]:
        soup = BeautifulSoup(src or "", "html.parser")
        findings: list[tuple[str, str]] = []
        for tag in soup.find_all(True):
            if (tag.name or "").lower() in {"style", "script", "svg", "path"}:
                continue
            role = role_for(tag)
            if not role:
                continue
            text = visible_text(tag)
            if not text:
                continue
            max_chars, max_words = budgets[role]
            words = re.findall(r"\b[\w.+-]+\b", text)
            if len(text) <= max_chars and len(words) <= max_words:
                continue
            key = f"{role}:{text.casefold()[:100]}"
            findings.append((
                key,
                f"{role} is over budget ({len(text)} chars, {len(words)} words): "
                f"\"{text[:120]}{'...' if len(text) > 120 else ''}\"",
            ))
        return findings

    baseline_keys = {key for key, _ in collect(baseline_html or "")}
    return [message for key, message in collect(html) if key not in baseline_keys]


def _visible_text_tokens(html: str) -> list:
    """Visible-text tokens (keeps %, ., ×, - inside a token so '94.75%' and
    'Set-to-set' stay whole)."""
    t = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL)
    t = re.sub(r'<script[^>]*>.*?</script>', '', t, flags=re.DOTALL)
    t = re.sub(r'<[^>]+>', ' ', t)
    t = re.sub(r'&[a-zA-Z]+;', ' ', t)
    return re.findall(r"[A-Za-z0-9][A-Za-z0-9.%×+\-/]*", t)


def _dropped_high_value_tokens(original: str, current: str, limit: int = 12) -> list:
    """High-value tokens (numbers/metrics, model/method/dataset names) present in
    the ORIGINAL but now TRULY ABSENT from current (not merely reduced in count).
    Used to name, in verify feedback, exactly which value-bearing content a
    deletion removed — so the agent can apply rule 8a ("never delete numbers or
    method names to save space") to specific items, not a vague word-count delta.
    Substring-checked against the current HTML so reworded/relocated tokens that
    still appear anywhere are NOT falsely reported as lost."""
    _STOP = frozenset("""the a an of to in on for and or but with from by as at is are
        was were be this that these those it its their our your we you they i he she""".split())
    # Deck chrome / section labels the prompt (rule re: structural labels) tells the
    # agent to DELETE — never report these as "high-value content lost", or the
    # feedback would push the agent to restore chrome it correctly pruned (a
    # content-preservation proxy that opposes the repair objective).
    _CHROME = frozenset("""one-slide refs ref fig fig. tbl table sec section slide
        paper overview agenda outline contents appendix summary takeaway takeaways
        motivation background intro introduction conclusion conclusions results
        result method methods approach evaluation experiment experiments discussion
        related abstract architecture scaling stability problem design main core
        benefits aspect example overview""".split())
    num_re = re.compile(r"\d")
    camel_re = re.compile(r"[a-z][A-Z]|[A-Z]{2,}")  # DINOv2, COLMAP, FID
    seen, out = set(), []
    cur_lower = current.lower()
    for tok in _visible_text_tokens(original):
        low = tok.lower().strip(".%×+-/")
        if not low or low in _STOP or low in _CHROME or low in seen:
            continue
        # high-value ONLY when number/metric, CamelCase/ACRONYM, or a hyphenated
        # method term. A bare Capitalised common word is NOT counted: it is almost
        # always a section label or sentence-initial word, and flagging it would
        # tell the agent to restore chrome (opposing repair). Precision > recall
        # here — losing the occasional single-word proper noun is the right trade.
        hv = (num_re.search(tok) or camel_re.search(tok)
              or ("-" in tok and len(tok) > 4))
        if not hv:
            continue
        seen.add(low)
        # Truly absent only if the token does not appear anywhere in current.
        if tok.lower() in cur_lower:
            continue
        # Compound rewording guard: a slash/hyphen compound (e.g. "image/video/3D",
        # "Canny-to-image") is frequently REWORDED rather than dropped — the slide
        # now says "image, video, and 3D tasks". Don't flag the compound as lost
        # when MOST of its meaningful pieces still appear on the slide; that would
        # tell the agent to restore content that is already there (just phrased
        # differently). Numbers are never softened this way — a missing metric is
        # always reported. Genuine partial loss (e.g. "multi-stage" → only "multi"
        # survives, "stage" gone) still trips because <half the pieces remain.
        if not num_re.search(tok) and re.search(r"[/-]", tok):
            pieces = [p for p in re.split(r"[/-]", low) if len(p) > 2]
            if pieces:
                present = sum(1 for p in pieces if p in cur_lower)
                if present * 2 >= len(pieces):  # majority of pieces still on slide
                    continue
        out.append(tok)
        if len(out) >= limit:
            break
    return out


_TEXT_LOSS_IGNORED = frozenset("""style width height color margin padding position
absolute relative hidden source page slide section overview summary takeaway
agenda outline introduction conclusion results result method methods background
redeck frame contract
display content border radius transform opacity important""".split())


def _meaningful_visible_words_lost(original: str, current: str) -> list[str]:
    """Return removed visible-word occurrences, excluding layout chrome."""

    before = Counter(token.lower().strip(".%×+-/") for token in _visible_text_tokens(original))
    after = Counter(token.lower().strip(".%×+-/") for token in _visible_text_tokens(current))
    lost: list[str] = []
    for token, count in before.items():
        if len(token) <= 4 or token in _TEXT_LOSS_IGNORED or not re.search(r"[a-z]", token):
            continue
        lost.extend([token] * max(0, count - after.get(token, 0)))
    return lost



@dataclass
class PlanStep:
    """A single step in the agent's repair plan."""
    text: str
    expected_outcome: str = ""   # What the agent expects after this step
    verify_criterion: str = ""   # What to check in verify_layout output
    status: str = "pending"      # pending | in_progress | done | skipped
    skip_reason: str = ""


@dataclass
class AgentState:
    """Mutable state tracked during the agent loop."""
    original_code: str
    current_code: str
    checkpoints: list[str]  # stack for rollback
    slide_id: int
    codegen_compiler: object
    case_dir: str
    bp_slide: "BlueprintSlide | None" = None  # for regen_slide tool
    # ── Macro-planning state (Route 1 improvement) ──
    has_plan: bool = False
    plan_steps: list[PlanStep] = field(default_factory=list)
    plan_summary: str = ""
    # ── Structured verify response state ──
    last_verify_result: dict | None = None  # structured verify output
    consecutive_verify_failures: int = 0     # for auto-rollback logic
    clean_checkpoint_idx: int = 0            # index into checkpoints of last clean state
    # ── Overlap stagnation tracking (A3 improvement) ──
    verify_overlap_history: list[set] = field(default_factory=list)  # sets of (a, b) overlap pairs
    # ── On-demand source search (evidence query tools) ──
    evidence: "EvidenceState | None" = None  # for search_source/lookup_table tools
    search_calls_used: int = 0               # cap searches to avoid wasting tool budget
    # ── Issue context (for tool-level decisions) ──
    issue_types: set = field(default_factory=set)  # set of issue_type strings for current repair
    # ── Rollback learning context ──
    rollback_history: list[str] = field(default_factory=list)  # brief summaries of failed attempts
    # ── Loop exit tracking ──
    submitted: bool = False                        # True if agent successfully submitted
    last_verified_code: str | None = None          # last code that passed verify_layout without regression
    attempted_code_change: bool = False            # true after any successful edit/regen/reflow attempt
    layout_revision: int = 0                       # increments whenever current_code changes
    last_verify_revision: int = -1                 # layout_revision measured by last verify_layout
    last_verify_stale_reason: str = ""             # why last verify no longer matches current_code
    last_code_read_revision: int = -1               # revision returned by get_current_code
    # ── Best-verified-state tracking (attacks "trajectory luck": the agent could
    #    reach a low-issue state then drift to a worse one before submitting). We
    #    remember the verified code with the FEWEST filtered (SSOT) issues and fall
    #    back to it on submit / loop-timeout instead of the last non-regressing one. ──
    best_verified_code: str | None = None
    best_verified_issues: int | None = None
    # Latest checkpoint that passed the hard content/media/scope guards.  Unlike
    # best_verified_code this may retain detector residuals: visual composition
    # repairs are allowed to trade one noisy DOM finding for a better rendered
    # arrangement when the agent has inspected that exact revision.
    latest_safe_verified_code: str | None = None
    latest_safe_verified_revision: int = -1
    latest_visual_checkpoint_code: str | None = None
    latest_visual_checkpoint_revision: int = -1
    latest_visual_checkpoint_hard_valid: bool = False
    latest_visual_checkpoint_targeted_issues: int | None = None
    last_verify_targeted_residual_counts: dict[str, int] = field(default_factory=dict)
    # ── Regen budget ──
    regen_attempts: int = 0                          # how many times regen_slide was called this session
    # ── Repair summary (generated by agent at end of repair) ──
    repair_summary: dict | None = None
    cumulative_words_lost: int = 0
    text_loss_budget: int = 4
    text_loss_locked: bool = False
    checkpoint_text_loss: list[int] = field(default_factory=lambda: [0])
    checkpoint_labels: list[str] = field(default_factory=lambda: ["original"])
    current_checkpoint_label: str = "original"
    allow_visible_text_change: bool = False
    allow_support_copy_compression: bool = False
    pending_edit_cluster: bool = False
    pending_edit_scopes: list[str] = field(default_factory=list)
    last_edit_scope: tuple[str, ...] = field(default_factory=tuple)
    active_cluster_start_code: str | None = None
    active_cluster_start_text_loss: int = 0
    active_cluster_start_label: str = ""
    last_cluster_start_code: str | None = None
    last_cluster_start_text_loss: int = 0
    last_cluster_start_label: str = ""
    spatial_regression_streak: int = 0
    last_spatial_regression_signature: tuple[str, ...] = field(default_factory=tuple)
    dashboard_verify_history: list[dict] = field(default_factory=list)
    trajectory_extensions: int = 0
    last_trajectory_extension_revision: int = -1
    _last_verify_visual_compression_failed: bool = False
    _last_verify_scope_failed: bool = False
    initial_spatial_state: object | None = None
    _run_dir: str | None = None
    _turn_index: int = 0


class AgentRepair:
    """Agent-based slide repair with tool-calling loop.

    The LLM autonomously plans edits, applies them, verifies spatial
    state, and can rollback on regression — all within a multi-turn
    conversation.
    """

    MAX_TOOL_CALLS_PER_ISSUE = 8   # budget per issue
    MAX_TOOL_CALLS_CAP = 40        # hard cap
    TRAJECTORY_EXTENSION_CALLS = 4
    MAX_NO_PROGRESS = 8   # abort if N turns without code change (raised from 4 to give more retry on parse errors)
    MAX_SEARCH_CALLS = 10  # max search_source/lookup_table calls per repair

    def __init__(self, llm: LLMClient, model: str = "gpt-5.5",
                 repair_config: dict | None = None):
        self.llm = llm
        self.model = model
        self._system_prompt = read_text(_REPAIR_PROMPT_PATH)
        self._last_retention = 0.0
        self._current_issues: list[Issue] = []
        self.last_repair_submitted = False
        self.last_repair_summary: dict | None = None
        self.last_repair_needs_composition_closure = False
        self.last_repair_has_valid_composition_closure = True
        self.last_repair_has_resolved_composition_closure = True
        self.last_repair_targeted_residual_total: int | None = None
        self.last_repair_best_verified_code: str | None = None
        self.last_repair_safe_checkpoint_current = False
        self.last_repair_visual_checkpoint_current = False

        # Configurable features (for ablation study)
        cfg = repair_config or {}
        self.LAYOUT_REPR_MODE = cfg.get(
            "layout_repr_mode", "elements_json",
        )
        self._enable_space_planning = cfg.get(
            "enable_space_planning", True,
        )
        self._enable_redistrib_guide = cfg.get(
            "enable_redistrib_guide", True,
        )
        self._enable_layout_preplan = cfg.get(
            "enable_layout_preplan", False,
        )
        self._enable_render_preview = cfg.get(
            "enable_render_preview", False,  # disabled by default; enable for visual verification
        )
        self._disable_step_render = cfg.get(
            "disable_step_render", False,  # ablation: disable verify_layout per-step feedback
        )
        self._enable_spatial_text_feedback = cfg.get(
            "enable_spatial_text_feedback", True,
        )
        self._enable_macro_planning = cfg.get(
            "enable_macro_planning", True,
        )
        self._text_loss_budget = max(0, int(cfg.get("text_loss_budget", 4)))
        self._max_edits_per_call = max(
            1, int(cfg.get("max_edits_per_call", 40)),
        )
        self._pending_actions: list[dict] = []
        self._last_parse_error_message = ""

    @staticmethod
    def _has_current_verify(state: AgentState) -> bool:
        """Whether verify_layout has measured the current code revision."""
        return (
            state.last_verify_result is not None
            and state.last_verify_revision == state.layout_revision
        )

    @staticmethod
    def _verify_needs_strategy_reconsideration(
        defect_history: list[int],
        signature_history: list[str],
    ) -> bool:
        """Return whether recent verifies show genuine non-convergence.

        Equal counts alone are not stagnation: a multi-step reflow can replace
        one residual with another while the structure is still settling. Ask
        for a causal review only when the same residual identities persist
        across three checks or the targeted count strictly worsens. Persistence
        alone does not prove that the current topology is wrong: an intervening
        edit may have restored hierarchy or worked on another owning region.
        """
        if len(defect_history) < 3 or defect_history[-1] <= 0:
            return False
        same_residuals = (
            len(signature_history) >= 3
            and signature_history[-1]
            == signature_history[-2]
            == signature_history[-3]
        )
        strictly_worsening = (
            defect_history[-3]
            < defect_history[-2]
            < defect_history[-1]
        )
        return same_residuals or strictly_worsening

    def _trajectory_continuation_message(
        self,
        state: AgentState,
        *,
        tool_name: str,
        code_changed: bool,
        tool_calls: int,
        soft_limit: int,
    ) -> str:
        """Offer a small, state-aware extension after late real progress.

        The normal budget remains the planning signal. This only prevents a
        meaningful final edit or verification from ending the trajectory before
        the agent can assess that exact revision and either submit or correct it.
        """
        if (
            tool_calls < soft_limit
            or soft_limit >= self.MAX_TOOL_CALLS_CAP
            or state.submitted
            or state.current_code == state.original_code
            or not (code_changed or tool_name == "verify_layout")
            or state.last_trajectory_extension_revision == state.layout_revision
        ):
            return ""

        has_current_verify = self._has_current_verify(state)
        residual_counts = getattr(
            state, "last_verify_targeted_residual_counts", {}
        ) or {}
        regression_total = int(
            getattr(state, "_last_verify_spatial_regression_total", 0) or 0
        )
        pending_steps = [
            step.text
            for step in getattr(state, "plan_steps", [])
            if step.status in {"pending", "in_progress"}
        ]

        if has_current_verify:
            residual_context = (
                ", ".join(
                    f"{category}={count}"
                    for category, count in sorted(residual_counts.items())
                )
                if residual_counts
                else "no targeted detector residuals recorded"
            )
            measurement = (
                f"The current revision has been measured: {residual_context}; "
                f"new spatial regressions={regression_total}."
            )
        else:
            measurement = (
                "The current revision changed after the last measurement and still "
                "needs verification before its visual/spatial effect is known."
            )

        if pending_steps:
            plan_context = (
                " Active plan work still marked open: "
                + " | ".join(pending_steps[:3])
                + (" | ..." if len(pending_steps) > 3 else "")
                + "."
            )
        else:
            plan_context = " No plan steps are currently marked open."

        cluster_context = (
            " The current edit cluster is explicitly unfinished."
            if state.pending_edit_cluster
            else ""
        )
        return (
            "\n\nTRAJECTORY BUDGET EXTENDED AFTER REAL PROGRESS\n"
            f"{measurement}{plan_context}{cluster_context} The extra calls are "
            "room to assess this exact revision, not an instruction to keep "
            "editing. Choose from the evidence: submit if the original issue is "
            "genuinely closed; continue the same coherent chain if a known dependent "
            "edit remains; or correct/rollback the attempt if it displaced pressure, "
            "damaged hierarchy, or made the topology less credible."
        )

    @staticmethod
    def _invalidate_verify_after_code_change(
        state: AgentState,
        reason: str,
    ) -> None:
        """Mark cached verify_layout feedback stale after current_code changes."""
        state.layout_revision += 1
        state.attempted_code_change = True
        state.last_verify_revision = -1
        state.last_verify_stale_reason = reason
        state.last_verify_result = None
        state._last_verify_text_regression = False
        state._last_verify_text_signal = False
        state._last_verify_text_signal_reason = ""
        state._last_verify_spatial_regression_total = 0
        state._last_verify_targeted_residual_total = 0
        state.last_verify_targeted_residual_counts = {}
        state._last_verify_compact_issues = 0
        state._last_verify_visual_compression_failed = False
        state.pending_edit_cluster = False
        state.pending_edit_scopes = []
        state.last_edit_scope = ()
        state._last_verify_scope_failed = False
        state._last_html_state = None
        state.spatial_regression_streak = 0
        state.last_spatial_regression_signature = ()

    @staticmethod
    def _registry_residual_categories_for_issues(issues: list[Issue]) -> set[str]:
        """Return categories explicitly owned by issue-type registry entries."""
        categories: set[str] = set()
        for issue in issues:
            issue_def = ISSUE_TYPE_DEFS.get(issue.issue_type)
            if issue_def:
                categories.update(issue_def.residual_categories)
        return categories

    @classmethod
    def _residual_categories_for_issues(cls, issues: list[Issue]) -> set[str]:
        """Return target-scoped DOM categories used for completion evidence."""
        categories = cls._registry_residual_categories_for_issues(issues)
        for issue in issues:
            if issue.issue_type in COMPOSITIONAL_SPATIAL_ISSUE_TYPES:
                categories.update(COMPOSITION_TARGET_HARD_SPATIAL_CATEGORIES)
        return categories

    def _targeted_residual_categories(self) -> set[str]:
        return self._residual_categories_for_issues(self._current_issues)

    @staticmethod
    def _issue_target_hints(issue: Issue) -> set[str]:
        """Extract strong selector/text hints from an evaluator repair target."""
        hints: set[str] = set()
        evidence = getattr(issue, "evidence", None)
        fix_detail = getattr(issue, "fix_detail", None)
        for ref in getattr(evidence, "object_refs", []) or []:
            normalized = re.sub(r"\s+", " ", str(ref)).strip().lower()
            if normalized:
                hints.add(normalized)
                # Evaluators commonly name a descendant selector such as
                # ``.hero .big`` while spatial blocks expose only the leaf
                # selector (``.big``). Keep the full reference and its simple
                # selector components so a precise target does not fall back
                # to category-wide residual tracking.
                hints.update(
                    selector.lower()
                    for selector in re.findall(
                        r"(?:[#.][A-Za-z_][\w-]+)", normalized,
                    )
                )

        target_texts = [
            getattr(fix_detail, "target_location", ""),
            getattr(evidence, "description", ""),
            getattr(issue, "planned_fix", ""),
        ]
        for value in target_texts:
            if not isinstance(value, str):
                continue
            for quoted in re.findall(r'["\u201c]([^"\u201d]{4,})["\u201d]', value):
                normalized = re.sub(r"\s+", " ", quoted).strip().lower()
                if normalized:
                    hints.add(normalized)
            for selector in re.findall(r"(?:[#.][A-Za-z_][\w-]+)", value):
                hints.add(selector.lower())
        return hints

    @staticmethod
    def _block_matches_target_hints(state, block_id: str, hints: set[str]) -> bool:
        block = next(
            (candidate for candidate in state.blocks if candidate.block_id == block_id),
            None,
        )
        if block is None:
            return False
        text = re.sub(r"\s+", " ", " ".join(block.text_lines)).strip().lower()
        blob = " ".join(
            str(value).lower() for value in (
                block.block_id,
                block.var_name,
                block.css_selector,
                " ".join(block.css_classes),
                block.dom_path,
                text,
            )
            if value
        )
        for hint in hints:
            if hint.startswith(".") and hint[1:] in {
                str(value).lower() for value in block.css_classes
            }:
                return True
            if hint.startswith("#") and (
                hint in str(block.css_selector or "").lower()
                or hint[1:] in str(block.dom_path or "").lower()
            ):
                return True
            if hint in blob:
                return True
            # Evaluators often quote a full sentence while the spatial block
            # contains a line fragment. Require a meaningful fragment length.
            if len(text) >= 12 and text in hint:
                return True
        return False

    def _targeted_significant_issues(self, baseline_state, current_state) -> dict:
        """Filter residual categories to baseline objects named by the issue."""
        from .html_spatial_state import (
            count_significant_issues,
            stable_block_identity,
            stable_pair_identity,
        )

        current = count_significant_issues(current_state)
        categories = self._targeted_residual_categories()
        if baseline_state is None:
            return {category: current.get(category, []) for category in categories}

        baseline = count_significant_issues(baseline_state)

        def identity(state, entry):
            if isinstance(entry, (list, tuple)) and len(entry) == 2:
                return stable_pair_identity(state, entry[0], entry[1])
            return stable_block_identity(state, entry)

        filtered = {}
        pair_categories = {"overlap", "occlusion"}
        for category in categories:
            relevant_issues = [
                issue for issue in self._current_issues
                if category in self._residual_categories_for_issues([issue])
            ]
            hints = set().union(*(
                self._issue_target_hints(issue) for issue in relevant_issues
            )) if relevant_issues else set()
            target_identities = set()
            if hints:
                candidate_categories = set()
                for issue in relevant_issues:
                    candidate_categories.update(
                        self._residual_categories_for_issues([issue])
                    )
                candidate_categories = {
                    candidate for candidate in candidate_categories
                    if (candidate in pair_categories) == (category in pair_categories)
                }
                for candidate in candidate_categories:
                    for entry in baseline.get(candidate, []):
                        block_ids = (
                            entry if isinstance(entry, (list, tuple)) else [entry]
                        )
                        if any(
                            self._block_matches_target_hints(
                                baseline_state, block_id, hints,
                            )
                            for block_id in block_ids
                        ):
                            target_identities.add(identity(baseline_state, entry))

            # Strong evaluator hints should narrow the target even when the
            # baseline identity cannot be recovered exactly. Keep current
            # blocks that match the hint plus genuinely new defects; broad
            # category-wide tracking is reserved for issues with no target hint.
            if hints and not target_identities:
                baseline_identities = {
                    identity(baseline_state, entry)
                    for entry in baseline.get(category, [])
                }
                filtered[category] = [
                    entry for entry in current.get(category, [])
                    if any(
                        self._block_matches_target_hints(
                            current_state, block_id, hints,
                        )
                        for block_id in (
                            entry if isinstance(entry, (list, tuple)) else [entry]
                        )
                    )
                    or identity(current_state, entry) not in baseline_identities
                ]
                continue
            if not target_identities:
                filtered[category] = current.get(category, [])
                continue
            filtered[category] = [
                entry for entry in current.get(category, [])
                if identity(current_state, entry) in target_identities
                or any(
                    self._block_matches_target_hints(
                        current_state, block_id, hints,
                    )
                    for block_id in (
                        entry if isinstance(entry, (list, tuple)) else [entry]
                    )
                )
            ]
        return filtered

    @staticmethod
    def _format_svg_text_overflow_residual(spatial_state, issue_id: str) -> str:
        from .svg_repair import format_svg_text_overflow_residual
        return format_svg_text_overflow_residual(spatial_state, issue_id)

    @staticmethod
    def _has_visual_issues(issues: list[Issue]) -> bool:
        """Return whether the repair batch contains a B-family visual issue."""
        for issue in issues:
            issue_def = ISSUE_TYPE_DEFS.get(issue.issue_type)
            if (issue_def and issue_def.is_b_series) or (
                issue.rubric_id or ""
            ).startswith("B"):
                return True
        return False

    @staticmethod
    def _needs_spatial_distribution(issues: list[Issue]) -> bool:
        """Return whether coarse occupancy/distribution data is relevant."""
        return any(
            issue.issue_type in COMPOSITIONAL_SPATIAL_ISSUE_TYPES
            for issue in issues
        )

    @staticmethod
    def _composition_closure_issues(issues: list[Issue]) -> list[Issue]:
        """Issues that need explicit issue-level composition self-assessment."""
        return [
            issue for issue in issues
            if issue.issue_type in COMPOSITION_CLOSURE_ISSUE_TYPES
        ]

    @classmethod
    def _needs_composition_closure(cls, issues: list[Issue]) -> bool:
        return bool(cls._composition_closure_issues(issues))

    @classmethod
    def _composition_issue_labels(
        cls, issues: list[Issue], *, limit: int = 8,
    ) -> list[str]:
        labels: list[str] = []
        for issue in cls._composition_closure_issues(issues)[:limit]:
            desc = (
                getattr(getattr(issue, "evidence", None), "description", "")
                or getattr(issue, "why_this_fails", "")
                or ""
            )
            desc = re.sub(r"\s+", " ", str(desc)).strip()
            if desc:
                desc = f" - {desc[:130]}"
            labels.append(
                f"{issue.rubric_id} {issue.issue_type} ({issue.issue_id}){desc}"
            )
        return labels

    @classmethod
    def _build_composition_closure_guidance(cls, issues: list[Issue]) -> str:
        labels = cls._composition_issue_labels(issues)
        if not labels:
            return ""
        return (
            "## Composition Self-Assessment Required\n"
            "These issues are about layout composition, spatial coherence, or "
            "figure inspectability. `verify_layout` and render tools provide "
            "evidence; they do not decide visual quality for you. For each listed "
            "issue, use the original evidence and planned_fix as a starting "
            "hypothesis, choose an appropriate repair strategy, then record your "
            "own assessment from the current revision and the spatial/render "
            "evidence available in this run.\n\n"
            "Applies to:\n"
            + "\n".join(f"- {label}" for label in labels)
            + "\n\nRequired self-assessment questions:\n"
            "- What exact void, side/corner imbalance, anchor mismatch, "
            "visual-weight mismatch, or raw-figure readability problem did the "
            "issue name?\n"
            "- Which repair family did you choose, and why: local resize/reposition, "
            "layout reflow of existing elements, source crop/recomposition, or an "
            "authorized exact-data redraw/source-grounded summary asset? Do not treat "
            "planned_fix as a literal command when the current evidence shows a "
            "different scale of change is needed.\n"
            "- Protect the content and semantic role of the title/header, slide "
            "number, source attribution, footer/takeaway bar, and any ReDeck "
            "frame contract; do not use them as filler for a body void. Their "
            "geometry is not frozen. If body pressure or lower-page voids persist, "
            "inspect the whole-slide budget, including whether title/header wrapping "
            "or an oversized frame track is an upstream constraint. Recalibrating "
            "frame width, wrapping, typography, padding, or occupied space is allowed "
            "when its content and role remain intact.\n"
            "- What existing substantive content, real visual content, or "
            "meaningful structure now uses or balances the originally problematic "
            "area? An empty frame, stretched container, moved caption, source note, "
            "citation, footer, or decorative border is not sufficient by itself. "
            "If those are the main change, treat the result as uncertain or choose "
            "a stronger reflow/adaptation strategy.\n"
            "- Which LAYOUT ANCHOR / RELATION MAP / SPACE MAP evidence and any "
            "available render observations support your assessment? Use measurements as evidence, "
            "not thresholds or automatic pass/fail rules. Do not optimize to a "
            "single proxy such as equal bottom edges, matching card heights, "
            "coverage, or a denser SPACE MAP; the rendered slide must still have "
            "a natural reading path and preserve each element's intended role.\n"
            "- Treat small-font and dense-content notices as prompts to inspect "
            "role-relative readability, not as unresolved verdicts. Support copy, "
            "labels, and repeated annotations may be smaller than focal titles, "
            "values, or conclusions. Check their wrapping, line boxes, contrast, "
            "separation, and repeated rhythm. Do not infer unreadability from one "
            "reported pixel size or total word count alone. A dense repeated-card "
            "or dashboard repair may pass when the original crowding/overflow is "
            "gone, focal and support roles remain distinct, all copy is visibly "
            "contained, and peers scan coherently; do not force another topology "
            "solely to clear generic typography warnings.\n"
            "- For B13/spatial-coherence issues, a shared anchor is not the whole "
            "goal. If the fix inflates table rows/cards, flattens natural hierarchy, "
            "creates a worse rhythm, or pushes auxiliary text close to content, "
            "prefer regrouping/body reflow or mark the result uncertain instead of "
            "calling the alignment resolved.\n"
            "- For image/chart/raw-figure issues, inspect the rendered image interior "
            "when preview is available, not only the outer box. A larger <img> bbox does not by itself prove "
            "that the chart, diagram, labels, or relevant panels became more useful. "
            "For SVG redraws/summaries, check internal labels, annotations, legends, "
            "cards, paths, and endpoints; revise or mark uncertain if text touches "
            "borders, crosses marks, clips, or competes with nearby labels.\n"
            "For quantitative charts/plots, also check fidelity to the original "
            "evidence. A handmade SVG summary that approximates curves, axes, "
            "legends, or tick values is a downgrade when the original chart is "
            "clean. Preserve the original, use a real crop/recomposition, or "
            "regenerate from exact source data; do not redraw a chart just because "
            "the surrounding layout is hard to fit.\n"
            "For raw figures displayed with object-fit: contain, do not use the "
            "outer <img> slot or media frame as completion evidence. If verification "
            "reports letterboxing, compare the rendered image content rect and inspect "
            "the internal labels/marks; a larger empty slot with the same or smaller "
            "chart content is unresolved.\n"
            "- Treat captions, source notes, citations, and footers as auxiliary "
            "labels. If they are the only reason a large blank region feels handled, "
            "say so and decide whether another strategy is more appropriate.\n"
            "- When a wide/shallow figure beside a text rail leaves lower-corner "
            "voids, ordinary vertical enlargement may only create letterbox. "
            "Consider body-only reflows such as a full-width/top evidence band "
            "with existing callouts grouped below, a shorter figure plus compact "
            "lower interpretation band made from existing body callouts, or "
            "regrouping the side rail into balanced blocks while preserving all "
            "visible strings, semantic roles, and a coherent reading path.\n"
            "- If your evidence leaves you unconvinced, choose another repair strategy "
            "or record a specific unresolved_concern. The next judge/render review "
            "will decide whether the issue is actually resolved. A specific "
            "uncertain/unresolved closure is useful traceability, but it is not "
            "a completed repair; if a credible stronger body-only reflow, regroup, "
            "source-preserving recomposition/crop, or faithful exact-data redraw remains, try it "
            "before submitting. If render_preview is unavailable, that alone is not "
            "a failure: a spatially explicit issue may pass when anchors, relations, "
            "space distribution, and hard guards directly address the complaint. "
            "Keep the submit confidence consistent with the verdict and self-assessment: "
            "unresolved concerns, an uncertain verdict, or a merely moderate improvement "
            "should not be labeled high confidence."
        )

    @classmethod
    def _build_composition_closure_verify_reminder(cls, issues: list[Issue]) -> str:
        labels = cls._composition_issue_labels(issues, limit=6)
        if not labels:
            return ""
        return (
            "\nCOMPOSITION SELF-ASSESSMENT REQUIRED:\n"
            "Before update_plan marks a composition step done or before submit, "
            "write an issue-level assessment for each target below: original "
            "failure -> chosen strategy -> current spatial/render evidence -> "
            "verdict or remaining uncertainty. "
            "Do not rely only on 'no hard deterministic defects'.\n"
            + "\n".join(f"  - {label}" for label in labels)
            + "\nFor blank-space issues, name what substantive content, visual content, "
            "or meaningful structure now uses or balances the original lower/side/"
            "corner area. Do not count an empty frame, stretched container, moved "
            "caption, source note, citation, footer, or decorative border as "
            "sufficient evidence by itself. If SPACE MAP / anchors still show the "
            "same void, mark the result uncertain/unresolved or change strategy. "
            "Small-font and dense-content notices are informational. Judge support "
            "copy by its role-relative wrapping, line boxes, contrast, separation, "
            "and repeated rhythm; do not mark an otherwise resolved checkpoint "
            "uncertain solely because support text is smaller than a generic body "
            "guideline or because the slide remains information-rich. "
            "Protect title/header, slide number, source attribution, footer/takeaway "
            "bar, and ReDeck frame content and roles; do not use them as filler for "
            "a body-composition void. Their geometry may still be recalibrated when "
            "the whole-slide evidence shows that frame space is part of the cause. "
            "For alignment issues, "
            "name the peer anchor or gap now shared and explain why the result still "
            "has a natural rhythm. Do not mark B13 done because one edge lines up "
            "if the fix stretched rows/cards, flattened hierarchy, or moved auxiliary "
            "text into visual competition with the body. For image/chart/raw-figure "
            "issues, inspect the image interior, not just the outer bbox. For SVG "
            "redraws/summaries, check internal labels, annotations, paths, and card "
            "bounds for collisions or clipping before claiming success. "
            "For quantitative charts/plots, verify that the repair preserves the "
            "original evidence fidelity: axes, legend, tick meanings, and curve "
            "relationships should remain source-grounded. Do not accept a handmade "
            "SVG approximation as a better chart merely because labels are bigger. "
            "If the image uses object-fit: contain, base your assessment on the "
            "rendered image content rect, not the outer slot/bbox; a larger letterboxed "
            "frame is not a raw-figure readability fix. "
            "If local "
            "resizing does not seem to solve the diagnosed composition, switch to "
            "a layout reflow, source crop/recomposition, or exact-data redraw when "
            "that better preserves the slide's meaning and evidence. If you remain uncertain, "
            "record that uncertainty explicitly and keep working when a credible "
            "stronger strategy remains instead of inventing certainty."
        )

    @staticmethod
    def _compact_section(
        compact: str,
        marker: str,
        stop_markers: tuple[str, ...] = (),
    ) -> str:
        """Extract one top-level section from compact spatial feedback."""
        start = compact.find(marker)
        if start < 0:
            return ""
        if compact[start] == "\n":
            start += 1
        end_candidates = [
            compact.find(stop, start + len(marker))
            for stop in stop_markers
        ]
        end_candidates = [end for end in end_candidates if end >= 0]
        end = min(end_candidates) if end_candidates else len(compact)
        return compact[start:end].strip()

    @staticmethod
    def _extract_low_contrast_findings(compact: str) -> str:
        """Return current deterministic low-contrast targets from compact DOM feedback."""
        lines = compact.splitlines()
        findings: list[str] = []
        i = 0
        while i < len(lines):
            line = lines[i]
            if "LOW CONTRAST" not in line:
                i += 1
                continue
            findings.append(line)
            i += 1
            while i < len(lines):
                continuation = lines[i]
                stripped = continuation.strip()
                if not stripped:
                    break
                if continuation.startswith("❌ ") or continuation.startswith("⚠️ "):
                    break
                if continuation.startswith("📐") or continuation.startswith("RELATION MAP") or continuation.startswith("SPACE MAP"):
                    break
                findings.append(continuation)
                i += 1
            continue
        return "\n".join(findings).strip()

    @staticmethod
    def _extract_svg_text_overflow_findings(compact: str) -> str:
        from .svg_repair import extract_svg_text_overflow_findings
        return extract_svg_text_overflow_findings(compact)

    @classmethod
    def _scope_spatial_context(cls, compact: str, issues: list[Issue]) -> str:
        """Layer spatial feedback by issue ownership and information type.

        Hard DOM-owned issues receive the complete report. Other visual issues
        receive neutral element geometry without unrelated deterministic defect
        lists. Composition/alignment issues additionally receive candidate peer
        relations and the coarse SPACE MAP. Content-only issues receive neither
        visual layer.
        """
        residual_categories = cls._registry_residual_categories_for_issues(issues)
        svg_overflow_findings = ""
        if residual_categories:
            # B20 is mostly VLM-owned, with a deterministic text-fit backstop
            # for both inline SVG labels and referenced SVG assets. Expose the
            # target measurements without turning unrelated DOM findings into
            # repair tasks.
            if residual_categories == {"svg_text_overflow"}:
                svg_overflow_findings = cls._extract_svg_text_overflow_findings(
                    compact,
                )
            else:
                return compact

        # When the assigned issues are purely visual (no residual_categories)
        # but the slide has significant spatial defects, expose the full report
        # so the agent can address them holistically. This covers B09/B20 issues
        # where the visual complaint is caused by underlying spatial pressure.
        n_spatial = compact.count("❌ ")
        if n_spatial >= 5 and not residual_categories:
            return compact

        header = compact.splitlines()[0] if compact else "DOM spatial state collected."
        parts = [
            header,
            "DOM findings are regression guards for this issue, not additional "
            "repair targets. Preserve or improve the baseline; resolve the listed "
            "issue using its evidence and render.",
        ]

        if cls._has_visual_issues(issues):
            if svg_overflow_findings:
                parts.extend([
                    "DETERMINISTIC SVG TEXT-FIT TARGETS: these are current "
                    "measurements for the B20 target. Use the vertical and "
                    "horizontal overflow values to choose the correct geometry "
                    "change; unrelated DOM findings remain regression context.",
                    svg_overflow_findings,
                ])

            if any(issue.issue_type == "low_contrast" for issue in issues):
                contrast_findings = cls._extract_low_contrast_findings(compact)
                if contrast_findings:
                    parts.extend([
                        "DETERMINISTIC LOW-CONTRAST TARGETS: these are the "
                        "current rendered text blocks failing WCAG contrast. "
                        "Prior issue wording may be stale after earlier repairs; "
                        "fix the blocks listed here first.",
                        contrast_findings,
                    ])
                else:
                    parts.append(
                        "DETERMINISTIC LOW-CONTRAST TARGETS: no current rendered "
                        "text block fails the contrast checker. Treat any old "
                        "low-contrast wording as likely stale unless it is visibly "
                        "obvious in the render."
                    )

            layout_anchor = cls._compact_section(
                compact,
                "\n📐 LAYOUT ANCHOR",
                ("\nRELATION MAP", "\nSPACE MAP"),
            )
            if layout_anchor:
                parts.extend([
                    "The LAYOUT ANCHOR below is neutral geometry for locating "
                    "the issue target; it is not a list of defects.",
                    layout_anchor,
                ])

        if any(
            issue.issue_type in RELATIONAL_SPATIAL_ISSUE_TYPES
            for issue in issues
        ):
            relation_map = cls._compact_section(
                compact,
                "\nRELATION MAP",
                ("\nSPACE MAP",),
            )
            if relation_map:
                parts.extend([
                    "The RELATION MAP contains conservative same-parent peer "
                    "hypotheses and raw spread measurements. Use only a group "
                    "that matches the issue's named target; it is not a defect "
                    "verdict.",
                    relation_map,
                ])

        if any(issue.issue_type == "alignment_inconsistency" for issue in issues):
            parts.append(
                "ALIGNMENT REVIEW: use the issue evidence to identify logical "
                "peers, then compare the matching RELATION MAP group and its "
                "LAYOUT ANCHOR values. Relation groups are hypotheses, not "
                "verdicts. Do not align unrelated elements, different hierarchy "
                "levels, or intentional asymmetry. Do not treat equalized edges "
                "as success if the repair creates stretched rows/cards, awkward "
                "rhythm, or auxiliary text competing with body content."
            )

        if cls._needs_spatial_distribution(issues):
            space_map = cls._compact_section(compact, "\nSPACE MAP")
            if space_map:
                parts.extend([
                    "The SPACE MAP is coarse distribution evidence only; it cannot "
                    "prove edge alignment or spacing regularity.",
                    space_map,
                ])

        return "\n\n".join(parts)

    @staticmethod
    def _alignment_issue_context(issue: Issue) -> str:
        """Collect the judge text that identifies an alignment target."""
        return " ".join(filter(None, (
            issue.evidence.description,
            issue.planned_fix,
            issue.why_this_fails,
            issue.fix_detail.target_location,
        )))

    @staticmethod
    def _alignment_metric_for_text(text: str) -> str | None:
        """Map an alignment issue description to one measured peer metric."""
        normalized = text.lower()
        if re.search(r"whitespace|empty space|internal slack|留白|空白", normalized):
            if re.search(
                r"(?:whitespace|empty space|internal slack).{0,24}(?:right|horizontal)"
                r"|(?:right|horizontal).{0,24}(?:whitespace|empty space|internal slack)"
                r"|右侧留白|右侧空白|横向留白|横向空白",
                normalized,
            ):
                return "internal_right_slack_spread_px"
            if re.search(r"left|top|左侧|顶部|上方", normalized):
                return None
            return "internal_bottom_slack_spread_px"

        metric_patterns = (
            ("gap_spread_px", r"\bgaps?\b|spacing|rhythm|间距|间隔"),
            ("width_spread_px", r"\bwidths?\b|equal[- ]width|等宽|宽度"),
            ("height_spread_px", r"\bheights?\b|equal[- ]height|等高|高度"),
            (
                "bottom_spread_px",
                r"bottom[- ]edge|bottom alignment|align(?:ed)? at the bottom|底边|下沿|底部对齐",
            ),
            (
                "top_spread_px",
                r"top[- ]edge|top alignment|align(?:ed)? at the top|顶边|上沿|顶部对齐",
            ),
            (
                "left_spread_px",
                r"left[- ]edge|left alignment|left[- ]aligned|左边缘|左对齐",
            ),
            (
                "right_spread_px",
                r"right[- ]edge|right alignment|right[- ]aligned|右边缘|右对齐",
            ),
        )
        for metric, pattern in metric_patterns:
            if re.search(pattern, normalized):
                return metric
        return None

    @staticmethod
    def _alignment_target_tokens(text: str) -> set[str]:
        """Return non-generic tokens suitable for matching named peer groups."""
        generic = {
            "align", "aligned", "alignment", "bottom", "top", "left", "right",
            "edge", "edges", "gap", "gaps", "spacing", "width", "widths",
            "height", "heights", "uneven", "inconsistent", "mismatch", "group",
            "groups", "element", "elements", "item", "items", "peer", "peers",
            "slide", "within", "between", "same", "make", "fix", "adjust",
            "visual", "logical", "position", "positions", "column", "columns",
            "row", "rows", "布局", "对齐", "间距", "边缘", "元素", "同级",
        }
        tokens = set(re.findall(r"[a-z0-9_-]{3,}|[\u4e00-\u9fff]{2,}", text.lower()))
        normalized = set()
        for token in tokens - generic:
            normalized.add(token)
            if token.endswith("s") and len(token) > 4:
                normalized.add(token[:-1])
        return normalized

    @classmethod
    def _format_alignment_relation_delta(
        cls,
        baseline_state,
        current_state,
        issues: list[Issue],
    ) -> str:
        """Compare only the peer relation named by each B13 issue.

        This is target-aware measurement feedback. It intentionally does not
        convert geometric improvement into a resolved/persisted verdict.
        """
        alignment_issues = [
            issue for issue in issues
            if issue.issue_type == "alignment_inconsistency"
        ]
        if not alignment_issues or baseline_state is None or current_state is None:
            return ""

        from .html_spatial_state import infer_layout_relations

        baseline_relations = infer_layout_relations(baseline_state)
        current_relations = infer_layout_relations(current_state)
        metric_labels = {
            "bottom_spread_px": "bottom-edge spread",
            "top_spread_px": "top-edge spread",
            "left_spread_px": "left-edge spread",
            "right_spread_px": "right-edge spread",
            "gap_spread_px": "gap spread",
            "width_spread_px": "width spread",
            "height_spread_px": "height spread",
            "internal_bottom_slack_spread_px": "internal bottom-slack spread",
            "internal_right_slack_spread_px": "internal right-slack spread",
        }
        lines = [
            "ALIGNMENT RELATION DELTA (measurement only; not a resolution verdict):"
        ]

        for issue in alignment_issues:
            context = cls._alignment_issue_context(issue)
            metric = cls._alignment_metric_for_text(context)
            tokens = cls._alignment_target_tokens(context)
            if metric is None:
                lines.append(
                    f"  {issue.issue_id}: no specific edge/gap/size/slack metric "
                    "could be mapped from the issue evidence. Inspect the named "
                    "peers in the render; do not infer success from global geometry."
                )
                continue

            candidates = []
            for relation in baseline_relations:
                searchable = " ".join(
                    relation.get("labels", [])
                    + relation.get("member_keys", [])
                    + [relation.get("parent", ""), relation.get("basis", "")]
                ).lower()
                searchable_tokens = set(re.findall(
                    r"[a-z0-9_-]{3,}|[\u4e00-\u9fff]{2,}",
                    searchable,
                ))
                hits = tokens & searchable_tokens
                if not hits:
                    continue
                score = len(hits) * 10
                if relation.get("confidence") == "high":
                    score += 2
                if metric in {"top_spread_px", "bottom_spread_px", "height_spread_px"}:
                    score += int(relation.get("orientation") == "row")
                elif metric in {"left_spread_px", "right_spread_px", "width_spread_px"}:
                    score += int(relation.get("orientation") == "column")
                candidates.append((score, relation))

            if not candidates:
                lines.append(
                    f"  {issue.issue_id}: no candidate peer group matched the "
                    "issue's named target. Do not use a global relation or whole-slide "
                    "average to claim the alignment repair succeeded."
                )
                continue

            baseline_relation = max(candidates, key=lambda item: item[0])[1]
            baseline_keys = set(baseline_relation.get("member_keys", []))
            current_matches = []
            for relation in current_relations:
                current_keys = set(relation.get("member_keys", []))
                if (
                    current_keys == baseline_keys
                    and relation.get("orientation") == baseline_relation.get("orientation")
                    and relation.get("role_key") == baseline_relation.get("role_key")
                ):
                    current_matches.append(relation)
            if not current_matches:
                lines.append(
                    f"  {issue.issue_id}: matched baseline "
                    f"{baseline_relation['relation_id']}, but no current peer group "
                    "retains the complete member set, role, and orientation. Verify "
                    "the target visually; splitting, moving, or dropping peers is not "
                    "proof of alignment repair."
                )
                continue

            current_relation = current_matches[0]
            if metric == "gap_spread_px" and (
                len(baseline_relation.get("gaps_px", [])) < 2
                or len(current_relation.get("gaps_px", [])) < 2
            ):
                lines.append(
                    f"  {issue.issue_id}: the matched group has fewer than three "
                    "peers, so gap regularity cannot be measured as a spread."
                )
                continue
            if metric.startswith("internal_"):
                baseline_slack = [
                    value for value in baseline_relation.get("internal_slack_px", [])
                    if value is not None
                ]
                current_slack = [
                    value for value in current_relation.get("internal_slack_px", [])
                    if value is not None
                ]
                if len(baseline_slack) < 2 or len(current_slack) < 2:
                    lines.append(
                        f"  {issue.issue_id}: the matched peer containers do not "
                        "have enough measurable descendants for internal-slack comparison."
                    )
                    continue

            before = baseline_relation["metrics"][metric]
            after = current_relation["metrics"][metric]
            if after < before:
                direction = "improved"
            elif after > before:
                direction = "worsened"
            else:
                direction = "unchanged"
            labels = ", ".join(baseline_relation.get("labels", [])[:3])
            lines.append(
                f"  {issue.issue_id}: {metric_labels[metric]} for matched peers "
                f"[{labels}] {before}px -> {after}px ({direction})."
            )

        lines.append(
            "  Lower spread supports regularity only for the matched logical peers. "
            "Use the rendered slide and original issue evidence to decide whether "
            "the visual problem is actually resolved."
        )
        return "\n".join(lines)

    @staticmethod
    def _edit_may_affect_svg(action: dict, code: str) -> bool:
        """Return whether an apply_edits action can change rendered SVG pixels."""
        svg_fragments = re.findall(
            r"<svg\b[^>]*>.*?</svg\s*>", code,
            flags=re.IGNORECASE | re.DOTALL,
        )
        if not svg_fragments:
            return False
        svg_source = "\n".join(svg_fragments)

        for edit in action.get("edits", []):
            search = str(edit.get("search", ""))
            replacement = str(edit.get("replace", ""))
            insert_after = str(edit.get("insert_after", ""))

            if search and search in svg_source:
                return True
            if insert_after and insert_after in svg_source:
                return True
            if re.search(
                r"</?(?:svg|g|path|line|polyline|polygon|rect|circle|"
                r"ellipse|marker|defs|text|tspan|foreignObject)\b",
                replacement,
                flags=re.IGNORECASE,
            ):
                return True

            # A CSS custom property declared outside the SVG can still control
            # its fills, strokes, or typography through var(--name).
            custom_properties = set(re.findall(r"--[A-Za-z0-9_-]+", search))
            custom_properties.update(
                re.findall(r"--[A-Za-z0-9_-]+", replacement)
            )
            if any(
                re.search(rf"var\(\s*{re.escape(name)}\s*[,)]", svg_source)
                for name in custom_properties
            ):
                return True

        return False

    # ================================================================
    # MAIN ENTRY
    # ================================================================

    def repair(
        self,
        slide_id: int,
        code: str,
        all_issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        codegen_compiler,
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int = 0,
        source_store=None,
        attempt: int = 0,
    ) -> str | None:
        """Run agent loop to repair a slide. Returns repaired code or None."""
        self.last_repair_submitted = False
        self.last_repair_summary = None
        self.last_repair_needs_composition_closure = False
        self.last_repair_has_valid_composition_closure = True
        self.last_repair_has_resolved_composition_closure = True
        self.last_repair_targeted_residual_total = None
        self.last_repair_best_verified_code = None
        self.last_repair_safe_checkpoint_current = False
        self.last_repair_visual_checkpoint_current = False

        # 0. Filter issues: remove cross-slide issues where this slide
        # is not the primary (first) affected slide.  Attempting to fix
        # cross-slide issues from a secondary slide's perspective causes
        # destructive changes (e.g. changing colors to match another
        # slide's theme without changing the background).
        # Also skip issue types that are known to cause more harm than
        # good when attempted (cross-slide styling, flow issues).
        SKIP_ISSUE_TYPES = UNSOLVABLE_TYPES
        filtered_issues = []
        for issue in all_issues:
            # Skip issue types that reliably cause regressions
            if issue.issue_type in SKIP_ISSUE_TYPES:
                logger.debug(
                    "Slide %d: skipping issue %s (%s) — issue type "
                    "is in SKIP_ISSUE_TYPES",
                    slide_id, issue.issue_id, issue.issue_type,
                )
                continue
            affected = issue.affected_slides
            if len(affected) > 1 and affected[0] != slide_id:
                logger.debug(
                    "Slide %d: skipping issue %s (%s) — this slide is "
                    "secondary in cross-slide issue (primary=%d)",
                    slide_id, issue.issue_id, issue.issue_type,
                    affected[0],
                )
                continue
            filtered_issues.append(issue)

        if not filtered_issues:
            logger.info(
                "Slide %d: no issues after cross-slide filtering "
                "(had %d before)", slide_id, len(all_issues),
            )
            return None

        if len(filtered_issues) < len(all_issues):
            n_skipped_type = sum(
                1 for i in all_issues
                if i.issue_type in SKIP_ISSUE_TYPES
            )
            n_skipped_cross = len(all_issues) - len(filtered_issues) - n_skipped_type
            logger.info(
                "Slide %d: filtered %d → %d issues "
                "(removed %d skip-type, %d cross-slide)",
                slide_id, len(all_issues), len(filtered_issues),
                n_skipped_type, max(0, n_skipped_cross),
            )
        all_issues = filtered_issues

        # Filter: allow B* visual issues plus content issues with an explicit
        # repair contract. Missing-point/evidence fixes may require local reflow
        # after the surgical content patch was spatially rejected.
        pre_filter_count = len(all_issues)
        all_issues = [
            i for i in all_issues
            if i.rubric_id.startswith("B")
            or i.issue_type in CONTENT_ACCURACY_ISSUE_TYPES
        ]
        if len(all_issues) < pre_filter_count:
            logger.info(
                "Slide %d: B+content filter: %d → %d issues "
                "(removed %d non-actionable issues)",
                slide_id, pre_filter_count, len(all_issues),
                pre_filter_count - len(all_issues),
            )

        # Suppress B16 (text_wall) when B4 (text_overflow) is already present.
        # Both issues target content density, but their fixes conflict:
        # B4 says "shrink CSS to fit", B16 says "delete content to reduce words".
        # When both exist, B4 is the actionable one (CSS-first). B16's word-count
        # pressure causes the agent to delete rows/bullets, creating empty space.
        has_overflow = any(
            i.issue_type in ("text_overflow", "content_overflow")
            for i in all_issues
        )
        if has_overflow:
            pre_suppress = len(all_issues)
            all_issues = [
                i for i in all_issues
                if i.issue_type != "text_wall"
            ]
            n_suppressed = pre_suppress - len(all_issues)
            if n_suppressed:
                logger.info(
                    "Slide %d: suppressed %d text_wall issues "
                    "(overflow already present, CSS-first strategy)",
                    slide_id, n_suppressed,
                )

        logger.info(
            "Slide %d: final issue count after filtering: %d",
            slide_id, len(all_issues),
        )

        if not all_issues:
            logger.info(
                "Slide %d: no issues after structural filtering "
                "(had %d before filtering)",
                slide_id, len(filtered_issues),
            )
            return None

        # 1. Prepare context — use bundle if available, else full evidence
        evidence_text = ""
        if source_store is not None:
            bundle = source_store.get_bundle(slide_id)
            if bundle and bundle.source_text:
                ev_parts = [bundle.source_text]
                for s in bundle.asset_summaries:
                    ev_parts.append(s)
                for s in bundle.table_summaries:
                    ev_parts.append(s)
                evidence_text = "\n\n".join(ev_parts)
        if not evidence_text:
            evidence_text = self._build_full_evidence_context(
                bp_slide, evidence, case_dir,
            )

        # Augment evidence: for fabricated/numeric_error issues, search
        # the full paper for the questioned claims.  This prevents the
        # agent from deleting content that is actually in the paper but
        # was not included in the per-slide evidence window.
        paper_search = self._search_paper_for_claims(
            all_issues, case_dir,
        )
        if paper_search:
            evidence_text += paper_search

        # ── Token optimization: skip evidence for spatial-only slides ──
        # If ALL issues are purely structural/layout (no content accuracy),
        # the agent doesn't need source material — truncate to save tokens.
        has_content_issue = any(
            i.issue_type in CONTENT_ACCURACY_ISSUE_TYPES
            or i.issue_type in CRITICAL_CONTENT_TYPES
            or i.issue_type in ("missing_entity", "unfaithful_compression",
                                "title_content_mismatch", "weak_closing",
                                "missing_section")
            for i in all_issues
        )
        if not has_content_issue and evidence_text:
            logger.info(
                "Slide %d: all %d issues are spatial/layout — "
                "truncating evidence text from %d to 0 chars",
                slide_id, len(all_issues), len(evidence_text),
            )
            evidence_text = ""

        # Extract content requirements
        must_contain, must_not = self._extract_content_requirements(
            code, all_issues,
        )

        # Debug log content requirements
        if must_not:
            logger.info(
                "Slide %d MUST NOT (%d entries): %s",
                slide_id, len(must_not),
                "; ".join(m[:80] for m in must_not[:5]),
            )
        if must_contain:
            logger.info(
                "Slide %d MUST CONTAIN (%d entries): %s",
                slide_id, len(must_contain),
                "; ".join(m[:80] for m in must_contain[:5]),
            )

        # Build content checklist
        content_checklist = self._build_content_checklist(all_issues)

        # Spatial state of original code
        if self._is_html_code(code):
            from .html_spatial_state import extract_html_slide_state, format_html_compact_state
            spatial_state = extract_html_slide_state(
                slide_id,
                code,
                html_base_dir=Path(case_dir),
                asset_base_dirs=self._html_asset_base_dirs(
                    case_dir, run_dir, turn_index,
                ),
            )
            spatial_info = self._scope_spatial_context(
                format_html_compact_state(spatial_state), all_issues,
            )
            # Provide overflow/leverage info when there are significant spatial issues
            # (not just for table-dashboard patterns — absolute-layout slides need it too)
            from .html_spatial_state import count_significant_issues as _count_sig
            _sig_issues = _count_sig(spatial_state)
            _sig_total = sum(len(v) for v in _sig_issues.values())
            _has_spatial_pressure = (
                self._looks_like_table_dashboard_pressure_from(
                    code, {issue.issue_type for issue in all_issues},
                )
                or _sig_total >= 5
            )
            if _has_spatial_pressure:
                spatial_info = self._dashboard_measurement_context(spatial_info)
                allocation_map = self._dashboard_allocation_map(spatial_state)
                if allocation_map:
                    spatial_info += "\n\n" + allocation_map
        else:
            spatial_state = extract_slide_state(slide_id, code)
            spatial_info = format_spatial_state(spatial_state)

        # Adjacent slide context for cross-slide issues
        cross_slide_issues = [
            i for i in all_issues
            if len(i.affected_slides) > 1
        ]
        adjacent_context = self._build_adjacent_context(
            slide_id, cross_slide_issues, codegen_compiler,
        )

        # Synthesize viz data if needed
        viz_data = None
        for issue in all_issues:
            if issue.issue_type == "missing_data_visualization":
                viz_data = self._synthesize_viz(
                    slide_id, all_issues, bp_slide, evidence, evidence_text,
                )
                break

        # Layout pre-planning for complex layout issues
        layout_plan = None
        if self._enable_layout_preplan:
            complex_layout_types = {
                "density_imbalance", "whitespace_imbalance",
                "layout_inappropriate",
            }
            has_complex_layout = (
                any(i.issue_type in complex_layout_types
                    for i in all_issues)
                and spatial_state
                and len(spatial_state.blocks) >= 5
            )
            if has_complex_layout:
                layout_plan = self._generate_layout_plan(
                    code, all_issues, spatial_state, bp_slide,
                )
                if layout_plan:
                    logger.info(
                        "Slide %d: generated layout plan with %d targets",
                        slide_id, len(layout_plan),
                    )

        # 2. Render T0 slide image (if render_preview enabled)
        t0_render_b64 = None
        if self._enable_render_preview:
            dummy_state = AgentState(
                original_code=code,
                current_code=code,
                checkpoints=[code],
                slide_id=slide_id,
                codegen_compiler=codegen_compiler,
                case_dir=case_dir,
            )
            t0_render_b64 = self._render_slide_to_base64(code, dummy_state)
            if t0_render_b64:
                logger.info("Slide %d: T0 render captured for visual context", slide_id)

        # 3. Build initial user message
        initial_msg = self._build_initial_message(
            code, all_issues, spatial_info, evidence_text,
            must_contain, must_not, content_checklist,
            bp_slide, viz_data, adjacent_context,
            spatial_state=spatial_state,
            layout_plan=layout_plan,
            t0_render_b64=t0_render_b64,
            evidence=evidence,
            task_brief=getattr(codegen_compiler, '_task_brief', ''),
        )

        # 3b. Load previous repair failures as bad-case context (turn > 0)
        prev_failures_ctx = ""
        if turn_index > 0 and run_dir:
            prev_failures_ctx = self._load_previous_repair_failures(
                run_dir, turn_index, slide_id,
            )

        # 4. Run single repair attempt
        return self._run_single_repair(
            slide_id=slide_id,
            code=code,
            all_issues=all_issues,
            must_not=must_not,
            must_contain=must_contain,
            initial_msg=initial_msg,
            state_template=dict(
                original_code=code,
                codegen_compiler=codegen_compiler,
                case_dir=case_dir,
                evidence=evidence,
                bp_slide=bp_slide,
                initial_spatial_state=spatial_state,
            ),
            attempt=attempt,
            prev_failures_ctx=prev_failures_ctx,
            run_dir=run_dir,
            turn_index=turn_index,
        )

    def _record_last_repair_result(
        self, state: AgentState, issues: list[Issue],
    ) -> None:
        """Expose final loop state to dispatcher-level acceptance gates."""
        self.last_repair_submitted = bool(state.submitted)
        self.last_repair_summary = state.repair_summary
        self.last_repair_needs_composition_closure = (
            self._needs_composition_closure(issues)
        )
        self.last_repair_has_valid_composition_closure = (
            not self.last_repair_needs_composition_closure
            or self._summary_has_composition_closure(state, issues)
        )
        self.last_repair_has_resolved_composition_closure = (
            not self.last_repair_needs_composition_closure
            or self._summary_has_resolved_composition_closure(state, issues)
        )
        residual_total = getattr(
            state, "_last_verify_targeted_residual_total", None,
        )
        if (
            state.best_verified_code is not None
            and state.current_code == state.best_verified_code
            and state.best_verified_issues is not None
        ):
            residual_total = state.best_verified_issues
        self.last_repair_targeted_residual_total = (
            int(residual_total) if residual_total is not None else None
        )
        self.last_repair_best_verified_code = state.best_verified_code
        self.last_repair_safe_checkpoint_current = bool(
            state.latest_safe_verified_code == state.current_code
            and state.latest_safe_verified_revision == state.layout_revision
        )
        self.last_repair_visual_checkpoint_current = bool(
            state.latest_visual_checkpoint_code == state.current_code
            and state.latest_visual_checkpoint_revision == state.layout_revision
            and state.latest_visual_checkpoint_hard_valid
        )

    def _run_single_repair(
        self,
        slide_id: int,
        code: str,
        all_issues: list[Issue],
        must_not: list[str],
        must_contain: list[str],
        initial_msg: str | list,
        state_template: dict,
        attempt: int = 0,
        prev_failures_ctx: str = "",
        run_dir: str | None = None,
        turn_index: int = 0,
    ) -> str | None:
        """Run one repair attempt. Returns repaired code or None."""
        # Store issues for verify_layout change-magnitude feedback
        self._current_issues = all_issues

        # Use slightly different temperature per attempt to diversify
        temperatures = [0.05, 0.30]
        temp = temperatures[attempt % len(temperatures)]

        state = AgentState(
            original_code=code,
            current_code=code,
            checkpoints=[code],
            slide_id=slide_id,
            codegen_compiler=state_template["codegen_compiler"],
            case_dir=state_template["case_dir"],
            evidence=state_template.get("evidence"),
            bp_slide=state_template.get("bp_slide"),
            issue_types={i.issue_type for i in all_issues},
            text_loss_budget=self._text_loss_budget,
            initial_spatial_state=state_template.get("initial_spatial_state"),
        )
        state._run_dir = run_dir
        self._pending_actions = []  # Clear between slides
        self._multi_action_ignored_count = 0
        state._turn_index = turn_index

        max_tool_calls = min(
            self.MAX_TOOL_CALLS_CAP,
            max(12, len(all_issues) * self.MAX_TOOL_CALLS_PER_ISSUE),
        )
        dashboard_pressure = self._looks_like_table_dashboard_pressure_from(
            code,
            {issue.issue_type for issue in all_issues},
        )
        if self._needs_composition_closure(all_issues) or dashboard_pressure:
            # Structural composition repairs may need two additional
            # edit/verify cycles after the broad reflow to close local residuals.
            # This expands execution time only; it does not prescribe a layout.
            max_tool_calls = min(self.MAX_TOOL_CALLS_CAP, max_tool_calls + 10)

        system_prompt = self._system_prompt

        # Post-repair whitespace self-check — applies to all repair sessions
        system_prompt += (
            "\n\n**POST-REPAIR WHITESPACE CHECK (ALWAYS DO THIS):**\n"
            "After fixing overflow/overlap issues, check the SPACE MAP in your "
            "`verify_layout` output. If it shows:\n"
            "- Coverage below 70%, OR\n"
            "- An empty band (Bottom/Middle X% has no content), OR\n"
            "- One side of the grid is mostly '.' while the other is '#'\n"
            "Then you MUST fix the whitespace by: reducing body/container padding, "
            "increasing element heights (figure max-height, row min-height, gaps), "
            "using justify-content:space-between on flex containers, or enlarging "
            "undersized figures. Do NOT add new text content — only redistribute "
            "existing elements to fill the space evenly.\n"
        )

        if all(issue.issue_type == "svg_visual_defect" for issue in all_issues):
            system_prompt += (
                "\n\n**PURE SVG REPAIR SCOPE:**\n"
                "- Change only SVG geometry, paths, markers, endpoints, and styling needed "
                "for the reported target.\n"
                "- If the defect is inside an external `<img src=\"*.svg\">` asset, "
                "use `create_svg_asset` to create a fixed turn-local SVG and then "
                "update only that target `<img src>` to the returned path. This is an "
                "allowed media-target replacement for `svg_visual_defect`.\n"
                "- Preserve all source-visible labels and claims. Keep the rendered "
                "reading path coherent if SVG geometry changes.\n"
                "- Preserve media references, SVG accessibility semantics, and the semantic "
                "DOM outside SVG except for the allowed target SVG asset replacement above. "
                "The acceptance gate rejects out-of-scope changes.\n"
            )

        # When blueprint is not available (spatial-only repair mode),
        # tell agent to skip regen_slide and use apply_edits directly
        if not state_template.get("bp_slide"):
            external_evidence = (
                "the current render and spatial evidence"
                if self._enable_render_preview
                else (
                    "the current HTML/CSS, LAYOUT ANCHOR, RELATION MAP, "
                    "SPACE MAP, and detector evidence"
                )
            )
            completion_evidence = (
                "the original issue and full render"
                if self._enable_render_preview
                else "the original issue and current spatial evidence"
            )
            system_prompt += (
                "\n\n**IMPORTANT — EXTERNAL HTML REPAIR MODE:**\n"
                "- `regen_slide` is NOT available. Do NOT attempt it.\n"
                "- You have full freedom to restructure the DOM, not just CSS tweaks, when "
                f"{external_evidence} shows that several symptoms share one spatial cause.\n"
                f"- Diagnose from {external_evidence} whether the issue is local or belongs to "
                "a shared body-space conflict. Inspect the whole fixed-canvas budget before "
                "choosing the scope of the edit.\n"
                "- A coherent restructure may require several edits. Continue through a "
                "recoverable same-region intermediate state while text, roles, and meaning "
                "remain intact and the approach is converging. Roll back or change direction "
                "when information is lost, unrelated regions regress, or pressure is merely "
                "being displaced.\n"
                "- Preserve all source-visible text, semantic relationships, information-bearing "
                "visuals, and meaning-dependent sequences. Do not use clipping, hidden overflow, "
                "or off-canvas placement as a fit solution.\n"
                "- After structural CSS or DOM changes, call verify_layout and judge the result "
                f"against {completion_evidence}.\n"
            )

        # Inject previous repair failures context if available
        user_content = initial_msg
        if prev_failures_ctx:
            if isinstance(initial_msg, list):
                user_content = [dict(block) for block in initial_msg]
                for block in reversed(user_content):
                    if block.get("type") == "text":
                        block["text"] = (
                            str(block.get("text", ""))
                            + "\n\n" + prev_failures_ctx
                        )
                        break
                else:
                    user_content.append({
                        "type": "text",
                        "text": prev_failures_ctx,
                    })
            else:
                user_content = initial_msg + "\n\n" + prev_failures_ctx

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content},
        ]

        tool_calls = 0
        no_progress = 0
        submit_bounced_content = 0   # MUST-NOT + content accuracy bounces
        submit_bounced_spatial = 0   # spatial (overflow/OOB/overlap) bounces
        submit_bounced_raw_figure = 0
        has_content_edit_since_verify = False  # track content edits without verify
        has_raw_figure_issue = any(
            issue.issue_type in {"raw_figure", "raw_table"}
            for issue in all_issues
        )
        has_composition_closure_issue = self._needs_composition_closure(all_issues)
        has_svg_visual_issue = (
            any(issue.issue_type == "svg_visual_defect" for issue in all_issues)
            or (
                self._enable_render_preview
                and self._is_html_code(code)
                and bool(re.search(r"<svg\b", code, re.IGNORECASE))
                and any(issue.rubric_id.startswith("B") for issue in all_issues)
            )
        )
        visual_only_repair = bool(all_issues) and all(
            issue.rubric_id.startswith("B") for issue in all_issues
        )
        visual_repair_may_replace_image = any(
            issue.issue_type in {"raw_figure", "raw_table", "svg_visual_defect"}
            for issue in all_issues
        )
        visual_repair_may_change_formatting = any(
            issue.issue_type == "formatting_error"
            for issue in all_issues
        )
        visual_repair_may_change_text = issues_allow_visible_text_change(all_issues)
        state.allow_visible_text_change = visual_repair_may_change_text
        state.allow_support_copy_compression = (
            issues_allow_support_copy_compression(all_issues)
        )

        while tool_calls < max_tool_calls:
            # Every action after the first must be planned from real tool output.
            # Never execute speculative multi-tool responses from an old state.
            pending = None
            if pending is None:
                # LLM turn
                try:
                    response = self.llm.call_multiturn(
                        messages=messages,
                        model=self.model,
                        module_name="agent_repair",
                        prompt_version="agent.v10",
                        max_tokens=4096,
                        temperature=temp,
                    )
                except Exception as e:
                    # Transient LLM errors (content-filter 400, rate-limit 429,
                    # 5xx/timeouts) often clear on retry. Retry a bounded number
                    # of times before giving up.
                    emsg = str(e)
                    is_transient = any(s in emsg.lower() for s in (
                        "content management policy", "content_filter", "filtered",
                        "rate limit", "429", "timeout", "timed out",
                        "temporarily", "503", "502", "500", "overloaded",
                    ))
                    n_llm_retries = getattr(state, "_llm_retry_count", 0)
                    MAX_LLM_RETRIES = 3
                    if is_transient and n_llm_retries < MAX_LLM_RETRIES:
                        state._llm_retry_count = n_llm_retries + 1
                        logger.warning(
                            "Agent repair slide %d turn %d: transient LLM error "
                            "(%s), retry %d/%d",
                            slide_id, tool_calls, emsg[:120],
                            n_llm_retries + 1, MAX_LLM_RETRIES,
                        )
                        # A content filter often trips on the LAST tool result we
                        # appended (e.g. a verify dump with flagged text). Nudge the
                        # model with a short neutral reminder rather than resending
                        # the identical context that just got filtered.
                        if "content" in emsg.lower() or "filter" in emsg.lower():
                            messages.append({
                                "role": "user",
                                "content": (
                                    "(Continue the repair. Respond with your next "
                                    "tool call as normal JSON — keep reasoning brief "
                                    "and focused on the CSS/layout fix.)"
                                ),
                            })
                        tool_calls += 1
                        continue
                    logger.warning(
                        "Agent repair slide %d LLM error at turn %d (giving up "
                        "after %d retries): %s",
                        slide_id, tool_calls, n_llm_retries, emsg[:200],
                    )
                    break

                messages.append({"role": "assistant", "content": response})

                # Parse action
                action = self._parse_action(response)
                if action is None:
                    # Empty responses are likely API transient errors — retry
                    # without counting as no-progress. Only count non-empty
                    # malformed responses as no-progress.
                    is_empty = not response or not response.strip()
                    logger.info(
                        "Agent slide %d turn %d: parse error (%s), response=%s",
                        slide_id, tool_calls,
                        "empty" if is_empty else "malformed",
                        response[:200],
                    )
                    if is_empty:
                        # Remove the empty assistant message to keep context clean
                        messages.pop()
                        tool_calls += 1
                        no_progress += 1
                        if no_progress >= self.MAX_NO_PROGRESS:
                            logger.info(
                                "Agent repair slide %d: %d empty/no-progress turns, aborting",
                                slide_id, no_progress,
                            )
                            break
                        continue
                    parse_error = getattr(self, "_last_parse_error_message", "")
                    messages.append({"role": "user", "content":
                        parse_error or (
                            "Error: could not parse your action. "
                            "Return a JSON object with a \"tool\" field. "
                            "Example: {\"tool\": \"apply_edits\", \"reason\": \"...\", \"edits\": [{\"search\": \"old\", \"replace\": \"new\"}]}. "
                            "Do NOT wrap in markdown fences. Just the raw JSON object."
                        )
                    })
                    tool_calls += 1
                    # Only count every other malformed response as no-progress
                    # to give the model more chances to recover
                    if tool_calls % 2 == 0:
                        no_progress += 1
                    if no_progress >= self.MAX_NO_PROGRESS:
                        logger.info(
                            "Agent repair slide %d: no progress for %d turns, aborting",
                            slide_id, no_progress,
                        )
                        break
                    continue

                if getattr(self, "_multi_action_ignored_count", 0):
                    messages[-1]["content"] = json.dumps(action, ensure_ascii=False)

                # Log agent action
                tool_name = action.get("tool", "?")
                reasoning = action.get("reasoning", "")[:150]
                logger.info(
                    "Agent slide %d turn %d: tool=%s reason=%s",
                    slide_id, tool_calls, tool_name, reasoning,
                )

            # ── Reasoning validation (Route 1.2) ──
            # Require reasoning for all tools except submit and plan
            # Only enforce once per session to avoid wasting tool budget
            if (self._enable_macro_planning
                    and tool_name not in ("submit", "plan", "update_plan")
                    and not action.get("reasoning")
                    and not getattr(state, '_reasoning_warned', False)):
                state._reasoning_warned = True
                messages.append({"role": "user", "content":
                    "Missing required 'reasoning' field. Every tool call "
                    "(except submit) MUST include a 'reasoning' field "
                    "explaining your analysis and why this action is the "
                    "right next step. Re-submit your action with reasoning."
                })
                tool_calls += 1
                continue

            # ── Macro plan enforcement (Route 1.1) ──
            # If macro planning is enabled and this is not the plan tool,
            # nudge the agent to submit a plan first. Auto-accept after
            # 2 nudges to avoid wasting budget.
            if (self._enable_macro_planning
                    and not state.has_plan
                    and tool_name not in ("plan", "update_plan", "get_current_code", "submit")):
                plan_nudge_count = getattr(state, '_plan_nudge_count', 0)
                if plan_nudge_count < 3:
                    state._plan_nudge_count = plan_nudge_count + 1
                    messages.append({"role": "user", "content":
                        "You must submit a repair plan FIRST before making any "
                        "edits. Call the plan tool with your overall strategy: "
                        "group issues by type, order from safe to risky, and "
                        "specify when to verify. Only then proceed with edits."
                    })
                    tool_calls += 1
                    continue
                else:
                    # Auto-accept to not waste more budget
                    state.has_plan = True
                    logger.info(
                        "Agent slide %d: auto-accepting plan after %d nudges",
                        slide_id, plan_nudge_count,
                    )

            # Submit = validate then exit
            if action.get("tool") == "submit":
                submit_bounced_total = (
                    submit_bounced_content + submit_bounced_spatial
                )

                # Zero-edit guard (C10): if agent hasn't made any edits
                # at all, bounce once to force at least one attempt
                if (state.current_code == state.original_code
                        and submit_bounced_total == 0
                        and self._enable_macro_planning
                        and not state.attempted_code_change):
                    submit_bounced_content += 1
                    messages.append({"role": "user", "content":
                        "⚠ SUBMIT BLOCKED — you haven't made any edits "
                        "yet. You must attempt at least one fix before "
                        "submitting. Review the issue list and attempt "
                        "the easiest content fix (e.g., text correction). "
                        "If ALL issues are genuinely unfixable, explain "
                        "why in your reasoning and submit again."
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced (zero edits)",
                        slide_id,
                    )
                    continue

                composition_closure_bounces = getattr(
                    state, "_composition_closure_bounces", 0,
                )
                if (
                    has_composition_closure_issue
                    and not self._summary_has_composition_closure(state, all_issues)
                ):
                    state._composition_closure_bounces = composition_closure_bounces + 1
                    messages.append({
                        "role": "user",
                        "content": self._composition_closure_block_message(
                            state, all_issues, for_submit=True,
                        ),
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced (missing composition self-assessment)",
                        slide_id,
                    )
                    continue

                if has_composition_closure_issue:
                    composition_unresolved_reasons = (
                        self._composition_closure_unresolved_reasons(
                            state, all_issues,
                        )
                    )
                    if composition_unresolved_reasons:
                        state._composition_closure_bounces = (
                            composition_closure_bounces + 1
                        )
                        messages.append({
                            "role": "user",
                            "content": self._composition_completion_block_message(
                                state, all_issues, for_submit=True,
                            ),
                        })
                        tool_calls += 1
                        logger.info(
                            "Agent slide %d: submit bounced "
                            "(composition unresolved: %s)",
                            slide_id,
                            "; ".join(composition_unresolved_reasons[:3]),
                        )
                        continue

                # SVG topology cannot be validated from DOM measurements. Once
                # an SVG repair changes code, require the same multimodal agent
                # to inspect the resulting pixels before it may submit.
                preview_bounces = getattr(state, "_svg_preview_bounces", 0)
                if (
                    self._enable_render_preview
                    and has_svg_visual_issue
                    and getattr(state, "_svg_edit_since_preview", False)
                    and preview_bounces < 2
                ):
                    state._svg_preview_bounces = preview_bounces + 1
                    messages.append({
                        "role": "user",
                        "content": (
                            "SUBMIT BLOCKED: the SVG changed after the last "
                            "visual inspection. Call render_preview, inspect "
                            "the current pixels and graph roles, then fix any "
                            "remaining defect or submit again."
                        ),
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced "
                        "(SVG edit without current preview)",
                        slide_id,
                    )
                    continue

                raw_css_crop_hints = (
                    html_image_css_crop_hints(state.current_code)
                    if has_raw_figure_issue and self._is_html_code(state.current_code)
                    else []
                )
                if raw_css_crop_hints and submit_bounced_raw_figure < 2:
                    submit_bounced_raw_figure += 1
                    messages.append({
                        "role": "user",
                        "content": (
                            "SUBMIT BLOCKED: this raw_figure/raw_table repair "
                            "still relies on CSS image-window cropping: "
                            f"{', '.join(raw_css_crop_hints)}. B17 requires a "
                            "presentation-adapted image source, exact-data generated "
                            "chart, or fidelity-preserving SVG summary displayed intact with object-fit: contain. "
                            "Create or use a real cropped/recomposed source asset first when possible, "
                            "replace the image src if needed, remove CSS crop "
                            "mechanisms such as object-view-box/object-fit:cover/"
                            "object-fit:none/negative offsets, then call "
                            "verify_layout and render_preview."
                        ),
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced "
                        "(raw figure CSS crop: %s)",
                        slide_id, ", ".join(raw_css_crop_hints),
                    )
                    continue

                raw_preview_bounces = getattr(state, "_raw_figure_preview_bounces", 0)
                if (
                    self._enable_render_preview
                    and has_raw_figure_issue
                    and getattr(state, "_raw_figure_edit_since_preview", False)
                    and raw_preview_bounces < 2
                ):
                    state._raw_figure_preview_bounces = raw_preview_bounces + 1
                    messages.append({
                        "role": "user",
                        "content": (
                            "SUBMIT BLOCKED: the raw figure/table region changed "
                            "after the last visual inspection. Call "
                            "render_preview and inspect whether the intended "
                            "subject remains complete, labels are readable, "
                            "key findings are visually guided, and no panel is "
                            "cut off. Fix any remaining defect before submitting."
                        ),
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced "
                        "(raw figure edit without current preview)",
                        slide_id,
                    )
                    continue

                # Spatial submit bounce: if significant overlaps or OOB remain,
                # bounce back to let agent continue fixing from current state.
                # Only bounce when budget allows meaningful follow-up (≥10 calls).
                _budget_left = max_tool_calls - tool_calls
                if submit_bounced_spatial < 2 and _budget_left >= 10 and self._is_html_code(state.current_code):
                    try:
                        from .html_spatial_state import extract_html_slide_state
                        _submit_state = extract_html_slide_state(
                            slide_id, state.current_code,
                            html_base_dir=Path(state.case_dir),
                            asset_base_dirs=self._html_asset_base_dirs(
                                state.case_dir, getattr(state, "_run_dir", None),
                            ),
                        )
                        _submit_sig = self._targeted_significant_issues(
                            getattr(state, '_t0_html_state', _submit_state),
                            _submit_state,
                        )
                        _submit_overlaps = len(_submit_sig.get("overlap", []))
                        _submit_oob = len(_submit_sig.get("out_of_bounds", []))
                        _submit_overflow = len(_submit_sig.get("text_overflow", []))
                        _submit_total = _submit_overlaps + _submit_oob + _submit_overflow
                        if _submit_total >= 3:
                            submit_bounced_spatial += 1
                            messages.append({
                                "role": "user",
                                "content": (
                                    f"SUBMIT BLOCKED: {_submit_total} significant "
                                    f"spatial issues remain ({_submit_overlaps} overlaps, "
                                    f"{_submit_oob} out-of-bounds, {_submit_overflow} "
                                    f"text overflow). Continue editing from the current "
                                    f"state to fix these — your edits so far are preserved. "
                                    f"Call verify_layout to see which elements still overlap."
                                ),
                            })
                            tool_calls += 1
                            logger.info(
                                "Agent slide %d: submit bounced "
                                "(residual significant issues: %s)",
                                slide_id, {k: len(v) for k, v in _submit_sig.items()},
                            )
                            continue
                    except Exception as _spatial_err:
                        logger.warning(
                            "Agent slide %d: spatial submit check failed: %s",
                            slide_id, str(_spatial_err)[:100],
                        )

                # Validate: check if MUST NOT strings still in code
                if submit_bounced_content < 2 and must_not:
                    warnings = self._validate_before_submit(
                        state.current_code, must_not,
                    )
                    if warnings:
                        submit_bounced_content += 1
                        messages.append({"role": "user", "content": warnings})
                        tool_calls += 1
                        logger.info(
                            "Agent slide %d: submit bounced "
                            "(MUST-NOT, content=%d spatial=%d)",
                            slide_id, submit_bounced_content,
                            submit_bounced_spatial,
                        )
                        continue

                # Word-retention hard floor: fires INDEPENDENTLY of the
                # spatial bounce counter. A slide that lost >30% of its
                # visible text is almost certainly over-deleted. The agent
                # may correctly remove a low-value support channel, but
                # deleting entire cards with method names / benchmarks /
                # metrics to solve overlap is the WRONG fix.
                # Bounce ONCE with strong restoration guidance; after that
                # let it ship (the agent may genuinely need the deletion).
                _wc_bounce_attr = '_wc_floor_bounced'
                if (not getattr(state, _wc_bounce_attr, False)
                        and state.current_code != code):
                    t0_wc = _count_html_words(state.original_code)
                    t1_wc = _count_html_words(state.current_code)
                    if t0_wc > 20:
                        wc_pct = round(100 * t1_wc / t0_wc)
                        if wc_pct < 70:
                            setattr(state, _wc_bounce_attr, True)
                            messages.append({
                                "role": "user",
                                "content": (
                                    f"🚨 SUBMIT BLOCKED — word retention "
                                    f"is only {wc_pct}% ({t0_wc}→{t1_wc} "
                                    f"words). You deleted too much content "
                                    f"while fixing spatial issues.\n\n"
                                    f"RESTORE the deleted content. The "
                                    f"correct fix for overlap/overflow is "
                                    f"to MOVE or RESIZE elements, not to "
                                    f"delete entire cards/sections. If a "
                                    f"region is too crowded:\n"
                                    f"  1. Reduce excessive padding/margins\n"
                                    f"  2. Adjust container dimensions or grid tracks\n"
                                    f"  3. Move the affected region as a unit\n"
                                    f"  4. Remove ONLY decorative accents\n"
                                    f"NEVER delete a card/section that "
                                    f"contains method names, metrics, "
                                    f"benchmark results, or key claims.\n"
                                    f"Rollback to restore, then try "
                                    f"CSS-only fixes. Call submit again."
                                ),
                            })
                            tool_calls += 1
                            logger.info(
                                "Agent slide %d: submit bounced "
                                "(word retention %d%%, floor=70%%)",
                                slide_id, wc_pct,
                            )
                            continue

                # This gate also applies to content repairs. Adding a required
                # sentence does not authorize collapsing the slide's dominant
                # hero/title or globally tightening typography to make it fit.
                if (
                    state.current_code != code
                    and self._is_html_code(state.current_code)
                ):
                    compression_ok, compression_reason = (
                        validate_repair_not_visual_compression(
                            state.original_code,
                            state.current_code,
                            allow_dominant_element_removal=(
                                issues_allow_dominant_element_removal(
                                    self._current_issues
                                )
                            ),
                        )
                    )
                    # Bypass compression gate when overflow still exists —
                    # compression is necessary to fit content within canvas.
                    _submit_overflow = getattr(state, '_last_verified_overflow_px', 0) or 0
                    if _submit_overflow > 30:
                        compression_ok = True
                    if not compression_ok:
                        state.current_code = (
                            state.best_verified_code
                            or state.last_verified_code
                            or state.original_code
                        )
                        self._invalidate_verify_after_code_change(
                            state,
                            "submit compression gate rolled back code",
                        )
                        submit_bounced_spatial += 1
                        messages.append({
                            "role": "user",
                            "content": (
                                "SUBMIT BLOCKED — visual compression shipment "
                                f"gate failed: {compression_reason}. The invalid "
                                "state was rolled back to the latest legal "
                                "checkpoint. Preserve the dominant hero/title. "
                                "Reallocate column width or container height and "
                                "reduce only local gaps/padding before reducing "
                                "body type; then verify again."
                            ),
                        })
                        tool_calls += 1
                        continue

                # Fixed-format text budget gate. Content repairs may need to
                # qualify a title or takeaway, but long source text belongs in
                # body/interpretation text, not in fixed title/footer chrome.
                fixed_budget_bounces = getattr(
                    state, "_fixed_format_budget_bounces", 0,
                )
                if (
                    state.current_code != code
                    and self._is_html_code(state.current_code)
                    and fixed_budget_bounces < 2
                ):
                    budget_warnings = _fixed_format_text_budget_warnings(
                        state.current_code,
                        baseline_html=state.original_code,
                    )
                    if budget_warnings:
                        state._fixed_format_budget_bounces = fixed_budget_bounces + 1
                        submit_bounced_content += 1
                        messages.append({
                            "role": "user",
                            "content": (
                                "SUBMIT BLOCKED — fixed-format text budget "
                                "failed:\n"
                                + "\n".join(f"- {w}" for w in budget_warnings[:4])
                                + "\n\nKeep titles/page headers concise. Do not put long "
                                  "source-backed corrections into full-width "
                                  "bottom bars, footers, or source notes. Move "
                                  "necessary qualifiers into the closest body/"
                                  "interpretation sentence, then call "
                                  "verify_layout before submitting."
                            ),
                        })
                        tool_calls += 1
                        logger.info(
                            "Agent slide %d: submit bounced "
                            "(fixed-format text budget: %s)",
                            slide_id, "; ".join(budget_warnings[:3]),
                        )
                        continue

                # Unconditional source/media visual-scope gate. Rendered text
                # sampling is intentionally advisory: it is brittle during
                # dashboard/table rhythm calibration and can flag useful
                # reveal/reflow states that preserve the source DOM.
                if (
                    visual_only_repair
                    and state.current_code != code
                    and self._is_html_code(state.current_code)
                ):
                    scope_ok, scope_reason = validate_visual_repair_scope(
                        state.original_code,
                        state.current_code,
                        allow_image_replacement=visual_repair_may_replace_image,
                        allow_text_formatting_change=visual_repair_may_change_formatting,
                        allow_text_content_change=visual_repair_may_change_text,
                    )
                    downgrade_ok, downgrade_reason = validate_repair_not_visual_downgrade(
                        state.original_code, state.current_code,
                    )
                    from .html_spatial_state import extract_html_slide_state
                    html_asset_base_dirs = self._html_asset_base_dirs(
                        state.case_dir,
                        getattr(state, "_run_dir", None),
                        getattr(state, "_turn_index", None),
                    )
                    t0_scope_state = getattr(state, '_t0_html_gate_state', None)
                    if t0_scope_state is None:
                        t0_scope_state = extract_html_slide_state(
                            slide_id,
                            state.original_code,
                            html_base_dir=Path(state.case_dir),
                            asset_base_dirs=html_asset_base_dirs,
                        )
                        state._t0_html_gate_state = t0_scope_state
                    t1_scope_state = extract_html_slide_state(
                        slide_id,
                        state.current_code,
                        html_base_dir=Path(state.case_dir),
                        asset_base_dirs=html_asset_base_dirs,
                    )
                    from .repair_utils import validate_rendered_text_preservation
                    rendered_ok, rendered_reason = validate_rendered_text_preservation(
                        t0_scope_state,
                        t1_scope_state,
                        allow_revealed_text=issues_allow_rendered_text_reveal(
                            self._current_issues,
                        ),
                        allow_text_formatting_change=(
                            visual_repair_may_change_formatting
                        ),
                        allow_text_content_change=visual_repair_may_change_text,
                    )
                    if not rendered_ok:
                        logger.info(
                            "Agent slide %d: submit rendered-text advisory: %s",
                            slide_id, rendered_reason,
                        )
                    if not scope_ok or not downgrade_ok:
                        reason = (
                            scope_reason if not scope_ok
                            else downgrade_reason
                        )
                        state.current_code = (
                            state.best_verified_code
                            or state.last_verified_code
                            or state.original_code
                        )
                        self._invalidate_verify_after_code_change(
                            state,
                            "submit visual-scope gate rolled back code",
                        )
                        submit_bounced_content += 1
                        messages.append({
                            "role": "user",
                            "content": (
                                "SUBMIT BLOCKED — visual-only repair shipment gate failed: "
                                f"{reason}. The invalid state was rolled back to the "
                                "latest legal checkpoint. Make a targeted layout/style "
                                "edit, verify it, then submit again."
                            ),
                        })
                        tool_calls += 1
                        continue

                # Pre-submit spatial check: overflow, OOB, AND overlaps.
                # This catches cases where the agent made edits without
                # calling verify_layout afterwards — the #1 source of
                # critical T1 issues (overlap + text_overflow).
                # Separate counter ensures spatial always gets checked.
                # For slides with critical content accuracy issues, reduce
                # bounce limit to avoid blocking content fixes.
                n_crit_content = sum(
                    1 for i in self._current_issues
                    if i.issue_type in CRITICAL_CONTENT_TYPES
                    and i.status.value == "open"
                )
                spatial_bounce_limit = 2 if n_crit_content >= 2 else 3
                if submit_bounced_spatial < spatial_bounce_limit and state.current_code != code:
                    try:
                        is_html = self._is_html_code(state.current_code)

                        if is_html:
                            # HTML mode: use Playwright DOM check (same as verify_layout)
                            from .html_spatial_state import (
                                extract_html_slide_state,
                                significant_issue_regressions,
                                stable_block_identity,
                            )
                            html_asset_base_dirs = self._html_asset_base_dirs(
                                state.case_dir,
                                getattr(state, "_run_dir", None),
                                getattr(state, "_turn_index", None),
                            )
                            t1_html_st = extract_html_slide_state(
                                slide_id,
                                state.current_code,
                                html_base_dir=Path(state.case_dir),
                                asset_base_dirs=html_asset_base_dirs,
                            )
                            if not hasattr(state, '_t0_html_gate_state'):
                                state._t0_html_gate_state = extract_html_slide_state(
                                    slide_id,
                                    state.original_code,
                                    html_base_dir=Path(state.case_dir),
                                    asset_base_dirs=html_asset_base_dirs,
                                )
                            t0_html_st = state._t0_html_gate_state

                            # Count per-category significant issues via the SINGLE
                            # SOURCE OF TRUTH (count_significant_issues) so the
                            # agent's notion of "clean" matches the external scorer
                            # exactly. The previous inline helpers
                            # (_significant_overflows/_count_real_overlaps/
                            # _significant_oob) diverged from count_issues AND were
                            # blind to clipped_blocks — a slide whose only residual
                            # was clipped text passed this gate while the harness
                            # still flagged it. The SSOT covers all six categories
                            # incl. clipped + canvas_truncation.
                            targeted_t1_cats = self._targeted_significant_issues(
                                t0_html_st, t1_html_st,
                            )
                            t1_counts = {
                                k: len(v) for k, v in targeted_t1_cats.items()
                            }
                            targeted_residual_categories = (
                                self._targeted_residual_categories()
                            )
                            residual_total = sum(
                                t1_counts.get(category, 0)
                                for category in targeted_residual_categories
                            )

                            # (1) Regression bounce. Compare stable physical
                            # defects rather than raw category counts: the same
                            # pair may move between overlap and occlusion, and
                            # the same block between visible overflow and clip.
                            regressions = significant_issue_regressions(
                                t0_html_st, t1_html_st,
                            )
                            if can_exempt_raw_figure_image_crop(
                                self._current_issues, state.current_code,
                            ):
                                regressions.pop("image_crop", None)
                            new_issues = {
                                category: len(items)
                                for category, items in regressions.items()
                            }

                            current_safe_checkpoint = bool(
                                state.latest_safe_verified_code == state.current_code
                                and state.latest_safe_verified_revision == state.layout_revision
                            )
                            detector_regressions_are_advisory = bool(
                                visual_only_repair
                                and current_safe_checkpoint
                            )

                            # NET improvement override: when total spatial
                            # issues dropped by ≥50%, a small number of new
                            # regressions (relocated content) should not block
                            # submission — the agent is moving in the right
                            # direction overall.
                            if (
                                new_issues
                                and not detector_regressions_are_advisory
                            ):
                                t0_sig = sum(
                                    len(v) for v in
                                    count_significant_issues(t0_html_st).values()
                                )
                                t1_sig = sum(
                                    len(v) for v in
                                    count_significant_issues(t1_html_st).values()
                                )
                                new_issue_total = sum(new_issues.values())
                                # Also check actual overflow px
                                _t1_overflow_px = compute_overflow_px(t1_html_st.blocks)
                                if (
                                    t0_sig > 0
                                    and t1_sig <= t0_sig * 0.5
                                    and new_issue_total <= 6
                                    and _t1_overflow_px <= 20
                                ):
                                    detector_regressions_are_advisory = True

                            # Build detailed rejection message.  For a pure visual
                            # repair, detector deltas are advisory once verify_layout
                            # has measured this exact hard-valid revision.  The DOM probe
                            # is deliberately conservative and the curated visual
                            # targets themselves contain residual hits, so it must not
                            # overrule a grounded visual judgment by itself.
                            if new_issues and not detector_regressions_are_advisory:
                                warn_parts = []
                                if regressions.get("interaction"):
                                    details = "\n".join(
                                        f"NEW {kind.upper()}: {first} ↔ {second}"
                                        for kind, (first, second) in regressions["interaction"]
                                    )
                                    warn_parts.append(
                                        f"🚨 SUBMIT BLOCKED — "
                                        f"{new_issues['interaction']} new element "
                                        f"interaction(s) introduced:\n{details}"
                                    )
                                if regressions.get("content_fit"):
                                    details = ", ".join(
                                        f"{kind}:{block_id}"
                                        for kind, block_id in regressions["content_fit"]
                                    )
                                    warn_parts.append(
                                        f"🚨 SUBMIT BLOCKED — "
                                        f"{new_issues['content_fit']} new content-fit "
                                        f"defect(s) introduced: {details}."
                                    )
                                if new_issues.get("out_of_bounds"):
                                    warn_parts.append(
                                        f"🚨 SUBMIT BLOCKED — "
                                        f"{new_issues['out_of_bounds']} new "
                                        f"out-of-bounds element(s) introduced."
                                    )
                                if new_issues.get("image_crop"):
                                    warn_parts.append(
                                        f"🚨 SUBMIT BLOCKED — "
                                        f"{new_issues['image_crop']} new excessive "
                                        f"image crop(s) introduced."
                                    )
                                warn_parts.append(
                                    "\nFix these issues before submitting. "
                                    "Use verify_layout for details. "
                                    "Then call submit again."
                                )
                                submit_bounced_spatial += 1
                                messages.append({
                                    "role": "user",
                                    "content": "\n".join(warn_parts),
                                })
                                tool_calls += 1
                                logger.info(
                                    "Agent slide %d: submit bounced "
                                    "(spatial check via Playwright DOM, "
                                    "content=%d spatial=%d): %s",
                                    slide_id, submit_bounced_content,
                                    submit_bounced_spatial, new_issues,
                                )
                                continue
                            elif new_issues:
                                logger.info(
                                    "Agent slide %d: allowing hard-valid verified "
                                    "B-family candidate with detector advisories: %s",
                                    slide_id, new_issues,
                                )

                            # (2) Residual backstop — BOUNDED, PROGRESS-AWARE,
                            # ESCAPABLE. The regression bounce above only blocks
                            # "worse than original"; it lets a slide submit with
                            # residual fixable issues that are simply fewer than
                            # baseline (e.g. 4→3). This is a known failure mode:
                            # clipped-text issues.
                            #
                            # A STRICT single-fire bounce proved too weak: on a
                            # hard multi-clip slide the agent
                            # burned its one bounce early, then later edits made
                            # clips worse and it submitted with residuals.
                            #
                            # Re-bounce when ALL of these hold (a clean or
                            # genuinely-justified trajectory always escapes):
                            #   • substantial residuals remain (> RESIDUAL_FLOOR),
                            #   • the agent has ample budget left,
                            #   • it has NOT justified them in its summary, and
                            #   • under a hard re-bounce ceiling.
                            # The agent keeps full authority to delete or
                            # restructure; this only prevents shipping residuals
                            # by inattention while budget remains.
                            RESIDUAL_FLOOR = 2          # ≤2 residuals → let it ship (don't nag tiny tails)
                            RESIDUAL_BOUNCE_CEIL = 3    # hard ceiling on re-bounces
                            budget_left = max_tool_calls - tool_calls
                            n_res_bounces = getattr(state, '_residual_bounce_count', 0)
                            justified = self._summary_justifies_residuals(state)
                            # Always fire at least once on any residual (preserves
                            # the prior thin backstop); re-fire only when the slide
                            # is materially dirty, budget is healthy, unjustified.
                            first_fire = (residual_total > 0 and n_res_bounces == 0)
                            refire = (
                                residual_total > RESIDUAL_FLOOR
                                and n_res_bounces < RESIDUAL_BOUNCE_CEIL
                                and budget_left >= 6
                                and not justified
                            )
                            if first_fire or refire:
                                setattr(state, '_residual_bounced', True)
                                state._residual_bounce_count = n_res_bounces + 1
                                # Enumerate the residuals with element ids + px so
                                # the agent knows exactly what is still open.
                                res_lines = []
                                for bid in (
                                    (targeted_t1_cats.get("clipped", []) if "clipped" in targeted_residual_categories else [])
                                    + (targeted_t1_cats.get("canvas_truncation", []) if "canvas_truncation" in targeted_residual_categories else [])
                                ):
                                    blk = next((b for b in t1_html_st.blocks if b.block_id == bid), None)
                                    if blk is not None:
                                        px = int(getattr(blk, 'clipped_bottom_px', 0) or 0)
                                        prev = " ".join(blk.text_lines)[:40] if blk.text_lines else bid
                                        res_lines.append(f"  • CLIPPED {bid}: \"{prev}\" ({px}px hidden)")
                                for bid in (
                                    targeted_t1_cats.get("text_overflow", [])
                                    if "text_overflow" in targeted_residual_categories else []
                                ):
                                    res_lines.append(f"  • OVERFLOW {bid}")
                                for issue_id in (
                                    targeted_t1_cats.get("svg_text_overflow", [])
                                    if "svg_text_overflow" in targeted_residual_categories else []
                                ):
                                    res_lines.append(
                                        self._format_svg_text_overflow_residual(
                                            t1_html_st, issue_id,
                                        )
                                    )
                                for a, b in (
                                    targeted_t1_cats.get("overlap", [])
                                    if "overlap" in targeted_residual_categories else []
                                ):
                                    res_lines.append(f"  • OVERLAP {a} ↔ {b}")
                                for bid in (
                                    targeted_t1_cats.get("out_of_bounds", [])
                                    if "out_of_bounds" in targeted_residual_categories else []
                                ):
                                    res_lines.append(f"  • OUT-OF-BOUNDS {bid}")
                                for a, b in (
                                    targeted_t1_cats.get("occlusion", [])
                                    if "occlusion" in targeted_residual_categories else []
                                ):
                                    res_lines.append(f"  • OCCLUDED {b} behind {a}")
                                detail = "\n".join(res_lines) if res_lines else f"  {residual_total} residual issue(s)"
                                # On a re-fire, the message is sharper: the agent
                                # has tried to submit a still-dirty slide more than
                                # once while holding budget — remind it clipped text
                                # is a COUNTED defect (not ignorable chrome) and that
                                # it has tool calls left to spend.
                                if first_fire:
                                    body = (
                                        f"⚠ {residual_total} spatial issue(s) still open on this "
                                        f"slide:\n{detail}\n\n"
                                        f"Compare these residuals with the original issue and the current "
                                        f"render. Decide whether they are isolated symptoms or evidence that "
                                        f"the owning composition is still under pressure, then continue the "
                                        f"current structural attempt or change direction accordingly. Preserve "
                                        f"all information-bearing content. If a residual is genuinely outside "
                                        f"scope or cannot be resolved without damage, explain that specific "
                                        f"judgment in submit_repair_summary. For SVG text overflow, repair the "
                                        f"asset's internal layout rather than hiding the label."
                                    )
                                else:
                                    body = (
                                        f"⚠ Still {residual_total} spatial issue(s) open — and you "
                                        f"have ~{budget_left} tool calls left to spend:\n{detail}\n\n"
                                        f"Clipped protected content is still unreadable. Use the current "
                                        f"render to decide whether the active approach can converge through "
                                        f"further coherent edits or should be replaced. Do not rollback only "
                                        f"because an unfinished same-region state is imperfect, and do not "
                                        f"continue merely by moving the same pressure around. Preserve visible "
                                        f"text, semantic roles, and information-bearing visuals. If a residual "
                                        f"is genuinely outside scope or structurally unresolved, identify the "
                                        f"specific reason in submit_repair_summary. For SVG text overflow, "
                                        f"repair the asset's internal layout rather than hiding the label."
                                    )
                                messages.append({
                                    "role": "user",
                                    "content": body,
                                })
                                tool_calls += 1
                                logger.info(
                                    "Agent slide %d: submit bounced "
                                    "(residual significant issues: %s, "
                                    "bounce #%d, budget_left=%d, first=%s)",
                                    slide_id, t1_counts,
                                    state._residual_bounce_count, budget_left,
                                    first_fire,
                                )
                                continue

                            # Coverage and absolute font size are diagnostic
                            # signals, not acceptance targets. Dense tables,
                            # rankings, legends, and compact support copy can use
                            # smaller role-relative type when the hierarchy and
                            # rendered information remain coherent. Concrete
                            # clipping, overlap, text loss, media loss, and visual
                            # hierarchy damage are handled by the existing gates.
                        else:
                            # PPTX mode: use extract_slide_state for detailed feedback
                            baseline_st = extract_slide_state(
                                slide_id, state.original_code,
                            )
                            current_st = extract_slide_state(
                                slide_id, state.current_code,
                            )
                            t0_of = set(baseline_st.overflow_blocks)
                            t1_of = set(current_st.overflow_blocks)
                            new_of = t1_of - t0_of
                            t0_ob = set(baseline_st.oob_blocks)
                            t1_ob = set(current_st.oob_blocks)
                            new_ob = t1_ob - t0_ob
                            # Overlap detection: find overlaps in T1
                            # that did not exist in T0
                            t0_overlap_set = {
                                (min(a, b), max(a, b))
                                for a, b, _ in baseline_st.overlap_pairs
                            }
                            # Container/SVG tags whose overlaps are structural, not defects
                            _CONTAINER_TAGS = frozenset({
                                "svg", "g", "path", "rect", "text", "line",
                                "circle", "ellipse", "polygon", "polyline",
                                "tspan", "use", "col", "colgroup",
                            })
                            def _is_container_blk(bid):
                                blk = next((b for b in current_st.blocks if b.block_id == bid), None)
                                if not blk:
                                    return False
                                sel = (blk.css_selector or "").lower()
                                tag = sel.split(".")[-1].split("[")[0].split(":")[0].strip()
                                return tag in _CONTAINER_TAGS or "svg" in bid.lower()

                            new_overlaps = []
                            for a, b, area in current_st.overlap_pairs:
                                key = (min(a, b), max(a, b))
                                if key not in t0_overlap_set:
                                    # Skip container/SVG element overlaps
                                    if _is_container_blk(a) or _is_container_blk(b):
                                        continue
                                    new_overlaps.append((a, b, area))
                            # Near-overflow detection: blocks that are
                            # overflowing per Playwright, regardless of
                            # whether text grew (layout changes can also
                            # cause overflow by shrinking containers).
                            near_overflow = []
                            for b in current_st.blocks:
                                if b.is_overflowing and b.text_chars > 20:
                                    # Check if this is a NEW overflow
                                    orig_b = next(
                                        (ob for ob in baseline_st.blocks
                                         if ob.var_name == b.var_name),
                                        None,
                                    )
                                    if orig_b and not orig_b.is_overflowing:
                                        near_overflow.append(b)
                            has_hard_issues = new_of or new_ob or new_overlaps
                            if has_hard_issues or near_overflow:
                                # Build warning message
                                warn_parts = []
                                if new_of:
                                    of_details = []
                                    for bid in new_of:
                                        for b in current_st.blocks:
                                            if b.block_id == bid:
                                                of_details.append(
                                                    f"  {b.var_name}: "
                                                    f"{b.text_chars}ch in "
                                                    f"{b.w:.1f}\"×{b.h:.1f}\" "
                                                    f"(overflow: {b.overflow_bottom_px}px bottom, {b.overflow_right_px}px right)"
                                                )
                                    warn_parts.append(
                                        "🚨 SUBMIT BLOCKED — your edits "
                                        "introduced NEW text overflow:\n"
                                        + "\n".join(of_details)
                                    )
                                if new_ob:
                                    ob_details = []
                                    for bid in new_ob:
                                        for b in current_st.blocks:
                                            if b.block_id == bid:
                                                bottom = b.y + b.h
                                                ob_details.append(
                                                    f"  {b.var_name}: "
                                                    f"bottom={bottom:.2f}\""
                                                )
                                    warn_parts.append(
                                        "🚨 SUBMIT BLOCKED — your edits "
                                        "pushed elements out-of-bounds:\n"
                                        + "\n".join(ob_details)
                                    )
                                if new_overlaps:
                                    ol_details = []
                                    for a, b, area in new_overlaps:
                                        a_blk = next(
                                            (bl for bl in current_st.blocks
                                             if bl.var_name == a), None,
                                        )
                                        b_blk = next(
                                            (bl for bl in current_st.blocks
                                             if bl.var_name == b), None,
                                        )
                                        a_desc = (
                                            f"{a} (y={a_blk.y:.2f}→"
                                            f"{a_blk.y+a_blk.h:.2f})"
                                            if a_blk else a
                                        )
                                        b_desc = (
                                            f"{b} (y={b_blk.y:.2f}→"
                                            f"{b_blk.y+b_blk.h:.2f})"
                                            if b_blk else b
                                        )
                                        ol_details.append(
                                            f"  {a_desc} ↔ {b_desc} "
                                            f"(overlap: {area:.2f} sq in)"
                                        )
                                    warn_parts.append(
                                        "🚨 SUBMIT BLOCKED — your edits "
                                        "introduced NEW element overlaps:\n"
                                        + "\n".join(ol_details)
                                        + "\n  Move or resize elements so "
                                        "they don't overlap."
                                    )
                                if near_overflow:
                                    no_details = []
                                    for b in near_overflow:
                                        no_details.append(
                                            f"  {b.var_name}: "
                                            f"{b.text_chars}ch in "
                                            f"{b.w:.1f}\"×{b.h:.1f}\" "
                                            f"(overflow: {b.overflow_bottom_px}px) "
                                            f"— newly overflowing"
                                        )
                                    warn_parts.append(
                                        "⚠ NEAR OVERFLOW WARNING — content "
                                        "grew and containers are near "
                                        "capacity:\n"
                                        + "\n".join(no_details)
                                        + "\n  Expand these boxes by at "
                                        "least 0.3\" height or shorten "
                                        "the text to avoid visual "
                                        "clipping."
                                    )
                                if has_hard_issues:
                                    # Check net improvement: if total issues decreased,
                                    # allow submit even with some new issues
                                    total_before = len(t0_of) + len(t0_ob) + len(t0_overlap_set)
                                    total_after = len(t1_of) + len(t1_ob) + len(current_st.overlap_pairs)
                                    net_improved = total_after < total_before
                                    if net_improved:
                                        # Net improvement — warn but allow submit
                                        logger.info(
                                            "Agent slide %d: submit allowed with net improvement "
                                            "(issues %d -> %d, %d new but %d fixed)",
                                            slide_id, total_before, total_after,
                                            len(new_of) + len(new_ob) + len(new_overlaps),
                                            total_before - total_after + len(new_of) + len(new_ob) + len(new_overlaps),
                                        )
                                        # Fall through to submit
                                    else:
                                        # No net improvement — block submit
                                        warn_parts.append(
                                            "\nFix these issues before submitting. "
                                            "Either shorten text, expand boxes, or "
                                            "move elements back within slide bounds. "
                                            "Then call submit again."
                                        )
                                        submit_bounced_spatial += 1
                                        messages.append({
                                            "role": "user",
                                            "content": "\n".join(warn_parts),
                                        })
                                        tool_calls += 1
                                        logger.info(
                                            "Agent slide %d: submit bounced "
                                            "(spatial check, content=%d spatial=%d): "
                                            "%d new overflow, %d new OOB, "
                                            "%d new overlaps, %d near-overflow",
                                            slide_id, submit_bounced_content,
                                            submit_bounced_spatial,
                                            len(new_of), len(new_ob),
                                            len(new_overlaps), len(near_overflow),
                                        )
                                        continue
                                else:
                                    # Near-overflow only: warn but allow submit
                                    logger.info(
                                        "Agent slide %d: submit accepted "
                                        "with near-overflow warning (%d blocks)",
                                        slide_id, len(near_overflow),
                                    )
                                    # Fall through to content accuracy check and submit
                    except Exception:
                        pass  # Don't block submit on analysis errors

                # Contrast regression gate: block submit if text becomes
                # invisible (contrast drops below 3:1 WCAG AA minimum).
                if state.current_code != code and submit_bounced_spatial < 1:
                    contrast_warning = self._check_contrast_regression(
                        state.original_code, state.current_code, slide_id
                    )
                    if contrast_warning and "SUBMIT BLOCKED" in contrast_warning:
                        submit_bounced_spatial += 1
                        logger.info(
                            "Agent slide %d: submit bounced (contrast regression, "
                            "bounce %d/%d)",
                            slide_id, submit_bounced_content,
                            submit_bounced_spatial,
                        )
                        result = contrast_warning
                        code_changed = False
                        messages.append({"role": "user", "content": result})
                        tool_calls += 1
                        continue
                    elif contrast_warning:
                        logger.info(
                            "Agent slide %d: contrast regression warning (soft, not blocking)",
                            slide_id,
                        )

                # --- Process-note / meta-content leak gate ---
                if submit_bounced_content < 3:
                    import re as _re_pn
                    _PROCESS_PATTERNS = [
                        r'Refocus the slide', r'This slide should',
                        r'Restructure the', r'Move the content',
                        r'Rewrite this', r'Compress the text',
                        r'Reduce the font', r'The goal is to',
                        r'Plan:', r'Step \d+:',
                        # Meta-content patterns (instruction text leaked into slide)
                        r'REMOVE\s*[-—–:]',
                        r'Replace\s+with\s*:',
                        r'Suggested\s+replacement',
                        r'No verified replacement',
                        r'Source-supported',
                        r'\bTODO\b',
                        r'\bFIXME\b',
                        r'\bPlaceholder\b',
                        # Evaluator judgment/correction language
                        r'corrected to source[- ]grounded',
                        r'not supported by (?:the )?(?:paper|source|evidence)',
                        r'no source support',
                        r'unsupported (?:claim|text|content|number)',
                        r'Correction:\s',
                    ]
                    bounce_msg = None
                    for pat in _PROCESS_PATTERNS:
                        if (_re_pn.search(pat, state.current_code, _re_pn.IGNORECASE)
                                and not _re_pn.search(pat, state.original_code, _re_pn.IGNORECASE)):
                            submit_bounced_content += 1
                            leak_match = _re_pn.search(pat, state.current_code, _re_pn.IGNORECASE).group()
                            bounce_msg = (
                                f"Meta-instruction text leaked into slide content: '{leak_match}'. "
                                f"This is NOT presentation content — remove it from the slide text."
                            )
                            messages.append({"role": "user", "content": bounce_msg})
                            break
                    if bounce_msg:
                        continue

                # Pre-submit content verification: if agent has fabricated/
                # incorrect_claim issues but never called search_source,
                # bounce once to force verification from paper.
                # EXCEPTION: if ALL such issues have judge-verified fix_detail,
                # the agent can use that data directly — no search needed.
                has_fabricated = any(
                    i.issue_type in CRITICAL_CONTENT_TYPES
                    and not (i.fix_detail and i.fix_detail.correct_content)
                    for i in self._current_issues
                    if hasattr(i, 'status') and i.status.value == "open"
                )
                if (state.current_code != state.original_code
                        and state.search_calls_used == 0
                        and self._is_html_code(state.current_code)):
                    orig_text = self._extract_visible_text(state.original_code)
                    curr_text = self._extract_visible_text(state.current_code)
                    orig_words = set(orig_text.split())
                    curr_words = set(curr_text.split())
                    added_words = curr_words - orig_words
                    if has_fabricated and submit_bounced_content < 2:
                        # Hard bounce for fabricated issues without source check
                        submit_bounced_content += 1
                        messages.append({"role": "user", "content":
                            "⚠ SUBMIT BLOCKED — you have fabricated/incorrect_claim "
                            "issues but never called search_source or lookup_table "
                            "to verify correct content from the paper. You MUST "
                            "call search_source to find the correct data before "
                            "submitting. Replacing numbers with descriptive text "
                            "is NOT acceptable — find the real numbers."
                        })
                        tool_calls += 1
                        logger.info(
                            "Agent slide %d: submit bounced "
                            "(fabricated without search_source)",
                            slide_id,
                        )
                        continue
                    elif (
                        len(added_words) >= 10
                        and submit_bounced_content < 2
                        and not state.allow_support_copy_compression
                        and (
                            state.allow_visible_text_change
                            or any(
                                not (i.rubric_id or "").startswith("B")
                                for i in self._current_issues
                            )
                        )
                    ):
                        submit_bounced_content += 1
                        messages.append({"role": "user", "content":
                            "⚠ SUBMIT BLOCKED — you added significant new "
                            "visible text but never called search_source or "
                            "lookup_table. For C/D/E fixes and B14/B12 text "
                            "exceptions, every new claim must come from the "
                            "paper or from the issue's exact correct_content. "
                            "Do not introduce limitation/scope/causal wording "
                            "such as 'limitation', 'only', 'depends on', or "
                            "'remains' unless the source explicitly says it. "
                            "Either call search_source and revise the text, or "
                            "remove the new unsupported wording and keep only "
                            "the source-backed retained content."
                        })
                        tool_calls += 1
                        logger.info(
                            "Agent slide %d: submit bounced "
                            "(significant new text without source search: %d words)",
                            slide_id, len(added_words),
                        )
                        continue
                    elif len(added_words) >= 10:
                        logger.info(
                            "Agent slide %d: content verification warning "
                            "(%d new words, 0 searches)",
                            slide_id, len(added_words),
                        )
                        # Soft warning only — do NOT bounce/block submit

                # Gate: force verify_layout after content edits
                # If agent made edits but never ran verify_layout afterwards,
                # bounce once to force spatial regression check.
                if (has_content_edit_since_verify
                        and submit_bounced_spatial < 2
                        and state.current_code != state.original_code
                        and self._is_html_code(state.current_code)):
                    submit_bounced_spatial += 1
                    messages.append({"role": "user", "content":
                        "⚠ SUBMIT BLOCKED — you made code edits but did not "
                        "call verify_layout afterwards. Content and layout "
                        "changes often cause text_overflow or density_imbalance "
                        "regressions. Call verify_layout now to check for "
                        "spatial issues, fix any problems found, then submit."
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced (no verify after edits)",
                        slide_id,
                    )
                    continue

                # (G10 deleted: high overflow gate used unreliable char-based
                # estimation with 2% precision. Playwright-based overflow
                # detection in G4 spatial regression check is accurate.)

                # (G11 deleted: density coverage gate contradicted eval's
                # container_contract_breach — forcing elements larger caused
                # new eval issues. Coverage is a soft metric, not a gate.)

                # ── Plan completion gate (soft bounce, once) ──
                pending_steps = [
                    s for s in state.plan_steps
                    if s.status == "pending"
                ]
                if (pending_steps
                        and not getattr(state, '_plan_submit_warned', False)):
                    state._plan_submit_warned = True
                    step_list = "\n".join(
                        f"  - Step {state.plan_steps.index(s) + 1}: {s.text}"
                        for s in pending_steps
                    )
                    messages.append({"role": "user", "content":
                        f"{len(pending_steps)} plan step(s) still pending:\n"
                        f"{step_list}\n\n"
                        "If these are intentionally unaddressed, call "
                        "update_plan to mark them as skipped (with a "
                        "reason), then submit again. If you still need "
                        "to work on them, continue with your plan."
                    })
                    tool_calls += 1
                    logger.info(
                        "Agent slide %d: submit bounced "
                        "(%d pending plan steps)",
                        slide_id, len(pending_steps),
                    )
                    continue

                logger.info(
                    "Agent repair slide %d: submitted after %d tool calls",
                    slide_id, tool_calls,
                )
                state.submitted = True
                break

            # Execute tool
            code_before_action = state.current_code
            result, code_changed = self._execute_tool(action, state)
            pure_svg_repair = all(
                issue.issue_type == "svg_visual_defect" for issue in all_issues
            )
            if (
                code_changed
                and has_svg_visual_issue
                and (
                    pure_svg_repair
                    or self._edit_may_affect_svg(action, code_before_action)
                )
            ):
                state._svg_edit_since_preview = True
            if code_changed and has_raw_figure_issue and self._is_html_code(state.current_code):
                state._raw_figure_edit_since_preview = True
            if tool_name in {"render_preview", "verify_layout"} and isinstance(result, list):
                state._svg_edit_since_preview = False
                state._raw_figure_edit_since_preview = False
            # result may be str (normal) or list (multimodal, e.g. render_preview)
            if isinstance(result, list):
                log_preview = "[multimodal: image + text]"
            else:
                log_preview = result[:200]

            ignored_actions = getattr(self, "_multi_action_ignored_count", 0)
            if ignored_actions:
                warning = (
                    f"\n\nSEQUENCING: ignored {ignored_actions} additional tool call(s) "
                    "from your previous response. Inspect this actual result, then return "
                    "exactly one next JSON tool call."
                )
                if isinstance(result, list):
                    result.append({"type": "text", "text": warning})
                else:
                    result += warning
                self._multi_action_ignored_count = 0
            logger.info(
                "Agent slide %d turn %d: result=%s changed=%s",
                slide_id, tool_calls, log_preview, code_changed,
            )

            # ── Inject plan progress after every tool result ──
            plan_progress = self._format_plan_progress(state)
            if plan_progress:
                if isinstance(result, list):
                    # Multimodal: append to last text block
                    for i in range(len(result) - 1, -1, -1):
                        if isinstance(result[i], dict) and result[i].get("type") == "text":
                            result[i]["text"] += plan_progress
                            break
                    else:
                        result.append({"type": "text", "text": plan_progress})
                else:
                    result += plan_progress

            # Sanitize content before appending — ensure list content has only dicts
            if isinstance(result, list):
                sanitized = []
                for item in result:
                    if isinstance(item, dict):
                        sanitized.append(item)
                    elif isinstance(item, str):
                        sanitized.append({"type": "text", "text": item})
                result = sanitized

            messages.append({"role": "user", "content": result})

            # ── Save HTML snapshot after each code-changing action ──
            if code_changed and run_dir and self._is_html_code(state.current_code):
                snap_dir = Path(run_dir) / f"turn_{turn_index:02d}" / "snapshots" / f"slide_{slide_id:02d}"
                snap_dir.mkdir(parents=True, exist_ok=True)
                snap_idx = len(list(snap_dir.glob("step_*.html")))
                snap_path = snap_dir / f"step_{snap_idx:02d}_{tool_name}.html"
                snap_path.write_text(state.current_code, encoding="utf-8")
                logger.debug("Snapshot saved: %s", snap_path)
            # regen_slide is high-cost: count as 5 tool calls
            # plan/update_plan/submit_repair_summary are admin — don't count toward budget
            admin_tools = {"plan", "update_plan", "submit_repair_summary"}
            if tool_name in admin_tools:
                pass  # don't increment tool_calls for admin actions
            elif tool_name == "regen_slide":
                tool_calls += 5
            else:
                tool_calls += 1

            # ── Structured verify response ──
            # After verify_layout, track defect trajectory for escalation.
            if tool_name == "verify_layout":
                has_content_edit_since_verify = False  # reset flag

                # Track verify defect counts for stagnation detection
                verify_defect_count = 0
                targeted_categories = self._targeted_residual_categories()
                significant = {}
                if hasattr(state, '_last_html_state') and state._last_html_state:
                    hs = state._last_html_state
                    # Use the FILTERED single-source-of-truth count, not raw len().
                    # Two reasons: (1) it includes clipped_blocks + canvas
                    # truncation, so a clip-only residual can still trigger
                    # escalation (a previously invisible failure mode);
                    # (2) raw len() counted sub-threshold noise the scorer ignores,
                    # which could keep the trajectory "hot" and fire the
                    # destructive condensation escalation on non-issues.
                    significant = self._targeted_significant_issues(
                        getattr(state, '_t0_html_state', None), hs,
                    )
                    verify_defect_count = sum(
                        len(significant.get(category, []))
                        for category in targeted_categories
                    )
                if not hasattr(state, '_verify_defect_history'):
                    state._verify_defect_history = []
                state._verify_defect_history.append(verify_defect_count)
                targeted_snapshot = {
                    category: significant.get(category, [])
                    for category in sorted(targeted_categories)
                }
                verify_signature = json.dumps(
                    targeted_snapshot,
                    ensure_ascii=True,
                    sort_keys=True,
                    default=str,
                )
                if not hasattr(state, '_verify_signature_history'):
                    state._verify_signature_history = []
                state._verify_signature_history.append(verify_signature)

                # Reconsider only on persistent identities or strict worsening.
                # Equal counts with changing targets can be a legitimate
                # intermediate state in a coupled multi-edit reflow.
                hist = state._verify_defect_history
                signature_hist = state._verify_signature_history
                if self._verify_needs_strategy_reconsideration(
                    hist, signature_hist,
                ):
                    reconsideration_key = tuple(signature_hist[-3:])
                    already_reported = getattr(
                        state, "_last_strategy_reconsideration_key", None,
                    ) == reconsideration_key
                    if not already_reported:
                        state._last_strategy_reconsideration_key = reconsideration_key
                        escalation_msg = (
                            "\n\nSTRATEGY RECONSIDERATION\n"
                            f"The recent verify_layout evidence indicates that the current "
                            f"execution needs a causal review (targeted trajectory: "
                            f"{hist[-3]}→{hist[-2]}→{hist[-1]}). This does not by itself "
                            f"show that the current topology is wrong. Use the current revision's "
                            f"layout anchors, relation/space map, detector details, and original "
                            f"issue to identify the unresolved owning region. First ask whether "
                            f"the edits between these verifies actually changed that region; an "
                            f"edit that restored hierarchy/content or handled another region may "
                            f"legitimately leave the same residual unchanged. If the reading path "
                            f"and grouping remain coherent, finish the relevant role-aware "
                            f"calibration before switching. Changing peer orientation or grouping "
                            f"(for example, a vertical stack into side-by-side columns) is reflow, "
                            f"not same-topology calibration. Choose it only when current spatial "
                            f"evidence shows that the existing organization itself causes or moves "
                            f"the pressure."
                        )
                        # Append to the last tool result message
                        if messages and messages[-1]["role"] == "user":
                            messages[-1]["content"] += escalation_msg
                        else:
                            messages.append({"role": "user", "content": escalation_msg})
                        logger.info(
                            "Agent slide %d: verify stagnation detected "
                            "(trajectory: %s), injecting escalation guidance",
                            slide_id, hist[-3:],
                        )

                signatures = state._verify_signature_history
                if (
                    verify_defect_count > 0
                    and len(signatures) >= 4
                    and len(set(signatures[-4:])) == 1
                ):
                    persistent_key = signatures[-1]
                    if getattr(
                        state, "_last_persistent_residual_advisory", None,
                    ) != persistent_key:
                        state._last_persistent_residual_advisory = persistent_key
                        advisory = (
                            "\n\nPERSISTENT RESIDUAL ADVISORY\n"
                            "Four verifies show the same targeted identities. This is "
                            "evidence that those objects are not yet closed, but it is not "
                            "evidence by itself that the current topology has failed and it is "
                            "not an automatic stop or rollback. Check whether the verified edits "
                            "actually changed the residual-owning region. If they restored "
                            "hierarchy/content, handled another region, or are an unfinished part "
                            "of one coherent allocation, complete and verify the relevant repair. "
                            "Otherwise reconsider the spatial cause instead of continuing smaller "
                            "versions of the same fit edit. Treat any change to peer orientation "
                            "or semantic grouping as reflow and justify it from current spatial "
                            "evidence."
                        )
                        if messages and messages[-1]["role"] == "user":
                            messages[-1]["content"] += advisory
                        else:
                            messages.append({"role": "user", "content": advisory})
                        logger.info(
                            "Agent slide %d: four identical targeted verify states; "
                            "injected advisory without forcing loop exit",
                            slide_id,
                        )

            # ── Reset verify failure counter on successful edits ──
            if tool_name in ("apply_edits", "delete_shape") and code_changed:
                has_content_edit_since_verify = True  # any edit needs verify
                # Don't reset counter here — wait for verify_layout
                # to confirm the edit was actually beneficial.
                pass

            # Track progress — only count apply_edits/delete_shape failures.
            # verify_layout, get_current_code, plan, submit, rollback are
            # all legitimate non-edit actions in a complex repair and should
            # not count as "no progress".
            edit_tools = {"apply_edits", "delete_shape"}
            tool_name = action.get("tool", "")
            if tool_name in edit_tools:
                if code_changed:
                    no_progress = 0
                else:
                    no_progress += 1

            if no_progress >= self.MAX_NO_PROGRESS:
                logger.info(
                    "Agent repair slide %d: no progress for %d turns, aborting",
                    slide_id, no_progress,
                )
                break

            continuation = self._trajectory_continuation_message(
                state,
                tool_name=tool_name,
                code_changed=code_changed,
                tool_calls=tool_calls,
                soft_limit=max_tool_calls,
            )
            if continuation:
                previous_limit = max_tool_calls
                max_tool_calls = min(
                    self.MAX_TOOL_CALLS_CAP,
                    max_tool_calls + self.TRAJECTORY_EXTENSION_CALLS,
                )
                state.trajectory_extensions += 1
                state.last_trajectory_extension_revision = state.layout_revision
                if messages and messages[-1]["role"] == "user":
                    messages[-1]["content"] += continuation
                else:
                    messages.append({"role": "user", "content": continuation})
                logger.info(
                    "Agent slide %d: extended active trajectory budget %d->%d "
                    "after %s at revision %d",
                    slide_id,
                    previous_limit,
                    max_tool_calls,
                    tool_name,
                    state.layout_revision,
                )

        logger.info(
            "Agent repair slide %d: loop ended after %d tool calls, "
            "code %s",
            slide_id, tool_calls,
            "changed" if state.current_code != state.original_code else "unchanged",
        )

        # The final allowed tool call can be a successful edit. Measure that
        # candidate before timeout fallback so a useful, non-regressing state is
        # not discarded merely because the model had no call left for verify.
        if (
            not state.submitted
            and state.current_code != state.original_code
            and not self._has_current_verify(state)
            and not self._disable_step_render
        ):
            try:
                self._tool_verify_layout(state)
                logger.info(
                    "Agent slide %d: deterministically verified final unmeasured "
                    "candidate before timeout fallback",
                    slide_id,
                )
            except Exception as exc:
                logger.warning(
                    "Agent slide %d: final deterministic verification failed: %s",
                    slide_id,
                    str(exc)[:160],
                )

        # Loop-timeout fallback preserves the latest hard-valid checkpoint.  A
        # detector-minimal state is not automatically visually better, and using
        # it as the first fallback erased coherent multi-edit composition work.
        if not state.submitted and state.current_code != state.original_code:
            rollback_target = None
            target_label = None
            current_is_safe = bool(
                state.latest_safe_verified_code == state.current_code
                and state.latest_safe_verified_revision == state.layout_revision
            )
            if current_is_safe:
                target_label = "current hard-valid verified checkpoint"
            elif state.latest_safe_verified_code is not None:
                rollback_target = state.latest_safe_verified_code
                target_label = "latest hard-valid verified checkpoint"
            elif state.last_verified_code is not None and state.last_verified_code != state.current_code:
                rollback_target = state.last_verified_code
                target_label = "last non-regressing verified"
            elif state.best_verified_code is not None and state.best_verified_code != state.current_code:
                rollback_target = state.best_verified_code
                target_label = f"best verified ({state.best_verified_issues} issues)"
            if rollback_target is not None:
                logger.info(
                    "Agent slide %d: loop timeout — rolling back to %s "
                    "(discarding %d chars of unverified edits)",
                    slide_id, target_label,
                    abs(len(state.current_code) - len(rollback_target)),
                )
                state.current_code = rollback_target
            elif not current_is_safe and (
                state.last_verified_code is None
                and state.best_verified_code is None
                and state.latest_safe_verified_code is None
            ):
                logger.info(
                    "Agent slide %d: loop timeout, no verified checkpoint "
                    "— rolling back to original code",
                    slide_id,
                )
                state.current_code = state.original_code

        # Best-state preference on SUBMIT: if the agent submitted a state that is
        # strictly worse than a verified state it had already achieved, ship the
        # better one. The agent's judgment about WHAT to keep is respected; we only
        # protect it from a final-edit regression it didn't intend. One extra count
        # here (on the shipped artifact) is cheap vs. the whole LLM loop.
        if (state.submitted
                and state.best_verified_code is not None
                and state.best_verified_code != state.current_code
                and not (
                    state.latest_safe_verified_code == state.current_code
                    and state.latest_safe_verified_revision == state.layout_revision
                )):
            try:
                from .html_spatial_state import extract_html_slide_state
                if self._is_html_code(state.current_code):
                    html_asset_base_dirs = self._html_asset_base_dirs(
                        state.case_dir,
                        getattr(state, "_run_dir", None),
                        getattr(state, "_turn_index", None),
                    )
                    cur_state = extract_html_slide_state(
                        slide_id,
                        state.current_code,
                        html_base_dir=Path(state.case_dir),
                        asset_base_dirs=html_asset_base_dirs,
                    )
                    baseline_state = getattr(state, '_t0_html_state', None)
                    if baseline_state is None:
                        baseline_state = extract_html_slide_state(
                            slide_id,
                            state.original_code,
                            html_base_dir=Path(state.case_dir),
                            asset_base_dirs=html_asset_base_dirs,
                        )
                    current_significant = self._targeted_significant_issues(
                        baseline_state, cur_state,
                    )
                    cur_issues = sum(len(items) for items in current_significant.values())
                    if (state.best_verified_issues is not None
                            and state.best_verified_issues < cur_issues):
                        logger.info(
                            "Agent slide %d: submitted state has %d issues but a "
                            "verified state had %d — shipping the better one.",
                            slide_id, cur_issues, state.best_verified_issues,
                        )
                        state.current_code = state.best_verified_code
            except Exception as e:
                logger.warning(
                    "Agent slide %d: best-state submit check failed: %s",
                    slide_id, str(e)[:120],
                )

        # Save conversation log for debugging
        # Priority: run_dir (always save), then REPAIR_LOG_DIR env var
        log_dir = None
        if run_dir:
            log_dir = str(Path(run_dir) / f"turn_{turn_index:02d}" / "repair_logs")
        if not log_dir:
            log_dir = os.environ.get("REPAIR_LOG_DIR")
        if log_dir:
            log_path = Path(log_dir) / f"slide_{slide_id:02d}_attempt_{attempt}.json"
            log_path.parent.mkdir(parents=True, exist_ok=True)
            # Redact base64 images for log readability
            log_messages = []
            for m in messages:
                content = m.get("content", "")
                if isinstance(content, list):
                    # Multimodal — redact image data
                    redacted = []
                    for block in content:
                        if isinstance(block, dict) and block.get("type") == "image_url":
                            redacted.append({"type": "image_url", "image_url": {"url": "[base64 image redacted]"}})
                        else:
                            redacted.append(block)
                    log_messages.append({"role": m["role"], "content": redacted})
                elif isinstance(content, str) and len(content) > 20000:
                    # Keep both the repair brief and the late spatial maps.  The
                    # latter are deliberately appended after compact DOM output,
                    # so prefix-only truncation makes a successful diagnostic
                    # delivery look absent when reviewing an agent trajectory.
                    head_chars = 12000
                    tail_chars = 8000
                    omitted = len(content) - head_chars - tail_chars
                    log_messages.append({
                        "role": m["role"],
                        "content": (
                            content[:head_chars]
                            + f"\n... [middle truncated, {omitted} chars omitted; "
                              f"total {len(content)} chars] ...\n"
                            + content[-tail_chars:]
                        ),
                    })
                else:
                    log_messages.append(m)
            try:
                with open(log_path, "w") as f:
                    json.dump(log_messages, f, indent=2, ensure_ascii=False, default=str)
                logger.info("Agent slide %d: saved conversation log to %s", slide_id, log_path)
            except Exception as e:
                logger.warning("Agent slide %d: failed to save log: %s", slide_id, str(e)[:100])

        # Render snapshots to PNG for visual debugging
        if run_dir:
            snap_dir = Path(run_dir) / f"turn_{turn_index:02d}" / "snapshots" / f"slide_{slide_id:02d}"
            if snap_dir.exists():
                html_files = sorted(snap_dir.glob("step_*.html"))
                if html_files:
                    try:
                        from ...render_backends.playwright_backend import PlaywrightRenderBackend
                        backend = PlaywrightRenderBackend()
                        for html_file in html_files:
                            png_file = html_file.with_suffix(".png")
                            backend.render_html_file_to_png(html_file, png_file)
                        backend.close()
                        logger.info("Agent slide %d: rendered %d snapshots", slide_id, len(html_files))
                    except Exception as e:
                        logger.warning("Agent slide %d: snapshot render failed: %s", slide_id, str(e)[:200])

        self._record_last_repair_result(state, all_issues)

        # 5. Validate result
        if state.current_code == state.original_code:
            logger.warning(
                "Agent repair slide %d: code unchanged after %d tool calls "
                "(current_code len=%d, original_code len=%d)",
                slide_id, tool_calls, len(state.current_code), len(state.original_code),
            )
            return None

        if (
            visual_only_repair
            and self._is_html_code(state.current_code)
        ):
            scope_ok, scope_reason = validate_visual_repair_scope(
                state.original_code,
                state.current_code,
                allow_image_replacement=visual_repair_may_replace_image,
                allow_text_formatting_change=(
                    visual_repair_may_change_formatting
                ),
                allow_text_content_change=visual_repair_may_change_text,
            )
            if not scope_ok:
                logger.info(
                    "Agent repair slide %d: visual scope violation after loop "
                    "(%s)", slide_id, scope_reason,
                )
                return None

        # Content retention check
        threshold = self._get_retention_threshold(all_issues)
        if not self._check_content_retention(
            state.original_code, state.current_code, threshold,
        ):
            logger.info(
                "Agent repair slide %d: content retention too low "
                "(%.0f%%, threshold %.0f%%)",
                slide_id, self._last_retention * 100, threshold * 100,
            )
            return None

        # 7. Aesthetic quality check — reject on severe regressions
        #    (color contrast violations, massive content loss), but allow
        #    through if issues are minor (single overflow, small gap).
        #    Skip for external HTML repair (bp_slide=None) — the aesthetic
        #    check uses PPTX coordinate system which gives false positives
        #    on HTML slides (wrong unit: inches vs pixels).
        if state.bp_slide is not None:
            aes_ok, aes_reason = self._aesthetic_quality_check(
                code, state.current_code, slide_id,
            )
        else:
            aes_ok, aes_reason = True, ""
        if not aes_ok:
            # Count severity: color contrast is a hard reject,
            # other issues are soft warnings
            contrast_issues = [
                line for line in aes_reason.split("\n")
                if "low contrast" in line
            ]
            content_loss = "CONTENT CHANGE" in aes_reason
            overflow_regressions = [
                line for line in aes_reason.split("\n")
                if "NEW TEXT OVERFLOW" in line
            ]
            oob_regressions = [
                line for line in aes_reason.split("\n")
                if "NEW OOB" in line
            ]
            overlap_regressions = [
                line for line in aes_reason.split("\n")
                if "NEW OVERLAP" in line
            ]
            if len(contrast_issues) >= 2:
                # Multiple text elements with unreadable contrast = reject
                logger.warning(
                    "Agent repair slide %d: REJECTED — severe color "
                    "contrast regression (%d elements): %s",
                    slide_id, len(contrast_issues),
                    aes_reason[:300],
                )
                return None
            if content_loss:
                # For slides with critical content accuracy issues,
                # significant content change is expected behavior — the
                # fabricated/incorrect content MUST be replaced.
                critical_content_types = CRITICAL_CONTENT_TYPES
                has_critical_content = any(
                    i.issue_type in critical_content_types
                    for i in all_issues
                    if i.status.value == "open"
                )
                if has_critical_content:
                    logger.info(
                        "Agent repair slide %d: content change accepted "
                        "(critical content accuracy issues present): %s",
                        slide_id, aes_reason[:200],
                    )
                else:
                    logger.warning(
                        "Agent repair slide %d: REJECTED — significant "
                        "content loss: %s",
                        slide_id, aes_reason[:300],
                    )
                    return None
            if len(overflow_regressions) >= 2:
                # New overflow blocks = hard reject (tightened from ≥3)
                logger.warning(
                    "Agent repair slide %d: REJECTED — %d new text "
                    "overflow blocks introduced: %s",
                    slide_id, len(overflow_regressions),
                    aes_reason[:300],
                )
                return None
            if oob_regressions:
                # Any new out-of-bounds element = reject
                logger.warning(
                    "Agent repair slide %d: REJECTED — new out-of-"
                    "bounds elements: %s",
                    slide_id, aes_reason[:300],
                )
                return None
            if overlap_regressions:
                # Any new overlaps = hard reject (tightened from ≥2).
                # Overlap is always a critical regression — better to
                # reject the repair than introduce overlapping elements.
                logger.warning(
                    "Agent repair slide %d: REJECTED — %d new element "
                    "overlaps introduced: %s",
                    slide_id, len(overlap_regressions),
                    aes_reason[:300],
                )
                return None
            # Other aesthetic issues: warn but accept
            logger.warning(
                "Agent repair slide %d: aesthetic issues in final "
                "code (accepted): %s",
                slide_id, aes_reason[:300],
            )

        # --- Diff ratio gate: reject if repair changed too much ---
        if state.current_code != state.original_code:
            import difflib
            orig_lines = state.original_code.splitlines()
            new_lines = state.current_code.splitlines()
            matcher = difflib.SequenceMatcher(None, orig_lines, new_lines)
            ratio = matcher.ratio()  # 1.0 = identical, 0.0 = completely different
            changed_pct = (1.0 - ratio) * 100
            if changed_pct > 60:
                logger.warning(
                    "Agent repair slide %d: REJECTED — diff ratio %.0f%% "
                    "exceeds 60%% threshold (repair changed too much). "
                    "Reverting to original code.",
                    slide_id, changed_pct,
                )
                return None
            elif changed_pct > 40:
                logger.warning(
                    "Agent repair slide %d: high diff ratio %.0f%% "
                    "(close to 60%% rejection threshold)",
                    slide_id, changed_pct,
                )

        return state.current_code

    # ================================================================
    # TOOL EXECUTION
    # ================================================================

    def _execute_tool(
        self, action: dict, state: AgentState,
    ) -> tuple[str | list, bool]:
        """Execute a tool call, return (result_text, code_changed).

        result_text may be a string or a list of content blocks (multimodal).
        """
        tool = action.get("tool", "")

        if tool == "apply_edits":
            # P0-1: Intercept bulk coordinate edits when layout issues exist.
            # The agent should use reflow_layout instead of manually editing
            # many Inches() values, which causes 70%+ regression rate.
            # Only block when edits are primarily coordinate-focused (≥3 changed
            # Inches values AND majority of edits involve Inches changes).
            if state.issue_types & LAYOUT_REPAIR_TYPES:
                import re as _re
                edits = action.get("edits", [])
                inches_changes = 0
                edits_with_inches = 0
                for edit in edits:
                    search = edit.get("search", "")
                    replace = edit.get("replace", "")
                    if search and replace:
                        old_inches = _re.findall(r'Inches\([\d.]+\)', search)
                        new_inches = _re.findall(r'Inches\([\d.]+\)', replace)
                        if old_inches and new_inches:
                            has_change = False
                            for o, n in zip(old_inches, new_inches):
                                if o != n:
                                    inches_changes += 1
                                    has_change = True
                            if has_change:
                                edits_with_inches += 1
                # Only reject if ≥3 Inches changes AND most edits are coordinate-focused
                if inches_changes >= 3 and edits_with_inches >= len(edits) * 0.5:
                    return (
                        "❌ REJECTED: This edit changes %d coordinate values. "
                        "For layout/density issues, use the `reflow_layout` tool "
                        "instead — it handles coordinate changes atomically and "
                        "prevents regressions. Example: "
                        '{"tool": "reflow_layout"}'
                        % inches_changes
                    ), False
            return self._tool_apply_edits(action, state)

        if tool == "apply_css_patch":
            return self._tool_apply_css_patch(action, state)

        # Dispatch table for simple tool routing
        _TOOL_DISPATCH: dict[str, tuple] = {
            "verify_layout": (self._tool_verify_layout, "state_only"),
            "rollback": (self._tool_rollback, "action_state"),
            "delete_shape": (self._tool_delete_shape, "action_state"),
            "reflow_layout": (self._tool_reflow_layout, "action_state"),
            "get_current_code": (self._tool_get_current_code, "state_only"),
            "measure_space": (self._tool_measure_space, "state_only"),
            "plan": (self._tool_plan, "action_state"),
            "update_plan": (self._tool_update_plan, "action_state"),
            "search_source": (self._tool_search_source, "action_state"),
            "lookup_table": (self._tool_lookup_table, "action_state"),
            "crop_image": (self._tool_crop_image, "action_state"),
            "compose_image_grid": (self._tool_compose_image_grid, "action_state"),
            "generate_chart": (self._tool_generate_chart, "action_state"),
            "create_svg_asset": (self._tool_create_svg_asset, "action_state"),
            "regen_slide": (self._tool_regen_slide, "action_state"),
        }

        if tool == "render_preview":
            if not self._enable_render_preview:
                return (
                    "render_preview is unavailable for this repair. "
                    "Use verify_layout for compile and spatial checks, but do not "
                    "claim a high-confidence visual/composition pass from this alone. "
                    "If the original issue is subjective, cite the strongest spatial "
                    "evidence you have or mark the result uncertain.",
                    False,
                )
            return self._tool_render_preview(state)

        if tool == "submit_repair_summary":
            return self._tool_submit_repair_summary(
                action, state, state.slide_id,
                getattr(state, '_run_dir', None),
                getattr(state, '_turn_index', 0),
            )

        handler = _TOOL_DISPATCH.get(tool)
        if handler:
            fn, sig = handler
            if sig == "state_only":
                return fn(state)
            else:
                return fn(action, state)

        available = ", ".join([
            "apply_edits", "apply_css_patch", *_TOOL_DISPATCH.keys(),
            "submit_repair_summary", "submit",
        ])
        return f"Unknown tool: {tool}. Available: {available}", False

    @staticmethod
    def _normalize_checkpoint_metadata(state: AgentState) -> None:
        """Keep optional checkpoint metadata aligned with legacy state fixtures."""
        if not hasattr(state, "checkpoint_text_loss"):
            state.checkpoint_text_loss = [0] * len(state.checkpoints)
        if not hasattr(state, "checkpoint_labels"):
            state.checkpoint_labels = ["legacy checkpoint"] * len(state.checkpoints)
        if not hasattr(state, "current_checkpoint_label"):
            state.current_checkpoint_label = "legacy checkpoint"
        defaults = {
            "pending_edit_cluster": False,
            "pending_edit_scopes": [],
            "last_edit_scope": (),
            "active_cluster_start_code": None,
            "active_cluster_start_text_loss": 0,
            "active_cluster_start_label": "",
            "last_cluster_start_code": None,
            "last_cluster_start_text_loss": 0,
            "last_cluster_start_label": "",
        }
        for name, value in defaults.items():
            if not hasattr(state, name):
                setattr(state, name, value.copy() if isinstance(value, list) else value)
        while len(state.checkpoint_text_loss) < len(state.checkpoints):
            state.checkpoint_text_loss.append(0)
        while len(state.checkpoint_labels) < len(state.checkpoints):
            state.checkpoint_labels.append("legacy checkpoint")

    def _commit_html_code_change(
        self,
        state: AgentState,
        new_code: str,
        *,
        projected_text_loss: int,
        edit_scopes: list[str],
        cluster_complete: bool,
        action_label: str,
    ) -> None:
        """Record one HTML edit while preserving a recoverable cluster boundary."""
        self._normalize_checkpoint_metadata(state)
        was_pending_cluster = state.pending_edit_cluster
        pre_edit_code = state.current_code
        pre_edit_text_loss = state.cumulative_words_lost
        pre_edit_label = state.current_checkpoint_label
        combined_scopes = list(getattr(state, "pending_edit_scopes", []))
        for scope in edit_scopes:
            if scope not in combined_scopes:
                combined_scopes.append(scope)

        if not was_pending_cluster:
            state.active_cluster_start_code = pre_edit_code
            state.active_cluster_start_text_loss = pre_edit_text_loss
            state.active_cluster_start_label = pre_edit_label

        state.checkpoints.append(pre_edit_code)
        state.checkpoint_text_loss.append(pre_edit_text_loss)
        state.checkpoint_labels.append(pre_edit_label)
        state.current_code = new_code
        state.cumulative_words_lost = projected_text_loss
        self._invalidate_verify_after_code_change(state, f"{action_label} changed code")
        state.last_edit_scope = tuple(combined_scopes)
        state.current_checkpoint_label = (
            f"{action_label} revision {state.layout_revision}; scopes: "
            + (", ".join(combined_scopes) if combined_scopes else "unclassified")
        )

        if cluster_complete:
            state.last_cluster_start_code = (
                state.active_cluster_start_code
                if state.active_cluster_start_code is not None
                else pre_edit_code
            )
            state.last_cluster_start_text_loss = state.active_cluster_start_text_loss
            state.last_cluster_start_label = (
                state.active_cluster_start_label or pre_edit_label
            )
            state.pending_edit_cluster = False
            state.pending_edit_scopes = []
            state.active_cluster_start_code = None
            state.active_cluster_start_text_loss = 0
            state.active_cluster_start_label = ""
        else:
            state.pending_edit_cluster = True
            state.pending_edit_scopes = combined_scopes

    @staticmethod
    def _introduced_unstyled_layout_classes(
        before_code: str,
        after_code: str,
    ) -> list[str]:
        """Find new layout-role classes whose closure CSS is not present yet.

        Structural HTML reflows often need one edit for wrappers and another for
        their CSS. Treating the wrapper-only state as a completed rollback unit
        can strand half of that reflow when the closure CSS is later abandoned.
        """
        if not (
            AgentRepair._is_html_code(before_code)
            and AgentRepair._is_html_code(after_code)
        ):
            return []

        try:
            from bs4 import BeautifulSoup

            def classes_in(code: str) -> set[str]:
                soup = BeautifulSoup(code, "html.parser")
                return {
                    str(class_name)
                    for tag in soup.find_all(class_=True)
                    for class_name in (tag.get("class") or [])
                    if class_name
                }

            introduced = classes_in(after_code) - classes_in(before_code)
            soup = BeautifulSoup(after_code, "html.parser")
            css = "\n".join(
                style.get_text("\n") for style in soup.find_all("style")
            )
        except Exception:
            return []

        layout_terms = (
            "stack", "grid", "wrap", "wrapper", "container", "row", "column",
            "col", "panel", "rail", "card", "zone", "region", "layout", "body",
            "main", "aside", "cluster", "group",
        )
        return [
            class_name
            for class_name in sorted(introduced)
            if any(term in class_name.lower() for term in layout_terms)
            and not re.search(rf"\.{re.escape(class_name)}(?![\w-])", css)
        ]

    def _tool_apply_edits(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Apply search/replace edits."""
        edits = action.get("edits", [])
        if not edits:
            return "No edits provided.", False

        max_edits_per_call = getattr(self, "_max_edits_per_call", 24)
        if len(edits) > max_edits_per_call:
            return (
                "EDIT BATCH NOT APPLIED: received "
                f"{len(edits)} edits, but the per-call limit is "
                f"{max_edits_per_call}. Consolidate selectors or split the "
                "work at a coherent cluster boundary. No partial prefix was "
                "applied, so the current code is unchanged."
            ), False

        cluster_complete = action.get("cluster_complete", True) is not False
        edit_scopes = self._edit_scope_labels(edits)

        # Reject ambiguous edits before applying any of the batch. Silent global
        # replacement is too risky in generated HTML/CSS where short fragments
        # commonly repeat across unrelated elements.
        edit_counts = []
        ambiguous = []
        exact_no_op_edits: list[int] = []
        for index, edit in enumerate(edits, 1):
            search = edit.get("search", "")
            insert_after = edit.get("insert_after", "")
            if search and search == edit.get("replace", ""):
                exact_no_op_edits.append(index)
            if search:
                count = state.current_code.count(search)
                occurrence = edit.get("occurrence")
                expected = edit.get("expected_matches")
                edit_counts.append((search, count, occurrence, expected, False))
                if occurrence is None and expected is None and count != 1:
                    ambiguous.append(
                        f"Edit {index}: search matches {count} times; make the search "
                        "unique, set occurrence (1-based), or set expected_matches."
                    )
                elif expected is not None:
                    try:
                        expected_int = int(expected)
                    except (TypeError, ValueError):
                        ambiguous.append(f"Edit {index}: expected_matches must be an integer.")
                    else:
                        if count != expected_int:
                            ambiguous.append(
                                f"Edit {index}: expected {expected_int} matches but found {count}."
                            )
                elif occurrence is not None:
                    try:
                        occurrence_int = int(occurrence)
                    except (TypeError, ValueError):
                        ambiguous.append(f"Edit {index}: occurrence must be an integer.")
                    else:
                        if occurrence_int < 1 or occurrence_int > count:
                            ambiguous.append(
                                f"Edit {index}: occurrence {occurrence_int} is outside 1..{count}."
                            )
            elif insert_after:
                count = state.current_code.count(insert_after)
                occurrence = edit.get("occurrence")
                edit_counts.append((f"[insert after] {insert_after[:40]}", count, occurrence, None, True))
                if occurrence is None and count != 1:
                    ambiguous.append(
                        f"Edit {index}: insert_after matches {count} times; make it unique "
                        "or set occurrence (1-based)."
                    )

        if ambiguous:
            refresh = (
                "Call get_current_code before retrying this structural batch and "
                "construct every exact search from the current revision"
            )
            if state.last_code_read_revision == state.layout_revision:
                refresh = (
                    "Re-read the current-code result and rebuild the failed exact "
                    "search around the current DOM structure before retrying"
                )
            return (
                "AMBIGUOUS EDITS - no changes applied:\n"
                + "\n".join(ambiguous)
                + f"\n{refresh}. Do not reuse a pre-edit wrapper/closing-tag pattern."
            ), False

        dashboard_strategy_warning = self._dashboard_local_first_html_edit_message(
            edits, state,
        )
        broad_structural_warning = self._broad_structural_html_edit_message(
            edits, state,
        )

        new_code = _apply_edits(state.current_code, edits)
        if new_code == state.current_code:
            # Find which edits failed to match
            failed = []
            for i, edit in enumerate(edits):
                search = edit.get("search", "")
                if search and search not in state.current_code:
                    failed.append(f"Edit {i+1}: search string not found in code: \"{search[:80]}...\"")
            if failed:
                return "No changes applied. Failed matches:\n" + "\n".join(failed), False
            return "No changes applied — search strings may not match current code.", False

        edit_diff = ""
        for e in action.get("edits", []):
            edit_diff += e.get("search", "") + e.get("replace", "")
        raw_slot_warning = self._raw_figure_dense_slot_edit_warning(state, edit_diff)
        if raw_slot_warning:
            return raw_slot_warning, False

        # Compile test
        ok = self._test_compile(
            new_code, state.codegen_compiler, state.case_dir, state.slide_id,
        )
        if not ok:
            return (
                "Compile error — edits reverted. The replacement code has a "
                "Python syntax or runtime error. Check indentation and variable names."
            ), False

        text_loss_warning = ""
        projected_text_loss = state.cumulative_words_lost
        allows_content_change = (
            bool(state.issue_types & CONTENT_ACCURACY_ISSUE_TYPES)
            or getattr(state, "allow_visible_text_change", False)
        )
        if (
            self._is_html_code(new_code)
            and getattr(state, "allow_support_copy_compression", False)
        ):
            semantic_role_regressions = self._support_compression_role_regressions(
                state.current_code,
                new_code,
            )
            if semantic_role_regressions:
                return (
                    "AUTO-ROLLBACK: Support-copy compression removed an "
                    "information-bearing structural role: "
                    + ", ".join(semantic_role_regressions)
                    + ". Compression may shorten the authorized explanatory "
                    "support copy, but it may not delete a repeated metric/KPI "
                    "item, dissolve a findings/takeaway branch into metric notes, "
                    "or merge distinct roles. Preserve each recognizable semantic "
                    "unit and solve its geometry separately. Shorten prose within "
                    "each role, then use the measured card extents to reserve a real "
                    "terminal region or choose another role-preserving reflow."
                ), False
            dropped_high_value = _dropped_high_value_tokens(
                state.current_code, new_code, limit=10,
            )
            if dropped_high_value:
                shown = ", ".join(f'"{token}"' for token in dropped_high_value)
                return (
                    "AUTO-ROLLBACK: Support-copy compression removed protected "
                    f"values or named terms: {shown}. Shorten only explanatory "
                    "wording while retaining every number, metric, entity, label, "
                    "finding, and source attribution."
                ), False
        elif self._is_html_code(new_code) and not allows_content_change:
            dropped_high_value = _dropped_high_value_tokens(state.current_code, new_code, limit=10)
            if dropped_high_value:
                shown = ", ".join(f'"{token}"' for token in dropped_high_value)
                return (
                    "AUTO-ROLLBACK: This layout edit removed value-bearing content: "
                    f"{shown}. Preserve numbers, metrics, model/method names, and dataset names; "
                    "resolve the issue with CSS or relocate the existing content. "
                    "For table/card pressure, do not keep only the currently visible "
                    "subset; instead give the content a real region, recalibrate the "
                    "existing grid/table/card tracks, use stable-DOM grid areas, or "
                    "move supporting cards into another explicit track while preserving "
                    "all strings."
                ), False

            ordinary_lost = _meaningful_visible_words_lost(state.current_code, new_code)
            projected_text_loss += len(ordinary_lost)
            if ordinary_lost and projected_text_loss > state.text_loss_budget:
                shown = ", ".join(ordinary_lost[:6])
                return (
                    "AUTO-ROLLBACK: Cumulative visible-text loss budget would be exceeded "
                    f"({projected_text_loss} > {state.text_loss_budget}; removed: {shown}). "
                    "This edit was rejected, but DOM reflow remains available. Retry with "
                    "the complete information-bearing elements preserved or relocated "
                    "atomically."
                ), False
            if ordinary_lost:
                text_loss_warning = (
                    f"\nText-loss budget: {projected_text_loss}/{state.text_loss_budget} "
                    f"after removing {len(ordinary_lost)} ordinary word(s)."
                )

        # A wrapper-only DOM reflow is not a coherent rollback boundary yet.
        # Keep the transaction open until its dependent CSS/DOM closure arrives.
        implicit_cluster_classes = self._introduced_unstyled_layout_classes(
            state.current_code,
            new_code,
        )
        implicit_cluster_note = ""
        if cluster_complete and implicit_cluster_classes:
            cluster_complete = False
            shown = ", ".join(f".{name}" for name in implicit_cluster_classes)
            implicit_cluster_note = (
                "\n\nSTRUCTURAL EDIT CLUSTER KEPT OPEN: this DOM batch introduced "
                f"layout wrapper class(es) without closure CSS: {shown}. Apply the "
                "dependent CSS/DOM closure next with cluster_complete=true. If the "
                "reflow is abandoned, rollback(scope=\"cluster\") will restore the "
                "state before these wrappers were introduced."
            )

        support_copy_checkpoint_note = self._dashboard_support_copy_checkpoint_note(
            state,
            state.current_code,
            new_code,
            cluster_complete=cluster_complete,
        )
        no_op_edit_note = ""
        if exact_no_op_edits:
            shown = ", ".join(str(index) for index in exact_no_op_edits)
            no_op_edit_note = (
                "\n\nNO-OP EDIT NOTE: edit entries "
                f"{shown} replace text with the identical text and therefore make "
                "no code or rendered-demand change. Do not count them as completed "
                "compression or repair work; omit them from the next batch and judge "
                "copy calibration by changed wrapping and containment."
            )

        active_cluster_start = getattr(state, "active_cluster_start_code", None)
        cluster_start_code = (
            active_cluster_start
            if getattr(state, "pending_edit_cluster", False)
            and active_cluster_start is not None
            else state.current_code
        )

        # Success — save checkpoint and update code.
        self._commit_html_code_change(
            state,
            new_code,
            projected_text_loss=projected_text_loss,
            edit_scopes=edit_scopes,
            cluster_complete=cluster_complete,
            action_label="apply_edits",
        )
        edit_cluster_coverage_note = self._edit_cluster_execution_coverage_note(
            state,
            before_code=cluster_start_code,
            after_code=new_code,
            action_text=str(action.get("reasoning", "") or ""),
            cluster_complete=cluster_complete,
        )

        # NOTE: No auto-rollback here. All spatial checking is deferred
        # to verify_layout, where the model can decide whether to rollback
        # or continue with follow-up edits to fix intermediate regressions.

        # Build detailed feedback with occurrence counts
        details = []
        for i, (search_str, count, occurrence, expected, is_insert) in enumerate(edit_counts):
            verb = "inserted after" if is_insert else "replaced"
            if occurrence is not None:
                details.append(
                    f"  Edit {i+1}: {verb} occurrence {occurrence} of {count} "
                    f"matches for \"{search_str[:60]}\""
                )
            elif expected is not None and count > 1:
                details.append(
                    f"  Edit {i+1}: {verb} {count} verified occurrences of "
                    f"\"{search_str[:60]}\""
                )
            elif count > 1:
                details.append(f"  Edit {i+1}: {verb} {count} occurrences of \"{search_str[:60]}\"")
            elif count == 1:
                details.append(f"  Edit {i+1}: {verb} 1 occurrence of \"{search_str[:60]}\"")
            else:
                details.append(f"  Edit {i+1}: search string not found (insert or no-op)")

        # Check for text length changes that might cause overflow
        overflow_warnings = []
        import re as _re
        for edit in edits:
            search = edit.get("search", "")
            replace = edit.get("replace", "")
            if not search or not replace:
                continue
            # Check for text content edits (various patterns)
            text_patterns = ['.text = "', ".text = '", '.text = f"',
                             'p.text = "', 'run.text = "']
            is_text_edit = any(
                pat in search or pat in replace
                for pat in text_patterns
            )
            if is_text_edit:
                # Estimate char count change from the text portion
                old_texts = _re.findall(
                    r'\.text\s*=\s*["\'](.+?)["\']', search,
                )
                new_texts = _re.findall(
                    r'\.text\s*=\s*["\'](.+?)["\']', replace,
                )
                old_len = sum(len(t) for t in old_texts)
                new_len = sum(len(t) for t in new_texts)
                char_delta = new_len - old_len
                if char_delta > 20:
                    overflow_warnings.append(
                        f"⚠ Text grew by {char_delta} chars "
                        f"({old_len}→{new_len}). "
                        f"VERIFY the container box is large enough. "
                        f"If needed, increase the box height in the "
                        f"SAME edit batch to prevent overflow."
                    )
            # Check for container shrinking without checking text
            if 'Inches(' in search and 'Inches(' in replace:
                # Look for height reductions
                old_inches = _re.findall(
                    r'Inches\(([\d.]+)\)', search,
                )
                new_inches = _re.findall(
                    r'Inches\(([\d.]+)\)', replace,
                )
                if (len(old_inches) >= 4 and len(new_inches) >= 4):
                    # Position args: x, y, w, h
                    old_h = float(old_inches[3])
                    new_h = float(new_inches[3])
                    if new_h < old_h * 0.7 and old_h - new_h > 0.3:
                        overflow_warnings.append(
                            f"⚠ Container height reduced from "
                            f"{old_h:.2f}\" to {new_h:.2f}\" "
                            f"({(1-new_h/old_h)*100:.0f}% smaller). "
                            f"If this box has text, it may overflow. "
                            f"Call verify_layout after."
                        )

        detail_str = "\n".join(details) if details else ""
        warn_str = "\n".join(overflow_warnings) if overflow_warnings else ""

        # IMMEDIATE high-value-deletion warning — fire IN THE EDIT RESULT, not
        # Warn when an edit deletes high-value content (benchmark numbers,
        # model names). Surfacing this on the edit lets the agent choose to
        # RELOCATE/SHRINK instead of delete before it commits.
        # Non-blocking (the deletion may be correct).
        # Compares this edit's before (checkpoints[-1]) vs after.
        hv_delete_warning = ""
        if state.checkpoints:
            just_dropped = _dropped_high_value_tokens(
                state.checkpoints[-1], new_code, limit=10)
            if just_dropped:
                shown = ", ".join(f'"{d}"' for d in just_dropped)
                hv_delete_warning = (
                    f"\n\n⚠ This edit removed value-bearing content from the slide: "
                    f"{shown}. These are numbers/metrics or model/method/dataset "
                    f"names (rule 8a: never delete these just to clear a layout "
                    f"issue). If you deleted a card/box only to resolve an overlap, "
                    f"prefer RELOCATING or SHRINKING it instead — move it to empty "
                    f"space, reduce its size/font, or fold the figures into a nearby "
                    f"box. If the content is genuinely redundant, fine; otherwise "
                    f"reinstate these specific values."
                )

        # Detect layout property changes and force verify reminder
        layout_verify_warning = ""
        is_html = self._is_html_code(state.current_code)
        if is_html:
            layout_keywords = ["left:", "top:", "width:", "height:", "position:", "margin", "padding", "flex", "grid"]
            if any(kw in edit_diff for kw in layout_keywords):
                if cluster_complete:
                    layout_verify_warning = (
                        "\n\nLayout checkpoint marked complete. Call verify_layout "
                        "before submitting or starting a different repair family."
                    )
                else:
                    layout_verify_warning = (
                        "\n\nThis batch is marked as part of an unfinished coupled "
                        "edit cluster. Apply the remaining dependent-region edits "
                        "before treating the layout as a checkpoint. If you verify "
                        "now for measurements, interpret same-cluster residuals as "
                        "an intermediate state rather than a finished strategy."
                    )
        else:
            if "Inches(" in edit_diff:
                if cluster_complete:
                    layout_verify_warning = (
                        "\n\nCoordinate checkpoint marked complete. Call "
                        "verify_layout before submitting or changing repair family."
                    )
                else:
                    layout_verify_warning = (
                        "\n\nThis coordinate batch is an unfinished coupled edit "
                        "cluster. Complete its dependent edits before final "
                        "verification."
                    )
        dashboard_table_warning = self._dashboard_table_outer_frame_warning_from_edits(
            edits, state,
        )
        dashboard_coupled_warning = self._dashboard_coupled_cluster_warning_from_edits(
            edits, state,
        )
        dashboard_descendant_warning = self._dashboard_parent_descendant_patch_warning(
            state,
            self._html_edit_blob(edits),
        )
        dashboard_plan_implementation_warning = (
            self._dashboard_plan_implementation_warning(
                state,
                action_text=str(action.get("reasoning", "") or ""),
                cluster_complete=cluster_complete,
            )
        )

        result_msg = (
            f"Applied {len(edits)} edit(s) successfully. Code compiles OK.\n"
            f"{detail_str}\n"
            f"{warn_str}"
            f"{hv_delete_warning}\n"
            f"{text_loss_warning}\n"
            f"Verify after the current coherent edit cluster is complete."
            f"{layout_verify_warning}"
            f"{dashboard_strategy_warning or ''}"
            f"{broad_structural_warning or ''}"
            f"{dashboard_table_warning}"
            f"{dashboard_coupled_warning}"
            f"{dashboard_descendant_warning}"
            f"{dashboard_plan_implementation_warning}"
            f"{support_copy_checkpoint_note}"
            f"{no_op_edit_note}"
            f"{implicit_cluster_note}"
            f"{edit_cluster_coverage_note}"
        )
        return result_msg, True

    def _tool_apply_css_patch(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Append or replace one cascade-late CSS repair block in HTML slides."""
        if not self._is_html_code(state.current_code):
            return "apply_css_patch is available only for HTML slides.", False

        css = str(action.get("css", "") or "").strip()
        if not css:
            return "No CSS provided.", False
        if "<style" in css.lower() or "</style" in css.lower():
            return (
                "Provide CSS declarations only; do not include <style> tags. "
                "The tool owns the repair patch wrapper."
            ), False

        mode = str(action.get("mode", "append") or "append").lower()
        if mode not in {"append", "replace"}:
            return "mode must be either 'append' or 'replace'.", False

        start_marker = "/* REDECK_REPAIR_PATCH_START */"
        end_marker = "/* REDECK_REPAIR_PATCH_END */"
        patch_re = re.compile(
            re.escape(start_marker) + r".*?" + re.escape(end_marker),
            re.DOTALL,
        )
        existing = patch_re.search(state.current_code)
        if existing and mode == "replace":
            patch = f"{start_marker}\n{css}\n{end_marker}"
            new_code = (
                state.current_code[:existing.start()]
                + patch
                + state.current_code[existing.end():]
            )
        elif existing:
            insertion = f"\n{css}\n"
            marker_start = existing.end() - len(end_marker)
            new_code = (
                state.current_code[:marker_start]
                + insertion
                + state.current_code[marker_start:]
            )
        else:
            style_end = state.current_code.lower().rfind("</style>")
            if style_end < 0:
                return "Could not find a closing </style> tag for the CSS patch.", False
            patch = f"\n{start_marker}\n{css}\n{end_marker}\n"
            new_code = state.current_code[:style_end] + patch + state.current_code[style_end:]

        raw_slot_warning = self._raw_figure_dense_slot_edit_warning(state, css)
        if raw_slot_warning:
            return raw_slot_warning, False
        if not self._test_compile(
            new_code, state.codegen_compiler, state.case_dir, state.slide_id,
        ):
            return "Compile error — CSS patch reverted. Check the CSS syntax.", False

        cluster_complete = action.get("cluster_complete", True) is not False
        edit_scopes = self._edit_scope_labels([{"search": "", "replace": css}])
        active_cluster_start = getattr(state, "active_cluster_start_code", None)
        cluster_start_code = (
            active_cluster_start
            if getattr(state, "pending_edit_cluster", False)
            and active_cluster_start is not None
            else state.current_code
        )
        self._commit_html_code_change(
            state,
            new_code,
            projected_text_loss=state.cumulative_words_lost,
            edit_scopes=edit_scopes,
            cluster_complete=cluster_complete,
            action_label="apply_css_patch",
        )
        edit_cluster_coverage_note = self._edit_cluster_execution_coverage_note(
            state,
            before_code=cluster_start_code,
            after_code=new_code,
            action_text=str(action.get("reasoning", "") or ""),
            cluster_complete=cluster_complete,
        )

        boundary = (
            "checkpoint marked complete; call verify_layout"
            if cluster_complete
            else "unfinished coupled edit cluster; apply the dependent patch/DOM edits next"
        )
        scopes = ", ".join(edit_scopes) if edit_scopes else "unclassified selectors"
        plan_implementation_warning = self._dashboard_plan_implementation_warning(
            state,
            action_text=str(action.get("reasoning", "") or ""),
            cluster_complete=cluster_complete,
        )
        allocation_warning = (
            ""
            if plan_implementation_warning
            else self._dashboard_terminal_support_patch_warning(state, css)
        )
        variable_track_warning = self._dashboard_variable_track_patch_warning(
            state,
            css,
        )
        descendant_warning = self._dashboard_parent_descendant_patch_warning(
            state,
            css,
        )
        owner_budget_warning = self._dashboard_repeated_owner_budget_warning(
            state,
            css,
        )
        return (
            f"CSS patch {mode} applied successfully for {scopes}; {boundary}. "
            "The patch cascades after the original stylesheet, so revise it with "
            "mode='replace' if the combined hypothesis is wrong instead of layering "
            "contradictory overrides."
            f"{allocation_warning}"
            f"{plan_implementation_warning}"
            f"{variable_track_warning}"
            f"{descendant_warning}"
            f"{owner_budget_warning}"
            f"{edit_cluster_coverage_note}"
        ), True

    def _tool_verify_layout(self, state: AgentState) -> tuple[str | list, bool]:
        """Render current HTML and check spatial layout via Playwright DOM.

        Returns compact spatial state with violation counts vs baseline.
        Single source of truth — no EMU/GeomChecks indirection.
        """
        # Ablation: disable step-level render feedback
        if self._disable_step_render:
            return ("verify_layout is disabled in this configuration. "
                    "Spatial feedback is only available after submission. "
                    "Proceed with your planned edits based on the issue description."), False

        from .html_spatial_state import extract_html_slide_state, format_html_compact_state

        asset_base_dirs = self._html_asset_base_dirs(
            state.case_dir,
            getattr(state, "_run_dir", None),
            getattr(state, "_turn_index", None),
        )

        def _extract_state(code: str, slide_id: int):
            """Render HTML and extract spatial state."""
            try:
                state_obj = extract_html_slide_state(
                    slide_id,
                    code,
                    html_base_dir=Path(state.case_dir),
                    asset_base_dirs=asset_base_dirs,
                )
                return state_obj, None
            except Exception as e:
                return None, str(e)[:200]

        # Check current code
        t1_state, err = _extract_state(state.current_code, state.slide_id)
        if err:
            return f"Verification error: {err}", False
        state._last_html_state = t1_state

        # Check baseline (cached after first call)
        if not hasattr(state, '_t0_html_state') or state._t0_html_state is None:
            t0_state, t0_err = _extract_state(state.original_code, state.slide_id)
            if t0_state:
                state._t0_html_state = t0_state
                # Cache T0 overflow for rollback protection
                state._t0_overflow_px = compute_overflow_px(t0_state.blocks)
                t0_compact = self._scope_spatial_context(
                    format_html_compact_state(t0_state), self._current_issues,
                )
                state._t0_compact_issues = t0_compact.count("❌ ")
                # Cache T0 space-map coverage for delta tracking
                import re as _re
                _cov_m = _re.search(r"Coverage: (\d+)%", t0_compact)
                state._t0_space_coverage = int(_cov_m.group(1)) if _cov_m else None

        # Format compact state (single, consistent output)
        compact = self._scope_spatial_context(
            format_html_compact_state(t1_state), self._current_issues,
        )
        if self._looks_like_table_dashboard_pressure(state):
            compact = self._dashboard_measurement_context(compact)
            allocation_map = self._dashboard_allocation_map(t1_state)
            if allocation_map:
                compact += "\n\n" + allocation_map
        compact_issues = compact.count("❌ ")

        # Authoritative residual count via the single source of truth — this is
        # what the external scorer AND the submit gate use, so the agent should
        # be told about the SAME issues (the compact "❌" tally diverges slightly
        # on thresholds and historically under-counted clipped text).
        from .html_spatial_state import (
            significant_issue_regressions,
        )
        targeted_categories = self._targeted_residual_categories()
        t0_state = getattr(state, '_t0_html_state', None)
        targeted_sig = self._targeted_significant_issues(t0_state, t1_state)
        sig_total = sum(len(value) for value in targeted_sig.values())
        state.last_verify_targeted_residual_counts = {
            category: len(items)
            for category, items in targeted_sig.items()
            if items
        }
        targeted_measurement_seen = False
        if targeted_categories:
            from .html_spatial_state import count_significant_issues
            current_all_sig = count_significant_issues(t1_state)
            baseline_all_sig = count_significant_issues(t0_state) if t0_state else {}
            targeted_measurement_seen = any(
                current_all_sig.get(category) or baseline_all_sig.get(category)
                for category in targeted_categories
            )

        lines = []
        if self._looks_like_table_dashboard_pressure(state):
            decision_summary = self._dashboard_decision_summary(state, t1_state)
            if decision_summary:
                lines.append(decision_summary)
        coupled_compression_note = self._dashboard_pending_coupled_compression_note(
            state,
            t1_state,
        )
        if coupled_compression_note:
            lines.append(coupled_compression_note)
        recent_scope = ", ".join(state.last_edit_scope) or "the edited region"
        if state.pending_edit_cluster:
            lines.append(
                "COUPLED EDIT CLUSTER IN PROGRESS: the latest revision was "
                f"explicitly marked unfinished. Recent edit scope: {recent_scope}. "
                "Use this verification as intermediate spatial evidence. Before "
                "rolling back, distinguish residuals owned by dependent regions "
                "that have not yet received their closure edit from damage to "
                "content, media, hierarchy, or unrelated regions. Continue the "
                "cluster when its remaining dependent edits have a credible direct "
                "closure path; abandon it when that path is no longer credible."
            )
        elif state.last_edit_scope:
            lines.append(
                "LATEST COHERENT EDIT SCOPE: "
                f"{recent_scope}. Judge new residuals by whether this checkpoint "
                "actually changed their owning region; persistence elsewhere does "
                "not by itself invalidate the checkpoint."
            )
        state._last_verify_text_regression = False
        state._last_verify_text_signal = False
        state._last_verify_text_signal_reason = ""
        state._last_verify_visual_compression_failed = False
        state._last_verify_scope_failed = False
        text_ok = True
        text_reason = ""

        compression_ok = True
        compression_reason = ""
        recoverable_dashboard_compression = False
        if state.current_code != state.original_code:
            # Skip visual compression check when large overflow exists —
            # compression is NECESSARY to fit content within 720px canvas.
            _current_overflow = compute_overflow_px(t1_state.blocks) if t1_state else 0
            # Also check stored overflow from last verify as fallback
            _stored_overflow = getattr(state, '_last_verified_overflow_px', 0) or 0
            _effective_overflow = max(_current_overflow, _stored_overflow)
            _large_overflow_bypass = _effective_overflow > 30

            compression_ok, compression_reason = (
                validate_repair_not_visual_compression(
                    state.original_code,
                    state.current_code,
                    allow_dominant_element_removal=(
                        issues_allow_dominant_element_removal(
                            self._current_issues
                        )
                    ),
                )
            )
            if _large_overflow_bypass:
                compression_ok = True  # allow compression when overflow demands it
                state._last_verify_visual_compression_failed = False
            if not compression_ok:
                recoverable_dashboard_compression = (
                    self._is_recoverable_dashboard_dominant_compression(
                        state, compression_reason,
                    )
                )
                if recoverable_dashboard_compression:
                    lines.append(
                        "VISUAL HIERARCHY CORRECTION REQUIRED: "
                        f"{compression_reason}. The coupled dashboard checkpoint "
                        "may still contain useful table, ranking, and summary "
                        "calibration, so do not roll back the whole batch solely "
                        "for this focal-scale error. Make the next edit a local "
                        "correction that restores only the primary hero/KPI to the "
                        "nearest clearly dominant fitting scale with a natural "
                        "line-height. Do not compensate by shrinking its description, "
                        "notes, or summary copy further. Choose the corrected scale "
                        "from the current full render rather than a memorized size. "
                        "Preserve the support-role calibration, then verify again. This checkpoint "
                        "cannot be submitted until the final compression gate passes."
                    )
                else:
                    state._last_verify_visual_compression_failed = True
                    lines.append(
                        "SHIPMENT GATE FAILED — VISUAL COMPRESSION: "
                        f"{compression_reason}. Evaluate whether this compression "
                        "is NECESSARY to resolve the overflow (if content still "
                        "extends past 720px, compression is justified — continue "
                        "editing to finish the fit). Roll back ONLY if the "
                        "compression damaged unrelated regions or compressed "
                        "elements that were already fitting fine. If overflow "
                        "remains, this checkpoint is a valid intermediate state."
                    )

        # Use the same stable physical-defect comparison as the submit gate.
        # Scoped compact output intentionally hides unrelated baseline noise,
        # but it must never hide a regression introduced by the current edit.
        visual_only_repair = bool(self._current_issues) and all(
            (issue.rubric_id or "").startswith("B")
            for issue in self._current_issues
        )
        if t0_state is not None and visual_only_repair:
            from .repair_utils import validate_rendered_text_preservation
            text_ok, text_reason = validate_rendered_text_preservation(
                t0_state, t1_state,
                allow_revealed_text=issues_allow_rendered_text_reveal(
                    self._current_issues,
                ),
                allow_text_formatting_change=any(
                    issue.issue_type == "formatting_error"
                    for issue in self._current_issues
                ),
                allow_text_content_change=issues_allow_visible_text_change(
                    self._current_issues
                ),
            )
            if not text_ok:
                state._last_verify_text_signal = True
                state._last_verify_text_signal_reason = text_reason
                lines.append(
                    "VISIBLE TEXT CHANGE SIGNAL (advisory): "
                    f"{text_reason}. This rendered-token comparison is a "
                    "risk signal, not a checkpoint or shipment gate by itself. "
                    "Use it to inspect whether source text was truly deleted, "
                    "hidden, or semantically reordered. If the source DOM, media, "
                    "roles, and deterministic spatial issues are acceptable, "
                    "continue the current closure/calibration chain instead of "
                    "rolling back solely because of this signal."
                )
        regressions = (
            significant_issue_regressions(t0_state, t1_state)
            if t0_state is not None else {}
        )
        target_aligned_crop = []
        raw_css_crop_hints = (
            html_image_css_crop_hints(state.current_code)
            if any(
                issue.issue_type in {"raw_figure", "raw_table"}
                for issue in self._current_issues
            ) and self._is_html_code(state.current_code)
            else []
        )
        if can_exempt_raw_figure_image_crop(self._current_issues, state.current_code):
            target_aligned_crop = regressions.pop("image_crop", [])
        regression_total = sum(len(items) for items in regressions.values())
        state._last_verify_spatial_regression_total = regression_total
        state._last_verify_targeted_residual_total = sig_total
        state._last_verify_compact_issues = compact_issues

        def _describe_block(block_id: str) -> str:
            block = next(
                (item for item in t1_state.blocks if item.block_id == block_id),
                None,
            )
            if block is None:
                return block_id
            text = " ".join(block.text_lines).strip()
            label = text[:56] if text else (block.css_selector or block.var_name)
            x, y, width, height = block.bbox_px
            return (
                f'{block_id} "{label}" '
                f'at ({x:.0f},{y:.0f},{width:.0f}x{height:.0f})'
            )

        if regression_total:
            lines.append(self._spatial_regression_policy_message(
                regression_total,
                hard_quality_failure=(
                    f"visual compression: {compression_reason}"
                    if not compression_ok and not recoverable_dashboard_compression
                    else ""
                ),
            ))
            regression_signature = tuple(sorted(
                f"{group}:{kind}"
                for group, items in regressions.items()
                for kind, _payload in items
            ))
            if regression_signature == state.last_spatial_regression_signature:
                state.spatial_regression_streak += 1
            else:
                state.last_spatial_regression_signature = regression_signature
                state.spatial_regression_streak = 1
            shown_regressions, omitted_regressions = (
                self._representative_spatial_regressions(regressions)
            )
            for group, kind, payload in shown_regressions:
                if group == "interaction":
                    first_id, second_id = payload
                    lines.append(
                        f"  NEW {kind.upper()}: "
                        f"{_describe_block(first_id)} <-> "
                        f"{_describe_block(second_id)}"
                    )
                else:
                    lines.append(
                        f"  NEW {kind.upper()}: "
                        f"{_describe_block(payload)}"
                    )
            if omitted_regressions:
                omitted_summary = ", ".join(
                    f"{group} +{count}"
                    for group, count in omitted_regressions.items()
                )
                lines.append(
                    "  REGRESSION DETAIL SUMMARY: representative root examples "
                    f"shown above; omitted repeated detector findings: {omitted_summary}. "
                    "Use the descendant, relation, and repeated-allocation maps to "
                    "reason about their shared spatial cause rather than treating "
                    "each nested pair as an independent failure."
                )
            if state.spatial_regression_streak >= 2:
                lines.append(
                    "STRATEGY ESCALATION: the same class of hard spatial "
                    "regression has appeared on consecutive verify_layout calls. "
                    "Use the render to decide whether this is an unfinished state "
                    "inside the same structural repair or evidence that the current "
                    "approach is not converging. Continue the same attempt when the "
                    "content and roles remain intact and the remaining defect has a "
                    "clear structural cause. Roll back or change direction when "
                    "information is damaged, unrelated regions regress, or the same "
                    "pressure is only being moved around."
                )
        else:
            state.spatial_regression_streak = 0
            state.last_spatial_regression_signature = ()
        if target_aligned_crop:
            lines.append(
                "TARGET-ALIGNED MEASUREMENT: the requested raw-figure crop "
                "changed visible image framing. This is not a regression by "
                "itself; use render_preview to confirm the intended subject "
                "remains complete and the unwanted source material is gone."
            )
        if raw_css_crop_hints:
            lines.append(
                "RAW-FIGURE CSS CROP WARNING: current code uses CSS image "
                f"windowing ({', '.join(raw_css_crop_hints)}). For B17, this "
                "is not a sufficient final repair because it can make DOM "
                "geometry pass while cutting away figure content. Replace the "
                "image with a real cropped/recomposed asset, exact-data generated chart, "
                "or fidelity-preserving SVG summary asset, "
                "display it intact with object-fit: contain, then call "
                "render_preview."
            )

        # Delta vs baseline — with specific fixed/new issue details
        t0_count = getattr(state, '_t0_compact_issues', 0)
        compact_delta = compact_issues - t0_count
        if regression_total and compact_delta < 0:
            reduction_pct = int(100 * (-compact_delta) / max(t0_count, 1))
            if reduction_pct >= 75:
                lines.append(
                    f"NET PROGRESS: strongly improved by {-compact_delta} hard "
                    f"defect(s) vs baseline ({t0_count}→{compact_issues}, "
                    f"{reduction_pct}% reduction), despite {regression_total} "
                    f"new finding(s)."
                )
            else:
                lines.append(
                    f"NET PROGRESS: improved by {-compact_delta} hard defect(s) "
                    f"vs baseline ({t0_count}→{compact_issues}), despite "
                    f"{regression_total} new finding(s)."
                )
            # If new findings include overlaps, call them out explicitly
            if regression_total > 0:
                lines.append(
                    f"The {regression_total} new finding(s) listed above "
                    f"(NEW OVERLAP / NEW CLIPPED) are real spatial defects "
                    f"that need fixing — they are not detector artifacts. "
                    f"Do not mark the plan step as done while these remain."
                )
        elif regression_total and compact_delta == 0:
            lines.append(
                f"NET NEUTRAL: total hard defect count unchanged vs baseline "
                f"({t0_count}→{compact_issues}). {regression_total} new "
                f"finding(s) appeared but an equal number of baseline findings "
                f"resolved. The strategy is making progress if the remaining "
                f"defects are in the region being actively repaired."
            )
        elif regression_total and compact_delta > 0:
            lines.append(
                f"NET REGRESSION: {compact_delta:+d} net hard defect(s) vs "
                f"baseline ({t0_count}→{compact_issues}). {regression_total} "
                f"new finding(s) appeared. Inspect whether new defects are "
                f"in the actively edited region (expected intermediate state) "
                f"or in unrelated regions (real regression requiring rollback)."
            )
        elif compact_delta < 0:
            lines.append(
                f"REGRESSION CHECK: improved by {-compact_delta} hard defect(s) "
                f"vs baseline ({t0_count}→{compact_issues}). Remaining baseline "
                "findings are context, not automatically new tasks. Do not dismiss "
                "one as unrelated when it shows protected content participating in "
                "the same spatial-pressure chain as the original issue."
            )
        elif compact_delta > 0:
            lines.append(
                f"REGRESSION CHECK: {compact_delta:+d} net hard defect(s) vs "
                f"baseline ({t0_count}→{compact_issues}). Inspect the diff and "
                "fix only defects introduced by the current edits."
            )
        elif compact_issues == 0 and sig_total == 0:
            lines.append(
                "REGRESSION CHECK: no new hard regression vs baseline and no "
                "target-region deterministic visibility defect. Resolve the "
                "original visual issue and use the space map only as supporting evidence."
            )
        elif compact_issues == 0:
            lines.append(
                "REGRESSION CHECK: no new hard regression vs baseline, but "
                f"{sig_total} deterministic visibility measurement(s) remain in "
                "the issue's named region. They are listed below and must not be "
                "described as a clean or fully visible checkpoint."
            )
        else:
            lines.append(
                f"REGRESSION CHECK: no net new hard defects; {compact_issues} "
                "finding(s) are unchanged from baseline. They are not automatically "
                "additional repair tasks, but they still matter when they expose the "
                "same named layout conflict or keep protected content invisible."
            )

        alignment_delta = self._format_alignment_relation_delta(
            t0_state,
            t1_state,
            self._current_issues,
        )
        if alignment_delta:
            lines.append(alignment_delta)

        lines.append("\n" + compact)

        # Residual significant-issue feedback (non-blocking, every verify call).
        # On every verify where significant issues remain, name them and ask
        # the agent to fix or justify — informing its judgment, never blocking it.
        if sig_total > 0:
            res_lines = []
            for bid in (
                targeted_sig.get("clipped", [])
                + targeted_sig.get("canvas_truncation", [])
            ):
                blk = next((b for b in t1_state.blocks if b.block_id == bid), None)
                if blk is not None:
                    px = int(getattr(blk, 'clipped_bottom_px', 0) or 0)
                    prev = " ".join(blk.text_lines)[:36] if blk.text_lines else bid
                    res_lines.append(f"  • CLIPPED {bid}: \"{prev}\"" + (f" ({px}px hidden)" if px else ""))
            for bid in targeted_sig.get("text_overflow", []):
                res_lines.append(f"  • OVERFLOW {bid}")
            for issue_id in targeted_sig.get("svg_text_overflow", []):
                res_lines.append(
                    self._format_svg_text_overflow_residual(t1_state, issue_id)
                )
            for a, b in targeted_sig.get("overlap", []):
                res_lines.append(f"  • OVERLAP {a} ↔ {b}")
            for bid in targeted_sig.get("out_of_bounds", []):
                res_lines.append(f"  • OUT-OF-BOUNDS {bid}")
            for a, b in targeted_sig.get("occlusion", []):
                res_lines.append(f"  • OCCLUDED {b} behind {a}")
            if res_lines:
                if self._looks_like_table_dashboard_pressure(state):
                    residual_advice = (
                        "For this dense dashboard, decide from the current revision evidence whether "
                        "the residual is local or reflects shared pressure across regions. "
                        "A role-aware calibration may be enough when topology and hierarchy "
                        "remain sound; regional reflow may be better when pressure keeps "
                        "moving or normal flow cannot expose all semantic groups. Repeated "
                        "rows and cards multiply small rhythm costs, so compare their cumulative "
                        "height and wrapping with one-time title, header, and KPI roles before "
                        "shrinking the frame again. Choose the next direction from that diagnosis, "
                        "not from a fixed sequence or detector count."
                    )
                else:
                    residual_advice = (
                        "If several remaining measurements belong to the same body, lower, "
                        "table/card, or footer-adjacent region, diagnose that shared spatial "
                        "pressure before the next edit. A repeated local font/padding/height "
                        "nudge is usually the wrong next move unless the region topology is "
                        "already sound; prefer a real regional/body reflow that gives the "
                        "affected semantic units visible space while preserving text and roles."
                    )
                lines.append(
                    f"\n{sig_total} target-category deterministic measurement(s) "
                    f"remain:\n"
                    + "\n".join(res_lines)
                    + "\n  Use these measurements only to assess the original issue and "
                    "its owning spatial region. Preserve unrelated baseline content, "
                    "but do not label clipped protected content unrelated when it shares "
                    "the same layout pressure. "
                    + residual_advice
                )

        # ── Quadrant fill data (spatial, non-blocking) ──
        # Reuse the exact structured measurement rendered in SPACE MAP.
        _qf = {}
        if self._needs_spatial_distribution(self._current_issues):
            from .html_spatial_state import measure_space_occupancy
            _space_measurement = measure_space_occupancy(t1_state.blocks)
            _qf = _space_measurement["quadrant_fill"]
            if not _space_measurement["significant_block_count"]:
                _qf = {}

        if self._needs_spatial_distribution(self._current_issues) and _qf:
            _max_q = max(_qf.values())
            _min_q = min(_qf.values())
            _sparse = [k for k, v in _qf.items() if v < 40]
            _dense = [k for k, v in _qf.items() if v >= 55]

            if _max_q >= 55 and _min_q < 40 and _sparse and _dense:
                lines.append(
                    f"\nSPACE BALANCE SIGNAL: {', '.join(_sparse)} region(s) only "
                    f"{', '.join(f'{_qf[q]}%' for q in _sparse)} filled. This is "
                    f"supporting evidence, not a repair task unless the original "
                    f"VLM issue identifies density or whitespace imbalance."
                )

        # Detailed diff: which issues were fixed, which are new
        if t0_state and t1_state and targeted_categories:
            t0_ovlp = {(min(a,b),max(a,b)) for a,b,_ in t0_state.overlap_pairs}
            t1_ovlp = {(min(a,b),max(a,b)) for a,b,_ in t1_state.overlap_pairs}
            t0_oob = set(t0_state.oob_blocks)
            t1_oob = set(t1_state.oob_blocks)
            t0_clip = set(t0_state.clipped_blocks)
            t1_clip = set(t1_state.clipped_blocks)
            fixed_ovlp = t0_ovlp - t1_ovlp
            new_ovlp = t1_ovlp - t0_ovlp
            fixed_oob = t0_oob - t1_oob
            new_oob = t1_oob - t0_oob
            fixed_clip = t0_clip - t1_clip
            new_clip = t1_clip - t0_clip
            # Image crop diff
            t0_img_crop = {b.block_id for b in t0_state.blocks
                           if getattr(b, 'shape_type', '') == 'picture' and getattr(b, 'img_crop_pct', 0) > 0.25}
            t1_img_crop = {b.block_id for b in t1_state.blocks
                           if getattr(b, 'shape_type', '') == 'picture' and getattr(b, 'img_crop_pct', 0) > 0.25}
            fixed_img_crop = t0_img_crop - t1_img_crop
            new_img_crop = t1_img_crop - t0_img_crop
            diff_parts = []
            if fixed_ovlp: diff_parts.append(f"FIXED {len(fixed_ovlp)} overlaps")
            if fixed_oob: diff_parts.append(f"FIXED {len(fixed_oob)} OOB")
            if fixed_clip: diff_parts.append(f"FIXED {len(fixed_clip)} clipped")
            if fixed_img_crop: diff_parts.append(f"FIXED {len(fixed_img_crop)} image crops")
            if new_ovlp: diff_parts.append(f"NEW {len(new_ovlp)} overlaps: {list(new_ovlp)[:3]}")
            if new_oob: diff_parts.append(f"NEW {len(new_oob)} OOB: {list(new_oob)[:3]}")
            if new_clip: diff_parts.append(f"NEW {len(new_clip)} clipped: {list(new_clip)[:3]}")
            if new_img_crop: diff_parts.append(f"NEW {len(new_img_crop)} image crops: {list(new_img_crop)[:3]}")
            if diff_parts:
                lines.append("\nDiff vs baseline: " + " | ".join(diff_parts))

        # Coverage delta vs T0 is measurement only. Bounding-box occupancy does
        # not establish whether whitespace is intentional or whether the focal
        # hierarchy improved; the VLM issue brief and render preview own that
        # judgment.
        t0_cov = getattr(state, '_t0_space_coverage', None)
        if (
            self._needs_spatial_distribution(self._current_issues)
            and t0_cov is not None
        ):
            import re as _re
            _cov_m = _re.search(r"Coverage: (\d+)%", compact)
            if _cov_m:
                cur_cov = int(_cov_m.group(1))
                cov_delta = cur_cov - t0_cov
                if abs(cov_delta) >= 5:
                    lines.append(
                        f"\nCoverage measurement: {cov_delta:+d}pp vs original "
                        f"({t0_cov}% → {cur_cov}%). This is not a pass/fail target; "
                        f"judge the full render against the diagnosed compositional "
                        f"failure."
                    )
                elif cur_cov < 50:
                    lines.append(
                        f"\nBaseline-aware coverage: {cur_cov}% (original "
                        f"{t0_cov}%). Low absolute coverage is not a repair task "
                        f"unless the original VLM issue identifies density or "
                        f"whitespace imbalance."
                    )

        closure_reminder = self._build_composition_closure_verify_reminder(
            self._current_issues,
        )
        if closure_reminder:
            lines.append(closure_reminder)

        # Broken image guidance — prevent agent from wasting tool budget on path guessing
        if t1_state and t1_state.broken_images:
            lines.append("")
            lines.append(
                "⚠️ BROKEN IMAGE GUIDANCE: The broken image(s) above are caused by "
                "path resolution — the rendering pipeline resolves relative paths at "
                "compile time. Do NOT attempt to fix image src paths (changing the path "
                "will not help and wastes tool budget). Instead, keep the original src "
                "unchanged and focus on other issues."
            )

        # Meta-content scan: detect instruction/editorial text in HTML
        meta_warnings = self._scan_meta_content(state.current_code)
        if meta_warnings:
            lines.append("")
            lines.append("🚫 META-CONTENT DETECTED (must remove before submit):")
            for mw in meta_warnings:
                lines.append(f"  • {mw}")

        # ── Quality signals (font degradation + word retention) ──
        # These are NOT gates — they are real-time feedback so the agent
        # can see the cost of its strategy before committing further.
        if state.current_code != state.original_code:
            # Font degradation: compare per-element font sizes vs original
            t0_st = getattr(state, '_t0_html_state', None)
            if t0_st and t1_state:
                t0_fonts = sorted(
                    [b.font_size_px for b in t0_st.blocks
                     if b.font_size_px > 0 and b.text_chars > 10
                     and b.shape_type not in ('picture', 'chart')],
                    reverse=True,
                )
                t1_fonts = sorted(
                    [b.font_size_px for b in t1_state.blocks
                     if b.font_size_px > 0 and b.text_chars > 10
                     and b.shape_type not in ('picture', 'chart')],
                    reverse=True,
                )
                if t0_fonts and t1_fonts:
                    dashboard_role_calibration = self._looks_like_table_dashboard_pressure(state)
                    t0_median = t0_fonts[len(t0_fonts) // 2]
                    t1_median = t1_fonts[len(t1_fonts) // 2]
                    if t1_median < t0_median - 2:
                        if dashboard_role_calibration:
                            lines.append(
                                f"\nDASHBOARD ROLE-SCALE CHECK: median text scale "
                                f"changed {t0_median:.0f}px → {t1_median:.0f}px. "
                                "This can be appropriate when the reduction is "
                                "concentrated in repeated table rows, ranking rows, "
                                "and summary support copy. Confirm that title, KPI, "
                                "card identities, and interpretation notes retain a "
                                "clear hierarchy; do not enlarge dense support roles "
                                "solely to restore the median."
                            )
                        else:
                            lines.append(
                                f"\n⚠ FONT DEGRADATION: median body font "
                                f"{t0_median:.0f}px → {t1_median:.0f}px. "
                                f"Consider restructuring instead of shrinking fonts."
                            )
                    t0_below_14 = [f for f in t0_fonts if f < 14]
                    below_14 = [f for f in t1_fonts if f < 14]
                    explicit_typography_target = any(
                        issue.issue_type == "typography_error"
                        for issue in self._current_issues
                    )
                    newly_small = max(0, len(below_14) - len(t0_below_14))
                    if below_14 and dashboard_role_calibration and newly_small:
                        lines.append(
                            f"ℹ Dense dashboard support scale: {len(below_14)} text "
                            f"element(s) are below 14px ({newly_small} introduced). "
                            "Small repeated data/support roles are allowed when they "
                            "remain readable and the focal/card hierarchy is intact; "
                            "treat this as inspection context, not an instruction to "
                            "inflate table or rail rows."
                        )
                    elif below_14 and (newly_small or explicit_typography_target):
                        lines.append(
                            f"🚨 {len(below_14)} text element(s) below 14px "
                            f"({newly_small} introduced by this repair). Increase "
                            f"font or reflow the existing content; do not delete or "
                            f"condense text unless a content issue explicitly requires it."
                        )
                    elif below_14:
                        lines.append(
                            f"ℹ Baseline contains {len(below_14)} text element(s) "
                            "below 14px. This repair did not introduce them; treat this "
                            "as context, not an instruction to alter unrelated content."
                        )

            # Word retention — surfaced as a SEMANTIC prompt, not a verdict.
            # A large drop is not automatically wrong (a section may have
            # genuinely needed to go); a small drop can still be wrong (a
            # dropped figure legend, an orphaned header). So we show the
            # delta and ask the agent to JUDGE what was removed, rather than
            # asserting "you deleted too much."
            t0_wc = _count_html_words(state.original_code)
            t1_wc = _count_html_words(state.current_code)
            if t0_wc > 10:
                wc_pct = round(100 * t1_wc / t0_wc)
                if wc_pct < 85:
                    lines.append(
                        f"\nℹ Word count: {t0_wc} → {t1_wc} ({wc_pct}% retained). "
                        f"This is a signal, not a verdict — judge what was removed:\n"
                        f"  • If you dropped a whole low-value support channel "
                        f"(decorative bullets, a redundant callout, an accent "
                        f"ribbon) to make crowded content fit — that's correct, keep it.\n"
                        f"  • If you deleted a paragraph but left its heading, "
                        f"emptied a card/pill, removed a figure's legend/caption, "
                        f"or dropped numbers/method names — that's a regression. "
                        f"Restore that specific content (shrink/relayout to fit it) "
                        f"or delete the now-orphaned wrapper too. Don't leave it half-removed."
                    )
                # Name the SPECIFIC high-value tokens that disappeared, so the
                # agent can apply rule 8a to concrete items rather than a vague
                # delta. Fires even at ≥85% retention, because a single dropped
                # model name / metric is a regression regardless of total words.
                dropped = _dropped_high_value_tokens(
                    state.original_code, state.current_code)
                if dropped:
                    shown = ", ".join(f'"{d}"' for d in dropped)
                    lines.append(
                        f"\n⚠ High-value content no longer on the slide: {shown}. "
                        f"These are numbers/metrics or model/method/dataset names — "
                        f"rule 8a says never drop these to save space. If one went "
                        f"out with a genuinely redundant lane, fine; otherwise restore "
                        f"it (shrink or relayout other things to make room). Do NOT "
                        f"invent values — only restore what was there."
                    )

        # Text diff summary: show what text changed vs original
        if state.current_code != state.original_code:
            diff_summary = self._text_diff_summary(
                state.original_code, state.current_code, state=state,
            )
            if diff_summary:
                lines.append("")
                lines.append(diff_summary)
                # If significant text was added, prompt agent to verify
                added_count = diff_summary.count("Added words")
                if added_count and "Added words" in diff_summary:
                    import re
                    m = re.search(r'Added words \((\d+)\)', diff_summary)
                    n_added = int(m.group(1)) if m else 0
                    guidance = self._text_diff_source_guidance(state, n_added)
                    if guidance:
                        lines.append(guidance)

        # NOTE: Density-target and coverage-drop feedback removed.
        # The bbox-union "coverage %" metric does not correlate with
        # visual-density assessment (contiguous empty regions) and
        # causes premature stopping.
        # The agent should rely on the issue's planned_fix and
        # verify_layout's spatial checks (overlap/overflow/OOB)
        # instead of a misleading numeric target.

        # Step verification reminder: show current step's expected outcome
        # so the agent self-checks against the LAYOUT ANCHOR data.
        dashboard_next_note = str(
            getattr(state, "_dashboard_next_strategy_note", "") or ""
        )
        if dashboard_next_note:
            lines.append("\n" + dashboard_next_note)

        if state.plan_steps:
            current_step = next(
                (s for s in state.plan_steps if s.status == "in_progress"),
                None,
            )
            if current_step and current_step.expected_outcome:
                lines.append(
                    f"\n🎯 Current step expected: {current_step.expected_outcome}"
                )
                if current_step.verify_criterion:
                    lines.append(
                        f"   ✅ Verify criterion: {current_step.verify_criterion}"
                    )
                lines.append(
                    "   ↳ Review the LAYOUT ANCHOR data above to confirm "
                    "this outcome is met before proceeding."
                )

        # Store for external use
        state.last_verify_result = {
            "t0_total": t0_count,
            "t1_total": compact_issues,
            "delta_total": compact_delta,
        }
        state.last_verify_revision = state.layout_revision
        state.last_verify_stale_reason = ""

        safe_checkpoint_valid = (
            compression_ok
            and not raw_css_crop_hints
        )
        if visual_only_repair:
            scope_ok, scope_reason = validate_visual_repair_scope(
                state.original_code,
                state.current_code,
                allow_image_replacement=any(
                    issue.issue_type in {"raw_figure", "raw_table", "svg_visual_defect"}
                    for issue in self._current_issues
                ),
                allow_text_formatting_change=any(
                    issue.issue_type == "formatting_error"
                    for issue in self._current_issues
                ),
                allow_text_content_change=issues_allow_visible_text_change(
                    self._current_issues
                ),
            )
            safe_checkpoint_valid = safe_checkpoint_valid and scope_ok
            if not scope_ok:
                state._last_verify_scope_failed = True
                lines.append(
                    f"CHECKPOINT INVALID: visual-only repair scope changed: {scope_reason}."
                )

        # Preserve content/media/scope-safe intermediate checkpoints even while
        # the named target still has measurable residuals. Final summary/submit
        # gates remain responsible for blocking unresolved objective visibility
        # failures; checkpointing must not erase a coherent multi-edit chain.
        if safe_checkpoint_valid:
            state.latest_safe_verified_code = state.current_code
            state.latest_safe_verified_revision = state.layout_revision

        checkpoint_valid = (
            regression_total == 0
            and safe_checkpoint_valid
            and sig_total == 0
        )
        if compact_delta <= 0 and checkpoint_valid:
            state.last_verified_code = state.current_code

        # Track BEST verified state — save any checkpoint that improves on the
        # previous best, even if residual issues remain. A state with 3 issues
        # is strictly better than the original with 23 issues; discarding it
        # because it's "not perfect" means the agent's work is wasted.
        if safe_checkpoint_valid and targeted_categories and targeted_measurement_seen and (
            state.best_verified_issues is None
            or sig_total < state.best_verified_issues
        ):
            state.best_verified_issues = sig_total
            state.best_verified_code = state.current_code

        # ── Overflow budget summary (clear actionable signal) ──
        if t1_state is not None:
            _overflow_px = compute_overflow_px(t1_state.blocks)
            _max_bottom = max(
                (b.bbox_px[1] + b.bbox_px[3] for b in t1_state.blocks),
                default=720,
            )
            # Stall detection: if overflow hasn't decreased, agent needs stronger compression
            _prev_overflow = getattr(state, '_last_verified_overflow_px', 0)
            _overflow_stalled = (
                _overflow_px > 20
                and _prev_overflow > 20
                and _overflow_px >= _prev_overflow - 5
            )
            state._last_verified_overflow_px = _overflow_px

            # Compute style match for spatial info (used by multiple branches below)
            import re as _re
            _style_match = _re.search(r'<style>(.*?)</style>', state.current_code or '', _re.DOTALL)

            if _overflow_px > 0:
                # Find containers where true content exceeds declared height
                _container_hints = []
                for blk in t1_state.blocks:
                    true_sh = getattr(blk, 'true_scroll_h_px', 0) or 0
                    client_h = getattr(blk, 'client_h_px', 0) or 0
                    if true_sh > client_h + 10 and client_h > 50:
                        css = blk.css_selector or blk.var_name
                        _container_hints.append(
                            (true_sh - client_h, css, client_h, true_sh)
                        )
                _container_hints.sort(reverse=True)

                # Find top-level regions and their heights
                _region_heights = []
                for blk in t1_state.blocks:
                    x, y, w, h = blk.bbox_px
                    css = blk.css_selector or blk.var_name
                    dom_depth = (getattr(blk, 'dom_path', '') or '').count('/')
                    if h > 80 and dom_depth <= 4 and w > 400:
                        _region_heights.append((h, y, css))
                _region_heights.sort(key=lambda x: x[1])

                _ratio_msg = ""
                if _style_match and _overflow_px > 10:
                    _css = _style_match.group(1)
                    _body = state.current_code[state.current_code.find('<body'):]

                    # Compute rendered vertical spacing (CSS value × element count)
                    _rendered_spacing = 0
                    # Find each CSS class with vertical spacing
                    for _cls_m in _re.finditer(r'\.([a-z][\w-]*)\s*\{([^}]*)\}', _css):
                        _cls_name = _cls_m.group(1)
                        _rules = _cls_m.group(2)
                        # Sum vertical spacing in this class
                        _vert = 0
                        for _sp_m in _re.finditer(r'(?:padding-top|padding-bottom|margin-top|margin-bottom)\s*:\s*(\d+)px', _rules):
                            _vert += int(_sp_m.group(1))
                        for _sp_m in _re.finditer(r'(?:padding|margin)\s*:\s*(\d+)px\s+\d+px\s+(\d+)px', _rules):
                            _vert += int(_sp_m.group(1)) + int(_sp_m.group(2))
                        for _sp_m in _re.finditer(r'gap\s*:\s*(\d+)px', _rules):
                            _vert += int(_sp_m.group(1))
                        if _vert > 0:
                            # Count how many DOM elements use this class
                            _count = max(1, len(_re.findall(rf'class="[^"]*{_cls_name}[^"]*"', _body)))
                            _rendered_spacing += _vert * _count

                    # Also add non-class rules (element selectors)
                    _rendered_spacing = max(_rendered_spacing, 100)  # floor

                    if _rendered_spacing > 0:
                        _ratio = _overflow_px / _rendered_spacing
                        _pct = int(_ratio * 100)
                        # Build leverage-based guidance instead of fixed tiers
                        _lev_items = []
                        # Check if this is an absolute-layout slide
                        _abs_count = _css.lower().count('position:absolute') + _css.lower().count('position: absolute')
                        _is_absolute_layout = _abs_count >= 8

                        for _cls_m2 in _re.finditer(r'\.([a-z][\w-]*)\s*\{([^}]*)\}', _css):
                            _cn = _cls_m2.group(1)
                            _rl = _cls_m2.group(2)
                            _v = 0
                            # Standard spacing properties
                            for _sp2 in _re.finditer(r'(padding|margin|gap)[^:]*:\s*([^;]+)', _rl):
                                _nums = _re.findall(r'(\d+)px', _sp2.group(2))
                                if _nums:
                                    _v = max(_v, max(int(n) for n in _nums))
                            # For absolute layouts, also consider top/height/font-size
                            if _is_absolute_layout:
                                for _sp2 in _re.finditer(r'(top|height|font-size)\s*:\s*(\d+)px', _rl):
                                    _prop = _sp2.group(1)
                                    _val = int(_sp2.group(2))
                                    if _prop == 'top' and _val > 100:
                                        _v = max(_v, _val // 10)  # top values are movable
                                    elif _prop == 'height' and _val > 50:
                                        _v = max(_v, _val // 5)  # height is shrinkable
                                    elif _prop == 'font-size' and _val > 14:
                                        _v = max(_v, _val)
                            if _v >= 8:
                                _ec = max(1, len(_re.findall(rf'class="[^"]*\b{_cn}\b[^"]*"', _body)))
                                _lev_items.append((f'.{_cn}', _v, _ec, _v * _ec))
                        for _tg in ('td', 'th', 'li'):
                            _tg_m = _re.search(rf'(?:^|\n)\s*{_tg}\s*\{{([^}}]+)\}}', _css)
                            if _tg_m:
                                _ec = max(1, len(_re.findall(rf'<{_tg}\b', _body)))
                                _v = 0
                                for _sp2 in _re.finditer(r'(padding|margin)[^:]*:\s*([^;]+)', _tg_m.group(1)):
                                    _nums = _re.findall(r'(\d+)px', _sp2.group(2))
                                    if _nums:
                                        _v = max(_v, max(int(n) for n in _nums))
                                if _v >= 6 and _ec >= 2:
                                    _lev_items.append((_tg, _v, _ec, _v * _ec))
                        _lev_items.sort(key=lambda x: -x[3])
                        _top3 = _lev_items[:3]
                        if _top3:
                            _lev_str = "; ".join(
                                f"{s}({v}px×{c}={p})" for s, v, c, p in _top3
                            )
                            _tier = f"highest leverage: {_lev_str}"
                        elif _ratio < 0.30:
                            _tier = "moderate compression needed"
                        else:
                            _tier = "aggressive compression needed"
                        _ratio_msg = (
                            f"\n  RATIO: overflow/total_spacing = {_pct}% → {_tier}"
                        )

                _details = ""
                if _container_hints:
                    _hint_lines = []
                    for _surplus, css, ch, tsh in _container_hints[:6]:
                        # Find the block to get position info
                        _blk_pos = ""
                        for _blk in t1_state.blocks:
                            _sel = _blk.css_selector or _blk.var_name
                            if _sel == css:
                                _bx, _by, _bw, _bh = _blk.bbox_px
                                _blk_pos = f" at ({int(_bx)}, {int(_by)}, {int(_bw)}×{int(_bh)}px)"
                                break
                        _hint_lines.append(
                            f"  {css}{_blk_pos}: needs {tsh:.0f}px, "
                            f"has {ch:.0f}px → compress {tsh-ch:.0f}px within THIS container"
                        )
                    _details = (
                        "\nPer-container overflow (fix each independently):\n"
                        + "\n".join(_hint_lines)
                        + "\n  → Compress content INSIDE each container to fit. "
                        "Do NOT uniformly shrink everything."
                    )
                elif _region_heights and _overflow_px > 30:
                    _region_lines = [
                        f"  {css}: y={y:.0f}px, height={h:.0f}px"
                        for h, y, css in _region_heights[:6]
                    ]
                    _total_h = sum(h for h, _, _ in _region_heights)
                    _details = (
                        f"\nRegion heights (total ~{_total_h:.0f}px, "
                        f"need to fit in 720px, reduce by {_overflow_px}px):\n"
                        + "\n".join(_region_lines)
                    )

                _stall_msg = ""
                if _overflow_stalled:
                    _stall_msg = (
                        " STALLED: overflow has not decreased in consecutive verifications."
                    )

                # Show which elements extend past canvas (all layout types)
                _element_overflow_details = ""
                if _overflow_px > 5 and t1_state:
                    _overflowing = []
                    for _blk in t1_state.blocks:
                        _bx, _by, _bw, _bh = _blk.bbox_px
                        _bottom = _by + _bh
                        if _bottom > 720 and _bh > 10 and _bh < 700 and _bw > 30:
                            _sel = getattr(_blk, 'css_selector', '') or getattr(_blk, 'var_name', '') or getattr(_blk, 'tag', '?')
                            _sel = _sel[:30]
                            _excess = int(_bottom - 720)
                            _overflowing.append((_sel, int(_by), int(_bh), _excess))
                    if _overflowing:
                        _overflowing.sort(key=lambda x: -x[3])
                        _ov_lines = [
                            f"  {sel}: top={y}px height={h}px → bottom={y+h}px ({exc}px past canvas)"
                            for sel, y, h, exc in _overflowing[:8]
                        ]
                        _element_overflow_details = (
                            "\nElements extending past 720px canvas:\n"
                            + "\n".join(_ov_lines)
                        )

                _deficit_msg = (
                    f"\n⚠️ OVERFLOW REMAINING: content extends to "
                    f"{int(_max_bottom)}px, which is {_overflow_px}px past "
                    f"the 720px canvas."
                    f"\n  DEFICIT: ~{_overflow_px}px to save."
                    f"{_ratio_msg}{_stall_msg}{_details}"
                    f"{_element_overflow_details}"
                )
                lines.append(_deficit_msg)
                logger.info(
                    "Slide %d verify: DEFICIT=%dpx, %s",
                    state.slide_id, _overflow_px,
                    _ratio_msg.strip() if _ratio_msg else "no leverage data",
                )
            elif sig_total > 0:
                # Canvas fits but internal clip/overlap remains
                _clip_count = len(
                    targeted_sig.get("clipped", [])
                    + targeted_sig.get("canvas_truncation", [])
                )
                _overlap_count = len(
                    targeted_sig.get("overlap", [])
                    + targeted_sig.get("occlusion", [])
                )
                _oob_count = len(targeted_sig.get("out_of_bounds", []))
                _text_ovf_count = len(targeted_sig.get("text_overflow", []))
                if _clip_count + _overlap_count + _oob_count + _text_ovf_count > 0:
                    # Build overlap details for absolute-layout slides
                    _overlap_detail = ""
                    if _overlap_count > 0 and _style_match if "_style_match" in dir() else _re.search(r"<style>(.*?)</style>", state.current_code, _re.DOTALL):
                        _abs_c = (
                            state.current_code.lower().count("position:absolute")
                            + state.current_code.lower().count("position: absolute")
                        )
                        if _abs_c >= 6 and hasattr(t1_state, 'overlap_pairs'):
                            _seen = set()
                            _ov_items = []
                            for a_id, b_id, _area in t1_state.overlap_pairs[:8]:
                                key = (a_id, b_id)
                                if key not in _seen:
                                    _seen.add(key)
                                    _a = next((b for b in t1_state.blocks if b.block_id == a_id), None)
                                    _b = next((b for b in t1_state.blocks if b.block_id == b_id), None)
                                    if _a and _b:
                                        _a_sel = (_a.css_selector or _a.var_name or "")[:25]
                                        _b_sel = (_b.css_selector or _b.var_name or "")[:25]
                                        _ax, _ay, _aw, _ah = _a.bbox_px
                                        _bx, _by, _bw, _bh = _b.bbox_px
                                        _ov_items.append(
                                            f"  {_a_sel}({int(_ax)},{int(_ay)} {int(_aw)}x{int(_ah)}) ↔ "
                                            f"{_b_sel}({int(_bx)},{int(_by)} {int(_bw)}x{int(_bh)})"
                                        )
                            if _ov_items:
                                _overlap_detail = "\nOverlapping element pairs:\n" + "\n".join(_ov_items[:6])

                    lines.append(
                        f"\n⚠️ {sig_total} SPATIAL ISSUES REMAIN: "
                        f"{_overlap_count} overlaps, {_clip_count} clipped, "
                        f"{_oob_count} out-of-bounds, {_text_ovf_count} text overflow."
                        f"{_overlap_detail}"
                    )
                else:
                    lines.append(
                        "\n✅ All content fits within 720px canvas. "
                        "No spatial overflow remaining."
                    )
            elif sig_total == 0:
                lines.append(
                    "\n✅ All content fits within 720px canvas. "
                    "No spatial overflow remaining."
                )

            # Container utilization: pure facts about fixed-height containers
            if _style_match:
                try:
                    from playwright.sync_api import sync_playwright as _sync_pw
                    _pw = _sync_pw().start()
                    _br = _pw.chromium.launch()
                    _pg = _br.new_page(viewport={"width": 1280, "height": 720})
                    _pg.set_content(state.current_code)
                    _util = _pg.evaluate('''() => {
                        const results = [];
                        const all = document.querySelectorAll('*');
                        for (const el of all) {
                            const style = getComputedStyle(el);
                            const h = el.getBoundingClientRect().height;
                            if (h > 50 && h < 650 && style.height.includes('px')) {
                                const origH = el.style.height;
                                const origMinH = el.style.minHeight;
                                const origOvf = el.style.overflow;
                                el.style.height = 'auto';
                                el.style.minHeight = '0';
                                el.style.overflow = 'visible';
                                const contentH = el.scrollHeight;
                                el.style.height = origH;
                                el.style.minHeight = origMinH;
                                el.style.overflow = origOvf;
                                const unused = Math.round(h - contentH);
                                if (Math.abs(unused) > 10) {
                                    const cls = el.className
                                        ? '.' + el.className.split(' ')[0]
                                        : el.tagName.toLowerCase();
                                    results.push({
                                        sel: cls,
                                        containerH: Math.round(h),
                                        contentH,
                                        unused
                                    });
                                }
                            }
                        }
                        return results.sort((a,b) => b.unused - a.unused).slice(0, 6);
                    }''')
                    _br.close()
                    _pw.stop()
                    if _util:
                        _util_lines = []
                        for w in _util:
                            if w['unused'] > 0:
                                _util_lines.append(
                                    f"  {w['sel']}: height={w['containerH']}px, "
                                    f"content={w['contentH']}px, "
                                    f"unused={w['unused']}px"
                                )
                            else:
                                _util_lines.append(
                                    f"  {w['sel']}: height={w['containerH']}px, "
                                    f"content={w['contentH']}px, "
                                    f"overflow={-w['unused']}px"
                                )
                        if _util_lines:
                            lines.append(
                                "\nCONTAINER UTILIZATION "
                                "(fixed-height containers vs actual content):\n"
                                + "\n".join(_util_lines)
                            )
                except Exception:
                    pass

        result_text = "\n".join(lines)
        if self._enable_render_preview:
            encoded = self._render_slide_to_base64(state.current_code, state)
            if encoded:
                state.latest_visual_checkpoint_code = state.current_code
                state.latest_visual_checkpoint_revision = state.layout_revision
                state.latest_visual_checkpoint_hard_valid = safe_checkpoint_valid
                state.latest_visual_checkpoint_targeted_issues = sig_total
                return ([
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{encoded}",
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            result_text
                            + "\n\nVISUAL VERIFICATION: the image above is the exact "
                            "revision measured by this verify_layout call. Judge the "
                            "original issue from the rendered composition first, using "
                            "detector findings as localized evidence rather than an "
                            "automatic score. Continue the current edit chain when the "
                            "remaining problem has a clear same-region closure; roll back "
                            "when content, media, roles, or an unrelated region is visibly "
                            "damaged."
                        ),
                    },
                ], False)

        return result_text, False

    # ── Meta-content / text-diff / alignment checks ──────────────

    _META_PATTERNS = [
        (r'\bREMOVE\s*[-—–:]', "REMOVE instruction"),
        (r'\bReplace\s+with\s*:', "Replace with: instruction"),
        (r'\bSuggested\s+replacement\b', "Suggested replacement note"),
        (r'\bTODO\b', "TODO note"),
        (r'\bFIXME\b', "FIXME note"),
        (r'\bNOTE:\s', "NOTE: annotation"),
        (r'\bPlaceholder\b', "Placeholder text"),
        (r'\bNo verified replacement\b', "Unverified placeholder"),
        (r'\bSource-supported\b', "Source-supported annotation"),
        (r'\b(?:insert|add|include)\s+(?:here|below|above)\b', "Editing instruction"),
        (r'\bUse source[\s-]+(?:grounded|backed|wording)\b', "Source instruction leak"),
        (r'\bUse source wording\b', "Source instruction leak"),
        (r'\bExisting solutions against\b.*\bcan be classified\b', "Pasted source text"),
        (r'\bUse source-grounded content only\b', "Source instruction leak"),
        (r'\bMUST contain\b', "Repair instruction leak"),
        (r'\bEnsure these appear\b', "Repair instruction leak"),
        (r'\bContent Verification Checklist\b', "Repair checklist leak"),
        # Evaluator judgment/correction language leaked into slide text
        (r'\bcorrected to source[- ]grounded\b', "Evaluator correction leak"),
        (r'\bnot supported by (?:the )?(?:paper|source|evidence)\b', "Evaluator judgment leak"),
        (r'\bno source support\b', "Evaluator judgment leak"),
        (r'\bunsupported (?:claim|text|content|number)\b', "Evaluator judgment leak"),
        (r'\bCorrection:\s', "Evaluator correction label"),
    ]

    @staticmethod
    def _extract_visible_text(html: str) -> str:
        """Extract visible text from HTML, ignoring tags and style/script."""
        import re
        # Remove style and script blocks
        text = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.S)
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.S)
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', ' ', text)
        # Decode HTML entities
        text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        text = text.replace('&quot;', '"').replace('&#39;', "'").replace('&nbsp;', ' ')
        return text

    def _scan_meta_content(self, code: str) -> list[str]:
        """Scan HTML for instruction/editorial text that should not appear."""
        import re
        visible = self._extract_visible_text(code)
        warnings = []
        for pattern, label in self._META_PATTERNS:
            matches = list(re.finditer(pattern, visible, re.IGNORECASE))
            if matches:
                # Show context around first match
                m = matches[0]
                start = max(0, m.start() - 20)
                end = min(len(visible), m.end() + 40)
                context = visible[start:end].strip().replace('\n', ' ')
                warnings.append(f'{label}: "...{context}..."')
        return warnings

    def _text_diff_summary(self, original: str, current: str, state=None) -> str:
        """Compare visible text before/after and summarize changes."""
        orig_text = self._extract_visible_text(original)
        curr_text = self._extract_visible_text(current)

        orig_words = orig_text.split()
        curr_words = curr_text.split()
        orig_set = set(orig_words)
        curr_set = set(curr_words)

        removed = orig_set - curr_set
        added = curr_set - orig_set

        if not removed and not added:
            return ""

        lines = ["📝 TEXT CHANGES vs original:"]
        if removed:
            sample = sorted(removed)[:8]
            lines.append(f"  Removed words ({len(removed)}): {', '.join(sample)}" +
                        (" ..." if len(removed) > 8 else ""))
        if added:
            sample = sorted(added)[:8]
            lines.append(f"  Added words ({len(added)}): {', '.join(sample)}" +
                        (" ..." if len(added) > 8 else ""))

        # Retention ratio
        if orig_words:
            retained = len([w for w in orig_words if w in curr_set])
            pct = retained * 100 // len(orig_words)
            lines.append(f"  Word retention: {pct}% ({retained}/{len(orig_words)})")

        # Inject rollback learning context
        if hasattr(state, 'rollback_history') and state.rollback_history:
            lines.append(
                "\nPREVIOUS FAILED ATTEMPTS (avoid repeating these mistakes):"
            )
            for j, reason in enumerate(state.rollback_history, 1):
                lines.append(f"  {j}. {reason}")

        return "\n".join(lines)

    @staticmethod
    def _text_diff_source_guidance(state: AgentState, added_words: int) -> str:
        """Explain when visible-text edits need external source lookup.

        Existing visible support copy is sufficient grounding for an explicitly
        authorized meaning-preserving shortening. Source lookup is needed only
        when an edit adds or changes factual content beyond that visible copy.
        """
        if added_words < 5:
            return ""
        if getattr(state, "allow_support_copy_compression", False):
            return (
                "\nℹ AUTHORIZED SUPPORT-COPY CALIBRATION: existing visible "
                "support text is the semantic source for a meaning-preserving "
                "shortening. Do not call search_source merely because concise "
                "wording uses different tokens. Search only if you add or change "
                "a factual proposition, number, named entity, technical term, or "
                "claim beyond what the original visible text states."
            )
        return (
            "\n💡 You added significant new text. Before submitting, use "
            "search_source to verify any new claims, numbers, or technical terms "
            "are supported by the paper. Fabricated content is a critical issue."
        )

    @staticmethod
    def _is_html_code(code: str) -> bool:
        """Check if code is HTML (vs python-pptx)."""
        return "<!DOCTYPE" in code or "<html" in code or ("<head" in code and "<style" in code)

    @staticmethod
    def _spatial_regression_policy_message(
        regression_total: int,
        hard_quality_failure: str = "",
    ) -> str:
        """Return agent-facing policy for spatial regressions.

        Spatial regressions need issue-level interpretation before submission,
        but they are not automatically a signal to roll back. During a body,
        table, or card reflow the first edit can create a same-cluster
        intermediate state that the next edit should close. Text loss and
        unrelated damage are different: those should be rolled back or
        explicitly reversed.
        """
        if hard_quality_failure:
            return (
                f"REGRESSION CHECK: {regression_total} new deterministic "
                "regression(s) vs baseline. These were introduced by the "
                "current edits and must be fixed before submit.\n"
                "INVALID INTERMEDIATE CHECKPOINT: a hard repair-quality "
                f"failure is also present ({hard_quality_failure}). Do not "
                "continue a same-cluster closure chain or design a new repair "
                "from this checkpoint. Restore a checkpoint without that hard "
                "failure first, verify the restored code, then continue from "
                "the cleanest useful state."
            )

        return (
            f"REGRESSION CHECK: {regression_total} new deterministic "
            "regression(s) vs baseline. They appeared after the current edits and "
            "require issue-level judgment before submit; the count alone does not "
            "make the checkpoint worse.\n"
            "CHECKPOINT NEEDS INTERPRETATION: compare each named measurement with "
            "the original issue, owning region, and current revision evidence. "
            "NEW OVERLAP and NEW CLIPPED findings are real spatial defects — "
            "fix them before marking the plan step as done. "
            "Roll back when source text, media, hierarchy, or information-bearing roles "
            "were damaged, an unrelated region regressed, or the current "
            "repair family has no credible information-preserving closure path. "
            "If this verify result also contains SHIPMENT GATE FAILED or "
            "CHECKPOINT INVALID, that hard quality failure takes precedence and "
            "the invalid checkpoint must be restored before continuing."
        )

    @staticmethod
    def _representative_spatial_regressions(
        regressions: dict[str, list[tuple[str, object]]],
        *,
        max_per_group: int = 6,
        max_total: int = 18,
    ) -> tuple[list[tuple[str, str, object]], dict[str, int]]:
        """Keep regression feedback causal without flooding the agent.

        Nested HTML structures can produce dozens of pairwise findings from one
        unfinished allocation. Preserve the full counts and regression signature,
        but show a bounded set of concrete examples from every detector group.
        """
        shown: list[tuple[str, str, object]] = []
        omitted: dict[str, int] = {}
        for group, items in regressions.items():
            group_limit = min(max_per_group, max_total - len(shown))
            selected = items[:max(0, group_limit)]
            shown.extend((group, kind, payload) for kind, payload in selected)
            remaining = len(items) - len(selected)
            if remaining:
                omitted[group] = remaining
            if len(shown) >= max_total:
                for later_group, later_items in list(regressions.items())[
                    list(regressions).index(group) + 1:
                ]:
                    if later_items:
                        omitted[later_group] = len(later_items)
                break
        return shown, omitted

    @staticmethod
    def _changed_css_properties_from_edits(edits: list[dict]) -> set[str]:
        """Best-effort list of CSS properties actually changed by edits."""
        def declarations_by_property(css: str) -> dict[str, Counter[str]]:
            cleaned = re.sub(r"/\*.*?\*/", "", str(css or ""), flags=re.DOTALL)
            declarations: dict[str, Counter[str]] = {}
            for match in re.finditer(
                r"(?:^|[;{])\s*([a-zA-Z-]+)\s*:\s*([^;{}]+)",
                cleaned,
            ):
                prop = match.group(1).lower()
                value = re.sub(r"\s+", " ", match.group(2).strip())
                if not value:
                    continue
                declarations.setdefault(prop, Counter())[value] += 1
            return declarations

        props: set[str] = set()
        for edit in edits:
            old_props = declarations_by_property(str(edit.get("search", "")))
            new_props = declarations_by_property(str(edit.get("replace", "")))
            for prop in set(old_props) | set(new_props):
                if old_props.get(prop, Counter()) != new_props.get(prop, Counter()):
                    props.add(prop)
        return props

    @staticmethod
    def _html_edit_blob(edits: list[dict]) -> str:
        """Concatenate edit snippets for lightweight selector heuristics."""
        parts: list[str] = []
        for edit in edits:
            parts.append(str(edit.get("search", "")))
            parts.append(str(edit.get("replace", "")))
            parts.append(str(edit.get("insert_after", "")))
        return "\n".join(parts).lower()

    @staticmethod
    def _edit_scope_labels(edits: list[dict], *, limit: int = 16) -> list[str]:
        """Best-effort semantic scope labels for a CSS/HTML edit batch."""
        labels: list[str] = []
        seen: set[str] = set()

        def add(label: str) -> None:
            if label and label not in seen and len(labels) < limit:
                seen.add(label)
                labels.append(label)

        for edit in edits:
            blob = "\n".join((
                str(edit.get("search", "")),
                str(edit.get("replace", "")),
                str(edit.get("insert_after", "")),
            ))
            for _, class_value in re.findall(
                r"class\s*=\s*(['\"])(.*?)\1",
                blob,
                flags=re.IGNORECASE | re.DOTALL,
            ):
                for name in re.split(r"\s+", class_value.strip()):
                    if re.fullmatch(r"[A-Za-z_][\w-]*", name):
                        add(f".{name}")
            for name in re.findall(r"\.([A-Za-z_][\w-]*)", blob):
                add(f".{name}")
            for name in re.findall(r"#([A-Za-z_][\w-]*)", blob):
                add(f"#{name}")
            for tag in re.findall(
                r"(?:^|[,{>+~\s])(table|thead|tbody|tr|th|td|svg|img|figure)(?=[\s.{:#>+~,{]|$)",
                blob,
                flags=re.IGNORECASE | re.MULTILINE,
            ):
                add(tag.lower())
        return labels

    @staticmethod
    def _mentions_class_or_selector(blob: str, name: str) -> bool:
        """Return whether an edit blob mentions a CSS class/selector name."""
        escaped = re.escape(name)
        return bool(
            re.search(rf"(?<![\w-])\.{escaped}(?![\w-])", blob)
            or re.search(rf"class\s*=\s*['\"][^'\"]*\b{escaped}\b", blob)
        )

    @classmethod
    def _looks_like_table_dashboard_pressure(cls, state: AgentState) -> bool:
        from .dashboard_heuristics import looks_like_table_dashboard_pressure
        return looks_like_table_dashboard_pressure(state)

    @classmethod
    def _looks_like_table_dashboard_pressure_from(
        cls,
        code: str,
        issue_types: set[str],
    ) -> bool:
        from .dashboard_heuristics import looks_like_table_dashboard_pressure_from
        return looks_like_table_dashboard_pressure_from(code, issue_types)

    @staticmethod
    def _support_compression_role_regressions(
        before_code: str,
        after_code: str,
    ) -> list[str]:
        """Find support-compression edits that delete information roles."""
        try:
            from bs4 import BeautifulSoup

            before = BeautifulSoup(before_code, "html.parser")
            after = BeautifulSoup(after_code, "html.parser")
        except Exception:
            return []

        regressions: list[str] = []
        for selector, label in ((".metric", "metric item"), (".kpi", "KPI item")):
            before_count = len(before.select(selector))
            after_count = len(after.select(selector))
            if after_count < before_count:
                regressions.append(
                    f"{label} count {before_count}->{after_count}"
                )

        # Copy compression can shorten prose inside a terminal support branch,
        # but it cannot dissolve a distinct takeaway/conclusion branch into a
        # metric note. Match semantic class families rather than one project
        # selector so a branch may be rewrapped or renamed without false failure.
        support_stems = ("finding", "takeaway", "conclusion", "recommendation")

        def semantic_support_count(soup) -> int:
            return sum(
                1
                for element in soup.find_all(class_=True)
                if any(
                    stem in str(class_name).lower()
                    for class_name in (element.get("class") or [])
                    for stem in support_stems
                )
            )

        before_support = semantic_support_count(before)
        after_support = semantic_support_count(after)
        if after_support < before_support:
            regressions.append(
                "terminal findings/takeaway branch count "
                f"{before_support}->{after_support}"
            )
        return regressions

    @staticmethod
    def _dashboard_measured_spatial_state(
        state: AgentState,
        *,
        allow_previous_revision: bool = False,
    ):
        from .dashboard_heuristics import dashboard_measured_spatial_state
        return dashboard_measured_spatial_state(state, allow_previous_revision=allow_previous_revision)

    @staticmethod
    def _dashboard_descendant_extent_measurements(spatial_state) -> list[dict]:
        from .dashboard_heuristics import dashboard_descendant_extent_measurements
        return dashboard_descendant_extent_measurements(spatial_state)

    @classmethod
    def _dashboard_descendant_plan_note(
        cls,
        steps: list[PlanStep],
        state: AgentState,
        planning_context: str = "",
    ) -> str:
        from .dashboard_heuristics import dashboard_descendant_plan_note
        return dashboard_descendant_plan_note(steps, state, planning_context)

    @classmethod
    def _dashboard_parent_descendant_patch_warning(
        cls,
        state: AgentState,
        edit_blob: str,
    ) -> str:
        from .dashboard_heuristics import dashboard_parent_descendant_patch_warning
        return dashboard_parent_descendant_patch_warning(state, edit_blob)

    @classmethod
    def _dashboard_repeated_owner_budget_warning(
        cls,
        state: AgentState,
        css: str,
    ) -> str:
        from .dashboard_heuristics import dashboard_repeated_owner_budget_warning
        return dashboard_repeated_owner_budget_warning(state, css)

    @classmethod
    def _dashboard_terminal_support_patch_warning(
        cls,
        state: AgentState,
        css: str,
    ) -> str:
        from .dashboard_heuristics import dashboard_terminal_support_patch_warning
        return dashboard_terminal_support_patch_warning(state, css)

    @classmethod
    def _names_role_demand_calibration(cls, text: str) -> bool:
        """Return whether a plan names a concrete same-topology fit hypothesis.

        A repeated-card repair can be falsifiable without inventing a new DOM
        topology. Coordinated calibration of the focal, repeated-detail, and
        support roles is a real allocation hypothesis when the plan names the
        roles it will change instead of merely promising that everything fits.
        """
        low = str(text or "").lower()
        names_calibration = any(
            term in low
            for term in (
                "recalibrate", "calibrate", "compact", "tighten", "reduce",
                "shorten", "condense", "role rhythm", "copy calibration",
            )
        )
        role_groups = (
            ("score", "focal", "kpi", "value scale", "value size"),
            ("metric", "detail row", "repeated row", "row rhythm"),
            (
                "support copy", "explanatory copy", "finding", "takeaway",
                "support rhythm", "terminal support",
            ),
        )
        named_groups = sum(
            any(term in low for term in group)
            for group in role_groups
        )
        return names_calibration and named_groups >= 2

    @classmethod
    def _dashboard_plan_implementation_warning(
        cls,
        state: AgentState,
        *,
        action_text: str = "",
        cluster_complete: bool = True,
    ) -> str:
        from .dashboard_heuristics import dashboard_plan_implementation_warning
        return dashboard_plan_implementation_warning(state, action_text=action_text, cluster_complete=cluster_complete)

    @classmethod
    def _edit_cluster_execution_coverage_note(
        cls,
        state: AgentState,
        *,
        before_code: str,
        after_code: str,
        action_text: str,
        cluster_complete: bool,
    ) -> str:
        from .dashboard_heuristics import edit_cluster_execution_coverage_note
        return edit_cluster_execution_coverage_note(state, before_code=before_code, after_code=after_code, action_text=action_text, cluster_complete=cluster_complete)

    @classmethod
    def _dashboard_variable_track_patch_warning(
        cls,
        state: AgentState,
        css: str,
    ) -> str:
        from .dashboard_heuristics import dashboard_variable_track_patch_warning
        return dashboard_variable_track_patch_warning(state, css)

    @classmethod
    def _dashboard_pending_coupled_compression_note(
        cls,
        state: AgentState,
        spatial_state,
    ) -> str:
        from .dashboard_heuristics import dashboard_pending_coupled_compression_note
        return dashboard_pending_coupled_compression_note(state, spatial_state)

    @classmethod
    def _dashboard_support_copy_checkpoint_note(
        cls,
        state: AgentState,
        before_code: str,
        after_code: str,
        *,
        cluster_complete: bool,
    ) -> str:
        from .dashboard_heuristics import dashboard_support_copy_checkpoint_note
        return dashboard_support_copy_checkpoint_note(state, before_code, after_code, cluster_complete=cluster_complete)

    @classmethod
    def _dashboard_coupled_cluster_guidance(
        cls, code: str = "", *, preview_enabled: bool = False,
    ) -> str:
        from .dashboard_heuristics import dashboard_coupled_cluster_guidance
        return dashboard_coupled_cluster_guidance(code, preview_enabled=preview_enabled)

    @staticmethod
    def _dashboard_measurement_context(spatial_text: str) -> str:
        from .dashboard_heuristics import dashboard_measurement_context
        return dashboard_measurement_context(spatial_text)

    @classmethod
    def _dashboard_decision_summary(
        cls,
        state: AgentState,
        spatial_state,
    ) -> str:
        from .dashboard_heuristics import dashboard_decision_summary
        return dashboard_decision_summary(state, spatial_state)

    @staticmethod
    def _dashboard_allocation_map(spatial_state) -> str:
        from .dashboard_heuristics import dashboard_allocation_map
        return dashboard_allocation_map(spatial_state)

    @staticmethod
    def _dashboard_fit_magnitude_cue(spatial_text: str) -> str:
        from .dashboard_heuristics import dashboard_fit_magnitude_cue
        return dashboard_fit_magnitude_cue(spatial_text)

    @classmethod
    def _dashboard_local_first_html_edit_message(
        cls,
        edits: list[dict],
        state: AgentState,
    ) -> str | None:
        from .dashboard_heuristics import dashboard_local_first_html_edit_message
        return dashboard_local_first_html_edit_message(edits, state)

    @classmethod
    def _dashboard_table_outer_frame_warning_from_edits(
        cls,
        edits: list[dict],
        state: AgentState,
    ) -> str:
        from .dashboard_heuristics import dashboard_table_outer_frame_warning_from_edits
        return dashboard_table_outer_frame_warning_from_edits(edits, state)

    @classmethod
    def _dashboard_coupled_cluster_warning_from_edits(
        cls,
        edits: list[dict],
        state: AgentState,
    ) -> str:
        from .dashboard_heuristics import dashboard_coupled_cluster_warning_from_edits
        return dashboard_coupled_cluster_warning_from_edits(edits, state)

    @classmethod
    def _broad_structural_html_edit_message(
        cls,
        edits: list[dict],
        state: AgentState,
    ) -> str | None:
        """Return a non-blocking debuggability warning for broad HTML edits."""
        if len(edits) <= 6 or not cls._is_html_code(state.current_code):
            return None
        spatial_issue_types = set(getattr(state, "issue_types", set())) - {"low_contrast"}
        if not (spatial_issue_types & STRUCTURAL_ISSUE_TYPES):
            return None
        changed_props = cls._changed_css_properties_from_edits(edits)
        structural_props = {
            "display", "position", "grid-template-columns",
            "grid-template-rows", "grid-template-areas", "grid-area",
            "flex", "flex-direction", "justify-content", "align-items",
            "gap", "row-gap", "column-gap", "width", "height",
            "min-width", "min-height", "max-width", "max-height",
            "padding", "padding-top", "padding-right", "padding-bottom",
            "padding-left", "margin", "margin-top", "margin-right",
            "margin-bottom", "margin-left", "top", "right", "bottom",
            "left", "font-size", "line-height", "overflow",
        }
        changed_structural = sorted(changed_props & structural_props)
        if not changed_structural:
            return None
        shown = ", ".join(changed_structural[:8])
        return (
            "\n\nBROAD EDIT ADVISORY: this accepted checkpoint changes several "
            f"geometry/typography properties ({shown}). Keep the causal hypothesis "
            "explicit and verify immediately so useful and harmful changes can be "
            "distinguished. A broad edit is appropriate when one coherent reflow "
            "requires it; unrelated changes are easier to diagnose as separate "
            "checkpoints. This is guidance, not a request to undo the edit."
        )

    def _verify_pptx_layout(self, code: str, slide_id: int, state: AgentState) -> dict:
        """Verify layout for python-pptx code."""
        import tempfile
        from pathlib import Path
        from ...modules.extractor import StructuralExtractor
        from ...modules.evaluators.geom_checks import DeterministicGeomChecks

        prs, error = state.codegen_compiler._compile_code(code)
        if error or prs is None:
            return {"compile_error": True}

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            prs.save(str(pptx_path))

            extractor = StructuralExtractor()
            extractions = extractor.extract(str(pptx_path))

            checker = DeterministicGeomChecks()
            issues = checker.check_all(extractions)

            counts = {}
            for iss in issues:
                counts[iss.issue_type] = counts.get(iss.issue_type, 0) + 1
            return counts

    def _verify_html_layout(self, code: str, slide_id: int) -> dict:
        """Verify layout for HTML code using Playwright DOM extraction.

        Renders HTML, extracts spatial state via DOM queries, and runs
        applicable geometry checks.
        """
        from .html_spatial_state import extract_html_slide_state
        from ...modules.evaluators.geom_checks import DeterministicGeomChecks
        from ...schemas.extraction import ExtractedObject, SlideExtraction

        # Extract spatial state from rendered HTML
        state_obj = extract_html_slide_state(slide_id, code)

        # Store for compact state display
        # (access via AgentState._last_html_state in _tool_verify_layout)
        self._last_html_spatial_state = state_obj

        # Convert to SlideExtraction for geom_checks compatibility
        objects = []
        for blk in state_obj.blocks:
            # Convert inches to EMU
            EMU = SlideDimensions.EMU_PER_INCH
            objects.append(ExtractedObject(
                object_id=blk.block_id,
                shape_name=blk.var_name,
                object_type=blk.shape_type if blk.shape_type != "title" else "text_box",
                bbox_emu=[
                    int(blk.x * EMU),
                    int(blk.y * EMU),
                    int(blk.w * EMU),
                    int(blk.h * EMU),
                ],
                text_content=" ".join(blk.text_lines),
                font_sizes_pt=[blk.font_size_pt] if blk.font_size_pt > 0 else [],
                has_image=(blk.shape_type == "picture"),
                z_order=0,
            ))

        extraction = SlideExtraction(
            slide_id=slide_id,
            slide_index=slide_id - 1,
            title=next((b.text_lines[0] for b in state_obj.blocks if b.shape_type == "title" and b.text_lines), ""),
            objects=objects,
            total_text_length=sum(b.text_chars for b in state_obj.blocks),
            total_objects=len(objects),
        )

        # Run geom checks — html_mode=True skips font/text_overflow/narrow_textbox
        # checks (CSS handles sizing), keeps overlap/OOB/meta-content/spelling
        checker = DeterministicGeomChecks(html_mode=True)
        issues = checker.check_all([extraction])

        counts = {}
        for iss in issues:
            counts[iss.issue_type] = counts.get(iss.issue_type, 0) + 1

        # Compute space utilization metrics for density feedback
        if objects:
            SLIDE_W, SLIDE_H = SlideDimensions.WIDTH_IN, SlideDimensions.HEIGHT_IN
            slide_area = SLIDE_W * SLIDE_H
            content_bottom = 0.0
            content_area = 0.0
            for blk in state_obj.blocks:
                bottom = blk.y + blk.h
                content_bottom = max(content_bottom, bottom)
                content_area += blk.w * blk.h
            coverage_pct = round(content_area / slide_area * 100, 1)
            bottom_empty_pct = round((1 - content_bottom / SLIDE_H) * 100, 1)
            counts["_space_coverage_pct"] = coverage_pct
            counts["_bottom_empty_pct"] = max(0, bottom_empty_pct)

        return counts

    def _tool_rollback(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Undo last N edits."""
        if len(state.checkpoints) <= 1:
            return "Nothing to rollback — already at original code.", False

        self._normalize_checkpoint_metadata(state)

        # Check if rollback would lose overflow progress
        _current_overflow = getattr(state, '_last_verified_overflow_px', 0)
        rollback_scope = str(action.get("scope", "") or "").lower()
        if rollback_scope == "cluster":
            target_code = (
                state.active_cluster_start_code
                if state.pending_edit_cluster and state.active_cluster_start_code is not None
                else state.last_cluster_start_code
            )
            if target_code is None:
                return "No recorded edit-cluster boundary is available to rollback.", False
            target_index = next(
                (
                    index for index in range(len(state.checkpoints) - 1, -1, -1)
                    if state.checkpoints[index] == target_code
                ),
                -1,
            )
            if target_index < 0:
                return "The recorded edit-cluster start is no longer in rollback history.", False
            steps = len(state.checkpoints) - target_index
        else:
            steps = action.get("steps", 1)
            steps = max(1, min(steps, len(state.checkpoints) - 1))
            target_index = len(state.checkpoints) - steps

        hard_repair_quality_failure = bool(
            getattr(state, "_last_verify_visual_compression_failed", False)
            or getattr(state, "_last_verify_scope_failed", False)
        )

        # Block rollback ONLY when it would lose overflow progress AND the
        # target is the original code. Rolling back to an intermediate checkpoint
        # (e.g., a verified KPI-only state) is always allowed — the agent needs
        # freedom to try different repair strategies from partial-progress states.
        #
        # Successful repair trajectories show the pattern:
        #   1. Fix region A (KPI) → checkpoint
        #   2. Try fixing region B (body) → fails → rollback to A checkpoint
        #   3. Try different strategy for B → succeeds
        # Blocking rollback to intermediate checkpoints prevents this pattern.
        _target_code = state.checkpoints[target_index]
        _t0_overflow = getattr(state, '_t0_overflow_px', 0)
        _target_is_original = (_target_code == state.original_code)
        if (
            _target_is_original
            and _current_overflow > 20
            and _current_overflow < _t0_overflow * 0.7
            and not hard_repair_quality_failure
        ):
            # Only block rollback to ORIGINAL when real progress exists.
            # Rolling back to intermediate checkpoints is always allowed.
            return (
                f"ROLLBACK BLOCKED: you would lose all progress (overflow reduced "
                f"from {_t0_overflow:.0f}px to {_current_overflow:.0f}px). "
                f"Roll back to a more recent checkpoint instead, or continue "
                f"editing from the current state."
            ), False

        state.current_code = state.checkpoints[target_index]
        restored_label = state.checkpoint_labels[target_index]

        # Record rollback context so future edits can learn from failures
        _rollback_reason = str(action.get("reasoning", action.get("reason", "")) or "")[:200]
        if _rollback_reason:
            state.rollback_history.append(_rollback_reason)
            # Keep only last 3 to avoid context bloat
            if len(state.rollback_history) > 3:
                state.rollback_history = state.rollback_history[-3:]
        if target_index < len(state.checkpoint_text_loss):
            state.cumulative_words_lost = state.checkpoint_text_loss[target_index]
        else:
            state.cumulative_words_lost = len(
                _meaningful_visible_words_lost(state.original_code, state.current_code)
            )
        state.checkpoints = state.checkpoints[:target_index]
        state.checkpoint_text_loss = state.checkpoint_text_loss[:target_index]
        state.checkpoint_labels = state.checkpoint_labels[:target_index]
        state.current_checkpoint_label = restored_label
        state.pending_edit_cluster = False
        state.pending_edit_scopes = []
        state.active_cluster_start_code = None
        state.active_cluster_start_text_loss = 0
        state.active_cluster_start_label = ""
        self._invalidate_verify_after_code_change(state, "rollback changed code")
        hard_failure_followup = ""
        if hard_repair_quality_failure:
            hard_failure_followup = (
                " This rollback left a hard-quality-failure checkpoint. "
                "Verify the restored checkpoint before any new edit. If the "
                "same hard failure is still reported, rollback another step "
                "instead of branching a new plan from that state."
            )
        return (
            f"Rolled back {steps} step(s)"
            + (" to the start of the last edit cluster" if rollback_scope == "cluster" else "")
            + f". Restored checkpoint: {restored_label}. "
            f"{len(state.checkpoints) - 1} checkpoint(s) remaining. "
            "Previous verify_layout feedback is now stale; call verify_layout "
            "before marking a plan step done or submitting. ROLLBACK DEPTH CHECK: "
            "a restored historical checkpoint is not automatically the clean "
            "boundary before an abandoned strategy. If verification shows that "
            "the rejected topology, hidden content, or its spatial damage is still "
            "present, roll back farther or replace that strategy's patch before "
            "branching into a new plan."
            f"{hard_failure_followup}"
        ), True

    def _tool_delete_shape(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Delete a shape by variable name."""
        var_name = action.get("var_name", "")
        if not var_name:
            return "No var_name provided.", False

        new_code = _delete_shape_code(state.current_code, var_name)
        if new_code == state.current_code:
            return (
                f"Could not find shape '{var_name}' in code. "
                f"Check the exact variable name."
            ), False

        ok = self._test_compile(
            new_code, state.codegen_compiler, state.case_dir, state.slide_id,
        )
        if not ok:
            return (
                f"Compile error after deleting '{var_name}'. "
                f"Deletion reverted — there may be remaining references."
            ), False

        prev_code = state.current_code
        projected_text_loss = state.cumulative_words_lost
        allows_content_change = (
            bool(state.issue_types & CONTENT_ACCURACY_ISSUE_TYPES)
            or state.allow_visible_text_change
        )
        if self._is_html_code(new_code) and not allows_content_change:
            dropped_high_value = _dropped_high_value_tokens(prev_code, new_code, limit=10)
            if dropped_high_value:
                shown = ", ".join(f'"{token}"' for token in dropped_high_value)
                return (
                    f"Deletion rejected because '{var_name}' contains value-bearing content: {shown}. "
                    "Relocate or restyle the element instead."
                ), False
            ordinary_lost = _meaningful_visible_words_lost(prev_code, new_code)
            projected_text_loss += len(ordinary_lost)
            if ordinary_lost and projected_text_loss > state.text_loss_budget:
                return (
                    f"Deletion rejected: visible-text loss budget would be {projected_text_loss}/"
                    f"{state.text_loss_budget}. Preserve or atomically relocate the "
                    "element, then retry; later DOM reflow remains allowed."
                ), False

        state.checkpoints.append(state.current_code)
        state.checkpoint_text_loss.append(state.cumulative_words_lost)
        state.current_code = new_code
        state.cumulative_words_lost = projected_text_loss
        self._invalidate_verify_after_code_change(state, "delete_shape changed code")
        # Immediate high-value-deletion warning (same rationale as apply_edits):
        # deleting a shape to clear an overlap can take benchmark numbers / model
        # names with it. Name them now so the agent can reconsider relocating.
        hv_warn = ""
        just_dropped = _dropped_high_value_tokens(prev_code, new_code, limit=10)
        if just_dropped:
            shown = ", ".join(f'"{d}"' for d in just_dropped)
            hv_warn = (
                f"\n⚠ Deleting '{var_name}' removed value-bearing content: {shown} "
                f"(numbers/metrics/method names — rule 8a). If you removed this only "
                f"to clear a layout issue, prefer relocating/shrinking it instead, or "
                f"reinstate these values elsewhere on the slide."
            )
        return (
            f"Deleted shape '{var_name}' and its code block. "
            f"Call verify_layout to check spatial state.{hv_warn}"
        ), True

    def _tool_reflow_layout(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """LLM-planned, deterministically-executed layout rearrangement.

        The agent provides a target layout as a JSON map of
        {var_name: {left, top, width, height}} for each element to move.
        The tool applies ALL coordinate changes atomically via AST-level
        code rewriting, compile-tests, and checks for overlaps.
        If no targets are provided, the tool uses _generate_layout_plan()
        to let the LLM compute optimal positions automatically.

        This is safer than manual apply_edits because:
        1. All coordinates change atomically (no intermediate broken states)
        2. Overlap check is automatic (revert if any overlap created)
        3. Handles both literal Inches() and add_textbox() constructor forms
        """
        # reflow_layout is pptx-specific; HTML uses CSS positions via apply_edits
        if self._is_html_code(state.current_code):
            return (
                "reflow_layout is not available for HTML slides. "
                "Use apply_edits to change CSS position/size values directly."
            ), False

        import re as _re

        targets = action.get("targets")  # Optional: agent-provided layout

        if targets and isinstance(targets, dict):
            # Agent explicitly provided target positions — validate them
            validated: dict[str, dict] = {}
            for var_name, pos in targets.items():
                if not isinstance(pos, dict):
                    continue
                entry = {}
                for k in ("left", "top", "width", "height"):
                    if k in pos:
                        try:
                            v = float(pos[k])
                            if 0.0 <= v <= 14.0:
                                entry[k] = round(v, 2)
                        except (ValueError, TypeError):
                            pass
                if "top" in entry:  # at minimum need a y position
                    validated[var_name] = entry
            if not validated:
                return "No valid targets provided. Each target needs at least {top: <y_inches>}.", False
            targets = validated
        else:
            # No explicit targets — use LLM layout planner
            spatial = extract_slide_state(state.slide_id, state.current_code)
            if not spatial or not spatial.blocks:
                return "Could not parse spatial state from current code.", False

            # Gather issues from the repair instance
            issues = self._current_issues
            bp_slide = getattr(state, 'bp_slide', None)

            plan = self._generate_layout_plan(
                state.current_code, issues, spatial, bp_slide,
            )
            if not plan:
                return (
                    "LLM layout planner could not generate a valid plan. "
                    "Try providing explicit targets: "
                    '{"tool": "reflow_layout", "targets": {"box_name": {"left": 0.5, "top": 1.5, "width": 12.0, "height": 2.0}}}'
                ), False
            targets = plan

        # === Apply targets to code ===
        # Strategy: find each var_name's coordinate assignments and rewrite
        new_code = state.current_code
        applied = 0
        failed_vars = []

        for var_name, pos in targets.items():
            var_applied = 0
            for attr in ("left", "top", "width", "height"):
                if attr not in pos:
                    continue
                val = pos[attr]

                # Pattern 1: var.left = Inches(X.XX)
                pat1 = _re.compile(
                    rf'({_re.escape(var_name)}\.{attr}\s*=\s*Inches\()[\d.]+(\))'
                )
                new_code, n = pat1.subn(rf'\g<1>{val}\2', new_code)
                var_applied += n

                # Pattern 2: add_textbox(Inches(...), Inches(...), Inches(...), Inches(...))
                # This is harder — we match the full constructor call for this var_name
                # and replace the Nth Inches() argument
                # Order: left=0, top=1, width=2, height=3
                attr_idx = {"left": 0, "top": 1, "width": 2, "height": 3}[attr]
                # Match: var_name = slide.shapes.add_textbox(Inches(A), Inches(B), ...)
                # or: var_name = slide.shapes.add_shape(..., Inches(A), Inches(B), ...)
                if var_applied == 0:
                    # Try to find constructor-style assignment
                    ctor_pattern = _re.compile(
                        rf'({_re.escape(var_name)}\s*=\s*slide\.shapes\.add_(?:textbox|shape)\([^)]*)'
                    )
                    match = ctor_pattern.search(new_code)
                    if match:
                        # Find all Inches() in this call and replace the right one
                        call_start = match.start()
                        # Find the closing paren
                        depth = 0
                        call_end = call_start
                        for ci in range(call_start, len(new_code)):
                            if new_code[ci] == '(':
                                depth += 1
                            elif new_code[ci] == ')':
                                depth -= 1
                                if depth == 0:
                                    call_end = ci + 1
                                    break
                        call_str = new_code[call_start:call_end]

                        # Find all Inches(X) in the call
                        inches_matches = list(_re.finditer(r'Inches\(([\d.]+)\)', call_str))
                        if len(inches_matches) > attr_idx:
                            m = inches_matches[attr_idx]
                            old_call = call_str
                            new_call = (
                                call_str[:m.start(1)]
                                + str(val)
                                + call_str[m.end(1):]
                            )
                            new_code = new_code[:call_start] + new_call + new_code[call_end:]
                            var_applied += 1

                applied += var_applied

            if var_applied == 0:
                failed_vars.append(var_name)

        if applied == 0:
            return (
                f"Could not find any Inches() assignments to modify. "
                f"Variables not found: {failed_vars}. "
                f"The code may use variables/expressions instead of literal Inches() values. "
                f"Use apply_edits for manual coordinate changes."
            ), False

        # Compile test
        ok = self._test_compile(
            new_code, state.codegen_compiler, state.case_dir, state.slide_id,
        )
        if not ok:
            return (
                f"Compile error after reflow — all changes reverted. "
                f"Modified {applied} coordinates for {len(targets) - len(failed_vars)} elements."
            ), False

        # Check for overlaps in reflowed layout
        new_spatial = extract_slide_state(state.slide_id, new_code)
        if new_spatial and new_spatial.overlap_pairs:
            overlap_desc = "; ".join(
                f"{a}↔{b}" for a, b, *_ in new_spatial.overlap_pairs[:5]
            )
            return (
                f"Reflow created {len(new_spatial.overlap_pairs)} overlap(s): {overlap_desc}. "
                f"All changes reverted. Provide adjusted targets with more spacing, "
                f"or use apply_edits for fine-grained control."
            ), False

        # Success — commit changes
        state.checkpoints.append(state.current_code)
        state.checkpoint_text_loss.append(state.cumulative_words_lost)
        state.current_code = new_code
        self._invalidate_verify_after_code_change(state, "reflow changed code")

        # Format feedback
        lines = [
            f"✓ Reflow applied: modified {applied} coordinates "
            f"across {len(targets) - len(failed_vars)}/{len(targets)} elements.",
        ]
        if failed_vars:
            lines.append(f"⚠ Could not find: {', '.join(failed_vars)}")
        lines.append("Call verify_layout to confirm no spatial violations.")
        if new_spatial:
            lines.append("\nNew layout:")
            for b in sorted(new_spatial.blocks, key=lambda b: b.y):
                lines.append(
                    f"  {b.var_name}: x={b.x:.2f} y={b.y:.2f} w={b.w:.2f} h={b.h:.2f}"
                )
        return "\n".join(lines), True

    @classmethod
    def _dashboard_coupled_plan_notes(
        cls,
        steps: list[PlanStep],
        state: AgentState,
        summary: str = "",
    ) -> list[str]:
        from .dashboard_heuristics import dashboard_coupled_plan_notes
        return dashboard_coupled_plan_notes(steps, state, summary)

    @classmethod
    def _is_recoverable_dashboard_dominant_compression(
        cls,
        state: AgentState,
        compression_reason: str,
    ) -> bool:
        from .dashboard_heuristics import is_recoverable_dashboard_dominant_compression
        return is_recoverable_dashboard_dominant_compression(state, compression_reason)

    @staticmethod
    def _strategy_fit_notes_for_step_text(step_text: str, step_idx: int) -> list[str]:
        """Return non-blocking feedback for mismatched repair-family plans."""
        low = (step_text or "").lower()
        notes: list[str] = []
        compression_terms = (
            "tighten", "shrink", "reduce", "smaller", "compact",
            "condense", "font-size", "line-height", "padding", "margin",
            "gap", "height budget", "spacing",
        )
        destructive_terms = (
            "delete row", "delete rows", "remove row", "remove rows",
            "hide row", "hide rows", "drop row", "drop rows",
            "only visible rows", "currently visible rows",
        )
        topology_terms = (
            "grid-template", "grid area", "grid-area", "track", "tracks",
            "stack", "regroup", "group", "move whole", "semantic",
            "two-zone", "band", "placement", "rebuild", "container",
            "column", "columns",
        )
        shared_support_terms = (
            "shared support band", "shared takeaway band", "shared note band",
            "support row below", "takeaway row below", "notes row below",
            "detach support", "detach takeaway", "move each card's",
            "move each card’s",
        )
        continuation_terms = (
            "continuation", "continued table", "split table",
            "split the table", "table segment", "table segments",
        )
        horizontal_reflow_terms = (
            "side-by-side", "side by side", "horizontalize", "horizontal layout",
            "horizontal grid", "two-column", "two column", "multi-column",
            "multiple columns", "split into columns", "grid columns",
            "column grid", "across the card", "across each card",
        )
        text_heavy_role_terms = (
            "metric", "support", "finding", "takeaway", "explanation",
            "interpretation", "note", "caption", "copy", "text block",
        )
        table_representation_terms = (
            "table representation", "representation problem",
            "readability itself", "table itself", "row/header",
            "header relationship", "every row", "all rows",
            "rendered text order", "reading order",
        )

        says_reflow = any(
            marker in low
            for marker in ("regional reflow", "body recompose", "recompose")
        )
        looks_like_compression = any(term in low for term in compression_terms)
        names_topology = any(term in low for term in topology_terms)
        if says_reflow and looks_like_compression and not names_topology:
            notes.append(
                f"Step {step_idx} is labeled reflow/recompose but mostly "
                "describes size/spacing compression. If that is the intended "
                "first move, relabel it local-fit/dashboard-fit; otherwise "
                "make the first edit change real region topology: tracks, "
                "grid areas, stacking, grouping, or whole-container placement."
            )
        if any(term in low for term in shared_support_terms):
            notes.append(
                f"Step {step_idx} detaches per-card support into a shared lower "
                "region. That region consumes body height while shortening every "
                "peer card, so do not infer extra capacity from the shorter cards "
                "alone. Compare the total peer-card-plus-support demand and the "
                "card-to-support reading relationship with reserving a terminal "
                "support zone inside each card before committing to this topology."
            )
        if (
            any(term in low for term in horizontal_reflow_terms)
            and any(term in low for term in text_heavy_role_terms)
        ):
            notes.append(
                f"Step {step_idx} places text-heavy metrics/support/takeaways "
                "side by side. Before committing, compare each item's actual "
                "available line width, resulting wrapping, and the region's total "
                "vertical demand with the current arrangement. A parent region "
                "spanning more card width does not mean each child becomes wider; "
                "multiple tracks may narrow every text item and increase height. "
                "Use the horizontal reflow when it improves the reading path and "
                "fit; otherwise preserve or revise the stack. This is diagnostic "
                "feedback, not a prohibition on columns."
            )
        if (
            says_reflow
            and any(term in low for term in ("table", "dashboard", "ranking", "summary"))
            and any(term in low for term in ("slide", "header", "footer", "content grid", "body grid", "whole"))
            and "same-topology" not in low
        ):
            notes.append(
                f"Step {step_idx} proposes whole-slide/body-grid reflow for a "
                "dashboard/table problem. Confirm from the render whether the "
                "existing topology and reading path are still coherent. If so, "
                "role-aware calibration is one lower-cost option; if not, the "
                "proposed reflow may be appropriate. This note does not impose "
                "a local-first ordering."
            )
        if "dashboard-fit" in low and any(
            marker in low
            for marker in ("body recompose", "regen", "delete", "remove rows")
        ):
            notes.append(
                f"Step {step_idx} is labeled dashboard-fit but mentions a "
                "larger/destructive strategy. Keep dashboard-fit local and "
                "reversible, or relabel the step before editing."
            )
        if (
            "dashboard-fit" in low
            and re.search(r"\btable\b", low)
            and not any(
                term in low
                for term in (
                    "notes", "right rail", "right-rail", "ranking", "summary",
                    "hero", "kpi", "sibling", "cluster", "card rhythm",
                )
            )
        ):
            notes.append(
                f"Step {step_idx} describes dashboard-fit as a table-only repair. "
                "Check whether nearby notes/cards or a KPI/summary rail compete "
                "for the same body space. Keep the edit local when evidence shows "
                "an isolated table defect; otherwise carry the shared-pressure "
                "hypothesis into the next edit and verification."
            )
        if any(term in low for term in destructive_terms):
            notes.append(
                f"Step {step_idx} mentions removing or keeping only visible "
                "table rows. For B-family spatial repairs, clipped or hidden "
                "HTML/source rows are still protected content. Prefer giving "
                "the table a different region, recalibrating the existing "
                "table/card tracks, or moving support cards into real tracks."
            )
        if any(term in low for term in continuation_terms) and not any(
            term in low for term in table_representation_terms
        ):
            notes.append(
                f"Step {step_idx} proposes a continuation/split-table strategy. "
                "That is a table-representation change, not a generic stronger "
                "layout strategy after a short slot or failed fit pass. Use it "
                "only when the table representation/readability itself is the "
                "diagnosed problem and the plan preserves every row, header "
                "relationship, and meaning-dependent row sequence. Otherwise prefer "
                "same-topology track calibration, grid-area placement of whole "
                "semantic units, or moving the table/support regions into "
                "real visible tracks."
            )
        return notes

    def _tool_plan(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Process the macro repair plan submitted by the agent.

        The plan tool is called at the start to lay out an overall strategy
        before diving into individual edits.  Steps are stored on state and
        progress is tracked throughout the repair session.  Re-submitting a
        plan replaces the previous one (useful when the agent needs to revise
        strategy mid-repair).
        """
        plan = action.get("plan")
        if not plan or not isinstance(plan, dict):
            return (
                "Invalid plan format. Provide a plan object with "
                "'summary', 'steps', and optionally 'skip' and "
                "'checkpoint_strategy' fields."
            ), False

        steps = plan.get("steps", [])
        if not steps:
            return (
                "Plan must include at least one step. Each step should have "
                "'action', 'type' (content/structural), and 'risk' fields."
            ), False

        candidate_steps: list[PlanStep] = []
        for s in steps:
            if isinstance(s, dict):
                text = s.get("action", s.get("text", str(s)))
                expected = s.get("expected_outcome", "")
                criterion = s.get("verify_criterion", "")
            else:
                text = str(s)
                expected = ""
                criterion = ""
            candidate_steps.append(
                PlanStep(
                    text=text,
                    expected_outcome=expected,
                    verify_criterion=criterion,
                )
            )

        summary = plan.get("summary", "")
        dashboard_plan_notes = self._dashboard_coupled_plan_notes(
            candidate_steps, state, summary,
        )
        descendant_plan_note = self._dashboard_descendant_plan_note(
            candidate_steps,
            state,
            " ".join((str(summary or ""), str(action.get("reasoning", "") or ""))),
        )
        if descendant_plan_note:
            dashboard_plan_notes.append(descendant_plan_note)

        replacing = state.has_plan
        if replacing:
            logger.info("Agent slide %d: plan replaced", state.slide_id)

        state.has_plan = True
        # Record clean checkpoint at plan time
        state.clean_checkpoint_idx = len(state.checkpoints) - 1

        # Store steps as PlanStep objects
        state.plan_steps = candidate_steps
        # Auto-mark first step as in_progress
        if state.plan_steps:
            state.plan_steps[0].status = "in_progress"

        state.plan_summary = summary
        skip = plan.get("skip", [])
        checkpoint_strategy = plan.get("checkpoint_strategy", "")

        feedback_parts = []
        if replacing:
            feedback_parts.append(
                f"Plan replaced ({len(steps)} steps). Previous plan discarded."
            )
        else:
            feedback_parts.append(f"Plan accepted ({len(steps)} steps).")
        if summary:
            feedback_parts.append(f"Summary: {summary}")
        if skip:
            feedback_parts.append(
                f"Skipping: {', '.join(skip) if isinstance(skip, list) else skip}"
            )
        if checkpoint_strategy:
            feedback_parts.append(f"Checkpoint strategy: {checkpoint_strategy}")
        if dashboard_plan_notes:
            feedback_parts.append(
                "DASHBOARD STRATEGY NOTE:\n"
                + "\n".join(f"- {note}" for note in dashboard_plan_notes)
            )

        strategy_fit_notes: list[str] = []
        for step_idx, step in enumerate(state.plan_steps, 1):
            strategy_fit_notes.extend(
                self._strategy_fit_notes_for_step_text(step.text, step_idx)
            )
        if strategy_fit_notes:
            feedback_parts.append(
                "PLAN-FAMILY FIT NOTE:\n"
                + "\n".join(f"- {note}" for note in strategy_fit_notes)
            )

        feedback_parts.append(
            "Now execute the plan in causal order. Checkpoints follow shared spatial "
            "causes, not merely slide membership: one coherent edit may advance "
            "several steps only when they share an owning region or allocation. Keep "
            "an independent localized defect in its own verified checkpoint so a "
            "later broad reflow or rollback cannot erase that success. Call "
            "update_plan as evidence makes a checkpoint genuinely complete or changes "
            "the strategy."
        )
        # Append current progress
        feedback_parts.append(self._format_plan_progress(state))
        return "\n".join(p for p in feedback_parts if p), False

    # ── Plan progress helpers ──────────────────────────────────────

    def _format_plan_progress(self, state: AgentState) -> str:
        """Format the current plan as a compact checklist.

        Appended to every tool result so the agent always sees its plan
        status.  Costs ~15-25 tokens for a typical 4-5 step plan.
        """
        if not state.plan_steps:
            return ""
        icons = {
            "done": "✓",
            "in_progress": "→",
            "pending": " ",
            "skipped": "⊘",
        }
        lines = ["", "PLAN PROGRESS:"]
        for i, step in enumerate(state.plan_steps, 1):
            icon = icons.get(step.status, " ")
            line = f"  [{icon}] {i}. {step.text}"
            if step.expected_outcome and step.status in ("in_progress", "pending"):
                line += f"\n       expect: {step.expected_outcome}"
            if step.status == "skipped" and step.skip_reason:
                line += f"  (skipped: {step.skip_reason})"
            lines.append(line)
        done = sum(
            1 for s in state.plan_steps
            if s.status in ("done", "skipped")
        )
        in_prog = sum(
            1 for s in state.plan_steps if s.status == "in_progress"
        )
        total = len(state.plan_steps)
        summary = f"  {done}/{total} completed"
        if in_prog:
            summary += f", {in_prog} in progress"
        lines.append(summary)
        return "\n".join(lines)

    def _tool_update_plan(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Update the repair plan: mark steps done/skipped, add or modify steps.

        Expected action format:
        {
            "tool": "update_plan",
            "updates": [
                {"step": 1, "status": "done"},
                {"step": 3, "status": "skipped", "reason": "..."},
                {"add": "New step text"},
                {"add": {"action": "New step", "expected_outcome": "...", "verify_criterion": "..."}},
                {"step": 2, "text": "Revised step description"}
            ],
            "new_steps": [{"action": "Structured replacement step", "expected_outcome": "...", "verify_criterion": "..."}]
        }
        """
        if not state.plan_steps:
            return (
                "No plan submitted yet. Call the plan tool first to "
                "create a repair plan."
            ), False

        raw_updates = action.get("updates")
        new_steps = action.get("new_steps")
        if raw_updates is not None and not isinstance(raw_updates, list):
            return (
                "Invalid update_plan format. 'updates' must be a list."
            ), False
        if new_steps is not None and not isinstance(new_steps, list):
            return (
                "Invalid update_plan format. 'new_steps' must be a list."
            ), False
        updates = list(raw_updates or [])
        updates.extend({"add": step} for step in (new_steps or []))
        if not updates:
            return (
                "Invalid update_plan format. Provide an 'updates' list or "
                "a 'new_steps' list. "
                "Each entry: {\"step\": N, \"status\": \"done\"} or "
                "{\"add\": \"new step text\"}."
            ), False

        VALID_STATUSES = {"done", "skipped", "in_progress", "pending"}
        results = []
        marked_done_steps: list[int] = []
        marked_skipped_steps: list[int] = []
        strategy_text_changed = False
        has_replacement_update = any(
            isinstance(upd, dict)
            and (
                "add" in upd
                or (isinstance(upd.get("text"), str) and bool(upd.get("text", "").strip()))
            )
            for upd in updates
        )
        for upd in updates:
            if not isinstance(upd, dict):
                results.append(f"Skipped non-dict entry: {upd}")
                continue

            # Add new step
            if "add" in upd:
                added = upd["add"]
                if isinstance(added, dict):
                    text = str(added.get("action", added.get("text", ""))).strip()
                    expected = str(added.get("expected_outcome", "") or "")
                    criterion = str(added.get("verify_criterion", "") or "")
                else:
                    text = str(added).strip()
                    expected = ""
                    criterion = ""
                if not text:
                    results.append("Skipped empty replacement step.")
                    continue
                state.plan_steps.append(PlanStep(
                    text=text,
                    expected_outcome=expected,
                    verify_criterion=criterion,
                ))
                strategy_text_changed = True
                results.append(
                    f"Added step {len(state.plan_steps)}: {text}"
                )
                continue

            # Update existing step
            step_num = upd.get("step")
            if not isinstance(step_num, int):
                results.append(
                    f"Invalid step number: {step_num}. "
                    f"Use 1-{len(state.plan_steps)}."
                )
                continue
            idx = step_num - 1  # 1-indexed → 0-indexed
            if idx < 0 or idx >= len(state.plan_steps):
                results.append(
                    f"Step {step_num} out of range "
                    f"(1-{len(state.plan_steps)})."
                )
                continue

            ps = state.plan_steps[idx]
            changed = []
            # Status update
            new_status = upd.get("status")
            if new_status:
                if new_status not in VALID_STATUSES:
                    results.append(
                        f"Step {step_num}: invalid status '{new_status}'. "
                        f"Use: {', '.join(sorted(VALID_STATUSES))}."
                    )
                    continue
                if (
                    new_status == "skipped"
                    and not self._has_current_verify(state)
                    and getattr(state, "last_verify_stale_reason", "")
                ):
                    results.append(
                        f"Step {step_num}: skip not applied because the current "
                        "checkpoint has not been measured by verify_layout after "
                        f"{state.last_verify_stale_reason}. Call verify_layout "
                        "first, then either continue this target or add a "
                        "replacement step with a different repair family."
                    )
                    continue
                if new_status == "skipped" and not has_replacement_update:
                    residual_total = int(
                        getattr(state, "_last_verify_targeted_residual_total", 0) or 0
                    )
                    core_terms = (
                        "table", "body", "cluster", "reflow", "dashboard",
                        "calibration", "fit", "overflow", "overlap",
                        "out-of-bounds", "lower", "footer",
                    )
                    if residual_total and any(term in ps.text.lower() for term in core_terms):
                        results.append(
                            f"Step {step_num}: skip not applied because the latest "
                            f"verify_layout still has {residual_total} targeted "
                            "residual finding(s) and this update does not add or "
                            "revise a replacement step for the same target. Add a "
                            "replacement step with a different credible repair "
                            "family, or keep this step in progress."
                        )
                        continue
                ps.status = new_status
                changed.append(f"status={new_status}")
                if new_status == "done":
                    marked_done_steps.append(step_num)
                if new_status == "skipped":
                    ps.skip_reason = upd.get("reason", "")
                    marked_skipped_steps.append(step_num)
            # Text update
            new_text = upd.get("text")
            if new_text and isinstance(new_text, str):
                ps.text = new_text
                strategy_text_changed = True
                changed.append("text updated")

            if changed:
                results.append(
                    f"Step {step_num}: {', '.join(changed)}"
                )

        feedback = "\n".join(results) if results else "No changes applied."
        if strategy_text_changed:
            strategy_fit_notes: list[str] = []
            for step_idx, step in enumerate(state.plan_steps, 1):
                strategy_fit_notes.extend(
                    self._strategy_fit_notes_for_step_text(step.text, step_idx)
                )
            if strategy_fit_notes:
                feedback += (
                    "\n\nPLAN-FAMILY FIT NOTE:\n"
                    + "\n".join(f"- {note}" for note in strategy_fit_notes)
                )
            dashboard_plan_notes = self._dashboard_coupled_plan_notes(
                state.plan_steps, state, getattr(state, "plan_summary", ""),
            )
            if dashboard_plan_notes:
                feedback += (
                    "\n\nDASHBOARD PLAN COUPLING NOTE:\n"
                    + "\n".join(f"- {note}" for note in dashboard_plan_notes)
                )
        if marked_done_steps:
            verify_notes: list[str] = []
            has_current_verify = self._has_current_verify(state)
            if state.current_code != state.original_code and not has_current_verify:
                stale_reason = state.last_verify_stale_reason or "latest code change"
                verify_notes.append(
                    "there is no current verify_layout after the latest code "
                    f"change ({stale_reason})"
                )
            elif getattr(state, "_last_verify_text_signal", False):
                verify_notes.append(
                    "the last verify_layout reported a visible-text change advisory"
                )
            if has_current_verify:
                regression_total = getattr(
                    state, "_last_verify_spatial_regression_total", 0,
                )
                if regression_total:
                    verify_notes.append(
                        f"the last verify_layout reported {regression_total} new "
                        "deterministic regression(s) vs baseline"
                    )
                residual_total = getattr(
                    state, "_last_verify_targeted_residual_total", 0,
                )
                if residual_total:
                    verify_notes.append(
                        f"{residual_total} targeted residual finding(s) remain in "
                        "the last verify_layout measurement"
                    )
            if verify_notes:
                feedback += (
                    "\n\nVERIFY-AWARE PLAN NOTE: Step(s) "
                    + ", ".join(str(n) for n in marked_done_steps)
                    + " were marked done, but "
                    + "; ".join(verify_notes)
                    + ". This is not a hard gate. If those findings belong to "
                    "later steps, continue; otherwise do not treat count "
                    "reduction or an applied edit as completion. A step is done "
                    "only when its verify_criterion and the original issue "
                    "evidence are actually satisfied."
                )
                if self._looks_like_table_dashboard_pressure(state):
                    feedback += (
                        " For a subregion in a shared fixed-canvas pressure chain, "
                        "local improvement is provisional until the neighboring roles "
                        "and frame coexist coherently. Keep that subregion available "
                        "for later calibration instead of treating done as frozen "
                        "geometry."
                    )
        if marked_skipped_steps:
            verify_notes: list[str] = []
            has_current_verify = self._has_current_verify(state)
            if state.current_code != state.original_code and not has_current_verify:
                stale_reason = state.last_verify_stale_reason or "latest code change"
                verify_notes.append(
                    "the current code has not been measured by verify_layout "
                    f"after {stale_reason}"
                )
            if has_current_verify:
                residual_total = getattr(
                    state, "_last_verify_targeted_residual_total", 0,
                )
                regression_total = getattr(
                    state, "_last_verify_spatial_regression_total", 0,
                )
                if residual_total:
                    verify_notes.append(
                        f"{residual_total} targeted residual finding(s) remain"
                    )
                if regression_total:
                    verify_notes.append(
                        f"{regression_total} new deterministic regression(s) remain"
                    )
                if getattr(state, "_last_verify_text_signal", False):
                    verify_notes.append("a visible-text change advisory remains")
            if verify_notes:
                feedback += (
                    "\n\nSKIPPED CORE TARGET NOTE: Step(s) "
                    + ", ".join(str(n) for n in marked_skipped_steps)
                    + " were marked skipped while "
                    + "; ".join(verify_notes)
                    + ". Skipping means a step is no longer applicable, not "
                    "that the original body/table/card/figure cluster can be "
                    "abandoned after one failed strategy. If the same cluster "
                    "still appears in verify_layout, add a replacement step that "
                    "uses a different information-preserving family: same-topology "
                    "track calibration, grid-area reflow, regrouped support cards, "
                    "body recompose, or source-preserving media adaptation."
                )
        feedback += "\n" + self._format_plan_progress(state)
        return feedback, False

    def _tool_search_source(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Search source evidence chunks by keyword relevance.

        Allows the repair agent to dynamically query original paper content
        when it encounters claims, numbers, or facts it needs to verify.
        Capped at MAX_SEARCH_CALLS to avoid wasting tool budget.
        """
        if state.search_calls_used >= self.MAX_SEARCH_CALLS:
            return (
                f"Search budget exhausted ({self.MAX_SEARCH_CALLS} calls used). "
                "Use the source evidence already provided in the initial "
                "context to guide your repairs."
            ), False

        query = action.get("query", "").strip()
        if not query:
            return "Error: 'query' field is required for search_source.", False

        top_k = action.get("top_k", 5)
        if not isinstance(top_k, int) or top_k < 1:
            top_k = 5
        top_k = min(top_k, 8)  # hard cap on results

        from ..evaluators.eval_tools import search_source
        result = search_source(query, state.evidence, top_k=top_k)
        state.search_calls_used += 1

        remaining = self.MAX_SEARCH_CALLS - state.search_calls_used
        result += f"\n\n[Search calls remaining: {remaining}/{self.MAX_SEARCH_CALLS}]"
        return result, False

    def _tool_lookup_table(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Look up tables in source materials by keyword.

        Allows the repair agent to find specific tables (comparison tables,
        result tables, ablation tables) from the original paper to verify
        numbers and data in the slide.
        Shares the search call budget with search_source.
        """
        if state.search_calls_used >= self.MAX_SEARCH_CALLS:
            return (
                f"Search budget exhausted ({self.MAX_SEARCH_CALLS} calls used). "
                "Use the source evidence already provided in the initial "
                "context to guide your repairs."
            ), False

        query = action.get("query", "").strip()
        if not query:
            return "Error: 'query' field is required for lookup_table.", False

        from ..evaluators.eval_tools import lookup_table
        result = lookup_table(query, state.evidence)
        state.search_calls_used += 1

        remaining = self.MAX_SEARCH_CALLS - state.search_calls_used
        result += f"\n\n[Search calls remaining: {remaining}/{self.MAX_SEARCH_CALLS}]"
        return result, False

    @staticmethod
    def _generated_asset_dir(state: AgentState) -> Path:
        """Return the directory where repair-created image assets should live."""
        if state._run_dir:
            return Path(state._run_dir) / f"turn_{state._turn_index:02d}" / "generated_assets"
        return Path(state.case_dir) / "repair_assets"

    @staticmethod
    def _html_asset_base_dirs(
        case_dir: str,
        run_dir: str | None = None,
        turn_index: int | None = None,
    ) -> list[Path]:
        """Directories used by verify_layout to resolve local image assets."""
        dirs: list[Path] = []

        def add(path: Path) -> None:
            try:
                resolved = path.resolve()
            except OSError:
                return
            if resolved not in dirs:
                dirs.append(resolved)

        case_path = Path(case_dir)
        add(case_path)
        add(case_path / "images")
        add(case_path / "source_pack" / "figures")
        add(case_path / "source_pack" / "tables")
        add(case_path / "source_pack" / "screenshots")
        if run_dir:
            run_path = Path(run_dir)
            if turn_index is not None:
                add(run_path / f"turn_{turn_index:02d}" / "generated_assets")
            for asset_dir in sorted(run_path.glob("turn_*/generated_assets")):
                add(asset_dir)
        return dirs

    @staticmethod
    def _safe_asset_name(
        raw_name: str,
        fallback: str,
        allowed_exts: tuple[str, ...] = (".png", ".jpg", ".jpeg", ".webp"),
    ) -> str:
        """Normalize a user/model-supplied filename for generated assets."""
        safe_name = re.sub(r"[^A-Za-z0-9_.-]+", "_", raw_name.strip()).strip("._")
        if not safe_name:
            safe_name = fallback
        if not safe_name.lower().endswith(allowed_exts):
            safe_name += allowed_exts[0]
        return safe_name

    @staticmethod
    def _svg_local_tag(tag: str) -> str:
        from .svg_repair import svg_local_tag
        return svg_local_tag(tag)

    @classmethod
    def _svg_visible_texts(cls, svg: str) -> list[str]:
        from .svg_repair import svg_visible_texts
        return svg_visible_texts(svg)

    @staticmethod
    def _svg_text_identity(text: str) -> str:
        from .svg_repair import svg_text_identity
        return svg_text_identity(text)

    @staticmethod
    def _svg_canvas_signature(svg: str) -> tuple[str, str, str] | None:
        from .svg_repair import svg_canvas_signature
        return svg_canvas_signature(svg)

    def _target_svg_asset_paths(self, state: AgentState) -> list[Path]:
        """Resolve external SVG assets named by current svg_visual_defect issues."""
        try:
            from .html_spatial_state import _resolve_local_img_src
        except Exception:
            return []

        issue_hints = set()
        for issue in self._current_issues:
            if issue.issue_type != "svg_visual_defect":
                continue
            evidence = getattr(issue, "evidence", None)
            for value in (
                getattr(evidence, "description", ""),
                getattr(issue, "planned_fix", ""),
                *list(getattr(evidence, "object_refs", []) or []),
            ):
                if isinstance(value, str) and value.strip():
                    issue_hints.add(value.lower())

        asset_base_dirs = self._html_asset_base_dirs(
            state.case_dir,
            getattr(state, "_run_dir", None),
            getattr(state, "_turn_index", None),
        )
        paths: list[Path] = []
        for src in re.findall(
            r'<img\s[^>]*src=["\']([^"\']+\.svg(?:[#?][^"\']*)?)["\']',
            state.original_code,
            flags=re.IGNORECASE,
        ):
            resolved = _resolve_local_img_src(
                src,
                html_base_dir=Path(state.case_dir),
                asset_base_dirs=asset_base_dirs,
            )
            if not resolved:
                continue
            asset_tokens = (src.lower(), resolved.name.lower(), str(resolved).lower())
            if issue_hints and not any(
                token in hint or hint in token
                for hint in issue_hints
                for token in asset_tokens
                if token
            ):
                continue
            if resolved not in paths:
                paths.append(resolved)
        return paths

    def _validate_svg_asset_repair_fidelity(
        self,
        svg: str,
        state: AgentState,
    ) -> str | None:
        """Reject external-SVG repairs that redraw instead of preserving content."""
        issue_types = set(getattr(state, "issue_types", set()) or set())
        if issue_types != {"svg_visual_defect"}:
            return None
        target_paths = self._target_svg_asset_paths(state)
        if not target_paths:
            return None
        try:
            original_svg = target_paths[0].read_text(encoding="utf-8", errors="replace")
        except Exception:
            return None

        original_canvas = self._svg_canvas_signature(original_svg)
        new_canvas = self._svg_canvas_signature(svg)
        if original_canvas and new_canvas and original_canvas != new_canvas:
            return (
                "SVG repair changed the asset canvas/viewBox. Preserve the original "
                f"viewBox/width/height {original_canvas}; fix only local label geometry."
            )

        original_texts = self._svg_visible_texts(original_svg)
        new_texts = self._svg_visible_texts(svg)
        original_blob = "".join(
            self._svg_text_identity(text) for text in original_texts
        )
        new_blob = "".join(self._svg_text_identity(text) for text in new_texts)
        missing = [
            text for text in original_texts
            if self._svg_text_identity(text)
            and self._svg_text_identity(text) not in new_blob
        ]
        added = [
            text for text in new_texts
            if self._svg_text_identity(text)
            and self._svg_text_identity(text) not in original_blob
        ]
        if missing:
            shown = "; ".join(missing[:5])
            return (
                "SVG repair dropped original visible SVG text. Preserve every "
                f"non-target label/caption exactly; missing: {shown}"
            )
        if added:
            shown = "; ".join(added[:5])
            return (
                "SVG repair added new visible SVG text. For svg_visual_defect, "
                f"wrap or reposition existing labels only; added: {shown}"
            )
        return None

    @staticmethod
    def _existing_image_paths_for_state(state: AgentState, limit: int = 12) -> list[Path]:
        """Return useful local image paths for tool error messages."""
        paths: list[Path] = []

        def add(path: Path | None) -> None:
            if path is None:
                return
            try:
                resolved = path.expanduser().resolve()
            except OSError:
                return
            if resolved.exists() and resolved.is_file() and resolved not in paths:
                paths.append(resolved)

        for src in re.findall(
            r"<img\b[^>]*\bsrc\s*=\s*['\"]([^'\"]+)['\"]",
            state.current_code,
            flags=re.IGNORECASE,
        ):
            add(AgentRepair._resolve_image_src_path(src, state))

        case_dir = Path(state.case_dir).resolve()
        for base in (
            case_dir / "source_pack" / "figures",
            case_dir / "source_pack" / "tables",
            case_dir / "source_pack" / "screenshots",
            case_dir / "images",
        ):
            if not base.exists():
                continue
            for suffix in ("*.png", "*.jpg", "*.jpeg", "*.webp"):
                for path in sorted(base.glob(suffix)):
                    add(path)

        if state._run_dir:
            run_dir = Path(state._run_dir)
            for path in sorted(run_dir.glob("turn_*/generated_assets/*")):
                if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}:
                    add(path)

        return paths[:limit]

    @staticmethod
    def _svg_text_width_estimate(text: str, font_size: float) -> float:
        from .svg_repair import svg_text_width_estimate
        return svg_text_width_estimate(text, font_size)

    @staticmethod
    def _svg_attr_float(el: ET.Element, name: str, default: float | None = None) -> float | None:
        from .svg_repair import svg_attr_float
        return svg_attr_float(el, name, default)

    @staticmethod
    def _svg_tag(el: ET.Element) -> str:
        from .svg_repair import svg_tag
        return svg_tag(el)

    @classmethod
    def _svg_text_fit_warning(cls, svg: str) -> str | None:
        from .svg_repair import svg_text_fit_warning
        return svg_text_fit_warning(svg)

    def _tool_create_svg_asset(self, action: dict, state: AgentState) -> tuple[str, bool]:
        """Create a constrained SVG asset for presentation-adapted figure repairs."""
        svg = action.get("svg") or action.get("content") or ""
        if not isinstance(svg, str) or not svg.strip():
            return "create_svg_asset failed: provide an SVG string in `svg`.", False
        svg = svg.strip()
        if not re.match(r"^<svg\b", svg, re.IGNORECASE):
            return "create_svg_asset failed: SVG content must start with <svg>.", False
        if re.search(r"<\s*(?:script|foreignObject)\b|\son[a-z]+\s*=", svg, re.IGNORECASE):
            return (
                "create_svg_asset failed: SVG may not contain script, "
                "foreignObject, or event-handler attributes."
            ), False
        if "xmlns=" not in svg[:250].lower():
            svg = re.sub(
                r"^<svg\b",
                '<svg xmlns="http://www.w3.org/2000/svg"',
                svg,
                count=1,
                flags=re.IGNORECASE,
            )
        fit_warning = self._svg_text_fit_warning(svg)
        if fit_warning:
            return f"create_svg_asset failed: {fit_warning}", False
        fidelity_warning = self._validate_svg_asset_repair_fidelity(svg, state)
        if fidelity_warning:
            return f"create_svg_asset failed: {fidelity_warning}", False
        safe_name = self._safe_asset_name(
            str(action.get("output_name") or "slide_figure_summary.svg"),
            "slide_figure_summary.svg",
            allowed_exts=(".svg",),
        )
        out_dir = self._generated_asset_dir(state)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / safe_name
        try:
            out_path.write_text(svg, encoding="utf-8")
        except Exception as exc:
            return f"create_svg_asset failed while writing {out_path}: {exc}", False
        turn_local_ref = f"../generated_assets/{safe_name}" if state._run_dir else str(out_path)
        return (
            f"create_svg_asset ok: {out_path}. Use `{turn_local_ref}` as the "
            "target raw-figure <img src> so the slide HTML stays turn-local; "
            "keep object-fit:contain, then call verify_layout and render_preview."
        ), False

    @staticmethod
    def _raw_figure_dense_slot_edit_warning(state: AgentState, edit_diff: str) -> str | None:
        """Block B17 fixes that steal space from dense adjacent text columns."""
        if not (state.issue_types & {"raw_figure", "raw_table"}):
            return None
        if not AgentRepair._is_html_code(state.current_code):
            return None
        if len(_visible_text_tokens(state.current_code)) <= 240:
            return None
        if not edit_diff:
            return None

        guarded_block = re.compile(
            r"\.(?:left|right|figure-wrap)\s*\{[^}]*\b(?:width|height|flex|grid|padding|margin)\s*:",
            re.IGNORECASE | re.DOTALL,
        )
        if not guarded_block.search(edit_diff):
            return None
        return (
            "EDIT BLOCKED: this dense slide has a raw_figure/raw_table issue, "
            "but the edit changes the left/right column or figure container size. "
            "Do not fix B17 by widening the figure slot or narrowing adjacent text. "
            "If the planned_fix authorizes a changed figure source, prefer a real "
            "crop/recomposition of the original source image inside the existing media "
            "slot. For quantitative charts/plots, do not escape this layout constraint "
            "by replacing a clean original chart with a hand-drawn SVG approximation; "
            "regenerate only from exact data or preserve the original image and mark "
            "the remaining layout limitation honestly. "
            "Handle separate B13/B04 issues with small targeted edits after the asset "
            "passes verify_layout and render_preview."
        )

    @staticmethod
    def _resolve_image_src_path(src: str, state: AgentState) -> Path | None:
        """Resolve an HTML image src to an existing local path."""
        if not src:
            return None
        raw = src.strip()
        if raw.startswith("file://"):
            raw = raw[7:]
        if raw.startswith("data:") or re.match(r"https?://", raw):
            return None

        candidate = Path(raw).expanduser()
        candidates = []
        case_dir = Path(state.case_dir).resolve()
        if candidate.is_absolute():
            candidates.append(candidate)
            parts = candidate.parts
            if "source_pack" in parts:
                idx = parts.index("source_pack")
                candidates.append(case_dir / Path(*parts[idx:]))
            if "generated_assets" in parts and state._run_dir:
                idx = parts.index("generated_assets")
                for turn_dir in sorted(Path(state._run_dir).glob("turn_*")):
                    candidates.append(turn_dir / Path(*parts[idx:]))
        else:
            candidates.extend([
                Path.cwd() / candidate,
                case_dir / candidate,
                case_dir.parent / candidate,
                case_dir.parent.parent / candidate,
            ])
            parts = candidate.parts
            if "source_pack" in parts:
                idx = parts.index("source_pack")
                candidates.append(case_dir / Path(*parts[idx:]))
            if "generated_assets" in parts and state._run_dir:
                idx = parts.index("generated_assets")
                for turn_dir in sorted(Path(state._run_dir).glob("turn_*")):
                    candidates.append(turn_dir / Path(*parts[idx:]))
        if candidate.name:
            for base in (
                case_dir / "source_pack" / "figures",
                case_dir / "source_pack" / "tables",
                case_dir / "source_pack" / "screenshots",
                case_dir / "images",
            ):
                candidates.append(base / candidate.name)
            if state._run_dir:
                candidates.extend(Path(state._run_dir).glob(f"turn_*/generated_assets/{candidate.name}"))
        for path in candidates:
            try:
                resolved = path.resolve()
            except OSError:
                continue
            if resolved.exists() and resolved.is_file():
                return resolved
        return None

    @staticmethod
    def _parse_crop_bbox(action: dict, img_size: tuple[int, int]) -> tuple[int, int, int, int] | None:
        """Parse a crop bbox as pixel left/top/right/bottom, clamped to image."""
        width, height = img_size
        bbox_pct = action.get("bbox_pct")
        if bbox_pct is not None:
            if isinstance(bbox_pct, dict):
                if {"left", "top", "right", "bottom"} <= set(bbox_pct):
                    values = [
                        bbox_pct["left"], bbox_pct["top"],
                        bbox_pct["right"], bbox_pct["bottom"],
                    ]
                elif {"x", "y", "width", "height"} <= set(bbox_pct):
                    x = float(bbox_pct["x"])
                    y = float(bbox_pct["y"])
                    values = [x, y, x + float(bbox_pct["width"]), y + float(bbox_pct["height"])]
                else:
                    return None
            elif isinstance(bbox_pct, list | tuple) and len(bbox_pct) == 4:
                values = list(bbox_pct)
            else:
                return None
            nums = [float(value) for value in values]
            # Accept either fractions [0, 1] or percentages [0, 100].
            if any(abs(value) > 1 for value in nums):
                nums = [value / 100 for value in nums]
            left, top, right, bottom = nums
            bbox = [left * width, top * height, right * width, bottom * height]
        else:
            bbox = action.get("bbox_px", action.get("bbox"))
        if isinstance(bbox, dict):
            if {"left", "top", "right", "bottom"} <= set(bbox):
                left = float(bbox["left"])
                top = float(bbox["top"])
                right = float(bbox["right"])
                bottom = float(bbox["bottom"])
            elif {"x", "y", "width", "height"} <= set(bbox):
                left = float(bbox["x"])
                top = float(bbox["y"])
                right = left + float(bbox["width"])
                bottom = top + float(bbox["height"])
            else:
                return None
        elif isinstance(bbox, list | tuple) and len(bbox) == 4:
            left, top, right, bottom = (float(value) for value in bbox)
        else:
            return None

        left_i = max(0, min(width, int(round(left))))
        top_i = max(0, min(height, int(round(top))))
        right_i = max(0, min(width, int(round(right))))
        bottom_i = max(0, min(height, int(round(bottom))))
        if right_i <= left_i or bottom_i <= top_i:
            return None
        return left_i, top_i, right_i, bottom_i

    @classmethod
    def _parse_crop_bboxes(
        cls,
        action: dict,
        img_size: tuple[int, int],
    ) -> list[tuple[int, int, int, int]] | None:
        """Parse one or more crop boxes from compose_image_grid input."""
        raw_bboxes = None
        pct_mode = False
        if action.get("bboxes_pct") is not None:
            raw_bboxes = action.get("bboxes_pct")
            pct_mode = True
        elif action.get("bboxes_px") is not None:
            raw_bboxes = action.get("bboxes_px")
        elif action.get("bboxes") is not None:
            raw_bboxes = action.get("bboxes")
        elif action.get("panels") is not None:
            raw_bboxes = []
            for panel in action.get("panels") or []:
                if not isinstance(panel, dict):
                    return None
                if panel.get("bbox_pct") is not None:
                    raw_bboxes.append({"bbox_pct": panel.get("bbox_pct")})
                elif panel.get("bbox_px") is not None:
                    raw_bboxes.append(panel.get("bbox_px"))
                elif panel.get("bbox") is not None:
                    raw_bboxes.append(panel.get("bbox"))
                else:
                    return None
        if not isinstance(raw_bboxes, list | tuple) or not raw_bboxes:
            return None

        bboxes: list[tuple[int, int, int, int]] = []
        for raw_bbox in raw_bboxes:
            if isinstance(raw_bbox, dict) and "bbox_pct" in raw_bbox:
                bbox = cls._parse_crop_bbox({"bbox_pct": raw_bbox["bbox_pct"]}, img_size)
            elif pct_mode:
                bbox = cls._parse_crop_bbox({"bbox_pct": raw_bbox}, img_size)
            else:
                bbox = cls._parse_crop_bbox({"bbox_px": raw_bbox}, img_size)
            if bbox is None:
                return None
            bboxes.append(bbox)
        return bboxes

    def _tool_crop_image(self, action: dict, state: AgentState) -> tuple[str, bool]:
        """Create a real cropped image asset from a local source image."""
        src = str(action.get("src", "")).strip()
        source_path = self._resolve_image_src_path(src, state)
        if source_path is None:
            available = self._existing_image_paths_for_state(state)
            available_text = "\n".join(f"- {path}" for path in available) or "- none found"
            return (
                "crop_image failed: provide a local src path from an existing "
                "<img>, e.g. cases/.../source_pack/figures/fig.png. Remote or "
                "data URLs are not supported. Available local image candidates:\n"
                f"{available_text}",
                False,
            )

        try:
            from PIL import Image
        except Exception as exc:
            return f"crop_image failed: Pillow is unavailable ({exc}).", False

        try:
            with Image.open(source_path) as img:
                bbox = self._parse_crop_bbox(action, img.size)
                if bbox is None:
                    return (
                        "crop_image failed: bbox_px must be [left, top, right, "
                        "bottom] pixels, bbox_pct must be [left, top, right, "
                        "bottom] as fractions/percentages, or use {x,y,width,height}. "
                        f"Source image size is {img.size[0]}x{img.size[1]}px.",
                        False,
                    )
                cropped = img.crop(bbox)
                if cropped.mode not in {"RGB", "RGBA"}:
                    cropped = cropped.convert("RGBA")
        except Exception as exc:
            return f"crop_image failed while reading {source_path}: {exc}", False

        safe_name = self._safe_asset_name(
            str(action.get("output_name", "")),
            f"slide_{state.slide_id:02d}_figure_excerpt_{len(state.checkpoints)}.png",
        )

        out_dir = self._generated_asset_dir(state)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = (out_dir / safe_name).resolve()
        try:
            cropped.save(out_path)
        except Exception as exc:
            return f"crop_image failed while writing {out_path}: {exc}", False

        left, top, right, bottom = bbox
        return (
            f"Cropped image saved to {out_path} ({cropped.width}x{cropped.height}px) "
            f"from {source_path} bbox [{left}, {top}, {right}, {bottom}].\n"
            "Next: replace the target <img src> with this path via apply_edits, "
            "set the image CSS to object-fit: contain, then call verify_layout "
            "and render_preview.",
            False,
        )

    def _tool_compose_image_grid(self, action: dict, state: AgentState) -> tuple[str, bool]:
        """Create a real multi-crop recomposed image asset from one source image."""
        src = str(action.get("src", "")).strip()
        source_path = self._resolve_image_src_path(src, state)
        if source_path is None:
            available = self._existing_image_paths_for_state(state)
            available_text = "\n".join(f"- {path}" for path in available) or "- none found"
            return (
                "compose_image_grid failed: provide a local src path from an "
                "existing <img> or source_pack image. Available local image "
                f"candidates:\n{available_text}",
                False,
            )

        try:
            from PIL import Image, ImageColor, ImageDraw
        except Exception as exc:
            return f"compose_image_grid failed: Pillow is unavailable ({exc}).", False

        try:
            with Image.open(source_path) as img:
                source = img.convert("RGBA")
                bboxes = self._parse_crop_bboxes(action, source.size)
                if not bboxes:
                    return (
                        "compose_image_grid failed: provide bboxes_px or "
                        "bboxes_pct as a non-empty list of [left, top, right, "
                        "bottom] boxes, or panels with bbox_px/bbox_pct. "
                        f"Source image size is {source.width}x{source.height}px.",
                        False,
                    )
                crops = [source.crop(bbox) for bbox in bboxes]
        except Exception as exc:
            return f"compose_image_grid failed while reading {source_path}: {exc}", False

        raw_layout = str(action.get("layout", "grid")).strip().lower()
        n = len(crops)
        if raw_layout in {"vertical", "stack", "column"}:
            columns = 1
        elif raw_layout in {"horizontal", "row"}:
            columns = n
        else:
            try:
                columns = int(action.get("columns") or 0)
            except (TypeError, ValueError):
                columns = 0
            if columns <= 0:
                columns = max(1, min(n, int(n ** 0.5 + 0.999)))
        columns = max(1, min(columns, n))
        rows = (n + columns - 1) // columns

        def _positive_int(name: str, default: int, min_value: int = 0) -> int:
            try:
                value = int(action.get(name, default))
            except (TypeError, ValueError):
                value = default
            return max(min_value, value)

        padding = _positive_int("padding_px", 24, 0)
        gap = _positive_int("gap_px", 18, 0)
        target_width = _positive_int("target_width_px", 1400, 320)
        max_cell_width = max(1, (target_width - 2 * padding - gap * (columns - 1)) // columns)

        scaled: list[Image.Image] = []
        for crop in crops:
            scale = min(1.0, max_cell_width / max(1, crop.width))
            if scale < 1.0:
                new_size = (
                    max(1, int(round(crop.width * scale))),
                    max(1, int(round(crop.height * scale))),
                )
                crop = crop.resize(new_size, Image.Resampling.LANCZOS)
            scaled.append(crop)

        cell_w = max(crop.width for crop in scaled)
        row_heights = []
        for row_idx in range(rows):
            row = scaled[row_idx * columns:(row_idx + 1) * columns]
            row_heights.append(max(crop.height for crop in row))
        canvas_w = padding * 2 + columns * cell_w + gap * (columns - 1)
        canvas_h = padding * 2 + sum(row_heights) + gap * (rows - 1)

        background = str(action.get("background", "#ffffff"))
        border = str(action.get("border", "#d5dde5"))
        try:
            bg_rgba = ImageColor.getcolor(background, "RGBA")
            border_rgba = ImageColor.getcolor(border, "RGBA")
        except ValueError:
            bg_rgba = (255, 255, 255, 255)
            border_rgba = (213, 221, 229, 255)

        composed = Image.new("RGBA", (canvas_w, canvas_h), bg_rgba)
        draw = ImageDraw.Draw(composed)
        y = padding
        for row_idx in range(rows):
            row = scaled[row_idx * columns:(row_idx + 1) * columns]
            row_h = row_heights[row_idx]
            for col_idx, crop in enumerate(row):
                x = padding + col_idx * (cell_w + gap) + max(0, (cell_w - crop.width) // 2)
                yy = y + max(0, (row_h - crop.height) // 2)
                composed.alpha_composite(crop, (x, yy))
                draw.rectangle(
                    [x, yy, x + crop.width - 1, yy + crop.height - 1],
                    outline=border_rgba,
                    width=1,
                )
            y += row_h + gap

        safe_name = self._safe_asset_name(
            str(action.get("output_name", "")),
            f"slide_{state.slide_id:02d}_figure_grid_{len(state.checkpoints)}.png",
        )
        out_dir = self._generated_asset_dir(state)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = (out_dir / safe_name).resolve()
        try:
            composed.convert("RGB").save(out_path)
        except Exception as exc:
            return f"compose_image_grid failed while writing {out_path}: {exc}", False

        bbox_text = "; ".join(
            f"[{left}, {top}, {right}, {bottom}]"
            for left, top, right, bottom in bboxes
        )
        return (
            f"Composed image grid saved to {out_path} ({canvas_w}x{canvas_h}px) "
            f"from {source_path}; panels={n}, layout={columns}x{rows}, "
            f"bboxes={bbox_text}.\n"
            "Next: replace the target <img src> with this path via apply_edits, "
            "set the image CSS to object-fit: contain, keep the original media "
            "slot/alt/ARIA semantics, then call verify_layout and render_preview.",
            False,
        )

    def _tool_generate_chart(self, action: dict, state) -> tuple[str, bool]:
        """Generate a chart image from viz_data."""
        viz_data = action.get("viz_data", {})
        chart_type = viz_data.get("chart_type", "")
        if chart_type == "flowchart":
            if not viz_data.get("nodes"):
                return "Flowchart viz_data must include 'nodes' (list of {id, label, type}).", False
        elif not viz_data.get("categories") or not viz_data.get("series"):
            return "viz_data must include 'categories' (list) and 'series' (list of {name, values}).", False

        try:
            from app.modules.chart_generator import ChartGenerator
            from pathlib import Path
            chart_path = Path(state.case_dir) / f"repair_chart_s{state.slide_id}_{len(state.checkpoints)}.png"
            result = ChartGenerator().generate_chart(viz_data, chart_path)
            if not result:
                return "Chart generation failed. Check viz_data format.", False
            abs_path = str(chart_path.resolve())
            return (
                f"Chart saved to {abs_path} ({chart_path.stat().st_size} bytes).\n"
                f"Insert it into the slide with apply_edits:\n"
                f'<img src="{abs_path}" style="max-width: 100%; border-radius: 2px;">'
            ), False
        except Exception as e:
            return f"Chart generation error: {e}", False

    def _tool_regen_slide(
        self, action: dict, state: AgentState,
    ) -> tuple[str, bool]:
        """Full slide regeneration from blueprint + evidence.

        Re-runs the codegen pipeline from scratch, injecting current issues
        as detailed negative examples with judge descriptions and fix content.
        Enhanced acceptance checks spatial quality, word count, structural
        change for layout issues, and must-contain content.

        This is a high-cost tool (counts as 5 tool calls).
        Limited to 1 attempt per session.
        """
        MAX_REGEN_PER_SESSION = 2
        if state.regen_attempts >= MAX_REGEN_PER_SESSION:
            return (
                "🚫 Regen budget exhausted (max 2 per session). "
                "Use apply_edits for targeted CSS/text fixes instead."
            ), False
        state.regen_attempts += 1

        if not state.bp_slide:
            return (
                "Cannot regen: blueprint slide not available. "
                "Use incremental apply_edits instead."
            ), False

        compiler = state.codegen_compiler
        if not hasattr(compiler, "_generate_code_only"):
            return (
                "Cannot regen: codegen_compiler does not support "
                "_generate_code_only. Use incremental apply_edits instead."
            ), False

        try:
            from pathlib import Path

            # Determine available images
            image_dir = Path(state.case_dir) / "images"
            available_images = []
            if image_dir.exists():
                available_images = [
                    f.name for f in image_dir.iterdir()
                    if f.suffix.lower() in {".png", ".jpg", ".jpeg", ".svg"}
                ]

            # Build rich regen context from issues — not just type names
            import copy
            bp_copy = copy.deepcopy(state.bp_slide)
            regen_note = "\n\n## CRITICAL CONSTRAINTS FOR REGENERATION\n"
            regen_note += "The previous version had these specific problems:\n"

            has_layout_issue = False
            has_density_issue = False
            must_contain_phrases: list[str] = []

            for issue in self._current_issues:
                if issue.status.value != "open":
                    continue
                desc = (
                    issue.evidence.description
                    or issue.why_this_fails
                    or ""
                )
                fix = issue.planned_fix or ""
                regen_note += f"- [{issue.issue_type}]: {desc[:200]}\n"
                if fix:
                    regen_note += f"  Fix: {fix[:200]}\n"
                # Inject fix_detail if available
                if issue.fix_detail and issue.fix_detail.correct_content:
                    content = normalize_correct_content_text(
                        issue.fix_detail.correct_content,
                    )[:300]
                    if content:
                        regen_note += f"  MUST INCLUDE (source-verified by judge): {content}\n"
                    # Only enforce must-contain for content accuracy issues;
                    # for density/layout issues, correct_content is often a
                    # repair instruction (e.g. "condense bullets") that won't
                    # appear verbatim in the new code.
                    if content and issue.issue_type in CONTENT_ACCURACY_ISSUE_TYPES:
                        must_contain_phrases.append(content)
                if issue.issue_type == "layout_inappropriate":
                    has_layout_issue = True
                if issue.issue_type in {
                    "density_imbalance", "whitespace_imbalance",
                }:
                    has_density_issue = True

            # Layout guidance — suggest a DIFFERENT template
            if has_layout_issue:
                regen_note += (
                    "\nThe layout pattern itself is wrong. Use a DIFFERENT "
                    "layout structure than the original. Consider: two-column "
                    "comparison, metric cards, image+text split, or whatever "
                    "best fits the content.\n"
                )

            # Preserve prior fixes — content that was corrected in prior turns
            for issue in self._current_issues:
                if (
                    issue.fix_detail
                    and issue.fix_detail.correct_content
                    and issue.resolved_at_turn is not None
                ):
                    content = normalize_correct_content_text(
                        issue.fix_detail.correct_content,
                    )
                    if content:
                        must_contain_phrases.append(content[:200])

            if must_contain_phrases:
                regen_note += "\n## MUST PRESERVE (from previous fixes):\n"
                for phrase in must_contain_phrases:
                    regen_note += f"- {phrase}\n"

            # Dynamic word budget based on current code
            current_words = len(re.findall(
                r'\b\w+\b',
                self._extract_visible_text(state.current_code),
            ))

            # Layout regeneration does not authorize new content. Missing
            # information is diagnosed and repaired by C-family issues.
            word_budget = min(current_words + 20, 120)

            regen_note += (
                f"\nHard constraints (EXCEEDING ANY WILL CAUSE REJECTION):\n"
                f"- Body text ≤{word_budget} words (current code has {current_words} words)\n"
                f"- ≤6 bullet points, ≤10 <li> elements total\n"
                f"- No overlapping elements\n"
                f"- Preserve the existing visible claims and numeric values unless a listed content issue explicitly requires correction\n"
                f"- Do not introduce new numbers, result values, or named claims without using search_source first\n"
                f"- Source attribution at bottom-right\n"
            )

            # Layout-specific guidance for whitespace issues
            if has_density_issue:
                regen_note += (
                    "\n## LAYOUT REQUIREMENTS (density/whitespace fix):\n"
                    "- Content MUST be distributed across the FULL slide height\n"
                    "- Use `min-height: 100%` on the main container\n"
                    "- Use `flex: 1` and `justify-content: space-between` to spread content vertically\n"
                    "- Do NOT huddle all content in the top 40-50% — this is the problem being fixed\n"
                    "- If needed, add larger padding, bigger images, or more detailed bullet points\n"
                    "- Target: content should occupy at least 60% of the slide area\n"
                )

            # Append to blueprint notes
            bp_copy.notes = (bp_copy.notes or "") + regen_note

            new_code, info = compiler._generate_code_only(
                bp_copy,
                state.evidence,
                image_dir if image_dir.exists() else None,
                available_images,
            )

            if not new_code:
                return "Regen failed: codegen returned no code.", False

            # Enhanced acceptance check
            old_spatial = self._quick_spatial_check(state.current_code, state)
            new_spatial = self._quick_spatial_check(new_code, state)

            accept, reason = self._regen_acceptance_check(
                state.current_code, new_code, state,
                old_spatial, new_spatial,
                has_layout_issue, must_contain_phrases,
            )

            if accept:
                projected_text_loss = state.cumulative_words_lost
                allows_content_change = (
                    bool(state.issue_types & CONTENT_ACCURACY_ISSUE_TYPES)
                    or state.allow_visible_text_change
                )
                if self._is_html_code(new_code) and not allows_content_change:
                    dropped_high_value = _dropped_high_value_tokens(state.current_code, new_code, limit=10)
                    if dropped_high_value:
                        return (
                            "Regen rejected because it removed value-bearing content: "
                            + ", ".join(dropped_high_value)
                        ), False
                    ordinary_lost = _meaningful_visible_words_lost(state.current_code, new_code)
                    projected_text_loss += len(ordinary_lost)
                    if ordinary_lost and projected_text_loss > state.text_loss_budget:
                        return (
                            "Regen rejected because cumulative visible-text loss would reach "
                            f"{projected_text_loss}/{state.text_loss_budget}. Retry with all "
                            "information-bearing content preserved; later DOM reflow remains allowed."
                        ), False
                state.checkpoints.append(state.current_code)
                state.checkpoint_text_loss.append(state.cumulative_words_lost)
                state.current_code = new_code
                state.cumulative_words_lost = projected_text_loss
                self._invalidate_verify_after_code_change(state, "regen_slide changed code")
                old_n = old_spatial.get("total_issues", "?") if old_spatial else "?"
                new_n = new_spatial.get("total_issues", "?") if new_spatial else "?"
                return (
                    f"✅ Regen accepted ({reason})! Spatial: {old_n}→{new_n}. "
                    f"Code replaced. Call verify_layout to confirm, then submit."
                ), True
            else:
                return (
                    f"🚫 Regen REJECTED: {reason}. "
                    f"DO NOT retry regen_slide — use apply_edits for "
                    f"targeted CSS/text fixes instead."
                ), False

        except Exception as e:
            logger.warning("regen_slide failed for slide %d: %s", state.slide_id, e)
            return f"Regen failed: {e}", False

    def _regen_acceptance_check(
        self,
        old_code: str,
        new_code: str,
        state: AgentState,
        old_spatial: dict | None,
        new_spatial: dict | None,
        has_layout_issue: bool,
        must_contain_phrases: list[str],
    ) -> tuple[bool, str]:
        """Multi-gate acceptance check for regen output.

        Returns (accept, reason) tuple.

        The gates are intentionally lenient compared to incremental edits:
        regen replaces the entire slide, so a small spatial regression is
        acceptable if the overall layout is substantially different and the
        issue count remains manageable.
        """
        old_issues = old_spatial.get("total_issues", 999) if old_spatial else 999
        new_issues = new_spatial.get("total_issues", 999) if new_spatial else 999

        # Gate 1: spatial must not BADLY degrade — allow small regression
        # Regen produces a fresh layout that may have 1-2 new spatial issues
        # but could fix the root cause (e.g., a fundamentally broken layout).
        max_allowed = old_issues + 2
        if new_issues > max_allowed:
            return False, f"spatial regression too large ({old_issues}→{new_issues}, max {max_allowed})"

        # Gate 2: word count — only reject extreme blowup
        old_words = len(re.findall(r'\b\w+\b', self._extract_visible_text(old_code)))
        new_words = len(re.findall(r'\b\w+\b', self._extract_visible_text(new_code)))
        if new_words > old_words + 50:
            return False, f"word count increased too much ({old_words}→{new_words})"

        has_content_issue = any(
            i.issue_type in CONTENT_ACCURACY_ISSUE_TYPES
            for i in self._current_issues
            if hasattr(i, "status") and i.status.value == "open"
        )
        if not has_content_issue:
            downgrade_ok, downgrade_reason = validate_repair_not_visual_downgrade(
                old_code, new_code,
            )
            if not downgrade_ok:
                return False, f"layout-only regen visual downgrade: {downgrade_reason}"

        # Gate 2b: layout-only regeneration must not rewrite the slide's
        # claims. In rich63 traces, regen fixed spatial geometry but introduced
        # unverified numbers/phrasing, which then regressed fidelity. For pure
        # B-family/layout repairs, require the visible text and numeric claims
        # to stay close to the original; content corrections should use
        # search_source + targeted edits instead.
        if not has_content_issue:
            old_text = self._extract_visible_text(old_code)
            new_text = self._extract_visible_text(new_code)

            def _meaningful_tokens(text: str) -> set[str]:
                return {
                    token.lower()
                    for token in re.findall(r"\b[A-Za-z][A-Za-z0-9\-]{3,}\b", text)
                    if token.lower() not in {
                        "source", "page", "figure", "slide", "with", "from",
                        "that", "this", "into", "than", "then", "only",
                    }
                }

            old_tokens = _meaningful_tokens(old_text)
            new_tokens = _meaningful_tokens(new_text)
            if len(old_tokens) >= 6:
                retention = len(old_tokens & new_tokens) / max(1, len(old_tokens))
                if retention < 0.60:
                    return False, f"layout-only regen rewrote too much text (token retention {retention:.0%})"

            number_re = r"(?<![\w.])\d+(?:\.\d+)?(?:[%×xX]|e[-+]?\d+)?"
            old_numbers = set(re.findall(number_re, old_text))
            new_numbers = set(re.findall(number_re, new_text))
            missing_numbers = old_numbers - new_numbers
            extra_numbers = new_numbers - old_numbers
            if missing_numbers:
                sample = ", ".join(sorted(missing_numbers)[:4])
                return False, f"layout-only regen dropped existing numeric claims: {sample}"
            if extra_numbers:
                sample = ", ".join(sorted(extra_numbers)[:4])
                return False, f"layout-only regen introduced unverified numeric claims: {sample}"

        # Gate 3: for layout_inappropriate, verify structural change via CSS fingerprint
        if has_layout_issue:
            def _layout_fingerprint(code: str) -> frozenset:
                props = re.findall(
                    r'(display|flex-direction|grid-template\S*|justify-content|align-items)\s*:\s*([^;}"]+)',
                    code,
                )
                return frozenset(props)

            if _layout_fingerprint(old_code) == _layout_fingerprint(new_code):
                return False, "CSS layout structure unchanged despite layout_inappropriate issue"

        # Gate 4: must-contain check for content accuracy
        for phrase in must_contain_phrases:
            # Check first 50 chars as key phrase
            key = phrase[:50].strip().lower()
            if key and key not in new_code.lower():
                return False, f"missing required content: {key[:30]}..."

        # Gate 5: density — prevent catastrophic bullet/word blowup
        bullet_count = len(re.findall(r'<li[^>]*>|<li>', new_code))
        if bullet_count > 12:
            return False, f"density too high ({bullet_count} bullets, max 12)"
        if new_words > 200:
            return False, f"absolute word count too high ({new_words}, max 200)"

        return True, "all gates passed"

    @staticmethod
    def _extract_visible_text(code: str) -> str:
        """Extract visible text from HTML code by stripping tags."""
        # Remove style and script blocks
        text = re.sub(r'<style[^>]*>.*?</style>', '', code, flags=re.DOTALL)
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
        # Remove HTML comments
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', ' ', text)
        # Collapse whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def _quick_spatial_check(
        self, code: str, state: AgentState,
    ) -> dict | None:
        """Run a quick compile + spatial extraction on code.

        Returns dict with total_issues count or None on failure.
        """
        try:
            compiler = state.codegen_compiler
            if self._is_html_code(code):
                from .html_spatial_state import extract_html_slide_state
                spatial = extract_html_slide_state(
                    state.slide_id,
                    code,
                    html_base_dir=Path(state.case_dir),
                    asset_base_dirs=self._html_asset_base_dirs(
                        state.case_dir,
                        getattr(state, "_run_dir", None),
                        getattr(state, "_turn_index", None),
                    ),
                )
            else:
                spatial = extract_slide_state(state.slide_id, code)

            if spatial is None:
                return None

            total = (
                len(getattr(spatial, "overflow_blocks", []))
                + len(getattr(spatial, "overlap_pairs", []))
                + len(getattr(spatial, "oob_blocks", []))
            )
            return {"total_issues": total}
        except Exception:
            return None

    def _tool_get_current_code(self, state: AgentState) -> tuple[str, bool]:
        """Return current code with line numbers."""
        state.last_code_read_revision = state.layout_revision
        lines = state.current_code.split("\n")
        numbered = "\n".join(
            f"{i+1:4d}: {line}" for i, line in enumerate(lines)
        )
        return f"```python\n{numbered}\n```", False

    def _measure_overflow_px(self, html_code: str, state: AgentState) -> float:
        """Return how many px content extends past 720px canvas."""
        if not self._is_html_code(html_code):
            return 0
        try:
            html_state = extract_html_slide_state(
                state.slide_id, html_code,
                html_base_dir=Path(state.case_dir),
            )
            return compute_overflow_px(html_state.blocks)
        except Exception:
            return 0

    def _tool_measure_space(self, state: AgentState) -> tuple[str, bool]:
        """Return vertical space budget analysis for the current slide."""
        if not self._is_html_code(state.current_code):
            return "measure_space requires HTML code.", False
        try:
            html_state = extract_html_slide_state(
                state.slide_id, state.current_code,
                html_base_dir=Path(state.case_dir),
            )
        except Exception as exc:
            return f"measure_space failed: {exc}", False

        canvas_h = 720
        # Find the max bottom edge of any content
        max_bottom = 0
        region_bottoms: dict[str, float] = {}
        for block in html_state.blocks:
            x, y, w, h = block.bbox_px
            bottom = y + h
            if bottom > max_bottom:
                max_bottom = bottom
            # Group by top-level region (first CSS class or var_name)
            region = block.css_selector.split()[0] if block.css_selector else block.var_name
            region = region.strip(".").split(":")[0][:20]
            if region not in region_bottoms or bottom > region_bottoms[region]:
                region_bottoms[region] = bottom

        overflow = max(0, max_bottom - canvas_h)
        # Find overflow:hidden containers
        import re as _re
        overflow_hidden_count = (
            state.current_code.count("overflow:hidden")
            + state.current_code.count("overflow: hidden")
        )

        lines = [
            f"SPACE BUDGET: canvas={canvas_h}px, content reaches {max_bottom:.0f}px, "
            f"overflow={overflow:.0f}px",
        ]
        if overflow > 0:
            lines.append(
                f"DEFICIT={overflow:.0f}px to save."
            )
        else:
            lines.append("Content fits within canvas — no overflow to fix.")

        if overflow_hidden_count:
            lines.append(
                f"Found {overflow_hidden_count} overflow:hidden declaration(s) — "
                f"releasing these may immediately resolve clipped content."
            )

        # Show top regions by bottom edge
        sorted_regions = sorted(
            region_bottoms.items(), key=lambda x: -x[1],
        )[:8]
        lines.append("Regions by bottom edge (highest pressure first):")
        for region, bottom in sorted_regions:
            over = bottom - canvas_h
            marker = f" ← {over:.0f}px past canvas" if over > 0 else ""
            lines.append(f"  .{region}: bottom={bottom:.0f}px{marker}")

        # CSS leverage analysis: which attributes can save the most space
        if overflow > 0:
            import re as _re2
            _style_m = _re2.search(r'<style>(.*?)</style>', state.current_code, _re2.DOTALL)
            if _style_m:
                _css = _style_m.group(1)
                _body = state.current_code[state.current_code.find('<body'):]
                _leverage = []
                # Check class selectors
                for _cls_m in _re2.finditer(r'\.([a-z][\w-]*)\s*\{([^}]*)\}', _css):
                    _cls_name = _cls_m.group(1)
                    _rules = _cls_m.group(2)
                    _count = max(1, len(_re2.findall(
                        rf'class="[^"]*\b{_cls_name}\b[^"]*"', _body)))
                    # Vertical spacing
                    _vert = 0
                    for _sp_m in _re2.finditer(
                        r'(padding|margin|gap)[^:]*:\s*([^;]+)', _rules):
                        _vals = _re2.findall(r'(\d+)px', _sp_m.group(2))
                        if _vals:
                            _vert = max(_vert, max(int(v) for v in _vals))
                    if _vert >= 8 and _count >= 1:
                        _leverage.append((f'.{_cls_name}', _vert, _count, _vert * _count))
                # Check element selectors (td, th, li, p)
                for _tag in ('td', 'th', 'li', 'tr'):
                    _tag_m = _re2.search(
                        rf'(?:^|\n)\s*{_tag}\s*\{{([^}}]+)\}}', _css)
                    if _tag_m:
                        _count = len(_re2.findall(rf'<{_tag}\b', _body))
                        _rules = _tag_m.group(1)
                        _vert = 0
                        for _sp_m in _re2.finditer(
                            r'(padding|margin|gap)[^:]*:\s*([^;]+)', _rules):
                            _vals = _re2.findall(r'(\d+)px', _sp_m.group(2))
                            if _vals:
                                _vert = max(_vert, max(int(v) for v in _vals))
                        if _vert >= 6 and _count >= 2:
                            _leverage.append((_tag, _vert, _count, _vert * _count))

                _leverage.sort(key=lambda x: -x[3])
                if _leverage:
                    lines.append(
                        f"\nHIGHEST LEVERAGE targets (spacing × element_count):"
                    )
                    for _sel, _val, _cnt, _pot in _leverage[:6]:
                        lines.append(
                            f"  {_sel}: spacing={_val}px × {_cnt} elements"
                            f" = {_pot}px potential savings"
                        )
                    lines.append(
                        "Compress these first (to 4-8px for td/th, "
                        "6-12px for containers). "
                        "Total savings should equal ~overflow, not more."
                    )

                # For absolute-layout slides, also suggest coordinate adjustments
                _abs_count = (
                    _css.lower().count('position:absolute')
                    + _css.lower().count('position: absolute')
                )
                if _abs_count >= 8 and overflow > 0:
                    lines.append(
                        f"\nLayout: {_abs_count} absolute-positioned elements."
                    )

        return "\n".join(lines), False

    def _summary_justifies_residuals(self, state: AgentState) -> bool:
        """True if the agent's repair summary genuinely flags a residual as
        structural / unfixable — the escape hatch for the residual re-bounce.

        We do NOT want the bounded re-bounce to trap a slide whose remaining
        clips really have no room (the agent made the right call to stop). A
        genuine justification must do more than list remaining defects. It must
        explicitly explain why no credible stronger strategy remains for the
        named residuals. A bare ``unresolved_concerns`` list, or a statement
        like "improved but still unresolved", means the agent noticed the
        defects; it does not mean the loop should stop trying.
        """
        summary = getattr(state, "repair_summary", None)
        if not summary:
            return False
        concerns = summary.get("unresolved_concerns") or []
        if isinstance(concerns, str):
            concern_text = concerns
        elif isinstance(concerns, list):
            concern_text = " ".join(str(c) for c in concerns)
        else:
            concern_text = str(concerns or "")
        text = " ".join([
            str(summary.get("self_assessment", "")),
            " ".join(str(x) for x in (summary.get("actions_taken") or [])),
            concern_text,
        ]).lower()
        # Phrases that indicate the agent consciously judged a residual
        # unfixable after considering stronger strategies — not a generic
        # "all clean" claim and not merely "there are unresolved concerns".
        STRUCTURAL_MARKERS = (
            "no credible", "no safe", "no safer", "no viable",
            "safely achievable", "cannot fit without", "can't fit without",
            "cannot be fit without", "unavoidable", "unfixable",
            "cannot be fixed", "can't be fixed", "no space left",
            "zero room", "irreducible", "must clip",
            "all viable", "all credible", "stronger strategy failed",
        )
        DESTRUCTIVE_ESCAPE_MARKERS = (
            "would require deleting", "would require hiding",
            "require deleting", "require hiding", "hiding/deleting",
            "delete rows", "delete row", "deleting rows", "deleting row",
            "remove rows", "remove row", "drop rows", "drop row",
            "would require over-compressing", "over-compressing",
        )
        PRESERVING_STRATEGY_MARKERS = (
            "information-preserving", "preserve all", "preserving all",
            "real visible space", "same-topology", "track calibration",
            "calibrated tracks", "grid-area", "grid area", "stable dom",
            "regrouped support", "support cards", "body recompose",
            "regional reflow", "semantic regroup", "source-preserving",
        )
        ATTEMPT_MARKERS = (
            "considered", "tried", "attempted", "failed", "unsafe",
            "would be unsafe", "rolled back", "rollback", "text regression",
            "new clipping", "new deterministic regression",
        )
        has_structural_marker = any(m in text for m in STRUCTURAL_MARKERS)
        has_destructive_escape = any(m in text for m in DESTRUCTIVE_ESCAPE_MARKERS)
        if has_destructive_escape:
            has_preserving_strategy = any(
                m in text for m in PRESERVING_STRATEGY_MARKERS
            )
            has_attempt_context = any(m in text for m in ATTEMPT_MARKERS)
            return (
                has_structural_marker
                and has_preserving_strategy
                and has_attempt_context
            )
        return has_structural_marker

    @staticmethod
    def _composition_entry_text(entry) -> str:
        if isinstance(entry, dict):
            return " ".join(
                str(value) for value in entry.values()
                if value is not None
            )
        return str(entry or "")

    @staticmethod
    def _composition_entry_issue_id(entry) -> str:
        if isinstance(entry, dict):
            for key in ("issue_id", "id", "target_issue_id"):
                value = str(entry.get(key, "")).strip()
                if value:
                    return value
        return ""

    @staticmethod
    def _composition_entry_has_any_field(
        entry: dict, fields: tuple[str, ...],
    ) -> bool:
        return any(str(entry.get(field, "")).strip() for field in fields)

    @classmethod
    def _composition_entry_verdict(cls, entry) -> str:
        if isinstance(entry, dict):
            for key in (
                "verdict", "status", "result", "completion", "outcome",
                "remaining_uncertainty",
            ):
                value = str(entry.get(key, "")).strip()
                if value:
                    return value
        return ""

    @classmethod
    def _composition_entry_is_meaningful(cls, entry) -> bool:
        text = re.sub(r"\s+", " ", cls._composition_entry_text(entry)).strip()
        if isinstance(entry, dict):
            has_failure = cls._composition_entry_has_any_field(
                entry,
                ("original_failure", "failure", "original_issue"),
            )
            has_evidence = cls._composition_entry_has_any_field(
                entry,
                (
                    "current_spatial_evidence", "spatial_evidence",
                    "current_render_evidence", "render_evidence", "evidence",
                    "current_state",
                ),
            )
            has_verdict = cls._composition_entry_has_any_field(
                entry,
                ("verdict", "status", "result", "remaining_uncertainty"),
            )
            return has_failure and has_evidence and has_verdict
        return len(text) >= 80

    @classmethod
    def _composition_closure_entries_from_summary(cls, summary: dict) -> list:
        closure = summary.get("composition_closure") if isinstance(summary, dict) else None
        if isinstance(closure, list):
            return closure
        if isinstance(closure, dict):
            if any(
                key in closure
                for key in (
                    "issue_id", "original_failure", "current_spatial_evidence",
                    "verdict",
                )
            ):
                return [closure]
            return [
                {
                    "issue_id": key,
                    "current_spatial_evidence": value,
                    "verdict": value,
                }
                for key, value in closure.items()
            ]
        if str(closure or "").strip():
            return [str(closure)]
        return []

    @classmethod
    def _meaningful_composition_entries(cls, state: AgentState) -> list:
        summary = getattr(state, "repair_summary", None)
        if not summary:
            return []
        entries = cls._composition_closure_entries_from_summary(summary)
        return [entry for entry in entries if cls._composition_entry_is_meaningful(entry)]

    @staticmethod
    def _composition_summary_concerns(summary: dict) -> list[str]:
        raw = summary.get("unresolved_concerns", []) if isinstance(summary, dict) else []
        if isinstance(raw, str):
            raw_items = [raw]
        elif isinstance(raw, list):
            raw_items = raw
        else:
            raw_items = [raw]
        concerns: list[str] = []
        for item in raw_items:
            text = re.sub(r"\s+", " ", str(item or "")).strip()
            if not text or text.lower() in {"none", "n/a", "na", "[]"}:
                continue
            concerns.append(text)
        return concerns

    @classmethod
    def _summary_has_composition_closure(
        cls, state: AgentState, issues: list[Issue] | None = None,
    ) -> bool:
        """Return whether the agent recorded composition self-assessment.

        This is a traceability check, not the completion check. It accepts
        specific uncertainty so the repair trace stays honest; successful
        shipment is decided by _summary_has_resolved_composition_closure().
        """
        summary = getattr(state, "repair_summary", None)
        if not summary:
            return False
        required = cls._composition_closure_issues(issues or [])
        required_ids = {
            str(issue.issue_id).strip()
            for issue in required
            if str(getattr(issue, "issue_id", "")).strip()
        }
        meaningful_entries = cls._meaningful_composition_entries(state)
        if not meaningful_entries:
            return False
        if not required_ids:
            return True

        covered_ids = {
            cls._composition_entry_issue_id(entry)
            for entry in meaningful_entries
        }
        covered_ids.discard("")
        if required_ids.issubset(covered_ids):
            return True

        # Text-only fallbacks are accepted only if they explicitly name every
        # target issue id. This keeps older traces readable while still forcing
        # issue-level self-assessment.
        combined = " ".join(
            cls._composition_entry_text(entry) for entry in meaningful_entries
        )
        return all(issue_id in combined for issue_id in required_ids)

    @classmethod
    def _composition_closure_unresolved_reasons(
        cls, state: AgentState, issues: list[Issue] | None = None,
    ) -> list[str]:
        """Explain why a composition trace is not ready to ship.

        This deliberately avoids geometric thresholds. It only checks whether
        the agent's own structured assessment is internally compatible with a
        completed repair.
        """
        if not cls._needs_composition_closure(issues or []):
            return []
        if not cls._summary_has_composition_closure(state, issues):
            return ["composition self-assessment is missing or incomplete"]

        summary = getattr(state, "repair_summary", None) or {}
        concerns = cls._composition_summary_concerns(summary)
        reasons: list[str] = []
        if concerns:
            reasons.append(
                "unresolved_concerns is non-empty: " + "; ".join(concerns[:2])
            )

        confidence = str(summary.get("confidence", "")).strip().lower()
        if confidence in {"low", "uncertain", "unresolved", "partial", "weak"}:
            reasons.append(f"summary confidence is {confidence!r}")

        required = cls._composition_closure_issues(issues or [])
        required_ids = {
            str(issue.issue_id).strip()
            for issue in required
            if str(getattr(issue, "issue_id", "")).strip()
        }
        entries = cls._meaningful_composition_entries(state)
        entries_by_id = {
            cls._composition_entry_issue_id(entry): entry
            for entry in entries
            if cls._composition_entry_issue_id(entry)
        }
        combined_entries = " ".join(cls._composition_entry_text(e) for e in entries)

        unresolved_status = (
            "uncertain", "unresolved", "partial", "weak", "failed", "fail",
            "blocked", "not fixed", "not resolved", "not addressed",
        )
        resolved_status = (
            "pass", "resolved", "fixed", "done", "complete", "completed",
            "success", "succeeded", "addressed",
        )
        unresolved_patterns: tuple[tuple[str, str], ...] = (
            (r"\bcannot claim\b|\bcan't claim\b", "assessment says it cannot claim success"),
            (r"\bnot (?:a )?(?:definitive|confident|high-confidence) (?:composition |visual )?pass\b", "assessment disclaims a confident visual pass"),
            (r"\b(?:only|merely) moderate\b|\bmoderate rather than\b|\bimprovement is moderate\b", "assessment says the improvement is only moderate"),
            (r"\bweak improvement\b|\bweak repair\b", "assessment describes a weak repair"),
            (r"\bnot (?:fully|visually|actually|definitively|confidently )?(?:resolved|fixed|addressed|handled|filled)\b", "assessment says the issue is not fully resolved"),
            (r"\b(?:void|blank space|blank region|empty band)\s+(?:still\s+)?remains?\b", "assessment says the named void remains"),
            (r"\bstill\s+(?:shows?\s+)?(?:the\s+)?(?:same\s+)?(?:lower|upper|left|right|corner|body)?\s*(?:void|blank space|blank region|empty band)\b", "assessment says the same blank region remains"),
            (r"\bmay still read\b", "assessment says the result may still read as defective"),
            (r"\bessentially unchanged\b|\bsame void\b", "assessment says the core spatial pattern is unchanged"),
        )

        def entries_for_issue(issue_id: str) -> list:
            if issue_id and issue_id in entries_by_id:
                return [entries_by_id[issue_id]]
            if issue_id and issue_id in combined_entries:
                return entries
            return entries if not required_ids else []

        for issue in required:
            issue_id = str(getattr(issue, "issue_id", "")).strip()
            matched_entries = entries_for_issue(issue_id)
            if not matched_entries:
                reasons.append(f"{issue_id or issue.issue_type}: no matching closure entry")
                continue
            for entry in matched_entries:
                if not isinstance(entry, dict):
                    reasons.append(
                        f"{issue_id or 'composition issue'}: completion needs a structured closure entry"
                    )
                    continue
                verdict = cls._composition_entry_verdict(entry).lower()
                if any(marker in verdict for marker in unresolved_status):
                    reasons.append(
                        f"{issue_id or 'composition issue'} verdict is not resolved: {verdict!r}"
                    )
                elif not any(marker in verdict for marker in resolved_status):
                    reasons.append(
                        f"{issue_id or 'composition issue'} verdict is not pass/resolved: {verdict!r}"
                    )
                strategy = str(entry.get("chosen_strategy", "")).strip().lower()
                if strategy in {"uncertain", "unknown", "none", "n/a"}:
                    reasons.append(
                        f"{issue_id or 'composition issue'} strategy is unresolved: {strategy!r}"
                    )
                assessment_text = " ".join(
                    [
                        cls._composition_entry_text(entry),
                        str(summary.get("self_assessment", "")),
                        " ".join(concerns),
                    ]
                ).lower()
                assessment_text = re.sub(r"\s+", " ", assessment_text)
                for pattern, reason in unresolved_patterns:
                    if re.search(pattern, assessment_text):
                        reasons.append(f"{issue_id or 'composition issue'}: {reason}")
                        break

                issue_context_parts = [
                    str(getattr(issue, "issue_type", "")),
                    str(getattr(issue, "sub_type", "")),
                    str(getattr(issue, "planned_fix", "")),
                    str(getattr(issue, "why_this_fails", "")),
                ]
                evidence = getattr(issue, "evidence", None)
                if evidence is not None:
                    issue_context_parts.append(str(getattr(evidence, "description", "")))
                fix_detail = getattr(issue, "fix_detail", None)
                if fix_detail is not None:
                    issue_context_parts.extend([
                        str(getattr(fix_detail, "target_location", "")),
                        str(getattr(fix_detail, "correct_content", "")),
                    ])
                issue_context = re.sub(
                    r"\s+", " ", " ".join(issue_context_parts).lower(),
                )
                if (
                    str(getattr(issue, "issue_type", "")) in {"raw_figure", "raw_table"}
                    and re.search(
                        r"\b(chart|plot|axis|axes|legend|tick|curve|line chart|"
                        r"training dynamics|validation accuracy|response length)\b",
                        issue_context + " " + assessment_text,
                    )
                    and re.search(
                        r"\b(redraw|svg summary|summary asset|generated svg|hand[- ]drawn)\b",
                        strategy + " " + assessment_text,
                    )
                    and not re.search(
                        r"\b(exact (?:source )?data|viz_data|generate_chart|"
                        r"data extracted|source-preserving crop|crop/recomposition|"
                        r"cropped/recomposed|original chart remains|preserv(?:e|ed) "
                        r"the original chart)\b",
                        assessment_text,
                    )
                ):
                    reasons.append(
                        f"{issue_id or 'composition issue'}: quantitative chart redraw lacks evidence-fidelity support"
                    )

                if (
                    str(getattr(issue, "issue_type", "")) in {"raw_figure", "raw_table"}
                    and re.search(
                        r"\b(?:rendered\s+)?(?:chart|figure|image|media)?\s*slot\s+(?:is|now|renders|became)|"
                        r"\bouter\s+(?:bbox|box|slot)|\b(?:img|image)\s+bbox\b",
                        assessment_text,
                    )
                    and not re.search(
                        r"\b(rendered image content|content rect|image content rect|"
                        r"letterbox|letterboxing)\b",
                        assessment_text,
                    )
                ):
                    reasons.append(
                        f"{issue_id or 'composition issue'}: raw-figure closure cites outer slot/bbox without rendered content-rect evidence"
                    )

        return list(dict.fromkeys(reasons))

    @classmethod
    def _summary_has_resolved_composition_closure(
        cls, state: AgentState, issues: list[Issue] | None = None,
    ) -> bool:
        return not cls._composition_closure_unresolved_reasons(state, issues)

    @classmethod
    def _composition_closure_block_message(
        cls, state: AgentState, issues: list[Issue], *, for_submit: bool,
    ) -> str:
        labels = cls._composition_issue_labels(issues, limit=6)
        summary = getattr(state, "repair_summary", None) or {}
        closure = summary.get("composition_closure") if isinstance(summary, dict) else None
        if not summary or not closure:
            headline = "SUBMIT BLOCKED" if for_submit else "Repair summary recorded"
            body = (
                "composition self-assessment is missing. Call submit_repair_summary with "
                "one structured `composition_closure` entry per target issue."
            )
        else:
            headline = "SUBMIT BLOCKED" if for_submit else "Repair summary recorded"
            body = (
                "composition self-assessment is incomplete. Each entry must name "
                "the original failure, cite current LAYOUT ANCHOR / RELATION MAP / "
                "SPACE MAP or render evidence, and give a verdict or remaining "
                "uncertainty. Unresolved verdicts are allowed as honest trace, "
                "but they are not completed repairs."
            )
        issue_list = "\n".join(f"  - {label}" for label in labels)
        return (
            f"{headline}: {body}\n"
            "This is a traceability requirement. Use spatial/render evidence to "
            "explain your judgment; if the judgment is uncertain or unresolved, "
            "continue with a stronger credible strategy before submitting.\n"
            f"{issue_list}"
        )

    @classmethod
    def _composition_completion_block_message(
        cls, state: AgentState, issues: list[Issue], *, for_submit: bool,
    ) -> str:
        labels = cls._composition_issue_labels(issues, limit=6)
        reasons = cls._composition_closure_unresolved_reasons(state, issues)
        headline = "SUBMIT BLOCKED" if for_submit else "Repair summary recorded"
        issue_list = "\n".join(f"  - {label}" for label in labels)
        reason_list = "\n".join(f"  - {reason}" for reason in reasons[:5])
        return (
            f"{headline}: the composition self-assessment does not support "
            "calling this repair complete yet. Specific uncertainty is good "
            "traceability, but it is not a successful composition repair.\n"
            f"Why it is not ready:\n{reason_list}\n"
            "Continue with a stronger body-composition strategy when one is "
            "credible: reflow existing elements, regroup the body, recompose/crop "
            "the source asset, or use an exact-data redraw/summary asset only when "
            "that preserves source fidelity. "
            "Do not fill a void with footer/source/caption movement, decorative "
            "frames, or stretched empty containers. After editing, call "
            "verify_layout/render_preview as appropriate, then record a new "
            "composition_closure whose verdict and concerns match the actual "
            "result.\n"
            f"Targets:\n{issue_list}"
        )

    def _tool_submit_repair_summary(
        self, action: dict, state: AgentState, slide_id: int,
        run_dir: str | None, turn_index: int,
    ) -> tuple[str, bool]:
        """Store a repair summary generated by the agent itself.

        The agent calls this tool after completing repairs to provide a
        structured self-assessment.  The summary is saved to
        ``run_dir/turn_XX/repair_summaries/slide_XX.json`` and injected
        into the triage prompt at next eval turn.
        """
        summary_data = {
            "slide_id": slide_id,
            "issues_targeted": action.get("issues_targeted", []),
            "actions_taken": action.get("actions_taken", []),
            "self_assessment": action.get("self_assessment", ""),
            "composition_closure": action.get("composition_closure", []),
            "confidence": action.get("confidence", "medium"),
            "unresolved_concerns": action.get("unresolved_concerns", []),
        }

        # Save to disk
        if run_dir:
            summary_dir = Path(run_dir) / f"turn_{turn_index:02d}" / "repair_summaries"
            summary_dir.mkdir(parents=True, exist_ok=True)
            summary_path = summary_dir / f"slide_{slide_id:02d}.json"
            try:
                with open(summary_path, "w") as f:
                    json.dump(summary_data, f, indent=2, ensure_ascii=False)
                logger.info(
                    "Agent slide %d: saved repair summary to %s",
                    slide_id, summary_path,
                )
            except Exception as e:
                logger.warning(
                    "Agent slide %d: failed to save repair summary: %s",
                    slide_id, str(e)[:100],
                )

        # Store on state for immediate access
        state.repair_summary = summary_data

        if (
            self._needs_composition_closure(self._current_issues)
            and not self._summary_has_composition_closure(state, self._current_issues)
        ):
            return (
                self._composition_closure_block_message(
                    state, self._current_issues, for_submit=False,
                ),
                False,
            )

        if self._needs_composition_closure(self._current_issues):
            completion_reasons = self._composition_closure_unresolved_reasons(
                state, self._current_issues,
            )
            if completion_reasons:
                return (
                    self._composition_completion_block_message(
                        state, self._current_issues, for_submit=False,
                    ),
                    False,
                )

        if state.current_code != state.original_code and not self._has_current_verify(state):
            stale_reason = state.last_verify_stale_reason or "latest code change"
            return (
                "Repair summary recorded, but the current code has not been "
                "measured by verify_layout after "
                f"{stale_reason}. Do not infer completion from an older "
                "measurement. Call verify_layout on this checkpoint, compare "
                "the spatial evidence with the original issue, then update the "
                "summary if the evidence changed.",
                False,
            )

        residual_total = int(
            getattr(state, "_last_verify_targeted_residual_total", 0) or 0
        )
        residual_counts = getattr(
            state, "last_verify_targeted_residual_counts", {},
        ) or {}
        hard_target_residuals = {
            category: count
            for category, count in residual_counts.items()
            if category in COMPOSITION_TARGET_HARD_SPATIAL_CATEGORIES
            and count
        }
        # Allow submit when residual count is small and the repair made
        # substantial progress.  Don't block a 79→4 or 25→3 result.
        baseline_spatial = getattr(state, "_t0_compact_issues", 0) or 0
        total_residual = sum(hard_target_residuals.values())
        if baseline_spatial > 8 and total_residual <= max(4, baseline_spatial * 0.10):
            hard_target_residuals = {}  # accept: substantial improvement
        elif total_residual <= 3 and baseline_spatial >= total_residual * 3:
            hard_target_residuals = {}  # accept: small absolute residual with clear improvement
        elif total_residual <= 6 and baseline_spatial > 0:
            hard_target_residuals = {}  # accept: small absolute residual is normal
        if hard_target_residuals:
            details = ", ".join(
                f"{category}={count}"
                for category, count in sorted(hard_target_residuals.items())
            )
            return (
                "Repair summary recorded, but the latest verify_layout still "
                "finds protected content clipped, covered, overflowing, or "
                "outside the canvas in the issue's named region "
                f"({details}). This is an objective visibility failure, not an "
                "aesthetic warning that can be justified away in the summary. "
                "Continue the current structural repair or choose a different "
                "layout strategy, then verify this exact revision again.",
                False,
            )
        if residual_total > 0 and not self._summary_justifies_residuals(state):
            return (
                "Repair summary recorded, but the latest verify_layout still "
                f"shows {residual_total} target-category residual measurement(s), "
                "and the summary does not yet establish closure. Compare them with "
                "the original issue and current revision evidence. If they share the same "
                "spatial cause, continue or change direction based on the evidence; "
                "if they are genuinely outside scope, identify why. Preserve all "
                "information-bearing content and do not hide or delete it to make a "
                "measurement disappear.",
                False,
            )

        return (
            "✅ Repair summary recorded. You may now call submit; the summary "
            "is consistent with the current completion assessment.",
            False,
        )

    def _render_slide_to_base64(
        self, code: str, state: AgentState,
    ) -> str | None:
        """Render slide code to PNG and return base64 string.

        For HTML: uses Playwright to render HTML → PNG.
        For pptx: compiles to PPTX → LibreOffice → PNG.
        Returns None on failure.
        """
        import base64
        import tempfile

        if self._is_html_code(code):
            backend = None
            try:
                from ...render_backends.playwright_backend import PlaywrightRenderBackend
                backend = PlaywrightRenderBackend()
                with tempfile.TemporaryDirectory() as tmpdir:
                    png_path = Path(tmpdir) / "slide.png"
                    success = backend.render_html_to_png(code, str(png_path))
                    if success and png_path.exists():
                        with open(png_path, "rb") as f:
                            return base64.b64encode(f.read()).decode()
                return None
            except Exception as e:
                logger.debug("HTML render_preview error: %s", str(e)[:200])
                return None
            finally:
                if backend is not None:
                    backend.close()

        try:
            from pptx import Presentation
            from pptx.util import Emu

            # Create single-slide PPTX
            prs = Presentation()
            prs.slide_width = Emu(SlideDimensions.WIDTH_EMU)
            prs.slide_height = Emu(SlideDimensions.HEIGHT_EMU)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            case_path = Path(state.case_dir)
            image_dir = str(
                state.codegen_compiler._find_image_dir(case_path)
                or case_path / "images"
            )

            from ...backends.python_pptx import code_executor
            success, error = code_executor.execute_code(
                code, prs, slide, image_dir,
            )
            if not success:
                logger.debug(
                    "render_preview compile failed: %s", error[:200],
                )
                return None

            with tempfile.TemporaryDirectory() as tmpdir:
                pptx_path = Path(tmpdir) / "preview.pptx"
                prs.save(str(pptx_path))

                from ...render_backends.libreoffice_backend import (
                    LibreOfficePdfRenderBackend,
                )
                renderer = LibreOfficePdfRenderBackend()
                result = renderer.render_pptx_to_pngs(
                    str(pptx_path), tmpdir,
                )

                if not result.png_paths:
                    logger.debug("render_preview: no PNGs generated")
                    return None

                # Pick the first (and only) slide PNG
                png_path = result.png_paths[0]
                with open(png_path, "rb") as f:
                    return base64.b64encode(f.read()).decode()

        except Exception as e:
            logger.debug("render_preview error: %s", str(e)[:200])
            return None

    def _tool_render_preview(self, state: AgentState) -> tuple[list[dict] | str, bool]:
        """Return the current rendered slide as multimodal repair feedback."""
        encoded = self._render_slide_to_base64(state.current_code, state)
        if not encoded:
            return (
                "Current slide render failed. Use verify_layout to inspect "
                "compile and spatial results before continuing.",
                False,
            )
        state.latest_visual_checkpoint_code = state.current_code
        state.latest_visual_checkpoint_revision = getattr(
            state, "layout_revision", 0,
        )
        if (
            getattr(state, "last_verify_result", None) is not None
            and getattr(state, "last_verify_revision", -1)
            == getattr(state, "layout_revision", 0)
        ):
            state.latest_visual_checkpoint_hard_valid = bool(
                state.latest_safe_verified_code == state.current_code
                and state.latest_safe_verified_revision
                == getattr(state, "layout_revision", 0)
            )
            state.latest_visual_checkpoint_targeted_issues = getattr(
                state, "_last_verify_targeted_residual_total", None,
            )
        return ([
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/png;base64,{encoded}",
                },
            },
            {
                "type": "text",
                "text": (
                    "Above is the current slide after your edits. Inspect the "
                    "reported target and the surrounding region. For SVG work, "
                    "check source/target boundaries, final shaft direction, "
                    "marker scale, label clearance, clipping, and whether "
                    "non-target graph roles remain unchanged. For raw figure "
                    "or raw table work, check that the intended subject remains "
                    "complete, no panel/row/axis is cut off, labels used by the "
                    "slide are readable, and existing adjacent text or source-image "
                    "labels guide the viewer to the claimed finding. Do not add "
                    "new visible callout text for B-family raw figure/table work. "
                    "If any check fails, edit "
                    "again before completing the plan step."
                ),
            },
        ], False)

    @staticmethod
    def _extract_fill_colors(code: str) -> dict[str, tuple[int, int, int]]:
        from .contrast_utils import extract_fill_colors
        return extract_fill_colors(code)

    def _check_color_contrast(
        self,
        blocks: list,
        fill_colors: dict[str, tuple[int, int, int]],
        code: str,
    ) -> list[str]:
        from .contrast_utils import check_color_contrast
        return check_color_contrast(blocks, fill_colors, code)

    def _calculate_luminance(self, r: int, g: int, b: int) -> float:
        from .contrast_utils import calculate_luminance
        return calculate_luminance(r, g, b)

    def _calculate_contrast_ratio(self, lum1: float, lum2: float) -> float:
        from .contrast_utils import calculate_contrast_ratio
        return calculate_contrast_ratio(lum1, lum2)

    def _extract_text_rgb(self, code: str, var_name: str) -> tuple[int, int, int]:
        from .contrast_utils import extract_text_rgb
        return extract_text_rgb(code, var_name)

    def _extract_text_brightness(
        self, code: str, var_name: str,
    ) -> float:
        from .contrast_utils import extract_text_brightness
        return extract_text_brightness(code, var_name)

    def _detect_whitespace_gaps(self, state) -> list[str]:
        """Detect large vertical gaps between elements.

        Returns list of warnings for gaps > 1.5 inches.
        """
        if not state.blocks:
            return []

        warnings: list[str] = []
        sorted_blocks = sorted(state.blocks, key=lambda b: b.y)

        # Check gap between top of usable area and first element
        first_y = sorted_blocks[0].y if sorted_blocks else 0
        if first_y > 1.8:  # > 1.8" from top
            warnings.append(
                f"  {sorted_blocks[0].var_name}: starts at y={first_y:.2f}\", "
                f"{first_y - 0.25:.1f}\" gap from slide top"
            )

        # Check gaps between consecutive elements
        for i in range(len(sorted_blocks) - 1):
            curr = sorted_blocks[i]
            nxt = sorted_blocks[i + 1]
            curr_bottom = curr.y + curr.h
            gap = nxt.y - curr_bottom

            if gap > 1.5:
                warnings.append(
                    f"  Between {curr.var_name} (bottom={curr_bottom:.2f}\") "
                    f"and {nxt.var_name} (top={nxt.y:.2f}\"): "
                    f"{gap:.1f}\" gap"
                )

        # Check gap between last element and bottom of usable area
        if sorted_blocks:
            last = sorted_blocks[-1]
            last_bottom = last.y + last.h
            bottom_gap = 7.20 - last_bottom  # USABLE_BOTTOM
            if bottom_gap > 2.0:
                warnings.append(
                    f"  After {last.var_name} (bottom={last_bottom:.2f}\"): "
                    f"{bottom_gap:.1f}\" unused space at bottom"
                )

        return warnings

    def _aesthetic_quality_check(
        self,
        original_code: str,
        repaired_code: str,
        slide_id: int,
    ) -> tuple[bool, str]:
        """Final aesthetic quality gate — catch visual regressions.

        Returns (ok, reason). If not ok, the repair should be rejected
        or the agent should iterate more.

        Checks:
        1. New text overflow introduced
        2. Color contrast violations
        3. Large whitespace gaps introduced
        4. Content loss (significant text removal)
        5. Image/chart deletion (embedded pictures removed)
        """
        issues: list[str] = []

        try:
            baseline = extract_slide_state(slide_id, original_code)
            current = extract_slide_state(slide_id, repaired_code)
        except Exception:
            return True, ""  # Can't check — allow

        # 1a. Text overflow regression — new overflow blocks
        # Use count_significant_issues to filter noise (≤8px, containers).
        from .html_spatial_state import (
            count_significant_issues as _csi,
            stable_block_identity,
            stable_pair_identity,
        )
        t0_sig_overflow = {
            stable_block_identity(baseline, bid)
            for bid in _csi(baseline).get("text_overflow", [])
        }
        t1_overflow_by_identity = {
            stable_block_identity(current, bid): bid
            for bid in _csi(current).get("text_overflow", [])
        }
        t1_sig_overflow = set(t1_overflow_by_identity)
        new_overflows = t1_sig_overflow - t0_sig_overflow
        if len(t1_sig_overflow) < len(t0_sig_overflow):
            new_overflows = set()
        if new_overflows:
            for identity in new_overflows:
                bid = t1_overflow_by_identity[identity]
                for b in current.blocks:
                    if b.block_id == bid:
                        issues.append(
                            f"NEW TEXT OVERFLOW in {b.var_name}: "
                            f"{b.text_chars}ch in {b.w:.1f}\"×{b.h:.1f}\" "
                            f"(overflow: {b.overflow_bottom_px}px bottom)"
                        )

        # 1b-2. Overlap regression — new element overlaps
        # Use count_significant_issues to filter out container/SVG element
        # overlaps that are not real defects (same filter the submit gate
        # and verify_layout use). Without this, the aesthetic gate rejects
        # repairs where the agent correctly resolved a text↔text overlap
        # but a new rect↔label pair appeared from layout restructuring.
        t0_sig_overlaps = {
            stable_pair_identity(baseline, first, second)
            for first, second in _csi(baseline).get("overlap", [])
        }
        t1_overlap_by_identity = {
            stable_pair_identity(current, first, second): (first, second)
            for first, second in _csi(current).get("overlap", [])
        }
        t1_sig_overlaps = set(t1_overlap_by_identity)
        new_sig_overlaps = t1_sig_overlaps - t0_sig_overlaps
        if len(t1_sig_overlaps) < len(t0_sig_overlaps):
            new_sig_overlaps = set()
        if new_sig_overlaps:
            # Look up readable names from the block list
            def _blk_name(bid):
                for b in current.blocks:
                    if b.block_id == bid or b.var_name == bid:
                        return bid
                return bid
            for pair_identity in new_sig_overlaps:
                a, b = t1_overlap_by_identity[pair_identity]
                # Find area from overlap_pairs
                area = 0.0
                for oa, ob, oarea in current.overlap_pairs:
                    if (min(oa,ob), max(oa,ob)) == (min(a,b), max(a,b)):
                        area = oarea
                        break
                issues.append(
                    f"NEW OVERLAP: {_blk_name(a)} ↔ {_blk_name(b)} "
                    f"(area={area:.2f} sq in)"
                )

        # 1c. Out-of-bounds regression — elements pushed past slide edge
        # Use count_significant_issues (same SSOT as overlap check above)
        # to filter out ≤5px OOB that are noise, not real defects.
        t0_sig_oob = {
            stable_block_identity(baseline, bid)
            for bid in _csi(baseline).get("out_of_bounds", [])
        }
        t1_oob_by_identity = {
            stable_block_identity(current, bid): bid
            for bid in _csi(current).get("out_of_bounds", [])
        }
        t1_sig_oob = set(t1_oob_by_identity)
        new_oob = t1_sig_oob - t0_sig_oob
        if len(t1_sig_oob) < len(t0_sig_oob):
            new_oob = set()
        if new_oob:
            for identity in new_oob:
                bid = t1_oob_by_identity[identity]
                for b in current.blocks:
                    if b.block_id == bid:
                        bottom = b.y + b.h
                        right = b.x + b.w
                        detail = (
                            f"NEW OOB in {b.var_name}: "
                            f"y={b.y:.2f}\" h={b.h:.1f}\" "
                            f"bottom={bottom:.2f}\""
                        )
                        if bottom > 7.5:
                            detail += (
                                f" ({bottom-7.5:.2f}\" past slide bottom)"
                            )
                        issues.append(detail)

        # 2. Color contrast
        current_colors = self._extract_fill_colors(repaired_code)
        contrast_issues = self._check_color_contrast(
            current.blocks, current_colors, repaired_code,
        )
        if contrast_issues:
            issues.extend(contrast_issues)

        # 3. Whitespace gaps introduced
        baseline_gaps = self._detect_whitespace_gaps(baseline)
        current_gaps = self._detect_whitespace_gaps(current)
        new_gap_count = len(current_gaps) - len(baseline_gaps)
        if new_gap_count > 0:
            issues.append(
                f"{new_gap_count} new whitespace gap(s) introduced"
            )
            issues.extend(current_gaps)

        # 4. Significant content loss
        t0_chars = sum(b.text_chars for b in baseline.blocks)
        t1_chars = sum(b.text_chars for b in current.blocks)
        if t0_chars > 50 and t1_chars < t0_chars * 0.55:
            issues.append(
                f"CONTENT CHANGE: {t0_chars}→{t1_chars} chars "
                f"({t1_chars/t0_chars:.0%} retained)"
            )

        # 5. Image/chart deletion — prevent removing embedded pictures
        #    or add_chart() calls.  The agent sometimes deletes images
        #    it cannot modify, causing worse visual regression than
        #    keeping them.  (e.g. Mamba slide 8: 2 benchmark charts
        #    deleted, leaving only text)
        t0_pictures = len(re.findall(
            r'slide\.shapes\.add_picture\(', original_code,
        ))
        t1_pictures = len(re.findall(
            r'slide\.shapes\.add_picture\(', repaired_code,
        ))
        if t0_pictures > 0 and t1_pictures < t0_pictures:
            issues.append(
                f"IMAGE CHANGE: {t0_pictures}→{t1_pictures} embedded "
                f"images (add_picture calls)"
            )

        t0_charts = len(re.findall(
            r'slide\.shapes\.add_chart\(', original_code,
        ))
        t1_charts = len(re.findall(
            r'slide\.shapes\.add_chart\(', repaired_code,
        ))
        if t0_charts > 0 and t1_charts < t0_charts:
            issues.append(
                f"CHART CHANGE: {t0_charts}→{t1_charts} chart shapes "
                f"(add_chart calls)"
            )

        # 6. Fill color identity changes — detect when large elements
        #    have their fill color changed from dark to light or vice
        #    versa, which alters the visual identity of the slide.
        t0_fills = self._extract_fill_colors(original_code)
        t1_fills = self._extract_fill_colors(repaired_code)
        for var_name, t0_rgb in t0_fills.items():
            t1_rgb = t1_fills.get(var_name)
            if t1_rgb is None:
                continue
            # Check if it's a large element (>5 sq in) by looking
            # at the spatial state
            block = None
            for b in current.blocks:
                if b.var_name == var_name:
                    block = b
                    break
            if block is None or block.w * block.h < 5.0:
                continue
            # Compute brightness of each fill
            t0_bright = (t0_rgb[0] * 0.299 + t0_rgb[1] * 0.587
                         + t0_rgb[2] * 0.114) / 255.0
            t1_bright = (t1_rgb[0] * 0.299 + t1_rgb[1] * 0.587
                         + t1_rgb[2] * 0.114) / 255.0
            # Flag if brightness category changed (dark↔light)
            t0_dark = t0_bright < 0.4
            t1_dark = t1_bright < 0.4
            if t0_dark != t1_dark:
                direction = "dark→light" if t0_dark else "light→dark"
                issues.append(
                    f"FILL COLOR CHANGE on {var_name}: "
                    f"({t0_rgb[0]},{t0_rgb[1]},{t0_rgb[2]}) → "
                    f"({t1_rgb[0]},{t1_rgb[1]},{t1_rgb[2]}) "
                    f"({direction}, brightness {t0_bright:.2f}→"
                    f"{t1_bright:.2f})"
                )

        if issues:
            return False, "\n".join(issues)
        return True, ""

    # ================================================================
    # MESSAGE BUILDING
    # ================================================================

    @staticmethod
    def _load_previous_repair_failures(
        run_dir: str, turn_index: int, slide_id: int,
    ) -> str:
        """Load failed repair attempts from previous turn as bad-case context.

        Extracts submit bounces, verify failures, and auto-rollbacks from the
        previous turn's repair log and formats them as cautionary examples.
        """
        prev_turn = turn_index - 1
        log_path = Path(run_dir) / f"turn_{prev_turn:02d}" / "repair_logs" / f"slide_{slide_id:02d}_attempt_0.json"
        if not log_path.exists():
            return ""

        try:
            with open(log_path) as f:
                log_messages = json.load(f)
        except Exception:
            return ""

        # Extract actual failure strings emitted by the repair loop. Keep the
        # first diagnostic line for each event so the next turn receives useful
        # strategy memory instead of an empty generic section.
        failures = []
        markers = (
            "SUBMIT BLOCKED", "RENDERED TEXT REGRESSION",
            "VISIBLE TEXT CHANGE SIGNAL", "AUTO-ROLLBACK",
            "Regen REJECTED", "REGEN REJECTED", "CHECKPOINT INVALID",
            "AMBIGUOUS EDITS", "Compile error", "REGRESSION CHECK",
            "STOPPING:", "loop timeout", "timed out",
        )
        for msg in log_messages:
            content = msg.get("content", "")
            if not isinstance(content, str):
                continue
            lines = [line.strip() for line in content.splitlines() if line.strip()]
            for index, line in enumerate(lines):
                if any(marker.lower() in line.lower() for marker in markers):
                    detail = " ".join(lines[index:index + 3])[:360]
                    failures.append(detail)
                    break
            if "submit bounced" in content.lower() and not any(
                "submit bounced" in failure.lower() for failure in failures
            ):
                failures.append(content[:360])

        if not failures:
            return ""

        # Cap at 10 most relevant failures
        failures = list(dict.fromkeys(failures))[:10]

        ctx = "## Previous Repair Failures (Turn {prev})\n".format(prev=prev_turn)
        ctx += "The following repair attempts FAILED on this slide in the previous turn. "
        ctx += "Learn from these mistakes — do NOT repeat the same approaches:\n\n"
        for i, f in enumerate(failures, 1):
            ctx += f"{i}. {f}\n"
        ctx += "\nTry a fundamentally different strategy for these issues."

        return ctx

    def _resolve_source_ref_to_path(
        self, source_ref: str, evidence: "EvidenceState | None",
    ) -> str | None:
        """Map a source_ref like '[tbl_p2_t0]' to its on-disk image path.

        Looks up evidence.tables / evidence.figures by id (with brackets
        stripped). Returns absolute path if image_path is set and file
        exists; else None.
        """
        if not evidence or not source_ref:
            return None
        ref = source_ref.strip().strip("[]").strip()
        if not ref:
            return None
        from pathlib import Path
        candidates: list[str] = []
        # Tables
        try:
            for t in (evidence.tables or []):
                if t.table_id == ref and t.image_path:
                    candidates.append(t.image_path)
        except Exception:
            pass
        # Figures
        try:
            for f in (evidence.figures or []):
                if f.figure_id == ref and f.image_path:
                    candidates.append(f.image_path)
        except Exception:
            pass
        for p in candidates:
            try:
                if Path(p).exists():
                    return str(Path(p).resolve())
            except Exception:
                continue
        return None

    @staticmethod
    def _table_row_specs_for_issue(issue) -> tuple[str, ...]:
        fd = getattr(issue, "fix_detail", None)
        if not fd:
            return ()
        return extract_table_row_specs_from_correct_content(
            getattr(fd, "correct_content", ""),
        )

    @classmethod
    def _is_table_row_issue(cls, issue) -> bool:
        fd = getattr(issue, "fix_detail", None)
        action_type = (getattr(fd, "action_type", "") or "").lower()
        target = (getattr(fd, "target_location", "") or "").lower()
        text = cls._issue_cluster_text(issue)
        if action_type == "add_data_row":
            return True
        return bool(cls._table_row_specs_for_issue(issue)) and (
            "table" in target or "row" in target or "table" in text
        )

    @staticmethod
    def _issue_cluster_text(issue) -> str:
        evidence = getattr(issue, "evidence", None)
        fix_detail = getattr(issue, "fix_detail", None)
        return " ".join(
            part for part in (
                getattr(issue, "issue_type", ""),
                getattr(issue, "rubric_id", ""),
                getattr(evidence, "description", ""),
                getattr(issue, "why_this_fails", ""),
                getattr(issue, "planned_fix", ""),
                getattr(fix_detail, "target_location", ""),
                getattr(fix_detail, "correct_content", ""),
            )
            if isinstance(part, str) and part.strip()
        ).lower()

    @classmethod
    def _issue_cluster_label(cls, issue) -> str:
        text = cls._issue_cluster_text(issue)
        issue_type = getattr(issue, "issue_type", "")
        if issue_type == "svg_visual_defect" or any(
            token in text for token in (
                "svg", "diagram", "connector", "arrow", "node", "path",
                "flowchart",
            )
        ):
            return "svg/diagram region"
        if any(token in text for token in (
            "footer", "bottom", "source note", "source attribution",
            "citation", "takeaway",
        )):
            return "footer/bottom region"
        if any(token in text for token in (
            "header", "title", "subtitle", "top bar", "top edge",
        )):
            return "header/title region"
        if any(token in text for token in ("right column", "right panel", "right side")):
            return "right column/panel"
        if any(token in text for token in ("left column", "left panel", "left side")):
            return "left column/panel"
        if any(token in text for token in ("table", "row", "cell", "grid")):
            return "table/grid region"
        if any(token in text for token in ("image", "figure", "chart", "plot")):
            return "figure/chart region"
        if issue_type in {
            "text_overflow", "out_of_bounds", "overlap",
            "alignment_inconsistency", "low_contrast",
        }:
            return "shared layout conflict"
        return issue_type or "miscellaneous"

    @classmethod
    def _build_issue_cluster_brief(cls, issues: list[Issue]) -> str:
        if len(issues) < 2:
            return ""

        clusters: dict[str, list[Issue]] = {}
        for issue in issues:
            clusters.setdefault(cls._issue_cluster_label(issue), []).append(issue)

        lines = [
            "## Issue Cluster Brief",
            "Group findings only when they share an owning region or spatial cause. "
            "Make one coherent regional edit for that causal group, then verify it "
            "before moving to an unrelated group.",
            "Different clusters are separate checkpoints even when they occur on the "
            "same slide. Do not fold a localized SVG/media defect into a risky body "
            "reflow, and do not let rollback of one composition experiment erase an "
            "independent verified repair.",
            "When many B03/B04 hard defects point to the same lower/body/table/card area, infer a shared spatial-pressure problem. Do not plan a separate pixel nudge for each clipped block; choose a repair family for the cluster and make the first edit match that family.",
        ]
        for label, grouped in sorted(
            clusters.items(), key=lambda item: (-len(item[1]), item[0]),
        )[:6]:
            issue_types = sorted({getattr(issue, "issue_type", "") for issue in grouped})
            target_hints = []
            for issue in grouped[:3]:
                evidence = getattr(issue, "evidence", None)
                fix_detail = getattr(issue, "fix_detail", None)
                hint = (
                    getattr(fix_detail, "target_location", "")
                    or getattr(evidence, "description", "")
                    or getattr(issue, "planned_fix", "")
                )
                hint = re.sub(r"\s+", " ", str(hint)).strip()
                if hint:
                    target_hints.append(hint[:90])
            hint_text = "; ".join(target_hints) if target_hints else "no target hint"
            lines.append(
                f"- {label}: {len(grouped)} finding(s), types={', '.join(issue_types)}. "
                f"Target hints: {hint_text}"
            )
        return "\n".join(lines)

    def _build_initial_message(
        self,
        code: str,
        all_issues: list[Issue],
        spatial_info: str,
        evidence_text: str,
        must_contain: list[str],
        must_not: list[str],
        content_checklist: str,
        bp_slide: BlueprintSlide | None,
        viz_data: dict | None,
        adjacent_context: str,
        spatial_state=None,
        layout_plan: dict | None = None,
        t0_render_b64: str | None = None,
        evidence: "EvidenceState | None" = None,
        task_brief: str = "",
    ) -> str | list:
        """Build the initial user message with all context."""
        parts = []

        if self._enable_render_preview:
            parts.append(
                "## Verification Mode\n\n"
                "Image preview is available in this run. Use `verify_layout` after "
                "structural edits and call `render_preview` only when pixel-level "
                "inspection is needed for SVG/image internals or visual topology.\n"
            )
        else:
            parts.append(
                "## Verification Mode\n\n"
                "This run is intentionally text/spatial only: `render_preview` is "
                "disabled and no slide image is available to you. Do not call it. "
                "Use the issue evidence, current HTML/CSS, `LAYOUT ANCHOR`, "
                "`RELATION MAP`, `SPACE MAP`, clipping/overlap details, and baseline "
                "delta to diagnose and verify the current revision. Detector counts "
                "are evidence, not a requirement to reach zero.\n"
            )

        # 1. Content requirements (MUST/MUST NOT) — front-loaded
        if must_not:
            parts.append("## MUST NOT contain (remove or replace these):\n")
            parts.extend(must_not)
            parts.append("")
        if must_contain:
            parts.append("## SHOULD contain (verify wording against source with search_source):\n")
            parts.extend(must_contain)
            parts.append("")

        # 2. Content checklist
        if content_checklist:
            parts.append(content_checklist)
            parts.append("")

        cluster_brief = self._build_issue_cluster_brief(all_issues)
        if cluster_brief:
            parts.append(cluster_brief)
            parts.append("")

        issue_types = {
            getattr(issue, "issue_type", "")
            for issue in all_issues
            if getattr(issue, "issue_type", "")
        }
        is_dashboard_pressure = self._looks_like_table_dashboard_pressure_from(
            code, issue_types,
        )
        dashboard_guidance = (
            self._dashboard_coupled_cluster_guidance(
                code, preview_enabled=self._enable_render_preview,
            )
            if is_dashboard_pressure
            else ""
        )

        composition_guidance = self._build_composition_closure_guidance(all_issues)
        if composition_guidance:
            parts.append(composition_guidance)
            parts.append("")

        # 3. Issues to fix — with code-line annotations
        code_lines = code.split("\n")
        parts.append("## Issues to Fix\n")
        for i, issue in enumerate(all_issues, 1):
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or ""
            )
            fix = issue.planned_fix or ""

            # For content accuracy issues, format as explicit edit command
            if issue.issue_type in CONTENT_ACCURACY_ISSUE_TYPES and fix:
                fd = issue.fix_detail
                correct_content = (
                    normalize_correct_content_text(fd.correct_content)
                    if fd and fd.correct_content else ""
                )

                # Distinguish between INSERT (missing_*) and REPLACE (incorrect/fabricated)
                is_missing = issue.issue_type.startswith("missing")
                row_specs = self._table_row_specs_for_issue(issue)
                is_table_row_insert = is_missing and self._is_table_row_issue(issue) and row_specs
                placeholder_context = " ".join(
                    part for part in (
                        desc,
                        fix,
                        getattr(fd, "target_location", "") if fd else "",
                        getattr(fd, "action_type", "") if fd else "",
                    )
                    if isinstance(part, str)
                ).lower()
                replaces_placeholder = is_missing and (
                    "placeholder" in placeholder_context
                    or "add:" in placeholder_context
                    or "add-note" in placeholder_context
                    or getattr(fd, "action_type", "") == "replace_text"
                )

                if is_table_row_insert:
                    rows_text = "\n".join(f"   - {row[:220]}" for row in row_specs)
                    fix_label = (
                        "   🎯 MANDATORY TABLE ROW INSERT — Add the missing "
                        "source rows as real table structure, not prose:\n"
                        f"{rows_text}\n"
                        "   Use apply_edits to insert <tr>/<td> rows inside "
                        "the existing <table>/<tbody>. Split each pipe-separated "
                        "row into cells. Do NOT add a <p>, <div>, footer note, "
                        "or visible editorial instruction sentence. "
                        "If space is tight, reduce table font/padding or reflow "
                        "the table region; do not move the row data into the "
                        "top prose or footer. Original judge instruction: "
                    )
                elif correct_content and replaces_placeholder:
                    # Missing-content issues sometimes arrive as placeholder
                    # cleanup tasks. Treating these as pure insertions preserves
                    # the placeholder and creates duplicate claims.
                    fix_label = (
                        "   🎯 MANDATORY PLACEHOLDER REPLACE — Replace the existing "
                        "placeholder/add-note line with finalized content:\n"
                        f"   Final text: \"{correct_content[:300]}\"\n"
                        "   Remove only the placeholder marker/text (for example "
                        "`Add:` or quoted instruction text). Preserve every other "
                        "visible word on the slide. Do NOT add a second copy next "
                        "to the placeholder.\n"
                        "   Original judge instruction: "
                    )
                elif correct_content and is_missing:
                    # C-family missing issues: include source content without
                    # forcing append-only edits that overload dense slides.
                    fix_label = (
                        "   🎯 MANDATORY INCLUDE — Make this source-backed content visible:\n"
                        f"   Required content: \"{correct_content[:300]}\"\n"
                        "   Prefer replacing/merging the closest same-topic sentence, list item, or paragraph with a combined sentence that preserves existing source-backed facts. Append with insert_after only when there is clear space. Do NOT duplicate the same idea in multiple places. Preserve all existing numbers, model names, and claims unless the judge flags them as wrong.\n"
                        "   Fixed-format budget: do NOT put long required content into a title, page header, full-width bottom bar, footer, or source note; merge longer qualifiers into body/interpretation text instead.\n"
                        "   Original judge instruction: "
                    )
                elif correct_content:
                    # D/E-family: precise replacement of the WRONG phrase only
                    target_label = _source_target_label(
                        issue.issue_type, correct_content, limit=300,
                    )
                    split_context = " ".join(
                        part for part in (fix, desc) if isinstance(part, str)
                    )
                    split_instruction = ""
                    if re.search(
                        r"\b(separate|split|distinct\s+bullets?|two\s+bullets?|separate\s+bullets?|separate\s+captions?)\b",
                        split_context,
                        re.IGNORECASE,
                    ):
                        split_instruction = (
                            "   STRUCTURE REQUIREMENT: the judge asked for separated takeaways. "
                            "Use separate existing/new sibling elements such as two <li> items or distinct captions; "
                            "do NOT leave both claims inside one bullet, one sentence, or one semicolon-separated paragraph.\n"
                        )
                    fix_label = (
                        "   🎯 MANDATORY REPLACE — Find the WRONG phrase and replace its meaning:\n"
                        f"   {target_label}\n"
                        "   ⚠️ Change ONLY the specific wrong phrase described above.\n"
                        f"{split_instruction}"
                        "   Use the shortest faithful slide wording that fixes the issue; if the target container is dense, merge with or replace the closest same-topic bullet instead of adding a new long sentence.\n"
                        "   Fixed-format budget: if the wrong phrase is in a title/header/footer/bottom bar and the source-backed correction is longer, keep that fixed region concise and put any longer qualifier into the closest body/interpretation sentence. Never replace a title with a paragraph or move a long correction into a footer.\n"
                        "   Original judge instruction: "
                    )
                else:
                    fix_label = (
                        "   ✅ VERIFIED FIX (judge extracted this from source — "
                        "use directly without re-verification): "
                    )
            elif issue.issue_type in ("title_content_mismatch", "weak_closing") and fix:
                fix_label = (
                    "   💡 SUGGESTED FIX (adapt wording as appropriate, "
                    "you may rephrase for better fit): "
                )
            elif issue.issue_type in ("text_overflow", "overlap"):
                if is_dashboard_pressure:
                    # Check if the HTML has overflow:hidden that could be released
                    _overflow_hidden_count = (
                        code.count("overflow:hidden")
                        + code.count("overflow: hidden")
                        if code
                        else 0
                    )
                    _overflow_hint = ""
                    if _overflow_hidden_count >= 2:
                        # Estimate total vertical spacing in CSS
                        import re as _re2
                        _total_spacing = sum(
                            int(m.group(1))
                            for m in _re2.finditer(r'(?:padding|margin|gap)\s*(?:-\w+)?\s*:\s*(\d+)px', code or '')
                            if int(m.group(1)) >= 4
                        )
                        _overflow_hint = (
                            f" This slide has {_overflow_hidden_count} "
                            "overflow:hidden containers that act as boundaries. "
                            "Compress content (font/padding/gap) to fit inside "
                            "each container. Only release overflow:hidden on a "
                            "specific container if compression alone cannot "
                            "resolve its clipping."
                        )
                        if _total_spacing > 0:
                            # The scrollHeight info is already in the spatial report
                            # Just remind agent to calibrate compression to the overflow amount
                            _overflow_hint += (
                                f" Total vertical spacing in CSS: ~{_total_spacing}px. "
                                f"Calibrate your reductions to cover the full overflow "
                                f"amount shown in the spatial report above."
                            )
                    fix_label = (
                        "   🔧 DASHBOARD FIT DIRECTION — preserve every visible "
                        "string and diagnose whether the overflow is local or caused "
                        "by several regions sharing one body constraint."
                        + _overflow_hint
                        + " Choose the actual targets "
                        "and scale from the current render; no selector set or edit "
                        "sequence is required. Suggested fix: "
                    )
                else:
                    fix_label = (
                        "   🔧 FIX OVERFLOW — use the rendered cause to choose among "
                        "local sizing/spacing, a larger owning region, or surrounding "
                        "layout reflow while preserving every visible string. Do not "
                        "optimize toward a fixed font-size recipe: "
                    )
            elif issue.issue_type in ("raw_figure", "raw_table"):
                preview_instruction = (
                    "For generated SVG/PNG summaries, inspect internal labels, "
                    "annotations, cards, paths, and bounds in render_preview; revise "
                    "or mark uncertain if they collide, clip, or become miniature "
                    "academic panels. After any changed figure source or structural "
                    "visual edit, call verify_layout and render_preview. "
                    if self._enable_render_preview
                    else
                    "This run has no image preview. Do not call render_preview. Use "
                    "verify_layout, image-content bounds, SVG text/shape diagnostics, "
                    "and source fidelity checks, and limit claims to what those signals "
                    "can establish. After any changed figure source or structural "
                    "visual edit, call verify_layout. "
                )
                fix_label = (
                    "   RAW SOURCE-FIGURE ADAPTATION — treat the judge's planned_fix as "
                    "a starting hypothesis, then choose the repair family that the current "
                    "render and spatial evidence support. First decide whether local frame, "
                    "size, or placement changes can make the source figure useful without "
                    "crowding protected text. If not, prefer a real source crop/recomposition "
                    "before any replacement. For quantitative charts/plots, preserve the "
                    "original chart or regenerate only from exact source data; do not use a "
                    "hand-drawn SVG summary that approximates curves, axes, legends, tick "
                    "values, or measured relationships. If the chart is good evidence and the "
                    "remaining failure is lower/side whitespace or awkward body composition, "
                    "repair that as layout reflow or mark the B17 result uncertain rather than "
                    "downgrading the chart. Use source-grounded SVG summary assets only for "
                    "conceptual diagrams, intrinsically unreadable structures, or explicitly "
                    "authorized presentation summaries. Do not automatically replace a visually "
                    "acceptable figure just because some internal labels are small. CSS-only image-window crops "
                    "such as object-view-box, object-fit:cover/none, negative offsets, "
                    "clip-path, or overflow-hidden enlarged images are NOT an acceptable "
                    "final repair. Preserve the existing media slot, alt/ARIA semantics, "
                    "and all visible slide text; labels inside generated assets must be "
                    "source-grounded and presentation-readable. "
                    f"{preview_instruction}"
                    "Suggested fix: "
                )
            elif issue.issue_type == "svg_visual_defect":
                fix_label = (
                    "   🔧 FIX SVG ASSET — the rendered image slot can fit while "
                    "text inside the referenced SVG asset is still clipped or crowded. "
                    "Do NOT solve this by enlarging the outer slide image box. Recreate "
                    "or update the SVG internals: split long labels with <tspan>, widen "
                    "only the local label shape, or reduce only the local SVG label font "
                    "while keeping it readable. If the SVG is referenced through "
                    "<img src>, use `create_svg_asset`, replace only that target img src "
                    "with the returned path, keep the same media slot and alt/ARIA "
                    "semantics, then call verify_layout. Suggested fix: "
                )
            elif (
                getattr(issue.fix_detail, "action_type", "")
                == "compress_support_copy"
            ):
                fix_label = (
                    "   AUTHORIZED SUPPORT-COPY COMPRESSION — first repair the "
                    "layout geometry. If the named explanatory copy still prevents "
                    "a readable fixed-canvas composition, shorten only those "
                    "support sentences or list items while preserving every factual "
                    "distinction, number, metric, named entity, label, conclusion, "
                    "and source attribution. Do not rewrite titles, KPI values, "
                    "table data, chart labels, or unrelated text. Suggested fix: "
                )
            elif issue.issue_type in STRUCTURAL_ISSUE_TYPES:
                if self._is_html_code(code):
                    fix_label = (
                        "   Suggested fix (LAYOUT ONLY — change CSS position/"
                        "size/color values, do NOT rewrite text content): "
                    )
                else:
                    fix_label = (
                        "   Suggested fix (LAYOUT ONLY — change Inches/Pt/"
                        "RGBColor values, do NOT rewrite text strings): "
                    )
            elif issue.issue_type in CRITICAL_CONTENT_TYPES:
                fix_label = (
                    "   Suggested fix (TEXT CONTENT ONLY — rewrite the "
                    "flagged text to match source evidence; do NOT change "
                    "CSS layout/position/size): "
                )
            elif issue.issue_type == "typography_error":
                # Typography errors may be truncation (CSS fix) or corruption
                # (text rewrite). Allow both approaches.
                fix_label = (
                    "   Suggested fix (may need CSS resize OR text rewrite "
                    "— if text is truncated/corrupted, you MAY rewrite the "
                    "text string to restore its full content): "
                )
            else:
                fix_label = "   Suggested fix: "

            # For spatial/layout issues (B-family), strip content-deletion
            # prescriptions from planned_fix. The judge often prescribes "remove
            # rows X and Y" which causes visual regression (empty space). We keep
            # only the spatial target (e.g., "fit within container") and let the
            # repair agent decide HOW (CSS-first per prompt rules).
            if issue.rubric_id.startswith("B") and issue.issue_type in (
                "text_overflow", "overlap", "density_imbalance",
                "text_wall", "content_overflow",
            ):
                import re as _re
                # Remove sentences that prescribe content deletion
                fix = _re.sub(
                    r'(?i)(remove|delete|drop|cut|reduce)\s+(the\s+)?\d[\w\s,/-]*rows?',
                    '[CSS-first: shrink font/padding instead]', fix,
                )
                fix = _re.sub(
                    r'(?i)keep only[^.;]*[.;]',
                    '', fix,
                )
                fix = _re.sub(
                    r'(?i)shorten the (paragraph|takeaway|text|sentence)[^.;]*[.;]',
                    '', fix,
                )

            # For alignment_inconsistency, constrain edits to the peer relation
            # and metric named by the judge. The required distance is contextual;
            # a fixed minimum move or global equalization can damage hierarchy.
            if issue.issue_type == "alignment_inconsistency":
                fix = (
                    "[TARGETED ALIGNMENT: identify the logical peer group from "
                    "the issue evidence and RELATION MAP, then change only the "
                    "named anchor, edge, or gap by the smallest amount that "
                    "removes the visible near-miss. Do not force unrelated "
                    "elements to equal width, equal height, or shared alignment. "
                    "Do not stretch table rows/cards or move caption/source/footer "
                    "text to fake balance; if a local anchor fix harms rhythm or "
                    "hierarchy, choose a body regroup/reflow instead.] "
                    + fix
                )
                align_context = " ".join([
                    getattr(issue, "description", "") or "",
                    getattr(issue, "why_this_fails", "") or "",
                    fix or "",
                ])
                if re.search(
                    r"(?i)(column rhythm|bottom[- ]edge|bottom alignment|extends lower|terminate|baseline|same bottom)",
                    align_context,
                ):
                    fix = (
                        "[COLUMN RHYTHM: when one column is longer than a shorter "
                        "visual/caption column, do not fix the mismatch by making "
                        "already-dense body text more cramped. Prefer using empty "
                        "space in the shorter column, adjusting non-text visual/caption "
                        "height or placement, or making a small grid/gap change. Keep "
                        "body text line-height and font size at least as readable as the "
                        "current slide unless the issue explicitly names text overflow.] "
                        + fix
                    )

            # For low_contrast issues, prefer the smallest contrast-safe color
            # change. Text on a light/tinted background usually needs darker
            # text; white text on an accent row may instead need a darker
            # text-bearing accent fill. Do not force one direction globally.
            if issue.issue_type == "low_contrast":
                import re as _re
                # White text is only valid when the current deterministic
                # target is a filled/accent row that needs a darker fill.
                if _re.search(r'(?i)\b(to|use|set)\s+(#[Ff]{3,6}|white)\b', fix) and not _re.search(
                    r'(?i)(row|filled|fill|background|accent)', fix,
                ):
                    fix = (
                        'Use a contrast-safe color pairing: darken text on '
                        'light/tinted backgrounds, or darken only the existing '
                        'text-bearing accent fill when preserving white text on '
                        'a highlighted row. NEVER set white text on a light or '
                        'transparent background.'
                    )
                fix = (
                    '[LOW CONTRAST: use the DETERMINISTIC LOW-CONTRAST TARGETS '
                    'section below as the current source of truth; old issue '
                    'wording may name elements already fixed. Preserve all text.] '
                    + fix
                )

            is_svg_visual = issue.issue_type == "svg_visual_defect"
            desc_limit = 1200 if is_svg_visual else 300
            fix_limit = 1200 if is_svg_visual else 250
            issue_block = (
                f"{i}. [{issue.severity.value}] {issue.rubric_id} "
                f"{issue.issue_type} ({issue.issue_id})\n"
                f"   Description: {desc[:desc_limit]}\n"
                f"{fix_label}{fix[:fix_limit]}"
            )

            if issue.issue_type in ("density_imbalance", "layout_inappropriate"):
                issue_block += (
                    "\n   STRATEGY CHOICE: before editing, decide whether the judge's "
                    "target can be solved by local resize/reposition or needs a "
                    "larger reflow of existing elements. Use the current render, "
                    "LAYOUT ANCHOR, RELATION MAP, and SPACE MAP as evidence. If local "
                    "changes do not appear to address the diagnosed reading path, "
                    "imbalance, or unused region, choose a reflow pattern that better "
                    "preserves the slide's meaning. Treat header/footer/source/frame "
                    "elements as layout constraints, not filler, unless this issue "
                    "explicitly names them as targets."
                )

            if issue.issue_type in COMPOSITION_CLOSURE_ISSUE_TYPES:
                issue_block += (
                    "\n   COMPOSITION SELF-ASSESSMENT REQUIRED: after verify_layout, "
                    "do not mark this issue done only because hard defects are "
                    "clean. Compare the original evidence and planned_fix to the "
                    "current LAYOUT ANCHOR / RELATION MAP / SPACE MAP and render; "
                    "state your chosen strategy, what changed, and any remaining "
                    "uncertainty for the next judge pass. A mechanical change that "
                    "only equalizes a metric, stretches a container, or moves auxiliary "
                    "text into a void should be reported as weak/uncertain unless it "
                    "also improves the rendered reading path."
                )

            if is_svg_visual:
                issue_block += (
                    "\n   Graph-preservation requirement: Before editing, "
                    "identify sequential stages versus parallel peers from "
                    "the current render and labels. Preserve all non-target "
                    "node roles and edges; repair the reported path, marker, "
                    "endpoint, or grouping with the smallest coherent change."
                )
                fd = getattr(issue, "fix_detail", None)
                if fd and fd.correct_content:
                    issue_block += (
                        "\n   Visual success contract: "
                        f"{fd.correct_content[:600]}"
                    )
                if fd and fd.target_location:
                    issue_block += (
                        "\n   Target location: "
                        f"{fd.target_location[:400]}"
                    )

            # Surface REGEN recommendation so agent knows to use regen_slide
            if (
                hasattr(issue, "recommended_action")
                and issue.recommended_action
                and str(issue.recommended_action.value) == "REGEN"
                and bool((getattr(issue, "action_rationale", "") or "").strip())
            ):
                issue_block += (
                    "\n   RECOMMENDED: consider regen_slide because the judge "
                    f"provided this rationale: {issue.action_rationale[:300]}"
                )

            # Surface judge-provided fix_detail for C/D issues
            # For B-family (layout) issues, suppress correct_content since it
            # often contains content-deletion prescriptions ("keep only X rows")
            # that cause visual regression. Layout issues should be fixed via CSS.
            if (
                hasattr(issue, "fix_detail")
                and issue.fix_detail
                and issue.fix_detail.correct_content
                and not issue.rubric_id.startswith("B")
            ):
                fd = issue.fix_detail
                content_text = normalize_correct_content_text(
                    fd.correct_content,
                )[:400]
                # Check for signs of truncated/incomplete source
                if content_text:
                    row_specs = self._table_row_specs_for_issue(issue)
                    if self._is_table_row_issue(issue) and row_specs:
                        rows_text = "\n".join(
                            f"   - {row[:220]}" for row in row_specs
                        )
                        issue_block += (
                            "\n   📋 SOURCE-VERIFIED TABLE ROWS (insert as "
                            "<tr>/<td> cells, not as visible prose):\n"
                            f"{rows_text}"
                        )
                    else:
                        truncation_warning = ""
                        if content_text.rstrip().endswith((",", "...", "…")) or "(source truncated)" in content_text:
                            truncation_warning = (
                                "\n   ⚠️ WARNING: This content may be based on a truncated source caption. "
                                "Use search_source to verify the exact wording before inserting. "
                                "Do NOT use content that you cannot confirm in the source."
                            )
                        target_label = _source_target_label(
                            issue.issue_type, content_text, limit=400,
                        )
                        issue_block += (
                            "\n   📋 SOURCE-VERIFIED CONTENT (from judge's direct reading of the source): "
                            f"{target_label}"
                            "\n   Fixed-format budget: do not paste long source text wholesale into a title/header/footer/bottom bar or a cramped list item; keep fixed regions concise and place necessary qualifiers in body/interpretation text."
                            f"{truncation_warning}"
                        )
                if fd.source_ref:
                    issue_block += f"\n   Source ref: {fd.source_ref}"
                    # Resolve source_ref → concrete asset path on disk so
                    # the agent can paste it directly into <img src="...">.
                    # Handles refs like "[tbl_p2_t0]", "tbl_p2_t0", "fig_p8_fig2".
                    asset_path = self._resolve_source_ref_to_path(
                        fd.source_ref, evidence,
                    )
                    if asset_path:
                        # Heuristic: warn when figure is likely multi-panel/dense
                        # (heavy raster, low display area → easily flagged raw_figure).
                        is_busy_fig = False
                        try:
                            for f in (evidence.figures or []):
                                if f.figure_id == (fd.source_ref or "").strip("[] "):
                                    # multi-panel heuristic: large source width or
                                    # caption mentions panels/comparison
                                    cap = (f.caption or "").lower()
                                    if (
                                        (f.width and f.width >= 800)
                                        or any(k in cap for k in (
                                            "comparison", "panel", "(a)", "(b)",
                                            "four", "multi", "subplot",
                                        ))
                                    ):
                                        is_busy_fig = True
                                    break
                        except Exception:
                            pass
                        busy_warn = (
                            "\n      ⚠️ This figure looks dense/multi-panel. "
                            "Insert with caution — prefer a real cropped/recomposed "
                            "asset via `crop_image`, `compose_image_grid`, `generate_chart`, or `create_svg_asset`, then "
                            "display it with `object-fit:contain`; CSS-only crops "
                            "are not accepted as final B17 repairs."
                            if is_busy_fig else ""
                        )
                        issue_block += (
                            f"\n   🖼️ AVAILABLE ASSET: {asset_path}"
                            f"\n      Insert directly: <img src=\"{asset_path}\" "
                            f"style=\"max-width:100%; max-height:100%; object-fit:contain;\">"
                            f"\n      (This file already exists on disk — no need to search/generate.)"
                            f"{busy_warn}"
                        )
                if fd.target_location:
                    issue_block += f"\n   Target location: {fd.target_location}"
                if fd.action_type:
                    issue_block += f"\n   Action: {fd.action_type}"

            # Annotate with specific code lines for content issues
            # so the agent knows EXACTLY where to make edits
            annotations = self._annotate_issue_locations(
                issue, code, code_lines,
            )
            if annotations:
                issue_block += "\n" + annotations

            # Enrich spatial issues with px bounding boxes for HTML
            if (
                self._is_html_code(code)
                and issue.issue_type in {
                    "overlap", "text_overflow", "out_of_bounds",
                }
                and spatial_state
                and spatial_state.blocks
            ):
                spatial_bbox = self._format_spatial_issue_with_px(
                    issue, spatial_state,
                )
                if spatial_bbox:
                    issue_block += "\n" + spatial_bbox

            # Enrich B2/B8 issues with spatial distribution facts
            if (issue.issue_type in {
                "density_imbalance", "layout_inappropriate",
                "whitespace_imbalance",
            } and spatial_state and spatial_state.blocks):
                spatial_ctx = self._compute_spatial_context(
                    spatial_state,
                )
                if spatial_ctx:
                    issue_block += "\n" + spatial_ctx

                # Route density_imbalance by sub_type for correct repair direction
                if issue.issue_type == "density_imbalance":
                    sub_type = getattr(issue, "sub_type", "") or ""
                    if not sub_type:
                        # Fallback: infer from text
                        for field in [issue.planned_fix or "", (issue.evidence.description if issue.evidence else "") or ""]:
                            if "cramped" in field or "content_overflow" in field:
                                sub_type = "cramped_content"
                                break
                            elif "sparse" in field or "underutilized" in field:
                                sub_type = "sparse_content"
                                break
                            elif "undersized" in field or "element" in field.lower():
                                sub_type = "element_undersized"
                                break
                            elif "column" in field.lower() and "mismatch" in field.lower():
                                sub_type = "column_height_mismatch"
                                break
                            elif "uneven" in field:
                                sub_type = "element_undersized"
                                break

                    if sub_type == "cramped_content":
                        issue_block += (
                            "\n   🔧 FIX OVERFLOW — diagnose whether the cause is local "
                            "sizing/spacing, an undersized owning region, or the "
                            "surrounding layout. Choose the least damaging direction "
                            "from the render. Only condense text as a last resort and "
                            "keep all distinct information points."
                        )
                    elif sub_type == "sparse_content":
                        issue_block += (
                            "\n   📐 RECOMPOSE EXISTING CONTENT — preserve every visible "
                            "string. Improve scale, grouping, focal hierarchy, and placement "
                            "so whitespace supports the message instead of looking accidental. "
                            "If the void sits below or beside a focal visual and local scaling "
                            "would leave the same structure, choose a body reflow of existing "
                            "callouts/notes/figure/caption relationships. Do not add source "
                            "material; true missing content belongs to C-family."
                        )
                    elif sub_type == "element_undersized":
                        issue_block += (
                            "\n   📐 RESIZE OR REFLOW THE FOCAL ELEMENT — increase the "
                            "usable rendered content only when that genuinely improves "
                            "inspection. Preserve image aspect ratio and avoid creating "
                            "letterbox, clipping, or overlap. For a complete wide/shallow "
                            "embedded image, giving it more height may not fix the actual "
                            "content; prefer more suitable horizontal span or a body reflow "
                            "that moves existing interpretation/callout content into a support "
                            "region. Do NOT spread text/bullet lists with space-between or "
                            "space-evenly — that creates ugly gaps. Text stays grouped."
                        )
                    elif sub_type == "column_height_mismatch":
                        issue_block += (
                            "\n   📐 ALIGN COLUMNS — if the shorter column has images/tables, "
                            "increase their CSS height. If it has only text/bullets, do NOT "
                            "spread them with space-between/space-evenly — keep items grouped "
                            "with flex-start. A compact text cluster looks better than "
                            "scattered items with huge gaps."
                        )

                    # NOTE: No quantified coverage metric is shown here.
                    # The bbox-union "space coverage %" does not correlate with
                    # visual-density judgment (large contiguous empty regions).
                    # A fixed threshold causes premature stopping. The issue's own
                    # planned_fix already contains concrete, per-element layout
                    # targets — the agent should follow those instead.

            parts.append(issue_block)
        parts.append("")

        svg_source_context = self._build_external_svg_source_context(
            all_issues,
            spatial_state,
        )
        if svg_source_context:
            parts.append(svg_source_context)
            parts.append("")

        # 4. Concrete fix suggestions for layout issues — BEFORE code
        # so the agent sees actionable copy-pasteable edits early
        layout_types = {
            "density_imbalance", "text_visual_imbalance",
            "layout_inappropriate", "text_overflow",
            "competing_focal_points", "whitespace_imbalance",
            "element_too_small", "text_cramped",
        }
        has_layout = any(
            i.issue_type in layout_types for i in all_issues
        )
        if has_layout and spatial_state:
            is_html = self._is_html_code(code)

            # For HTML: the compact state already provides rich layout info
            # via format_html_compact_state (in spatial_info at step 5).
            # Skip the pptx-specific layout representations that parse
            # Inches/RGBColor/Pt from python-pptx code.
            if not is_html:
                # Redistribution analysis for density/distribution issues
                if self._enable_redistrib_guide:
                    redistrib = self._compute_redistribution_guide(
                        code, spatial_state, all_issues,
                    )
                    if redistrib:
                        parts.append(redistrib)
                        parts.append("")

                # Layout representation
                mode = self.LAYOUT_REPR_MODE
                if mode == "diagnostic":
                    diag = self._build_layout_diagnostic(
                        code, spatial_state, all_issues,
                    )
                    if diag:
                        parts.append(diag)
                elif mode == "elements_json":
                    elems = self._build_elements_json(code, spatial_state)
                    if elems:
                        parts.append(elems)
                else:
                    if "dominance" in mode or mode == "all":
                        layout_guide = self._build_layout_guide(
                            code, spatial_state, all_issues,
                        )
                        if layout_guide:
                            parts.append(layout_guide)
                    if "ascii" in mode or mode == "all":
                        ascii_grid = self._build_ascii_grid(spatial_state)
                        if ascii_grid:
                            parts.append(ascii_grid)
                    if "relations" in mode or mode == "all":
                        relations = self._build_pairwise_relations(
                            spatial_state, code,
                        )
                        if relations:
                            parts.append(relations)
                    if "strip" in mode or mode == "all":
                        strip = self._build_vertical_strip(spatial_state)
                        if strip:
                            parts.append(strip)

        # 5. Current spatial state
        parts.append(spatial_info)
        parts.append("")

        # 6. Current code (line-numbered)
        numbered = "\n".join(
            f"{i+1:4d}: {line}" for i, line in enumerate(code_lines)
        )
        is_html = self._is_html_code(code)
        code_lang = "html" if is_html else "python"
        parts.append(f"## Current Slide Code\n\n```{code_lang}\n{numbered}\n```\n")

        # 6b. Embedded image context note
        if is_html and "<img" in code:
            parts.append(
                "**NOTE**: This slide has embedded images (`<img>` tags). "
                "Content inside these images (labels, axes, text, chart data) "
                "cannot be edited — you can only resize, reposition, or "
                "replace the `<img>` element.\n"
                "\n"
                "**CRITICAL**: NEVER delete or replace `<img>` tags with text "
                "placeholders. If an image appears broken or fails to render, "
                "the issue is likely a path resolution problem — try fixing "
                "the `src` attribute path instead. Removing figures degrades "
                "the slide quality severely.\n"
            )
        elif not is_html and "add_picture" in code:
            parts.append(
                "**NOTE**: This slide has embedded bitmap images "
                "(add_picture). Content inside these images (labels, "
                "axes, text, chart data) cannot be edited by code "
                "changes — you can only resize or reposition the "
                "image frame.\n"
            )

        # 6c. Deck color palette reminder — extract from current code
        palette_note = self._extract_palette_note(code, is_html)
        if palette_note:
            parts.append(palette_note)

        # 7. Chart data (if available)
        if viz_data:
            if is_html:
                parts.append(
                    f"## Chart Data (for CSS bar chart)\n"
                    f"```json\n{json.dumps(viz_data, indent=2)}\n```\n"
                )
            else:
                parts.append(
                    f"## Chart Data (for add_chart)\n"
                    f"```json\n{json.dumps(viz_data, indent=2)}\n```\n"
                )
            # Compute suggested placement for the chart
            if spatial_state and spatial_state.blocks:
                placement = self._find_available_slot(spatial_state)
                if placement:
                    if is_html:
                        # Convert inches to px (1 inch ≈ 96px at 1280/13.333)
                        px_l = int(placement[0] * 96)
                        px_t = int(placement[1] * 96)
                        px_w = int(placement[2] * 96)
                        px_h = int(placement[3] * 96)
                        parts.append(
                            f"**Suggested chart placement**: "
                            f"left: {px_l}px; top: {px_t}px; "
                            f"width: {px_w}px; height: {px_h}px — "
                            f"adapt position and size to fit.\n"
                        )
                    else:
                        parts.append(
                            f"**Suggested chart placement**: "
                            f"Inches({placement[0]:.2f}), "
                            f"Inches({placement[1]:.2f}), "
                            f"Inches({placement[2]:.2f}), "
                            f"Inches({placement[3]:.2f}) — "
                            f"adapt position and size to fit the slide's content.\n"
                            f"**Important**: Do NOT delete existing embedded "
                            f"images to make room. Place the chart in available "
                            f"space or reduce another element's height.\n"
                        )

        # 8. Adjacent slide context
        if adjacent_context:
            parts.append(adjacent_context)

        # 9. Slide context
        if bp_slide:
            parts.append(
                f"## Slide Context\n"
                f"Role: {bp_slide.role}\n"
                f"Primary proposition: {bp_slide.primary_proposition}\n"
            )

        # 10. Source evidence — use full context to match evaluator
        if evidence_text:
            # evidence_text already has its own header and 48K budget;
            # pass through without additional truncation so the agent
            # sees the same material the evaluator does.
            parts.append(evidence_text)

        # 10a. Global task brief / instructions
        if task_brief:
            parts.append(f"\n## Global Instructions (from task brief)\n\n{task_brief}")

        # 10b. Layout pre-plan (computed target positions)
        if layout_plan:
            plan_lines = [
                "## Layout Redistribution Plan\n",
                "A target layout has been computed to fix the density/"
                "distribution issues. Treat it as a design target for the "
                "HTML/CSS reflow, not as a requirement to preserve the "
                "current skeleton:\n",
            ]
            for var_name, target in sorted(
                layout_plan.items(),
                key=lambda kv: kv[1].get("top", 0),
            ):
                left = target.get("left", "?")
                top = target.get("top", "?")
                height = target.get("height", "?")
                width = target.get("width")
                line = f"  {var_name:20s}: left={left}, top={top}, height={height}"
                if width:
                    line += f", width={width}"
                plan_lines.append(line)
            plan_lines.append(
                "\nUse apply_edits to update CSS grid/flex/positioning so the "
                "visual regions move toward this target while preserving the "
                "original visible text, image assets, and reading path.\n"
            )
            parts.append("\n".join(plan_lines))

        # 11. Instructions — with conflict-aware prioritization
        additive_types = {
            "missing_evidence", "missing_data_visualization",
            "missing_entity", "missing_point",
        }
        layout_types = {
            "text_overflow", "text_wall", "density_imbalance",
        }
        additive_issues = [i for i in all_issues if i.issue_type in additive_types]
        layout_issues = [i for i in all_issues if i.issue_type in layout_types]

        conflict_guidance = ""
        if additive_issues and layout_issues:
            conflict_guidance = (
                "\n⚠️ **CONFLICT WARNING**: This slide has BOTH additive issues "
                "(missing content) and layout issues (overflow/density). "
                "Priorities:\n"
                "1. Apply the specifically authorized C-family insertion/correction.\n"
                "2. Then solve the B-family layout using geometry and grouping while "
                "preserving the resulting visible text exactly.\n"
                "3. If the payload cannot fit, report the unresolved structural constraint; "
                "do not silently skip required content or delete existing content.\n"
            )

        # For persistent text_overflow, mandate condensation
        persistent_overflow_guidance = ""
        persistent_overflows = [
            i for i in all_issues
            if i.issue_type == "text_overflow"
            and i.planned_fix
            and "[PERSISTED]" in str(i.planned_fix)
        ]
        if persistent_overflows:
            persistent_overflow_guidance = (
                "\n🔴 **PERSISTENT OVERFLOW**: Some text_overflow issues have "
                "persisted from previous turns. CSS-only fixes (font-size, "
                "container resize) have already been tried and FAILED. Change the "
                "layout family while preserving visible text: regroup containers, change "
                "grid/flex tracks, or use a structurally different arrangement. Do NOT "
                "retry the same CSS nudges or condense the text.\n"
            )

        persistent_layout_guidance = ""
        persistent_layouts = [
            i for i in all_issues
            if i.issue_type == "layout_inappropriate"
            and (
                getattr(i, "persisted_turns", 0) >= 1
                or "[PERSISTED]" in str(i.planned_fix or "")
                or "[PARTIALLY_MITIGATED]" in str(i.planned_fix or "")
            )
        ]
        if persistent_layouts:
            persistent_layout_guidance = (
                "\n🔴 **PERSISTENT LAYOUT MISMATCH**: At least one "
                "layout_inappropriate issue survived a previous repair. Do not "
                "repeat a footer-only color/height tweak unless the current fix "
                "plan names only the footer. Follow the current planned_fix's "
                "named hierarchy targets: metric cards, hero rows, captions, "
                "right/left panels, figure/table regions, and top summary cards "
                "must be resized/regrouped when named. Make one coherent "
                "structural rebalancing edit, then verify.\n"
            )

        parts.append(
            "## Instructions\n"
            f"Fix all {len(all_issues)} issues listed above. "
            "When done, call submit.\n"
            "You MUST first submit a plan with the `plan` tool, listing "
            "how you will fix each issue step by step. Then execute your "
            "plan, calling `verify_layout` after each structural change.\n"
            "\n"
            "**IMPORTANT**: After completing content edits (text rewrites, "
            "adding/removing text, fixing fabricated claims), you MUST call "
            "`verify_layout` to check for spatial regressions (text_overflow, "
            "density_imbalance). Content changes often cause overflow or "
            "density issues. If verify_layout reveals new spatial problems, "
            "fix them (adjust CSS font-size, container dimensions, padding) "
            "BEFORE submitting.\n"
            f"{conflict_guidance}"
            f"{persistent_overflow_guidance}"
            f"{persistent_layout_guidance}"
        )
        if dashboard_guidance:
            # Keep the current-case decision guidance close to the final action
            # instruction. The issue/code/evidence packet can be very long, and
            # placing this only near the top made the model miss the coupled
            # dashboard trajectory and execute serial local edits instead.
            parts.append(dashboard_guidance)

        # 12. Space Planning requirement for layout issues
        # Forces the agent to reason about spatial budget before
        # making structural edits — prevents the "add chart without
        # shrinking other elements" failure pattern.
        if self._enable_space_planning:
            layout_issue_types = {
                "density_imbalance", "text_visual_imbalance",
                "layout_inappropriate", "competing_focal_points",
                "whitespace_imbalance", "missing_data_visualization",
                "element_too_small",
            }
            layout_issues = [
                i for i in all_issues
                if i.issue_type in layout_issue_types
            ]
            if layout_issues and spatial_state and spatial_state.blocks:
                # Build current element y-range map for reference
                sorted_blocks = sorted(
                    spatial_state.blocks,
                    key=lambda b: b.y,
                )
                element_map_lines = []
                for b in sorted_blocks:
                    bottom = b.y + b.h
                    element_map_lines.append(
                        f"  {b.var_name:20s}: y={b.y:.2f}→{bottom:.2f} "
                        f"(h={b.h:.2f}\") w={b.w:.2f}\""
                    )
                element_map = "\n".join(element_map_lines)

                parts.append(
                    "## REQUIRED: Space Planning Before Structural Edits\n\n"
                    "You have layout issues to fix. Before making ANY "
                    "structural change (resize, move, delete, add), you "
                    "MUST first reason through a space plan in your "
                    "`reasoning` field.\n\n"
                    "Current element y-ranges:\n"
                    f"{element_map}\n\n"
                    f"Usable vertical space: y=0.25 to y=7.20 "
                    f"(total: 6.95\")\n\n"
                    "For EACH structural change, answer in your reasoning:\n"
                    "1. **CURRENT**: Which elements occupy which y-ranges?\n"
                    "2. **TARGET**: After the fix, what y-range should each "
                    "affected element occupy?\n"
                    "3. **SPACE MATH**: Does the sum of all element heights "
                    "+ gaps fit within 6.95\"?\n"
                    "4. **CONFLICTS**: Will any element in the target state "
                    "overlap another? Check both y-ranges and x-ranges.\n"
                    "5. **COUPLED MOVES**: Which other elements need to "
                    "move/resize to accommodate this change? Include ALL "
                    "of them in ONE apply_edits call.\n\n"
                    "Protected frame regions (header/title, slide number, source "
                    "attribution, footer/takeaway bar, and ReDeck frame contract) "
                    "are constraints, not spare body space, unless the current issue "
                    "explicitly targets that region. If your planned move makes them "
                    "collide or changes their reading order, reshape the body content "
                    "instead of moving the protected frame/source element.\n\n"
                    "Example planning:\n"
                    "  \"Issue B9 wants a chart. Current: bullets "
                    "y=3.00→5.55 (h=2.55). Available below: "
                    "5.55→7.20 = 1.65\". Chart needs ~1.5\" min.\n"
                    "   Plan: shrink bullets from h=2.55 to h=1.80, "
                    "bottom moves to 4.80. Place chart at y=5.00, "
                    "h=1.50, bottom=6.50. No conflicts.\"\n\n"
                    "**DO NOT skip this step.** The #1 cause of broken "
                    "repairs is adding/moving elements without checking "
                    "whether they fit.\n"
                )

        text_content = "\n".join(parts)

        # If T0 render is available, return multimodal content
        if t0_render_b64:
            return [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{t0_render_b64}",
                    },
                },
                {
                    "type": "text",
                    "text": (
                        "Above is the current slide render (before repair). "
                        "Use this as visual reference while fixing the issues "
                        "below. You can call render_preview at any time to see "
                        "the slide after your edits.\n\n"
                        + text_content
                    ),
                },
            ]

        return text_content
    # ================================================================

    def _parse_action(self, response: str) -> dict | None:
        """Parse agent response into tool call dict.

        If the response contains multiple JSON objects, execute only the first
        valid tool call and discard the rest. Later actions are speculative
        because they were authored before seeing the actual result of the first
        tool call.
        """
        self._last_parse_error_message = ""
        data = _extract_json(response)
        if data and "tool" in data:
            if _has_extra_json(response):
                all_actions = _extract_all_json(response)
                if len(all_actions) > 1:
                    first_action = next(
                        (
                            item for item in all_actions
                            if isinstance(item, dict) and item.get("tool")
                        ),
                        None,
                    )
                    if first_action is not None:
                        ignored = max(0, len(all_actions) - 1)
                        logger.info(
                            "Agent emitted %d JSON tool calls in one message; "
                            "executing the first and ignoring %d speculative call(s).",
                            len(all_actions),
                            ignored,
                        )
                        self._pending_actions = []
                        self._multi_action_ignored_count = ignored
                        return first_action
                    logger.info(
                        "Agent emitted %d JSON tool calls in one message; "
                        "rejecting speculative sequence.",
                        len(all_actions),
                    )
                    self._pending_actions = []
                    self._multi_action_ignored_count = 0
                    self._last_parse_error_message = (
                        "Error: your response contained multiple JSON tool calls. "
                        "The repair loop executes one tool call per turn so each "
                        "next action can use the actual tool result. Return exactly "
                        "one JSON object now, with one `tool` field, and no text or "
                        "additional JSON before or after it. If the rejected message "
                        "included a plan plus edits, send only the next required tool "
                        "call now; do not replay the whole sequence."
                    )
                    return None
            return data
        return None

    def _next_pending_action(self) -> dict | None:
        """Return next queued action from a multi-JSON response, or None."""
        if self._pending_actions:
            return self._pending_actions.pop(0)
        return None

    def _validate_before_submit(
        self, code: str, must_not: list[str],
        must_contain: list[str] | None = None,
    ) -> str | None:
        """Check if MUST NOT strings are still in the code.

        Returns a warning message if problems found, or None if clean.
        This is a generic gate — it only checks the content requirements
        that were already communicated to the agent.
        """
        # Extract banned substrings from MUST NOT entries
        banned: list[str] = []
        for mn in must_not:
            for m in re.findall(r'`([^`]{8,})`', mn):
                banned.append(m)

        # Check which are still present
        still_present = []
        code_lower = code.lower()
        for b in banned:
            if b.lower() in code_lower:
                # Find line numbers where it appears
                lines = code.split("\n")
                line_nums = []
                for i, line in enumerate(lines, 1):
                    if b.lower() in line.lower():
                        line_nums.append(str(i))
                still_present.append(
                    f"  - `{b[:80]}` still at line(s): {', '.join(line_nums[:5])}"
                )

        if not still_present:
            return None

        return (
            "SUBMIT BLOCKED: The following MUST NOT strings are still "
            "present in the code. You must remove or replace them before "
            "submitting:\n"
            + "\n".join(still_present)
            + "\n\nUse apply_edits to fix these, then submit again."
        )

    def _check_contrast_regression(
        self, original_code: str, current_code: str, slide_id: int
    ) -> str | None:
        from .contrast_utils import check_contrast_regression
        return check_contrast_regression(original_code, current_code, slide_id)

    def _check_html_contrast_regression(
        self, original_code: str, current_code: str, slide_id: int
    ) -> str | None:
        from .contrast_utils import check_html_contrast_regression
        return check_html_contrast_regression(original_code, current_code, slide_id)

    def _check_pptx_contrast_regression(
        self, original_code: str, current_code: str, slide_id: int
    ) -> str | None:
        from .contrast_utils import check_pptx_contrast_regression
        return check_pptx_contrast_regression(original_code, current_code, slide_id)

    def _calculate_element_contrast_ratio(
        self, code: str, var_name: str, fill_colors: dict[str, tuple[int, int, int]]
    ) -> float:
        from .contrast_utils import calculate_element_contrast_ratio
        return calculate_element_contrast_ratio(code, var_name, fill_colors)

    def _extract_slide_background_luminance(self, code: str) -> float:
        from .contrast_utils import extract_slide_background_luminance
        return extract_slide_background_luminance(code)

    # ================================================================
    # CONTENT REQUIREMENT EXTRACTION
    # ================================================================

    def _extract_content_requirements(
        self, code: str, all_issues: list[Issue],
    ) -> tuple[list[str], list[str]]:
        """Extract MUST contain and MUST NOT contain from issues."""
        must_contain = []
        must_not_contain = []

        for issue in all_issues:
            fix = issue.planned_fix or ""
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or ""
            )

            # MUST NOT — problematic strings to remove
            if issue.issue_type in {
                "fabricated", "incorrect_claim", "numeric_error",
                "entity_error", "chart_misinterpretation",
            }:
                bad_strings = self._find_problematic_strings(code, desc, fix)
                for bs in bad_strings:
                    must_not_contain.append(
                        f"- `{bs}` -- {issue.issue_type}: {desc[:80]}"
                    )
                # Also check for fabricated data in chart/viz data
                chart_warnings = self._detect_fabricated_in_chart_data(
                    code, desc,
                )
                must_not_contain.extend(chart_warnings)
            elif issue.issue_type == "unfaithful_compression":
                bad_strings = self._find_problematic_strings(code, desc, fix)
                if bad_strings:
                    must_not_contain.append(
                        f"- BANNED CLAIM [{issue.issue_id}]: The following "
                        f"phrasing is over-broad and must NOT appear "
                        f"ANYWHERE in the final code — not in existing "
                        f"text, not in new elements you add, not in any "
                        f"form:"
                    )
                    for bs in bad_strings:
                        must_not_contain.append(f"    `{bs}`")
                    if fix:
                        must_not_contain.append(
                            f"  → REPLACE ALL instances with narrower "
                            f"phrasing: {fix[:250]}"
                        )

            # MUST contain — required content
            if issue.issue_type in CONTENT_ACCURACY_ISSUE_TYPES and fix:
                # Prefer fix_detail.correct_content (exact text from source)
                # over extracting quoted strings from planned_fix (which may
                # quote the WRONG text that should be removed).
                row_specs = self._table_row_specs_for_issue(issue)
                if self._is_table_row_issue(issue) and row_specs:
                    for row in row_specs:
                        must_contain.append(
                            f'- Table row with cells: "{row[:140]}" '
                            f"(from [{issue.issue_type}])"
                        )
                    continue
                correct = (
                    normalize_correct_content_text(
                        issue.fix_detail.correct_content,
                    )
                    if issue.fix_detail and issue.fix_detail.correct_content
                    else ""
                )
                if correct and len(correct) >= 8:
                    label = _source_target_label(issue.issue_type, correct, limit=180)
                    must_contain.append(
                        f"- {label} (from [{issue.issue_type}])"
                    )
                elif issue.issue_type == "missing_entity":
                    must_contain.append(
                        f"- REQUIRED content: {fix[:250]}"
                    )
                elif issue.issue_type == "missing_context":
                    # Only use the planned_fix as guidance, NOT quotes from
                    # the description.  Description quotes are examples of
                    # what the issue CRITICISES (too-generic text), not what
                    # should be added — using them as MUST CONTAIN directly
                    # contradicts unfaithful_compression MUST NOT entries.
                    must_contain.append(
                        f"- REQUIRED context: {fix[:250]}"
                    )
                else:
                    # Fallback: extract quoted strings from planned_fix,
                    # but only those that look like replacement text
                    # (after "to", "with", "→", "->")
                    replace_quoted = re.findall(
                        r'(?:to|with|→|->)\s*"([^"]{8,})"', fix,
                    )
                    if replace_quoted:
                        must_contain.append(
                            f'- Text containing: "{replace_quoted[0]}" '
                            f"(from [{issue.issue_type}])"
                        )
                    else:
                        must_contain.append(
                            f"- Apply: {fix[:200]}"
                        )

        # ── Dedup MUST CONTAIN: merge entries with similar quoted text ──
        # When multiple issues (e.g. incorrect_claim + fabricated) reference
        # the same source claim, keep only one entry to avoid the agent
        # copy-pasting the same sentence into every element.
        if len(must_contain) > 1:
            deduped: list[str] = []
            seen_keys: list[str] = []
            for mc in must_contain:
                quoted = re.findall(r'"([^"]{8,})"', mc)
                key = quoted[0].lower().strip()[:80] if quoted else mc.lower()[:80]
                is_dup = False
                for sk in seen_keys:
                    # Substring match or high word overlap
                    if sk in key or key in sk:
                        is_dup = True
                        break
                    sk_w = set(sk.split())
                    k_w = set(key.split())
                    if sk_w and k_w:
                        overlap = len(sk_w & k_w) / min(len(sk_w), len(k_w))
                        if overlap >= 0.6:
                            is_dup = True
                            break
                if not is_dup:
                    deduped.append(mc)
                    seen_keys.append(key)
                else:
                    logger.debug("Deduped SHOULD CONTAIN: %s", mc[:100])
            must_contain = deduped

        # ── Conflict resolution: MUST NOT trumps MUST CONTAIN ──
        # Cross-slide missing_context issues can add quoted text to
        # MUST CONTAIN that directly contradicts unfaithful_compression
        # MUST NOT entries.  The MUST NOT is more specific (evidence-
        # grounded narrowing) and must win.
        if must_not_contain and must_contain:
            # Collect all banned substrings (lowercase) from MUST NOT
            banned_phrases: list[str] = []
            for mn in must_not_contain:
                for m in re.findall(r'`([^`]{5,})`', mn):
                    banned_phrases.append(m.lower())

            # Filter out any MUST CONTAIN entry whose quoted text
            # substantially overlaps with a banned phrase
            filtered_must_contain: list[str] = []
            for mc in must_contain:
                mc_quoted = re.findall(r'"([^"]{5,})"', mc)
                conflict = False
                for mq in mc_quoted:
                    mq_lower = mq.lower()
                    for bp in banned_phrases:
                        # Check bidirectional substring overlap
                        if bp in mq_lower or mq_lower in bp:
                            conflict = True
                            break
                        # Also check significant word overlap (≥60%)
                        bp_words = set(bp.split())
                        mq_words = set(mq_lower.split())
                        if bp_words and mq_words:
                            overlap = len(bp_words & mq_words)
                            ratio = overlap / min(len(bp_words), len(mq_words))
                            if ratio >= 0.6:
                                conflict = True
                                break
                    if conflict:
                        break
                if not conflict:
                    filtered_must_contain.append(mc)
                else:
                    logger.debug(
                        "Dropped conflicting MUST CONTAIN: %s", mc[:100],
                    )
            must_contain = filtered_must_contain

        return must_contain, must_not_contain

    def _find_problematic_strings(
        self, code: str, desc: str, fix: str,
    ) -> list[str]:
        """Find exact strings in code that an issue identifies as wrong."""
        # Find text strings in code — both .text = "..." and function args
        # Use separate patterns for double/single quotes to handle
        # strings containing the other quote type (e.g. "Mamba's")
        code_texts = re.findall(
            r'\.text\s*=\s*[f]?"(.+?)"', code,
        )
        code_texts += re.findall(
            r"\.text\s*=\s*[f]?'(.+?)'", code,
        )
        # Also find string literals in function calls (add_card, etc.)
        code_texts += re.findall(
            r'(?:^|[,(])\s*"([^"]{10,})"', code,
        )
        code_texts += re.findall(
            r"(?:^|[,(])\s*'([^']{10,})'", code,
        )
        if not code_texts:
            return []

        quoted_in_desc = re.findall(r'"([^"]{5,})"', desc)
        quoted_in_desc += re.findall(r"'([^']{5,})'", desc)
        # Also look in the fix text for quoted strings
        if fix:
            quoted_in_fix = re.findall(r'"([^"]{5,})"', fix)
            quoted_in_fix += re.findall(r"'([^']{5,})'", fix)

        # Clean up: strip bullet markers and leading/trailing whitespace
        cleaned_quoted = []
        for q in quoted_in_desc:
            # Strip leading bullet markers (•, -, *, >) and spaces
            cleaned = re.sub(r'^[•\-\*\>\s]+', '', q).strip()
            if len(cleaned) >= 5:
                cleaned_quoted.append(cleaned)
            # Keep original too for exact match
            cleaned_quoted.append(q)

        matches = []
        for text in code_texts:
            text_lower = text.lower()
            for q in cleaned_quoted:
                if q.lower() in text_lower:
                    matches.append(text)
                    break

        return matches[:5]

    @staticmethod
    def _detect_fabricated_in_chart_data(
        code: str, desc: str,
    ) -> list[str]:
        """Detect if fabricated values appear in chart data (JSON or add_series).

        Returns warning lines to add to MUST NOT section when the
        problematic data is embedded in chart/viz data structures rather
        than in text strings.
        """
        # Extract numbers from the issue description
        desc_numbers = set()
        for m in re.finditer(r'\b(\d+\.?\d+)\b', desc):
            val = m.group(1)
            if len(val) >= 3:  # At least 3 chars (e.g. "61.1")
                desc_numbers.add(val)

        if not desc_numbers:
            return []

        warnings = []
        # Check viz_data_str JSON strings
        for m in re.finditer(
            r'viz_data_str\s*=\s*[\'"](.+?)[\'"]', code, re.DOTALL,
        ):
            json_str = m.group(1)
            found_in_json = [n for n in desc_numbers if n in json_str]
            if found_in_json:
                warnings.append(
                    f"- CHART DATA (viz_data_str) contains values "
                    f"flagged as fabricated: {found_in_json}"
                )

        # Check add_series calls
        for m in re.finditer(
            r'add_series\([^)]+\)', code,
        ):
            series_str = m.group(0)
            found_in_series = [n for n in desc_numbers if n in series_str]
            if found_in_series:
                warnings.append(
                    f"- chart add_series() contains values flagged as "
                    f"fabricated: {found_in_series}"
                )

        # Note about embedded images
        if 'add_picture' in code:
            warnings.append(
                "- NOTE: Slide has embedded bitmap images. Content "
                "inside these images cannot be edited by code changes."
            )

        return warnings

    def _build_content_checklist(
        self, all_issues: list[Issue],
    ) -> str:
        """Build content verification checklist."""
        items = []

        if any(issue.issue_type in CONTENT_ACCURACY_ISSUE_TYPES for issue in all_issues):
            items.append(
                "[ ] Fixed-format budget: titles/page headers stay concise; "
                "do not put long source-backed corrections into full-width "
                "bottom bars, footers, or source notes. Merge longer qualifiers "
                "into body/interpretation text and verify layout after length changes."
            )

        for issue in all_issues:
            fix = issue.planned_fix or ""
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or ""
            )
            issue_id = getattr(issue, 'issue_id', '') or ''
            normalized_correct = ""
            if issue.fix_detail and issue.fix_detail.correct_content:
                normalized_correct = normalize_correct_content_text(
                    issue.fix_detail.correct_content,
                )

            if issue.issue_type in {
                "fabricated", "incorrect_claim", "numeric_error",
                "entity_error", "chart_misinterpretation",
                "unfaithful_compression",
            } and fix:
                quoted = [normalized_correct] if normalized_correct else re.findall(r'"([^"]{8,})"', fix)
                if quoted:
                    target_label = _source_target_label(
                        issue.issue_type, quoted[0], limit=220,
                    )
                    items.append(
                        f'[ ] [{issue.issue_type}] ({issue_id}) '
                        f'Must cover: {target_label}'
                    )
                else:
                    items.append(
                        f"[ ] [{issue.issue_type}] ({issue_id}) "
                        f"Fix: {fix[:150]}"
                    )
            elif issue.issue_type == "missing_entity" and fix:
                row_specs = self._table_row_specs_for_issue(issue)
                if self._is_table_row_issue(issue) and row_specs:
                    for row in row_specs:
                        items.append(
                            f"[ ] [{issue.issue_type}] ({issue_id}) "
                            f"REQUIRED TABLE ROW: {row[:180]} "
                            "-- add as <tr>/<td> cells inside the existing table, "
                            "not as prose text."
                        )
                else:
                    items.append(
                        f"[ ] [{issue.issue_type}] ({issue_id}) "
                        f"REQUIRED: {(normalized_correct or fix)[:200]}"
                    )
            elif issue.issue_type == "missing_context" and fix:
                items.append(
                    f"[ ] [{issue.issue_type}] ({issue_id}) "
                    f"Must add context: {(normalized_correct or fix)[:200]}"
                )
            elif issue.issue_type in {
                "missing_point", "missing_evidence", "missing_conclusion",
                "misleading_omission",
            } and fix:
                items.append(
                    f"[ ] [{issue.issue_type}] ({issue_id}) "
                    f"Must include: {(normalized_correct or fix)[:150]}"
                )
            elif issue.issue_type == "missing_data_visualization":
                items.append(
                    f"[ ] [data_viz] ({issue_id}) "
                    f"Must include an add_chart() call"
                )
            elif issue.issue_type == "overlap" and fix:
                items.append(
                    f"[ ] [overlap] ({issue_id}) {fix[:150]}. "
                    f"No shape may overlap a figure."
                )

        if not items:
            return ""

        return (
            "\n## Content Verification Checklist\n\n"
            "Every item below must be addressed in the repair.\n\n"
            + "\n".join(items)
        )

    @staticmethod
    def _compute_redistribution_guide(
        code: str,
        spatial_state,
        all_issues: list[Issue],
    ) -> str:
        """Compute spatial distribution facts for density/distribution issues.

        Reports computed facts about how elements are distributed on the
        slide: center of mass, top/bottom balance, available space.
        The LLM decides what to do with these facts.
        """
        distribution_types = {
            "density_imbalance", "layout_inappropriate",
            "whitespace_imbalance",
        }
        has_distribution = any(
            i.issue_type in distribution_types for i in all_issues
        )
        if not has_distribution:
            return ""

        if not spatial_state or not spatial_state.blocks:
            return ""

        # Filter to substantive elements (skip tiny decorative ones)
        blocks = [
            b for b in spatial_state.blocks
            if b.w > 0.5 and b.h > 0.3
        ]
        if len(blocks) < 2:
            return ""

        # Sort by y position
        blocks_sorted = sorted(blocks, key=lambda b: b.y)

        # Compute center of mass (weighted by area)
        total_area = sum(b.w * b.h for b in blocks_sorted)
        if total_area < 0.1:
            return ""
        center_y = sum(
            (b.y + b.h / 2) * b.w * b.h for b in blocks_sorted
        ) / total_area
        center_x = sum(
            (b.x + b.w / 2) * b.w * b.h for b in blocks_sorted
        ) / total_area

        slide_mid_y = 3.75  # middle of 7.5" slide
        slide_mid_x = 6.67  # middle of 13.333" slide

        # Detect top-bottom imbalance
        top_blocks = [b for b in blocks_sorted if b.y + b.h / 2 < slide_mid_y]
        bottom_blocks = [b for b in blocks_sorted if b.y + b.h / 2 >= slide_mid_y]
        top_area = sum(b.w * b.h for b in top_blocks)
        bottom_area = sum(b.w * b.h for b in bottom_blocks)
        tb_imbalance = top_area / (bottom_area + 0.01)

        # Detect left-right imbalance
        left_blocks = [b for b in blocks_sorted if b.x + b.w / 2 < slide_mid_x]
        right_blocks = [b for b in blocks_sorted if b.x + b.w / 2 >= slide_mid_x]
        left_area = sum(b.w * b.h for b in left_blocks)
        right_area = sum(b.w * b.h for b in right_blocks)
        lr_imbalance = left_area / (right_area + 0.01)

        # Only report if significantly imbalanced in at least one axis
        tb_imbalanced = not (0.4 < tb_imbalance < 2.5)
        lr_imbalanced = not (0.4 < lr_imbalance < 2.5)
        if not tb_imbalanced and not lr_imbalanced:
            return ""  # Reasonably balanced

        lines = [
            "## Spatial Distribution Analysis\n",
        ]

        # Report current distribution facts
        if tb_imbalanced:
            occupied_top = f"{len(top_blocks)} elements, {top_area:.1f} sq.in"
            occupied_bot = f"{len(bottom_blocks)} elements, {bottom_area:.1f} sq.in"
            target_ratio = "0.6:1 to 1.7:1"
            lines.append(
                f"Vertical center of mass: y={center_y:.1f}\" "
                f"(slide midpoint: {slide_mid_y}\")\n"
                f"- Top half (y < {slide_mid_y}\"): {occupied_top}\n"
                f"- Bottom half (y ≥ {slide_mid_y}\"): {occupied_bot}\n"
                f"- Top/bottom ratio: {tb_imbalance:.1f}:1 "
                f"(target: {target_ratio})\n"
            )

        if lr_imbalanced:
            occupied_left = f"{len(left_blocks)} elements, {left_area:.1f} sq.in"
            occupied_right = f"{len(right_blocks)} elements, {right_area:.1f} sq.in"
            target_ratio = "0.6:1 to 1.7:1"
            lines.append(
                f"Horizontal center of mass: x={center_x:.1f}\" "
                f"(slide midpoint: {slide_mid_x}\")\n"
                f"- Left half (x < {slide_mid_x}\"): {occupied_left}\n"
                f"- Right half (x ≥ {slide_mid_x}\"): {occupied_right}\n"
                f"- Left/right ratio: {lr_imbalance:.1f}:1 "
                f"(target: {target_ratio})\n"
            )

        # Report available space
        last_bottom = max(b.y + b.h for b in blocks_sorted)
        available_below = USABLE_BOTTOM - last_bottom
        if available_below > 1.0:
            lines.append(
                f"Available space below last element: "
                f"{available_below:.1f}\" (y={last_bottom:.2f}\" to "
                f"y={USABLE_BOTTOM:.2f}\")\n"
            )

        # Quantify the minimum change needed
        if tb_imbalanced or lr_imbalanced:
            # Find the two most imbalanced element pairs
            dominant = sorted(
                blocks_sorted,
                key=lambda b: b.w * b.h,
                reverse=True,
            )
            if len(dominant) >= 2:
                biggest = dominant[0]
                biggest_area = biggest.w * biggest.h
                avg_area = total_area / len(blocks_sorted)
                if biggest_area > avg_area * 2.5:
                    lines.append(
                        f"Largest element ({biggest.var_name}): "
                        f"{biggest.w:.1f}\"×{biggest.h:.1f}\" = "
                        f"{biggest_area:.1f} sq.in "
                        f"({biggest_area/total_area*100:.0f}% of total)\n"
                    )

        return "\n".join(lines)

    @classmethod
    def _build_external_svg_source_context(
        cls,
        issues: list[Issue],
        spatial_state,
    ) -> str:
        """Expose source SVG code for external SVG visual repairs.

        Browser DOM extraction can only see an external SVG as an <img>. When
        the repair target is inside that SVG, the agent needs the actual asset
        source to make a surgical label/geometry fix instead of redrawing from a
        vague visual summary.
        """
        if not spatial_state or not any(
            issue.issue_type == "svg_visual_defect" for issue in issues
        ):
            return ""

        issue_text = "\n".join(
            " ".join(filter(None, (
                getattr(getattr(issue, "evidence", None), "description", ""),
                getattr(issue, "planned_fix", ""),
                " ".join(getattr(getattr(issue, "evidence", None), "object_refs", []) or []),
            )))
            for issue in issues
            if issue.issue_type == "svg_visual_defect"
        ).lower()

        paths: list[Path] = []
        for svg_issue in getattr(spatial_state, "svg_asset_issues", []) or []:
            asset_path = svg_issue.get("asset_path")
            if not asset_path:
                continue
            path = Path(str(asset_path))
            haystack = f"{path.name} {asset_path} {svg_issue.get('label', '')}".lower()
            if issue_text and not any(
                token and token in issue_text
                for token in (path.name.lower(), str(svg_issue.get("label", "")).lower())
            ) and haystack not in issue_text:
                continue
            if path.exists() and path not in paths:
                paths.append(path)

        if not paths:
            return ""

        parts = [
            "## External SVG Asset Source",
            "For `svg_visual_defect` repairs on external `<img src=\"*.svg\">` assets, preserve the original SVG canvas, all visible SVG text, and all non-target graph/node/connector roles. Make the smallest local SVG source change needed for the named defect. Do not redraw, simplify, rename labels, change window/model names, or replace captions.",
        ]
        for path in paths[:3]:
            try:
                svg = path.read_text(encoding="utf-8", errors="replace")
            except Exception as exc:
                parts.append(f"### {path}\nCould not read SVG asset: {exc}")
                continue
            texts = cls._svg_visible_texts(svg)
            if len(svg) > 16000:
                shown_svg = svg[:16000] + "\n<!-- truncated -->"
            else:
                shown_svg = svg
            parts.append(
                f"### {path}\n"
                f"Visible SVG text to preserve: {json.dumps(texts, ensure_ascii=False)}\n"
                f"```xml\n{shown_svg}\n```"
            )
        return "\n\n".join(parts)

    def _generate_layout_plan(
        self,
        code: str,
        all_issues: list[Issue],
        spatial_state,
        bp_slide: BlueprintSlide | None,
    ) -> dict | None:
        """Generate a target layout plan for complex layout issues.

        Makes a focused LLM call to compute target positions for all
        elements, producing a concrete {var_name: {left, top, width, height}}
        plan that the reflow_layout tool executes atomically.

        The prompt provides rich context:
        - Full element list with text content preview + shape type
        - Current spatial problems (overlaps, gaps, imbalance)
        - Slide purpose from the blueprint
        - Strict design rules (gaps, margins, proportions)

        Returns dict[str, dict] or None on failure.
        """
        if not spatial_state or not spatial_state.blocks:
            return None

        # Build element list with rich context
        sorted_blocks = sorted(spatial_state.blocks, key=lambda b: b.y)
        element_lines = []
        for b in sorted_blocks:
            if b.w < 0.3 and b.h < 0.1:
                continue
            bottom = b.y + b.h
            # Show text content preview for semantic understanding
            text_preview = ""
            if b.text_lines:
                preview_text = " | ".join(l[:50] for l in b.text_lines[:3])
                text_preview = f'  content: "{preview_text}"'
            element_lines.append(
                f"  {b.var_name} ({b.shape_type}):"
                f" x={b.x:.2f}, y={b.y:.2f}, w={b.w:.2f}, h={b.h:.2f},"
                f" bottom={bottom:.2f}"
                f"{text_preview}"
            )
        elements_str = "\n".join(element_lines)

        # Detect current spatial problems
        spatial_diagnosis = []
        if spatial_state.overlap_pairs:
            spatial_diagnosis.append(
                f"OVERLAPS: {', '.join(f'{a}↔{b}' for a, b, *_ in spatial_state.overlap_pairs[:5])}"
            )

        # Detect large empty gaps
        gaps = []
        for i in range(len(sorted_blocks) - 1):
            b1 = sorted_blocks[i]
            b2 = sorted_blocks[i + 1]
            gap = b2.y - (b1.y + b1.h)
            if gap > 0.5:
                gaps.append(f"{b1.var_name}→{b2.var_name}: {gap:.1f}\" gap")
        if gaps:
            spatial_diagnosis.append(f"LARGE GAPS: {'; '.join(gaps[:5])}")

        # Check if content is squeezed into top half
        usable_bottom = 7.20
        last_bottom = max((b.y + b.h for b in sorted_blocks), default=0)
        unused_bottom = usable_bottom - last_bottom
        if unused_bottom > 2.0:
            spatial_diagnosis.append(
                f"WASTED SPACE: {unused_bottom:.1f}\" unused below last element "
                f"(last element ends at y={last_bottom:.1f}\")"
            )

        diagnosis_str = "\n".join(f"  • {d}" for d in spatial_diagnosis) if spatial_diagnosis else "  (none detected)"

        # Build issue descriptions
        issue_descs = []
        for issue in all_issues:
            desc = issue.evidence.description or issue.why_this_fails or ""
            issue_descs.append(f"  [{issue.issue_type}] {desc[:250]}")
        issues_str = "\n".join(issue_descs)

        # Slide purpose from blueprint
        slide_context = ""
        if bp_slide:
            slide_context = (
                f"\nSlide purpose: {bp_slide.primary_proposition or 'untitled'}\n"
                f"Slide role: {bp_slide.role or 'unknown'}\n"
                f"Layout hint: {bp_slide.layout_hint or 'unspecified'}\n"
            )

        prompt = f"""You are an expert slide layout designer. Analyze the current slide layout and compute optimal positions for all elements.

## Slide Info{slide_context}

## Current Elements
{elements_str}

## Spatial Problems
{diagnosis_str}

## Issues to Fix
{issues_str}

## Design Rules
1. Usable area: x=0.50 to 13.00, y=0.25 to 7.20 (below title)
2. Minimum gap between ANY two elements: 0.12" vertical, 0.15" horizontal
3. No overlaps allowed
4. Title elements (y < 1.0) should stay in place
5. Source footnote (last element with "Source:" text) should be at bottom (y ≈ 6.80)
6. Content should fill the usable area evenly — no large empty regions
7. Related elements (e.g., a heading and its bullet list) should stay close together
8. If elements are semantically paired (left-right comparison, before-after), keep them side by side
9. Charts/figures should be at least 2.5" tall to be readable
10. Text boxes should be wide enough for their content (min 3" for bullet lists)

## Output Format
Output a JSON object mapping EVERY content element's var_name to its target position.
Include ALL four coordinates: left, top, width, height (in inches).
Do NOT include title elements or source footnotes unless they need to move.

Example:
{{"bullet_box": {{"left": 0.50, "top": 1.50, "width": 5.80, "height": 3.00}}, "chart_panel": {{"left": 6.75, "top": 1.50, "width": 5.75, "height": 3.00}}}}

Think step by step:
1. What is the slide's content structure? (bullet list, comparison, flowchart, data display)
2. What layout pattern best fits this content? (single column, two column, grid, split)
3. Compute specific coordinates that eliminate the spatial problems above.

Output ONLY the JSON object, no other text."""

        try:
            response = self.llm.call_text(
                system_prompt="You are an expert presentation layout designer. Output only valid JSON.",
                user_content=prompt,
                model=self.model,
                module_name="layout_preplan",
                prompt_version="v2",
                max_tokens=3000,
                temperature=0.1,
            )

            # Parse JSON from response
            data = _extract_json(response)
            if not data or not isinstance(data, dict):
                logger.warning("Layout pre-plan: failed to parse response")
                return None

            # Validate: each value should have at least top
            plan: dict[str, dict] = {}
            for var_name, target in data.items():
                if not isinstance(target, dict):
                    continue
                entry = {}
                for k in ("left", "top", "width", "height"):
                    if k in target:
                        try:
                            val = float(target[k])
                            if k in ("left", "width") and 0.0 <= val <= 14.0:
                                entry[k] = round(val, 2)
                            elif k in ("top", "height") and 0.0 <= val <= 8.0:
                                entry[k] = round(val, 2)
                        except (ValueError, TypeError):
                            pass
                # Need at least top+height or left+top to be useful
                if len(entry) >= 2:
                    plan[var_name] = entry

            if not plan:
                logger.warning("Layout pre-plan: no valid targets")
                return None

            logger.info(
                "Layout pre-plan: %d targets computed (elements: %s)",
                len(plan),
                ", ".join(f"{k}→({v.get('left','?')},{v.get('top','?')})" for k,v in plan.items()),
            )
            return plan

        except Exception as e:
            logger.warning("Layout pre-plan failed: %s", str(e)[:200])
            return None

    @staticmethod
    def _build_layout_diagnostic(
        code: str,
        spatial_state,
        all_issues: list[Issue],
    ) -> str:
        """Build a compact, high-density layout diagnostic.

        Single unified representation that combines:
        - Per-element card: position, size, fill, font, visual weight, role
        - Pairwise gaps (only significant ones)
        - Issue-to-element mapping
        - Concrete fix prescriptions

        Format designed for LLM processing efficiency:
        - Structured key=value pairs (high signal density)
        - Sorted by visual dominance (action priority)
        - Fix prescriptions are specific code changes, not abstract advice
        """
        if not spatial_state or not spatial_state.blocks:
            return ""

        usable_w = USABLE_RIGHT - USABLE_LEFT
        usable_h = USABLE_BOTTOM - USABLE_TOP
        usable_area = usable_w * usable_h

        # Extract colors and compute per-block metadata
        fill_colors: dict[str, tuple[int, int, int]] = {}
        for m in re.finditer(
            r'(\w+)\.fill\.fore_color\.rgb\s*=\s*RGBColor\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)',
            code,
        ):
            fill_colors[m.group(1)] = (int(m.group(2)), int(m.group(3)), int(m.group(4)))
        theme_dict: dict[str, tuple[int, int, int]] = {}
        for m in re.finditer(
            r'"(\w+)"\s*:\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)',
            code,
        ):
            theme_dict[m.group(1)] = (int(m.group(2)), int(m.group(3)), int(m.group(4)))
        for m in re.finditer(
            r'(\w+)\.fill\.fore_color\.rgb\s*=\s*RGBColor\(\s*\*\s*theme_colors\[\s*"(\w+)"\s*\]\s*\)',
            code,
        ):
            key = m.group(2)
            if key in theme_dict:
                fill_colors[m.group(1)] = theme_dict[key]

        # Build element entries
        elements = []
        for b in spatial_state.blocks:
            if b.w < 0.3 and b.h < 0.1:
                continue
            area = b.w * b.h
            pct = area / usable_area * 100

            brightness = 0.95
            fill_str = "white"
            if b.var_name in fill_colors:
                r, g, bb = fill_colors[b.var_name]
                brightness = (0.299 * r + 0.587 * g + 0.114 * bb) / 255
                fill_str = f"rgb({r},{g},{bb})"

            font_pt = b.font_size_pt if b.font_size_pt > 0 else 14
            # Visual weight = area × (1-brightness) × (font/16)
            vw = area * max(0.1, 1.0 - brightness) * (font_pt / 16.0)

            text_preview = ""
            if b.text_lines:
                text_preview = b.text_lines[0][:35]

            elements.append({
                "var": b.var_name,
                "x": b.x, "y": b.y, "w": b.w, "h": b.h,
                "area_pct": pct,
                "fill": fill_str,
                "brightness": brightness,
                "font_pt": font_pt,
                "vw": vw,
                "text": text_preview,
                "bottom": b.y + b.h,
            })

        elements.sort(key=lambda e: e["vw"], reverse=True)
        max_vw = elements[0]["vw"] if elements else 1

        lines = ["## Layout Diagnostic\n"]

        # Section 1: Element cards (sorted by visual weight)
        lines.append("### Elements (sorted by visual dominance)\n")
        lines.append(
            "```"
        )
        for i, e in enumerate(elements):
            rel_dom = e["vw"] / max_vw * 100 if max_vw > 0 else 0
            dark_label = "DARK" if e["brightness"] < 0.4 else "MID" if e["brightness"] < 0.7 else "LIGHT"
            lines.append(
                f"{i+1}. {e['var']:18s} | pos=({e['x']:.1f},{e['y']:.1f}) "
                f"size={e['w']:.1f}x{e['h']:.1f} ({e['area_pct']:.0f}%) | "
                f"fill={e['fill']:18s} {dark_label:5s} | "
                f"{e['font_pt']:.0f}pt | "
                f"dominance={rel_dom:.0f}% | "
                f"\"{e['text']}\""
            )
        lines.append("```\n")

        # Section 2: Vertical scan (gaps and density)
        sorted_by_y = sorted(elements, key=lambda e: e["y"])
        lines.append("### Vertical Scan (top→bottom)\n")
        prev_bottom = 0.0
        for e in sorted_by_y:
            gap = e["y"] - prev_bottom
            if gap > 0.3:
                lines.append(f"  GAP {gap:.1f}\" (y={prev_bottom:.1f}→{e['y']:.1f})")
            lines.append(
                f"  [{e['var']}] y={e['y']:.1f}→{e['bottom']:.1f} "
                f"({e['h']:.1f}\" tall, {e['area_pct']:.0f}% area)"
            )
            prev_bottom = max(prev_bottom, e["bottom"])
        remaining = 7.2 - prev_bottom
        if remaining > 0.5:
            lines.append(f"  GAP {remaining:.1f}\" (y={prev_bottom:.1f}→7.2) UNUSED SPACE")
        lines.append("")

        # Section 3: Fix prescriptions
        dominant = [e for e in elements if e["vw"] / max_vw > 0.5]
        secondary = [e for e in elements if e["vw"] / max_vw < 0.25 and e["area_pct"] > 1]

        if dominant:
            lines.append("### Fix Prescriptions\n")
            for d in dominant:
                rel = d["vw"] / max_vw * 100
                fixes = []

                # Size fix
                if d["area_pct"] > 10:
                    new_h = max(0.4, d["h"] * 0.5)
                    fixes.append(
                        f"RESIZE: Inches({d['h']:.2f}) → Inches({new_h:.2f}) "
                        f"(halve height, frees {d['h'] - new_h:.1f}\")"
                    )

                # Color fix
                if d["brightness"] < 0.6:
                    fixes.append(
                        f"RECOLOR: change {d['fill']} to a lighter shade "
                        f"(e.g., RGBColor(230, 230, 230) or white)"
                    )

                # Font fix
                if d["font_pt"] >= 20:
                    new_pt = max(16, int(d["font_pt"] * 0.7))
                    fixes.append(
                        f"REFONT: Pt({d['font_pt']:.0f}) → Pt({new_pt}) "
                        f"(reduce visual prominence)"
                    )

                if fixes:
                    lines.append(
                        f"  ▼ REDUCE '{d['var']}' (dominance={rel:.0f}%): "
                        f"apply {len(fixes)} changes:"
                    )
                    for fix in fixes:
                        lines.append(f"    • {fix}")

            if secondary:
                for s in secondary[:2]:
                    new_h = min(s["h"] + 1.5, 4.0)
                    lines.append(
                        f"  ▲ EXPAND '{s['var']}': "
                        f"Inches({s['h']:.2f}) → Inches({new_h:.2f}), "
                        f"move up to fill freed space"
                    )

            lines.append(
                "\n  RULE: Changes must affect ≥2 of {size, color, font} "
                "per dominant element. Size-only changes are invisible."
            )

        lines.append("")
        return "\n".join(lines)

    @staticmethod
    def _build_layout_guide(
        code: str,
        spatial_state,
        all_issues: list[Issue],
    ) -> str:
        """Build quantitative layout restructuring guide.

        For layout issues (B2/B5/B8), compute area percentages, visual
        weight (color brightness + font size + area), and suggest
        specific multi-dimensional fixes:
        - Resize (Inches changes)
        - Recolor (RGBColor changes to reduce contrast dominance)
        - Refont (Pt changes to reduce visual prominence)

        The verifier compares BEFORE/AFTER screenshots. Changes must
        affect the VISUAL APPEARANCE, not just coordinates. A dark
        full-width banner at 20pt still dominates even if shrunk by 0.3".
        """
        if not spatial_state or not spatial_state.blocks:
            return ""

        usable_w = USABLE_RIGHT - USABLE_LEFT  # ~12.5"
        usable_h = USABLE_BOTTOM - USABLE_TOP  # ~6.95"
        usable_area = usable_w * usable_h

        # Extract fill colors and font sizes from code
        fill_colors: dict[str, tuple[int, int, int]] = {}
        # Direct RGBColor
        for m in re.finditer(
            r'(\w+)\.fill\.fore_color\.rgb\s*=\s*RGBColor\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)',
            code,
        ):
            fill_colors[m.group(1)] = (int(m.group(2)), int(m.group(3)), int(m.group(4)))

        # Theme-based RGBColor
        theme_dict: dict[str, tuple[int, int, int]] = {}
        for m in re.finditer(
            r'"(\w+)"\s*:\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)',
            code,
        ):
            theme_dict[m.group(1)] = (int(m.group(2)), int(m.group(3)), int(m.group(4)))
        for m in re.finditer(
            r'(\w+)\.fill\.fore_color\.rgb\s*=\s*RGBColor\(\s*\*\s*theme_colors\[\s*"(\w+)"\s*\]\s*\)',
            code,
        ):
            key = m.group(2)
            if key in theme_dict:
                fill_colors[m.group(1)] = theme_dict[key]

        # Extract font sizes by shape variable
        font_sizes: dict[str, list[int]] = {}
        # Find set_text_style calls or direct Pt() assignments near shape vars
        for b in spatial_state.blocks:
            sizes = []
            if b.font_size_pt > 0:
                sizes.append(int(b.font_size_pt))
            font_sizes[b.var_name] = sizes

        # Compute visual weight score: area × darkness × font_scale
        elements = []
        for b in sorted(spatial_state.blocks, key=lambda b: b.w * b.h, reverse=True):
            area = b.w * b.h
            pct = area / usable_area * 100

            # Brightness (0=black, 1=white)
            brightness = 1.0
            if b.var_name in fill_colors:
                r, g, bb = fill_colors[b.var_name]
                brightness = (0.299 * r + 0.587 * g + 0.114 * bb) / 255

            # Darkness factor: dark elements draw more attention
            darkness_factor = max(0.3, 1.0 - brightness)

            # Font scale: larger fonts draw more attention
            font_pt = b.font_size_pt if b.font_size_pt > 0 else 14
            font_factor = font_pt / 16.0  # normalize to 16pt baseline

            # Visual weight = area × darkness × font_scale
            visual_weight = area * darkness_factor * font_factor

            color_str = ""
            weight_label = "LIGHT"
            if b.var_name in fill_colors:
                r, g, bb = fill_colors[b.var_name]
                weight_label = "DARK" if brightness < 0.4 else "MID" if brightness < 0.7 else "LIGHT"
                color_str = f"fill=({r},{g},{bb})"

            elements.append({
                "var": b.var_name,
                "w": b.w, "h": b.h,
                "area": area, "pct": pct,
                "brightness": brightness,
                "weight_label": weight_label,
                "color_str": color_str,
                "font_pt": font_pt,
                "visual_weight": visual_weight,
                "text_preview": (b.text_lines[0][:40] if b.text_lines else ""),
            })

        if not elements:
            return ""

        # Sort by visual weight (what the human eye sees)
        elements.sort(key=lambda e: e["visual_weight"], reverse=True)
        max_vw = elements[0]["visual_weight"] if elements else 1

        lines = [
            "## Layout Analysis & Visual Dominance\n",
            "The verifier compares screenshots. Visual dominance comes from "
            "THREE factors: size + dark color + large font. You must change "
            "MULTIPLE factors to shift visual balance — resizing alone is often insufficient.\n",
            "",
            "Elements ranked by visual dominance (what the eye sees):",
        ]

        dominant_elements = []
        secondary_elements = []

        for e in elements:
            rel_weight = e["visual_weight"] / max_vw * 100 if max_vw > 0 else 0
            font_str = f"{e['font_pt']:.0f}pt" if e["font_pt"] > 0 else ""
            lines.append(
                f"  {e['var']:20s}  {e['w']:.1f}\" x {e['h']:.1f}\" "
                f"({e['pct']:.0f}% area)  {e['weight_label']:5s} "
                f"{e['color_str']:20s} {font_str:6s} "
                f"dominance={rel_weight:.0f}%"
                f"  {e['text_preview']}"
            )
            if rel_weight > 60:
                dominant_elements.append(e)
            elif rel_weight < 30:
                secondary_elements.append(e)

        # Generate specific fix recommendations
        if dominant_elements and secondary_elements:
            lines.append("")
            lines.append("### Recommended fixes (do ALL that apply):\n")

            for de in dominant_elements:
                fix_actions = []
                if de["pct"] > 15:
                    target_h = max(0.5, de["h"] * 0.5)
                    fix_actions.append(
                        f"RESIZE: reduce height from {de['h']:.1f}\" to ~{target_h:.1f}\" "
                        f"(saves {de['h'] - target_h:.1f}\")"
                    )
                if de["brightness"] < 0.5:
                    fix_actions.append(
                        "RECOLOR: change dark fill to a lighter shade "
                        "(e.g., replace dark accent with light_bg or white)"
                    )
                if de["font_pt"] >= 20:
                    target_pt = max(16, int(de["font_pt"] * 0.7))
                    fix_actions.append(
                        f"REFONT: reduce from {de['font_pt']:.0f}pt to {target_pt}pt"
                    )

                if fix_actions:
                    lines.append(
                        f"  REDUCE dominance of '{de['var']}' "
                        f"(currently {de['visual_weight'] / max_vw * 100:.0f}%):"
                    )
                    for fa in fix_actions:
                        lines.append(f"    - {fa}")

            lines.append("")
            for se in secondary_elements[:2]:
                target_h = min(se["h"] + 1.5, 5.0)
                lines.append(
                    f"  INCREASE presence of '{se['var']}': "
                    f"expand height from {se['h']:.1f}\" to ~{target_h:.1f}\", "
                    f"move up into freed space"
                )

            lines.append(
                "\n  KEY PRINCIPLE: A dark full-width banner at 24pt still dominates "
                "even if shrunk by 0.3\". You must ALSO lighten the color and reduce "
                "the font size to actually shift visual balance."
            )

        lines.append("")
        return "\n".join(lines)

    @staticmethod
    def _build_ascii_grid(spatial_state) -> str:
        from .spatial_helpers import build_ascii_grid
        return build_ascii_grid(spatial_state)

    @staticmethod
    def _build_pairwise_relations(spatial_state, code: str) -> str:
        from .spatial_helpers import build_pairwise_relations
        return build_pairwise_relations(spatial_state, code)

    @staticmethod
    def _build_vertical_strip(spatial_state) -> str:
        from .spatial_helpers import build_vertical_strip
        return build_vertical_strip(spatial_state)

    @staticmethod
    def _build_elements_json(code: str, spatial_state) -> str:
        from .spatial_helpers import build_elements_json
        return build_elements_json(code, spatial_state)

    @staticmethod
    def _find_available_slot(
        spatial_state,
    ) -> tuple[float, float, float, float] | None:
        from .spatial_helpers import find_available_slot
        return find_available_slot(spatial_state)

    @staticmethod
    def _compute_spatial_context(spatial_state) -> str:
        from .spatial_helpers import compute_spatial_context
        return compute_spatial_context(spatial_state)

    @staticmethod
    def _compute_coverage_pct(spatial_state) -> float:
        from .spatial_helpers import compute_coverage_pct
        return compute_coverage_pct(spatial_state)

    @staticmethod
    def _format_spatial_issue_with_px(
        issue: "Issue",
        spatial_state: "SlideState",
    ) -> str:
        from .spatial_helpers import format_spatial_issue_with_px
        return format_spatial_issue_with_px(issue, spatial_state)

    @staticmethod
    def _extract_palette_note(code: str, is_html: bool) -> str:
        from .spatial_helpers import extract_palette_note
        return extract_palette_note(code, is_html)

    @staticmethod
    def _annotate_issue_locations(
        issue: Issue,
        code: str,
        code_lines: list[str],
    ) -> str:
        from .spatial_helpers import annotate_issue_locations
        return annotate_issue_locations(issue, code, code_lines)

    def _build_adjacent_context(
        self,
        slide_id: int,
        cross_slide_issues: list[Issue],
        codegen_compiler,
    ) -> str:
        """Build adjacent slide context for cross-slide issues."""
        adjacent_ids = set()
        for issue in cross_slide_issues:
            if hasattr(issue, 'affected_slides') and issue.affected_slides:
                for sid in issue.affected_slides:
                    if sid != slide_id:
                        adjacent_ids.add(sid)

        adjacent_ids.add(slide_id - 1)
        adjacent_ids.add(slide_id + 1)
        adjacent_ids.discard(0)

        if not adjacent_ids:
            return ""

        parts = ["\n## Adjacent Slides (avoid topic overlap with neighbors)\n"]
        found_any = False

        for adj_id in sorted(adjacent_ids):
            adj_code = codegen_compiler.slide_codes.get(adj_id)
            if not adj_code:
                continue

            # Extract title and key text from adjacent slides
            # Support both python-pptx (.text =) and HTML code
            texts = re.findall(
                r'\.text\s*=\s*[f]?["\'](.+?)["\']', adj_code,
            )
            # HTML: extract title from <h1>/<h2> and first few text elements
            if not texts:
                titles = re.findall(r'<h[12][^>]*>([^<]+)</h[12]>', adj_code)
                paragraphs = re.findall(r'<(?:p|li|span)[^>]*>([^<]{5,})</(?:p|li|span)>', adj_code)
                texts = titles + paragraphs[:5]
            if not texts:
                continue

            found_any = True
            parts.append(f"### Slide {adj_id} text:")
            for t in texts[:10]:
                if len(t.strip()) > 3:
                    parts.append(f"  - {t[:150]}")
            parts.append("")

        for issue in cross_slide_issues:
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or ""
            )
            fix = issue.planned_fix or ""
            parts.append(
                f"**Cross-slide action required:** {desc[:250]}"
            )
            if fix:
                parts.append(f"-> Fix: {fix[:200]}")
            parts.append("")

        return "\n".join(parts) if found_any else ""

    # ================================================================
    # VIZ SYNTHESIS
    # ================================================================

    def _synthesize_viz(
        self,
        slide_id: int,
        issues: list[Issue],
        bp_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        evidence_text: str,
    ) -> dict | None:
        """Extract structured chart data from evidence tables."""
        data_viz_issue = None
        for issue in issues:
            if issue.issue_type == "missing_data_visualization":
                data_viz_issue = issue
                break

        parts = [
            "Given this evidence data, produce a viz_data JSON for a "
            "PowerPoint chart.\n",
        ]

        if evidence and evidence.tables:
            parts.append("## Evidence Tables\n")
            for table in evidence.tables[:5]:
                parts.append(
                    f"### {table.table_id}: {table.caption}\n"
                    f"{table.content[:2000]}\n"
                )

        if evidence and evidence.numeric_facts:
            parts.append("## Numeric Facts\n")
            for fact in evidence.numeric_facts[:20]:
                parts.append(
                    f"- {fact.context}: {fact.value} {fact.unit or ''}"
                )

        if bp_slide:
            parts.append(f"\n## Slide Goal\n{bp_slide.primary_proposition}")

        if data_viz_issue:
            desc = (
                data_viz_issue.evidence.description
                or data_viz_issue.why_this_fails
                or ""
            )
            parts.append(f"\n## What to Visualize\n{desc}")

        parts.append(
            f"\n## Output\n"
            f"Return ONLY a JSON object:\n"
            f'{{\n'
            f'  "chart_type": "column_clustered",\n'
            f'  "categories": ["Label1", "Label2", ...],\n'
            f'  "series": [\n'
            f'    {{"name": "Series1", "values": [1.0, 2.0, ...]}}\n'
            f'  ]\n'
            f'}}\n\n'
            f"Rules: ONLY numbers from evidence. Max 8 categories, 4 series."
        )

        try:
            response = self.llm.call_text(
                system_prompt=(
                    "You extract structured chart data from evidence text. "
                    "Return ONLY valid JSON."
                ),
                user_content="\n".join(parts),
                model=self.model,
                module_name="agent_repair_synth_viz",
                prompt_version="synth_viz.v1",
                max_tokens=1024,
                temperature=0.1,
            )
            return _parse_viz_data(response)
        except Exception as e:
            logger.warning(
                "Synth viz slide %d failed: %s",
                slide_id, str(e)[:200],
            )
            return None

    # ================================================================
    # VERIFICATION UTILITIES
    # ================================================================
    # EVIDENCE CONTEXT
    # ================================================================

    @staticmethod
    def _build_full_evidence_context(
        bp_slide,
        evidence,
        case_dir: str,
    ) -> str:
        """Build full source evidence context matching the evaluator.

        The evaluator sees up to 48K chars of source material. The
        repair agent must see the same information to make correct
        content edits. This replaces the old per-slide evidence window
        that only included ~20K chars of linked chunks.
        """
        if not evidence:
            return ""

        parts = []
        budget = SpatialThresholds.SOURCE_BUDGET_CHARS
        used = 0

        # Pass 1: if we have a blueprint slide, prioritize its linked
        # evidence to put the most relevant material first.
        linked_ids = set()
        if bp_slide:
            linked_ids = set(bp_slide.linked_evidence_ids or [])
            for chunk in evidence.chunks:
                if chunk.chunk_id in linked_ids:
                    section = chunk.metadata.get(
                        "section",
                        chunk.metadata.get("heading", ""),
                    )
                    header = (
                        f"[{chunk.chunk_id}] {section}"
                        if section
                        else f"[{chunk.chunk_id}]"
                    )
                    content = chunk.content.strip()
                    max_chunk = min(4000, budget - used - len(header) - 10)
                    if max_chunk <= 100:
                        break
                    if len(content) > max_chunk:
                        content = content[:max_chunk] + "..."
                    entry = f"### {header}\n{content}"
                    parts.append(entry)
                    used += len(entry)

        # Pass 2: fill remaining budget with all other chunks
        for chunk in evidence.chunks:
            if chunk.chunk_id in linked_ids:
                continue
            section = chunk.metadata.get(
                "section",
                chunk.metadata.get("heading", ""),
            )
            header = (
                f"[{chunk.chunk_id}] {section}"
                if section
                else f"[{chunk.chunk_id}]"
            )
            content = chunk.content.strip()
            max_chunk = min(4000, budget - used - len(header) - 10)
            if max_chunk <= 100:
                break
            if len(content) > max_chunk:
                content = content[:max_chunk] + "..."
            entry = f"### {header}\n{content}"
            parts.append(entry)
            used += len(entry)

        # Pass 3: tables
        if evidence.tables and used < budget - 500:
            parts.append("\n## Tables")
            for tbl in evidence.tables[:8]:
                cap = tbl.caption or tbl.table_id
                desc = f" — {tbl.description}" if tbl.description else ""
                tbl_content = (tbl.content or "")[:600]
                entry = f"[{tbl.table_id}] {cap}{desc}\n{tbl_content}"
                if used + len(entry) > budget:
                    break
                parts.append(entry)
                used += len(entry)

        if parts:
            return (
                "\n\n## Source Evidence (use ONLY data from here, "
                "do NOT invent numbers)\n\n"
                + "\n\n".join(parts)
            )
        return ""

    # ================================================================

    @staticmethod
    def _search_paper_for_claims(
        all_issues: list[Issue], case_dir: str,
    ) -> str:
        """Search the full paper for claims flagged as fabricated/numeric_error.

        When the per-slide evidence window doesn't include a passage, the
        evaluator may incorrectly flag a real paper claim as fabricated.
        This method searches the full paper text for key terms from such
        issues and returns matching excerpts so the repair agent can make
        an informed decision about whether to keep or remove the content.
        """
        # Find issues that question specific numbers or claims
        suspect_types = {
            "fabricated", "numeric_error", "entity_error",
            "chart_misinterpretation", "incorrect_claim",
        }
        search_terms: list[str] = []
        for issue in all_issues:
            if issue.issue_type not in suspect_types:
                continue
            desc = (
                issue.evidence.description
                or issue.why_this_fails
                or ""
            )
            # Extract numbers and short quoted phrases from the description
            numbers = re.findall(r'\d+\.?\d*', desc)
            # Only use distinctive numbers (not 1, 2, 3, etc.)
            for n in numbers:
                if len(n) >= 3 or '.' in n:
                    search_terms.append(n)
            # Extract quoted strings
            for q in re.findall(r'"([^"]{5,40})"', desc):
                search_terms.append(q)

        if not search_terms:
            return ""

        # Find the paper full text
        case_path = Path(case_dir)
        paper_paths = list(case_path.glob("source_pack/paper_full.md"))
        if not paper_paths:
            paper_paths = list(case_path.glob("source_pack/*.md"))
        if not paper_paths:
            return ""

        try:
            paper_text = paper_paths[0].read_text(errors="replace")
        except Exception:
            return ""

        # Search for each term and collect surrounding context
        results: list[str] = []
        seen_positions: set[int] = set()
        for term in search_terms[:15]:  # cap to avoid huge output
            pos = paper_text.lower().find(term.lower())
            if pos < 0:
                continue
            # Deduplicate nearby matches
            if any(abs(pos - s) < 200 for s in seen_positions):
                continue
            seen_positions.add(pos)
            # Extract context window
            start = max(0, pos - 150)
            end = min(len(paper_text), pos + len(term) + 150)
            excerpt = paper_text[start:end].replace("\n", " ").strip()
            results.append(f"  Found \"{term}\" → ...{excerpt}...")

        if not results:
            return ""

        return (
            "\n\n## Paper Search Results\n"
            "The following terms from content-accuracy issues were found "
            "in the full paper. If a claim was flagged as fabricated but "
            "appears here, it may be valid — do NOT delete it blindly.\n\n"
            + "\n".join(results[:10])
        )

    def _test_compile(
        self, code: str, codegen_compiler, case_dir: str, slide_id: int,
    ) -> bool:
        """Test-compile code to check for errors."""
        # HTML mode: just validate basic structure
        if self._is_html_code(code):
            return bool(code.strip())

        try:
            from pptx import Presentation
            from pptx.util import Emu

            test_prs = Presentation()
            test_prs.slide_width = Emu(SlideDimensions.WIDTH_EMU)
            test_prs.slide_height = Emu(SlideDimensions.HEIGHT_EMU)
            test_slide = test_prs.slides.add_slide(test_prs.slide_layouts[6])

            case_path = Path(case_dir)
            test_image_dir = str(
                codegen_compiler._find_image_dir(case_path)
                or case_path / "images"
            )

            from ...backends.python_pptx import code_executor
            success, error = code_executor.execute_code(
                code, test_prs, test_slide, test_image_dir,
            )
            if not success:
                logger.debug(
                    "Compile slide %d: %s", slide_id, error[:300],
                )
            return success
        except Exception as e:
            logger.debug(
                "Compile slide %d exception: %s",
                slide_id, str(e)[:200],
            )
            return False

    def _get_retention_threshold(self, all_issues: list[Issue]) -> float:
        """Determine content retention threshold based on issue types."""
        issue_types = {i.issue_type for i in all_issues}

        if "missing_data_visualization" in issue_types:
            return 0.10  # Major structural changes expected

        # Critical content accuracy issues require full content rewrite —
        # fabricated/incorrect content MUST be replaced, so low retention
        # is expected and correct behavior.
        critical_content = issue_types & {
            "fabricated", "incorrect_claim", "unfaithful_compression",
        }
        if len(critical_content) >= 2:
            return 0.03  # Near-total rewrite expected
        if critical_content:
            return 0.05  # Significant rewrite expected

        if issue_types & STRUCTURAL_ISSUE_TYPES:
            return 0.12  # Structural changes

        return 0.20  # Text-only fixes

    def _check_content_retention(
        self, original_code: str, new_code: str,
        threshold: float = 0.20,
    ) -> bool:
        """Ensure repaired slide retains key content terms."""
        is_html = self._is_html_code(original_code)

        if is_html:
            # For HTML: extract visible text by stripping tags
            original_text = self._extract_visible_text(original_code)
            new_text = self._extract_visible_text(new_code)
            original_texts = [original_text]
            new_texts = [new_text]
        else:
            original_texts = re.findall(
                r'\.text\s*=\s*[f]?["\'](.+?)["\']', original_code,
            )
            new_texts = re.findall(
                r'\.text\s*=\s*[f]?["\'](.+?)["\']', new_code,
            )

        original_terms = set()
        for text in original_texts:
            for word in text.split():
                cleaned = re.sub(r'[^a-zA-Z0-9]', '', word).lower()
                if len(cleaned) > 3:
                    original_terms.add(cleaned)

        new_terms = set()
        for text in new_texts:
            for word in text.split():
                cleaned = re.sub(r'[^a-zA-Z0-9]', '', word).lower()
                if len(cleaned) > 3:
                    new_terms.add(cleaned)

        if not original_terms:
            self._last_retention = 1.0
            return True

        retention = len(original_terms & new_terms) / len(original_terms)
        self._last_retention = retention
        if retention < threshold:
            logger.debug(
                "Content retention: %.0f%% (%d/%d terms, threshold=%.0f%%)",
                retention * 100,
                len(original_terms & new_terms),
                len(original_terms),
                threshold * 100,
            )
            return False

        return True


# ================================================================
# Module-level helpers
# ================================================================


def _delete_shape_code(code: str, var_name: str) -> str:
    """Delete all code lines belonging to a shape variable.

    Finds the line where var_name is assigned (typically from
    slide.shapes.add_*), then deletes everything from that line
    until the next shape assignment or end of function.
    """
    lines = code.split("\n")

    # Find the assignment line for this variable
    assignment_pattern = re.compile(
        rf'^\s+{re.escape(var_name)}\s*=\s*slide\.shapes\.add_'
    )
    start_idx = None
    for i, line in enumerate(lines):
        if assignment_pattern.match(line):
            start_idx = i
            break

    if start_idx is None:
        return code  # Variable not found

    # Find the end: next shape assignment, next function/class def,
    # or a blank line followed by a comment block for the next section
    end_idx = start_idx + 1
    indent = len(lines[start_idx]) - len(lines[start_idx].lstrip())

    while end_idx < len(lines):
        line = lines[end_idx]
        stripped = line.strip()

        # Skip blank lines
        if not stripped:
            end_idx += 1
            continue

        # Check if this line is at the same or lower indentation
        # and starts a new shape or section
        line_indent = len(line) - len(line.lstrip())
        if line_indent <= indent:
            # New top-level statement at same or lower indent — stop
            if re.match(r'\s*(#|def |class |return )', line):
                break
            # Another variable assignment at same indent level
            if re.match(r'\s+\w+\s*=\s*slide\.shapes\.add_', line):
                break
            # A variable assignment that looks like a new section
            if re.match(r'\s+\w+\s*=\s*\w', line) and line_indent == indent:
                # Check if it references our var_name — if so, keep going
                if var_name not in stripped:
                    break

        # Lines that reference var_name — include in deletion
        # Lines that are indented deeper — include in deletion
        # (they're setting properties on our shape)
        end_idx += 1

    # Now find all lines that reference var_name after end_idx
    # and mark them for deletion too
    ref_pattern = re.compile(rf'\b{re.escape(var_name)}\b')
    extra_deletions = set()
    for i in range(end_idx, len(lines)):
        if ref_pattern.search(lines[i]):
            extra_deletions.add(i)

    # Build new code without the deleted lines
    keep = []
    for i, line in enumerate(lines):
        if start_idx <= i < end_idx:
            continue
        if i in extra_deletions:
            continue
        keep.append(line)

    return "\n".join(keep)
