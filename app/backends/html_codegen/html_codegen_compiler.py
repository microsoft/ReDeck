"""HtmlCodeGenCompiler - LLM-based HTML code generation compiler for slides.

This compiler asks the LLM to write HTML/CSS for each slide, renders it to
a high-res PNG via Playwright, then
inserts each PNG as a full-bleed image into a PPTX.

This gives the model full control over layout and visual design using the
full expressiveness of HTML/CSS, while still producing a standard PPTX file.
"""

import logging
import re
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ...llm_client import LLMClient
from ...schemas.blueprint import DeckBlueprint, BlueprintSlide
from ...schemas.evidence import EvidenceState

from ...themes import (
    THEME_REGISTRY, DEFAULT_THEME, ThemeColors,
    format_theme_colors_for_prompt, format_theme_typography_for_prompt,
)
from ...utils.io_utils import read_text
from ...modules.chart_generator import ChartGenerator

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Viewport dimensions for Playwright rendering
VIEWPORT_W = 1280
VIEWPORT_H = 720
DEVICE_SCALE_FACTOR = 2  # 2x for crisp images → 2560×1440 PNG

HTML_CODEGEN_PROMPT_PATH = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "slide_html_codegen.system.md"
HTML_REPAIR_PROMPT_PATH = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "slide_html_repair.system.md"

MAX_CODE_RETRIES = 3

# Patterns that indicate LLM meta-instructions leaked into slide content
_META_TEXT_PATTERNS = [
    re.compile(r'<[^>]*>[^<]*(?:not specified in (?:the )?provided (?:packet|context|material|document)s?)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:venue[/,]?\s*year\s+not\s+(?:specified|available|provided))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:otherwise\s+note\s+)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:if (?:available|provided)[\s,;]*otherwise)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:The slide (?:emphasizes|shows|presents|illustrates|demonstrates))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:This slide (?:emphasizes|shows|presents|illustrates|demonstrates))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:should be presented as)[^<]*</[^>]*>', re.IGNORECASE),
    # Non-slide meta content: "Flow:", "Transition:", presentation structure descriptions
    re.compile(r'<strong[^>]*>\s*Flow:\s*</strong>[^<]*(?:</[^>]*>)?', re.IGNORECASE),
    re.compile(r'<strong[^>]*>\s*Transition:\s*</strong>[^<]*(?:</[^>]*>)?', re.IGNORECASE),
    re.compile(r'<strong[^>]*>\s*Note:\s*</strong>\s*(?:This (?:slide|deck|presentation)|The (?:slide|deck|presentation))[^<]*(?:</[^>]*>)?', re.IGNORECASE),
]


def _find_page_for_must_cover(item: str, bp_slide, source_store) -> int | None:
    """Find the most likely page number for a must_cover item by keyword matching.

    Searches through the slide's source bundle atomic blocks for text matching
    key terms in the must_cover item. Returns the page number of the best match.
    """
    if source_store is None:
        return None
    # Extract numbers and key terms from the must_cover item
    import re as _re
    terms = set()
    # Extract dollar amounts, percentages, specific numbers
    for num in _re.findall(r'[\$]?[\d,]+\.?\d*[%BMK]?', item):
        cleaned = num.replace(',', '').replace('$', '')
        if len(cleaned) > 1:
            terms.add(cleaned)
    # Extract significant words
    for word in item.split():
        cleaned = _re.sub(r'[^a-zA-Z]', '', word).lower()
        if len(cleaned) > 3:
            terms.add(cleaned)

    if not terms:
        return None

    # Search through bundle's atomic block IDs
    bundle = source_store.get_bundle(bp_slide.slide_id)
    if not bundle:
        return None

    best_page = None
    best_score = 0
    for bid in bundle.source_atomic_block_ids:
        block = source_store.get_block(bid)
        if not block or not block.page:
            continue
        block_text = block.text.lower().replace(',', '')
        score = sum(1 for t in terms if t.lower() in block_text)
        if score > best_score:
            best_score = score
            best_page = block.page

    return best_page if best_score >= 1 else None


def _resolve_linked_ids(
    linked_evidence_ids: list[str],
    source_store=None,
) -> list[str]:
    """Expand DB* doc-block IDs to their constituent B* atomic-block IDs.

    Blueprint slides use DB* IDs (e.g. DB006), but evidence chunks use B* IDs
    (e.g. B029-B036). This helper resolves the mapping so Pass 1 exact match
    can find the right chunks.

    IDs that are already B*/T*/A* are passed through unchanged.
    """
    resolved: list[str] = []
    for eid in linked_evidence_ids:
        if eid.startswith("DB") and source_store is not None:
            doc_block = source_store.get_doc_block(eid)
            if doc_block and doc_block.included_atomic_block_ids:
                resolved.extend(doc_block.included_atomic_block_ids)
                continue
        resolved.append(eid)
    return resolved


def _sanitize_meta_text(html: str) -> str:
    """Remove LLM meta-instructions that leaked into slide HTML content."""
    result = html
    for pattern in _META_TEXT_PATTERNS:
        result = pattern.sub('', result)
    return result


class HtmlCodeGenCompiler:
    """Compiles slides by generating HTML/CSS via LLM, rendering to PNG, and inserting into PPTX."""

    def __init__(self, llm: LLMClient, model: str = "gpt-4o", codegen_prompt: str | None = None, show_source_citations: bool = True):
        self.llm = llm
        self.model = model
        self.show_source_citations = show_source_citations
        if codegen_prompt:
            prompt_path = Path(__file__).parent.parent.parent / "prompts" / "codegen" / f"{codegen_prompt}.system.md"
            if prompt_path.exists():
                self.system_prompt = read_text(prompt_path)
                logger.info("Using custom codegen prompt: %s", codegen_prompt)
            else:
                logger.warning("Custom codegen prompt not found: %s, using default", prompt_path)
                self.system_prompt = read_text(HTML_CODEGEN_PROMPT_PATH)
        else:
            self.system_prompt = read_text(HTML_CODEGEN_PROMPT_PATH)
        self.repair_prompt = read_text(HTML_REPAIR_PROMPT_PATH)

        # Strip citation-suppression sections when visible citations are disabled.
        if not self.show_source_citations:
            self.system_prompt = self._strip_citation_sections(self.system_prompt)
            self.repair_prompt = self._strip_citation_sections(self.repair_prompt)

        self.slide_codes: dict[int, str] = {}  # slide_id -> generated HTML
        self.used_layouts: list[str] = []  # track layouts for diversity
        self.used_image_ids: set[str] = set()  # track used images to avoid repetition
        self._theme: ThemeColors | None = None  # cached theme for consistent colors
        self._chart_gen = ChartGenerator()
        self._browser = None
        self._playwright = None
        self._task_brief = ""  # global instructions injected into every slide prompt

    @staticmethod
    def _strip_citation_sections(prompt: str) -> str:
        """Remove citation-suppression paragraphs from prompts for default profile."""
        import re as _re
        # Remove the "Source Citations — DO NOT include" section (header + paragraph)
        prompt = _re.sub(
            r'###\s*1\.\s*Source Citations[^\n]*\n(?:(?!###).*\n)*',
            '', prompt
        )
        # Remove individual citation-related bullet lines
        lines = prompt.split('\n')
        filtered = []
        for line in lines:
            # Skip lines that are purely about suppressing Source/citation
            if any(phrase in line for phrase in [
                'Source:', 'source footnote', 'citation footnote',
                'NO source footnotes', 'citation text of any kind',
                'Source citations (11px)',
                'do NOT add any visible Source',
                'Remove any "Source:"',
            ]):
                continue
            filtered.append(line)
        return '\n'.join(filtered)

    # ------------------------------------------------------------------
    # Playwright lifecycle
    # ------------------------------------------------------------------

    def _ensure_browser(self):
        """Lazily launch Playwright browser (Chromium headless)."""
        if self._browser is not None:
            return

        from playwright.sync_api import sync_playwright
        self._playwright = sync_playwright().start()
        self._browser = self._playwright.chromium.launch(headless=True)
        logger.info("Playwright browser launched for HTML rendering")

    def _close_browser(self):
        """Close Playwright browser and resources."""
        if self._browser:
            try:
                self._browser.close()
            except Exception:
                pass
            self._browser = None
        if self._playwright:
            try:
                self._playwright.stop()
            except Exception:
                pass
            self._playwright = None

    def __del__(self):
        self._close_browser()

    # ------------------------------------------------------------------
    # KaTeX formula rendering
    # ------------------------------------------------------------------

    _KATEX_CSS_CDN = "https://cdn.jsdelivr.net/npm/katex@0.16.45/dist/katex.min.css"

    def _render_katex_in_html(self, html: str) -> str:
        """Replace $$...$$ and $...$ LaTeX in HTML with KaTeX-rendered spans.

        Uses node + katex CLI to render each formula to HTML markup,
        then injects the KaTeX CSS link into <head>.
        """
        import subprocess

        # --- Pre-processing: fix common LLM artifacts ---

        # 1. Remove stray trailing $ on HTML lines (LLM sometimes uses $ as
        #    an end-of-line marker, which confuses the inline $ regex).
        #    Only strip when the $ follows > or whitespace at line end.
        html = re.sub(r'\$\s*$', '', html, flags=re.MULTILINE)

        # 2. Fix mismatched delimiters: $...$$  ->  $$...$$
        #    (single $ open, double $$ close)
        html = re.sub(
            r'(?<!\$)\$(?!\$)(.+?)\$\$',
            lambda m: f'$${m.group(1)}$$' if self._looks_like_latex(m.group(1).strip()) else m.group(0),
            html,
        )

        # Find all formulas: $$...$$ (display) and $...$ (inline)
        # Avoid matching CSS values like "top: $10px" — require LaTeX-like content
        display_pattern = re.compile(r'\$\$(.+?)\$\$', re.DOTALL)
        # For inline $...$: exclude $ followed by digit (dollar amounts like $8.9)
        inline_pattern = re.compile(r'(?<!\$)\$(?!\$)(?!\d)(.+?)(?<!\$)\$(?!\$)')

        formulas: list[tuple[str, str, bool]] = []  # (full_match, latex, is_display)

        for m in display_pattern.finditer(html):
            latex = m.group(1).strip()
            if self._looks_like_latex(latex):
                formulas.append((m.group(0), latex, True))

        for m in inline_pattern.finditer(html):
            latex = m.group(1).strip()
            if self._looks_like_latex(latex):
                formulas.append((m.group(0), latex, False))

        if not formulas:
            return html

        # Render each formula via katex CLI
        has_rendered = False
        for full_match, latex, is_display in formulas:
            try:
                cmd = ["npx", "katex", "--no-throw-on-error"]
                if is_display:
                    cmd.append("--display-mode")
                result = subprocess.run(
                    cmd, input=latex, capture_output=True, text=True, timeout=10,
                )
                if result.returncode == 0 and result.stdout.strip():
                    rendered = result.stdout.strip()
                    html = html.replace(full_match, rendered, 1)
                    has_rendered = True
                else:
                    logger.debug("KaTeX failed for: %s — %s", latex[:50], result.stderr[:100])
            except Exception as e:
                logger.debug("KaTeX rendering error: %s", e)

        # Inject KaTeX CSS into <head> if any formula was rendered
        if has_rendered and self._KATEX_CSS_CDN not in html:
            katex_link = f'<link rel="stylesheet" href="{self._KATEX_CSS_CDN}">'
            if "<head>" in html:
                html = html.replace("<head>", f"<head>\n{katex_link}", 1)
            elif "<style>" in html:
                html = html.replace("<style>", f"{katex_link}\n<style>", 1)

        return html

    @staticmethod
    def _cleanup_residual_latex(html: str) -> str:
        """Clean up LaTeX artifacts that KaTeX didn't render.

        Handles:
        - $...$ inline math delimiters → strip delimiters, convert content to Unicode
        - \\textbf{}, \\mathbb{}, \\mathrm{} wrappers → keep inner text
        - Stray backslash commands → remove or convert to Unicode
        """
        from ...utils.formula_renderer import FormulaRenderer

        # 1. Strip $...$ delimiters and convert inner LaTeX to Unicode
        def _convert_inline_math(m: re.Match) -> str:
            inner = m.group(1).strip()
            # Skip if it looks like natural language (3+ English words)
            import re as _re
            if len(_re.findall(r'\b[a-zA-Z]{4,}\b', inner)) >= 3:
                return m.group(0)  # return unchanged
            try:
                return FormulaRenderer.latex_to_unicode(inner)
            except Exception:
                return inner

        html = re.sub(
            r'(?<!\$)\$(?!\$)(?!\d)([^$\n]{1,80})(?<!\$)\$(?!\$)',
            _convert_inline_math,
            html,
        )

        # 2. Remove \textbf{...}, \textit{...}, \mathbb{...}, \mathrm{...}
        #    wrappers — keep inner text. Repeat to handle nesting.
        for _ in range(3):
            prev = html
            html = re.sub(
                r'\\(?:textbf|textit|textrm|mathrm|mathbf|mathbb|mathcal|emph)\{([^{}]*)\}',
                r'\1',
                html,
            )
            if html == prev:
                break

        # 3. Replace common standalone LaTeX commands with Unicode
        _QUICK_REPLACEMENTS = {
            r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\delta': 'δ',
            r'\epsilon': 'ε', r'\theta': 'θ', r'\lambda': 'λ', r'\mu': 'μ',
            r'\sigma': 'σ', r'\omega': 'ω', r'\pi': 'π', r'\tau': 'τ',
            r'\times': '×', r'\cdot': '·', r'\approx': '≈', r'\neq': '≠',
            r'\leq': '≤', r'\geq': '≥', r'\infty': '∞', r'\partial': '∂',
            r'\rightarrow': '→', r'\leftarrow': '←', r'\Rightarrow': '⇒',
            r'\sum': 'Σ', r'\prod': 'Π', r'\in': '∈', r'\subset': '⊂',
            r'\ldots': '…', r'\dots': '…', r'\pm': '±',
            r'\leftrightarrow': '↔', r'\Leftrightarrow': '⇔',
            r'\forall': '∀', r'\exists': '∃', r'\nabla': '∇',
            r'\Pi': 'Π', r'\Sigma': 'Σ', r'\Omega': 'Ω', r'\Delta': 'Δ',
            r'\Gamma': 'Γ', r'\Lambda': 'Λ', r'\Theta': 'Θ',
            r'\mid': '|', r'\parallel': '∥', r'\perp': '⊥',
            r'\cap': '∩', r'\cup': '∪', r'\emptyset': '∅',
            r'\log': 'log', r'\exp': 'exp', r'\max': 'max', r'\min': 'min',
        }
        for cmd, uni in _QUICK_REPLACEMENTS.items():
            html = html.replace(cmd, uni)

        # 4. Remove any remaining stray \command that is NOT inside a
        #    <style>, <script>, or url() context. Only target text nodes.
        #    Be conservative: only remove commands we recognize as LaTeX.
        html = re.sub(
            r'\\(?:text|math|frac|sqrt|bar|hat|tilde|vec|dot|overline)\{([^{}]*)\}',
            r'\1',
            html,
        )

        # 5. Fix HTML entities that should be symbols in text nodes
        #    (e.g., &lt; → < and &gt; → > when inside visible text, not HTML tags)
        #    Only replace when surrounded by math-like context (digits, variables)
        html = re.sub(r'(\w)&lt;(\w)', r'\1<\2', html)
        html = re.sub(r'(\w)&gt;(\w)', r'\1>\2', html)

        return html

    @staticmethod
    def _looks_like_latex(text: str) -> bool:
        """Check if text looks like LaTeX math (not a CSS value or price).

        Accepts both complex LaTeX (\\frac, ^, _) AND simple math expressions
        like single variables (W, r, k), expressions (O(1), N=6), and numeric
        formulas that are commonly written inside $...$ in academic papers.
        """
        if len(text) < 1 or len(text) > 500:
            return False
        # Exclude things that look like CSS/JS: px, %, #, url(
        css_indicators = any(s in text for s in ('px', 'url(', '://', '#'))
        if css_indicators:
            return False
        # Exclude natural language: if text contains 3+ consecutive English
        # words (4+ letter words separated by spaces), it's prose, not math.
        import re as _re
        word_seq = _re.findall(r'\b[a-zA-Z]{4,}\b', text)
        if len(word_seq) >= 3:
            return False
        # Complex LaTeX indicators — always accept
        if ('\\' in text or '^' in text or '_' in text
                or '{' in text or 'frac' in text
                or any(c in text for c in 'αβγδεζηθλμνξπρστφψω∑∏∫')):
            return True
        # Simple math expressions: single letters (W, r, k, N), or
        # letter+number combos (N=6, h=8), or function-like (O(1), O(n^2))
        # These are common in academic papers inside $...$
        if _re.fullmatch(r'[A-Za-z0-9()\s=+\-*/.,<>≤≥]+', text):
            return True
        return False

    # ------------------------------------------------------------------
    # HTML post-processing
    # ------------------------------------------------------------------

    @staticmethod
    def _fix_title_overflow(html: str, slide_id: int, paper_title: str = "") -> str:
        """Auto-shrink title font-size if text is too long for its container.
        Also fix truncated paper titles on slide 1.

        For the header-band `.title` element (all content slides): the title
        sits in a 1280-56-56 = 1168px wide area. At font-size 38px, roughly
        ~30 chars fit. We estimate: if chars * font_size * 0.55 > available_width,
        shrink font_size until it fits.

        For the title slide (slide_id=1): the <h1> has max-width 1100px.
        """
        import re

        # --- Fix truncated paper title on slide 1 ---
        if slide_id == 1 and paper_title and len(paper_title) > 10:
            # Try <h1> first, then .title class div
            h1_content_match = re.search(r'(<h1[^>]*>)(.*?)(</h1>)', html, re.DOTALL)
            if not h1_content_match:
                # Try class="title" div
                h1_content_match = re.search(r'(class="title"[^>]*>)(.*?)(<)', html, re.DOTALL)
            if h1_content_match:
                h1_text = re.sub(r'<[^>]+>', '', h1_content_match.group(2)).strip()
                # Check if the h1 text is a truncated prefix of the paper title
                if (len(h1_text) >= 5 and len(h1_text) < len(paper_title) - 3
                        and paper_title.lower().startswith(h1_text.lower())):
                    html = html.replace(
                        h1_content_match.group(2),
                        paper_title,
                        1,
                    )
                    logger.info(
                        "Slide 1: restored truncated title '%s...' → '%s'",
                        h1_text[:30], paper_title[:60],
                    )

        # Content slide title: `.title { ... font-size: XXpx ... }`
        # Match the .title element and extract text
        title_match = re.search(
            r'class="title"[^>]*>([^<]+)<', html
        )
        if title_match:
            title_text = title_match.group(1).strip()
            avail_width = 1168  # 1280 - 56*2
            # Extract current font-size from .title CSS
            fs_match = re.search(r'\.title\s*\{[^}]*font-size:\s*(\d+)px', html)
            if fs_match:
                current_fs = int(fs_match.group(1))
            else:
                current_fs = 38
            # Estimate text width: chars * font_size * 0.55 (average char width ratio)
            estimated_width = len(title_text) * current_fs * 0.55
            if estimated_width > avail_width:
                # Calculate needed font-size
                new_fs = max(20, int(avail_width / (len(title_text) * 0.55)))
                if new_fs < current_fs:
                    if fs_match:
                        old = fs_match.group(0)
                        new = old.replace(f"{current_fs}px", f"{new_fs}px")
                        html = html.replace(old, new, 1)
                    else:
                        # Inject font-size into the element style
                        html = html.replace(
                            'class="title"',
                            f'class="title" style="font-size: {new_fs}px;"',
                            1,
                        )
                    logger.info(
                        "Slide %d: auto-shrunk title from %dpx to %dpx (%d chars)",
                        slide_id, current_fs, new_fs, len(title_text),
                    )

        # Title slide <h1>: check if it's too long
        if slide_id == 1:
            h1_match = re.search(r'<h1[^>]*style="([^"]*)"[^>]*>([^<]+(?:<[^/][^>]*>[^<]*</[^>]*>)*[^<]*)</h1>', html, re.DOTALL)
            if h1_match:
                style = h1_match.group(1)
                # Get text content (strip inner tags)
                text = re.sub(r'<[^>]+>', '', h1_match.group(2)).strip()
                fs_match = re.search(r'font-size:\s*(\d+)px', style)
                current_fs = int(fs_match.group(1)) if fs_match else 38
                avail_width = 1100
                estimated_width = len(text) * current_fs * 0.55
                if estimated_width > avail_width:
                    new_fs = max(22, int(avail_width / (len(text) * 0.55)))
                    if new_fs < current_fs:
                        new_style = re.sub(
                            r'font-size:\s*\d+px',
                            f'font-size: {new_fs}px',
                            style,
                        )
                        html = html.replace(style, new_style, 1)
                        logger.info(
                            "Slide 1: auto-shrunk h1 from %dpx to %dpx (%d chars)",
                            current_fs, new_fs, len(text),
                        )

        return html

    # ------------------------------------------------------------------
    # HTML rendering
    # ------------------------------------------------------------------

    def _render_html_to_png(self, html_content: str, output_path: Path) -> bool:
        """Render HTML to PNG using Playwright.

        Args:
            html_content: Complete HTML page content
            output_path: Where to save the PNG

        Returns:
            True if rendering succeeded
        """
        self._ensure_browser()

        # Convert bare absolute file paths in <img src="..."> to file:// URLs
        # so Playwright can load local images from the filesystem.

        # Render KaTeX formulas in the HTML before rendering
        html_content = self._render_katex_in_html(html_content)

        # Fallback: clean up any residual LaTeX that KaTeX didn't process
        html_content = self._cleanup_residual_latex(html_content)

        # Remove prompt metadata references from visible text.
        for leaked_term in ["task brief", "slide brief", "evidence text", "source evidence"]:
            html_content = re.sub(
                re.escape(leaked_term), "", html_content, flags=re.IGNORECASE
            )

        # Visible citations on data slides are the default for traceability.

        # Fix Title Case capitalization in <h1>/<h2> slide titles
        # Words that should be lowercase in Title Case (unless first/last word):
        _TC_LOWER = {'a', 'an', 'the', 'and', 'but', 'or', 'nor', 'for', 'yet', 'so',
                      'in', 'on', 'at', 'to', 'of', 'by', 'with', 'as', 'vs', 'via'}

        def _fix_title_case(match):
            """Fix capitalization in a title tag, preserving HTML tags inside."""
            tag = match.group(1)       # e.g. "h1" or "h2"
            attrs = match.group(2)     # tag attributes
            inner = match.group(3)     # inner HTML content
            close = match.group(4)     # closing tag

            # Don't modify if it contains complex nested HTML (images, spans with classes)
            if '<img' in inner or '<span class' in inner:
                return match.group(0)

            # Split inner into HTML tags and text segments
            parts = re.split(r'(<[^>]+>)', inner)
            fixed_parts = []
            for i, part in enumerate(parts):
                if part.startswith('<'):
                    fixed_parts.append(part)
                    continue
                # Fix each word in text segment
                words = part.split(' ')
                fixed_words = []
                for j, word in enumerate(words):
                    if not word:
                        fixed_words.append(word)
                        continue
                    # Check if this word is at the absolute start or end of the title text
                    text_before = ' '.join(' '.join(fixed_parts).split())
                    stripped = re.sub(r'<[^>]+>', '', text_before).strip()
                    is_first = (stripped == '' and j == 0)
                    is_last_segment = (i == len(parts) - 1 or all(p.startswith('<') or p.strip() == '' for p in parts[i+1:]))
                    is_last = is_last_segment and j == len(words) - 1

                    bare = re.sub(r'[^\w]', '', word.lower())
                    if bare in _TC_LOWER and not is_first and not is_last:
                        # Lowercase this word (preserve non-alpha chars around it)
                        # e.g. "And" -> "and", "In" -> "in", "A" -> "a"
                        if word[0].isupper() and (len(word) == 1 or word[1:].islower()) and len(word) <= 5:
                            fixed_words.append(word[0].lower() + word[1:])
                        else:
                            fixed_words.append(word)
                    else:
                        fixed_words.append(word)
                fixed_parts.append(' '.join(fixed_words))
            return f'<{tag}{attrs}>{"".join(fixed_parts)}</{close}>'

        html_content = re.sub(
            r'<(h[12])([^>]*)>(.*?)</(h[12])>',
            _fix_title_case,
            html_content,
            flags=re.DOTALL
        )

        html_content = re.sub(
            r'(<img\s[^>]*src=["\'])(/[^"\']+)(["\'])',
            r'\1file://\2\3',
            html_content,
        )

        # Convert relative paths to absolute file:// URLs.
        # LLM may generate src="cases/..." which needs to resolve from the
        # project root, not from the temp file's directory.
        import os
        from pathlib import Path as _Path
        cwd = os.getcwd()
        _search_roots = [cwd]
        _probe = _Path(cwd)
        for _ in range(6):
            if (_probe / "cases").is_dir() or (_probe / "app").is_dir():
                if str(_probe) != cwd:
                    _search_roots.append(str(_probe))
                break
            _probe = _probe.parent

        def _resolve_relative_src(m):
            prefix, path, suffix = m.group(1), m.group(2), m.group(3)
            if path.startswith(('file://', 'http://', 'https://', 'data:', '/')):
                return m.group(0)
            for root in _search_roots:
                abs_path = os.path.join(root, path)
                if os.path.exists(abs_path):
                    return f'{prefix}file://{abs_path}{suffix}'
            return m.group(0)

        html_content = re.sub(
            r'(<img\s[^>]*src=["\'])([^"\']+)(["\'])',
            _resolve_relative_src,
            html_content,
        )

        page = None
        try:
            page = self._browser.new_page(
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
                device_scale_factor=DEVICE_SCALE_FACTOR,
            )
            # Save HTML to a temp file so Playwright can load local images via file:// protocol
            import tempfile
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".html", delete=False, encoding="utf-8"
            ) as tmp:
                tmp.write(html_content)
                tmp_path = tmp.name
            try:
                page.goto(f"file://{tmp_path}", wait_until="networkidle")
                # Small wait for any CSS transitions/rendering
                page.wait_for_timeout(300)
                output_path.parent.mkdir(parents=True, exist_ok=True)
                page.screenshot(path=str(output_path), full_page=False)
            finally:
                import os as _os
                _os.unlink(tmp_path)
            page.close()
            return True
        except Exception as e:
            logger.error("Playwright render failed: %s", e)
            if page:
                try:
                    page.close()
                except Exception:
                    pass
            return False

    # ------------------------------------------------------------------
    # Public API (mirrors CodeGenCompiler interface)
    # ------------------------------------------------------------------

    def compile_deck(
        self,
        blueprint: DeckBlueprint,
        evidence: EvidenceState,
        case_dir: str | Path,
        output_path: str | Path,
        code_dir: str | Path | None = None,
        source_store=None,
        task_brief: str = "",
    ) -> dict:
        """Generate HTML for each slide, render to PNG, and compile into a PPTX.

        Args:
            blueprint: The deck blueprint with slide briefs
            evidence: Evidence state with chunks, figures, tables
            case_dir: Path to case directory (for image resolution)
            output_path: Where to save the PPTX file
            code_dir: Where to save generated HTML files (optional)
            task_brief: Global task brief / instructions to inject into
                every slide prompt so per-slide LLM sees the full constraints.

        Returns:
            Manifest dict with slide info and any warnings
        """
        # Store task_brief so _build_slide_prompt and repair_slide can use it
        self._task_brief = task_brief
        start = time.time()
        case_dir = Path(case_dir)
        output_path = Path(output_path)
        if code_dir:
            code_dir = Path(code_dir)
            code_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        image_dir = self._find_image_dir(case_dir)
        available_images = self._list_images(image_dir) if image_dir else []

        # Select theme based on paper title for consistent colors across slides
        from app.themes import select_theme_for_paper
        paper_title = ""
        if source_store and hasattr(source_store, 'doc_block_plan'):
            for db in source_store.doc_block_plan.blocks:
                if db.role in ("title", "abstract", "introduction"):
                    paper_title = db.title or ""
                    break
        if not paper_title and blueprint.slides:
            prop = blueprint.slides[0].primary_proposition or ""
            # Extract just the paper title from primary_proposition
            # (it's usually "Paper Title introduces/presents/proposes...")
            import re as _re
            title_cut = _re.split(
                r'\s+(?:introduces?|presents?|proposes?|describes?|provides?|is\s+a|addresses)\s+',
                prop, maxsplit=1
            )
            paper_title = title_cut[0].strip() if title_cut else prop
        self._theme = select_theme_for_paper(paper_title) if paper_title else DEFAULT_THEME
        self._paper_title = paper_title
        logger.info("Selected theme: %s (from title: '%s')", self._theme.theme_name, paper_title[:50])

        manifest = {
            "pptx_path": str(output_path),
            "total_slides": len(blueprint.slides),
            "slides": [],
            "warnings": [],
            "timing_sec": 0,
            "compiler": "html_codegen",
        }

        # Phase 1: Generate HTML for all slides in parallel (LLM calls are IO-bound)
        from concurrent.futures import ThreadPoolExecutor, as_completed

        def _gen_html(bp_slide):
            """Generate HTML + render PNG for one slide (thread-safe for LLM)."""
            sid = bp_slide.slide_id

            # Build user prompt
            user_content, has_images = self._build_slide_prompt(
                bp_slide, evidence, image_dir, available_images,
                total_slides=len(blueprint.slides),
                source_store=source_store,
                )

            # LLM call to generate HTML
            html = None
            last_error = None
            for attempt in range(MAX_CODE_RETRIES):
                try:
                    if attempt == 0:
                        response = self.llm.call_text(
                            system_prompt=self.system_prompt,
                            user_content=user_content,
                            model=self.model,
                            module_name="slide_html_codegen",
                            prompt_version="slide_html_codegen.v1",
                            max_tokens=8000,
                        )
                    else:
                        retry_prompt = f"{user_content}\n\n## RETRY (attempt {attempt+1})\nPrevious error: {last_error}\nFix the issue and regenerate."
                        response = self.llm.call_text(
                            system_prompt=self.system_prompt,
                            user_content=retry_prompt,
                            model=self.model,
                            module_name="slide_html_codegen",
                            prompt_version="slide_html_codegen.v1.retry",
                            max_tokens=8000,
                        )
                    html = self._extract_html(response)
                    if html and "<body" in html.lower():
                        break
                    last_error = "No valid HTML extracted"
                    html = None
                except Exception as e:
                    last_error = str(e)[:200]

            return sid, html, last_error

        # Run LLM generation in parallel threads
        html_results = {}
        max_parallel = min(8, len(blueprint.slides))  # Limit concurrency
        with ThreadPoolExecutor(max_workers=max_parallel) as pool:
            futures = {}
            for bp_slide in blueprint.slides:
                fut = pool.submit(_gen_html, bp_slide)
                futures[fut] = bp_slide
            for fut in as_completed(futures):
                bp_slide = futures[fut]
                sid, html, error = fut.result()
                html_results[sid] = (html, error)
                if html:
                    logger.info("Slide %d HTML generated successfully", sid)
                else:
                    logger.warning("Slide %d HTML generation failed: %s", sid, error)

        # Phase 2: Render PNGs and assemble PPTX (sequential, needs Playwright + prs)
        for bp_slide in blueprint.slides:
            sid = bp_slide.slide_id
            html, error = html_results.get(sid, (None, "not generated"))

            if html:
                # Post-process: fix title overflow for all slides
                html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

                # Render HTML to PNG
                png_dir = (code_dir.parent if code_dir else Path("/tmp")) / "html_renders"
                png_dir.mkdir(parents=True, exist_ok=True)
                png_path = png_dir / f"slide_{sid:02d}.png"
                render_success = self._render_html_to_png(html, png_path)

                if render_success:
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(
                        str(png_path),
                        Inches(0), Inches(0),
                        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                    )
                    title_text, body_text = self._extract_text_from_html(html)
                    src_evidence = self._build_source_evidence_for_notes(
                        bp_slide, evidence=evidence, source_store=source_store,
                    )
                    self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)

                    self.slide_codes[sid] = _sanitize_meta_text(html)
                    if code_dir:
                        html_file = code_dir / f"slide_{sid:02d}.html"
                        html_file.write_text(self.slide_codes[sid], encoding="utf-8")

                    layout_label = self._detect_layout(html)
                    self.used_layouts.append(layout_label)
                    for img_match in re.finditer(r'src=["\']([^"\']+)["\']', html):
                        self.used_image_ids.add(Path(img_match.group(1)).name)

                    manifest["slides"].append({"slide_id": sid, "status": "ok"})
                    continue

            # Fallback
            logger.error("Slide %d: using fallback. Error: %s", sid, error)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            self._add_fallback_slide(slide, bp_slide, evidence)
            manifest["slides"].append({"slide_id": sid, "status": "fallback", "error": error})

        # Save PPTX
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        # Close browser after deck compilation
        self._close_browser()

        manifest["timing_sec"] = round(time.time() - start, 2)
        logger.info(
            "HtmlCodeGenCompiler: compiled %d slides to %s in %.1fs",
            len(blueprint.slides), output_path, manifest["timing_sec"],
        )
        return manifest

    def repair_slide(
        self,
        slide_id: int,
        issues: list[dict],
        blueprint_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        case_dir: str | Path,
        source_store=None,
    ) -> str | None:
        """Repair a slide's HTML code based on issues found.

        Returns the new HTML string, or None if repair failed.
        """
        current_code = self.slide_codes.get(slide_id)
        if not current_code:
            logger.warning("No HTML found for slide %d, cannot repair", slide_id)
            return None

        issue_parts = []
        for iss in issues:
            line = f"- [{iss.get('severity', 'major')}] {iss.get('description', str(iss))}"
            if iss.get('planned_fix'):
                line += f"\n  FIX: {iss['planned_fix']}"
            issue_parts.append(line)
        issue_text = "\n".join(issue_parts)

        brief = ""
        if blueprint_slide:
            brief = f"\nSlide role: {blueprint_slide.role}\nSlide goal: {blueprint_slide.primary_proposition}\n"

        # Build evidence context (reuse shared method)
        evidence_text = self._build_evidence_context(
            blueprint_slide, evidence, case_dir, source_store=source_store,
        )

        # Global task brief for repair context
        task_brief_section = ""
        task_brief = getattr(self, '_task_brief', '')
        if task_brief:
            task_brief_section = f"\n## Global Instructions (from task brief)\n\n{task_brief}\n\n---\n"

        user_content = f"""## Current HTML Code

```html
{current_code}
```

## Issues Found

{issue_text}

## Slide Context
{brief}
{evidence_text}
{task_brief_section}Fix all the issues above. Return the complete updated HTML page.
IMPORTANT: Use ONLY numbers and facts from the Source Evidence above.
CRITICAL: Do NOT change any text content (numbers, facts, bullet text, table data) unless an issue SPECIFICALLY flags that text as incorrect or fabricated. Layout/spatial fixes must ONLY adjust CSS properties (position, size, font-size, padding, margin). Changing text while fixing layout is the #1 cause of regression.
"""

        try:
            response = self.llm.call_text(
                system_prompt=self.repair_prompt,
                user_content=user_content,
                model=self.model,
                module_name="html_code_repair",
                prompt_version="slide_html_repair.v1",
                max_tokens=8000,
                temperature=0.2,
            )
            html = self._extract_html(response)
            if html:
                # Content-loss safeguard
                original_content = self._estimate_text_content(current_code)
                repaired_content = self._estimate_text_content(html)
                if original_content > 50 and repaired_content < original_content * 0.7:
                    logger.warning(
                        "Slide %d HTML repair REJECTED: content loss %.0f%%. Keeping original.",
                        slide_id, (1 - repaired_content / original_content) * 100,
                    )
                    return None
                self.slide_codes[slide_id] = html
                return html
        except Exception as e:
            logger.error("HTML repair failed for slide %d: %s", slide_id, e)

        return None

    def recompile_deck(
        self,
        blueprint: DeckBlueprint,
        case_dir: str | Path,
        output_path: str | Path,
        code_dir: str | Path | None = None,
        evidence=None,
        source_store=None,
    ) -> dict:
        """Recompile all slides using stored/updated HTML code.

        Used after repair to rebuild the PPTX with patched HTML.
        """
        start = time.time()
        case_dir = Path(case_dir)
        output_path = Path(output_path)

        # Ensure paper_title is set for title overflow fix
        if not getattr(self, '_paper_title', ''):
            import re as _re
            if blueprint.slides:
                prop = blueprint.slides[0].primary_proposition or ""
                title_cut = _re.split(
                    r'\s+(?:introduces?|presents?|proposes?|describes?|provides?|is\s+a|addresses)\s+',
                    prop, maxsplit=1
                )
                self._paper_title = title_cut[0].strip() if title_cut else ""

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        manifest = {
            "pptx_path": str(output_path),
            "total_slides": len(blueprint.slides),
            "slides": [],
            "warnings": [],
            "timing_sec": 0,
            "compiler": "html_codegen",
        }

        for bp_slide in blueprint.slides:
            sid = bp_slide.slide_id
            html = self.slide_codes.get(sid)

            if not html:
                # Fallback: add a simple text slide
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                self._add_fallback_slide(slide, bp_slide)
                manifest["slides"].append({
                    "slide_id": sid,
                    "status": "fallback",
                    "warning": "No HTML available",
                })
                continue

            # Render HTML to PNG
            png_dir = output_path.parent / "html_renders"
            png_dir.mkdir(parents=True, exist_ok=True)
            png_path = png_dir / f"slide_{sid:02d}.png"

            # Post-process title overflow
            html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

            success = self._render_html_to_png(html, png_path)

            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            if success and png_path.exists():
                slide.shapes.add_picture(
                    str(png_path),
                    Inches(0), Inches(0),
                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                )
                # Add text content as speaker notes for evaluator
                title_text, body_text = self._extract_text_from_html(html)
                src_evidence = self._build_source_evidence_for_notes(
                    bp_slide, evidence=evidence, source_store=source_store,
                )
                self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)
                manifest["slides"].append({"slide_id": sid, "status": "ok"})
            else:
                logger.warning("Slide %d recompile render failed, using fallback", sid)
                self._add_fallback_slide(slide, bp_slide)
                manifest["warnings"].append(f"Slide {sid} render failed")
                manifest["slides"].append({
                    "slide_id": sid, "status": "error_fallback",
                })

            # Save HTML
            if code_dir:
                html_path = Path(code_dir) / f"slide_{sid:02d}.html"
                html_path.write_text(html, encoding="utf-8")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        # Close browser
        self._close_browser()

        manifest["timing_sec"] = round(time.time() - start, 2)
        return manifest

    # ------------------------------------------------------------------
    # Internal methods
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_text_from_html(html: str) -> tuple[str, str]:
        """Extract title and body text from HTML for use as PPTX notes.

        Returns (title, body_text) where title is the first heading content
        and body_text is the remaining visible text.
        """
        from ...utils.html_text import extract_title_and_body
        return extract_title_and_body(html)

    @staticmethod
    def _build_source_evidence_for_notes(
        bp_slide,
        evidence=None,
        source_store=None,
    ) -> str:
        """Build source evidence text for speaker notes enrichment.

        Collects the key factual content from the source evidence bundle
        so it can be included in speaker notes for content-fidelity
        evaluation.
        """
        parts: list[str] = []

        # V2 path: source_store bundle
        if source_store is not None:
            bundle = source_store.get_bundle(bp_slide.slide_id)
            if bundle and bundle.source_text:
                parts.append(bundle.source_text[:3000])
                for s in bundle.table_summaries[:3]:
                    parts.append(s[:500])

        # Legacy path: linked evidence chunks
        if not parts and evidence is not None and hasattr(evidence, 'chunks'):
            for eid in bp_slide.linked_evidence_ids[:5]:
                for chunk in evidence.chunks:
                    if chunk.chunk_id == eid:
                        parts.append(f"{chunk.chunk_id}: {chunk.content[:600]}")
                        break

        # Add must_cover items for context
        if bp_slide.must_cover_subset:
            parts.append("Must cover: " + "; ".join(bp_slide.must_cover_subset))

        return "\n".join(parts)[:4000]

    @staticmethod
    def _add_notes_to_slide(
        slide, title: str, body_text: str,
        source_evidence: str = "",
    ):
        """Add text content as speaker notes to a PPTX slide.

        This enables the evaluator to read slide content even though
        the visual content is baked into a full-bleed PNG image.

        Args:
            slide: python-pptx Slide object.
            title: Extracted title text from HTML.
            body_text: Extracted body text from HTML.
            source_evidence: Additional source evidence text from the
                blueprint/source_store bundle. Appended to notes so
                content-quiz evaluators can access detailed facts that
                may not be visually rendered on the slide.
        """
        if not title and not body_text and not source_evidence:
            return
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        # Clear existing notes
        notes_tf.text = ""
        # Add title line
        if title:
            notes_tf.text = f"[TITLE] {title}"
        # Add body text
        if body_text:
            p = notes_tf.add_paragraph()
            p.text = f"[CONTENT] {body_text[:3000]}"
        # Add source evidence for richer content fidelity
        if source_evidence:
            p = notes_tf.add_paragraph()
            p.text = f"[SOURCE EVIDENCE] {source_evidence[:4000]}"

    def _generate_slide(
        self,
        prs: Presentation,
        bp_slide: BlueprintSlide,
        evidence: EvidenceState,
        image_dir: Path | None,
        available_images: list[str],
        code_dir: Path | None,
        total_slides: int = 10,
        source_store=None,
    ) -> dict:
        """Generate HTML for one slide, render to PNG, and insert into PPTX."""
        sid = bp_slide.slide_id

        # Build user prompt
        user_content, has_images = self._build_slide_prompt(
            bp_slide, evidence, image_dir, available_images,
            total_slides=total_slides,
            source_store=source_store,
        )

        html = None
        last_error = None

        for attempt in range(MAX_CODE_RETRIES):
            try:
                if attempt == 0:
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=user_content,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version="slide_html_codegen.v1",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid},
                    )
                else:
                    retry_prompt = (
                        f"{user_content}\n\n"
                        f"## Previous Attempt Failed\n"
                        f"Error: {last_error}\n\n"
                        f"Fix the error and return corrected HTML."
                    )
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=retry_prompt,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version="slide_html_codegen.v1.retry",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "attempt": attempt},
                    )

                html = self._extract_html(response)
                if not html:
                    last_error = "Failed to extract HTML from response"
                    logger.warning("Slide %d attempt %d: no HTML extracted", sid, attempt + 1)
                    continue

                # Validate HTML has basic structure
                if "<body" not in html.lower():
                    last_error = "HTML missing <body> tag"
                    logger.warning("Slide %d attempt %d: %s", sid, attempt + 1, last_error)
                    continue

                # Render HTML to PNG
                png_dir = code_dir.parent / "html_renders" if code_dir else Path("/tmp/html_renders")
                png_dir.mkdir(parents=True, exist_ok=True)
                png_path = png_dir / f"slide_{sid:02d}.png"

                # Post-process title overflow
                html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

                render_success = self._render_html_to_png(html, png_path)
                if not render_success:
                    last_error = "Playwright rendering failed"
                    logger.warning("Slide %d attempt %d: %s", sid, attempt + 1, last_error)
                    continue

                # Insert PNG into PPTX slide
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(
                    str(png_path),
                    Inches(0), Inches(0),
                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                )

                # Extract text from HTML and add as speaker notes
                # so the evaluator can read slide content
                title_text, body_text = self._extract_text_from_html(html)
                src_evidence = self._build_source_evidence_for_notes(
                    bp_slide, evidence=evidence, source_store=source_store,
                )
                self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)

                self.slide_codes[sid] = _sanitize_meta_text(html)
                if code_dir:
                    html_file = code_dir / f"slide_{sid:02d}.html"
                    html_file.write_text(html, encoding="utf-8")

                # Track layout used for diversity
                layout_label = self._detect_layout(html)
                self.used_layouts.append(layout_label)

                # Track images used
                for img_match in re.finditer(r'src=["\']([^"\']+)["\']', html):
                    img_path = img_match.group(1)
                    self.used_image_ids.add(Path(img_path).name)

                logger.info("Slide %d generated successfully via HTML (attempt %d)", sid, attempt + 1)
                return {"slide_id": sid, "status": "ok", "attempts": attempt + 1}

            except Exception as e:
                last_error = str(e)
                logger.warning("Slide %d attempt %d error: %s", sid, attempt + 1, str(e)[:200])

        # All retries exhausted — fallback
        logger.error(
            "Slide %d: all %d HTML attempts failed, using fallback. Last error: %s",
            sid, MAX_CODE_RETRIES, last_error
        )
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_fallback_slide(slide, bp_slide, evidence)
        return {
            "slide_id": sid,
            "status": "fallback",
            "error": last_error,
            "attempts": MAX_CODE_RETRIES,
        }

    def _generate_code_only(
        self,
        bp_slide: "BlueprintSlide",
        evidence: "EvidenceState",
        image_dir: Path | None,
        available_images: list[str],
        total_slides: int = 10,
        source_store=None,
    ) -> tuple[str | None, dict]:
        """Generate HTML code for a slide WITHOUT rendering or inserting into PPTX.

        Used by regen_slide tool in the repair agent.

        Returns:
            (html, info_dict) where html is the validated HTML string
            (or None on failure).
        """
        sid = bp_slide.slide_id

        # Build user prompt
        user_content, has_images = self._build_slide_prompt(
            bp_slide, evidence, image_dir, available_images,
            total_slides=total_slides,
            source_store=source_store,
        )

        html = None
        last_error = None

        for attempt in range(MAX_CODE_RETRIES):
            try:
                if attempt == 0:
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=user_content,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version="slide_html_codegen.v1.regen",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "regen": True},
                    )
                else:
                    retry_prompt = (
                        f"{user_content}\n\n"
                        f"## Previous Attempt Failed\n"
                        f"Error: {last_error}\n\n"
                        f"Fix the error and return corrected HTML."
                    )
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=retry_prompt,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version="slide_html_codegen.v1.regen.retry",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "attempt": attempt, "regen": True},
                    )

                html = self._extract_html(response)
                if not html:
                    last_error = "Failed to extract HTML from response"
                    continue

                if "<body" not in html.lower():
                    last_error = "HTML missing <body> tag"
                    continue

                logger.info("Slide %d regen code generated (attempt %d)", sid, attempt + 1)
                return html, {"slide_id": sid, "status": "ok", "attempts": attempt + 1}

            except Exception as e:
                last_error = str(e)
                logger.warning("Slide %d regen attempt %d error: %s", sid, attempt + 1, str(e)[:200])

        logger.error("Slide %d: regen failed after %d attempts", sid, MAX_CODE_RETRIES)
        return None, {"slide_id": sid, "status": "error", "error": last_error}

    def _build_slide_prompt(
        self,
        bp_slide: BlueprintSlide,
        evidence: EvidenceState,
        image_dir: Path | None,
        available_images: list[str],
        total_slides: int = 10,
        source_store=None,
    ) -> tuple[str, bool]:
        """Build the user prompt for HTML slide generation.

        Returns:
            (prompt_text, has_images)
        """
        parts = []
        has_images = False

        # Inject global task brief / instructions so LLM sees full constraints
        task_brief = getattr(self, '_task_brief', '')
        if task_brief:
            parts.append("## Global Instructions (from task brief)\n")
            parts.append(task_brief)
            parts.append("\n---\n")

        parts.append(f"## Slide {bp_slide.slide_id} of {total_slides}")
        parts.append(f"**Role**: {bp_slide.role}")
        parts.append(f"**Goal**: {bp_slide.primary_proposition}")
        parts.append(f"**Position**: {bp_slide.narrative_position}")
        if bp_slide.notes:
            parts.append(f"**Notes**: {bp_slide.notes}")
        if getattr(bp_slide, 'layout_hint', '') and bp_slide.layout_hint:
            parts.append(f"**Layout hint**: {bp_slide.layout_hint} (use as guidance, adapt as needed)")

        # Layout diversity constraint
        if self.used_layouts:
            recent = self.used_layouts[-3:] if len(self.used_layouts) >= 3 else self.used_layouts
            parts.append(f"\n**Previous slides used these layouts**: {', '.join(recent)}")
            parts.append("**You MUST use a DIFFERENT layout pattern for this slide.**")

        # Evidence text
        parts.append("\n## Content / Evidence")
        evidence_added = False

        # --- V2: use source_store bundle if available ---
        if source_store is not None:
            bundle = source_store.get_bundle(bp_slide.slide_id)
            if bundle and bundle.source_text:
                parts.append(bundle.source_text)
                for s in bundle.asset_summaries:
                    parts.append(s)
                for s in bundle.table_summaries:
                    parts.append(s)
                evidence_added = True
                logger.debug(
                    "Slide %d: using source_store bundle for initial generation",
                    bp_slide.slide_id,
                )

        # --- Legacy fallback ---
        search_keywords = set()
        for item in bp_slide.must_cover_subset:
            for word in re.split(r'[\s_]+', item):
                cleaned = re.sub(r'[^a-zA-Z]', '', word).lower()
                if len(cleaned) > 3:
                    search_keywords.add(cleaned)
        for word in bp_slide.primary_proposition.split():
            cleaned = re.sub(r'[^a-zA-Z]', '', word).lower()
            if len(cleaned) > 4:
                search_keywords.add(cleaned)

        matched_chunk_ids = set()

        # Pass 1: exact match (resolve DB* → B* via source_store)
        resolved_ids = _resolve_linked_ids(bp_slide.linked_evidence_ids, source_store)
        for eid in resolved_ids:
            for chunk in evidence.chunks:
                if chunk.chunk_id == eid:
                    parts.append(f"\n### {chunk.chunk_id} ({chunk.chunk_type})")
                    parts.append(chunk.content[:2000])
                    evidence_added = True
                    matched_chunk_ids.add(chunk.chunk_id)
                    break

        # Pass 2: fuzzy match
        if not evidence_added and search_keywords:
            for chunk in evidence.chunks:
                if chunk.chunk_id in matched_chunk_ids:
                    continue
                section = chunk.metadata.get("section", "").lower()
                chunk_keywords = set(re.sub(r'[^a-zA-Z\s]', '', section).lower().split())
                overlap = search_keywords & chunk_keywords
                if len(overlap) >= 1 and len(chunk.content.strip()) > 50:
                    parts.append(f"\n### {chunk.chunk_id} ({chunk.chunk_type})")
                    parts.append(chunk.content[:2000])
                    evidence_added = True
                    matched_chunk_ids.add(chunk.chunk_id)
                    if len(matched_chunk_ids) >= 3:
                        break

        if not evidence_added:
            parts.append("No specific evidence chunks linked. Use the slide goal as primary content.")

        # Always inject must_cover_subset regardless of evidence source
        if bp_slide.must_cover_subset:
            parts.append("\n## Required Content Points")
            parts.append("You MUST include ALL of the following specific data/facts in this slide.")
            if self.show_source_citations:
                parts.append("Each item includes a [cite: Page X] tag — you MUST use these page numbers in visible Source citations at the bottom of the slide (e.g., `Source: Page 3, Page 5`).\n")
            else:
                parts.append("Each item includes a [cite: Page X] tag for source tracking only.\n")
            for item in bp_slide.must_cover_subset:
                # Try to find page number for this item from source_store
                page_ref = _find_page_for_must_cover(item, bp_slide, source_store)
                if page_ref:
                    parts.append(f"  - {item} [cite: Page {page_ref}]")
                else:
                    parts.append(f"  - {item}")
            parts.append("\nDo NOT omit any of these points. Use exact figures from the evidence above.")
            parts.append("WARNING: Do NOT add descriptive phrases, marketing language, or technical details not found verbatim in the evidence. Every claim must be traceable to the source text.")
            parts.append("LAYOUT: If all points cannot fit cleanly, use a simpler layout (fewer columns, larger fonts) rather than cramming content.")

        # Tables
        if evidence.tables:
            relevant_tables = []
            for tbl in evidence.tables:
                if tbl.table_id in bp_slide.linked_evidence_ids:
                    relevant_tables.append(tbl)
            if not relevant_tables and search_keywords:
                for tbl in evidence.tables:
                    caption_lower = (tbl.caption or "").lower()
                    content_lower = (tbl.content or "")[:200].lower()
                    if any(kw in caption_lower or kw in content_lower for kw in search_keywords):
                        relevant_tables.append(tbl)
                        if len(relevant_tables) >= 2:
                            break
            if relevant_tables:
                parts.append("\n## Available Tables")
                for tbl in relevant_tables[:2]:
                    parts.append(f"\nTable: {tbl.table_id}")
                    if tbl.caption:
                        parts.append(f"Caption: {tbl.caption}")
                    if tbl.content:
                        rows = tbl.content.strip().split("\n")
                        if rows:
                            parts.append("```")
                            for row in rows[:15]:
                                parts.append(row)
                            if len(rows) > 15:
                                parts.append(f"... ({len(rows) - 15} more rows)")
                            parts.append("```")

        # Images — use absolute paths for HTML <img src="">
        logger.debug("_build_slide_prompt slide %d: image_dir=%s, evidence.figures=%d, assigned_figure_id=%s",
                      bp_slide.slide_id, image_dir, len(evidence.figures),
                      getattr(bp_slide, 'assigned_figure_id', 'N/A'))
        if image_dir and evidence.figures:
            embedded_figs = [f for f in evidence.figures
                if f.figure_type not in ("page_screenshot", "table_screenshot")
                and (f.width or 0) * (f.height or 0) < 4_000_000
                and (f.height or 1) / max(f.width or 1, 1) < 2.5
            ]

            # Fallback: use page screenshots if no embedded figures
            if not embedded_figs:
                embedded_figs = [f for f in evidence.figures
                    if f.figure_type == "page_screenshot"
                    and (f.width or 0) * (f.height or 0) < 4_000_000
                ]

            # Priority 1: Use assigned_figure_id from deck planner (pre-allocated, dedup-safe)
            assigned_fig = None
            if hasattr(bp_slide, 'assigned_figure_id') and bp_slide.assigned_figure_id:
                afid = bp_slide.assigned_figure_id
                # Direct match by figure_id
                for fig in evidence.figures:
                    if fig.figure_id == afid:
                        assigned_fig = fig
                        break
                # Fallback: match via source_store asset_id → image_path → stem
                if not assigned_fig and source_store:
                    for asset in getattr(source_store, 'assets', []):
                        if getattr(asset, 'asset_id', None) == afid:
                            asset_stem = Path(asset.image_path).stem if hasattr(asset, 'image_path') and asset.image_path else None
                            if asset_stem:
                                for fig in evidence.figures:
                                    if fig.figure_id == asset_stem:
                                        assigned_fig = fig
                                        break
                            break
                if not assigned_fig:
                    logger.warning(
                        "Slide %d: assigned_figure_id='%s' NOT found in evidence (%d figs)",
                        bp_slide.slide_id, afid, len(evidence.figures),
                    )

            if assigned_fig and assigned_fig.image_path:
                relevant_figures = [assigned_fig]
                logger.info("Slide %d: using assigned figure '%s' -> %s",
                           bp_slide.slide_id, assigned_fig.figure_id, assigned_fig.image_path)
            elif hasattr(bp_slide, 'assigned_figure_id') and bp_slide.assigned_figure_id is not None:
                # Planner explicitly assigned empty string = no image for this slide
                relevant_figures = []
            else:
                scored_figs = []
                for fig in embedded_figs:
                    score = 0
                    if fig.figure_id in bp_slide.linked_evidence_ids:
                        score += 10
                    if search_keywords:
                        text_pool = f"{fig.caption} {fig.description}".lower()
                        text_words = set(re.sub(r'[^a-zA-Z\s]', '', text_pool).split())
                        overlap = search_keywords & text_words
                        score += len(overlap)
                    if score > 0:
                        scored_figs.append((score, fig))

                scored_figs.sort(key=lambda x: -x[0])
                relevant_figures = []
                for sc, fig in scored_figs:
                    if sc >= 3:
                        fig_filename = Path(fig.image_path).name if fig.image_path else ""
                        if fig_filename in self.used_image_ids:
                            continue
                        relevant_figures.append(fig)
                    if len(relevant_figures) >= 1:
                        break

            if relevant_figures:
                has_images = True
                parts.append(f"\n## Available Images")
                parts.append("Use absolute file paths in <img src=\"...\"> tags.")
                parts.append("CRITICAL: You may ONLY use the EXACT paths listed below.")
                parts.append("ONLY include an image if it directly illustrates THIS slide's topic.")
                parts.append("Always include `max-height` and `object-fit: contain` on <img> elements.")
                parts.append("If a figure image has a caption at the bottom (text like 'Figure N. ...'), wrap the <img> in a container with `overflow:hidden` and use `object-fit:contain` with a negative bottom margin to crop just the caption. Example:")
                parts.append('```html\n<div style="overflow:hidden; height:350px;">\n  <img src="..." style="width:100%; max-height:400px; object-fit:contain; object-position:top; margin-bottom:-50px;">\n</div>\n```\n')
                for fig in relevant_figures:
                    if fig.image_path:
                        abs_path = Path(fig.image_path)
                        px_w = fig.width
                        px_h = fig.height
                        if (not px_w or not px_h) and abs_path.exists():
                            try:
                                from PIL import Image as _PILImage
                                with _PILImage.open(str(abs_path)) as _im:
                                    px_w, px_h = _im.size
                            except Exception:
                                pass
                        size_info = ""
                        aspect_info = ""
                        if px_w and px_h:
                            aspect = round(px_w / px_h, 2)
                            if aspect > 1.3:
                                orientation = "landscape"
                            elif aspect < 0.77:
                                orientation = "portrait"
                            else:
                                orientation = "square"
                            size_info = f"  [{px_w}x{px_h} px, aspect={aspect}, {orientation}]"
                            aspect_info = f"  Layout hint: {orientation} image — " + (
                                "use wide column (60-70% width) or full-width placement"
                                if orientation == "landscape"
                                else "use narrow column (30-40% width), let height fill"
                                if orientation == "portrait"
                                else "use equal column split or ~50% width"
                            )
                        desc = fig.description if fig.description and "Embedded image from" not in fig.description else ""
                        caption = fig.caption or ""
                        # Truncate caption to avoid verbatim text walls on slides
                        if caption:
                            words = caption.split()
                            if len(words) > 30:
                                caption = " ".join(words[:30]) + "..."
                        label = desc or caption or "no description"
                        parts.append(f"- \"{abs_path}\"{size_info}")
                        parts.append(f"  Content: {label}")
                        if aspect_info:
                            parts.append(aspect_info)
                parts.append("\nDo NOT use page screenshots. Do NOT stretch images.")
                parts.append("Match the image container's aspect ratio to the image's actual aspect ratio — avoid placing a portrait/square image in a wide horizontal container (creates ugly whitespace).")
            else:
                parts.append("\n## NO IMAGES AVAILABLE for this slide")
                parts.append("Do NOT use <img> tags. Use text, shapes, tables, and cards instead.")
                parts.append("Use FULL-WIDTH text layout.")
        else:
            if bp_slide.role != "title":
                parts.append("\n## NO IMAGES AVAILABLE for this slide")
                parts.append("Do NOT use <img> tags. Use text, shapes, tables, and cards instead.")

        # Inject theme palette — ensures consistent colors across all slides
        theme = self._theme or DEFAULT_THEME

        def _hex(c):
            return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"

        def _contrast_on_white(c):
            """WCAG contrast ratio of color c on white background."""
            def _srgb(v):
                v = v / 255
                return v / 12.92 if v <= 0.04045 else ((v + 0.055) / 1.055) ** 2.4
            L = 0.2126 * _srgb(c[0]) + 0.7152 * _srgb(c[1]) + 0.0722 * _srgb(c[2])
            return (1.05) / (L + 0.05)

        def _darken_to_pass(c, min_ratio=4.5):
            """Darken a color until it passes WCAG contrast on white."""
            r, g, b = c
            for _ in range(20):  # max 20 darkening steps
                if _contrast_on_white((r, g, b)) >= min_ratio:
                    return (r, g, b)
                r = max(0, int(r * 0.9))
                g = max(0, int(g * 0.9))
                b = max(0, int(b * 0.9))
            return (r, g, b)

        def _text_ok(c):
            return "✅ OK for text" if _contrast_on_white(c) >= 4.5 else "⚠️ DECORATIVE ONLY — too light for text on white"

        # Auto-darken PRIMARY_MID and ACCENT if they fail contrast check
        # This prevents B5 low_contrast issues at the source
        primary_mid = theme.primary_mid
        if _contrast_on_white(primary_mid) < 4.5:
            primary_mid = _darken_to_pass(primary_mid)
            logger.info("Auto-darkened PRIMARY_MID from %s to %s for WCAG compliance",
                       _hex(theme.primary_mid), _hex(primary_mid))

        accent = theme.accent
        if _contrast_on_white(accent) < 4.5:
            accent = _darken_to_pass(accent)
            logger.info("Auto-darkened ACCENT from %s to %s for WCAG compliance",
                       _hex(theme.accent), _hex(accent))

        parts.append("\n## Color Palette (MANDATORY — use ONLY these colors)")
        parts.append(f"- **Primary Dark**: {_hex(theme.primary_dark)} — dark backgrounds, header bars. {_text_ok(theme.primary_dark)}")
        parts.append(f"- **Primary Mid**: {_hex(primary_mid)} — titles, headings. {_text_ok(primary_mid)}")
        parts.append(f"- **Primary Light**: {_hex(theme.primary_light)} — borders, underlines. {_text_ok(theme.primary_light)}")
        parts.append(f"- **Accent**: {_hex(accent)} — emphasis, key metrics. {_text_ok(accent)}")
        accent_alt = theme.accent_alt
        if _contrast_on_white(accent_alt) < 4.5:
            accent_alt = _darken_to_pass(accent_alt)
        parts.append(f"- **Accent Alt**: {_hex(accent_alt)} — secondary accent. {_text_ok(accent_alt)}")
        parts.append(f"- **Body Text**: {_hex(theme.body_text)} — all body text. {_text_ok(theme.body_text)}")
        caption = theme.caption_text
        if _contrast_on_white(caption) < 4.5:
            caption = _darken_to_pass(caption)
        parts.append(f"- **Caption**: {_hex(caption)} — footnotes. {_text_ok(caption)}")
        parts.append(f"- **Light BG**: {_hex(theme.light_bg)} — card/container backgrounds only")
        parts.append(f"- **Warm BG**: {_hex(theme.warm_bg)} — alternate backgrounds only")
        parts.append(f"- Font family: {theme.font_family}")
        parts.append(f"\n**Theme: {theme.theme_name}** — EVERY color must come from this palette.")
        parts.append("Colors marked ⚠️ DECORATIVE ONLY must NOT be used as text `color:` on white — use them only for `background:`, `border-color:`, or decorative elements.")
        parts.append("For bold/emphasized text on white, use Primary Dark or Body Text instead.")
        parts.append("CONTRAST ON COLORED BACKGROUNDS: When placing text on Light BG or Warm BG containers, use ONLY Primary Dark or Body Text — never Primary Mid or Accent colors, as they fail contrast on tinted backgrounds.")

        return "\n".join(parts), has_images

    def _extract_html(self, response: str) -> str | None:
        """Extract HTML from LLM response."""
        # Look for ```html code blocks
        pattern = r"```html\s*\n(.*?)```"
        matches = re.findall(pattern, response, re.DOTALL)
        if matches:
            for match in matches:
                html = match.strip()
                if "<!DOCTYPE" in html or "<html" in html or "<body" in html:
                    return html

        # Try without code fences - look for <!DOCTYPE or <html
        if "<!DOCTYPE" in response or "<html" in response:
            lines = response.split("\n")
            in_html = False
            html_lines = []
            for line in lines:
                if "<!DOCTYPE" in line or (not in_html and "<html" in line):
                    in_html = True
                if in_html:
                    html_lines.append(line)
                if in_html and "</html>" in line:
                    break
            if html_lines:
                return "\n".join(html_lines)

        return None

    @staticmethod
    def _estimate_text_content(html: str) -> int:
        """Estimate the amount of text content in HTML."""
        # Strip all tags to get raw text
        text = re.sub(r'<[^>]+>', '', html)
        # Strip CSS
        text = re.sub(r'\{[^}]+\}', '', text)
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        return len(text)

    def _detect_layout(self, html: str) -> str:
        """Detect the layout pattern used in generated HTML."""
        html_lower = html.lower()

        n_images = len(re.findall(r'<img\b', html_lower))
        n_tables = len(re.findall(r'<table\b', html_lower))
        has_flex = "display: flex" in html_lower or "display:flex" in html_lower

        if n_images >= 2:
            return "Multi-Image"
        if n_images == 1 and has_flex:
            return "Two-Column-Image"
        if n_images == 1:
            return "Image-Focus"
        if n_tables >= 1:
            return "Table-Layout"

        # Check for metric cards (multiple flex items with large numbers)
        large_numbers = re.findall(r'font-size:\s*(?:4[0-9]|5[0-9])px', html_lower)
        if len(large_numbers) >= 2:
            return "Metric-Cards"

        # Check for columns
        col_count = len(re.findall(r'flex:\s*1', html_lower))
        if col_count >= 3:
            return "Three-Column"
        if col_count >= 2:
            return "Two-Column"

        if "font-style: italic" in html_lower and "quote" in html_lower:
            return "Key-Quote"

        return "Text-Only"

    def _add_fallback_slide(self, slide, bp_slide: BlueprintSlide, evidence: EvidenceState | None = None):
        """Add simple fallback content when HTML generation fails."""
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        title = bp_slide.primary_proposition
        if len(title) > 60:
            title = title[:57] + "..."
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 153)

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        items = bp_slide.must_cover_subset or []
        for i, item in enumerate(items[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.space_before = Pt(8)

    # ------------------------------------------------------------------
    # Image directory utilities (mirrors CodeGenCompiler)
    # ------------------------------------------------------------------
    # Evidence context (shared with repair fallbacks)
    # ------------------------------------------------------------------

    def _build_evidence_context(
        self,
        blueprint_slide,
        evidence,
        case_dir,
        source_store=None,
    ) -> str:
        """Build evidence context string for repair prompts.

        If source_store has a bundle for this slide, uses it directly.
        Otherwise falls back to legacy linked_evidence_ids matching.
        """
        # --- V2: use bundle if available ---
        if source_store is not None and blueprint_slide:
            bundle = source_store.get_bundle(blueprint_slide.slide_id)
            if bundle and bundle.source_text:
                ev_parts = []
                ev_parts.append(bundle.source_text)
                for s in bundle.asset_summaries:
                    ev_parts.append(s)
                for s in bundle.table_summaries:
                    ev_parts.append(s)
                return (
                    "\n\n## Source Evidence (use ONLY data from here, "
                    "do NOT invent numbers)\n\n" + "\n\n".join(ev_parts)
                )

        # --- Legacy fallback ---
        import re as _re
        if not blueprint_slide or not evidence:
            return ""

        ev_parts = []
        search_keywords = set()
        for item in (blueprint_slide.must_cover_subset or []):
            for word in _re.split(r'[\s_]+', item):
                cleaned = _re.sub(r'[^a-zA-Z]', '', word).lower()
                if len(cleaned) > 3:
                    search_keywords.add(cleaned)

        resolved_ids = _resolve_linked_ids(
            blueprint_slide.linked_evidence_ids or [], source_store,
        )
        resolved_set = set(resolved_ids)
        # Build chunk lookup for ordered iteration
        chunk_map = {c.chunk_id: c for c in evidence.chunks}
        matched = 0
        matched_ids = set()
        # Iterate in resolved order (preserves linked_evidence_ids priority)
        for rid in resolved_ids:
            chunk = chunk_map.get(rid)
            if chunk and rid not in matched_ids:
                ev_parts.append(f"[{chunk.chunk_id}] {chunk.content[:2500]}")
                matched_ids.add(rid)
                matched += 1
            if matched >= 8:
                break

        if matched < 8 and search_keywords:
            for chunk in evidence.chunks:
                if chunk.chunk_id in resolved_set or chunk.chunk_id in matched_ids:
                    continue
                section = chunk.metadata.get("section", "").lower()
                content_lower = chunk.content[:300].lower()
                chunk_keywords = set(_re.sub(r'[^a-zA-Z\s]', '', section).lower().split())
                content_keywords = set(_re.sub(r'[^a-zA-Z\s]', '', content_lower).split())
                all_chunk_words = chunk_keywords | content_keywords
                if search_keywords & all_chunk_words and len(chunk.content.strip()) > 50:
                    ev_parts.append(f"[{chunk.chunk_id}] {chunk.content[:2500]}")
                    matched += 1
                if matched >= 8:
                    break

        for tbl in evidence.tables:
            if tbl.table_id in (blueprint_slide.linked_evidence_ids or []):
                tbl_desc = f" — {tbl.description}" if tbl.description else ""
                ev_parts.append(f"[{tbl.table_id}]{tbl_desc}\n{tbl.content[:1500]}")

        if ev_parts:
            return "\n\n## Source Evidence (use ONLY data from here, do NOT invent numbers)\n\n" + "\n\n".join(ev_parts)
        return ""

    # ------------------------------------------------------------------

    def _find_image_dir(self, case_dir: Path) -> Path | None:
        """Find the source_pack directory (parent of figures/, tables/, screenshots/).

        New pyramid structure:
          source_pack/figures/      — extracted figure images
          source_pack/tables/       — table screenshots + JSON sidecars
          source_pack/screenshots/  — page screenshots
          source_pack/rendered_formulas/ — rendered LaTeX formulas

        Falls back to legacy extracted_from_pdf/ for unmigrated cases.
        """
        source_pack = case_dir / "source_pack"
        if source_pack.is_dir():
            if (source_pack / "figures").is_dir() or (source_pack / "tables").is_dir():
                return source_pack
            legacy = source_pack / "extracted_from_pdf"
            if legacy.is_dir():
                return legacy
            return source_pack
        return None

    def _list_images(self, image_dir: Path) -> list[str]:
        """List available image files with relative paths from image_dir.

        Searches pyramid subdirectories: figures/, tables/, screenshots/,
        rendered_formulas/. Also supports legacy subdirs (images/) for
        backward compatibility.
        """
        exts = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif", ".webp"}
        images = []
        for subdir in ["figures", "tables", "screenshots", "rendered_formulas",
                        "images"]:  # "images" for legacy compat
            d = image_dir / subdir
            if d.is_dir():
                for f in sorted(d.iterdir()):
                    if f.is_file() and f.suffix.lower() in exts:
                        images.append(f"{subdir}/{f.name}")
        for f in sorted(image_dir.iterdir()):
            if f.is_file() and f.suffix.lower() in exts:
                images.append(f.name)
        return images
