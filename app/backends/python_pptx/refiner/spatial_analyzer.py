"""Font-aware overflow and overlap analysis for PPTX files."""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .font_metrics import FontMetrics, OverflowResult, EMU_PER_INCH
from ....schemas.extraction import ExtractedObject, SlideExtraction

_DEFAULT_FONT = "Arial"
_DEFAULT_FONT_SIZE_EMU = 12 * 12700  # 12 pt in EMU
_OVERLAP_THRESHOLD_SQ_IN = 0.3
_BACKGROUND_AREA_RATIO = 0.85


# ── Dataclasses ───────────────────────────────────────────────────────────────

@dataclass
class ShapeOverflow:
    slide_index: int
    shape_name: str
    shape_index: int
    text_snippet: str
    overflow: OverflowResult
    font_name: str
    font_size_emu: float
    box_width_emu: float
    box_height_emu: float


@dataclass
class OverlapPair:
    slide_index: int
    shape_a_name: str
    shape_a_index: int
    shape_b_name: str
    shape_b_index: int
    overlap_area_sq_in: float


@dataclass
class SmallFont:
    """A text run with font size below the readability threshold."""
    slide_index: int
    shape_name: str
    shape_index: int
    text_snippet: str
    font_size_pt: float
    min_readable_pt: float  # the threshold it violated


@dataclass
class SpatialReport:
    extractions: list[SlideExtraction] = field(default_factory=list)
    overflows: list[ShapeOverflow] = field(default_factory=list)
    overlaps: list[OverlapPair] = field(default_factory=list)
    small_fonts: list[SmallFont] = field(default_factory=list)
    slide_width_emu: int = 0
    slide_height_emu: int = 0


# ── Analyzer ─────────────────────────────────────────────────────────────────

class PptxSpatialAnalyzer:
    def __init__(self, font_metrics: FontMetrics):
        self._fm = font_metrics

    def analyze(self, pptx_path) -> SpatialReport:
        prs = Presentation(str(pptx_path))
        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)

        report = SpatialReport(slide_width_emu=slide_w, slide_height_emu=slide_h)

        for slide_index, slide in enumerate(prs.slides):
            slide_id = slide.slide_id
            extraction, overflows, overlaps, small_fonts = self._analyze_slide(
                slide, slide_index, slide_id, slide_w, slide_h
            )
            report.extractions.append(extraction)
            report.overflows.extend(overflows)
            report.overlaps.extend(overlaps)
            report.small_fonts.extend(small_fonts)

        return report

    def _analyze_slide(self, slide, slide_index: int, slide_id: int,
                       slide_w: int, slide_h: int):
        slide_area = slide_w * slide_h

        objects: list[ExtractedObject] = []
        overflows: list[ShapeOverflow] = []
        overlaps: list[OverlapPair] = []

        shapes = list(slide.shapes)

        # Extract objects and check overflows
        for shape_index, shape in enumerate(shapes):
            obj = self._extract_shape(shape, slide_index, shape_index)
            objects.append(obj)

            overflow = self._check_shape_overflow(shape, slide_index, shape_index)
            if overflow is not None:
                overflows.append(overflow)

        # Overlap detection — skip background shapes
        non_bg: list[tuple[int, object, ExtractedObject]] = []
        for shape_index, (shape, obj) in enumerate(zip(shapes, objects)):
            if len(obj.bbox_emu) == 4:
                w, h = obj.bbox_emu[2], obj.bbox_emu[3]
                area = w * h
                if slide_area > 0 and area / slide_area > _BACKGROUND_AREA_RATIO:
                    continue
            non_bg.append((shape_index, shape, obj))

        for i in range(len(non_bg)):
            for j in range(i + 1, len(non_bg)):
                idx_a, _, obj_a = non_bg[i]
                idx_b, _, obj_b = non_bg[j]
                overlap = self._compute_overlap(obj_a, obj_b)
                if overlap > _OVERLAP_THRESHOLD_SQ_IN:
                    overlaps.append(OverlapPair(
                        slide_index=slide_index,
                        shape_a_name=obj_a.shape_name,
                        shape_a_index=idx_a,
                        shape_b_name=obj_b.shape_name,
                        shape_b_index=idx_b,
                        overlap_area_sq_in=overlap,
                    ))

        # Build SlideExtraction
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text
        except Exception:
            pass

        extraction = SlideExtraction(
            slide_id=slide_id,
            slide_index=slide_index,
            title=title,
            objects=objects,
            total_text_length=sum(len(o.text_content) for o in objects),
            total_objects=len(objects),
        )

        # Small font detection (< 12pt body text is hard to read)
        MIN_READABLE_PT = 12.0
        # Thresholds to skip decorative/icon labels
        MIN_BOX_AREA_SQ_IN = 0.5    # boxes smaller than 0.5 sq in are likely decorative
        MIN_TEXT_LENGTH = 5          # very short text (1-4 chars) is likely a label/icon
        MIN_BOX_WIDTH_IN = 1.5      # boxes narrower than 1.5in can't hold enlarged text
        small_fonts: list[SmallFont] = []
        for shape_index, shape in enumerate(shapes):
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue

            # Skip small decorative boxes
            try:
                box_w_in = shape.width / 914400
                box_h_in = shape.height / 914400
                box_area_sq_in = box_w_in * box_h_in
            except Exception:
                box_w_in = box_h_in = box_area_sq_in = 0
            full_text = shape.text_frame.text.strip()

            # Skip if: tiny box AND short text (decorative label)
            if box_area_sq_in < MIN_BOX_AREA_SQ_IN and len(full_text) < MIN_TEXT_LENGTH:
                continue

            # Skip narrow boxes where enlarging would just cause overflow
            if box_w_in < MIN_BOX_WIDTH_IN and len(full_text) > 3:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    pt = run.font.size / 12700
                    if pt < MIN_READABLE_PT and run.text.strip():
                        small_fonts.append(SmallFont(
                            slide_index=slide_index,
                            shape_name=shape.name or f"shape_{shape_index}",
                            shape_index=shape_index,
                            text_snippet=run.text[:50],
                            font_size_pt=round(pt, 1),
                            min_readable_pt=MIN_READABLE_PT,
                        ))

        return extraction, overflows, overlaps, small_fonts

    def _extract_shape(self, shape, slide_index: int, shape_index: int) -> ExtractedObject:
        try:
            left = int(shape.left or 0)
            top = int(shape.top or 0)
            width = int(shape.width or 0)
            height = int(shape.height or 0)
            bbox = [left, top, width, height]
        except Exception:
            bbox = []

        object_type = "shape"
        text_content = ""
        font_sizes: list[float] = []

        try:
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.TEXT_BOX:
                object_type = "text_box"
            elif st == MSO_SHAPE_TYPE.PICTURE:
                object_type = "picture"
            elif st == MSO_SHAPE_TYPE.TABLE:
                object_type = "table"
            elif st == MSO_SHAPE_TYPE.GROUP:
                object_type = "group"
        except Exception:
            pass

        try:
            if shape.has_text_frame:
                object_type = object_type if object_type != "shape" else "text_box"
                text_content = shape.text_frame.text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size / 12700.0)
        except Exception:
            pass

        return ExtractedObject(
            object_id=f"s{slide_index}_sh{shape_index}",
            shape_name=shape.name,
            object_type=object_type,
            bbox_emu=bbox,
            text_content=text_content,
            font_sizes_pt=font_sizes,
            z_order=shape_index,
        )

    def _check_shape_overflow(self, shape, slide_index: int,
                               shape_index: int) -> Optional[ShapeOverflow]:
        try:
            if not shape.has_text_frame:
                return None
        except Exception:
            return None

        try:
            text = shape.text_frame.text
            if not text or not text.strip():
                return None

            box_w = int(shape.width or 0)
            box_h = int(shape.height or 0)
            if box_w <= 0 or box_h <= 0:
                return None

            # Get dominant font from first run of first paragraph
            font_name = _DEFAULT_FONT
            font_size_emu = _DEFAULT_FONT_SIZE_EMU
            bold = False

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        font_name = run.font.name
                    if run.font.size:
                        font_size_emu = int(run.font.size)
                    if run.font.bold:
                        bold = True
                    break
                break

            overflow = self._fm.check_overflow(
                text=text,
                font_name=font_name,
                size_emu=font_size_emu,
                bold=bold,
                box_width_emu=box_w,
                box_height_emu=box_h,
            )

            if not (overflow.overflows_width or overflow.overflows_height):
                return None

            snippet = text[:60] + ("..." if len(text) > 60 else "")
            return ShapeOverflow(
                slide_index=slide_index,
                shape_name=shape.name,
                shape_index=shape_index,
                text_snippet=snippet,
                overflow=overflow,
                font_name=font_name,
                font_size_emu=font_size_emu,
                box_width_emu=box_w,
                box_height_emu=box_h,
            )
        except Exception:
            return None

    def _compute_overlap(self, a: ExtractedObject, b: ExtractedObject) -> float:
        """Return AABB overlap in square inches. Returns 0 if no overlap."""
        if len(a.bbox_emu) < 4 or len(b.bbox_emu) < 4:
            return 0.0

        al, at, aw, ah = a.bbox_emu
        bl, bt, bw, bh = b.bbox_emu

        ar, ab_ = al + aw, at + ah
        br, bb_ = bl + bw, bt + bh

        ox = max(0, min(ar, br) - max(al, bl))
        oy = max(0, min(ab_, bb_) - max(at, bt))

        overlap_emu2 = ox * oy
        sq_in = overlap_emu2 / (EMU_PER_INCH ** 2)
        return sq_in
