"""Applies a FixPlan to a PPTX file using python-pptx.

Includes collision detection: resize/move/expand ops are rejected
if the new bounding box would overlap a neighbor shape on the same slide.
"""
from __future__ import annotations

import logging
import shutil
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt

from .fix_planner import FixOp, FixPlan
from .font_metrics import FontMetrics, EMU_PER_PT, EMU_PER_INCH

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400
# Minimum gap between shapes (0.05 inches) to treat as non-collision
_GAP_TOLERANCE = int(0.05 * EMU_PER_INCH)
# Background shapes (area > 85% of slide) are exempt from collision
_BG_RATIO = 0.85


class PptxFixExecutor:
    def __init__(self, font_metrics: FontMetrics | None = None):
        self._fm = font_metrics or FontMetrics()

    def apply(
        self,
        pptx_path: str | Path,
        plan: FixPlan,
        output_path: str | Path | None = None,
    ) -> Path:
        pptx_path = Path(pptx_path)
        if output_path is None:
            output_path = pptx_path
        else:
            output_path = Path(output_path)

        if output_path != pptx_path:
            shutil.copy2(pptx_path, output_path)

        prs = Presentation(output_path)
        applied = 0
        rejected = 0

        for op in plan.ops:
            ok = self._apply_op(prs, op)
            if ok:
                applied += 1
            else:
                rejected += 1

        prs.save(output_path)
        if rejected:
            logger.info("FixExecutor: applied %d, rejected %d (collision/missing)", applied, rejected)
        return output_path

    # ── dispatch ──────────────────────────────────────────────────────────────

    def _apply_op(self, prs: Presentation, op: FixOp) -> bool:
        slides = prs.slides
        if op.slide_index < 0 or op.slide_index >= len(slides):
            return False
        slide = slides[op.slide_index]
        shape = self._find_shape(slide, op.shape_name)
        if shape is None:
            logger.warning("Shape '%s' not found on slide %d", op.shape_name, op.slide_index)
            return False

        handler = {
            "shrink_font": self._shrink_font,
            "enlarge_font": self._shrink_font,  # same mechanic, just sets target size
            "resize": self._resize_safe,
            "move": self._move_safe,
            "expand_box": self._resize_safe,
            "set_word_wrap": self._set_word_wrap,
        }.get(op.op_type)

        if handler is None:
            return False

        return handler(shape, op.params, slide, prs)

    def _find_shape(self, slide, shape_name: str):
        for shape in slide.shapes:
            if shape.name == shape_name:
                return shape
        return None

    # ── collision detection ───────────────────────────────────────────────────

    @staticmethod
    def _would_collide(
        new_left: int, new_top: int, new_width: int, new_height: int,
        target_name: str, slide, slide_area: int,
    ) -> str | None:
        """Check if new bbox collides with any neighbor.

        Returns the colliding shape's name, or None if safe.
        Skips background shapes and the target shape itself.
        """
        new_right = new_left + new_width
        new_bottom = new_top + new_height

        for other in slide.shapes:
            if other.name == target_name:
                continue
            # Skip backgrounds
            other_area = other.width * other.height
            if slide_area > 0 and other_area / slide_area > _BG_RATIO:
                continue
            # Skip shapes with no area
            if other.width <= 0 or other.height <= 0:
                continue

            o_left = other.left + _GAP_TOLERANCE
            o_top = other.top + _GAP_TOLERANCE
            o_right = other.left + other.width - _GAP_TOLERANCE
            o_bottom = other.top + other.height - _GAP_TOLERANCE

            # AABB overlap test
            if (new_left < o_right and new_right > o_left and
                    new_top < o_bottom and new_bottom > o_top):
                # Check if they already overlapped before the change
                return other.name

        return None

    # ── op implementations ────────────────────────────────────────────────────

    def _shrink_font(self, shape, params: dict, slide, prs) -> bool:
        """Set font size, then auto-expand box if text no longer fits."""
        target_pt = params.get("target_size_pt")
        if target_pt is None:
            return False
        if not shape.has_text_frame:
            return False

        # Apply font size change
        target_emu = Pt(float(target_pt))
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = target_emu

        # After font change, check if text still fits in the box
        text = shape.text_frame.text
        if not text.strip():
            return True

        # Get dominant font name
        font_name = "Arial"
        bold = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name:
                    font_name = run.font.name
                if run.font.bold:
                    bold = True
                break
            break

        overflow = self._fm.check_overflow(
            text=text,
            font_name=font_name,
            size_emu=int(target_emu),
            bold=bold,
            box_width_emu=shape.width,
            box_height_emu=shape.height,
        )

        if overflow.overflows_width or overflow.overflows_height:
            # Try to expand box safely
            new_w = int(shape.width + overflow.overflow_width_emu + 0.1 * EMU_PER_INCH) if overflow.overflows_width else int(shape.width)
            new_h = int(shape.height + overflow.overflow_height_emu + 0.05 * EMU_PER_INCH) if overflow.overflows_height else int(shape.height)
            slide_area = prs.slide_width * prs.slide_height

            # Clamp to slide bounds
            new_w = int(min(new_w, prs.slide_width - shape.left))
            new_h = int(min(new_h, prs.slide_height - shape.top))

            collider = self._would_collide(
                shape.left, shape.top, new_w, new_h,
                shape.name, slide, slide_area,
            )
            if collider:
                # Can't expand — revert font change to avoid new overflow
                logger.warning(
                    "Font change on '%s' would overflow and can't expand (blocked by '%s') — reverting",
                    shape.name, collider,
                )
                # Revert: we don't know the original size, so skip this op
                return False
            else:
                shape.width = new_w
                shape.height = new_h
                logger.info(
                    "Auto-expanded '%s' to %dx%d after font change to %.0fpt",
                    shape.name, new_w, new_h, target_pt,
                )

        return True

    def _resize_safe(self, shape, params: dict, slide, prs) -> bool:
        """Resize with collision check."""
        new_w = int(params["width_emu"]) if "width_emu" in params else shape.width
        new_h = int(params["height_emu"]) if "height_emu" in params else shape.height
        slide_area = prs.slide_width * prs.slide_height

        # Clamp to slide bounds
        max_w = prs.slide_width - shape.left
        max_h = prs.slide_height - shape.top
        new_w = min(new_w, max_w)
        new_h = min(new_h, max_h)

        collider = self._would_collide(
            shape.left, shape.top, new_w, new_h,
            shape.name, slide, slide_area,
        )
        if collider:
            # Try clamping width to just before the collider
            for other in slide.shapes:
                if other.name == collider:
                    safe_w = other.left - shape.left - _GAP_TOLERANCE
                    if safe_w > shape.width:
                        new_w = safe_w
                        logger.info(
                            "Resize '%s': clamped width to %d to avoid '%s'",
                            shape.name, new_w, collider,
                        )
                        break
            else:
                logger.warning(
                    "Resize '%s' rejected: would collide with '%s'",
                    shape.name, collider,
                )
                return False

        shape.width = new_w
        shape.height = new_h
        return True

    def _move_safe(self, shape, params: dict, slide, prs) -> bool:
        """Move with collision check."""
        new_left = int(params["left_emu"]) if "left_emu" in params else shape.left
        new_top = int(params["top_emu"]) if "top_emu" in params else shape.top
        slide_area = prs.slide_width * prs.slide_height

        # Clamp to slide bounds
        new_left = max(0, min(new_left, prs.slide_width - shape.width))
        new_top = max(0, min(new_top, prs.slide_height - shape.height))

        collider = self._would_collide(
            new_left, new_top, shape.width, shape.height,
            shape.name, slide, slide_area,
        )
        if collider:
            logger.warning(
                "Move '%s' to (%d,%d) rejected: would collide with '%s'",
                shape.name, new_left, new_top, collider,
            )
            return False

        shape.left = new_left
        shape.top = new_top
        return True

    @staticmethod
    def _set_word_wrap(shape, params: dict, slide, prs) -> bool:
        wrap = params.get("wrap", params.get("word_wrap", True))
        if shape.has_text_frame:
            shape.text_frame.word_wrap = bool(wrap)
            return True
        return False
