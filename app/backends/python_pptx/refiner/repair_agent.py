"""PptxRepairAgent — Agentic repair loop for PPTX spatial issues.

Mirrors ReDeck's AgentRepair architecture but operates on python-pptx
shapes instead of HTML code. The LLM agent has tools to:
  - move_shape / resize_shape / set_font_size
  - verify_layout (re-render + return screenshot)
  - get_shape_info (read current positions)
  - submit (finalize)

The agent sees the slide screenshot, issues, and shape data.
It makes tool calls to fix issues one at a time, verifying after each.
"""

import json
import logging
import shutil
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu

from .font_metrics import FontMetrics, EMU_PER_INCH, EMU_PER_PT
from .spatial_analyzer import PptxSpatialAnalyzer, SpatialReport
from ....llm_client import LLMClient
from ....schemas.issue import Issue
from ....utils.image_ops import image_to_base64

logger = logging.getLogger(__name__)

# Tool definitions for the LLM
TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "move_shape",
            "description": "Move a text box to a new position. Provide left_emu and/or top_emu. Omit a field to keep current value.",
            "parameters": {
                "type": "object",
                "properties": {
                    "slide_index": {"type": "integer", "description": "0-based slide index"},
                    "shape_name": {"type": "string", "description": "Name of the shape to move"},
                    "left_emu": {"type": "integer", "description": "New left position in EMU (914400 EMU = 1 inch)"},
                    "top_emu": {"type": "integer", "description": "New top position in EMU"},
                },
                "required": ["slide_index", "shape_name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "resize_shape",
            "description": "Resize a text box. Provide width_emu and/or height_emu. Omit a field to keep current value.",
            "parameters": {
                "type": "object",
                "properties": {
                    "slide_index": {"type": "integer", "description": "0-based slide index"},
                    "shape_name": {"type": "string", "description": "Name of the shape"},
                    "width_emu": {"type": "integer", "description": "New width in EMU"},
                    "height_emu": {"type": "integer", "description": "New height in EMU"},
                },
                "required": ["slide_index", "shape_name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_font_size",
            "description": "Set font size for all text in a shape.",
            "parameters": {
                "type": "object",
                "properties": {
                    "slide_index": {"type": "integer", "description": "0-based slide index"},
                    "shape_name": {"type": "string", "description": "Name of the shape"},
                    "target_size_pt": {"type": "number", "description": "Target font size in points"},
                },
                "required": ["slide_index", "shape_name", "target_size_pt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_shape_info",
            "description": "Get current position, size, and text of a shape.",
            "parameters": {
                "type": "object",
                "properties": {
                    "slide_index": {"type": "integer", "description": "0-based slide index"},
                    "shape_name": {"type": "string", "description": "Name of the shape"},
                },
                "required": ["slide_index", "shape_name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "verify_layout",
            "description": "Re-render the current slide to PNG and return the image for visual verification. Call this after making changes to see the result.",
            "parameters": {
                "type": "object",
                "properties": {
                    "slide_index": {"type": "integer", "description": "0-based slide index"},
                },
                "required": ["slide_index"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "submit",
            "description": "Finalize repairs for this slide. Call when all issues are fixed or no further improvements can be made.",
            "parameters": {
                "type": "object",
                "properties": {
                    "summary": {"type": "string", "description": "Brief summary of changes made"},
                },
                "required": ["summary"],
            },
        },
    },
]

SYSTEM_PROMPT = """You are a PPTX slide layout repair agent. You fix spatial issues by moving and resizing text boxes.

IMPORTANT CONTEXT:
- The slide background is a SINGLE IMAGE containing all icons, borders, decorative elements
- Only the text boxes listed in the shape data are movable — you CANNOT move background elements
- Your goal: position text boxes so they align properly with background icons and don't overlap background elements

WORKFLOW — MUST follow this exact process:

**Phase 1: Fix all issues**
For EACH reported issue, apply the fix (move_shape / resize_shape / set_font_size).

**Phase 2: First verify**
Call verify_layout. You will see the updated slide image.

**Phase 3: Critical self-check** — go through this checklist on the NEW image:
For EVERY issue that was reported, ask yourself:
  - "Can I STILL see this problem in the new image?"
  - "Is the padding between text and border AT LEAST 0.25 inches (a visible gap)?"
  - "Is the number truly CENTERED in the circle, or still off to one side?"
  - "Is the text clearly separated from the icon, or still overlapping?"

If you answer "still a problem" for ANY issue → the fix was not bold enough.
DOUBLE the move amount and apply again. Then verify_layout again.

**Phase 4: Second verify and final check**
After the second round of fixes, verify_layout one more time.
Only call submit when you are genuinely satisfied that every issue is resolved.

CRITICAL RULES:
- 914400 EMU = 1 inch. 0.3 inches = 274320 EMU.
- **Your first attempt is ALWAYS too small.** Plan for at least 2 rounds of fixes.
- **Minimum meaningful move: 0.2 inches (182880 EMU).** Anything smaller is invisible.
- **For centering numbers in circles:** The circle is about 0.4-0.5 inches in diameter. The number text box is about 0.3-0.4 inches. To center it, the number should be at (circle_center - text_width/2, circle_center_y - text_height/2). If unsure, move by 0.15in in each direction and verify.
- **For text-border padding:** Move AT LEAST 0.25 inches away from any border line. If after verify you can still see the text near the border, move another 0.15in.
- Minimum font size: 10pt
- Don't move shapes outside slide bounds
- Maximum 30 tool calls per slide"""


class PptxRepairAgent:
    """Agent-based repair for PPTX slides using tool-calling loop."""

    def __init__(
        self,
        llm: LLMClient,
        font_metrics: FontMetrics,
        render_fn,  # callable(pptx_path, slide_index) -> png_path
        model: str | None = None,
        max_tool_calls: int = 30,
    ):
        self._llm = llm
        self._fm = font_metrics
        self._render_fn = render_fn
        self._model = model
        self._max_tool_calls = max_tool_calls

    def repair_slide(
        self,
        pptx_path: Path,
        slide_index: int,
        issues: list[Issue],
        png_path: str,
        shape_data: str,
    ) -> int:
        """Repair one slide using agentic tool-calling loop.

        Returns number of tool calls made.
        """
        # Build initial message
        issue_text = "\n".join(
            f"- [{iss.issue_type}] {iss.why_this_fails}\n  SUGGESTED FIX: {iss.planned_fix}"
            for iss in issues
        )

        # Load initial screenshot
        try:
            img_b64 = image_to_base64(png_path)
        except Exception as e:
            logger.error("Failed to load PNG for slide %d: %s", slide_index, e)
            return 0

        user_msg = f"""Fix the spatial issues on slide {slide_index + 1}.

## Issues to fix:
{issue_text}

## Current shape positions:
{shape_data}

Look at the slide image carefully. Fix each issue by moving/resizing the text boxes.
After each fix, call verify_layout to check the result."""

        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": [
                {"type": "text", "text": user_msg},
                {"type": "image_url", "image_url": {"url": img_b64}},
            ]},
        ]

        # Tool-calling loop
        tool_calls_made = 0
        prs = Presentation(str(pptx_path))

        for _ in range(self._max_tool_calls):
            try:
                response = self._llm.client.chat.completions.create(
                    model=self._model or self._llm.default_model,
                    messages=messages,
                    tools=TOOLS,
                    tool_choice="auto",
                    max_completion_tokens=2048,
                    temperature=0.1,
                )
            except Exception as e:
                logger.error("LLM call failed in repair agent: %s", e)
                break

            msg = response.choices[0].message

            # If no tool calls, agent is done
            if not msg.tool_calls:
                break

            # Process tool calls
            messages.append(msg)  # Add assistant message with tool calls

            for tc in msg.tool_calls:
                tool_calls_made += 1
                fn_name = tc.function.name
                try:
                    args = json.loads(tc.function.arguments)
                except json.JSONDecodeError:
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": "Error: invalid JSON arguments",
                    })
                    continue

                result = self._execute_tool(
                    fn_name, args, prs, pptx_path, slide_index,
                )

                # verify_layout returns a special dict with image data
                if fn_name == "verify_layout" and isinstance(result, dict) and "image_b64" in result:
                    # Tool response is text-only
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": result["message"],
                    })
                    # Then add a user message with the actual image
                    messages.append({
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Here is the updated slide after your changes. Compare with the original image carefully. Is the padding sufficient? Are numbers centered in circles? Are texts aligned with icons? If any issue is NOT fully fixed, continue adjusting."},
                            {"type": "image_url", "image_url": {"url": result["image_b64"]}},
                        ],
                    })
                else:
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": result if isinstance(result, str) else json.dumps(result),
                    })

                # If submit was called, save and exit
                if fn_name == "submit":
                    prs.save(str(pptx_path))
                    logger.info(
                        "Slide %d repair done (%d tool calls): %s",
                        slide_index, tool_calls_made, args.get("summary", ""),
                    )
                    return tool_calls_made

        # Save even if we hit max calls
        prs.save(str(pptx_path))
        logger.info("Slide %d repair ended after %d tool calls (max reached)", slide_index, tool_calls_made)
        return tool_calls_made

    def _execute_tool(self, fn_name: str, args: dict, prs, pptx_path: Path, current_slide_idx: int):
        """Execute a single tool call."""
        if fn_name == "move_shape":
            return self._tool_move(prs, args)
        elif fn_name == "resize_shape":
            return self._tool_resize(prs, args)
        elif fn_name == "set_font_size":
            return self._tool_set_font(prs, args)
        elif fn_name == "get_shape_info":
            return self._tool_get_info(prs, args)
        elif fn_name == "verify_layout":
            # Save current state, re-render, return image as base64
            prs.save(str(pptx_path))
            si = args.get("slide_index", current_slide_idx)
            try:
                png_path = self._render_fn(str(pptx_path), si)
                img_b64 = image_to_base64(png_path)
                return {
                    "message": f"Slide {si+1} re-rendered. The updated image is shown below. Compare carefully with the original.",
                    "image_b64": img_b64,
                }
            except Exception as e:
                return {"status": "error", "message": str(e)}
        elif fn_name == "submit":
            return {"status": "ok", "message": "Repair submitted."}
        else:
            return {"status": "error", "message": f"Unknown tool: {fn_name}"}

    def _tool_move(self, prs, args: dict):
        slide_idx = args["slide_index"]
        shape_name = args["shape_name"]
        shape = self._find_shape(prs, slide_idx, shape_name)
        if not shape:
            return f"Error: shape '{shape_name}' not found on slide {slide_idx}"

        old_left, old_top = shape.left, shape.top
        if "left_emu" in args:
            shape.left = int(args["left_emu"])
        if "top_emu" in args:
            shape.top = int(args["top_emu"])

        return (
            f"Moved '{shape_name}' from ({old_left/914400:.2f}, {old_top/914400:.2f})in "
            f"to ({shape.left/914400:.2f}, {shape.top/914400:.2f})in"
        )

    def _tool_resize(self, prs, args: dict):
        slide_idx = args["slide_index"]
        shape_name = args["shape_name"]
        shape = self._find_shape(prs, slide_idx, shape_name)
        if not shape:
            return f"Error: shape '{shape_name}' not found on slide {slide_idx}"

        old_w, old_h = shape.width, shape.height
        if "width_emu" in args:
            shape.width = int(args["width_emu"])
        if "height_emu" in args:
            shape.height = int(args["height_emu"])

        return (
            f"Resized '{shape_name}' from ({old_w/914400:.2f}x{old_h/914400:.2f})in "
            f"to ({shape.width/914400:.2f}x{shape.height/914400:.2f})in"
        )

    def _tool_set_font(self, prs, args: dict):
        slide_idx = args["slide_index"]
        shape_name = args["shape_name"]
        shape = self._find_shape(prs, slide_idx, shape_name)
        if not shape:
            return f"Error: shape '{shape_name}' not found on slide {slide_idx}"
        if not shape.has_text_frame:
            return f"Error: '{shape_name}' has no text frame"

        target_pt = args["target_size_pt"]
        if target_pt < 10:
            return "Error: minimum font size is 10pt"

        target_emu = Pt(target_pt)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = target_emu

        return f"Set font size of '{shape_name}' to {target_pt}pt"

    def _tool_get_info(self, prs, args: dict):
        slide_idx = args["slide_index"]
        shape_name = args["shape_name"]
        shape = self._find_shape(prs, slide_idx, shape_name)
        if not shape:
            return f"Error: shape '{shape_name}' not found on slide {slide_idx}"

        info = {
            "name": shape.name,
            "left_in": round(shape.left / 914400, 2),
            "top_in": round(shape.top / 914400, 2),
            "width_in": round(shape.width / 914400, 2),
            "height_in": round(shape.height / 914400, 2),
            "left_emu": shape.left,
            "top_emu": shape.top,
            "width_emu": shape.width,
            "height_emu": shape.height,
        }
        if shape.has_text_frame:
            info["text"] = shape.text_frame.text[:100]
            fonts = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fonts.append(round(run.font.size / 12700, 1))
            if fonts:
                info["font_sizes_pt"] = fonts

        return info

    @staticmethod
    def _find_shape(prs, slide_index: int, shape_name: str):
        if slide_index < 0 or slide_index >= len(prs.slides):
            return None
        for shape in prs.slides[slide_index].shapes:
            if shape.name == shape_name:
                return shape
        return None
