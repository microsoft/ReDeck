"""PptxBuilder - low-level python-pptx construction utilities."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import design_tokens as dt

logger = logging.getLogger(__name__)


class PptxBuilder:
    """Low-level builder for creating PPTX slides using python-pptx."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = dt.SLIDE_WIDTH
        self.prs.slide_height = dt.SLIDE_HEIGHT

    def add_blank_slide(self):
        """Add a blank slide layout."""
        layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    def add_text_box(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        font_size: Pt | None = None,
        font_color: RGBColor | None = None,
        font_bold: bool = False,
        font_name: str | None = None,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        name: str = "",
    ):
        """Add a text box to a slide."""
        txBox = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        if name:
            txBox.name = name

        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment

        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = text
        run.font.size = font_size or dt.FONT_SIZE_BODY
        run.font.color.rgb = font_color or dt.COLOR_TEXT_DARK
        run.font.bold = font_bold
        run.font.name = font_name or dt.FONT_FAMILY

        return txBox

    def add_title(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        name: str = "title",
    ):
        """Add a title text box."""
        return self.add_text_box(
            slide, left, top, width, height, text,
            font_size=dt.FONT_SIZE_TITLE,
            font_color=dt.COLOR_PRIMARY,
            font_bold=True,
            name=name,
        )

    def add_bullet_list(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        items: list[str],
        font_size: Pt | None = None,
        name: str = "",
    ):
        """Add a bullet list text box."""
        txBox = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        if name:
            txBox.name = name

        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = dt.PARAGRAPH_SPACING_PT
            for run in p.runs:
                run.font.size = font_size or dt.FONT_SIZE_BODY
                run.font.name = dt.FONT_FAMILY
                run.font.color.rgb = dt.COLOR_TEXT_DARK

        return txBox

    def add_table(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        headers: list[str],
        rows: list[list[str]],
        name: str = "",
    ):
        """Add a table to a slide."""
        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)
        row_height = height // n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        if name:
            table_shape.name = name

        table = table_shape.table

        # Header row
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            for p in cell.text_frame.paragraphs:
                p.font.size = dt.FONT_SIZE_CAPTION
                p.font.bold = True
                p.font.color.rgb = dt.COLOR_BG_WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = dt.COLOR_PRIMARY

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = dt.FONT_SIZE_CAPTION
                    p.font.color.rgb = dt.COLOR_TEXT_DARK

        return table_shape

    def add_shape_with_text(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        fill_color: RGBColor | None = None,
        font_size: Pt | None = None,
        name: str = "",
    ):
        """Add a rounded rectangle shape with text (for metric chips, callouts)."""
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        if name:
            shape.name = name

        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color or dt.COLOR_BG_LIGHT

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER

        for run in p.runs:
            run.font.size = font_size or dt.FONT_SIZE_BODY
            run.font.color.rgb = dt.COLOR_TEXT_DARK
            run.font.name = dt.FONT_FAMILY

        return shape

    def add_image(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        image_path: str,
        name: str = "",
        maintain_aspect: bool = True,
    ):
        """Add an image to a slide.

        Args:
            slide: The slide to add the image to.
            left, top, width, height: Bounding box in EMU.
            image_path: Path to the image file on disk.
            name: Optional shape name.
            maintain_aspect: If True, scale image to fit within bbox
                while preserving aspect ratio (centered).
        """
        from PIL import Image as PILImage
        from ..python_pptx import design_tokens as _dt  # avoid circular

        img_path = Path(image_path)
        if not img_path.exists():
            logger.warning("Image not found: %s, adding placeholder", image_path)
            return self.add_text_box(
                slide, left, top, width, height,
                f"[Image: {img_path.name}]",
                name=name,
            )

        if maintain_aspect:
            try:
                img = PILImage.open(image_path)
                img_w, img_h = img.size
                img.close()

                aspect = img_w / img_h
                bbox_aspect = width / height if height > 0 else 1

                if bbox_aspect > aspect:
                    # Height-constrained
                    new_h = height
                    new_w = int(height * aspect)
                else:
                    # Width-constrained
                    new_w = width
                    new_h = int(width / aspect)

                # Center within bbox
                left += (width - new_w) // 2
                top += (height - new_h) // 2
                width, height = new_w, new_h
            except Exception as e:
                logger.warning("Could not read image for aspect ratio: %s", e)

        pic = slide.shapes.add_picture(
            str(image_path),
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        if name:
            pic.name = name
        return pic

    def save(self, path: str | Path) -> None:
        """Save the presentation to disk."""
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        logger.info("Saved PPTX to %s", path)
