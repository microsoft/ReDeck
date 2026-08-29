"""Safety-net geometry post-processing for python-pptx slides.

Lightweight fixes that catch common LLM code-generation format issues:
- Raw-inch coordinates (not wrapped in Inches())
- Shapes overflowing slide boundaries
- Sub-minimum font sizes

Overlap resolution, text overflow, and accent conflict handling are
delegated to the LLM via prompt-injected geometry rules.
"""

import logging

from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)


def fix_raw_coordinates(slide):
    """Fix coordinates that were set as raw inches/floats instead of EMU.

    When LLM code does `slide.shapes.add_textbox(Inches(0.5), 1.9, ...)`
    (mixing Inches() and raw floats), python-pptx writes the raw float
    directly to XML as e.g. y="1.9000000000000001", which crashes when
    the PPTX is later read back (int() on a float string).

    Detection heuristic: EMU values for a 13.333x7.5" slide range from 0
    to ~12,192,000. If a position/size attribute is > 0 and < 100, it's
    almost certainly a raw inch value (0-13.333) or a raw Pt value, not
    a valid EMU. Convert by multiplying by 914400 (EMU_PER_INCH).

    Also fix any float values that python-pptx might have stored.
    """
    EMU_PER_INCH = 914400
    MAX_RAW_INCHES = 15  # max plausible raw-inch value

    for shape in slide.shapes:
        try:
            for attr in ("left", "top", "width", "height"):
                val = getattr(shape, attr, None)
                if val is None:
                    continue
                # Check for float values (should always be int EMU)
                if isinstance(val, float):
                    if val < MAX_RAW_INCHES and val > 0:
                        setattr(shape, attr, round(val * EMU_PER_INCH))
                    else:
                        setattr(shape, attr, round(val))
                elif isinstance(val, int) and 0 < val < 100:
                    # Suspiciously small — likely raw inches as int
                    # (e.g., int(1.9) = 1)
                    setattr(shape, attr, val * EMU_PER_INCH)
        except Exception:
            pass

    # Also fix via XML directly — catches table frames and grouped shapes
    # whose attributes may not be exposed via shape.left/top
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for elem in slide._element.iter():
        tag = etree.QName(elem.tag).localname if isinstance(elem.tag, str) else ""
        if tag in ("off", "ext"):
            for coord in ("x", "y", "cx", "cy"):
                val_str = elem.get(coord)
                if val_str is None:
                    continue
                try:
                    # Already a valid integer string
                    int(val_str)
                except ValueError:
                    # It's a float string like "1.9000000000000001"
                    try:
                        float_val = float(val_str)
                        if float_val < MAX_RAW_INCHES:
                            elem.set(coord, str(round(float_val * EMU_PER_INCH)))
                        else:
                            elem.set(coord, str(round(float_val)))
                    except ValueError:
                        pass


def clamp_shapes(slide):
    """Clamp all shapes to stay within slide boundaries.

    Fixes the common issue where metric cards or text boxes extend
    beyond the right edge (x+width > 13.333") or bottom (y+height > 7.5").
    """
    max_right = Inches(13.0)   # leave 0.333" right margin
    max_bottom = Inches(7.2)   # leave 0.3" bottom margin
    slide_w = Inches(13.333)
    slide_h = Inches(7.5)
    min_dim = Inches(0.3)

    for shape in slide.shapes:
        try:
            right = shape.left + shape.width
            bottom = shape.top + shape.height

            # Clamp right overflow
            if right > slide_w:
                overflow = right - max_right
                if overflow > 0:
                    # First try to reduce width
                    new_width = shape.width - overflow
                    if new_width > Inches(0.5):
                        shape.width = round(new_width)
                    else:
                        # Move shape left
                        new_left = max_right - shape.width
                        shape.left = max(0, round(new_left))

            # Clamp bottom overflow
            if bottom > slide_h:
                overflow = bottom - max_bottom
                if overflow > 0:
                    new_height = shape.height - overflow
                    if new_height > min_dim:
                        shape.height = round(new_height)
                    else:
                        # Move shape up
                        new_top = max_bottom - shape.height
                        shape.top = max(0, round(new_top))
        except Exception:
            pass  # skip shapes that can't be modified (e.g. grouped shapes)


def enforce_min_font_size(slide, min_pt: int = 14):
    """Bump any text font below min_pt to min_pt.

    Evaluator flags fonts below 14pt as 'font_too_small'. This safety net
    catches cases where the LLM uses sub-14pt fonts despite prompt rules.
    Skips chart objects (their internal axis/legend labels are acceptable
    at smaller sizes).
    """
    from pptx.util import Pt as _Pt

    min_emu = _Pt(min_pt)
    for shape in slide.shapes:
        # Skip chart shapes — their axis/legend labels can be < 14pt
        if shape.has_chart:
            continue
        if not shape.has_text_frame:
            continue
        try:
            for paragraph in shape.text_frame.paragraphs:
                # Check paragraph-level font
                if paragraph.font.size is not None and paragraph.font.size < min_emu:
                    paragraph.font.size = min_emu
                # Check run-level fonts
                for run in paragraph.runs:
                    if run.font.size is not None and run.font.size < min_emu:
                        run.font.size = min_emu
        except Exception:
            pass


def apply_all(slide):
    """Apply safety-net geometry fixes only.

    Overlap resolution, text overflow, and accent conflicts are now
    handled by the LLM during codegen/repair via prompt-injected rules.
    This function only performs non-destructive format corrections.
    """
    fix_raw_coordinates(slide)
    clamp_shapes(slide)
    enforce_min_font_size(slide)
