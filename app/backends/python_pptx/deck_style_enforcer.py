"""Deck-level style enforcement for cross-slide visual consistency.

This module addresses the `visual_inconsistency` issue type (rubric B1) that
cannot be fixed by per-slide repair. It operates on the saved PPTX after all
slides have been generated, scanning for the "majority style" and normalizing
deviants.

Key normalizations:
1. Title zone: consistent font size, color, position, bold/not across all slides
2. Body text: consistent font size across content paragraphs
3. Background: slide 1 is cover (may be dark), slides 2+ use majority background
4. Accent bars: consistent color for thin accent/divider shapes
5. Font family: enforce the theme font across all text

This runs between codegen/recompile and render in the pipeline.
"""

import logging
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from ...schemas.issue_types import SlideDimensions

logger = logging.getLogger(__name__)

# Constants
EMU_PER_INCH = 914400
SLIDE_W = int(13.333 * EMU_PER_INCH)
SLIDE_H = int(7.5 * EMU_PER_INCH)

# Title zone: shapes in the top 1.5" with width > 8"
TITLE_TOP_MAX = int(1.5 * EMU_PER_INCH)
TITLE_MIN_WIDTH = int(8.0 * EMU_PER_INCH)

# Accent bar: thin shapes (height < 0.15" or width < 0.15") spanning > 3"
ACCENT_MAX_THIN = int(0.15 * EMU_PER_INCH)
ACCENT_MIN_SPAN = int(3.0 * EMU_PER_INCH)

# Background detection: shapes covering > 40% of slide area with dark fill
# are treated as "dark background" (catches header bands, half-slide fills)
BG_AREA_THRESHOLD = 0.40

# Body text: shapes below title zone, with text
BODY_TOP_MIN = int(1.0 * EMU_PER_INCH)

# Font size tolerance for "majority" detection (EMU)
FONT_SIZE_TOLERANCE = Pt(2)  # ±2pt is acceptable


def _rgb_tuple(color) -> tuple[int, int, int] | None:
    """Extract (r, g, b) tuple from a pptx color, or None if not solid RGB."""
    try:
        if color is None:
            return None
        if hasattr(color, 'rgb') and color.rgb is not None:
            rgb = color.rgb
            return (rgb[0], rgb[1], rgb[2]) if rgb else None
        if hasattr(color, 'type') and color.type is not None:
            from pptx.enum.dml import MSO_THEME_COLOR
            # Theme color — we can't easily resolve, skip
            return None
    except Exception:
        return None
    return None


def _shape_area(shape) -> int:
    """Return shape area in EMU²."""
    try:
        return shape.width * shape.height
    except Exception:
        return 0


def _is_title_shape(shape) -> bool:
    """Check if shape is likely a title (top zone, wide, has text)."""
    try:
        if not shape.has_text_frame:
            return False
        if shape.top > TITLE_TOP_MAX:
            return False
        if shape.width < TITLE_MIN_WIDTH:
            return False
        if not shape.text_frame.text.strip():
            return False
        return True
    except Exception:
        return False


def _is_accent_bar(shape) -> bool:
    """Check if shape is a thin accent/divider bar."""
    try:
        w, h = shape.width, shape.height
        thin_dim = min(w, h)
        long_dim = max(w, h)
        if thin_dim > ACCENT_MAX_THIN:
            return False
        if long_dim < ACCENT_MIN_SPAN:
            return False
        return True
    except Exception:
        return False


def _is_bg_shape(shape) -> bool:
    """Check if shape is a full-slide background rectangle."""
    try:
        area = _shape_area(shape)
        slide_area = SLIDE_W * SLIDE_H
        if area / slide_area < BG_AREA_THRESHOLD:
            return False
        name = shape.name.lower()
        # Usually rectangles used as backgrounds
        return True
    except Exception:
        return False


def _get_fill_color(shape) -> tuple[int, int, int] | None:
    """Get the solid fill color of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None:
            # Check for solid fill
            from pptx.enum.dml import MSO_FILL
            if fill.type == MSO_FILL.SOLID:
                return _rgb_tuple(fill.fore_color)
    except Exception:
        pass
    return None


def _get_title_font_info(shape) -> dict | None:
    """Extract title font properties from a title shape."""
    try:
        if not shape.has_text_frame:
            return None
        tf = shape.text_frame
        # Get the first non-empty paragraph
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Get font size from runs or paragraph
            font_size = None
            font_color = None
            font_bold = None
            font_name = None

            if para.font.size:
                font_size = para.font.size
            if para.font.bold is not None:
                font_bold = para.font.bold
            if para.font.name:
                font_name = para.font.name
            font_color = _rgb_tuple(para.font.color)

            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size
                if run.font.bold is not None:
                    font_bold = run.font.bold
                if run.font.name:
                    font_name = run.font.name
                c = _rgb_tuple(run.font.color)
                if c:
                    font_color = c
                break  # first run is enough

            if font_size:
                return {
                    "size": font_size,
                    "color": font_color,
                    "bold": font_bold,
                    "name": font_name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                }
    except Exception:
        pass
    return None


def _set_title_font(shape, target: dict):
    """Set title font properties to match target."""
    try:
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        for para in tf.paragraphs:
            if not para.text.strip():
                continue
            # Set paragraph-level font
            if target.get("size"):
                para.font.size = target["size"]
            if target.get("bold") is not None:
                para.font.bold = target["bold"]
            if target.get("name"):
                para.font.name = target["name"]
            if target.get("color"):
                r, g, b = target["color"]
                para.font.color.rgb = RGBColor(r, g, b)
            # Set run-level font
            for run in para.runs:
                if target.get("size"):
                    run.font.size = target["size"]
                if target.get("bold") is not None:
                    run.font.bold = target["bold"]
                if target.get("name"):
                    run.font.name = target["name"]
                if target.get("color"):
                    r, g, b = target["color"]
                    run.font.color.rgb = RGBColor(r, g, b)
    except Exception as e:
        logger.debug("Failed to set title font: %s", e)


def _slide_has_dark_bg(slide) -> bool:
    """Check if a slide has a dark background (either fill or bg shape)."""
    try:
        # Check slide background fill
        bg = slide.background
        fill = bg.fill
        from pptx.enum.dml import MSO_FILL
        if fill.type == MSO_FILL.SOLID:
            color = _rgb_tuple(fill.fore_color)
            if color:
                r, g, b = color
                luminance = 0.299 * r + 0.587 * g + 0.114 * b
                if luminance < 100:
                    return True
    except Exception:
        pass

    # Check for large background shapes with dark fill
    for shape in slide.shapes:
        if _is_bg_shape(shape):
            color = _get_fill_color(shape)
            if color:
                r, g, b = color
                luminance = 0.299 * r + 0.587 * g + 0.114 * b
                if luminance < 100:
                    return True
    return False


def _collect_body_font_sizes(slide) -> list[int]:
    """Collect body text font sizes from a slide (below title zone)."""
    sizes = []
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            if shape.top < BODY_TOP_MIN:
                continue
            if _is_accent_bar(shape):
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or len(text) < 5:
                    continue
                # Get font size
                fs = None
                if para.font.size:
                    fs = para.font.size
                for run in para.runs:
                    if run.font.size:
                        fs = run.font.size
                        break
                if fs and Pt(10) <= fs <= Pt(24):
                    sizes.append(fs)
        except Exception:
            continue
    return sizes


def _enforce_font_family(slide, font_family: str):
    """Set all text in a slide to the specified font family."""
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.font.name and para.font.name != font_family:
                    para.font.name = font_family
                for run in para.runs:
                    if run.font.name and run.font.name != font_family:
                        run.font.name = font_family
        except Exception:
            continue


def enforce(pptx_path: str | Path, font_family: str | None = None) -> dict:
    """Enforce cross-slide visual consistency on a saved PPTX.

    Args:
        pptx_path: Path to the PPTX file (will be overwritten)
        font_family: Optional font family to enforce (from theme)

    Returns:
        dict with enforcement statistics
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        logger.warning("PPTX not found: %s", pptx_path)
        return {"status": "skipped", "reason": "file_not_found"}

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    n_slides = len(slides)

    if n_slides < 3:
        logger.info("DeckStyleEnforcer: only %d slides, skipping", n_slides)
        return {"status": "skipped", "reason": "too_few_slides"}

    stats = {
        "titles_normalized": 0,
        "accent_bars_normalized": 0,
        "bg_shapes_removed": 0,
        "font_family_fixed": 0,
        "body_font_normalized": 0,
    }

    # ===================================================================
    # Phase 1: Collect majority styles from slides 2+ (skip cover slide)
    # ===================================================================
    content_slides = slides[1:]  # slides 2, 3, ..., N (0-indexed: slides[1:])

    # 1a. Collect title font info from content slides
    title_infos = []
    title_info_by_slide = {}
    for idx, slide in enumerate(content_slides):
        for shape in slide.shapes:
            if _is_title_shape(shape):
                info = _get_title_font_info(shape)
                if info:
                    title_infos.append(info)
                    title_info_by_slide[idx + 1] = info  # 1-indexed slide number
                break  # only first title per slide

    # 1b. Determine majority title style
    majority_title = None
    if len(title_infos) >= 2:
        # Vote on font size (most common)
        size_counter = Counter(info["size"] for info in title_infos if info.get("size"))
        if size_counter:
            majority_size = size_counter.most_common(1)[0][0]
        else:
            majority_size = Pt(30)

        # Vote on color (most common among content slides)
        color_counter = Counter(
            info["color"] for info in title_infos if info.get("color")
        )
        majority_color = color_counter.most_common(1)[0][0] if color_counter else None

        # Vote on bold
        bold_counter = Counter(
            info["bold"] for info in title_infos if info.get("bold") is not None
        )
        majority_bold = bold_counter.most_common(1)[0][0] if bold_counter else True

        # Vote on font name
        name_counter = Counter(
            info["name"] for info in title_infos if info.get("name")
        )
        majority_name = name_counter.most_common(1)[0][0] if name_counter else font_family

        majority_title = {
            "size": majority_size,
            "color": majority_color,
            "bold": majority_bold,
            "name": majority_name or font_family,
        }
        logger.info(
            "DeckStyleEnforcer: majority title — size=%s, color=%s, bold=%s, font=%s",
            majority_size, majority_color, majority_bold, majority_name,
        )

    # 1c. Collect accent bar colors from content slides
    accent_colors = []
    for slide in content_slides:
        for shape in slide.shapes:
            if _is_accent_bar(shape):
                color = _get_fill_color(shape)
                if color:
                    accent_colors.append(color)

    accent_color_counter = Counter(accent_colors)
    majority_accent = (
        accent_color_counter.most_common(1)[0][0]
        if accent_color_counter else None
    )

    # 1d. Determine majority background type (dark vs light) for content slides
    dark_count = sum(1 for s in content_slides if _slide_has_dark_bg(s))
    light_count = len(content_slides) - dark_count
    majority_is_light = light_count >= dark_count

    # 1e. Collect body font sizes
    all_body_sizes = []
    for slide in content_slides:
        all_body_sizes.extend(_collect_body_font_sizes(slide))

    body_size_counter = Counter(all_body_sizes)
    majority_body_size = (
        body_size_counter.most_common(1)[0][0]
        if body_size_counter else None
    )

    logger.info(
        "DeckStyleEnforcer: majority bg=%s, accent=%s, body_size=%s, %d title infos",
        "light" if majority_is_light else "dark",
        majority_accent,
        majority_body_size,
        len(title_infos),
    )

    # ===================================================================
    # Phase 2: Enforce majority style on all content slides (2+)
    # ===================================================================

    for slide_idx, slide in enumerate(content_slides):
        slide_num = slide_idx + 2  # 1-indexed, starting from slide 2

        # 2a. Enforce title font consistency
        if majority_title:
            for shape in slide.shapes:
                if _is_title_shape(shape):
                    current = _get_title_font_info(shape)
                    if current:
                        needs_fix = False
                        # Check if title style deviates from majority
                        if current.get("size") and majority_title.get("size"):
                            diff = abs(current["size"] - majority_title["size"])
                            if diff > FONT_SIZE_TOLERANCE:
                                needs_fix = True
                        if current.get("color") != majority_title.get("color"):
                            needs_fix = True
                        if current.get("bold") != majority_title.get("bold"):
                            needs_fix = True

                        if needs_fix:
                            _set_title_font(shape, majority_title)
                            stats["titles_normalized"] += 1
                            logger.debug(
                                "Slide %d: normalized title font (was size=%s color=%s)",
                                slide_num, current.get("size"), current.get("color"),
                            )
                    break

        # 2b. Enforce accent bar color consistency
        if majority_accent:
            for shape in slide.shapes:
                if _is_accent_bar(shape):
                    current_color = _get_fill_color(shape)
                    if current_color and current_color != majority_accent:
                        try:
                            r, g, b = majority_accent
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(r, g, b)
                            stats["accent_bars_normalized"] += 1
                            logger.debug(
                                "Slide %d: normalized accent bar from %s to %s",
                                slide_num, current_color, majority_accent,
                            )
                        except Exception as e:
                            logger.debug("Failed to set accent color: %s", e)

        # 2c. Handle rogue dark backgrounds on content slides
        # If majority is light but this slide has a dark bg shape, remove it
        if majority_is_light and _slide_has_dark_bg(slide):
            # Remove large dark background shapes (not the slide background itself)
            shapes_to_check = list(slide.shapes)
            for shape in shapes_to_check:
                if _is_bg_shape(shape):
                    color = _get_fill_color(shape)
                    if color:
                        r, g, b = color
                        luminance = 0.299 * r + 0.587 * g + 0.114 * b
                        if luminance < 100:
                            # This is a dark bg shape on a light-majority deck
                            # Convert to light_bg instead of removing
                            try:
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor(236, 240, 241)
                                stats["bg_shapes_removed"] += 1
                                logger.info(
                                    "Slide %d: converted dark bg shape to light_bg",
                                    slide_num,
                                )
                                # Fix title color since it was likely white for dark bg
                                for s in slide.shapes:
                                    if _is_title_shape(s) and majority_title:
                                        _set_title_font(s, majority_title)
                                # Fix white body text that was designed for dark bg
                                for s in slide.shapes:
                                    if not s.has_text_frame:
                                        continue
                                    for para in s.text_frame.paragraphs:
                                        _fix_white_text(para)
                            except Exception as e:
                                logger.debug("Failed to fix dark bg: %s", e)

            # Also check slide-level background fill
            try:
                bg = slide.background
                fill = bg.fill
                from pptx.enum.dml import MSO_FILL
                if fill.type == MSO_FILL.SOLID:
                    color = _rgb_tuple(fill.fore_color)
                    if color:
                        r, g, b = color
                        luminance = 0.299 * r + 0.587 * g + 0.114 * b
                        if luminance < 100:
                            # Dark slide background on light-majority deck
                            fill.background()  # reset to no fill (transparent/white)
                            stats["bg_shapes_removed"] += 1
                            logger.info(
                                "Slide %d: cleared dark slide background",
                                slide_num,
                            )
                            # Fix title color
                            for s in slide.shapes:
                                if _is_title_shape(s) and majority_title:
                                    _set_title_font(s, majority_title)
                            # Fix white body text that was for dark bg
                            for s in slide.shapes:
                                if not s.has_text_frame:
                                    continue
                                for para in s.text_frame.paragraphs:
                                    _fix_white_text(para)
            except Exception as e:
                logger.debug("Failed to check slide bg: %s", e)

        # 2d. Enforce font family
        if font_family:
            _enforce_font_family(slide, font_family)

    # ===================================================================
    # Phase 3: Enforce title style on cover slide (slide 1) too
    # ===================================================================
    # Cover slide may be dark — that's OK, but title font size and bold should match
    if majority_title and len(slides) > 0:
        cover_slide = slides[0]
        for shape in cover_slide.shapes:
            if _is_title_shape(shape):
                current = _get_title_font_info(shape)
                if current and current.get("size"):
                    diff = abs(current["size"] - majority_title["size"])
                    if diff > Pt(4):
                        # Only enforce size and bold on cover (not color — it may be white on dark)
                        cover_target = {
                            "size": majority_title["size"],
                            "bold": majority_title["bold"],
                            "name": majority_title.get("name"),
                        }
                        _set_title_font(shape, cover_target)
                        stats["titles_normalized"] += 1
                break

        # Enforce font family on cover slide too
        if font_family:
            _enforce_font_family(cover_slide, font_family)

    # ===================================================================
    # Phase 4: Save
    # ===================================================================
    if sum(stats.values()) > 0:
        prs.save(str(pptx_path))
        logger.info(
            "DeckStyleEnforcer: saved %s — %d titles, %d accents, %d bgs normalized",
            pptx_path.name,
            stats["titles_normalized"],
            stats["accent_bars_normalized"],
            stats["bg_shapes_removed"],
        )
    else:
        logger.info("DeckStyleEnforcer: no changes needed for %s", pptx_path.name)

    stats["status"] = "ok"
    return stats


def audit_shape_diversity(pptx_path: str | Path) -> dict:
    """Audit shape diversity across slides to detect overuse of ROUNDED_RECTANGLE.

    Args:
        pptx_path: Path to the PPTX file to audit

    Returns:
        dict with shape counts and warnings for monitoring
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        logger.warning("PPTX not found for audit: %s", pptx_path)
        return {"status": "skipped", "reason": "file_not_found"}

    try:
        prs = Presentation(str(pptx_path))
        slides = list(prs.slides)

        audit_results = {
            "status": "ok",
            "total_slides": len(slides),
            "slides_with_violations": 0,
            "shape_counts": {},
            "per_slide_counts": [],
            "warnings": []
        }

        for slide_idx, slide in enumerate(slides):
            slide_num = slide_idx + 1
            slide_shapes = {
                "ROUNDED_RECTANGLE": 0,
                "RECTANGLE": 0,
                "OVAL": 0,
                "LINE": 0,
                "OTHER": 0
            }

            for shape in slide.shapes:
                try:
                    # Check shape type by auto_shape_type
                    shape_name = "OTHER"
                    if hasattr(shape, 'auto_shape_type'):
                        from pptx.enum.shapes import MSO_SHAPE
                        if shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                            shape_name = "ROUNDED_RECTANGLE"
                        elif shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                            shape_name = "RECTANGLE"
                        elif shape.auto_shape_type == MSO_SHAPE.OVAL:
                            shape_name = "OVAL"
                        elif hasattr(MSO_SHAPE, 'LINE') and shape.auto_shape_type == MSO_SHAPE.LINE:
                            shape_name = "LINE"

                    slide_shapes[shape_name] += 1

                except Exception:
                    # If we can't determine shape type, count as OTHER
                    slide_shapes["OTHER"] += 1

            # Check for violations (>2 ROUNDED_RECTANGLE per slide)
            rounded_rect_count = slide_shapes["ROUNDED_RECTANGLE"]
            if rounded_rect_count > 2:
                audit_results["slides_with_violations"] += 1
                audit_results["warnings"].append(
                    f"Slide {slide_num}: {rounded_rect_count} ROUNDED_RECTANGLE shapes "
                    f"(limit: 2). Consider using RECTANGLE, LINE, OVAL, or pure textboxes for variety."
                )
                logger.warning(
                    "Shape diversity violation on slide %d: %d ROUNDED_RECTANGLE shapes (max 2 recommended)",
                    slide_num, rounded_rect_count
                )

            audit_results["per_slide_counts"].append({
                "slide": slide_num,
                "shapes": slide_shapes.copy()
            })

            # Accumulate total counts
            for shape_type, count in slide_shapes.items():
                audit_results["shape_counts"][shape_type] = audit_results["shape_counts"].get(shape_type, 0) + count

        # Log summary
        total_rounded = audit_results["shape_counts"].get("ROUNDED_RECTANGLE", 0)
        logger.info(
            "Shape diversity audit: %d slides, %d violations, %d total ROUNDED_RECTANGLE shapes",
            len(slides), audit_results["slides_with_violations"], total_rounded
        )

        return audit_results

    except Exception as e:
        logger.error("Failed to audit shape diversity: %s", e)
        return {"status": "error", "reason": str(e)}


def audit_chart_dimensions(pptx_path: str | Path) -> dict:
    """Audit chart dimensions across slides to detect undersized charts.

    Checks:
    - Charts below 3.0" height (minimum readable height)
    - Charts with aspect ratio > 3.0 (too flat, bars compressed)
    - Charts below 7.0" width (minimum recommended)

    Args:
        pptx_path: Path to the PPTX file to audit

    Returns:
        dict with chart info and warnings
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        logger.warning("PPTX not found for chart audit: %s", pptx_path)
        return {"status": "skipped", "reason": "file_not_found"}

    try:
        prs = Presentation(str(pptx_path))
        slides = list(prs.slides)

        audit_results = {
            "status": "ok",
            "total_charts": 0,
            "violations": 0,
            "charts": [],
            "warnings": [],
        }

        for slide_idx, slide in enumerate(slides):
            slide_num = slide_idx + 1
            for shape in slide.shapes:
                if not shape.has_chart:
                    continue

                audit_results["total_charts"] += 1
                # Convert EMU to inches
                w_in = shape.width / SlideDimensions.EMU_PER_INCH
                h_in = shape.height / SlideDimensions.EMU_PER_INCH
                aspect = w_in / h_in if h_in > 0 else 999

                chart_info = {
                    "slide": slide_num,
                    "width": round(w_in, 2),
                    "height": round(h_in, 2),
                    "aspect_ratio": round(aspect, 2),
                    "issues": [],
                }

                if h_in < 3.0:
                    chart_info["issues"].append("height_too_short")
                    audit_results["violations"] += 1
                    audit_results["warnings"].append(
                        f"Slide {slide_num}: Chart height {h_in:.1f}\" < 3.0\" minimum. "
                        f"Chart appears squashed and hard to read."
                    )
                    logger.warning(
                        "Chart height violation on slide %d: %.1f\" (min 3.0\")",
                        slide_num, h_in,
                    )

                if aspect > 3.0:
                    chart_info["issues"].append("aspect_too_flat")
                    audit_results["violations"] += 1
                    audit_results["warnings"].append(
                        f"Slide {slide_num}: Chart aspect ratio {aspect:.1f}:1 > 3.0:1 max. "
                        f"Chart is too wide and flat — increase height or decrease width."
                    )

                if w_in < 7.0:
                    chart_info["issues"].append("width_too_narrow")
                    audit_results["warnings"].append(
                        f"Slide {slide_num}: Chart width {w_in:.1f}\" < 7.0\" recommended minimum."
                    )

                audit_results["charts"].append(chart_info)

        logger.info(
            "Chart dimension audit: %d charts, %d violations",
            audit_results["total_charts"], audit_results["violations"],
        )
        return audit_results

    except Exception as e:
        logger.error("Failed to audit chart dimensions: %s", e)
        return {"status": "error", "reason": str(e)}


def _fix_white_text(paragraph):
    """Fix white text (from dark bg) to dark text."""
    try:
        for run in paragraph.runs:
            c = _rgb_tuple(run.font.color)
            if c:
                r, g, b = c
                luminance = 0.299 * r + 0.587 * g + 0.114 * b
                if luminance > 240:
                    # White or near-white text — change to dark
                    run.font.color.rgb = RGBColor(44, 62, 80)
        # Also check paragraph-level
        c = _rgb_tuple(paragraph.font.color)
        if c:
            r, g, b = c
            luminance = 0.299 * r + 0.587 * g + 0.114 * b
            if luminance > 240:
                paragraph.font.color.rgb = RGBColor(44, 62, 80)
    except Exception:
        pass
