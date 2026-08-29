"""HtmlCodeGenCompiler - LLM-based HTML code generation compiler for slides.

Instead of generating python-pptx code, this compiler asks the LLM to write
HTML/CSS for each slide, renders it to a high-res PNG via Playwright, then
inserts each PNG as a full-bleed image into a PPTX.

This gives the model full control over layout and visual design using the
full expressiveness of HTML/CSS, while still producing a standard PPTX file.
"""

import logging
import os
import re
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ...llm_client import LLMClient
from ...schemas.blueprint import DeckBlueprint, BlueprintSlide
from ...schemas.evidence import EvidenceState

from ...themes import (
    COMPOSITION_VARIANT_VERSION,
    THEME_FAMILIES,
    THEME_REGISTRY,
    DEFAULT_THEME,
    CompositionVariant,
    LayoutGrammar,
    ThemeColors,
    format_composition_variant_contract,
    format_html_design_contract,
    format_theme_colors_for_prompt,
    format_theme_typography_for_prompt,
    select_composition_variant,
    select_layout_grammar,
)
from ..html_codegen.visual_skills import (
    VISUAL_SKILL_LIBRARY_VERSION,
    VisualSkill,
    format_visual_skill_references,
    select_visual_skills,
)
from ..html_codegen.deck_frame_contract import enforce_html_slide_frame_contract
from ...utils.io_utils import read_text
from ...modules.chart_generator import ChartGenerator

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Viewport dimensions for Playwright rendering
VIEWPORT_W = 1280
VIEWPORT_H = 720
DEVICE_SCALE_FACTOR = 2  # 2x for crisp images → 2560×1440 PNG

HTML_CODEGEN_PROMPT_PATH = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "slide_html_codegen.system.md"
HTML_REPAIR_PROMPT_PATH = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "slide_html_repair.system.md"
HTML_CODEGEN_PROMPT_VERSION = "slide_html_codegen.hybrid.v1"

MAX_CODE_RETRIES = 3

# Patterns that indicate LLM meta-instructions leaked into slide content
_META_TEXT_PATTERNS = [
    re.compile(r'<[^>]*>[^<]*(?:not specified in (?:the )?provided (?:packet|context|material|document)s?)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:venue[/,]?\s*year\s+not\s+(?:specified|available|provided))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:otherwise\s+note\s+)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:if (?:available|provided)[\s,;]*otherwise)[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:The slide (?:emphasizes|shows|presents|illustrates|demonstrates))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:This slide (?:emphasizes|shows|presents|illustrates|demonstrates))[^<]*</[^>]*>', re.IGNORECASE),
    re.compile(r'<[^>]*>[^<]*(?:should be presented as)[^<]*</[^>]*>', re.IGNORECASE),
    # Non-slide meta content: "Flow:", "Transition:", presentation structure descriptions
    re.compile(r'<strong[^>]*>\s*Flow:\s*</strong>[^<]*(?:</[^>]*>)?', re.IGNORECASE),
    re.compile(r'<strong[^>]*>\s*Transition:\s*</strong>[^<]*(?:</[^>]*>)?', re.IGNORECASE),
    re.compile(r'<strong[^>]*>\s*Note:\s*</strong>\s*(?:This (?:slide|deck|presentation)|The (?:slide|deck|presentation))[^<]*(?:</[^>]*>)?', re.IGNORECASE),
]


def _find_page_for_must_cover(item: str, bp_slide, source_store) -> int | None:
    """Find the most likely page number for a must_cover item by keyword matching.

    Searches through the slide's source bundle atomic blocks for text matching
    key terms in the must_cover item. Returns the page number of the best match.
    """
    if source_store is None:
        return None
    # Extract numbers and key terms from the must_cover item
    import re as _re
    terms = set()
    # Extract dollar amounts, percentages, specific numbers
    for num in _re.findall(r'[\$]?[\d,]+\.?\d*[%BMK]?', item):
        cleaned = num.replace(',', '').replace('$', '')
        if len(cleaned) > 1:
            terms.add(cleaned)
    # Extract significant words
    for word in item.split():
        cleaned = _re.sub(r'[^a-zA-Z]', '', word).lower()
        if len(cleaned) > 3:
            terms.add(cleaned)

    if not terms:
        return None

    # Search through bundle's atomic block IDs
    bundle = source_store.get_bundle(bp_slide.slide_id)
    if not bundle:
        return None

    best_page = None
    best_score = 0
    for bid in bundle.source_atomic_block_ids:
        block = source_store.get_block(bid)
        if not block or not block.page:
            continue
        block_text = block.text.lower().replace(',', '')
        score = sum(1 for t in terms if t.lower() in block_text)
        if score > best_score:
            best_score = score
            best_page = block.page

    return best_page if best_score >= 1 else None


def _resolve_linked_ids(
    linked_evidence_ids: list[str],
    source_store=None,
) -> list[str]:
    """Expand DB* doc-block IDs to their constituent B* atomic-block IDs.

    Blueprint slides use DB* IDs (e.g. DB006), but evidence chunks use B* IDs
    (e.g. B029-B036). This helper resolves the mapping so Pass 1 exact match
    can find the right chunks.

    IDs that are already B*/T*/A* are passed through unchanged.
    """
    resolved: list[str] = []
    for eid in linked_evidence_ids:
        if eid.startswith("DB") and source_store is not None:
            doc_block = source_store.get_doc_block(eid)
            if doc_block and doc_block.included_atomic_block_ids:
                resolved.extend(doc_block.included_atomic_block_ids)
                continue
        resolved.append(eid)
    return resolved


def _sanitize_meta_text(html: str) -> str:
    """Remove LLM meta-instructions that leaked into slide HTML content."""
    result = html
    for pattern in _META_TEXT_PATTERNS:
        result = pattern.sub('', result)
    return result


def _theme_with_palette_override(base: ThemeColors, palette: dict | None) -> ThemeColors:
    """Return a ThemeColors view matching an explicit six-role palette."""

    if not palette:
        return base

    def _rgb(value: str, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
        match = re.fullmatch(r"#?([0-9a-fA-F]{6})", str(value or "").strip())
        if not match:
            return fallback
        raw = match.group(1)
        return tuple(int(raw[index:index + 2], 16) for index in (0, 2, 4))  # type: ignore[return-value]

    canvas = _rgb(palette.get("canvas"), base.canvas_color)
    ink = _rgb(palette.get("ink"), base.ink_color)
    primary = _rgb(palette.get("primary"), base.primary_color)
    secondary = _rgb(palette.get("secondary"), base.secondary_color)
    accent = _rgb(palette.get("accent"), base.accent)
    support = _rgb(palette.get("support"), base.support_color)
    return ThemeColors(
        primary_dark=ink,
        primary_mid=primary,
        primary_light=secondary,
        accent=accent,
        accent_alt=support,
        positive=base.positive,
        body_text=ink,
        caption_text=base.caption_text,
        light_bg=canvas,
        warm_bg=support,
        font_family=base.font_family,
        theme_id="demo_palette_override",
        theme_name="Explicit Demo Palette",
        canvas=canvas,
        ink=ink,
        secondary=secondary,
        support=support,
        style_family="demo_curated",
    )


class HtmlCodeGenCompiler:
    """Compiles slides by generating HTML/CSS via LLM, rendering to PNG, and inserting into PPTX."""

    def __init__(
        self,
        llm: LLMClient,
        model: str = "gpt-5.5",
        codegen_prompt: str | None = None,
        show_source_citations: bool = True,
        theme_id: str | None = None,
        demo_palette: dict | None = None,
        repair_model: str | None = None,
        style_pattern: str | None = None,
    ):
        self.llm = llm
        self.model = model
        self.repair_model = repair_model or model
        self.show_source_citations = show_source_citations
        self._theme_id_override = theme_id
        self._demo_palette = demo_palette
        self._style_pattern_arg = style_pattern
        self._selected_pattern: dict | None = None
        # When style-pattern is active, use the trimmed prompt (no conflicting design rules)
        if style_pattern and not codegen_prompt:
            self._codegen_prompt_name = "slide_html_codegen_pattern"
        else:
            self._codegen_prompt_name = codegen_prompt or "slide_html_codegen"
        if codegen_prompt:
            prompt_path = Path(__file__).parent.parent.parent / "prompts" / "codegen" / f"{codegen_prompt}.system.md"
            if prompt_path.exists():
                self.system_prompt = read_text(prompt_path)
                logger.info("Using custom codegen prompt: %s", codegen_prompt)
            else:
                logger.warning("Custom codegen prompt not found: %s, using default", prompt_path)
                self.system_prompt = read_text(HTML_CODEGEN_PROMPT_PATH)
        elif style_pattern:
            # Pattern mode: use minimal prompt without conflicting design rules
            pattern_prompt_path = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "slide_html_codegen_pattern.system.md"
            if pattern_prompt_path.exists():
                self.system_prompt = read_text(pattern_prompt_path)
                logger.info("Using pattern-mode codegen prompt (minimal, %d chars)", len(self.system_prompt))
            else:
                self.system_prompt = read_text(HTML_CODEGEN_PROMPT_PATH)
        else:
            self.system_prompt = read_text(HTML_CODEGEN_PROMPT_PATH)
        self.repair_prompt = read_text(HTML_REPAIR_PROMPT_PATH)

        # For default profile, strip citation-suppression sections from prompts
        # (PresentBench profile keeps them to match existing experiment data)
        if not self.show_source_citations:
            self.system_prompt = self._strip_citation_sections(self.system_prompt)
            self.repair_prompt = self._strip_citation_sections(self.repair_prompt)

        self.slide_codes: dict[int, str] = {}  # slide_id -> generated HTML
        self.used_layouts: list[str] = []  # track layouts for diversity
        self.used_image_ids: set[str] = set()  # track used images to avoid repetition
        self._theme: ThemeColors | None = None  # cached theme for consistent colors
        self._slide_grammars: dict[int, LayoutGrammar] = {}
        self._slide_variants: dict[int, CompositionVariant] = {}
        self._slide_skills: dict[int, list[VisualSkill]] = {}
        self._chart_gen = ChartGenerator()
        self._browser = None
        self._playwright = None
        self._task_brief = ""  # global instructions injected into every slide prompt

    @staticmethod
    def _strip_citation_sections(prompt: str) -> str:
        """Remove citation-suppression paragraphs from prompts for default profile."""
        import re as _re
        # Remove the "Source Citations — DO NOT include" section (header + paragraph)
        prompt = _re.sub(
            r'###\s*1\.\s*Source Citations[^\n]*\n(?:(?!###).*\n)*',
            '', prompt
        )
        # Remove individual citation-related bullet lines
        lines = prompt.split('\n')
        filtered = []
        for line in lines:
            # Skip lines that are purely about suppressing Source/citation
            if any(phrase in line for phrase in [
                'Source:', 'source footnote', 'citation footnote',
                'NO source footnotes', 'citation text of any kind',
                'Source citations (11px)',
                'do NOT add any visible Source',
                'Remove any "Source:"',
            ]):
                continue
            filtered.append(line)
        return '\n'.join(filtered)

    def _apply_frame_contract(self, html: str, slide_id: int) -> str:
        """Normalize full-width HTML header/footer colors before rendering."""

        if not html:
            return html
        # Skip frame contract when style pattern is active — pattern controls all colors
        if self._selected_pattern:
            return html
        theme = self._theme
        if theme is None and not self._demo_palette:
            from app.themes import match_theme_from_html
            theme = match_theme_from_html(html)
            if theme is not None:
                self._theme = theme
                logger.info("Recovered HTML theme from slide %d: %s", slide_id, theme.theme_name)
        if theme is None and not self._demo_palette:
            logger.debug(
                "Slide %d: skipping HTML frame contract because no theme could be recovered",
                slide_id,
            )
            return html

        design_theme = _theme_with_palette_override(theme or DEFAULT_THEME, self._demo_palette)
        normalized, report = enforce_html_slide_frame_contract(
            html,
            design_theme,
            slide_id=slide_id,
        )
        if report.changed and report.treatment != "none":
            logger.info(
                "Slide %d: applied HTML frame contract (%s header)",
                slide_id,
                report.treatment,
            )
        return normalized

    # ------------------------------------------------------------------
    # Playwright lifecycle
    # ------------------------------------------------------------------

    def _ensure_browser(self):
        """Lazily launch Playwright browser (Chromium headless)."""
        if self._browser is not None:
            return

        from playwright.sync_api import sync_playwright
        self._playwright = sync_playwright().start()
        self._browser = self._playwright.chromium.launch(headless=True)
        logger.info("Playwright browser launched for HTML rendering")

    def _close_browser(self):
        """Close Playwright browser and resources."""
        if self._browser:
            try:
                self._browser.close()
            except Exception:
                pass
            self._browser = None
        if self._playwright:
            try:
                self._playwright.stop()
            except Exception:
                pass
            self._playwright = None

    def __del__(self):
        try:
            self._close_browser()
        except AttributeError:
            pass

    # ------------------------------------------------------------------
    # KaTeX formula rendering
    # ------------------------------------------------------------------

    _KATEX_CSS_CDN = "https://cdn.jsdelivr.net/npm/katex@0.17.0/dist/katex.min.css"

    def _render_katex_in_html(self, html: str) -> str:
        """Replace $$...$$ and $...$ LaTeX in HTML with KaTeX-rendered spans.

        Uses node + katex CLI to render each formula to HTML markup,
        then injects the KaTeX CSS link into <head>.
        """
        import subprocess

        # --- Pre-processing: fix common LLM artifacts ---

        # 1. Remove stray trailing $ on HTML lines (LLM sometimes uses $ as
        #    an end-of-line marker, which confuses the inline $ regex).
        #    Only strip when the $ follows > or whitespace at line end.
        html = re.sub(r'\$\s*$', '', html, flags=re.MULTILINE)

        # 2. Fix mismatched delimiters: $...$$  ->  $$...$$
        #    (single $ open, double $$ close)
        html = re.sub(
            r'(?<!\$)\$(?!\$)(.+?)\$\$',
            lambda m: f'$${m.group(1)}$$' if self._looks_like_latex(m.group(1).strip()) else m.group(0),
            html,
        )

        # Find all formulas: $$...$$ (display) and $...$ (inline)
        # Avoid matching CSS values like "top: $10px" — require LaTeX-like content
        display_pattern = re.compile(r'\$\$(.+?)\$\$', re.DOTALL)
        # For inline $...$: exclude $ followed by digit (dollar amounts like $8.9)
        inline_pattern = re.compile(r'(?<!\$)\$(?!\$)(?!\d)(.+?)(?<!\$)\$(?!\$)')

        formulas: list[tuple[str, str, bool]] = []  # (full_match, latex, is_display)

        for m in display_pattern.finditer(html):
            latex = m.group(1).strip()
            if self._looks_like_latex(latex):
                formulas.append((m.group(0), latex, True))

        for m in inline_pattern.finditer(html):
            latex = m.group(1).strip()
            if self._looks_like_latex(latex):
                formulas.append((m.group(0), latex, False))

        if not formulas:
            return html

        # Render each formula via katex CLI
        has_rendered = False
        for full_match, latex, is_display in formulas:
            try:
                # Decode HTML entities before passing to KaTeX
                # LLM sometimes generates &gt; &lt; &amp; inside math
                import html as html_module
                latex_decoded = html_module.unescape(latex)

                cmd = ["npx", "katex", "--no-throw-on-error"]
                if is_display:
                    cmd.append("--display-mode")
                result = subprocess.run(
                    cmd, input=latex_decoded, capture_output=True, text=True, timeout=10,
                )
                if result.returncode == 0 and result.stdout.strip():
                    rendered = result.stdout.strip()
                    html = html.replace(full_match, rendered, 1)
                    has_rendered = True
                else:
                    logger.debug("KaTeX failed for: %s — %s", latex[:50], result.stderr[:100])
            except Exception as e:
                logger.debug("KaTeX rendering error: %s", e)

        # Inject KaTeX CSS into <head> if any formula was rendered
        if has_rendered and self._KATEX_CSS_CDN not in html:
            katex_link = f'<link rel="stylesheet" href="{self._KATEX_CSS_CDN}">'
            if "<head>" in html:
                html = html.replace("<head>", f"<head>\n{katex_link}", 1)
            elif "<style>" in html:
                html = html.replace("<style>", f"{katex_link}\n<style>", 1)

        return html

    @staticmethod
    def _cleanup_residual_latex(html: str) -> str:
        """Clean up LaTeX artifacts that KaTeX didn't render.

        Handles:
        - $...$ inline math delimiters → strip delimiters, convert content to Unicode
        - \\textbf{}, \\mathbb{}, \\mathrm{} wrappers → keep inner text
        - Stray backslash commands → remove or convert to Unicode
        """
        from ...utils.formula_renderer import FormulaRenderer

        # 1. Strip $...$ delimiters and convert inner LaTeX to Unicode
        def _convert_inline_math(m: re.Match) -> str:
            inner = m.group(1).strip()
            # Skip if it looks like natural language (3+ English words)
            import re as _re
            if len(_re.findall(r'\b[a-zA-Z]{4,}\b', inner)) >= 3:
                return m.group(0)  # return unchanged
            try:
                return FormulaRenderer.latex_to_unicode(inner)
            except Exception:
                return inner

        html = re.sub(
            r'(?<!\$)\$(?!\$)(?!\d)([^$\n]{1,80})(?<!\$)\$(?!\$)',
            _convert_inline_math,
            html,
        )

        # 2. Remove \textbf{...}, \textit{...}, \mathbb{...}, \mathrm{...}
        #    wrappers — keep inner text. Repeat to handle nesting.
        for _ in range(3):
            prev = html
            html = re.sub(
                r'\\(?:textbf|textit|textrm|mathrm|mathbf|mathbb|mathcal|emph)\{([^{}]*)\}',
                r'\1',
                html,
            )
            if html == prev:
                break

        # 3. Replace common standalone LaTeX commands with Unicode
        _QUICK_REPLACEMENTS = {
            r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\delta': 'δ',
            r'\epsilon': 'ε', r'\theta': 'θ', r'\lambda': 'λ', r'\mu': 'μ',
            r'\sigma': 'σ', r'\omega': 'ω', r'\pi': 'π', r'\tau': 'τ',
            r'\times': '×', r'\cdot': '·', r'\approx': '≈', r'\neq': '≠',
            r'\leq': '≤', r'\geq': '≥', r'\infty': '∞', r'\partial': '∂',
            r'\rightarrow': '→', r'\leftarrow': '←', r'\Rightarrow': '⇒',
            r'\sum': 'Σ', r'\prod': 'Π', r'\in': '∈', r'\subset': '⊂',
            r'\ldots': '…', r'\dots': '…', r'\pm': '±',
            r'\leftrightarrow': '↔', r'\Leftrightarrow': '⇔',
            r'\forall': '∀', r'\exists': '∃', r'\nabla': '∇',
            r'\Pi': 'Π', r'\Sigma': 'Σ', r'\Omega': 'Ω', r'\Delta': 'Δ',
            r'\Gamma': 'Γ', r'\Lambda': 'Λ', r'\Theta': 'Θ',
            r'\mid': '|', r'\parallel': '∥', r'\perp': '⊥',
            r'\cap': '∩', r'\cup': '∪', r'\emptyset': '∅',
            r'\log': 'log', r'\exp': 'exp', r'\max': 'max', r'\min': 'min',
        }
        for cmd, uni in _QUICK_REPLACEMENTS.items():
            html = html.replace(cmd, uni)

        # 4. Remove any remaining stray \command that is NOT inside a
        #    <style>, <script>, or url() context. Only target text nodes.
        #    Be conservative: only remove commands we recognize as LaTeX.
        html = re.sub(
            r'\\(?:text|math|frac|sqrt|bar|hat|tilde|vec|dot|overline)\{([^{}]*)\}',
            r'\1',
            html,
        )

        # 5. Fix HTML entities that should be symbols in text nodes
        #    (e.g., &lt; → < and &gt; → > when inside visible text, not HTML tags)
        #    Only replace when surrounded by math-like context (digits, variables)
        html = re.sub(r'(\w)&lt;(\w)', r'\1<\2', html)
        html = re.sub(r'(\w)&gt;(\w)', r'\1>\2', html)

        return html

    @staticmethod
    def _looks_like_latex(text: str) -> bool:
        """Check if text looks like LaTeX math (not a CSS value or price).

        Accepts both complex LaTeX (\\frac, ^, _) AND simple math expressions
        like single variables (W, r, k), expressions (O(1), N=6), and numeric
        formulas that are commonly written inside $...$ in academic papers.
        """
        if len(text) < 1 or len(text) > 500:
            return False
        # Exclude things that look like CSS/JS: px, %, #, url(
        css_indicators = any(s in text for s in ('px', 'url(', '://', '#'))
        if css_indicators:
            return False
        # Exclude natural language: if text contains 3+ consecutive English
        # words (4+ letter words separated by spaces), it's prose, not math.
        import re as _re
        word_seq = _re.findall(r'\b[a-zA-Z]{4,}\b', text)
        if len(word_seq) >= 3:
            return False
        # Complex LaTeX indicators — always accept
        if ('\\' in text or '^' in text or '_' in text
                or '{' in text or 'frac' in text
                or any(c in text for c in 'αβγδεζηθλμνξπρστφψω∑∏∫')):
            return True
        # Simple math expressions: single letters (W, r, k, N), or
        # letter+number combos (N=6, h=8), or function-like (O(1), O(n^2))
        # These are common in academic papers inside $...$
        if _re.fullmatch(r'[A-Za-z0-9()\s=+\-*/.,<>≤≥]+', text):
            return True
        return False

    # ------------------------------------------------------------------
    # Solid-color card background softener
    # ------------------------------------------------------------------

    @staticmethod
    def _soften_solid_card_backgrounds(html: str) -> str:
        """Convert garish solid-color card backgrounds to subtle rgba tints.

        Strategy: ANY inline style with background:#XXXXXX gets softened to
        rgba(R,G,B,0.10-0.12) UNLESS it is clearly a structural element
        (header band, badge, near-white) OR an intentional dark container
        whose children have explicit white text.

        Protected (NOT softened):
        - Gradient backgrounds (linear-gradient, radial-gradient)
        - Full-width structural bands (width:1280 or left:0+right:0)
        - Small badges/pills (border-radius: 50% in same style)
        - White/near-white backgrounds (all channels > 220)
        - Very small elements: width ≤ 56px (icon badges defined inline)
        - Dark containers (luminance < 120) with children that have explicit
          white text — signals intentional dark-bg + light-text design
        """
        import re as _re
        _WHITE_TEXT_IN_CHILD = _re.compile(
            r'style="[^"]*color:\s*#[Ff]{3}(?:[Ff]{3})?\b'
        )
        _BG_PATTERN = _re.compile(
            r'style="[^"]*?background:\s*#([0-9a-fA-F]{6})[^"]*?"'
        )

        result_parts = []
        last_end = 0

        for match in _BG_PATTERN.finditer(html):
            result_parts.append(html[last_end:match.start()])
            full_style = match.group(0)
            hex_color = match.group(1)

            # Skip if it's a gradient context
            if 'gradient' in full_style:
                result_parts.append(full_style)
                last_end = match.end()
                continue
            # Skip if it looks like a circular badge (border-radius:50%)
            if 'border-radius:50%' in full_style or 'border-radius: 50%' in full_style:
                result_parts.append(full_style)
                last_end = match.end()
                continue

            r_val = int(hex_color[0:2], 16)
            g_val = int(hex_color[2:4], 16)
            b_val = int(hex_color[4:6], 16)

            # Skip white/near-white (luminance > 220 is already light enough)
            if r_val > 220 and g_val > 220 and b_val > 220:
                result_parts.append(full_style)
                last_end = match.end()
                continue

            # Skip small inline icon badges (width ≤ 56px in same style)
            width_m = _re.search(r'width:\s*(\d+)px', full_style)
            if width_m and int(width_m.group(1)) <= 56:
                result_parts.append(full_style)
                last_end = match.end()
                continue

            # Skip full-width structural bands (header/footer)
            is_full_width = ('width:1280' in full_style or 'width: 1280' in full_style
                            or ('left:0' in full_style and ('right:0' in full_style or 'width:100%' in full_style)))
            luminance = 0.299 * r_val + 0.587 * g_val + 0.114 * b_val
            if is_full_width and luminance < 120:
                result_parts.append(full_style)
                last_end = match.end()
                continue

            # Skip pill/badge overrides: style only sets bg (no layout props
            # like padding, width, height, flex). These are small inline
            # color overrides whose text color comes from a CSS class —
            # softening would make the bg nearly invisible and orphan
            # class-level white text.
            style_inner = full_style[7:-1]  # strip style=" ... "
            has_layout = any(kw in style_inner for kw in (
                'padding', 'width', 'height', 'flex', 'display',
                'position', 'top:', 'left:', 'right:', 'bottom:',
            ))
            if not has_layout:
                result_parts.append(full_style)
                last_end = match.end()
                continue

            # Skip intentional dark containers (lum < 120) whose children
            # have explicit white text — softening would orphan white text
            # on a now-light bg, creating unreadable text.
            if luminance < 120:
                lookahead_end = min(match.end() + 600, len(html))
                lookahead = html[match.end():lookahead_end]
                if _WHITE_TEXT_IN_CHILD.search(lookahead):
                    result_parts.append(full_style)
                    last_end = match.end()
                    continue

            # Everything else: soften to rgba tint
            alpha = 0.15 if luminance < 100 else 0.18
            rgba = f"rgba({r_val},{g_val},{b_val},{alpha})"
            new_style = full_style.replace(f"#{hex_color}", rgba)
            # Fix white text that was paired with the solid background
            new_style = _re.sub(
                r'color:\s*#[Ff]{3,6}\b',
                'color:#2d2d2d',
                new_style,
            )
            result_parts.append(new_style)
            last_end = match.end()

        result_parts.append(html[last_end:])
        return ''.join(result_parts)

    @staticmethod
    def _fix_white_text_on_light_bg(html: str) -> str:
        """Fix white text (#fff/#ffffff) on light rgba backgrounds.

        When a style contains background: rgba(R,G,B, alpha<=0.3) AND
        color: #fff/#ffffff, the text is invisible. Convert the text color
        to #2d2d2d (safe dark text).

        Only targets inline style attributes (style="...").
        Does NOT touch header bands, badges, or solid dark backgrounds.
        """
        def _fix_white_on_rgba(match):
            full_style = match.group(0)

            # Only fix if it has an rgba background with alpha <= 0.3
            rgba_match = re.search(
                r'background:\s*rgba\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([\d.]+)\s*\)',
                full_style,
            )
            if not rgba_match:
                return full_style
            alpha = float(rgba_match.group(1))
            if alpha > 0.3:
                return full_style  # dark enough to have white text

            # Skip header bands and full-width elements
            if 'width:1280' in full_style or 'width: 1280' in full_style:
                return full_style
            # Skip small badge-like elements
            if 'border-radius:50%' in full_style or 'border-radius: 50%' in full_style:
                return full_style

            # Replace white text with safe dark
            new_style = re.sub(
                r'color:\s*#[Ff]{3,6}\b',
                'color:#2d2d2d',
                full_style,
            )
            return new_style

        # Match style attributes that contain both rgba background and white color
        html = re.sub(
            r'style="[^"]*?background:\s*rgba\([^)]+\)[^"]*?color:\s*#[Ff]{3,6}[^"]*?"',
            _fix_white_on_rgba,
            html,
        )
        # Also match when color comes before background
        html = re.sub(
            r'style="[^"]*?color:\s*#[Ff]{3,6}[^"]*?background:\s*rgba\([^)]+\)[^"]*?"',
            _fix_white_on_rgba,
            html,
        )
        return html

    @staticmethod
    def _strip_card_shadows(html: str) -> str:
        """Remove box-shadow from inline styles to prevent card proliferation.

        Keeps box-shadow ONLY on:
        - Elements with class "chart-area" (chart containers need subtle lift)
        - Elements inside <svg> (SVG filter shadows)
        - The takeaway/conclusion bar (typically has PRIMARY_DARK background)

        Everything else loses its box-shadow, forcing flat/clean PPT style.
        """
        import re as _re

        def _remove_shadow(match):
            full_style = match.group(0)
            # Don't strip from chart areas (detected by #fafbfd or chart-area nearby)
            # We can't see class from style attr alone, so preserve if bg is chart-like
            if '#fafbfd' in full_style or '#f8f9fb' in full_style:
                return full_style
            # Don't strip from dark background elements (takeaway bars, headers)
            if _re.search(r'background:\s*(?:PRIMARY_DARK|#[0-3])', full_style):
                return full_style
            # Strip box-shadow
            cleaned = _re.sub(
                r'box-shadow:\s*[^;]+;?\s*',
                '',
                full_style,
            )
            return cleaned

        # Match style attributes that contain box-shadow
        html = _re.sub(
            r'style="[^"]*?box-shadow:[^"]*?"',
            _remove_shadow,
            html,
        )
        return html

    # ------------------------------------------------------------------
    # HTML post-processing
    # ------------------------------------------------------------------

    @staticmethod
    def _fix_title_overflow(html: str, slide_id: int, paper_title: str = "") -> str:
        """Auto-shrink title font-size if text is too long for its container.
        Also fix truncated paper titles on slide 1.

        For the header-band `.title` element (all content slides): the title
        sits in a 1280-56-56 = 1168px wide area. At font-size 38px, roughly
        ~30 chars fit. We estimate: if chars * font_size * 0.55 > available_width,
        shrink font_size until it fits.

        For the title slide (slide_id=1): the <h1> has max-width 1100px.
        """
        import re

        # --- Fix truncated paper title on slide 1 ---
        if slide_id == 1 and paper_title and len(paper_title) > 10:
            # Try <h1> first, then .title class div
            h1_content_match = re.search(r'(<h1[^>]*>)(.*?)(</h1>)', html, re.DOTALL)
            if not h1_content_match:
                # Try class="title" div
                h1_content_match = re.search(r'(class="title"[^>]*>)(.*?)(<)', html, re.DOTALL)
            if h1_content_match:
                h1_text = re.sub(r'<[^>]+>', '', h1_content_match.group(2)).strip()
                # Check if the h1 text is a truncated prefix of the paper title
                if (len(h1_text) >= 5 and len(h1_text) < len(paper_title) - 3
                        and paper_title.lower().startswith(h1_text.lower())):
                    html = html.replace(
                        h1_content_match.group(2),
                        paper_title,
                        1,
                    )
                    logger.info(
                        "Slide 1: restored truncated title '%s...' → '%s'",
                        h1_text[:30], paper_title[:60],
                    )

        # Content slide title: `.title { ... font-size: XXpx ... }`
        # Match the .title element and extract text
        title_match = re.search(
            r'class="title"[^>]*>([^<]+)<', html
        )
        if title_match:
            title_text = title_match.group(1).strip()
            avail_width = 1168  # 1280 - 56*2
            # Extract current font-size from .title CSS
            fs_match = re.search(r'\.title\s*\{[^}]*font-size:\s*(\d+)px', html)
            if fs_match:
                current_fs = int(fs_match.group(1))
            else:
                current_fs = 38
            # Estimate text width: chars * font_size * 0.55 (average char width ratio)
            estimated_width = len(title_text) * current_fs * 0.55
            if estimated_width > avail_width:
                # Calculate needed font-size
                new_fs = max(30 if slide_id == 1 else 20, int(avail_width / (len(title_text) * 0.55)))
                if new_fs < current_fs:
                    if fs_match:
                        old = fs_match.group(0)
                        new = old.replace(f"{current_fs}px", f"{new_fs}px")
                        html = html.replace(old, new, 1)
                    else:
                        # Inject font-size into the element style
                        html = html.replace(
                            'class="title"',
                            f'class="title" style="font-size: {new_fs}px;"',
                            1,
                        )
                    logger.info(
                        "Slide %d: auto-shrunk title from %dpx to %dpx (%d chars)",
                        slide_id, current_fs, new_fs, len(title_text),
                    )

        # Title slide <h1>: check if it's too long
        if slide_id == 1:
            h1_match = re.search(r'<h1[^>]*style="([^"]*)"[^>]*>([^<]+(?:<[^/][^>]*>[^<]*</[^>]*>)*[^<]*)</h1>', html, re.DOTALL)
            if h1_match:
                style = h1_match.group(1)
                # Get text content (strip inner tags)
                text = re.sub(r'<[^>]+>', '', h1_match.group(2)).strip()
                fs_match = re.search(r'font-size:\s*(\d+)px', style)
                current_fs = int(fs_match.group(1)) if fs_match else 38
                avail_width = 1100
                estimated_width = len(text) * current_fs * 0.55
                if estimated_width > avail_width:
                    new_fs = max(22, int(avail_width / (len(text) * 0.55)))
                    if new_fs < current_fs:
                        new_style = re.sub(
                            r'font-size:\s*\d+px',
                            f'font-size: {new_fs}px',
                            style,
                        )
                        html = html.replace(style, new_style, 1)
                        logger.info(
                            "Slide 1: auto-shrunk h1 from %dpx to %dpx (%d chars)",
                            current_fs, new_fs, len(text),
                        )

        # --- Enforce minimum paper title size on slide 1 ---
        if slide_id == 1:
            # Find .title or .paper-title CSS rule and check font-size
            for cls_name in ('paper-title', 'title'):
                pat = re.compile(
                    r'(\.' + cls_name + r'\s*\{[^}]*font-size:\s*)(\d+)(px)',
                    re.DOTALL,
                )
                m = pat.search(html)
                if m:
                    fs = int(m.group(2))
                    if fs < 30:
                        html = html[:m.start(2)] + '34' + html[m.end(2):]
                        logger.info(
                            "Slide 1: enforced min paper-title from %dpx to 34px",
                            fs,
                        )
                    break  # only fix the first matching class

        return html

    # ------------------------------------------------------------------
    # HTML rendering
    # ------------------------------------------------------------------

    def _render_html_to_png(self, html_content: str, output_path: Path) -> bool:
        """Render HTML to PNG using Playwright.

        Args:
            html_content: Complete HTML page content
            output_path: Where to save the PNG

        Returns:
            True if rendering succeeded
        """
        self._ensure_browser()

        # Convert bare absolute file paths in <img src="..."> to file:// URLs
        # so Playwright can load local images from the filesystem.

        # Render KaTeX formulas in the HTML before rendering
        html_content = self._render_katex_in_html(html_content)

        # Fallback: clean up any residual LaTeX that KaTeX didn't process
        html_content = self._cleanup_residual_latex(html_content)

        # Soften garish solid-color card backgrounds to rgba tints.
        # Matches background:#XXXXXX on elements that look like cards (have padding/border-radius)
        # but NOT on headers, badges, or small pills.
        html_content = self._soften_solid_card_backgrounds(html_content)

        # Fix white text on light rgba backgrounds (invisible text safety net)
        html_content = self._fix_white_text_on_light_bg(html_content)

        # v20f: Strip box-shadow from non-chart containers to prevent card proliferation
        html_content = self._strip_card_shadows(html_content)

        # Remove leaked internal prompt references from visible text
        for leaked_term in ["task brief", "slide brief", "evidence text", "source evidence"]:
            html_content = re.sub(
                re.escape(leaked_term), "", html_content, flags=re.IGNORECASE
            )

        # NOTE: Source citation stripping for benchmark compatibility was removed.
        # Visible citations on data slides are now the default (good practice for
        # academic and financial presentations).

        # Fix Title Case capitalization in <h1>/<h2> slide titles
        # Words that should be lowercase in Title Case (unless first/last word):
        _TC_LOWER = {'a', 'an', 'the', 'and', 'but', 'or', 'nor', 'for', 'yet', 'so',
                      'in', 'on', 'at', 'to', 'of', 'by', 'with', 'as', 'vs', 'via'}

        def _fix_title_case(match):
            """Fix capitalization in a title tag, preserving HTML tags inside."""
            tag = match.group(1)       # e.g. "h1" or "h2"
            attrs = match.group(2)     # tag attributes
            inner = match.group(3)     # inner HTML content
            close = match.group(4)     # closing tag

            # Don't modify if it contains complex nested HTML (images, spans with classes)
            if '<img' in inner or '<span class' in inner:
                return match.group(0)

            # Split inner into HTML tags and text segments
            parts = re.split(r'(<[^>]+>)', inner)
            fixed_parts = []
            for i, part in enumerate(parts):
                if part.startswith('<'):
                    fixed_parts.append(part)
                    continue
                # Fix each word in text segment
                words = part.split(' ')
                fixed_words = []
                for j, word in enumerate(words):
                    if not word:
                        fixed_words.append(word)
                        continue
                    # Check if this word is at the absolute start or end of the title text
                    text_before = ' '.join(' '.join(fixed_parts).split())
                    stripped = re.sub(r'<[^>]+>', '', text_before).strip()
                    is_first = (stripped == '' and j == 0)
                    is_last_segment = (i == len(parts) - 1 or all(p.startswith('<') or p.strip() == '' for p in parts[i+1:]))
                    is_last = is_last_segment and j == len(words) - 1

                    bare = re.sub(r'[^\w]', '', word.lower())
                    if bare in _TC_LOWER and not is_first and not is_last:
                        # Lowercase this word (preserve non-alpha chars around it)
                        # e.g. "And" -> "and", "In" -> "in", "A" -> "a"
                        if word[0].isupper() and (len(word) == 1 or word[1:].islower()) and len(word) <= 5:
                            fixed_words.append(word[0].lower() + word[1:])
                        else:
                            fixed_words.append(word)
                    else:
                        fixed_words.append(word)
                fixed_parts.append(' '.join(fixed_words))
            return f'<{tag}{attrs}>{"".join(fixed_parts)}</{close}>'

        html_content = re.sub(
            r'<(h[12])([^>]*)>(.*?)</(h[12])>',
            _fix_title_case,
            html_content,
            flags=re.DOTALL
        )

        html_content = re.sub(
            r'(<img\s[^>]*src=["\'])(/[^"\']+)(["\'])',
            r'\1file://\2\3',
            html_content,
        )

        # Convert relative paths to absolute file:// URLs.
        # LLM may generate src="cases/..." which needs to resolve from the
        # project root, not from the temp file's directory.
        import os
        from pathlib import Path as _Path
        cwd = os.getcwd()
        _search_roots = [cwd]
        _probe = _Path(cwd)
        for _ in range(6):
            if (_probe / "cases").is_dir() or (_probe / "app").is_dir():
                if str(_probe) != cwd:
                    _search_roots.append(str(_probe))
                break
            _probe = _probe.parent

        def _resolve_relative_src(m):
            prefix, path, suffix = m.group(1), m.group(2), m.group(3)
            if path.startswith(('file://', 'http://', 'https://', 'data:', '/')):
                return m.group(0)
            path_parts = _Path(path).parts
            if 'generated_assets' in path_parts:
                turn_dir = _Path(output_path).parent.parent
                asset_name = _Path(path).name
                asset_candidates = [
                    turn_dir / 'generated_assets' / asset_name,
                    turn_dir / path,
                ]
                for abs_path in asset_candidates:
                    if abs_path.exists():
                        return f'{prefix}file://{abs_path.resolve()}{suffix}'
            for root in _search_roots:
                abs_path = os.path.join(root, path)
                if os.path.exists(abs_path):
                    return f'{prefix}file://{abs_path}{suffix}'
            return m.group(0)

        html_content = re.sub(
            r'(<img\s[^>]*src=["\'])([^"\']+)(["\'])',
            _resolve_relative_src,
            html_content,
        )

        page = None
        try:
            page = self._browser.new_page(
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
                device_scale_factor=DEVICE_SCALE_FACTOR,
            )
            # Save HTML to a temp file so Playwright can load local images via file:// protocol
            import tempfile
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".html", delete=False, encoding="utf-8"
            ) as tmp:
                tmp.write(html_content)
                tmp_path = tmp.name
            try:
                page.goto(f"file://{tmp_path}", wait_until="networkidle")
                # Small wait for any CSS transitions/rendering
                page.wait_for_timeout(300)
                output_path.parent.mkdir(parents=True, exist_ok=True)
                page.screenshot(path=str(output_path), full_page=False)
            finally:
                import os as _os
                _os.unlink(tmp_path)
            page.close()
            return True
        except Exception as e:
            logger.error("Playwright render failed: %s", e)
            if page:
                try:
                    page.close()
                except Exception:
                    pass
            return False

    # ------------------------------------------------------------------
    # Public API (mirrors CodeGenCompiler interface)
    # ------------------------------------------------------------------

    def compile_deck(
        self,
        blueprint: DeckBlueprint,
        evidence: EvidenceState,
        case_dir: str | Path,
        output_path: str | Path,
        code_dir: str | Path | None = None,
        source_store=None,
        task_brief: str = "",
    ) -> dict:
        """Generate HTML for each slide, render to PNG, and compile into a PPTX.

        Args:
            blueprint: The deck blueprint with slide briefs
            evidence: Evidence state with chunks, figures, tables
            case_dir: Path to case directory (for image resolution)
            output_path: Where to save the PPTX file
            code_dir: Where to save generated HTML files (optional)
            task_brief: Global task brief / instructions to inject into
                every slide prompt so per-slide LLM sees the full constraints.

        Returns:
            Manifest dict with slide info and any warnings
        """
        # Store task_brief so _build_slide_prompt and repair_slide can use it
        self._task_brief = task_brief
        start = time.time()
        case_dir = Path(case_dir)
        output_path = Path(output_path)
        if code_dir:
            code_dir = Path(code_dir)
            code_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # Reset per-deck state for visual rhythm tracking
        self.used_layouts.clear()
        self.used_image_ids.clear()
        self._slide_grammars.clear()
        self._slide_variants.clear()
        self._slide_skills.clear()

        image_dir = self._find_image_dir(case_dir)
        available_images = self._list_images(image_dir) if image_dir else []

        # Select theme based on paper title for consistent colors across slides
        from app.themes import select_theme_for_paper
        paper_title = ""
        if source_store and hasattr(source_store, 'doc_block_plan'):
            for db in source_store.doc_block_plan.blocks:
                if db.role in ("title", "abstract", "introduction"):
                    paper_title = db.title or ""
                    break
        if not paper_title and blueprint.slides:
            prop = blueprint.slides[0].primary_proposition or ""
            # Extract just the paper title from primary_proposition
            # (it's usually "Paper Title introduces/presents/proposes...")
            import re as _re
            title_cut = _re.split(
                r'\s+(?:introduces?|presents?|proposes?|describes?|provides?|is\s+a|addresses)\s+',
                prop, maxsplit=1
            )
            paper_title = title_cut[0].strip() if title_cut else prop
        self._theme = (
            select_theme_for_paper(paper_title, self._theme_id_override)
            if paper_title or self._theme_id_override else DEFAULT_THEME
        )

        # Try to match theme to figure colors for visual harmony
        if not self._theme_id_override and image_dir:
            try:
                from app.figure_theme_matcher import select_theme_by_figures
                fig_paths = [image_dir / img for img in available_images[:8]]
                matched_theme = select_theme_by_figures(fig_paths)
                if matched_theme and matched_theme in THEME_REGISTRY:
                    self._theme = THEME_REGISTRY[matched_theme]
                    logger.info("Theme overridden by figure-color match: %s", matched_theme)
            except Exception as e:
                logger.debug("Figure-theme matching skipped: %s", e)
        if self._theme_id_override:
            if self._theme_id_override in THEME_REGISTRY:
                logger.info("Theme overridden: %s", self._theme.theme_name)
            elif self._theme_id_override in THEME_FAMILIES:
                logger.info(
                    "Theme family overridden: %s -> %s",
                    self._theme_id_override,
                    self._theme.theme_name,
                )
            else:
                logger.warning(
                    "Unknown theme_id override %r; using selected theme %s",
                    self._theme_id_override,
                    self._theme.theme_name,
                )

        # Select style pattern for deck-level design consistency
        if self._style_pattern_arg:
            if self._style_pattern_arg == "vocab":
                # Vocab mode: compose from independent element categories
                from app.style_patterns.vocab_composer import format_vocab_style_contract
                self._vocab_contract = format_vocab_style_contract(paper_title)
                self._selected_pattern = {"id": "vocab", "dims": {"mode": "vocab-composed"}}
                logger.info("Style pattern: vocab composition mode (3229 elements × 9 categories)")
            else:
                from app.style_patterns import select_pattern_for_deck, get_pattern_by_name
                if self._style_pattern_arg == "auto":
                    self._selected_pattern = select_pattern_for_deck(paper_title)
                    logger.info("Style pattern auto-selected: %s (dims: %s)",
                                self._selected_pattern.get("id"),
                                self._selected_pattern.get("dims", {}))
                else:
                    self._selected_pattern = get_pattern_by_name(self._style_pattern_arg)
                    if self._selected_pattern:
                        logger.info("Style pattern specified: %s", self._selected_pattern.get("id"))
                    else:
                        logger.warning("Style pattern %r not found; proceeding without", self._style_pattern_arg)
        self._paper_title = paper_title
        logger.info("Selected theme: %s (from title: '%s')", self._theme.theme_name, paper_title[:50])

        manifest = {
            "pptx_path": str(output_path),
            "total_slides": len(blueprint.slides),
            "slides": [],
            "warnings": [],
            "timing_sec": 0,
            "compiler": "html_codegen",
            "design_profile": {
                "selection": self._theme_id_override or "auto",
                "theme_id": self._theme.theme_id,
                "theme_name": self._theme.theme_name,
                "style_family": self._theme.style_family,
                "codegen_prompt": self._codegen_prompt_name,
                "codegen_prompt_version": HTML_CODEGEN_PROMPT_VERSION,
                "visual_skill_library_version": VISUAL_SKILL_LIBRARY_VERSION,
                "composition_variant_version": COMPOSITION_VARIANT_VERSION,
            },
        }

        # Phase 0: Pre-generate charts for slides that have viz_data
        self._chart_paths: dict[int, str] = {}
        chart_dir = output_path.parent / "charts"
        theme = self._theme or DEFAULT_THEME

        def _theme_hex(c):
            return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"

        if getattr(theme, "style_family", "").startswith("demo_curated"):
            theme_hex = [
                _theme_hex(theme.primary_color),
                _theme_hex(theme.accent),
                _theme_hex(theme.secondary_color),
                _theme_hex(theme.support_color),
                _theme_hex(theme.positive),
                _theme_hex(theme.ink_color),
            ]
        else:
            theme_hex = [
                _theme_hex(theme.primary_dark),
                _theme_hex(theme.accent),
                _theme_hex(theme.primary_mid),
                _theme_hex(theme.accent_alt),
                _theme_hex(theme.positive),
                _theme_hex(theme.body_text),
            ]
        for bp_slide in blueprint.slides:
            viz = getattr(bp_slide, 'viz_data', None) or {}
            if viz and viz.get("chart_type"):
                chart_path = chart_dir / f"chart_slide_{bp_slide.slide_id}.png"
                result = self._chart_gen.generate_chart(viz, chart_path, theme_hex)
                if result:
                    self._chart_paths[bp_slide.slide_id] = str(result)
                    logger.info("Pre-generated chart for slide %d: %s", bp_slide.slide_id, result)

        # Phase 1: Generate HTML for slides
        # Strategy depends on prompt type:
        # - imgseed/freeform: SERIAL generation. Slide 1 establishes the visual language,
        #   subsequent slides receive its CSS variables as a consistency anchor.
        # - default prompt: PARALLEL generation (already has rigid style system).
        from concurrent.futures import ThreadPoolExecutor, as_completed

        def _gen_html(bp_slide, style_anchor: str = ""):
            """Generate HTML + render PNG for one slide (thread-safe for LLM)."""
            sid = bp_slide.slide_id

            # Build user prompt
            user_content, has_images = self._build_slide_prompt(
                bp_slide, evidence, image_dir, available_images,
                total_slides=len(blueprint.slides),
                source_store=source_store,
                )

            # Inject style anchor from slide 1 for consistency
            if style_anchor:
                if self._selected_pattern:
                    # Vocab/pattern mode: DON'T use slide 1 anchor — it may diverge
                    # from the vocab contract. The vocab contract itself provides
                    # consistent color instructions on every slide.
                    pass
                else:
                    user_content = (
                        f"## Deck Color & Typography (from slide 1 — reuse these variables)\n"
                        f"```css\n{style_anchor}\n```\n"
                        f"Reuse these CSS variables for color palette, font sizes, and border-radius. "
                        f"BUT you have COMPLETE FREEDOM over page structure: whether to use a header band, "
                        f"sidebar, full-bleed layout, split columns, centered card, etc. "
                        f"Do NOT copy the same header/sidebar skeleton from previous slides — "
                        f"vary the structural bones while keeping colors consistent.\n\n"
                        + user_content
                    )

            # LLM call to generate HTML
            html = None
            last_error = None
            for attempt in range(MAX_CODE_RETRIES):
                try:
                    if attempt == 0:
                        response = self.llm.call_text(
                            system_prompt=self.system_prompt,
                            user_content=user_content,
                            model=self.model,
                            module_name="slide_html_codegen",
                            prompt_version=HTML_CODEGEN_PROMPT_VERSION,
                            max_tokens=8000,
                        )
                    else:
                        retry_prompt = f"{user_content}\n\n## RETRY (attempt {attempt+1})\nPrevious error: {last_error}\nFix the issue and regenerate."
                        response = self.llm.call_text(
                            system_prompt=self.system_prompt,
                            user_content=retry_prompt,
                            model=self.model,
                            module_name="slide_html_codegen",
                            prompt_version=f"{HTML_CODEGEN_PROMPT_VERSION}.retry",
                            max_tokens=8000,
                        )
                    html = self._extract_html(response)
                    if html and "<body" in html.lower():
                        break
                    last_error = "No valid HTML extracted"
                    html = None
                except Exception as e:
                    last_error = str(e)[:200]

            return sid, html, last_error

        def _extract_style_anchor(html: str) -> str:
            """Extract :root CSS variables block from the first slide's HTML."""
            import re
            # Match :root { ... } block
            m = re.search(r':root\s*\{([^}]+)\}', html)
            if m:
                return f":root {{{m.group(1)}}}"
            # Fallback: extract all --var declarations from <style>
            style_match = re.search(r'<style[^>]*>(.*?)</style>', html, re.DOTALL)
            if style_match:
                vars_lines = [l.strip() for l in style_match.group(1).split('\n')
                              if '--' in l and ':' in l]
                if vars_lines:
                    return ":root {\n  " + "\n  ".join(vars_lines) + "\n}"
            return ""

        # Run generation
        _is_freeform = 'imgseed' in (self._codegen_prompt_name or '')
        _is_pattern_mode = bool(self._selected_pattern)
        html_results = {}

        def _describe_layout(html: str) -> str:
            """Produce a short structural description of a generated slide.

            Detects both structural elements AND content distribution pattern,
            so the diversity tracker can prevent repetitive spatial arrangements.
            """
            html_lower = html.lower()
            features = []
            # Header type
            if 'background' in html_lower and ('linear-gradient' in html_lower or '--primary-dark' in html_lower):
                if re.search(r'\.header\b.*?height:\s*(8[0-9]|9[0-9]|1[0-2][0-9])px', html_lower):
                    features.append("tall dark header band")
                else:
                    features.append("header band")
            elif re.search(r'border-bottom|accent.*?bar|\.bar.*?height:\s*[4-9]px', html_lower):
                features.append("thin accent bar")
            # Grid / columns
            if 'grid-template-columns' in html_lower:
                cols = re.findall(r'grid-template-columns\s*:\s*([^;]+)', html_lower)
                if cols:
                    features.append(f"CSS grid ({cols[0].strip()[:40]})")
            elif html_lower.count('flex: 1') >= 2 or html_lower.count('flex:1') >= 2:
                features.append("flexbox columns")
            # Sidebar
            if re.search(r'width:\s*(25|28|30|32|35)%', html_lower):
                features.append("sidebar")
            # Image count & placement (critical for diversity)
            n_images = len(re.findall(r'<img\b', html_lower))
            if n_images:
                # Detect image placement pattern
                # Check if image is in the first/left portion of a flex/grid container
                img_left = bool(re.search(
                    r'(?:flex|grid).*?<img|<img.*?(?:order:\s*-?\d|grid-column:\s*1)',
                    html_lower, re.DOTALL
                ))
                # Check for large image container (>50% width)
                img_large = bool(re.search(
                    r'(?:width:\s*(?:5[0-9]|6[0-9]|7[0-9]|8[0-9])%|width:\s*(?:6[0-9]{2}|7[0-9]{2}|8[0-9]{2}|9[0-9]{2}|1[0-2][0-9]{2})px).*?<img',
                    html_lower, re.DOTALL
                ))
                if img_large and img_left:
                    features.append(f"LARGE IMAGE LEFT + text right")
                elif img_large:
                    features.append(f"large image ({n_images})")
                elif n_images > 1:
                    features.append(f"{n_images} images")
                else:
                    features.append("1 image")
            # Table
            if '<table' in html_lower:
                features.append("data table")
            # Cards
            card_count = len(re.findall(r'border-radius.*?padding', html_lower))
            if card_count >= 3:
                features.append(f"{card_count} cards")
            # Hero number
            if re.search(r'font-size:\s*(?:4[5-9]|[5-9][0-9])px.*?font-weight:\s*[78]00', html_lower):
                features.append("hero number")
            # Content flow direction
            has_2col = bool(re.search(
                r'(?:display:\s*flex|display:\s*grid).*?(?:flex-direction:\s*row|grid-template-columns)',
                html_lower, re.DOTALL
            ))
            if has_2col and not any('grid' in f.lower() for f in features):
                features.append("horizontal 2-panel split")
            return " + ".join(features) if features else "minimal"

        if _is_freeform or _is_pattern_mode:
            # SERIAL: slide 1 first, then rest with style anchor
            # Pattern mode uses serial to enforce color consistency across slides
            style_anchor = ""
            layout_descriptions = []  # track structure of previous slides
            sorted_slides = sorted(blueprint.slides, key=lambda s: s.slide_id)
            for bp_slide in sorted_slides:
                # Build layout diversity context from previous slides
                # Skip in vocab/pattern mode — layout skeleton handles diversity
                layout_context = ""
                if layout_descriptions and not self._selected_pattern:
                    recent = layout_descriptions[-5:]
                    # Detect overused patterns
                    recent_joined = " ".join(recent).lower()
                    img_left_count = recent_joined.count("image left")
                    two_panel_count = recent_joined.count("2-panel") + recent_joined.count("flexbox columns")

                    # Build dynamic avoidance based on what's overused
                    avoid_lines = []
                    if img_left_count >= 2:
                        avoid_lines.append(
                            "  ⛔ 'LARGE IMAGE LEFT + text blocks right' is BANNED — "
                            "already used {n}× in this deck. Place images differently: "
                            "IMAGE ON TOP (full-width, short), IMAGE RIGHT + text left, "
                            "image as SMALL THUMBNAIL in a card, image BELOW the explanation.".format(n=img_left_count)
                        )
                    if two_panel_count >= 2:
                        avoid_lines.append(
                            "  ⛔ Horizontal 2-panel split is overused ({n}×). Try: "
                            "VERTICAL stack (top section + bottom section), "
                            "3-column layout, single centered content block, "
                            "asymmetric L-shaped layout.".format(n=two_panel_count)
                        )

                    layout_context = (
                        "\n## Layout Diversity (CRITICAL — do NOT repeat bones)\n"
                        "Previous slides used these structures:\n"
                        + "\n".join(f"  - Slide {i+1}: {desc}" for i, desc in enumerate(recent))
                        + "\n"
                        + ("\n".join(avoid_lines) + "\n" if avoid_lines else "")
                        + "\n**Design a COMPLETELY DIFFERENT page skeleton.** "
                        "Structural alternatives you SHOULD consider:\n"
                        "  • Image on TOP (full-width banner) + text below in cards\n"
                        "  • TEXT LEFT (wide) + small visual/chart RIGHT\n"
                        "  • Centered hero content (no sidebars, no 2-column)\n"
                        "  • Vertical 3-row stack: headline → visual → data strip\n"
                        "  • Grid of 4-6 equal cards (no large single image)\n"
                        "  • Full-width chart/table with caption below\n"
                        "\nVary the SPATIAL DISTRIBUTION of content, not just "
                        "the structural CSS. If you've been putting large visuals LEFT, "
                        "try RIGHT or TOP or CENTERED.\n"
                    )

                sid, html, error = _gen_html(bp_slide, style_anchor=style_anchor + layout_context)
                html_results[sid] = (html, error)
                # After first successful slide, extract its style as anchor
                if not style_anchor and html:
                    style_anchor = _extract_style_anchor(html)
                    if style_anchor:
                        logger.info("Style anchor extracted from slide %d (%d chars)",
                                    sid, len(style_anchor))
                # Track layout immediately so next slide gets diversity constraint
                if html:
                    layout_label = self._detect_layout(html)
                    self.used_layouts.append(layout_label)
                    layout_descriptions.append(_describe_layout(html))
        else:
            # PARALLEL: default prompt has rigid style system
            try:
                configured_parallel = int(
                    os.environ.get("HTML_CODEGEN_MAX_PARALLEL", "16")
                )
            except ValueError:
                logger.warning(
                    "Invalid HTML_CODEGEN_MAX_PARALLEL=%r; using 16",
                    os.environ.get("HTML_CODEGEN_MAX_PARALLEL"),
                )
                configured_parallel = 16
            max_parallel = max(
                1, min(configured_parallel, 16, len(blueprint.slides))
            )
            logger.info(
                "Generating %d HTML slides with max_parallel=%d",
                len(blueprint.slides),
                max_parallel,
            )
            with ThreadPoolExecutor(max_workers=max_parallel) as pool:
                futures = {}
                for bp_slide in blueprint.slides:
                    fut = pool.submit(_gen_html, bp_slide)
                    futures[fut] = bp_slide
                for fut in as_completed(futures):
                    bp_slide = futures[fut]
                    sid, html, error = fut.result()
                    html_results[sid] = (html, error)
                if html:
                    logger.info("Slide %d HTML generated successfully", sid)
                else:
                    logger.warning("Slide %d HTML generation failed: %s", sid, error)

        # Phase 2: Render PNGs and assemble PPTX (sequential, needs Playwright + prs)

        # Deck-level frame contract: normalize ALL slides to same header treatment
        # before rendering, to prevent mixed light/filled headers across the deck.
        if not self._selected_pattern and self._theme is not None:
            from ..html_codegen.deck_frame_contract import enforce_html_deck_frame_contract
            deck_codes = {
                bp_slide.slide_id: html_results.get(bp_slide.slide_id, (None, None))[0]
                for bp_slide in blueprint.slides
                if html_results.get(bp_slide.slide_id, (None, None))[0]
            }
            if deck_codes:
                normalized = enforce_html_deck_frame_contract(deck_codes, self._theme)
                for sid, html in normalized.items():
                    old_html, error = html_results[sid]
                    html_results[sid] = (html, error)
                logger.info("Deck frame contract applied: %d slides normalized to majority treatment", len(normalized))

        for bp_slide in blueprint.slides:
            sid = bp_slide.slide_id
            html, error = html_results.get(sid, (None, "not generated"))

            if html:
                # Skip per-slide frame contract — already done at deck level above
                # Only apply for pattern mode or if deck-level wasn't run
                if self._selected_pattern or self._theme is None:
                    pass  # pattern mode skips frame contract entirely
                # (deck-level already applied above for default mode)

                # Post-process: fix title overflow for all slides
                html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

                # Render HTML to PNG
                png_dir = (code_dir.parent if code_dir else Path("/tmp")) / "html_renders"
                png_dir.mkdir(parents=True, exist_ok=True)
                png_path = png_dir / f"slide_{sid:02d}.png"
                render_success = self._render_html_to_png(html, png_path)

                if render_success:
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(
                        str(png_path),
                        Inches(0), Inches(0),
                        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                    )
                    title_text, body_text = self._extract_text_from_html(html)
                    src_evidence = self._build_source_evidence_for_notes(
                        bp_slide, evidence=evidence, source_store=source_store,
                    )
                    self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)

                    self.slide_codes[sid] = _sanitize_meta_text(html)
                    if code_dir:
                        html_file = code_dir / f"slide_{sid:02d}.html"
                        html_file.write_text(self.slide_codes[sid], encoding="utf-8")

                    layout_label = self._detect_layout(html)
                    self.used_layouts.append(layout_label)
                    for img_match in re.finditer(r'src=["\']([^"\']+)["\']', html):
                        self.used_image_ids.add(Path(img_match.group(1)).name)

                    grammar = self._slide_grammars.get(sid)
                    variant = self._slide_variants.get(sid)
                    skills = self._slide_skills.get(sid, [])
                    manifest["slides"].append({
                        "slide_id": sid,
                        "status": "ok",
                        "layout_grammar": grammar.grammar_id if grammar else "unknown",
                        "layout_grammar_name": grammar.name if grammar else "Unknown",
                        "composition_variant": variant.variant_id if variant else "unknown",
                        "composition_variant_name": variant.name if variant else "Unknown",
                        "visual_skills": [skill.skill_id for skill in skills],
                    })
                    continue

            # Fallback
            logger.error("Slide %d: using fallback. Error: %s", sid, error)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            self._add_fallback_slide(slide, bp_slide, evidence)
            manifest["slides"].append({"slide_id": sid, "status": "fallback", "error": error})

        # Save PPTX
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        # Close browser after deck compilation
        self._close_browser()

        manifest["timing_sec"] = round(time.time() - start, 2)
        logger.info(
            "HtmlCodeGenCompiler: compiled %d slides to %s in %.1fs",
            len(blueprint.slides), output_path, manifest["timing_sec"],
        )
        return manifest

    def repair_slide(
        self,
        slide_id: int,
        issues: list[dict],
        blueprint_slide: BlueprintSlide | None,
        evidence: EvidenceState,
        case_dir: str | Path,
        source_store=None,
    ) -> str | None:
        """Repair a slide's HTML code based on issues found.

        Returns the new HTML string, or None if repair failed.
        """
        current_code = self.slide_codes.get(slide_id)
        if not current_code:
            logger.warning("No HTML found for slide %d, cannot repair", slide_id)
            return None

        issue_parts = []
        for iss in issues:
            line = f"- [{iss.get('severity', 'major')}] {iss.get('description', str(iss))}"
            if iss.get('planned_fix'):
                line += f"\n  FIX: {iss['planned_fix']}"
            issue_parts.append(line)
        issue_text = "\n".join(issue_parts)

        brief = ""
        if blueprint_slide:
            brief = f"\nSlide role: {blueprint_slide.role}\nSlide goal: {blueprint_slide.primary_proposition}\n"

        grammar = self._slide_grammars.get(slide_id)
        content_text = ""
        if blueprint_slide is not None:
            content_text = " ".join([
                blueprint_slide.primary_proposition or "",
                getattr(blueprint_slide, "layout_hint", "") or "",
                " ".join(blueprint_slide.must_cover_subset or []),
            ])
            if grammar is None:
                grammar = select_layout_grammar(
                    blueprint_slide.role,
                    layout_hint=getattr(blueprint_slide, "layout_hint", "") or "",
                    has_images="<img" in current_code.lower(),
                    has_table="<table" in current_code.lower(),
                    content_text=content_text,
                )
                self._slide_grammars[slide_id] = grammar

        design_contract = ""
        if grammar is not None:
            if self._selected_pattern:
                # Vocab/pattern mode: skip all old design system elements
                # The vocab contract handles layout + color + decoration
                design_contract = ""
            elif self._theme is not None or self._demo_palette:
                design_theme = _theme_with_palette_override(self._theme or DEFAULT_THEME, self._demo_palette)
                design_contract = format_html_design_contract(design_theme, grammar)
            else:
                design_contract = (
                    f"## Layout Grammar - {grammar.name} (`{grammar.grammar_id}`)\n\n"
                    f"- Composition: {grammar.composition}\n"
                    "- Preserve the existing HTML palette and CSS variables exactly."
                )

        # Skip composition variant and visual skills in vocab/pattern mode
        # — the vocab contract's layout skeleton replaces these
        if self._selected_pattern:
            variant_context = ""
            skill_context = ""
        else:
            variant = self._slide_variants.get(slide_id)
            if variant is None and grammar is not None and blueprint_slide is not None:
                variant = select_composition_variant(
                    grammar,
                    slide_role=blueprint_slide.role,
                    slide_index=slide_id,
                    layout_hint=getattr(blueprint_slide, "layout_hint", "") or "",
                    has_images="<img" in current_code.lower(),
                    has_table="<table" in current_code.lower(),
                    content_text=content_text,
                )
                self._slide_variants[slide_id] = variant
            variant_context = format_composition_variant_contract(variant) if variant else ""

            skills = self._slide_skills.get(slide_id, [])
            if not skills and grammar is not None and blueprint_slide is not None:
                skills = select_visual_skills(
                    grammar,
                    slide_role=blueprint_slide.role,
                    content_text=content_text,
                    has_images="<img" in current_code.lower(),
                    has_table="<table" in current_code.lower(),
                    evidence_item_count=len(blueprint_slide.must_cover_subset or []),
                    limit=3,
                )
                self._slide_skills[slide_id] = skills
            skill_context = format_visual_skill_references(skills, include_code=False)

        # Build evidence context (reuse shared method)
        evidence_text = self._build_evidence_context(
            blueprint_slide, evidence, case_dir, source_store=source_store,
        )

        # Global task brief for repair context
        task_brief_section = ""
        task_brief = getattr(self, '_task_brief', '')
        if task_brief:
            task_brief_section = f"\n## Global Instructions (from task brief)\n\n{task_brief}\n\n---\n"

        user_content = f"""## Current HTML Code

```html
{current_code}
```

## Issues Found

{issue_text}

## Slide Context
{brief}
{evidence_text}
{design_contract}
{self._format_pattern_contract()}{variant_context}
{skill_context}
{task_brief_section}Fix all the issues above. Return the complete updated HTML page.
IMPORTANT: Use ONLY numbers and facts from the Source Evidence above.
CRITICAL: Do NOT change any text content (numbers, facts, bullet text, table data) unless an issue SPECIFICALLY flags that text as incorrect or fabricated. Layout/spatial fixes must ONLY adjust CSS properties (position, size, font-size, padding, margin). Changing text while fixing layout is the #1 cause of regression.
PRESERVE STYLE: Keep the semantic palette CSS variables, selected layout grammar, composition variant, and visual idioms. Do not replace the slide with a generic template or a new card grid.
"""

        try:
            response = self.llm.call_text(
                system_prompt=self.repair_prompt,
                user_content=user_content,
                model=self.repair_model,
                module_name="html_code_repair",
                prompt_version="slide_html_repair.v1",
                max_tokens=8000,
                temperature=0.2,
            )
            html = self._extract_html(response)
            if html:
                # Content-loss safeguard
                original_content = self._estimate_text_content(current_code)
                repaired_content = self._estimate_text_content(html)
                if original_content > 50 and repaired_content < original_content * 0.7:
                    logger.warning(
                        "Slide %d HTML repair REJECTED: content loss %.0f%%. Keeping original.",
                        slide_id, (1 - repaired_content / original_content) * 100,
                    )
                    return None
                html = self._apply_frame_contract(html, slide_id)
                self.slide_codes[slide_id] = html
                return html
        except Exception as e:
            logger.error("HTML repair failed for slide %d: %s", slide_id, e)

        return None

    def recompile_deck(
        self,
        blueprint: DeckBlueprint,
        case_dir: str | Path,
        output_path: str | Path,
        code_dir: str | Path | None = None,
        evidence=None,
        source_store=None,
    ) -> dict:
        """Recompile all slides using stored/updated HTML code.

        Used after repair to rebuild the PPTX with patched HTML.
        """
        start = time.time()
        case_dir = Path(case_dir)
        output_path = Path(output_path)

        # Ensure paper_title is set for title overflow fix
        if not getattr(self, '_paper_title', ''):
            import re as _re
            if blueprint.slides:
                prop = blueprint.slides[0].primary_proposition or ""
                title_cut = _re.split(
                    r'\s+(?:introduces?|presents?|proposes?|describes?|provides?|is\s+a|addresses)\s+',
                    prop, maxsplit=1
                )
                self._paper_title = title_cut[0].strip() if title_cut else ""

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        manifest = {
            "pptx_path": str(output_path),
            "total_slides": len(blueprint.slides),
            "slides": [],
            "warnings": [],
            "timing_sec": 0,
            "compiler": "html_codegen",
            "design_profile": (
                {
                    "selection": self._theme_id_override or "auto",
                    "theme_id": self._theme.theme_id,
                    "theme_name": self._theme.theme_name,
                    "style_family": self._theme.style_family,
                    "codegen_prompt": self._codegen_prompt_name,
                    "codegen_prompt_version": HTML_CODEGEN_PROMPT_VERSION,
                    "visual_skill_library_version": VISUAL_SKILL_LIBRARY_VERSION,
                    "composition_variant_version": COMPOSITION_VARIANT_VERSION,
                }
                if self._theme is not None else None
            ),
        }

        for bp_slide in blueprint.slides:
            sid = bp_slide.slide_id
            html = self.slide_codes.get(sid)

            if not html:
                # Fallback: add a simple text slide
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                self._add_fallback_slide(slide, bp_slide)
                manifest["slides"].append({
                    "slide_id": sid,
                    "status": "fallback",
                    "warning": "No HTML available",
                })
                continue

            # Render HTML to PNG
            png_dir = output_path.parent / "html_renders"
            png_dir.mkdir(parents=True, exist_ok=True)
            png_path = png_dir / f"slide_{sid:02d}.png"

            # Keep repaired/prebuilt HTML on the same deck-frame contract.
            html = self._apply_frame_contract(html, sid)
            self.slide_codes[sid] = html

            # Post-process title overflow
            html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

            success = self._render_html_to_png(html, png_path)

            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            if success and png_path.exists():
                slide.shapes.add_picture(
                    str(png_path),
                    Inches(0), Inches(0),
                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                )
                # Add text content as speaker notes for evaluator
                title_text, body_text = self._extract_text_from_html(html)
                src_evidence = self._build_source_evidence_for_notes(
                    bp_slide, evidence=evidence, source_store=source_store,
                )
                self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)
                grammar = self._slide_grammars.get(sid)
                variant = self._slide_variants.get(sid)
                skills = self._slide_skills.get(sid, [])
                manifest["slides"].append({
                    "slide_id": sid,
                    "status": "ok",
                    "layout_grammar": grammar.grammar_id if grammar else "unknown",
                    "layout_grammar_name": grammar.name if grammar else "Unknown",
                    "composition_variant": variant.variant_id if variant else "unknown",
                    "composition_variant_name": variant.name if variant else "Unknown",
                    "visual_skills": [skill.skill_id for skill in skills],
                })
            else:
                logger.warning("Slide %d recompile render failed, using fallback", sid)
                self._add_fallback_slide(slide, bp_slide)
                manifest["warnings"].append(f"Slide {sid} render failed")
                manifest["slides"].append({
                    "slide_id": sid, "status": "error_fallback",
                })

            # Save HTML
            if code_dir:
                html_path = Path(code_dir) / f"slide_{sid:02d}.html"
                html_path.write_text(html, encoding="utf-8")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        # Close browser
        self._close_browser()

        manifest["timing_sec"] = round(time.time() - start, 2)
        return manifest

    # ------------------------------------------------------------------
    # Internal methods
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_text_from_html(html: str) -> tuple[str, str]:
        """Extract title and body text from HTML for use as PPTX notes.

        Returns (title, body_text) where title is the first heading content
        and body_text is the remaining visible text.
        """
        from ...utils.html_text import extract_title_and_body
        return extract_title_and_body(html)

    @staticmethod
    def _build_source_evidence_for_notes(
        bp_slide,
        evidence=None,
        source_store=None,
    ) -> str:
        """Build source evidence text for speaker notes enrichment.

        Collects the key factual content from the source evidence bundle
        so it can be included in speaker notes for content-fidelity
        evaluation.
        """
        parts: list[str] = []

        # V2 path: source_store bundle
        if source_store is not None:
            bundle = source_store.get_bundle(bp_slide.slide_id)
            if bundle and bundle.source_text:
                parts.append(bundle.source_text[:3000])
                for s in bundle.table_summaries[:3]:
                    parts.append(s[:500])

        # Legacy path: linked evidence chunks
        if not parts and evidence is not None and hasattr(evidence, 'chunks'):
            for eid in bp_slide.linked_evidence_ids[:5]:
                for chunk in evidence.chunks:
                    if chunk.chunk_id == eid:
                        parts.append(f"{chunk.chunk_id}: {chunk.content[:600]}")
                        break

        # Add must_cover items for context
        if bp_slide.must_cover_subset:
            parts.append("Must cover: " + "; ".join(bp_slide.must_cover_subset))

        return "\n".join(parts)[:4000]

    @staticmethod
    def _add_notes_to_slide(
        slide, title: str, body_text: str,
        source_evidence: str = "",
    ):
        """Add text content as speaker notes to a PPTX slide.

        This enables the evaluator to read slide content even though
        the visual content is baked into a full-bleed PNG image.

        Args:
            slide: python-pptx Slide object.
            title: Extracted title text from HTML.
            body_text: Extracted body text from HTML.
            source_evidence: Additional source evidence text from the
                blueprint/source_store bundle. Appended to notes so
                content-quiz evaluators can access detailed facts that
                may not be visually rendered on the slide.
        """
        if not title and not body_text and not source_evidence:
            return
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        # Clear existing notes
        notes_tf.text = ""
        # Add title line
        if title:
            notes_tf.text = f"[TITLE] {title}"
        # Add body text
        if body_text:
            p = notes_tf.add_paragraph()
            p.text = f"[CONTENT] {body_text[:3000]}"
        # Add source evidence for richer content fidelity
        if source_evidence:
            p = notes_tf.add_paragraph()
            p.text = f"[SOURCE EVIDENCE] {source_evidence[:4000]}"

    def _generate_slide(
        self,
        prs: Presentation,
        bp_slide: BlueprintSlide,
        evidence: EvidenceState,
        image_dir: Path | None,
        available_images: list[str],
        code_dir: Path | None,
        total_slides: int = 10,
        source_store=None,
    ) -> dict:
        """Generate HTML for one slide, render to PNG, and insert into PPTX."""
        sid = bp_slide.slide_id

        # Build user prompt
        user_content, has_images = self._build_slide_prompt(
            bp_slide, evidence, image_dir, available_images,
            total_slides=total_slides,
            source_store=source_store,
        )

        html = None
        last_error = None

        for attempt in range(MAX_CODE_RETRIES):
            try:
                if attempt == 0:
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=user_content,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version=HTML_CODEGEN_PROMPT_VERSION,
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid},
                    )
                else:
                    retry_prompt = (
                        f"{user_content}\n\n"
                        f"## Previous Attempt Failed\n"
                        f"Error: {last_error}\n\n"
                        f"Fix the error and return corrected HTML."
                    )
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=retry_prompt,
                        model=self.model,
                        module_name="slide_html_codegen",
                        prompt_version=f"{HTML_CODEGEN_PROMPT_VERSION}.retry",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "attempt": attempt},
                    )

                html = self._extract_html(response)
                if not html:
                    last_error = "Failed to extract HTML from response"
                    logger.warning("Slide %d attempt %d: no HTML extracted", sid, attempt + 1)
                    continue

                # Validate HTML has basic structure
                if "<body" not in html.lower():
                    last_error = "HTML missing <body> tag"
                    logger.warning("Slide %d attempt %d: %s", sid, attempt + 1, last_error)
                    continue

                # Render HTML to PNG
                png_dir = code_dir.parent / "html_renders" if code_dir else Path("/tmp/html_renders")
                png_dir.mkdir(parents=True, exist_ok=True)
                png_path = png_dir / f"slide_{sid:02d}.png"

                # Apply the same deck-frame color contract in the legacy
                # single-slide generation path.
                html = self._apply_frame_contract(html, sid)

                # Post-process title overflow
                html = self._fix_title_overflow(html, sid, getattr(self, '_paper_title', ''))

                render_success = self._render_html_to_png(html, png_path)
                if not render_success:
                    last_error = "Playwright rendering failed"
                    logger.warning("Slide %d attempt %d: %s", sid, attempt + 1, last_error)
                    continue

                # Insert PNG into PPTX slide
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(
                    str(png_path),
                    Inches(0), Inches(0),
                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
                )

                # Extract text from HTML and add as speaker notes
                # so the evaluator can read slide content
                title_text, body_text = self._extract_text_from_html(html)
                src_evidence = self._build_source_evidence_for_notes(
                    bp_slide, evidence=evidence, source_store=source_store,
                )
                self._add_notes_to_slide(slide, title_text, body_text, source_evidence=src_evidence)

                self.slide_codes[sid] = _sanitize_meta_text(html)
                if code_dir:
                    html_file = code_dir / f"slide_{sid:02d}.html"
                    html_file.write_text(html, encoding="utf-8")

                # Track layout used for diversity
                layout_label = self._detect_layout(html)
                self.used_layouts.append(layout_label)

                # Track images used
                for img_match in re.finditer(r'src=["\']([^"\']+)["\']', html):
                    img_path = img_match.group(1)
                    self.used_image_ids.add(Path(img_path).name)

                logger.info("Slide %d generated successfully via HTML (attempt %d)", sid, attempt + 1)
                return {"slide_id": sid, "status": "ok", "attempts": attempt + 1}

            except Exception as e:
                last_error = str(e)
                logger.warning("Slide %d attempt %d error: %s", sid, attempt + 1, str(e)[:200])

        # All retries exhausted — fallback
        logger.error(
            "Slide %d: all %d HTML attempts failed, using fallback. Last error: %s",
            sid, MAX_CODE_RETRIES, last_error
        )
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_fallback_slide(slide, bp_slide, evidence)
        return {
            "slide_id": sid,
            "status": "fallback",
            "error": last_error,
            "attempts": MAX_CODE_RETRIES,
        }

    def _generate_code_only(
        self,
        bp_slide: "BlueprintSlide",
        evidence: "EvidenceState",
        image_dir: Path | None,
        available_images: list[str],
        total_slides: int = 10,
        source_store=None,
    ) -> tuple[str | None, dict]:
        """Generate HTML code for a slide WITHOUT rendering or inserting into PPTX.

        Used by regen_slide tool in the repair agent.

        Returns:
            (html, info_dict) where html is the validated HTML string
            (or None on failure).
        """
        sid = bp_slide.slide_id

        # Build user prompt
        user_content, has_images = self._build_slide_prompt(
            bp_slide, evidence, image_dir, available_images,
            total_slides=total_slides,
            source_store=source_store,
        )

        html = None
        last_error = None

        for attempt in range(MAX_CODE_RETRIES):
            try:
                if attempt == 0:
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=user_content,
                        model=self.repair_model,
                        module_name="slide_html_regen",
                        prompt_version=f"{HTML_CODEGEN_PROMPT_VERSION}.regen",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "regen": True},
                    )
                else:
                    retry_prompt = (
                        f"{user_content}\n\n"
                        f"## Previous Attempt Failed\n"
                        f"Error: {last_error}\n\n"
                        f"Fix the error and return corrected HTML."
                    )
                    response = self.llm.call_text(
                        system_prompt=self.system_prompt,
                        user_content=retry_prompt,
                        model=self.repair_model,
                        module_name="slide_html_regen",
                        prompt_version=f"{HTML_CODEGEN_PROMPT_VERSION}.regen.retry",
                        max_tokens=8000,
                        temperature=0.3,
                        input_packet={"slide_id": sid, "attempt": attempt, "regen": True},
                    )

                html = self._extract_html(response)
                if not html:
                    last_error = "Failed to extract HTML from response"
                    continue

                if "<body" not in html.lower():
                    last_error = "HTML missing <body> tag"
                    continue

                logger.info("Slide %d regen code generated (attempt %d)", sid, attempt + 1)
                return html, {"slide_id": sid, "status": "ok", "attempts": attempt + 1}

            except Exception as e:
                last_error = str(e)
                logger.warning("Slide %d regen attempt %d error: %s", sid, attempt + 1, str(e)[:200])

        logger.error("Slide %d: regen failed after %d attempts", sid, MAX_CODE_RETRIES)
        return None, {"slide_id": sid, "status": "error", "error": last_error}

    # ------------------------------------------------------------------
    # Title-slide enrichment helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_method_summary_for_title(
        source_store,
        max_chars: int = 1200,
    ) -> str:
        """Pull a concise method-overview blurb from the SourceStore.

        Strategy:
        1. Prefer the *summary* field of the first ``role=method`` doc block.
        2. Fall back to the first ``role=overview`` block's summary (which
           typically contains the abstract and may sketch the method).
        3. If summaries are empty, concatenate the first few atomic blocks
           of the method doc block (truncated).

        Returns an empty string when nothing useful is found.
        """
        if source_store is None:
            return ""

        plan = getattr(source_store, 'doc_block_plan', None)
        if plan is None:
            return ""

        method_db = None
        overview_db = None
        for db in plan.blocks:
            if method_db is None and db.role == "method":
                method_db = db
            if overview_db is None and db.role == "overview":
                overview_db = db
            if method_db and overview_db:
                break

        # Try method summary first
        if method_db and method_db.summary and len(method_db.summary.strip()) > 40:
            return method_db.summary.strip()[:max_chars]

        # Try overview summary (abstract) — often sketches the method
        if overview_db and overview_db.summary and len(overview_db.summary.strip()) > 40:
            return overview_db.summary.strip()[:max_chars]

        # Last resort: splice the first few atomic blocks of the method DB
        if method_db and method_db.included_atomic_block_ids:
            parts: list[str] = []
            total = 0
            for bid in method_db.included_atomic_block_ids[:8]:
                block = source_store.get_block(bid)
                if block:
                    parts.append(block.text)
                    total += len(block.text)
                    if total > max_chars:
                        break
            if parts:
                return "\n".join(parts)[:max_chars]

        return ""

    def _format_pattern_contract(self, slide_role: str | None = None) -> str:
        """Format selected style pattern as deck design language block."""
        if not self._selected_pattern:
            return ""
        # Vocab mode: generate per-slide contract with role-aware layout
        if hasattr(self, '_vocab_contract') and self._vocab_contract:
            from app.style_patterns.vocab_composer import format_vocab_style_contract
            return format_vocab_style_contract(
                self._paper_title,
                slide_role=slide_role,
            )
        from app.style_patterns import format_deck_style_contract
        return format_deck_style_contract(self._selected_pattern)

    _deco_seeds_cache: dict | None = None

    def _inject_deco_seeds(self, parts: list, slide_id: int, role: str = "data"):
        """Inject 4-6 decorative CSS seed elements based on slide role.

        Role-aware selection:
        - data/results: accent bars, table highlights, chart decorations
        - method/process: timeline connectors, step dots, flow arrows
        - title/conclusion: large shapes, background gradients, hero accents
        - comparison: divider lines, column separators, badges
        """
        import json, re, random
        from pathlib import Path

        # Lazy load and categorize
        if self._deco_seeds_cache is None:
            vocab_path = Path(__file__).parent.parent / "style_patterns" / "element_vocab.json"
            if not vocab_path.exists():
                self.__class__._deco_seeds_cache = {"line": [], "shape": [], "gradient": [], "svg": []}
                return

            vocab = json.loads(vocab_path.read_text())
            cache = {"line": [], "shape": [], "gradient": [], "svg": []}
            for e in vocab:
                if e['category'] not in ('decoration', 'separator', 'background'):
                    continue
                code = e['code']
                if e['category'] == 'background':
                    if len(code) > 200 or 'gradient' not in code:
                        continue
                    cache["gradient"].append(code)
                elif len(code) > 250:
                    continue
                elif re.search(r'>[A-Z][a-z]', code):
                    continue
                elif re.search(r'width:\s*[5-9]\d\dpx', code):
                    continue
                elif len(re.findall(r'#[0-9a-fA-F]{6}', code)) > 3:
                    continue
                elif '<svg' in code:
                    cache["svg"].append(code)
                elif any(x in code for x in ['width: 1px', 'height: 1px', 'height: 2px', 'height: 3px', 'width: 2px', 'divider', 'rule', 'line']):
                    cache["line"].append(code)
                else:
                    cache["shape"].append(code)

            self.__class__._deco_seeds_cache = cache
            total = sum(len(v) for v in cache.values())
            logger.info("Loaded %d decorative seeds (line=%d, shape=%d, gradient=%d, svg=%d)",
                        total, len(cache["line"]), len(cache["shape"]), len(cache["gradient"]), len(cache["svg"]))

        cache = self._deco_seeds_cache
        if not any(cache.values()):
            return

        # Role-based selection weights
        role_lower = (role or "data").lower()
        rng = random.Random(slide_id * 7919)

        if role_lower in ('title', 'conclusion', 'motivation'):
            # Big impact: gradients + large shapes + SVG icons
            pool = [("gradient", 2), ("shape", 2), ("svg", 1)]
        elif role_lower in ('method', 'architecture', 'process', 'roadmap'):
            # Connectors: lines + SVG + shapes
            pool = [("line", 2), ("svg", 2), ("shape", 1)]
        elif role_lower in ('results', 'data', 'evaluation', 'comparison'):
            # Data accents: lines + shapes + gradient subtle
            pool = [("line", 2), ("shape", 2), ("gradient", 1)]
        else:
            # Default: balanced mix
            pool = [("line", 1), ("shape", 2), ("gradient", 1), ("svg", 1)]

        # Build selection
        seeds = []
        for cat, count in pool:
            available = cache.get(cat, [])
            if available:
                seeds.extend(rng.sample(available, min(count, len(available))))

        if not seeds:
            return

        parts.append(f"\n**DECORATIVE ELEMENTS for {role_lower} slide** (adapt colors to your palette):")
        parts.append("```css")
        for seed in seeds[:6]:
            parts.append(seed.strip())
        parts.append("```")
        parts.append("Include these decorative elements adapted to `var(--primary)` / `var(--accent)`. They add visual depth.\n")

    def _build_slide_prompt(
        self,
        bp_slide: BlueprintSlide,
        evidence: EvidenceState,
        image_dir: Path | None,
        available_images: list[str],
        total_slides: int = 10,
        source_store=None,
    ) -> tuple[str, bool]:
        """Build the user prompt for HTML slide generation.

        Returns:
            (prompt_text, has_images)
        """
        parts = []
        has_images = False
        has_table = False
        primary_image_aspect: float | None = None

        # Inject style pattern at the very top — must be first thing LLM sees
        slide_role = getattr(bp_slide, 'role', None) or getattr(bp_slide, 'slide_role', None)
        pattern_contract = self._format_pattern_contract(slide_role=slide_role)
        if pattern_contract:
            parts.append(pattern_contract)
            parts.append("\n---\n")

        # Inject fixed color palette as literal :root block — prevents per-slide color drift
        if self._theme and not self._selected_pattern:
            t = self._theme
            def _h(c): return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"
            parts.append("## MANDATORY COLOR PALETTE (copy this EXACT `:root` block into your CSS)")
            parts.append("```css")
            parts.append(":root {")
            parts.append(f"  --canvas: {_h(t.canvas_color)};")
            parts.append(f"  --ink: {_h(t.ink_color)};")
            parts.append(f"  --primary: {_h(t.primary_color)};")
            parts.append(f"  --secondary: {_h(t.secondary_color)};")
            parts.append(f"  --accent: {_h(t.accent)};")
            parts.append(f"  --support: {_h(t.support_color)};")
            parts.append("}")
            parts.append("```")
            parts.append("Use ONLY these `var(--xxx)` colors. Background: `var(--canvas)`. Text: `var(--ink)`. ALL slides MUST use this identical palette.\n---\n")

        # Inject global task brief / instructions so LLM sees full constraints
        task_brief = getattr(self, '_task_brief', '')
        if task_brief:
            parts.append("## Global Instructions (from task brief)\n")
            parts.append(task_brief)
            parts.append("\n---\n")

        parts.append(f"## Slide {bp_slide.slide_id} of {total_slides}")
        # Detect roadmap/outline slides even if planner assigned a generic role
        effective_role = bp_slide.role
        roadmap_cues = ('roadmap', 'outline', 'agenda', 'talk proceeds', 'overview of the talk')
        _goal_lower = (bp_slide.primary_proposition or '').lower()
        _notes_lower = (bp_slide.notes or '').lower()
        _covers = ' '.join(bp_slide.must_cover_subset or []).lower()
        if effective_role not in ('roadmap', 'outline', 'agenda'):
            if any(c in _goal_lower or c in _notes_lower or c in _covers for c in roadmap_cues):
                effective_role = 'roadmap'
        parts.append(f"**Role**: {effective_role}")
        if effective_role != bp_slide.role:
            parts.append(f"**⚠ This slide is a ROADMAP** — use Pattern 14 (numbered circles), NOT cards.")

        # Mandatory layout assignment — rotate through diverse layouts
        _LAYOUT_POOL = [
            ("full-bleed", "No header band. Title top-left 38px. Hero number 80px+ at right. Dense data table mid-slide. Mini-KPI strip at bottom."),
            ("left sidebar", "Full-height sidebar (25% width) with title + eyebrow + 3 stacked mini-KPIs. Main content on right 75%: dense table or chart + footnotes."),
            ("top-heavy 40/60", "Large top zone (40%) with title + hero number + eyebrow. Bottom 60% packed with data table + SVG chart + source footnotes."),
            ("full-width table", "Thin 50px title strip only. Rest is a massive 15-row data table spanning full width with alternating rows and bold highlight row."),
            ("three-column", "Three equal vertical panels separated by 1px rules. Each panel: section title + different data type (metrics | chart | bullets)."),
            ("dashboard grid", "Top row: 3-4 KPI metric panels (48px number each). Middle: full-width SVG bar chart. Bottom: annotation strip + source."),
            ("timeline/process", "Horizontal timeline with 4-5 connected nodes. Each: numbered circle + bold label + 2-line description. Connected by lines/arrows."),
            ("stepped rows", "5 full-width horizontal strips with alternating tints, each containing: label + value + inline bar + delta indicator."),
            ("center hero", "One massive metric (92px+) or diagram centered (60% of area). 4 mini-KPI panels in corners. Source strip at bottom."),
            ("L-shape", "Connected header band (top) + sidebar (left) forming an L-shape. Content fills remaining bottom-right rectangle with table + chart."),
        ]
        # Select layout based on slide index (deterministic rotation, skip title slides)
        if effective_role in ('title', 'roadmap', 'outline', 'agenda'):
            pass  # Title/roadmap slides keep their special treatment
        else:
            _layout_idx = (bp_slide.slide_id - 1) % len(_LAYOUT_POOL)
            _layout_name, _layout_desc = _LAYOUT_POOL[_layout_idx]
            parts.append(f"\n**MANDATORY LAYOUT — {_layout_name}**:")
            parts.append(f"{_layout_desc}")
            parts.append(f"Follow this spatial structure. Use the SAME background color as all other slides in this deck (from the palette). Do NOT switch between light and dark backgrounds.\n")

        # Inject decorative seed elements for visual diversity (role-aware)
        self._inject_deco_seeds(parts, bp_slide.slide_id, role=effective_role)
        parts.append(f"**Goal**: {bp_slide.primary_proposition}")
        parts.append(f"**Position**: {bp_slide.narrative_position}")
        if bp_slide.notes:
            parts.append(f"**Notes**: {bp_slide.notes}")
        if getattr(bp_slide, 'layout_hint', '') and bp_slide.layout_hint:
            parts.append(f"**Layout hint**: {bp_slide.layout_hint} (use as guidance, adapt as needed)")

        # Layout diversity constraint
        # Layout diversity constraint (skip for imgseed — handled via serial layout_context)
        # Also skip for vocab/pattern mode — layout skeleton handles diversity
        _is_freeform_prompt = 'imgseed' in (getattr(self, '_codegen_prompt_name', '') or '')
        if self.used_layouts and not _is_freeform_prompt and not self._selected_pattern:
            recent = self.used_layouts[-3:] if len(self.used_layouts) >= 3 else self.used_layouts
            parts.append(f"\n**Previous slides used these layouts**: {', '.join(recent)}")
            parts.append("**You MUST use a DIFFERENT layout pattern for this slide.**")

        # Header style rotation — assign a recommended header style per slide
        # to prevent every slide from having the same dark gradient header.
        # NOTE: Skip when style pattern is active — pattern controls all styling
        _skip_header_rotation = bool(self._selected_pattern)
        _is_freeform_prompt = 'imgseed' in (getattr(self, '_codegen_prompt_name', '') or '')
        if not _is_freeform_prompt and not _skip_header_rotation:
            _theme_for_header = self._theme or DEFAULT_THEME
            _is_curated_theme = getattr(_theme_for_header, "style_family", "").startswith("demo_curated")
            if self._demo_palette or _is_curated_theme:
                # Demo palette: vary visual weight while keeping one structural
                # header hue across the full deck. Accent is never a title-band fill.
                _header_cycle = {
                    "title": "primary",
                    "roadmap": "primary",
                    "outline": "light",
                    "motivation": "light",
                    "problem": "light",
                    "context": "light",
                    "background": "light",
                    "method": "primary",
                    "architecture": "primary",
                    "results": "light",
                    "evaluation": "light",
                    "comparison": "light",
                    "conclusion": "primary",
                    "discussion": "light",
                    "limitations": "light",
                }
                if self._demo_palette:
                    dp = self._demo_palette
                else:
                    def _h(c):
                        return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"
                    dp = {
                        "canvas": _h(_theme_for_header.canvas_color),
                        "ink": _h(_theme_for_header.ink_color),
                        "primary": _h(_theme_for_header.primary_color),
                        "secondary": _h(_theme_for_header.secondary_color),
                        "accent": _h(_theme_for_header.accent),
                        "support": _h(_theme_for_header.support_color),
                    }
                _primary_header_text = (
                    f"contrast-safe Canvas text ({dp['canvas']})"
                    if getattr(_theme_for_header, "style_family", "") == "demo_curated_dark"
                    else "white text"
                )
                _header_desc = {
                    "primary": f"Solid PRIMARY header (background: {dp['primary']}, {_primary_header_text}) - the deck's only filled title-band hue",
                    "light": f"Light/editorial header (background: {dp['canvas']}, border-bottom: 4px solid {dp['primary']}, text color: {dp['ink']})",
                }
                _header_default = "primary"
                _available_header_styles = ["primary", "light"]
            else:
                _header_cycle = {
                    "title": "dark",
                    "roadmap": "dark",
                    "outline": "dark",
                    "motivation": "muted_tint",
                    "problem": "muted_tint",
                    "context": "light",
                    "background": "light",
                    "method": "dark",
                    "architecture": "dark",
                    "results": "accent",
                    "evaluation": "light",
                    "comparison": "light",
                    "conclusion": "accent",
                    "discussion": "muted_tint",
                    "limitations": "muted_tint",
                }
                _header_desc = {
                    "dark": "Dark gradient (PRIMARY_DARK to PRIMARY_MID, white text)",
                    "light": "Light/white header with colored bottom border (border-bottom: 4px solid ACCENT, dark text)",
                    "accent": "Accent-dominant header (ACCENT background, white text)",
                    "muted_tint": "Tinted header (LIGHT_BG background with left ACCENT bar, dark text)",
                }
                _header_default = "dark"
                _available_header_styles = ["dark", "light", "accent", "muted_tint"]
            rec_header = _header_cycle.get(effective_role, _header_default)
            # Prompt construction is parallel, so header selection is purely
            # role-driven and never depends on shared mutation or thread order.
            parts.append(f"\n**Recommended header style**: {rec_header} — {_header_desc.get(rec_header, '')}")
            if effective_role == "title":
                footer_desc = "No takeaway footer. Use only a thin accent rule or compact venue/date label near the bottom."
            elif rec_header in ("light", "muted_tint"):
                if self._demo_palette or _is_curated_theme:
                    footer_desc = (
                        f"Quiet editorial footer: background {dp['canvas']} or a subtle {dp['support']} tint, "
                        f"3-4px top rule in {dp['primary']}, text {dp['ink']}. "
                        "Do NOT use Ink/PRIMARY_DARK as a solid full-width footer fill."
                    )
                else:
                    footer_desc = (
                        "Quiet editorial footer: LIGHT_BG background, 3-4px PRIMARY_MID or ACCENT top rule, "
                        "BODY_TEXT/PRIMARY_DARK text. Do not use a solid dark footer."
                    )
            elif rec_header == "accent":
                if self._demo_palette or _is_curated_theme:
                    footer_desc = (
                        f"Quiet footer: background {dp['canvas']} or a subtle {dp['support']} tint, "
                        f"3-4px top rule in {dp['accent']}, text {dp['ink']}. "
                        "Do not repeat the Accent as a second full-width block."
                    )
                else:
                    footer_desc = (
                        "Quiet footer: LIGHT_BG background, 3-4px ACCENT top rule, BODY_TEXT text. "
                        "Do not repeat the Accent as a second full-width block."
                    )
            else:
                if self._demo_palette or _is_curated_theme:
                    footer_desc = (
                        f"Quiet footer: background {dp['canvas']} or a subtle {dp['primary']} tint, "
                        f"3-4px top rule in {dp['primary']}, text {dp['ink']}. "
                        "Do not repeat the filled header as a second full-width block."
                    )
                else:
                    footer_desc = (
                        "Quiet footer: LIGHT_BG background, 3-4px PRIMARY_MID top rule, "
                        "BODY_TEXT/PRIMARY_DARK text. Do not repeat the filled header "
                        "as a second full-width block."
                    )
            parts.append(f"**Required footer treatment**: {footer_desc}")
            parts.append("Follow the paired header/footer treatment from the design system. This keeps visual rhythm without an unrelated dark block at the bottom.")

        # Evidence text
        parts.append("\n## Content / Evidence")
        evidence_added = False

        # --- V2: use source_store bundle if available ---
        if source_store is not None:
            bundle = source_store.get_bundle(bp_slide.slide_id)
            if bundle and bundle.source_text:
                parts.append(bundle.source_text)

                # Title slides: inject method summary so codegen can build
                # the technical pipeline diagram requested by the system prompt.
                # The title bundle typically only contains front-matter (authors,
                # affiliations, venue).  Without a method overview the LLM falls
                # back to a metadata-only layout.
                if bp_slide.role == "title":
                    _method_blurb = self._extract_method_summary_for_title(
                        source_store, max_chars=1200,
                    )
                    if _method_blurb:
                        parts.append("\n## Method Summary (optional visual reference)")
                        parts.append(
                            "The following describes the paper's core method. "
                            "You may use it to inform a visual element on the title slide "
                            "(a brief process indication, key terms, or a subtle diagram) "
                            "— or focus on the title itself as the hero element. "
                            "Do NOT create a row of equal-sized pipeline cards with icon badges."
                        )
                        parts.append(_method_blurb)

                # For results/data slides with large tables, add row selection reminder
                if bp_slide.role in ("results", "comparison", "evaluation"):
                    table_lines = [l for l in bundle.source_text.split('\n')
                                   if '\t' in l or '|' in l or (l.strip() and
                                   any(c.isdigit() for c in l) and
                                   len(l.split()) >= 3)]
                    if len(table_lines) > 10:
                        parts.append("\n⚠️ The evidence above contains a large table. "
                                     "For the slide, select only the 8–10 most representative rows "
                                     "(e.g., 2–3 per category + best result). "
                                     "Including all rows will cause vertical overflow.")
                for s in bundle.asset_summaries:
                    parts.append(s)
                for s in bundle.table_summaries:
                    parts.append(s)
                evidence_added = True
                logger.debug(
                    "Slide %d: using source_store bundle for initial generation",
                    bp_slide.slide_id,
                )

        # --- Legacy fallback ---
        search_keywords = set()
        for item in bp_slide.must_cover_subset:
            for word in re.split(r'[\s_]+', item):
                cleaned = re.sub(r'[^a-zA-Z]', '', word).lower()
                if len(cleaned) > 3:
                    search_keywords.add(cleaned)
        for word in bp_slide.primary_proposition.split():
            cleaned = re.sub(r'[^a-zA-Z]', '', word).lower()
            if len(cleaned) > 4:
                search_keywords.add(cleaned)

        matched_chunk_ids = set()

        # Pass 1: exact match (resolve DB* → B* via source_store)
        resolved_ids = _resolve_linked_ids(bp_slide.linked_evidence_ids, source_store)
        for eid in resolved_ids:
            for chunk in evidence.chunks:
                if chunk.chunk_id == eid:
                    parts.append(f"\n### {chunk.chunk_id} ({chunk.chunk_type})")
                    parts.append(chunk.content[:2000])
                    evidence_added = True
                    matched_chunk_ids.add(chunk.chunk_id)
                    break

        # Pass 2: fuzzy match
        if not evidence_added and search_keywords:
            for chunk in evidence.chunks:
                if chunk.chunk_id in matched_chunk_ids:
                    continue
                section = chunk.metadata.get("section", "").lower()
                chunk_keywords = set(re.sub(r'[^a-zA-Z\s]', '', section).lower().split())
                overlap = search_keywords & chunk_keywords
                if len(overlap) >= 1 and len(chunk.content.strip()) > 50:
                    parts.append(f"\n### {chunk.chunk_id} ({chunk.chunk_type})")
                    parts.append(chunk.content[:2000])
                    evidence_added = True
                    matched_chunk_ids.add(chunk.chunk_id)
                    if len(matched_chunk_ids) >= 3:
                        break

        if not evidence_added:
            parts.append("No specific evidence chunks linked. Use the slide goal as primary content.")

        # Always inject must_cover_subset regardless of evidence source
        if bp_slide.must_cover_subset:
            parts.append("\n## Required Content Points")
            parts.append("You MUST include ALL of the following specific data/facts in this slide.")
            if self.show_source_citations:
                parts.append("Each item includes a [cite: Page X] tag — you MUST use these page numbers in visible Source citations at the bottom of the slide (e.g., `Source: Page 3, Page 5`).\n")
            else:
                parts.append("Each item includes a [cite: Page X] tag for internal reference only.\n")
            for item in bp_slide.must_cover_subset:
                # Try to find page number for this item from source_store
                page_ref = _find_page_for_must_cover(item, bp_slide, source_store)
                if page_ref:
                    parts.append(f"  - {item} [cite: Page {page_ref}]")
                else:
                    parts.append(f"  - {item}")
            parts.append("\nDo NOT omit any of these points. Use exact figures from the evidence above.")
            parts.append("WARNING: Do NOT add descriptive phrases, marketing language, or technical details not found verbatim in the evidence. Every claim must be traceable to the source text.")
            parts.append("LAYOUT: If all points cannot fit cleanly, use a simpler layout (fewer columns, larger fonts) rather than cramming content.")

        # Tables
        if evidence.tables:
            relevant_tables = []
            for tbl in evidence.tables:
                if tbl.table_id in bp_slide.linked_evidence_ids:
                    relevant_tables.append(tbl)
            if not relevant_tables and search_keywords:
                for tbl in evidence.tables:
                    caption_lower = (tbl.caption or "").lower()
                    content_lower = (tbl.content or "")[:200].lower()
                    if any(kw in caption_lower or kw in content_lower for kw in search_keywords):
                        relevant_tables.append(tbl)
                        if len(relevant_tables) >= 2:
                            break
            if relevant_tables:
                has_table = True
                parts.append("\n## Available Tables")
                # For results/comparison slides, add row selection guidance
                is_data_slide = bp_slide.role in ("results", "comparison", "evaluation")
                for tbl in relevant_tables[:2]:
                    parts.append(f"\nTable: {tbl.table_id}")
                    if tbl.caption:
                        parts.append(f"Caption: {tbl.caption}")
                    if tbl.content:
                        rows = tbl.content.strip().split("\n")
                        if rows:
                            parts.append("```")
                            # Cap at 8 rows for results slides to prevent overflow
                            max_rows = 8 if is_data_slide else 12
                            for row in rows[:max_rows]:
                                parts.append(row)
                            if len(rows) > max_rows:
                                parts.append(f"... ({len(rows) - max_rows} more rows)")
                            parts.append("```")
                    if is_data_slide and tbl.content:
                        row_count = len([r for r in tbl.content.strip().split("\n")
                                        if r.strip() and '---' not in r])
                        if row_count > 8:
                            parts.append(f"NOTE: This table has {row_count} rows. Select the 6–8 most representative entries for the slide (e.g., 2 per category + the best result). Do NOT include all rows — the slide will overflow.")

        # Images — use absolute paths for HTML <img src="">
        relevant_figures = []
        logger.debug("_build_slide_prompt slide %d: image_dir=%s, evidence.figures=%d, assigned_figure_id=%s",
                      bp_slide.slide_id, image_dir, len(evidence.figures),
                      getattr(bp_slide, 'assigned_figure_id', 'N/A'))
        if image_dir and evidence.figures:
            embedded_figs = [f for f in evidence.figures
                if f.figure_type not in ("page_screenshot", "table_screenshot")
                and (f.width or 0) * (f.height or 0) < 4_000_000
                and (f.height or 1) / max(f.width or 1, 1) < 2.5
            ]

            # Fallback: use page screenshots if no embedded figures
            if not embedded_figs:
                embedded_figs = [f for f in evidence.figures
                    if f.figure_type == "page_screenshot"
                    and (f.width or 0) * (f.height or 0) < 4_000_000
                ]

            # Priority 1: Use assigned_figure_id from deck planner (pre-allocated, dedup-safe)
            assigned_fig = None
            if hasattr(bp_slide, 'assigned_figure_id') and bp_slide.assigned_figure_id:
                afid = bp_slide.assigned_figure_id
                # Direct match by figure_id
                for fig in evidence.figures:
                    if fig.figure_id == afid:
                        assigned_fig = fig
                        break
                # Fallback: match via source_store asset_id → image_path → stem
                if not assigned_fig and source_store:
                    for asset in getattr(source_store, 'assets', []):
                        if getattr(asset, 'asset_id', None) == afid:
                            asset_stem = Path(asset.image_path).stem if hasattr(asset, 'image_path') and asset.image_path else None
                            if asset_stem:
                                for fig in evidence.figures:
                                    if fig.figure_id == asset_stem:
                                        assigned_fig = fig
                                        break
                            break
                if not assigned_fig:
                    logger.warning(
                        "Slide %d: assigned_figure_id='%s' NOT found in evidence (%d figs)",
                        bp_slide.slide_id, afid, len(evidence.figures),
                    )

            if assigned_fig and assigned_fig.image_path:
                relevant_figures = [assigned_fig]
                logger.info("Slide %d: using assigned figure '%s' -> %s",
                           bp_slide.slide_id, assigned_fig.figure_id, assigned_fig.image_path)
            elif hasattr(bp_slide, 'assigned_figure_id') and bp_slide.assigned_figure_id is not None:
                # Planner explicitly assigned empty string = no image for this slide
                relevant_figures = []
            else:
                scored_figs = []
                for fig in embedded_figs:
                    score = 0
                    if fig.figure_id in bp_slide.linked_evidence_ids:
                        score += 10
                    if search_keywords:
                        text_pool = f"{fig.caption} {fig.description}".lower()
                        text_words = set(re.sub(r'[^a-zA-Z\s]', '', text_pool).split())
                        overlap = search_keywords & text_words
                        score += len(overlap)
                    if score > 0:
                        scored_figs.append((score, fig))

                scored_figs.sort(key=lambda x: -x[0])
                relevant_figures = []
                for sc, fig in scored_figs:
                    if sc >= 3:
                        fig_filename = Path(fig.image_path).name if fig.image_path else ""
                        if fig_filename in self.used_image_ids:
                            continue
                        relevant_figures.append(fig)
                    if len(relevant_figures) >= 1:
                        break

            if relevant_figures:
                has_images = True
                parts.append(f"\n## Available Images")
                parts.append("Use absolute file paths in <img src=\"...\"> tags.")
                parts.append("CRITICAL: You may ONLY use the EXACT paths listed below.")
                parts.append("ONLY include an image if it directly illustrates THIS slide's topic.")
                parts.append("Use `object-fit:contain` on <img> elements. NEVER use negative margins or object-position:top — these crop the figure content.")
                for fig in relevant_figures:
                    if fig.image_path:
                        abs_path = Path(fig.image_path)
                        px_w = fig.width
                        px_h = fig.height
                        if (not px_w or not px_h) and abs_path.exists():
                            try:
                                from PIL import Image as _PILImage
                                with _PILImage.open(str(abs_path)) as _im:
                                    px_w, px_h = _im.size
                            except Exception:
                                pass
                        size_info = ""
                        aspect_info = ""
                        if px_w and px_h:
                            aspect = round(px_w / px_h, 2)
                            if primary_image_aspect is None:
                                primary_image_aspect = aspect
                            if aspect > 1.3:
                                orientation = "landscape"
                            elif aspect < 0.77:
                                orientation = "portrait"
                            else:
                                orientation = "square"
                            size_info = f"  [{px_w}x{px_h} px, aspect={aspect:.2f}:1, {orientation}]"
                            if aspect > 3.0:
                                aspect_info = (
                                    f"  ⚠️ ULTRA-WIDE image (aspect {aspect:.1f}:1). "
                                    f"Container height = container_width / {aspect:.1f}. "
                                    f"E.g. if container is 800px wide, height must be ~{int(800/aspect)}px. "
                                    f"NEVER use flex:1 or height:100% — compute height from width and this ratio."
                                )
                            elif aspect < 0.4:
                                aspect_info = (
                                    f"  ⚠️ ULTRA-TALL image (aspect {aspect:.2f}:1). "
                                    f"Container width = container_height × {aspect:.2f}. "
                                    f"E.g. if container is 450px tall, width must be ~{int(450*aspect)}px. "
                                    f"NEVER use flex:1 or width:100% — compute width from height and this ratio."
                                )
                            else:
                                aspect_info = (
                                    f"  Layout: {orientation} image (aspect {aspect:.2f}:1). "
                                    f"Size the container so that width/height ≈ {aspect:.2f}. "
                                    + (
                                        "Wide placement recommended (60-70% width)."
                                        if orientation == "landscape"
                                        else "Narrow column recommended (30-40% width)."
                                        if orientation == "portrait"
                                        else "Roughly equal split or ~50% width."
                                    )
                                )
                        desc = fig.description if fig.description and "Embedded image from" not in fig.description else ""
                        caption = fig.caption or ""
                        # Truncate caption to avoid verbatim text walls on slides
                        if caption:
                            words = caption.split()
                            if len(words) > 30:
                                caption = " ".join(words[:30]) + "..."
                        label = desc or caption or "no description"
                        parts.append(f"- \"{abs_path}\"{size_info}")
                        parts.append(f"  Content: {label}")
                        if aspect_info:
                            parts.append(aspect_info)
                parts.append("\nDo NOT use page screenshots. Do NOT stretch images.")
                parts.append("Match the image container's aspect ratio to the image's actual aspect ratio — avoid placing a portrait/square image in a wide horizontal container (creates ugly whitespace).")
                parts.append("⚠️ FIGURE SLIDES OVERFLOW RULE: When a figure occupies the left panel, the body zone has NO room for extra cards below. Structure: [figure-left + sidebar-right] in ONE row. Do NOT add process cards, explanation rows, or summary panels below the figure. The sidebar (max 2 cards) is your ONLY space for supplementary content.")
            else:
                parts.append("\n## NO IMAGES AVAILABLE for this slide")
                parts.append("Do NOT use <img> tags.")
                parts.append("Default to concise editorial typography, aligned evidence groups, and at most two restrained rules or bands.")
                parts.append("Do NOT compensate for missing images with an invented SVG motif, node-link diagram, connector web, decorative geometry, or rows of cards.")
                parts.append("Use SVG only if the supplied evidence explicitly defines a process, sequence, topology, or quantitative relationship.")
        else:
            if bp_slide.role != "title":
                parts.append("\n## NO IMAGES AVAILABLE for this slide")
                parts.append("Do NOT use <img> tags. Prefer concise editorial typography and directly aligned evidence.")
                parts.append("Do NOT invent abstract diagrams, decorative connector systems, or card grids to replace a missing image.")

        # Inject pre-generated chart image if available
        chart_path = getattr(self, '_chart_paths', {}).get(bp_slide.slide_id)
        if chart_path and Path(chart_path).exists():
            has_images = True
            # On results/comparison slides, chart is optional — HTML table should dominate
            is_data_slide = bp_slide.role in ("results", "comparison", "evaluation")
            if is_data_slide:
                parts.append(f"\n## Pre-Generated Chart (OPTIONAL — table is primary)")
                parts.append(f"A chart image has been pre-generated for this slide, but an HTML table should be the PRIMARY visual element.")
                parts.append(f"Build a full-width HTML data table first. You may optionally place this chart in a small supplementary position (max 30% width).")
                parts.append(f"- \"{chart_path}\"")
                parts.append(f"  Content: {getattr(bp_slide, 'viz_data', {}).get('title', 'Chart')}")
            else:
                parts.append(f"\n## Pre-Generated Chart (MANDATORY)")
                parts.append(f"A publication-quality chart has been pre-generated for this slide.")
                parts.append(f"You MUST embed it using <img> — do NOT draw your own CSS/SVG chart.")
                parts.append(f"- \"{chart_path}\"")
                parts.append(f"  Content: {getattr(bp_slide, 'viz_data', {}).get('title', 'Chart')}")
                parts.append(f"  Use: `<img src=\"{chart_path}\" style=\"max-width:100%; max-height:480px; object-fit:contain;\">`")
                parts.append(f"Place the chart as the primary visual element of the slide.")

        # Inject theme palette — ensures consistent colors across all slides
        # Demo palette override: use direct 6-role colors with bright/clean instructions
        if self._demo_palette:
            dp = self._demo_palette
            parts.append("\n## Color Palette (MANDATORY — use ONLY these colors)")
            parts.append(f"- **Canvas**: {dp['canvas']} — slide background. Use this as the main `background-color` for the slide.")
            parts.append(f"- **Ink**: {dp['ink']} — all body text, titles, and labels.")
            parts.append(f"- **Primary**: {dp['primary']} — the only full-width title framing hue. Use for filled title bands, title rules, section dividers, and primary chart marks; do not repeat it as a dominant filled footer.")
            parts.append(f"- **Secondary**: {dp['secondary']} — contrasting companion. Use for secondary chart series, supporting panels, borders, and icon backgrounds; never full-width title framing.")
            parts.append(f"- **Accent**: {dp['accent']} — small focal emphasis for key metrics, important numbers, and active states; never a full-width title band, title rule, or footer.")
            parts.append(f"- **Support**: {dp['support']} — low-emphasis fills. Use for subtle backgrounds, auxiliary chart series, muted containers.")
            parts.append(f"\n**COLOR USAGE RULES — READ CAREFULLY:**")
            parts.append(f"1. **Canvas is the slide background** — the slide should feel BRIGHT and CLEAN, not dark or muddy.")
            parts.append(f"2. **Primary and Secondary are SATURATED structural colors** — use them at full saturation for chart elements and local structure. Accent stays small-area only.")
            parts.append(f"3. **60/30/10 area ratio**: 60% Canvas (white/near-white), 30% Primary+Secondary (structural color), 10% Accent (focal pop).")
            parts.append(f"4. **ONE title hue across the deck** — use Primary as the only solid header background, OR Canvas with a Primary bottom border. Never switch full-width title framing to Secondary, Accent, or Support.")
            parts.append(f"5. **White text on Primary** is fine. Use Ink color for text on Canvas/Support backgrounds.")
            parts.append(f"6. **Avoid muddy/grayish mixes** — keep colors clean and distinct. If you need a lighter shade, use opacity (rgba) rather than mixing with gray.")
            parts.append(f"7. **Takeaway footer stays subordinate**: use Canvas/subtle Primary tint with a Primary top rule and Ink text, even when the title header is filled Primary. The footer is a supporting summary, not a second title band. Never use Ink as a solid footer below a light header.")
            parts.append(f"\n**Color Aliases** (for prompts that reference these names):")
            parts.append(f"- PRIMARY_DARK = Ink ({dp['ink']})")
            parts.append(f"- PRIMARY_MID = Primary ({dp['primary']})")
            parts.append(f"- PRIMARY_LIGHT = Secondary ({dp['secondary']})")
            parts.append(f"- ACCENT = Accent ({dp['accent']})")
            parts.append(f"- LIGHT_BG = Canvas ({dp['canvas']})")
            parts.append(f"- WARM_BG = Support ({dp['support']})")
            parts.append(f"Wherever the system prompt says PRIMARY_DARK, use {dp['ink']}. PRIMARY_MID → {dp['primary']}. ACCENT → {dp['accent']}.")
            parts.append("Exception: on light curated palettes, PRIMARY_DARK/Ink must not be used as a large structural fill for a footer paired with a light header.")
        else:
            theme = self._theme or DEFAULT_THEME

            def _hex(c):
                return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"

            if getattr(theme, "style_family", "").startswith("demo_curated"):
                dp = {
                    "canvas": _hex(theme.canvas_color),
                    "ink": _hex(theme.ink_color),
                    "primary": _hex(theme.primary_color),
                    "secondary": _hex(theme.secondary_color),
                    "accent": _hex(theme.accent),
                    "support": _hex(theme.support_color),
                }
                dark_palette = theme.style_family == "demo_curated_dark"
                parts.append("\n## Color Palette (MANDATORY - use ONLY these colors)")
                parts.append(f"- **Canvas**: {dp['canvas']} - slide background and broad quiet space.")
                parts.append(f"- **Ink**: {dp['ink']} - all body text, titles, labels, and table text.")
                parts.append(f"- **Primary**: {dp['primary']} - dominant structural color for headers, rules, diagram paths, and primary chart marks.")
                parts.append(f"- **Secondary**: {dp['secondary']} - companion structure, secondary chart marks, small panels, and borders; never full-width title framing.")
                parts.append(f"- **Accent**: {dp['accent']} - one small focal metric or callout per slide; never a full-width title band, title rule, or footer.")
                parts.append(f"- **Support**: {dp['support']} - soft fills and auxiliary chart series.")
                parts.append(f"- Font family: {theme.font_family}")
                parts.append(f"\n**Theme: {theme.theme_name}** - EVERY visible color must come from this palette.")
                if dark_palette:
                    parts.append("This is a dark contrast palette: use Canvas for the full slide background, Ink for readable text, and keep Support/Accent as small highlights.")
                else:
                    parts.append("Use a 60/30/10 area ratio: 60% Canvas, 30% Primary/Secondary structure, 10% Accent. Keep the slide bright, clean, and editorial.")
                parts.append("Do not invent gray-blue gradients, muddy desaturated variants, or repeated dark header bands.")
                parts.append("Across this deck, every full-width title band and title rule uses Primary only. Vary filled versus Canvas treatments, never the structural hue.")
                parts.append("Prefer one authored composition per slide over rows of identical cards. Use color to structure hierarchy, not as decoration.")
                if dark_palette:
                    parts.append("Pair headers and footers within the dark palette: use quiet Canvas/subtle Primary-tint footers with a Primary top rule. Do not repeat a filled header as a second dominant footer band.")
                else:
                    parts.append("Pair headers and footers by visual weight: use quiet Canvas/subtle Primary-tint footers with a Primary top rule and Ink text, including after filled headers. Never use Ink or Primary as a solid full-width footer fill.")
                parts.append("\n**Color Aliases** (for prompts that reference legacy names):")
                parts.append(f"- PRIMARY_DARK = Ink ({dp['ink']}) on light palettes; Canvas ({dp['canvas']}) only for dark full-slide backgrounds")
                parts.append(f"- PRIMARY_MID = Primary ({dp['primary']})")
                parts.append(f"- PRIMARY_LIGHT = Secondary ({dp['secondary']})")
                parts.append(f"- ACCENT = Accent ({dp['accent']})")
                parts.append(f"- LIGHT_BG = Canvas ({dp['canvas']})")
                parts.append(f"- WARM_BG = Support ({dp['support']})")
            else:
                def _contrast_on_white(c):
                    """WCAG contrast ratio of color c on white background."""
                    def _srgb(v):
                        v = v / 255
                        return v / 12.92 if v <= 0.04045 else ((v + 0.055) / 1.055) ** 2.4
                    L = 0.2126 * _srgb(c[0]) + 0.7152 * _srgb(c[1]) + 0.0722 * _srgb(c[2])
                    return (1.05) / (L + 0.05)

                def _darken_to_pass(c, min_ratio=4.5):
                    """Darken a color until it passes WCAG contrast on white."""
                    r, g, b = c
                    for _ in range(20):
                        if _contrast_on_white((r, g, b)) >= min_ratio:
                            return (r, g, b)
                        r = max(0, int(r * 0.9))
                        g = max(0, int(g * 0.9))
                        b = max(0, int(b * 0.9))
                    return (r, g, b)

                def _text_ok(c):
                    return "OK for text" if _contrast_on_white(c) >= 4.5 else "DECORATIVE ONLY - too light for text on white"

                primary_mid = theme.primary_mid
                if _contrast_on_white(primary_mid) < 4.5:
                    primary_mid = _darken_to_pass(primary_mid)
                    logger.info("Auto-darkened PRIMARY_MID from %s to %s for WCAG compliance",
                               _hex(theme.primary_mid), _hex(primary_mid))

                accent = theme.accent
                if _contrast_on_white(accent) < 4.5:
                    accent = _darken_to_pass(accent)
                    logger.info("Auto-darkened ACCENT from %s to %s for WCAG compliance",
                               _hex(theme.accent), _hex(accent))

                parts.append("\n## Color Palette (MANDATORY - use ONLY these colors)")
                parts.append(f"- **Primary Dark**: {_hex(theme.primary_dark)} - dark backgrounds, header bars. {_text_ok(theme.primary_dark)}")
                parts.append(f"- **Primary Mid**: {_hex(primary_mid)} - titles, headings. {_text_ok(primary_mid)}")
                parts.append(f"- **Primary Light**: {_hex(theme.primary_light)} - borders, underlines. {_text_ok(theme.primary_light)}")
                parts.append(f"- **Accent**: {_hex(accent)} - emphasis, key metrics. {_text_ok(accent)}")
                accent_alt = theme.accent_alt
                if _contrast_on_white(accent_alt) < 4.5:
                    accent_alt = _darken_to_pass(accent_alt)
                parts.append(f"- **Accent Alt**: {_hex(accent_alt)} - secondary accent. {_text_ok(accent_alt)}")
                parts.append(f"- **Body Text**: {_hex(theme.body_text)} - all body text. {_text_ok(theme.body_text)}")
                caption = theme.caption_text
                if _contrast_on_white(caption) < 4.5:
                    caption = _darken_to_pass(caption)
                parts.append(f"- **Caption**: {_hex(caption)} - footnotes. {_text_ok(caption)}")
                parts.append(f"- **Light BG**: {_hex(theme.light_bg)} - card/container backgrounds only")
                parts.append(f"- **Warm BG**: {_hex(theme.warm_bg)} - alternate backgrounds only")
                parts.append(f"- Font family: {theme.font_family}")
                parts.append(f"\n**Theme: {theme.theme_name}** - EVERY color must come from this palette.")
                parts.append("Colors marked DECORATIVE ONLY must not be used as text color on white; use them only for background, border-color, or decorative elements.")
                parts.append("For bold/emphasized text on white, use Primary Dark or Body Text instead.")
                parts.append("CONTRAST ON COLORED BACKGROUNDS: When placing text on Light BG or Warm BG containers, use ONLY Primary Dark or Body Text.")

        # Select one content-driven macro grammar and a small set of compatible
        # visual idioms. This is the canonical hybrid mode: deck-level palette
        # consistency from ReDeck, composition diversity from the public
        # design contracts, and anti-card restraint from the default prompt.
        content_text = " ".join([
            bp_slide.primary_proposition or "",
            getattr(bp_slide, "layout_hint", "") or "",
            " ".join(bp_slide.must_cover_subset or []),
        ])
        grammar = select_layout_grammar(
            effective_role,
            layout_hint=getattr(bp_slide, "layout_hint", "") or "",
            has_images=has_images,
            has_table=has_table,
            content_text=content_text,
        )
        self._slide_grammars[bp_slide.slide_id] = grammar
        # Skip ThemeColors design contract when pattern mode is active
        # (pattern CSS provides its own color system; theme colors would conflict)
        if not self._selected_pattern:
            design_theme = _theme_with_palette_override(self._theme or DEFAULT_THEME, self._demo_palette)
            parts.append("\n" + format_html_design_contract(design_theme, grammar))

        variant = select_composition_variant(
            grammar,
            slide_role=effective_role,
            slide_index=bp_slide.slide_id,
            total_slides=total_slides,
            layout_hint=getattr(bp_slide, "layout_hint", "") or "",
            has_images=has_images,
            has_table=has_table,
            content_text=content_text,
            image_aspect=primary_image_aspect,
        )
        self._slide_variants[bp_slide.slide_id] = variant
        parts.append("\n" + format_composition_variant_contract(variant))

        body_word_budget = 55 if has_images else 80 if has_table else 65
        if variant.variant_id == "data_source_chart_signal":
            body_word_budget = 42
        elif variant.variant_id in {"figure_hero_band", "figure_center_stage", "figure_strip_stack"}:
            body_word_budget = min(body_word_budget, 45)
        if effective_role in {"context", "agenda", "roadmap", "outline"} and bp_slide.narrative_position == "opening":
            body_word_budget = 50
        elif effective_role == "conclusion":
            body_word_budget = 55 if has_images else 65
        parts.append(f"""
## Slide-Specific Complexity Budget

- Use no more than {body_word_budget} visible body words, excluding the title, author line, and compact source footer.
- Express each Required Content Point once as a short label plus one concise clause.
- Use at most three primary body groups and one subordinate annotation/takeaway region.
- Keep body text in CSS Grid/Flex normal flow. Do not absolutely position multiple prose blocks.
- If the content does not fit at body text >=16px, simplify the composition before writing HTML; never conceal text with `overflow:hidden`.
""")
        if not has_images and not has_table:
            parts.append("""
## No-Image Editorial Discipline

- This slide has no real visual artifact and no table. Do not manufacture visual complexity.
- Do not use inline SVG, node-link diagrams, connector webs, timeline axes, milestone dots, phase paths, or abstract geometry unless the evidence explicitly defines that relationship.
- Do not use `<table>` or matrix grids unless the prompt supplies actual tabular data under Available Tables.
- Use an editorial composition: strong type hierarchy, one aligned evidence list or ranked evidence field, one compact note/metric only when directly supported, and deliberate whitespace.
- Agenda, outline, setup, and conclusion slides should feel like authored navigation or synthesis, not process diagrams.
""")
        if has_images:
            parts.append(f"""
## Figure-Led Restraint

- Follow the selected Composition Variant (`{variant.variant_id}`) for the macro layout.
- The real source figure is the visual explanation. Pair it with only two or three concise observations.
- Do not add a second chart, comparison table, metric-card grid, or decorative diagram unless required evidence cannot be expressed by annotating the figure.
- The figure wrapper must hug the image at its natural aspect ratio; unused blank wrapper area is a layout failure.
- Use exactly one annotation system: side rail, bottom band, perimeter labels, or metric edge.
- If the variant is not `figure_sidecar_left`, do not fall back to the common top-title + left-figure + right-observation-rail composition.
""")

        selected_visual_skills = select_visual_skills(
            grammar,
            slide_role=effective_role,
            content_text=content_text,
            has_images=has_images,
            has_table=has_table,
            evidence_item_count=len(bp_slide.must_cover_subset or []),
            limit=3,
        )
        self._slide_skills[bp_slide.slide_id] = selected_visual_skills

        # ── HTML skill injection ───────────────────────────
        skill_dir = Path(__file__).parent.parent.parent / "prompts" / "codegen" / "skills"
        _is_freeform_prompt = 'imgseed' in (getattr(self, '_codegen_prompt_name', '') or '')
        num_images = len(relevant_figures)

        if not _is_freeform_prompt:
            parts.append("\n" + format_visual_skill_references(selected_visual_skills))

        elif _is_freeform_prompt:
            # For imgseed/freeform: inject GENERAL component skills based on content
            # Not role-based templates, but reusable design patterns
            # Each skill teaches HOW a component should look, not WHERE it goes
            #
            # SKILL LIBRARY (30 active, 4 on-disk reserve):
            #   Layout:     text_layout, stacked_sidebar, two_panel, grid, split_comparison
            #   Data viz:   table, chart, metric, hero_row, stat_dashboard, ablation_strip
            #   Diagrams:   pipeline, timeline, numbered_steps, architecture
            #   Text/anno:  callout, section_badge, formula, definition_list, quote_block, contribution_list
            #   Figures:    figure, annotated_figure, figure_gallery, matrix_layout(v18)
            #   Evidence:   evidence_card, key_finding(v18)
            #   Contrast:   problem_solution, comparison_bar(v18)
            #   Process:    icon_process(v18)
            #   Reserve (v19, on disk but not wired — prompt bloat regressed issue count):
            #     grouped_bar, radar_chart, waterfall, highlight_table
            selected_skills = []

            # ── Universal base skills (always injected) ──
            selected_skills.append("text_layout")
            selected_skills.append("callout")
            selected_skills.append("section_badge")

            # ── Content-driven skill selection ──
            if has_images:
                selected_skills.append("figure")
                selected_skills.append("annotated_figure")
                selected_skills.append("stacked_sidebar")  # figure left + sidebar right is the #1 pattern
                if num_images >= 2:
                    selected_skills.append("figure_gallery")  # multi-panel for 2+ images
                if num_images >= 4:
                    selected_skills.append("matrix_layout")   # v18: 2×2+ grid for many images

            if bp_slide.role in ('results', 'comparison', 'evaluation'):
                selected_skills.append("table")
                selected_skills.append("chart")
                selected_skills.append("metric")
                selected_skills.append("hero_row")
                selected_skills.append("stacked_sidebar")
                selected_skills.append("evidence_card")
                selected_skills.append("key_finding")        # v18: hero result + evidence
                selected_skills.append("comparison_bar")     # v18: win/lose stacked bars
                if bp_slide.role == 'comparison':
                    selected_skills.append("split_comparison")
            elif bp_slide.role in ('method', 'architecture'):
                selected_skills.append("pipeline")
                selected_skills.append("chart")              # v20: method slides can show performance context
                selected_skills.append("formula")
                selected_skills.append("two_panel")
                selected_skills.append("definition_list")
                selected_skills.append("numbered_steps")
                selected_skills.append("architecture")
                selected_skills.append("icon_process")       # v18: high-level conceptual flow
            elif bp_slide.role in ('context', 'background', 'motivation'):
                selected_skills.append("chart")              # v20: context slides often have data
                selected_skills.append("split_comparison")
                selected_skills.append("grid")
                selected_skills.append("definition_list")
                selected_skills.append("contribution_list")
                selected_skills.append("problem_solution")
                selected_skills.append("hero_row")            # v20: hero metrics for context
            elif bp_slide.role in ('conclusion', 'discussion'):
                selected_skills.append("chart")              # v20d: conclusion often summarizes key results
                selected_skills.append("metric")
                selected_skills.append("two_panel")
                selected_skills.append("hero_row")
                selected_skills.append("quote_block")
                selected_skills.append("key_finding")        # v18: spotlight one key conclusion
            elif bp_slide.role in ('title', 'opening'):
                selected_skills.insert(0, "title")  # Title skill FIRST for opening slides
                selected_skills.append("metric")
                selected_skills.append("hero_row")
                selected_skills.append("contribution_list")
            elif bp_slide.role in ('roadmap', 'outline', 'overview'):
                selected_skills.append("timeline")
                selected_skills.append("grid")
                selected_skills.append("numbered_steps")
                selected_skills.append("icon_process")       # v18: conceptual stage overview
            elif bp_slide.role in ('ablation',):
                selected_skills.append("table")
                selected_skills.append("chart")
                selected_skills.append("metric")
                selected_skills.append("hero_row")
                selected_skills.append("ablation_strip")
                selected_skills.append("comparison_bar")     # v18: ablation win/lose bars
            elif bp_slide.role in ('setup', 'experiment', 'experiments'):
                selected_skills.append("chart")              # v20d: experiment setup has dataset/baseline stats
                selected_skills.append("table")
                selected_skills.append("definition_list")
                selected_skills.append("hero_row")
                selected_skills.append("stacked_sidebar")
                selected_skills.append("stat_dashboard")
            elif bp_slide.role in ('qualitative', 'case_study', 'analysis'):
                selected_skills.append("matrix_layout")      # v18: small-multiple grid
                selected_skills.append("figure_gallery")
                selected_skills.append("evidence_card")
                selected_skills.append("key_finding")

            # Deduplicate while preserving order
            seen = set()
            deduped = []
            for s in selected_skills:
                if s not in seen:
                    seen.add(s)
                    deduped.append(s)
            selected_skills = deduped

            # Load and inject selected skills
            skill_contents = []
            for skill_name in selected_skills:
                spath = skill_dir / f"{skill_name}.html"
                if spath.exists():
                    skill_contents.append(spath.read_text(encoding="utf-8").strip())

            if skill_contents:
                parts.append("\n## Design Skills (general patterns — combine freely, do NOT copy layout verbatim)")
                parts.append("These are CSS/HTML techniques for common components. Use them as building blocks:")
                parts.append("- Adapt colors to match CSS variables from the deck palette")
                parts.append("- You decide the overall layout — these just show how individual components should look")
                parts.append("- Key principle: content sits DIRECTLY ON WHITE. Use typography + spacing + colored left-borders for structure. NO card wrappers (no box-shadow + background + border-radius combos)")
                parts.append("- CHARTS: When numeric data is present, ALWAYS use an SVG bar/line chart as the primary visual. Tables are a fallback for 8+ rows only.")
                for sc in skill_contents:
                    parts.append(f"```html\n{sc}\n```")
        else:
            # Original role-based skill injection for non-imgseed prompts
            skill_role = effective_role if effective_role in ('roadmap', 'outline') else bp_slide.role
            skill_path = skill_dir / f"{skill_role}.html"
            if skill_path.exists():
                skill_content = skill_path.read_text(encoding="utf-8").strip()
                if skill_content:
                    parts.append("\n## HTML Skill Reference (ADAPT to this slide's content — do NOT copy verbatim)")
                    parts.append("Below is a structural pattern for this slide role. Use it as a layout/style reference:")
                    parts.append("- Match the COLOR variables to your palette above")
                    parts.append("- Replace ALL text content with actual evidence")
                    parts.append("- You may adjust proportions, card count, and visual elements")
                    if not has_images and skill_role in ('method', 'results', 'comparison'):
                        parts.append("- ★ DATA VISUALIZATION: This slide has NO paper figure. You SHOULD use an SVG chart or diagram as the dominant left-panel visual. For method: SVG flow diagram with 3-4 connected boxes. For results/comparison: SVG horizontal bar chart showing model scores. See the skill template for SVG structure. A slide with only cards and text looks like a document — add a chart to make it an infographic.")
                    parts.append("- ⚠️ CONTAINMENT: body height=560px max. Max 4 visible containers. Remove content if it doesn't fit — never let it overflow.")
                    parts.append("- ⚠️ TEXT DENSITY: Total body ≤120 words. Each card ≤20 words. Prefer noun phrases over sentences.")
                    parts.append("- ⚠️ NO OVERLAP: Use flex/grid only. Never use position:absolute for content.")
                    parts.append("- ⚠️ CONTRAST: White text on dark/accent backgrounds. Dark text on light backgrounds.")
                    parts.append(f"```html\n{skill_content}\n```")

        return "\n".join(parts), has_images

    def _extract_html(self, response: str) -> str | None:
        """Extract HTML from LLM response."""
        # Look for ```html code blocks
        pattern = r"```html\s*\n(.*?)```"
        matches = re.findall(pattern, response, re.DOTALL)
        if matches:
            for match in matches:
                html = match.strip()
                if "<!DOCTYPE" in html or "<html" in html or "<body" in html:
                    return html

        # Try without code fences - look for <!DOCTYPE or <html
        if "<!DOCTYPE" in response or "<html" in response:
            lines = response.split("\n")
            in_html = False
            html_lines = []
            for line in lines:
                if "<!DOCTYPE" in line or (not in_html and "<html" in line):
                    in_html = True
                if in_html:
                    html_lines.append(line)
                if in_html and "</html>" in line:
                    break
            if html_lines:
                return "\n".join(html_lines)

        return None

    @staticmethod
    def _estimate_text_content(html: str) -> int:
        """Estimate the amount of text content in HTML."""
        # Strip all tags to get raw text
        text = re.sub(r'<[^>]+>', '', html)
        # Strip CSS
        text = re.sub(r'\{[^}]+\}', '', text)
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        return len(text)

    def _detect_layout(self, html: str) -> str:
        """Detect the layout pattern used in generated HTML."""
        html_lower = html.lower()

        n_images = len(re.findall(r'<img\b', html_lower))
        n_tables = len(re.findall(r'<table\b', html_lower))
        has_flex = "display: flex" in html_lower or "display:flex" in html_lower

        if n_images >= 2:
            return "Multi-Image"
        if n_images == 1 and has_flex:
            return "Two-Column-Image"
        if n_images == 1:
            return "Image-Focus"
        if n_tables >= 1:
            return "Table-Layout"

        # Check for metric cards (multiple flex items with large numbers)
        large_numbers = re.findall(r'font-size:\s*(?:4[0-9]|5[0-9])px', html_lower)
        if len(large_numbers) >= 2:
            return "Metric-Cards"

        # Check for columns
        col_count = len(re.findall(r'flex:\s*1', html_lower))
        if col_count >= 3:
            return "Three-Column"
        if col_count >= 2:
            return "Two-Column"

        if "font-style: italic" in html_lower and "quote" in html_lower:
            return "Key-Quote"

        return "Text-Only"

    def _add_fallback_slide(self, slide, bp_slide: BlueprintSlide, evidence: EvidenceState | None = None):
        """Add simple fallback content when HTML generation fails."""
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        title = bp_slide.primary_proposition
        if len(title) > 60:
            title = title[:57] + "..."
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 153)

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        items = bp_slide.must_cover_subset or []
        for i, item in enumerate(items[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.space_before = Pt(8)

    # ------------------------------------------------------------------
    # Image directory utilities (mirrors CodeGenCompiler)
    # ------------------------------------------------------------------
    # Evidence context (shared with repair fallbacks)
    # ------------------------------------------------------------------

    def _build_evidence_context(
        self,
        blueprint_slide,
        evidence,
        case_dir,
        source_store=None,
    ) -> str:
        """Build evidence context string for repair prompts.

        If source_store has a bundle for this slide, uses it directly.
        Otherwise falls back to legacy linked_evidence_ids matching.
        """
        # --- V2: use bundle if available ---
        if source_store is not None and blueprint_slide:
            bundle = source_store.get_bundle(blueprint_slide.slide_id)
            if bundle and bundle.source_text:
                ev_parts = []
                ev_parts.append(bundle.source_text)
                for s in bundle.asset_summaries:
                    ev_parts.append(s)
                for s in bundle.table_summaries:
                    ev_parts.append(s)
                return (
                    "\n\n## Source Evidence (use ONLY data from here, "
                    "do NOT invent numbers)\n\n" + "\n\n".join(ev_parts)
                )

        # --- Legacy fallback ---
        import re as _re
        if not blueprint_slide or not evidence:
            return ""

        ev_parts = []
        search_keywords = set()
        for item in (blueprint_slide.must_cover_subset or []):
            for word in _re.split(r'[\s_]+', item):
                cleaned = _re.sub(r'[^a-zA-Z]', '', word).lower()
                if len(cleaned) > 3:
                    search_keywords.add(cleaned)

        resolved_ids = _resolve_linked_ids(
            blueprint_slide.linked_evidence_ids or [], source_store,
        )
        resolved_set = set(resolved_ids)
        # Build chunk lookup for ordered iteration
        chunk_map = {c.chunk_id: c for c in evidence.chunks}
        matched = 0
        matched_ids = set()
        # Iterate in resolved order (preserves linked_evidence_ids priority)
        for rid in resolved_ids:
            chunk = chunk_map.get(rid)
            if chunk and rid not in matched_ids:
                ev_parts.append(f"[{chunk.chunk_id}] {chunk.content[:2500]}")
                matched_ids.add(rid)
                matched += 1
            if matched >= 8:
                break

        if matched < 8 and search_keywords:
            for chunk in evidence.chunks:
                if chunk.chunk_id in resolved_set or chunk.chunk_id in matched_ids:
                    continue
                section = chunk.metadata.get("section", "").lower()
                content_lower = chunk.content[:300].lower()
                chunk_keywords = set(_re.sub(r'[^a-zA-Z\s]', '', section).lower().split())
                content_keywords = set(_re.sub(r'[^a-zA-Z\s]', '', content_lower).split())
                all_chunk_words = chunk_keywords | content_keywords
                if search_keywords & all_chunk_words and len(chunk.content.strip()) > 50:
                    ev_parts.append(f"[{chunk.chunk_id}] {chunk.content[:2500]}")
                    matched += 1
                if matched >= 8:
                    break

        for tbl in evidence.tables:
            if tbl.table_id in (blueprint_slide.linked_evidence_ids or []):
                tbl_desc = f" — {tbl.description}" if tbl.description else ""
                ev_parts.append(f"[{tbl.table_id}]{tbl_desc}\n{tbl.content[:1500]}")

        if ev_parts:
            return "\n\n## Source Evidence (use ONLY data from here, do NOT invent numbers)\n\n" + "\n\n".join(ev_parts)
        return ""

    # ------------------------------------------------------------------

    def _find_image_dir(self, case_dir: Path) -> Path | None:
        """Find the source_pack directory (parent of figures/, tables/, screenshots/).

        New pyramid structure:
          source_pack/figures/      — extracted figure images
          source_pack/tables/       — table screenshots + JSON sidecars
          source_pack/screenshots/  — page screenshots
          source_pack/rendered_formulas/ — rendered LaTeX formulas

        Falls back to legacy extracted_from_pdf/ for unmigrated cases.
        """
        source_pack = case_dir / "source_pack"
        if source_pack.is_dir():
            if (source_pack / "figures").is_dir() or (source_pack / "tables").is_dir():
                return source_pack
            legacy = source_pack / "extracted_from_pdf"
            if legacy.is_dir():
                return legacy
            return source_pack
        return None

    def _list_images(self, image_dir: Path) -> list[str]:
        """List available image files with relative paths from image_dir.

        Searches pyramid subdirectories: figures/, tables/, screenshots/,
        rendered_formulas/. Also supports legacy subdirs (images/) for
        backward compatibility.
        """
        exts = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif", ".webp"}
        images = []
        for subdir in ["figures", "tables", "screenshots", "rendered_formulas",
                        "images"]:  # "images" for legacy compat
            d = image_dir / subdir
            if d.is_dir():
                for f in sorted(d.iterdir()):
                    if f.is_file() and f.suffix.lower() in exts:
                        images.append(f"{subdir}/{f.name}")
        for f in sorted(image_dir.iterdir()):
            if f.is_file() and f.suffix.lower() in exts:
                images.append(f.name)
        return images
